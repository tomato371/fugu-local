"""モデル呼び出しなしの高速回帰テスト。Ollama 不要・数秒で完走する。
実行: python test_fugu_offline.py
fugu_local / eval_fugu の純粋ロジック（プラン検証・JSON抽出・思考除去・言語判定・
アグリゲータのフォールバック・採点関数）を合成入力で検証する。
"""
import contextlib
import copy
import io
import json
import sys
import types
import urllib.error
import urllib.parse
import urllib.request

import fugu_local as f
import eval_fugu as e
import bench_queue as bq

_FAILS = []


def check(name, cond):
    print(f"[{'OK' if cond else 'NG'}] {name}")
    if not cond:
        _FAILS.append(name)


# ---------- extract_json ----------
check("json: 素のJSON", f.extract_json('{"a": 1}') == {"a": 1})
check("json: コードフェンス", f.extract_json('x\n```json\n{"a": 1}\n```\ny') == {"a": 1})
check("json: think混入", f.extract_json('<think>ignore {"b":9}</think>{"a": 1}') == {"a": 1})
check("json: 地の文に埋没", f.extract_json('The plan is {"a": 1} as follows') == {"a": 1})
check("json: 抽出不能はNone", f.extract_json("no json here") is None)
check("json: 空はNone", f.extract_json("") is None)
# 2026-07-22: 貪欲re.searchの over-capture 回帰防止（先頭オブジェクト消失バグ）
check("json: 末尾に余分な波括弧があっても先頭を抽出",
      f.extract_json('Sure! {"mode":"single"} note {x}') == {"mode": "single"})
check("json: 先行する集合記法{1,2,3}に惑わされない",
      f.extract_json('The set {1,2,3} then {"a": 1}') == {"a": 1})
check("json: 2つ目の有効オブジェクトがあっても最初を返す",
      f.extract_json('{"a": 1} and also {"b": 2}') == {"a": 1})
check("json: 文字列値中の}に惑わされない(深さカウントの文字列認識)",
      f.extract_json('x {"s": "a}b", "n": 2} y') == {"s": "a}b", "n": 2})
check("json: 閉じない{単体はクラッシュせずNone",
      f.extract_json("prefix { unbalanced no closing brace") is None)

# ---------- extract_json: dict-or-None契約の強制 (2026-07-25) ----------
# 背景: ステップ1 `return json.loads(text)` は json.loads が成功しさえすれば
# 型を問わず結果を素通しで返していた。docstring は「最初の JSON オブジェクト」
# (=dict)を約束し、ステップ2/3は実際に dict-or-None 契約を守っているのに、
# ステップ1だけはトップレベルが妥当なJSONでありさえすれば list/int/float/bool/str
# も丸ごと返してしまっていた。呼び出し側の一部(_critic_judge/second_opinion/
# _sd_prompt_from_request/plan_pptx_images)は `extract_json(raw) or {}` の後に
# 無条件で `.get(...)` するため、モデルが `[{"ok": true}]` や `true`/`42`/`"text"`
# のようなトップレベル非objectな(しかし妥当な)JSONを出力すると、truthyな
# 非dict値がそのまま通り抜けて `.get` で AttributeError を送出していた
# (iteration 103/111/112/113 と同種の「妥当だが型が想定と違うモデル出力」対策)。
# ここでは extract_json 自体が常に dict か None のどちらかしか返さないことを
# 直接ロックする。
for _ej_bad in ('[1,2,3]', '42', '3.14', 'true', 'false', 'null', '"hello"'):
    check(f"json: トップレベル非object({_ej_bad!r})はNone",
          f.extract_json(_ej_bad) is None)
check("json: ネストしたobjectは3)の波括弧スキャナで回収できる([{...}]->中身のdict)",
      f.extract_json('[{"ok": true}]') == {"ok": True})
for _ej_case in ('{"a": 1}', '[1,2,3]', '42', '3.14', 'true', 'false', 'null',
                 '"hello"', '[{"ok": true}]', 'no json here', ''):
    _ej_r = f.extract_json(_ej_case)
    check(f"json: 戻り値は常にdictかNone(入力={_ej_case!r})",
          _ej_r is None or isinstance(_ej_r, dict))

# ---------- strip_think ----------
check("strip: think除去", f.strip_think("<think>x</think>answer") == "answer")
check("strip: thinking除去", f.strip_think("<THINKING>x</THINKING>ans") == "ans")
check("strip: 対象なしは素通し", f.strip_think("plain") == "plain")
check("strip: None耐性", f.strip_think(None) is None)
# 2026-07-22: num_predict 打ち切りで閉じタグの無い '<think>...' が丸ごと
# 「回答」として漏れる既知の失敗モード（gotcha #2 / #7）の回帰防止。
check("strip: 閉じタグ無しのthinkは末尾まで丸ごと除去",
      f.strip_think("<think>Let me compute... maybe 17... no, 42").strip() == "")
check("strip: 開始タグより前のテキストは保持",
      f.strip_think("answer is 5 <think>double-checking then cut off") == "answer is 5")
check("strip: 閉じタグ無しのTHINKINGも大小文字問わず除去",
      f.strip_think("<THINKING>still going with no closer").strip() == "")
check("strip: 対応の取れた既存ペアは従来通り除去(回帰)",
      f.strip_think("<think>x</think>answer") == "answer")
check("strip: 孤立した</think>閉じタグ単体は無視される",
      f.strip_think("no opener here </think> tail") == "no opener here </think> tail")
check("final_answer: 打ち切りthinkの中間値を誤って投票しない(E2E)",
      f.extract_final_answer(
          "<think>...intermediate value 17 then output cut off", "math") is None)

# ---------- validate_plan（新スキーマ: mode single|moa / selected_proposers 他） ----------
f.PROPOSERS = ["qwen3:4b", "phi4-mini", "gemma4:e2b-it-qat"]
f.AGGREGATOR = "deepseek-r1:7b"
f.CONDUCTOR = "qwen3:4b"

p = f.validate_plan({"mode": "moa",
                     "selected_proposers": ["qwen3:4b", "存在しないモデル"],
                     "rounds": 99, "use_image_generation": False,
                     "search_required": True})
check("plan: rounds を MAX_ROUNDS に丸める", p["rounds"] == f.MAX_ROUNDS)
check("plan: 未知プロポーザーを除外", p["selected_proposers"] == ["qwen3:4b"])
check("plan: search_required を反映", p["search_required"] is True)
check("plan: 不正mode は moa",
      f.validate_plan({"mode": "weird", "selected_proposers": ["qwen3:4b"]})["mode"] == "moa")
check("plan: dict以外はフォールバック", f.validate_plan(None).get("_fallback") is True)
check("plan: single は先頭1体のみ",
      f.validate_plan({"mode": "single",
                       "selected_proposers": ["phi4-mini", "qwen3:4b"]})["selected_proposers"]
      == ["phi4-mini"])
check("plan: 空selected は既定へフォールバック",
      len(f.validate_plan({"mode": "moa", "selected_proposers": []})["selected_proposers"]) >= 1)
check("plan: use_image_generation は mode を強制しない(非排他)",
      f.validate_plan({"mode": "moa", "selected_proposers": ["qwen3:4b"],
                       "use_image_generation": True})["mode"] == "moa")
check("plan: 画像生成フラグを反映",
      f.validate_plan({"use_image_generation": True,
                       "selected_proposers": ["qwen3:4b"]})["use_image_generation"] is True)
check("plan: image_only を反映(use_image_generation と整合時)",
      f.validate_plan({"image_only": True, "use_image_generation": True,
                       "selected_proposers": []})["image_only"] is True)
check("plan: make_pptx を反映し image_only を無効化",
      (lambda p: p["make_pptx"] is True and p["image_only"] is False)(
          f.validate_plan({"make_pptx": True, "image_only": True,
                           "selected_proposers": ["qwen3:4b"]})))
check("plan: use_image_generation=False なら image_only を強制的に無効化(矛盾解消)",
      f.validate_plan({"image_only": True, "use_image_generation": False,
                       "selected_proposers": []})["image_only"] is False)
check("plan: 矛盾したmath plan で SC投票ゲートの3フラグが全てFalseになる",
      (lambda p: p["image_only"] is False and p["make_pptx"] is False
       and p["use_image_generation"] is False)(
          f.validate_plan({"task_type": "math", "image_only": True,
                           "use_image_generation": False,
                           "selected_proposers": []})))
check("plan: 矛盾解消後は selected_proposers が通常moa分岐(image panelでない)",
      len(f.validate_plan({"task_type": "math", "image_only": True,
                           "use_image_generation": False,
                           "selected_proposers": []})["selected_proposers"]) >= 2)

# ---------- ペルソナ解決（selected_proposers のペルソナ名→実モデル） ----------
_op_persona = f.PROPOSERS
f.PROPOSERS = ["gpt-oss:20b", "qwen3-coder:30b", "gemma4:26b", "qwen3.6:35b"]
check("persona: 'Proposer A' → gpt-oss", f._resolve_proposer("Proposer A") == "gpt-oss:20b")
check("persona: 緩い 'a' → gpt-oss", f._resolve_proposer("a") == "gpt-oss:20b")
check("persona: モデル名直指定を許容", f._resolve_proposer("qwen3.6:35b") == "qwen3.6:35b")
check("persona: 未知は None", f._resolve_proposer("Proposer Z") is None)
check("persona: validate がペルソナ名を実モデルへ解決",
      f.validate_plan({"mode": "moa",
                       "selected_proposers": ["Proposer C", "Proposer D"]})["selected_proposers"]
      == ["gemma4:26b", "qwen3.6:35b"])
f.PROPOSERS = ["gpt-oss:20b", "qwen3-coder:30b", "phi4"]  # gemma4:26b 未導入シナリオ
check("persona: 未導入モデルのペルソナは None", f._resolve_proposer("Proposer C") is None)
check("persona: validate は未導入ペルソナを除外",
      f.validate_plan({"mode": "moa",
                       "selected_proposers": ["Proposer C", "Proposer B"]})["selected_proposers"]
      == ["qwen3-coder:30b"])

# ---------- _resolve_proposer: 非文字列/unhashable要素での例外耐性 (2026-07-25) ----------
# CONDUCTOR_SCHEMA の items:{type:string} 強制は完全ではない（本ファイル各所が既知の
# 通り「スキーマ強制でも稀に JSON が崩れる」）。selected_proposers に list/dict の
# ような非文字列・非ハッシュ可能な要素が紛れ込むと、旧実装は _resolve_proposer 冒頭の
# `name in PERSONA_MODELS`（dict メンバーシップテスト）で TypeError: unhashable type
# を送出し、唯一の呼び出し元 validate_plan、さらに conduct -> ask_fugu/fugu_answer まで
# 無捕捉で伝播してターン全体をクラッシュさせ、計算済みの回答を失っていた。
# iter 103 (_ddg_instant の非list RelatedTopics)、iter 111 (plan_pptx_images の非list
# images)、iter 112 (research_search の非list queries)、iter 113 (_read_ipynb の
# 非dict/非list cells) と同じ「壊れたスキーマ制約付きプランは例外を出さず既定値へ
# フォールバックさせる」作法での回帰防止。
_op_persona2 = f.PROPOSERS
try:
    f.PROPOSERS = ["gpt-oss:20b", "qwen3-coder:30b", "gemma4:26b", "qwen3.6:35b"]

    # --- 非文字列/unhashable要素は例外を出さず None (修正前は [] / {} で TypeError) ---
    check("resolve: 空listはNone(旧: unhashable TypeError)", f._resolve_proposer([]) is None)
    check("resolve: 空dictはNone(旧: unhashable TypeError)", f._resolve_proposer({}) is None)
    check("resolve: Noneは例外なくNone", f._resolve_proposer(None) is None)
    check("resolve: intはNone", f._resolve_proposer(5) is None)
    check("resolve: 非空listもNone(unhashable)", f._resolve_proposer(["nested"]) is None)
    check("resolve: 非空dictもNone(unhashable)", f._resolve_proposer({"a": 1}) is None)

    # --- 有効な文字列入力の回帰: 修正前と完全に同一の解決結果であること ---
    check("resolve回帰: 'Proposer A'(完全表記)",
          f._resolve_proposer("Proposer A") == "gpt-oss:20b")
    check("resolve回帰: 'A'(短縮・大文字)", f._resolve_proposer("A") == "gpt-oss:20b")
    check("resolve回帰: 'proposer a'(小文字・完全表記)",
          f._resolve_proposer("proposer a") == "gpt-oss:20b")
    check("resolve回帰: 実モデル名を直接指定",
          f._resolve_proposer("qwen3.6:35b") == "qwen3.6:35b")
    check("resolve回帰: 未知の文字列はNone", f._resolve_proposer("Proposer Z") is None)

    # --- validate_plan: 不正要素と有効な文字列が混在しても例外を出さず有効分のみ解決 ---
    p_mixed = f.validate_plan({"mode": "moa",
                               "selected_proposers": [[], {}, "Proposer A"]})
    check("validate: 不正要素混在でも例外なし・有効な1件のみ解決",
          p_mixed["selected_proposers"] == ["gpt-oss:20b"])

    # --- validate_plan: 全要素が不正なら例外を出さず既定値へフォールバック ---
    p_all_bad_moa = f.validate_plan({"mode": "moa",
                                     "selected_proposers": [[], {}, 123]})
    check("validate: 全要素不正(moa)でも例外なし・PROPOSERS[:3]へフォールバック",
          p_all_bad_moa["selected_proposers"] == f.PROPOSERS[:3])
    p_all_bad_single = f.validate_plan({"mode": "single",
                                        "selected_proposers": [[], {}, 123]})
    check("validate: 全要素不正(single)でも例外なし・PROPOSERS[:1]へフォールバック",
          p_all_bad_single["selected_proposers"] == f.PROPOSERS[:1])
    p_all_bad_image = f.validate_plan({"mode": "moa", "image_only": True,
                                       "use_image_generation": True,
                                       "selected_proposers": [[], {}, 123]})
    check("validate: 全要素不正(image_only)でも例外なし・空panelへ"
          "(呼び出し側の 'or PROPOSERS[:...]' が下流で担保)",
          p_all_bad_image["selected_proposers"] == [])
finally:
    f.PROPOSERS = _op_persona2

# ---------- proposer_sys_for / build_proposer_desc / _persona_str:
#            プロポーザーpersona配線 (offline coverage, 2026-07-26) ----------
# 3つとも純粋関数で、モジュールグローバル(PERSONA_MODELS/PERSONA_IDENTITY/
# PROPOSER_PROFILES/PROPOSERS/MODEL_TO_PERSONA/PROPOSER_SYS)のみを読む。
# proposer_sys_for は毎回のproposer呼び出し(get_single_proposal/単体モード
# fugu_answer)を通り、人格プロンプトの前置有無を左右する(精度ガードレール#7
# が拾うSC投票の入力に効く)。build_proposer_desc は Conductor プロンプトに
# 埋め込まれる導入済みpersona一覧を、_persona_str はプラン表示ラベルを作る。
# いずれもこれまでオフラインテストのカバレッジが皆無だったため特性化する
# (fugu_local.py は変更しない。テスト専用の追加)。
_op_persona_models = f.PERSONA_MODELS
_op_persona_identity = f.PERSONA_IDENTITY
_op_proposer_profiles = f.PROPOSER_PROFILES
_op_proposers_seam = f.PROPOSERS
_op_model_to_persona = f.MODEL_TO_PERSONA
_op_proposer_sys = f.PROPOSER_SYS
try:
    # --- proposer_sys_for: identityの前置/非前置/falsy(空文字)分岐 ---
    f.PERSONA_IDENTITY = {
        "model-with-id": "あなたはテスト用人格Aです。",
        "model-empty-id": "",
    }
    check("proposer_sys_for: identityありはPROPOSER_SYSの前に単一改行'\\n'で前置",
          f.proposer_sys_for("model-with-id")
          == "あなたはテスト用人格Aです。\n" + f.PROPOSER_SYS)
    check("proposer_sys_for: identityありの結果はPROPOSER_SYSを末尾に含む(接尾辞)",
          f.proposer_sys_for("model-with-id").endswith(f.PROPOSER_SYS))
    check("proposer_sys_for: identityなし(キー不在)はPROPOSER_SYSそのまま(前置なし)",
          f.proposer_sys_for("model-without-id") == f.PROPOSER_SYS)
    check("proposer_sys_for: 空文字identity(falsy)もPROPOSER_SYSそのまま",
          f.proposer_sys_for("model-empty-id") == f.PROPOSER_SYS)

    # --- build_proposer_desc: PERSONA_MODELS挿入順・導入済みのみ列挙・
    #     PROPOSER_PROFILES既定'汎用'フォールバック ---
    f.PERSONA_MODELS = {
        "Proposer A": "model-a",
        "Proposer B": "model-b",
        "Proposer C": "model-c",
    }
    f.MODEL_TO_PERSONA = {v: k for k, v in f.PERSONA_MODELS.items()}
    f.PROPOSER_PROFILES = {"model-a": "得意分野Aの説明"}  # model-b/model-cは未設定
    # 導入順(PROPOSERSの並び)はPERSONA_MODELS順とわざと変えて、出力順が
    # PERSONA_MODELS の挿入順に従うこと(PROPOSERSの並び順に引きずられない
    # こと)を検証する。Proposer Cは未導入のまま(除外を確認)。
    f.PROPOSERS = ["model-b", "model-a"]
    _pdesc = f.build_proposer_desc()
    check("build_proposer_desc: 導入済みのみをPERSONA_MODELS挿入順・"
          "'- {label} ({model}): {profile}'形式で改行結合",
          _pdesc == "- Proposer A (model-a): 得意分野Aの説明\n"
                    "- Proposer B (model-b): 汎用")
    check("build_proposer_desc: 未導入(Proposer C/model-c)は出力に含まれない",
          "Proposer C" not in _pdesc and "model-c" not in _pdesc)
    check("build_proposer_desc: PROPOSER_PROFILES未設定モデルは既定'汎用'",
          "(model-b): 汎用" in _pdesc)
    check("build_proposer_desc: PROPOSER_PROFILES設定済みモデルはその値を使う",
          "(model-a): 得意分野Aの説明" in _pdesc)

    f.PROPOSERS = []
    check("build_proposer_desc: 導入済みpersonaモデルが1つも無ければ空文字列",
          f.build_proposer_desc() == "")

    # --- _persona_str: 派生辞書MODEL_TO_PERSONA(PERSONA_MODELSではない)
    #     ベースの整形/フォールバック。PERSONA_MODELSを差し替える際は
    #     MODEL_TO_PERSONAも整合させて再構築する(でないと古い/不整合な
    #     マッピングを特性化してしまう) ---
    f.PERSONA_MODELS = {"Proposer X": "model-x"}
    f.MODEL_TO_PERSONA = {v: k for k, v in f.PERSONA_MODELS.items()}
    check("_persona_str: MODEL_TO_PERSONAにあれば'persona (model)'形式",
          f._persona_str("model-x") == "Proposer X (model-x)")
    check("_persona_str: MODEL_TO_PERSONAに無ければモデル名そのまま(str()化)",
          f._persona_str("unknown-model") == "unknown-model")
    check("_persona_str: 非文字列(int)モデルもstr()化されて返る",
          f._persona_str(123) == "123")
finally:
    f.PERSONA_MODELS = _op_persona_models
    f.PERSONA_IDENTITY = _op_persona_identity
    f.PROPOSER_PROFILES = _op_proposer_profiles
    f.PROPOSERS = _op_proposers_seam
    f.MODEL_TO_PERSONA = _op_model_to_persona
    f.PROPOSER_SYS = _op_proposer_sys

check("proposer persona配線: テスト後にPERSONA_MODELSが復元されている",
      f.PERSONA_MODELS == _op_persona_models)
check("proposer persona配線: テスト後にPERSONA_IDENTITYが復元されている",
      f.PERSONA_IDENTITY == _op_persona_identity)
check("proposer persona配線: テスト後にPROPOSER_PROFILESが復元されている",
      f.PROPOSER_PROFILES == _op_proposer_profiles)
check("proposer persona配線: テスト後にPROPOSERSが復元されている",
      f.PROPOSERS == _op_proposers_seam)
check("proposer persona配線: テスト後にMODEL_TO_PERSONAが復元されている",
      f.MODEL_TO_PERSONA == _op_model_to_persona)
check("proposer persona配線: テスト後にPROPOSER_SYSが復元されている",
      f.PROPOSER_SYS == _op_proposer_sys)

# ---------- 精度ガードレール（code/proof を single→moa へ格上げ） ----------
f.PROPOSERS = ["gpt-oss:20b", "qwen3-coder:30b", "gemma4:26b", "phi4"]


def _single_plan():
    return {"mode": "single", "selected_proposers": ["gpt-oss:20b"], "rounds": 1,
            "use_image_generation": False, "search_required": False,
            "reason": "r", "_fallback": False}


check("guard: コード質問は moa へ格上げ",
      f._apply_accuracy_guardrails("Pythonで実装して", _single_plan())["mode"] == "moa")
check("guard: 証明質問は moa へ格上げ",
      f._apply_accuracy_guardrails("背理法で証明せよ", _single_plan())["mode"] == "moa")
check("guard: 格上げ時は複数体を割当",
      len(f._apply_accuracy_guardrails("コードを書いて", _single_plan())["selected_proposers"]) >= 2)
check("guard: 平易な質問は single のまま",
      f._apply_accuracy_guardrails("日本の首都は？", _single_plan())["mode"] == "single")
check("guard: image_only は格上げ対象外",
      f._apply_accuracy_guardrails(
          "コードを実装して",
          {"mode": "single", "selected_proposers": [],
           "image_only": True})["mode"] == "single")
check("guard: イラスト付き(image_only=False)のコードは格上げ",
      f._apply_accuracy_guardrails(
          "コードを実装して",
          {"mode": "single", "selected_proposers": ["gpt-oss:20b"],
           "use_image_generation": True, "image_only": False})["mode"] == "moa")

# ---------- スライド分解（PowerPoint 用） ----------
_slides = f._parse_slides("## 概要\n- 要点1\n- 要点2\n\n## 詳細\n本文の段落です。\n1. 手順A\n2. 手順B")
check("pptx: 見出しでスライド分割", len(_slides) == 2)
check("pptx: 箇条書き記号を除去", _slides[0]["bullets"] == ["要点1", "要点2"])
check("pptx: タイトルは見出し由来", _slides[1]["title"] == "詳細")
check("pptx: 見出し無しは概要1枚",
      len(f._parse_slides("ただの文章その1。\nその2。")) == 1)

# 2026-07-23: 未対応(奇数個)の```フェンスが in_code を文書末尾まで固定し、
# 以降の見出しが新規スライドを作れず literal な "## ..." 箇条書きとして
# 1枚に押し込まれてしまう回帰の防止（extract_boxed iter11 / strip_think
# iter16 / _save_as_html iter37 と同種のバグクラス）。

# 回帰: 対になったコードブロック内の '# コメント' は見出しに昇格せず、
# コード本文も強調記号除去されない（従来通り）。
_bal_code = f._parse_slides(
    "## セクション1\n```python\n# これはコメント\nx = 1\n```\n## セクション2\n本文")
check("pptx回帰: 対になったコードフェンス内の#コメントは見出しに昇格しない",
      len(_bal_code) == 2
      and _bal_code[0]["title"] == "セクション1"
      and "# これはコメント" in _bal_code[0]["bullets"]
      and "x = 1" in _bal_code[0]["bullets"]
      and _bal_code[1]["title"] == "セクション2")

# 回帰: 対になったコードブロック2個(フェンス4本)+見出し混在は従来通り分割される。
_bal_two = f._parse_slides(
    "## A\n```python\n# a\n```\n## B\n```python\n# b\n```\n## C\n末尾の文")
check("pptx回帰: 対になったコードブロック2個でも見出し分割は不変",
      [s["title"] for s in _bal_two] == ["A", "B", "C"]
      and "# a" in _bal_two[0]["bullets"] and "# b" in _bal_two[1]["bullets"])

# 新規: 閉じ忘れの奇数フェンスの後に来る見出しは、literalな"## ..."箇条書きに
# ならず、新規スライドの見出しとして正しく分割される。
# (フェンス直後の本文行はあえて '#' で始めない。中断されたコードの地の文が
#  たまたま '#' で始まれば新仕様では見出し候補になるのは意図通りだが、
#  ここでは「フェンス後も後続の本当の見出しでちゃんと分割される」ことだけを
#  単純な形で確認する。)
_unterminated = f._parse_slides(
    "## 導入\n本文1\n```python\nx = 1\ny = 2\n## 結果\n本文2\n## まとめ\n本文3")
check("pptx新規: 閉じ忘れフェンス後の見出しは新規スライドになる",
      [s["title"] for s in _unterminated] == ["導入", "結果", "まとめ"])
check("pptx新規: 閉じ忘れフェンス後の見出しは#が除去されタイトルになる(literal箇条書きではない)",
      all("#" not in s["title"] for s in _unterminated)
      and not any(b.startswith("##") for s in _unterminated for b in s["bullets"]))
check("pptx新規: 閉じ忘れフェンスでも本文は失われない",
      "本文1" in _unterminated[0]["bullets"]
      and "本文2" in _unterminated[1]["bullets"]
      and "本文3" in _unterminated[2]["bullets"])

# 新規: 混在ケース(先に対になったブロックがあり、末尾だけ閉じ忘れ)でも、
# 先行する対になったブロックはコードモードを維持し#コメントを昇格させない。
_mixed = f._parse_slides(
    "## 前半\n```python\n# 対になっているコメント\n```\n## 後半\n```python\nx = 1\n## 次\n本文")
check("pptx混在: 先行する対になったブロックは見出し昇格させない",
      _mixed[0]["title"] == "前半" and "# 対になっているコメント" in _mixed[0]["bullets"])
check("pptx混在: 末尾の閉じ忘れフェンス後の見出しは新規スライドになる",
      [s["title"] for s in _mixed] == ["前半", "後半", "次"])

# 2026-07-23: 旧実装 re.sub(r"[*_`#]+", "", t) が * _ ` # を無差別に全削除して
# いた過剰マッチの回帰防止（詳細は fugu_local.py 内の該当コメント参照）。
# 数式の演算子/指数、識別子中のアンダースコア、インラインコード中の識別子が
# 破壊されず、対のバッククォートのみが除去されることを確認する。
_slides_math = f._parse_slides("## 式\n- E = m*c_0**2")
check("slides: 数式の演算子/指数はそのまま残る(mc02に化けない)",
      any("m*c_0**2" in b for b in _slides_math[0]["bullets"]))

_slides_code_id = f._parse_slides("- 呼び出しは `do_thing()` です")
check("slides: インラインコード識別子はアンダースコア込みで残る(dothingにならない)",
      any("do_thing()" in b for b in _slides_code_id[0]["bullets"]))
check("slides: 対になったバッククォート自体は除去される",
      all("`" not in b for b in _slides_code_id[0]["bullets"]))

_slides_uscore = f._parse_slides("- 係数は a_1 + a_2")
check("slides: 地の文中のアンダースコア識別子はそのまま残る(a_1/a_2がa1/a2にならない)",
      any("a_1" in b and "a_2" in b for b in _slides_uscore[0]["bullets"]))

check("slides: 見出し中のアンダースコアはそのまま残る(a_bがabにならない)",
      "a_b" in f._parse_slides("## a_b の説明")[0]["title"])

# 安全な方向(曖昧なら削除しない)の確認: 対になっていない単独のバッククォートや
# アスタリスクは装飾か演算子/識別子か判別できないため、削除せずそのまま残す。
_slides_stray = f._parse_slides("- value: `x と **重要**")
check("slides: 対になっていない単独バッククォートは残す(安全側)",
      any("`x" in b for b in _slides_stray[0]["bullets"]))
check("slides: 太字装飾の**もそのまま残る(削除より保持を優先)",
      any("**重要**" in b for b in _slides_stray[0]["bullets"]))

check("pptx: deck_title は短い質問を採用", f._deck_title("犬の紹介", _slides) == "犬の紹介")
# 2026-07-22: 空白のみの質問は if question で truthy のまま素通りし、
# strip()後に splitlines() が [] を返して [0] が IndexError になっていた回帰。
check("pptx: deck_title 空白のみ質問+スライド無しは既定値",
      f._deck_title("   ", []) == "プレゼンテーション")
check("pptx: deck_title 空白のみ質問はスライドタイトルへフォールバック",
      f._deck_title("\n\n", [{"title": "概要", "bullets": []}]) == "概要")
check("pptx: deck_title 空白のみ質問+無題スライドは既定値",
      f._deck_title("\t \n", [{"title": "", "bullets": []}]) == "プレゼンテーション")
check("pptx: deck_title 複数行質問は先頭行を採用",
      f._deck_title("\n  タイトル行\n本文\n", []) == "タイトル行")
check("pptx: deck_title 40字超はスライドタイトルへフォールバック",
      f._deck_title("あ" * 41, [{"title": "見出し", "bullets": []}]) == "見出し")
check("pptx: deck_title 空文字はスライドタイトルへフォールバック",
      f._deck_title("", [{"title": "見出し", "bullets": []}]) == "見出し")
check("pptx: deck_title None質問は既定値",
      f._deck_title(None, []) == "プレゼンテーション")

# ---------- 出力形態ルーティングガードレール ----------
f.PROPOSERS = ["gpt-oss:20b", "qwen3-coder:30b", "gemma4:26b", "phi4"]


def _base_plan():
    return {"mode": "single", "selected_proposers": ["gpt-oss:20b"], "rounds": 1,
            "use_image_generation": False, "image_only": False, "make_pptx": False,
            "search_required": False, "reason": "r", "_fallback": False}


_r = f._apply_routing_guardrails("機械学習入門のスライドを作って", _base_plan())
check("route: スライド→make_pptx+moa", _r["make_pptx"] is True and _r["mode"] == "moa")
_r = f._apply_routing_guardrails("かわいい柴犬のイラストを描いて", _base_plan())
check("route: イラストのみ→image_only", _r["use_image_generation"] is True and _r["image_only"] is True)
_r = f._apply_routing_guardrails("PINN洪水モデルを説明して図も作って", _base_plan())
check("route: 説明+図→イラスト付き(image_only=False)",
      _r["use_image_generation"] is True and _r["image_only"] is False)
_r = f._apply_routing_guardrails("日本の首都は？", _base_plan())
check("route: 通常質問は据え置き",
      _r["make_pptx"] is False and _r["use_image_generation"] is False)

# 2026-07-24: validate_plan(iteration 40) が確立した make_pptx=True⇒image_only=False
# の不変条件を、conduct() 内で validate_plan の「後」に呼ばれる
# _apply_routing_guardrails の画像分岐が再び壊さないことの回帰テスト。
# 例: コンダクタが '発表資料'（_PPTX_SIGNALS 非一致のPPTX同義語）から
# make_pptx=True を立てた plan に、_IMAGE_SIGNALS には一致するが
# _TEXT_TASK_SIGNALS には一致しない質問が来た場合。
_pptx_already_true_plan = _base_plan()
_pptx_already_true_plan["make_pptx"] = True
_r = f._apply_routing_guardrails("かわいい柴犬のイラスト付き発表資料を作成して", _pptx_already_true_plan)
check("route: make_pptx=True済みプランは画像分岐でimage_only=Trueにならない(不変条件維持)",
      _r["make_pptx"] is True and _r["image_only"] is False)
check("route: make_pptx かつ image_only が同時にTrueにはならない",
      not (_r["make_pptx"] and _r["image_only"]))
f.PROPOSERS = _op_persona

# ---------- Office添付ガードレール (_apply_office_guardrail, 2026-07-26追加) ----------
# CONDUCTOR_SYS 特殊ルーティング指示3 と conduct() の自然文ヒント（下の conduct() 節にある
# office_attached ヒント文言テスト参照）はプロンプトのみの安全網であり、他のガードレール
# (_apply_routing_guardrails/_apply_accuracy_guardrails/_apply_tasktype_guardrails)と同様に
# 小型 Conductor がヒントを無視しても結果が変わらないことを、決定的ガードレール単体で
# (f.PROPOSERS を mutate/restore する try/finally の中で) 検証する。
# PERSONA_MODELS['Proposer C'] 経由でモデル名を参照し、'gemma4:26b' を直接ハードコードしない。
_office_c_model = f.PERSONA_MODELS["Proposer C"]
_orig_props_office = f.PROPOSERS


def _office_plan(**overrides):
    plan = {"mode": "single", "task_type": "knowledge",
            "selected_proposers": ["some-other-model"], "rounds": 1,
            "use_image_generation": False, "image_only": False, "make_pptx": False,
            "search_required": False, "reason": "r", "_fallback": False}
    plan.update(overrides)
    return plan


try:
    f.PROPOSERS = ["some-other-model", _office_c_model, "another-model"]

    # (1) office_attached=False は完全な no-op(呼び出し前のdeep copyとbyte-for-byte一致)
    _p_off = _office_plan()
    _p_off_snapshot = copy.deepcopy(_p_off)
    _r_off = f._apply_office_guardrail(_p_off, False)
    check("office guard: office_attached=Falseはplanを一切変更しない(no-op)",
          _r_off == _p_off_snapshot and _p_off == _p_off_snapshot)

    # (2) office_attached=True・Proposer C導入済み・mode='single'でCが未選択
    #     -> moaへ格上げしProposer Cのモデルを注入する
    _p_on = _office_plan(mode="single", selected_proposers=["some-other-model"])
    _r_on = f._apply_office_guardrail(_p_on, True)
    check("office guard: office_attached=Trueでmode='moa'へ強制",
          _r_on["mode"] == "moa")
    check("office guard: 導入済みならProposer Cのモデルがselected_proposersに含まれる",
          _office_c_model in _r_on["selected_proposers"])

    # 冪等性: 同じplanに2回目を適用しても結果(mode/selected_proposers/reason)が変わらない
    _r_on_2 = f._apply_office_guardrail(_r_on, True)
    check("office guard: 2回適用しても結果は不変(冪等)",
          _r_on_2["mode"] == "moa"
          and _r_on_2["selected_proposers"] == _r_on["selected_proposers"]
          and _r_on_2["reason"] == _r_on["reason"])

    # (3) office_attached=True・Proposer Cのモデルが未導入(PROPOSERSに存在しない)
    #     -> mode='moa'へは強制するが例外を出さず、Cは注入しない(未導入モデルへは絶対に
    #     ルーティングしない = PPTX/精度ガードレールと同じ PROPOSERS 参照パターン)
    f.PROPOSERS = ["some-other-model", "another-model"]  # Proposer Cのモデルは未導入
    _p_noc = _office_plan(mode="single", selected_proposers=["some-other-model"])
    _r_noc = f._apply_office_guardrail(_p_noc, True)
    check("office guard: Proposer C未導入でも例外を出さずmode='moa'へ強制",
          _r_noc["mode"] == "moa")
    check("office guard: Proposer C未導入ならselected_proposersは非空のまま・Cは注入されない",
          bool(_r_noc["selected_proposers"]) and _office_c_model not in _r_noc["selected_proposers"])
    f.PROPOSERS = ["some-other-model", _office_c_model, "another-model"]

    # (4) image_only=True は no-op(mode/image_only/selected_proposersとも不変。
    #     _apply_accuracy_guardrails の image_only 早期returnと同じ理由:
    #     画像のみの回答にテキストパネルを割り当てても無意味なため)
    _p_img = _office_plan(mode="single", image_only=True,
                           selected_proposers=["some-other-model"])
    _p_img_snapshot = copy.deepcopy(_p_img)
    _r_img = f._apply_office_guardrail(_p_img, True)
    check("office guard: image_only=Trueはno-op(mode/image_only/selected_proposers不変)",
          _r_img["mode"] == _p_img_snapshot["mode"]
          and _r_img["image_only"] == _p_img_snapshot["image_only"]
          and _r_img["selected_proposers"] == _p_img_snapshot["selected_proposers"])

    # (5) selected_proposersが既に4件(Proposer C含まず)・Cは導入済み
    #     -> Proposer Cが主軸(先頭)に立ち、4件上限を超えず・重複無し・2回適用しても
    #     同じリスト(冪等)。validate_plan の models[:4] 上限と揃える。
    f.PROPOSERS = ["m1", "m2", "m3", "m4", _office_c_model]
    _p_full = _office_plan(mode="moa", selected_proposers=["m1", "m2", "m3", "m4"])
    _r_full = f._apply_office_guardrail(_p_full, True)
    check("office guard: 4件満杯でもProposer Cが主軸(先頭)に立つ",
          _r_full["selected_proposers"][0] == _office_c_model)
    check("office guard: 4件上限を超えない", len(_r_full["selected_proposers"]) <= 4)
    check("office guard: 重複無し",
          len(_r_full["selected_proposers"]) == len(set(_r_full["selected_proposers"])))
    _r_full_2 = f._apply_office_guardrail(_r_full, True)
    check("office guard: 4件満杯ケースも2回適用で同じリスト(冪等)",
          _r_full_2["selected_proposers"] == _r_full["selected_proposers"])
finally:
    f.PROPOSERS = _orig_props_office
check("office guard: テスト後にPROPOSERSが元へ復元されている",
      f.PROPOSERS == _orig_props_office)

# ---------- conduct(): プランオーケストレーション ----------
# conduct() は f.ask のみをモックし、extract_json/validate_plan/_apply_*_guardrails/
# build_proposer_desc は実物をそのまま通す（本物の合成ロジックを検証するため）。
# PERSONA_MODELS のキー(Proposer A〜D)は fugu_local.py の定義に合わせた実在値。
_orig_ask_conduct = f.ask
_orig_props_conduct = f.PROPOSERS
f.PROPOSERS = ["gpt-oss:20b", "qwen3-coder:30b", "gemma4:26b", "qwen3.6:35b"]


def _valid_plan_json(**overrides):
    plan = {
        "mode": "moa",
        "task_type": "knowledge",
        "selected_proposers": ["Proposer A", "Proposer B"],
        "rounds": 1,
        "use_image_generation": False,
        "image_only": False,
        "make_pptx": False,
        "search_required": False,
        "reason": "test plan",
    }
    plan.update(overrides)
    return json.dumps(plan)


def _make_conduct_ask(responses):
    """呼び出し毎に responses を順に返す f.ask 互換フェイク。
    どのモデルの ask() 呼び出しでも messages と label を記録する。"""
    calls = []

    def _fake(model, messages, temperature, **kwargs):
        calls.append({"messages": messages, "label": kwargs.get("label")})
        idx = min(len(calls) - 1, len(responses) - 1)
        return responses[idx]

    return _fake, calls


# --- 正常系: 1回目でJSONが取れればask()は1回だけ・引き直し無し ---
try:
    _resp_happy = _valid_plan_json(task_type="chat", reason="happy path")
    f.ask, _calls_happy = _make_conduct_ask([_resp_happy])
    _plan_happy, _raw_happy = f.conduct("こんにちは")
finally:
    f.ask = _orig_ask_conduct

check("conduct: 正常系はask()を1回だけ呼ぶ(引き直し無し)", len(_calls_happy) == 1)
check("conduct: 正常系のrawは最初のレスポンスそのもの", _raw_happy == _resp_happy)
check("conduct: 正常系のplanはJSONのmode/task_typeをvalidate_plan経由で反映",
      _plan_happy["mode"] == "moa" and _plan_happy["task_type"] == "chat")
check("conduct: 正常系のselected_proposersはペルソナ名から実モデル名へ解決",
      _plan_happy["selected_proposers"] == ["gpt-oss:20b", "qwen3-coder:30b"])

# --- 引き直し系: 1回目がJSON抽出失敗 -> 同一messagesで1回だけ引き直し、2回目を採用 ---
try:
    _resp_redraw = _valid_plan_json(task_type="writing", reason="second draw")
    f.ask, _calls_redraw = _make_conduct_ask(["not json at all, sorry", _resp_redraw])
    _plan_redraw, _raw_redraw = f.conduct("これは何ですか")
finally:
    f.ask = _orig_ask_conduct

check("conduct: 1回目JSON抽出失敗なら合計2回askする", len(_calls_redraw) == 2)
check("conduct: 引き直し時のmessagesは1回目と同一",
      _calls_redraw[0]["messages"] == _calls_redraw[1]["messages"])
check("conduct: 引き直し後のplanは2回目のJSONを反映",
      _plan_redraw["task_type"] == "writing" and "second draw" in _plan_redraw["reason"])
check("conduct: 引き直し後のrawは2回目のレスポンス", _raw_redraw == _resp_redraw)

# --- 二重失敗系: 両方JSON抽出失敗でも例外を送出せずdefault_plan()ベースへ劣化する ---
try:
    f.ask, _calls_double = _make_conduct_ask(["nonsense", "still nonsense"])
    _plan_fail, _raw_fail = f.conduct("42 を計算しなさい。答えを求めよ。")
finally:
    f.ask = _orig_ask_conduct

check("conduct: 両方失敗でもask()は2回で打ち切る(それ以上引き直さない)", len(_calls_double) == 2)
check("conduct: 両方失敗時のplanはdefault_planベースのフォールバック",
      _plan_fail.get("_fallback") is True)
check("conduct: フォールバックでも数学キーワードのtask_typeガードレールは効く",
      _plan_fail["task_type"] == "math")
check("conduct: 両方失敗時のrawは2回目の生レスポンス", _raw_fail == "still nonsense")

# --- ガードレール適用順の一貫性: validate_plan -> routing -> accuracy -> tasktype ---
try:
    _resp_code = _valid_plan_json(mode="single", task_type="knowledge",
                                   selected_proposers=["Proposer A"], reason="raw single")
    f.ask, _calls_code = _make_conduct_ask([_resp_code])
    _plan_code, _ = f.conduct("再帰関数を実装してください")
finally:
    f.ask = _orig_ask_conduct
check("conduct: コード質問はraw planがsingleでもmoaへ格上げ(精度ガードレール)",
      _plan_code["mode"] == "moa")

try:
    _resp_math = _valid_plan_json(mode="moa", task_type="knowledge", reason="raw math")
    f.ask, _calls_math = _make_conduct_ask([_resp_math])
    _plan_math, _ = f.conduct("100 の階乗を 7 で割った余りを求めよ")
finally:
    f.ask = _orig_ask_conduct
check("conduct: 数学キーワードはConductorの誤分類(knowledge)をmathへ上書き",
      _plan_math["task_type"] == "math")

try:
    _resp_pptx = _valid_plan_json(mode="single", task_type="writing",
                                   selected_proposers=["Proposer A"],
                                   image_only=True, use_image_generation=True,
                                   reason="raw pptx-ish")
    f.ask, _calls_pptx = _make_conduct_ask([_resp_pptx])
    _plan_pptx, _ = f.conduct("来週の会議用にパワポのスライド資料を作って")
finally:
    f.ask = _orig_ask_conduct
check("conduct: パワポ要求はraw planの中身に関わらずmake_pptx+moa+image_only=Falseに確定",
      _plan_pptx["make_pptx"] is True and _plan_pptx["mode"] == "moa"
      and _plan_pptx["image_only"] is False)

# --- Office添付ガードレール: conduct()の最終returnにも反映される(_apply_office_guardrail、
#     2026-07-26追加)。ここでは f.ask のみをモックし、Conductor が CONDUCTOR_SYS 特殊
#     ルーティング指示3/conduct()の自然文ヒントを無視してmode='single'かつProposer C不在の
#     rawプランを返した想定で、conduct()の最終出力が決定的ガードレールで補正されることを見る。 ---
try:
    _resp_office_guard = _valid_plan_json(mode="single", task_type="knowledge",
                                           selected_proposers=["Proposer A"],
                                           reason="raw single, ignores office hint")
    f.ask, _calls_office_guard = _make_conduct_ask([_resp_office_guard])
    _plan_office_guard, _ = f.conduct("この契約書を分析して", office_attached=True)
finally:
    f.ask = _orig_ask_conduct
check("conduct: office_attached=TrueならConductorがsingleを返してもmoaへ強制される",
      _plan_office_guard["mode"] == "moa")
check("conduct: office_attached=TrueならProposer Cのモデルが最終planに含まれる",
      f.PERSONA_MODELS["Proposer C"] in _plan_office_guard["selected_proposers"])

try:
    _resp_no_office_guard = _valid_plan_json(mode="single", task_type="knowledge",
                                              selected_proposers=["Proposer A"],
                                              reason="raw single, no office")
    f.ask, _calls_no_office_guard = _make_conduct_ask([_resp_no_office_guard])
    _plan_no_office_guard, _ = f.conduct("この契約書を分析して", office_attached=False)
finally:
    f.ask = _orig_ask_conduct
check("conduct: office_attached=Falseならsingleのまま"
      "(このガードレールはoffice_attached=Falseの通常planにmoaを強制しない)",
      _plan_no_office_guard["mode"] == "single")

# --- プロンプト構築: office_attached のヒント文言 ---
try:
    _resp_office = _valid_plan_json(reason="office test")
    f.ask, _calls_office = _make_conduct_ask([_resp_office])
    f.conduct("この資料の要点をまとめて", office_attached=True)
finally:
    f.ask = _orig_ask_conduct
_office_user_msg = _calls_office[0]["messages"][1]["content"]
check("conduct: office_attached=Trueでヒント文言(Proposer C / 特殊ルーティング指示#2)が挿入される",
      "Proposer C" in _office_user_msg and "特殊ルーティング指示 #2" in _office_user_msg)

try:
    _resp_no_office = _valid_plan_json(reason="no office test")
    f.ask, _calls_no_office = _make_conduct_ask([_resp_no_office])
    f.conduct("この資料の要点をまとめて", office_attached=False)
finally:
    f.ask = _orig_ask_conduct
_no_office_user_msg = _calls_no_office[0]["messages"][1]["content"]
check("conduct: office_attached=Falseではヒント文言が挿入されない",
      "特殊ルーティング指示 #2" not in _no_office_user_msg)

# --- プロンプト構築: history の直近の会話ノート(200字超は切り詰め) ---
try:
    _long_content = "あ" * 250
    _history = [
        {"role": "user", "content": "最古の質問(除外されるはず)"},
        {"role": "assistant", "content": "最古の回答(除外されるはず)"},
        {"role": "user", "content": "短い質問1"},
        {"role": "assistant", "content": "短い回答1"},
        {"role": "user", "content": _long_content},
        {"role": "assistant", "content": "直近の回答2"},
    ]
    _resp_hist = _valid_plan_json(reason="history test")
    f.ask, _calls_hist = _make_conduct_ask([_resp_hist])
    _plan_hist, _ = f.conduct("続きを教えて", history=_history)
finally:
    f.ask = _orig_ask_conduct
_hist_user_msg = _calls_hist[0]["messages"][1]["content"]
check("conduct: history指定時は直近の会話ノートが埋め込まれ例外なく完走する",
      "直近の会話" in _hist_user_msg)
check("conduct: 200字超のhistory本文は200字に切り詰められ...が付く",
      ("あ" * 200 + "...") in _hist_user_msg and ("あ" * 250) not in _hist_user_msg)
check("conduct: history[-4:]より古い発言は含まれない",
      "除外されるはず" not in _hist_user_msg)

f.PROPOSERS = _orig_props_conduct

# ---------- 自己一貫性投票（答え抽出・正規化・同値判定・投票） ----------
check("sc: boxed 抽出", f.extract_boxed("thus \\boxed{42}") == "42")
check("sc: boxed 入れ子", f.extract_boxed("\\boxed{\\frac{1}{2}}") == "\\frac{1}{2}")
check("sc: boxed 最後を採用", f.extract_boxed("\\boxed{1} then \\boxed{2}") == "2")
check("sc: boxed 無しは None", f.extract_boxed("no box") is None)
# 2026-07-22: \boxed{ が閉じられないまま出力が打ち切られた場合（thinking モデルの
# num_predict 打ち切り等、gotcha #2 の既知の失敗モード）は、切れた残骸を答えとして
# 返さず None（無投票）を返すことを検証する。
check("sc: boxed 未閉じは None（打ち切り）",
      f.extract_boxed("thus \\boxed{42 and then the response was cut off") is None)
check("sc: boxed 未閉じ・入れ子未対応も None",
      f.extract_boxed("\\boxed{\\frac{1}{2") is None)
check("sc: boxed 閉じ括弧後に散文があっても正しく抽出",
      f.extract_boxed("\\boxed{7} because it is prime") == "7")
check("sc: boxed 二重入れ子", f.extract_boxed("\\boxed{\\boxed{5}}") == "5")

# 2026-07-22 (iteration 12): 先に確定した \boxed{回答} があり、後続の
# \boxed{...} だけが打ち切られている場合は、手前の閉じた票を救出する
# （iteration 11 / gotcha #2, #7 参照。詳細は extract_boxed 本体のコメント）。
check("sc: boxed 後続が未閉じでも手前の確定票を救出",
      f.extract_boxed("\\boxed{42} then \\boxed{the next attempt got cut off") == "42")
check("sc: boxed 後続が入れ子ごと未閉じでも手前の確定票を救出",
      f.extract_boxed("\\boxed{7} ... \\boxed{\\frac{1}{2") == "7")

# 2026-07-22 (iteration 25, iteration 11/23 の続き): 末尾の \boxed{} が閉じては
# いるが中身が空/空白のみの場合は「無投票」として扱い、手前にある閉じた非空の
# 票まで遡って救出する（gotcha #2, #7 参照。詳細は extract_boxed 本体のコメント）。
# これをしないと extract_final_answer の math 分岐が None を受けて散文中の
# 数値を拾うフォールバックへ落ち、無投票のはずが誤投票に変わってしまう。
check("sc: boxed 末尾が空でも手前の確定票を救出",
      f.extract_boxed("\\boxed{5} then \\boxed{}") == "5")
check("sc: boxed 末尾が空白のみでも手前の確定票を救出",
      f.extract_boxed("\\boxed{42} then \\boxed{ }") == "42")
check("sc: boxed 単独の空は None（救出対象なし）",
      f.extract_boxed("\\boxed{}") is None)
check("sc: boxed 単独の空白のみも None（救出対象なし）",
      f.extract_boxed("\\boxed{ }") is None)
check("sc: boxed 末尾が空でも last-wins は非空同士で維持",
      f.extract_boxed("\\boxed{1} then \\boxed{2}") == "2")
check("sc: boxed 末尾空スキップ後もextract_final_answerが散文の数値を誤採用しない",
      f.extract_final_answer("\\boxed{5} then the loop ran 10 times \\boxed{}", "math") == "5")

check("sc: 正規化 全角→半角", f.normalize_answer("１２３") == "123")
check("sc: 正規化 桁区切り除去", f.normalize_answer("12,345") == "12345")
check("sc: 正規化 空白入り桁区切り", f.normalize_answer("11,\\! 111,\\! 111,\\! 100") == "11111111100")
check("sc: 正規化 前置き除去", f.normalize_answer("Answer: 700") == "700")
check("sc: 正規化 text外殻", f.normalize_answer("\\text{391}") == "391")
# 2026-07-23: \text/\mathrm 以外の体裁マクロ（\textbf/\mathbf/\boldsymbol 等）も
# 値を変えないので剥がす。mcq で \boxed{\textbf{B}} が票落ちしたり、math で
# \boxed{\mathbf{42}} が素の "42" 票と別クラスに分裂したりするのを防ぐ
# （normalize_answer 側のコメント参照、自己整合性投票 gotcha #7）。
check("sc: 正規化 textbf外殻", f.normalize_answer("\\textbf{B}") == "B")
check("sc: 正規化 mathbf外殻", f.normalize_answer("\\mathbf{42}") == "42")
check("sc: 正規化 boldsymbol外殻", f.normalize_answer("\\boldsymbol{x}") == "x")
check("sc: 正規化 入れ子体裁マクロ", f.normalize_answer("\\mathbf{\\text{D}}") == "D")
# 2026-07-24 (iteration 122): 単純な数値のみの \frac{a}{b} は "a/b" に正規化されるように
# なった（iteration 108 で3回スタックした修正のリトライ、詳細は normalize_answer 内コメント
# 参照）。変数を含む場合など非数値の \frac は引き続き対象外（下のテスト群で別途確認）。
check("sc: 正規化 数値fracはa/bに変換される", f.normalize_answer("\\frac{1}{2}") == "1/2")
# 2026-07-22: _FW_TRANS 拡張分（Unicode MINUS SIGN / 全角句点・読点・スラッシュ）の
# 正規化。CJK 寄りのプロポーザ (qwen/gemma 系) がこれらを出力し、正規化しないと
# vote_answers で本来同値な答えが2系統の票に割れてしまう（詳細は _FW_TRANS 定義部の
# コメント参照）。
check("sc: 正規化 U+2212マイナス", f.normalize_answer("−5") == "-5")
check("sc: 正規化 全角ハイフンマイナス", f.normalize_answer("－5") == "-5")
check("sc: 正規化 全角数字+全角句点(小数)", f.normalize_answer("３．１４") == "3.14")
check("sc: 正規化 全角スラッシュ(分数)", f.normalize_answer("１／２") == "1/2")
check("sc: 正規化 全角カンマ 桁区切り", f.normalize_answer("1，234") == "1234")
# 回帰: 既存ASCII表記は一切変わらないこと
check("sc: 正規化 ASCII -5 不変", f.normalize_answer("-5") == "-5")
check("sc: 正規化 ASCII 1/2 不変", f.normalize_answer("1/2") == "1/2")
check("sc: 正規化 ASCII 1,234 不変", f.normalize_answer("1,234") == "1234")
# 2026-07-22: 末尾カンマ除去（extract_final_answer の数値抽出正規表現 [\d,]* が
# 桁区切りでない末尾カンマまで貪欲に飲み込む問題への対処、normalize_answer 側のコメント参照）
check("sc: 正規化 末尾カンマ除去", f.normalize_answer("42,") == "42")
check("sc: 正規化 桁区切り+末尾カンマ除去", f.normalize_answer("1234,") == "1234")

# 2026-07-25: LaTeX でバレの '\%' と全角パーセント '％'(U+FF05) を '%' へ正規化する。
# 百分率の答えは \boxed{50\%}（LaTeX エスケープ、行儀の良いモデルが好む）/ \boxed{50%}
# （素の散文）/ \boxed{５０％}（CJK 寄りの qwen/gemma 系）の3系統の綴りに分裂しうるが、
# いずれも値としては同じ50%であり、投票を分裂させるべきではない（自己整合性投票
# gotcha #7、iteration 13/22/78/122 と同系統の姉妹修正。詳細は normalize_answer /
# _FW_TRANS 定義部のコメント参照）。パーセント記号自体は落とさない（50% を 50 という
# 別の値に変えてしまうと精度劣化になる）。
check("sc: 正規化 LaTeXエスケープ%", f.normalize_answer("50\\%") == "50%")
check("sc: 正規化 LaTeXエスケープ% + 末尾句点", f.normalize_answer("50\\%.") == "50%")
check("sc: 正規化 全角パーセント(全角数字込み)", f.normalize_answer("５０％") == "50%")
# 回帰: 素の '%' は従来どおり変化しない（%自体を落とさない・全角以外の桁区切り等）
check("sc: 正規化 素の%は不変", f.normalize_answer("50%") == "50%")
check("sc: 正規化 %以外の既存回帰は不変",
      f.normalize_answer("50") == "50"
      and f.normalize_answer("1/2") == "1/2"
      and f.normalize_answer("-5") == "-5"
      and f.normalize_answer(r"\frac{1}{2}") == "1/2"
      and f.normalize_answer("\\textbf{B}") == "B")

check("sc: 抽出 boxedのLaTeXエスケープ%とASCII%が同じ票に正規化される",
      f.extract_final_answer("\\boxed{50\\%}", "math")
      == f.extract_final_answer("\\boxed{50%}", "math")
      == "50%")

# 2026-07-25 (iteration 136): iteration 134 は normalize_answer/boxed 経路の '%' 保持を
# 修正したが、extract_final_answer の宣言ブランチ・最後の数値フォールバックの数値コア
# 正規表現は '%' を含めておらず、非boxedの百分率回答が "50%"（boxed票）と "50"
# （宣言/フォールバック票）という別クラスに分裂したまま残っていた（このタスクが仕上げる
# iteration 134 の積み残し）。宣言ブランチ・フォールバックの両方が boxed と同じ "50%" を
# 返し、3経路が1つの投票クラスに合流することを確認する。
check("sc: 抽出 宣言分岐の%が落ちない（boxed票との合流）",
      f.extract_final_answer("The final answer is 50%", "math") == "50%")
check("sc: 抽出 最後の数値フォールバックの%が落ちない（boxed票との合流）",
      f.extract_final_answer("so the probability is 50%.", "math") == "50%")
check("sc: 抽出 boxed/宣言/フォールバックの3経路が同一の50%票に合流",
      f.extract_final_answer("The final answer is 50%", "math")
      == f.extract_final_answer("so the probability is 50%.", "math")
      == f.extract_final_answer("\\boxed{50\\%}", "math")
      == "50%")
_pct_votes = [
    f.extract_final_answer("The final answer is 50%", "math"),
    f.extract_final_answer("so the probability is 50%.", "math"),
    f.extract_final_answer("\\boxed{50\\%}", "math"),
]
_pct_top, _pct_count, _pct_classes = f.vote_answers(_pct_votes)
check("sc: 投票 50%の3票が単一クラスに集約される（票割れしない）",
      _pct_top == "50%" and _pct_count == 3 and len(_pct_classes) == 1)
# 回帰: 途中に出てくるpercentは末尾の別の数値へ誤って付与されない
check("sc: 抽出 最後の数値フォールバック 途中percentは末尾数値に伝播しない",
      f.extract_final_answer("increased by 50% to reach 75", "math") == "75")

# answers_equivalent: math_verify を「呼ばれたら必ず例外」なスタブに差し替えても
# 高速パス（na.lower() 一致）だけで正しく判定できることを確認する（下の
# U+2212/全角スラッシュ用スタブブロック、L~763 と同じ swap-and-restore パターンだが、
# そちらの _fake_math_verify はまだ未定義のためここではローカルに同等のスタブを作る）。
def _mv_must_not_be_called_pct(*_a, **_kw):
    raise RuntimeError("math_verify should not be needed for these fast-path cases (percent)")


_fake_math_verify_pct = types.ModuleType("math_verify")
_fake_math_verify_pct.parse = _mv_must_not_be_called_pct
_fake_math_verify_pct.verify = _mv_must_not_be_called_pct
_orig_math_verify_mod_pct = sys.modules.get("math_verify")
sys.modules["math_verify"] = _fake_math_verify_pct
try:
    check("sc: 同値 LaTeXエスケープ%とASCII%（math_verify不要）",
          f.answers_equivalent("50\\%", "50%"))
    check("sc: 同値 全角パーセントとASCII%（math_verify不要）",
          f.answers_equivalent("５０％", "50%"))
    # 回帰: %を落として値まで変えてしまう誤修正でないことの確認。50% と 50 は
    # 異なる値であり、math_verify が例外を返す（=利用不能）環境でも高速パスのみで
    # 正しく非同値のまま（誤って併合されない）。
    check("sc: 非同値 50%と50は別クラスのまま（math_verifyスタブが例外を返しても誤併合しない）",
          not f.answers_equivalent("50%", "50"))
finally:
    if _orig_math_verify_mod_pct is not None:
        sys.modules["math_verify"] = _orig_math_verify_mod_pct
    else:
        del sys.modules["math_verify"]

# 2026-07-25: \(...\)（inline 数式）/ \[...\]（display 数式）はエスケープ済みの LaTeX
# 数式モード区切り文字であり、上の '$'（同じ数式モードの別綴り）/ '\!'/'\,' の除去と
# 対称的な、値を持たない体裁トークン。normalize_answer がこれらを剥がさないままだと
# \boxed{\(5\)} や \[x+1\]、「answer is \(42\)」のような区切り文字付きの答えが素の
# "5"/"x+1"/"42" とは別の投票クラスに分裂し（自己整合性投票 gotcha #7 の票割れ）、
# iteration 13/22/78/122/134/136 と同系統の姉妹修正としてここで剥がす。
check("sc: 正規化 \\(...\\) 剥がし", f.normalize_answer("\\(5\\)") == "5")
check("sc: 正規化 \\[...\\] 剥がし", f.normalize_answer("\\[x+1\\]") == "x+1")
check("sc: 正規化 \\(\\frac{}{}\\) は区切り剥がし後にfrac正規化される",
      f.normalize_answer("\\(\\frac{1}{2}\\)") == "1/2")
# 回帰: $ 剥がし・'\!' 桁区切り剥がし・素の値は不変のまま
check("sc: 正規化 \\(...\\)/\\[...\\] 以外の既存回帰は不変",
      f.normalize_answer("$5$") == "5"
      and f.normalize_answer("5") == "5"
      and f.normalize_answer("11,\\! 111,\\! 111,\\! 100") == "11111111100")

check("sc: 抽出 宣言分岐の\\(...\\)が剥がれる",
      f.extract_final_answer("The final answer is \\(42\\)", "math") == "42")
check("sc: 抽出 boxed分岐の\\(...\\)が剥がれる",
      f.extract_final_answer("\\boxed{\\(7\\)}", "math") == "7")

# answers_equivalent: math_verify を「呼ばれたら必ず例外」なスタブに差し替えても、
# \(5\) と 5 が normalize_answer の高速パス（na.lower() 一致）だけで合流することを確認する
# （上の百分率ブロックと同じ swap-and-restore パターン、iteration 122/129/134 と同系統）。
def _mv_must_not_be_called_delim(*_a, **_kw):
    raise RuntimeError("math_verify should not be needed for these fast-path cases (delimiters)")


_fake_math_verify_delim = types.ModuleType("math_verify")
_fake_math_verify_delim.parse = _mv_must_not_be_called_delim
_fake_math_verify_delim.verify = _mv_must_not_be_called_delim
_orig_math_verify_mod_delim = sys.modules.get("math_verify")
sys.modules["math_verify"] = _fake_math_verify_delim
try:
    check("sc: 同値 \\(5\\)と5（math_verify不要）", f.answers_equivalent("\\(5\\)", "5"))
    _delim_top, _delim_count, _delim_classes = f.vote_answers(["\\(5\\)", "5", "5"])
    check("sc: 投票 \\(5\\)/5/5の3票が単一クラスに集約される（票割れしない）",
          _delim_count == 3 and len(_delim_classes) == 1)
finally:
    if _orig_math_verify_mod_delim is not None:
        sys.modules["math_verify"] = _orig_math_verify_mod_delim
    else:
        del sys.modules["math_verify"]

check("sc: 抽出 答え宣言", f.extract_final_answer("計算すると、答えは 700 円です") == "700")
check("sc: 抽出 最後の数値", f.extract_final_answer("17 * 23 = 391") == "391")
check("sc: 抽出 無しは None", f.extract_final_answer("わかりません") is None)
# 2026-07-22: 末尾カンマを伴う抽出（宣言分岐・最後の数値フォールバックの両方）
check("sc: 抽出 最後の数値 末尾カンマ",
      f.extract_final_answer("so in total we get 42,", "math") == "42")
check("sc: 抽出 答え宣言 桁区切り+末尾カンマ",
      f.extract_final_answer("the final answer is 1,234,", "math") == "1234")
check("sc: 抽出 最後の数値(boxedなし) 末尾カンマ",
      f.extract_final_answer("17 * 23 = 391,", "math") == "391")
# 2026-07-22: 最後の数値フォールバックの符号クラスに Unicode マイナス(U+2212)/
# 全角ハイフンマイナス(U+FF0D)を追加した回帰確認。\boxed{} も「答え」宣言もない
# 終端数値のみのケースで、CJK プロポーザが出しがちな全角/Unicode 符号付き負数が
# 正の値として誤投票されないことを検証する（extract_final_answer 内のコメント参照）。
check("sc: 抽出 最後の数値 U+2212マイナス(boxed/宣言なし)",
      f.extract_final_answer("計算の結果は −5", "math") == "-5")
# 注: 「答え/正解/answer」を含む文言だと宣言ブランチ(2318行目)が先に拾ってしまい
# ここで検証したい「最後の数値フォールバック」に到達しないため、あえてそれらの
# キーワードを含まない文言を使う。
check("sc: 抽出 最後の数値 全角ハイフンマイナス(boxed/宣言なし)",
      f.extract_final_answer("結論としては、最終的な値は －5である", "math") == "-5")
check("sc: 抽出 最後の数値 U+2212マイナスとASCIIの投票クラス一致",
      f.answers_equivalent(f.extract_final_answer("最終値は −5", "math"), "-5"))
# 2026-07-22: 最後の数値フォールバック（および宣言ブランチの数値部抽出）の整数部
# 文字クラスを [\d,]* から「桁区切りとして妥当なカンマのみ許容」に厳格化した回帰確認
# （iteration 13/22/24 と同じ抽出経路の姉妹修正、fugu_local.py 側のコメント参照）。
# \boxed{} も「答え/正解/answer」宣言も無く、桁区切りとして不正なカンマ区切りの
# 数値列で終わる文では、1トークンに誤結合された "1,2,3" ではなく最後の数値のみを拾う。
check("sc: 抽出 最後の数値 不正なカンマ区切り列は結合されない",
      f.extract_final_answer("the roots are 1,2,3", "math") == "3")
check("sc: 抽出 最後の数値 座標のカンマ区切りは結合されない",
      f.extract_final_answer("the point is (1,2)", "math") == "2")
check("sc: 抽出 最後の数値 不正カンマ区切り列がそのまま誤投票票にならない",
      f.answers_equivalent(f.extract_final_answer("the roots are 1,2,3", "math"), "3"))
# 桁区切りとして正当なカンマ(3桁区切り)は引き続き1トークンとして丸ごと拾う回帰確認。
check("sc: 抽出 最後の数値(boxed/宣言なし) 桁区切り1234",
      f.extract_final_answer("in total we counted up to 1,234", "math") == "1234")
check("sc: 抽出 最後の数値(boxed/宣言なし) 桁区切り1234567",
      f.extract_final_answer("in total we counted up to 1,234,567", "math") == "1234567")
check("sc: 抽出 最後の数値(boxed/宣言なし) 桁区切り12345",
      f.extract_final_answer("in total we counted up to 12,345", "math") == "12345")
# 2026-07-23: \boxed{\mathbf{42}} を mathbf 外殻ごと剥がさないと文字列 "\mathbf{42}" の
# ままとなり、answers_equivalent の na.lower()/Fraction 系ファストパスに乗らず素の
# "42" 票と別の投票クラスに分裂していた（自己整合性投票 gotcha #7 の票割れ回帰確認）。
_mathbf_boxed = f.extract_final_answer("\\boxed{\\mathbf{42}}", "math")
check("sc: 抽出 boxed mathbf外殻（票割れ回帰）", _mathbf_boxed == "42")
check("sc: 同値 boxed mathbf外殻が素の42と合流（票割れ回帰）",
      f.answers_equivalent(_mathbf_boxed, "42"))
check("sc: mcq boxed", f.extract_final_answer("\\boxed{B}", "mcq") == "B")
check("sc: mcq 宣言", f.extract_final_answer("正解は (C) です", "mcq") == "C")
check("sc: mcq 無しは None", f.extract_final_answer("どれも違う", "mcq") is None)
check("sc: mcq boxed 散文混じりは先頭文字",
      f.extract_final_answer("reasoning...\\boxed{C, because it is the largest}", "mcq") == "C")
check("sc: mcq boxed 散文のみは誤答せず None",
      f.extract_final_answer("\\boxed{None of the above}", "mcq") is None)
check("sc: mcq boxed 括弧付き先頭文字", f.extract_final_answer("\\boxed{(A)}", "mcq") == "A")
check("sc: mcq boxed text外殻付き先頭文字", f.extract_final_answer("\\boxed{\\text{D}}", "mcq") == "D")
# 2026-07-23: \boxed{\textbf{B}} は textbf を剥がさないと先頭選択肢文字の正規表現が
# "\TEXTBF{B}" を見てマッチせず None（票落ち）になっていた回帰確認。
check("sc: mcq boxed textbf外殻付き先頭文字（票落ち回帰）",
      f.extract_final_answer("\\boxed{\\textbf{B}}", "mcq") == "B")
check("sc: mcq boxed 選択肢+本文", f.extract_final_answer("\\boxed{A) 5}", "mcq") == "A")
check("sc: mcq 宣言 ディストラクタ言及に釣られない",
      f.extract_final_answer(
          "The correct answer is B. Note that answer A was a common distractor.",
          "mcq") == "B")
check("sc: mcq 宣言 訂正で文字が競合したら棄権",
      f.extract_final_answer(
          "The answer is B; oh wait, the answer: A", "mcq") is None)
check("sc: mcq 宣言 同一文字の繰り返しは誤棄権しない",
      f.extract_final_answer(
          "The answer is D. Restating: the answer is D.", "mcq") == "D")

# 2026-07-24: _FW_TRANS（iter 13）は全角数字/A-E/マイナス/小数点/スラッシュ/カンマは
# 正規化するが全角括弧（U+FF08/FF09）は対象外のまま。CJK寄りのプロポーザーが
# \boxed{（A）}や『答えは（B）です』のように選択肢文字を全角括弧で囲むと、iter 3/26/78
# の mcq 抽出正規表現が ASCII 括弧しか許容していなかったため None（無投票）になり、
# 自己整合性投票（gotcha #7）から正当な1票が静かに失われていた。[(（]?/[)）]? へ
# 括弧クラスを広げてこの票落ちを回復する回帰確認（無投票より正しい1票）。
check("sc: mcq boxed 全角括弧付き先頭文字（票落ち回帰）",
      f.extract_final_answer("\\boxed{（A）}", "mcq") == "A")
check("sc: mcq boxed 全角文字+全角括弧（票落ち回帰）",
      f.extract_final_answer("\\boxed{（Ａ）}", "mcq") == "A")
check("sc: mcq 宣言 全角括弧付き（票落ち回帰）",
      f.extract_final_answer("したがって答えは（B）です。", "mcq") == "B")
check("sc: mcq 単独行 全角括弧のみ（票落ち回帰）",
      f.extract_final_answer("（C）", "mcq") == "C")
# 既存のASCII括弧・散文混じり・散文のみのケースが全角対応後も従来通りであることの回帰確認
check("sc: mcq boxed ASCII括弧は従来通り（全角対応の副作用なし）",
      f.extract_final_answer("\\boxed{(A)}", "mcq") == "A")
check("sc: mcq boxed 散文のみは全角対応後もNone（iter 3ガード維持）",
      f.extract_final_answer("\\boxed{None of the above}", "mcq") is None)
check("sc: mcq boxed 散文混じりは全角対応後も先頭文字（iter 3ガード維持）",
      f.extract_final_answer("\\boxed{C, because it is the largest}", "mcq") == "C")
check("sc: mcq 宣言 全角対応後も訂正で文字が競合したら棄権（iter 26ガード維持）",
      f.extract_final_answer(
          "The answer is B; oh wait, the answer: A", "mcq") is None)

# 2026-07-24 (iter 109): 上の全角括弧対応時（iter 102）に宣言パターンの文字クラスは
# [A-EＡ-Ｅ] へ広げたが、単独行パターンは [A-E] のままASCII限定に取り残されていた。
# そのためCJK寄りのプロポーザー（qwen/gemma系、iter 102 が想定した対象そのもの）が
# \boxed{} を無視して「Ｃ」「（Ｃ）」のように全角の選択肢文字だけを単独行で答えると、
# boxed分岐（該当なし）・宣言パターン（answer/答え/正解の連結詞なし）・単独行パターン
# （[A-E] が U+FF23 等の全角文字に不一致）のいずれにも拾われず None となり、
# 自己整合性投票（gotcha #7）の正当な1票が静かに失われていた。宣言パターン（iter 102）
# と同じ [A-EＡ-Ｅ] に揃えてこの票落ちを解消する。iter 13 の _FW_TRANS が全角→ASCII
# 正規化を担うため誤投票のリスクはなく、iter 3 の boxed 先頭文字ガード・iter 26 の
# 複数文字競合時の棄権ロジックもそのまま維持される。
check("sc: mcq 単独行 全角括弧+全角文字（票落ち回帰, iter109）",
      f.extract_final_answer("（Ｃ）", "mcq") == "C")
check("sc: mcq 単独行 括弧無し全角文字のみ（票落ち回帰, iter109）",
      f.extract_final_answer("Ｂ", "mcq") == "B")
check("sc: mcq 単独行 全角文字+が正解サフィックス（票落ち回帰, iter109）",
      f.extract_final_answer("Ｄが正解", "mcq") == "D")
# 回帰: 全角括弧+ASCII文字（iter 102 で既に対応済み）は今回の変更後も不変
check("sc: mcq 単独行 全角括弧+ASCII文字は従来通り（iter102回帰）",
      f.extract_final_answer("（C）", "mcq") == "C")
# 回帰: 素のASCII単独行文字（括弧・連結詞なし）と小文字（re.IGNORECASE）は不変
check("sc: mcq 単独行 ASCII文字のみ大文字は従来通り",
      f.extract_final_answer("C", "mcq") == "C")
check("sc: mcq 単独行 ASCII文字のみ小文字はIGNORECASEでCへ",
      f.extract_final_answer("c", "mcq") == "C")
# 回帰(iter 26棄権ロジック維持): 全角文字同士でも複数行で異なる文字が競合したらNone
check("sc: mcq 単独行 全角文字が複数行で競合したら棄権（iter26ガード維持）",
      f.extract_final_answer("Ａ\nＢ", "mcq") is None)
# 回帰: 文中(行の一部)に埋没した全角文字は ^...$ アンカーでスコープ外のままNone
check("sc: mcq 単独行 文中に埋没した全角文字は捕捉しない(^...$維持)",
      f.extract_final_answer("選択肢Ｃが気になるが自信はない", "mcq") is None)

# 2026-07-25: SC_PROMPT_MCQ は \boxed{} 指示だが、CJK寄りのプロポーザー（qwen/gemma系、
# iter 3/26/102/109 が繰り返し記録している対象そのもの）が無視して散文で答えを書き、
# その答えを Markdown の強調記号（**太字**/*斜体*/__太字__）で装飾するのは LLM の
# よくある癖。strip_think/normalize_answer は '*'/'_' を除去しないため、装飾された
# 選択肢文字が boxed 先頭文字・宣言・単独行の3分岐いずれにもマッチせず None（票落ち）に
# なっていた。文字に直接隣接する強調記号のみを許容して票を回復する回帰確認。
check("sc: mcq 宣言 太字強調（票落ち回帰, iter173）",
      f.extract_final_answer("The answer is **B**", "mcq") == "B")
check("sc: mcq 宣言 斜体強調（票落ち回帰, iter173）",
      f.extract_final_answer("Answer: *C*", "mcq") == "C")
check("sc: mcq 単独行 太字強調（票落ち回帰, iter173）",
      f.extract_final_answer("**D**", "mcq") == "D")
check("sc: mcq 単独行 アンダースコア太字強調（票落ち回帰, iter173）",
      f.extract_final_answer("__A__", "mcq") == "A")
check("sc: mcq boxed 太字強調（票落ち回帰, iter173）",
      f.extract_final_answer("\\boxed{**B**}", "mcq") == "B")
# 回帰: 強調記号なしの既存ケースは今回の変更後も不変
check("sc: mcq 単独行 全角括弧のみは強調対応後も従来通り（回帰）",
      f.extract_final_answer("答えは（B）です", "mcq") == "B")
check("sc: mcq 単独行 全角括弧は強調対応後も従来通り（回帰）",
      f.extract_final_answer("（C）", "mcq") == "C")
check("sc: mcq boxed ASCII括弧は強調対応後も従来通り（回帰）",
      f.extract_final_answer("\\boxed{(A)}", "mcq") == "A")
check("sc: mcq boxed 散文混じりは強調対応後も先頭文字（回帰）",
      f.extract_final_answer("\\boxed{C, because it is the largest}", "mcq") == "C")
check("sc: mcq boxed 散文のみは強調対応後もNone（回帰）",
      f.extract_final_answer("\\boxed{None of the above}", "mcq") is None)
# 競合棄権(iter 26ガード維持): 強調記号付き文字と素の文字が競合したらNone
check("sc: mcq 宣言 強調記号付きと素の文字が競合したら棄権（iter26ガード維持）",
      f.extract_final_answer(
          "the answer is **A**\nthe answer is B", "mcq") is None)
# 誤爆防止: 強調された単語の中から文字を誤って拾わない（iter 3 の \b／(?![A-Za-z]) ガード維持）
check("sc: mcq 宣言 太字語の中から誤爆しない（'**Bee**' が B にならない）",
      f.extract_final_answer("the answer is **Bee**", "mcq") is None)
check("sc: mcq 単独行 太字語の中から誤爆しない（'**Bee**' が B にならない）",
      f.extract_final_answer("**Bee**", "mcq") is None)
# 誤爆防止: 太字見出し行は単独行パターンの ^...$ アンカーでスコープ外のまま(iter109ガード維持)
check("sc: mcq 単独行 太字見出し行は捕捉しない(^...$維持)",
      f.extract_final_answer("**A. Introduction**", "mcq") is None)
# 非mcqパス(math)は今回の変更で一切触れていないことの確認: bold付き文言でも従来通り
# 数値抽出ロジックのみが働き、'*'/'_' がそのまま残った文字列を返す(グローバル除去なし)。
check("sc: math 抽出は強調記号対応の影響を受けない（byte-for-byte不変）",
      f.extract_final_answer("\\boxed{**42**}", "math") == "**42**")

# 2026-07-22: iteration 28 — math 宣言ブランチにも iteration 26 の MCQ 修正
# （複数宣言が競合したら無投票=None）を対称に適用した回帰確認。
# 注: 宣言抽出の捕獲グループ ([^\n]{1,60}) は改行を跨がず貪欲マッチするため、同一行に
# 複数の「answer is」があると最初のマッチが行末まで飲み込み2件目の宣言が独立して
# 検出されない。複数宣言を意図的に分離検出させるため、ここでは改行で区切って書く
# （実際の LLM 出力でも言い直し・訂正は改行/文区切りを伴うことが多い）。
check("sc: 抽出 答え宣言 訂正で数値が競合したら棄権",
      f.extract_final_answer(
          "The answer is 5.\nOn second thought, the answer is 7.", "math") is None)
_eq_restate = f.extract_final_answer(
    "The answer is 1/2.\nEquivalently, the answer is 0.5.", "math")
check("sc: 抽出 答え宣言 同値な言い直しは棄権しない",
      _eq_restate is not None and f.answers_equivalent(_eq_restate, "0.5"))
check("sc: 抽出 答え宣言 同一値の繰り返しは誤棄権しない",
      f.extract_final_answer(
          "The answer is 42.\nRestating: the answer is 42.", "math") == "42")
check("sc: 抽出 答え宣言 空の宣言候補は競合と数えない",
      f.extract_final_answer(
          "The answer is .\nThe answer is 9.", "math") == "9")

check("sc: 同値 完全一致", f.answers_equivalent("42", "42"))
check("sc: 同値 分数=小数", f.answers_equivalent("1/2", "0.5"))
check("sc: 同値 桁区切り", f.answers_equivalent("12,345", "12345"))
check("sc: 非同値", not f.answers_equivalent("41", "42"))
check("sc: 空は非同値", not f.answers_equivalent("", "42"))

# 2026-07-22: Unicode マイナス/全角スラッシュの同値判定が na.lower()/Fraction の
# 高速パスだけで完結し、math_verify に頼らないことを検証する。math_verify を
# 「呼ばれたら必ず例外」なスタブに差し替えても answers_equivalent が True を
# 返せることを確認し、フォールバック依存になっていないことを保証する。
def _mv_must_not_be_called(*_a, **_kw):
    raise RuntimeError("math_verify should not be needed for these fast-path cases")


_fake_math_verify = types.ModuleType("math_verify")
_fake_math_verify.parse = _mv_must_not_be_called
_fake_math_verify.verify = _mv_must_not_be_called
_orig_math_verify_mod = sys.modules.get("math_verify")
sys.modules["math_verify"] = _fake_math_verify
try:
    check("sc: 同値 U+2212マイナス（math_verify不要）", f.answers_equivalent("−5", "-5"))
    check("sc: 同値 全角スラッシュ分数（math_verify不要）", f.answers_equivalent("１／２", "1/2"))
finally:
    if _orig_math_verify_mod is not None:
        sys.modules["math_verify"] = _orig_math_verify_mod
    else:
        del sys.modules["math_verify"]

# 2026-07-22: math_verify フォールバック分岐そのもの（fugu_local.py ~L2433-2444）の
# 直接カバレッジ。上のテストは「フォールバックに頼らない」ことの証明であり、フォールバック
# 分岐自体には一度も入っていない。ここでは高速パス（na.lower()一致 / Fraction一致）を
# 意図的に迂回する入力（\frac{1}{2} vs 0.5、どちらも正規化後は非空・lower()不一致・
# Fraction変換失敗）を使い、記録スタブで (1) parse/verify が実際に呼ばれたこと、
# (2) gotcha #6 の parsing_timeout=None / timeout_seconds=None が渡っていること
# （Windows で math_verify 既定タイムアウトがハンドルエラーを撒く実測不具合への回帰防止）、
# (3) verify の戻り値がそのまま answers_equivalent の戻り値になること、
# (4) parse/verify が例外を送出しても except で握り潰され False になり例外が外に漏れない
# ことを検証する。


def _make_recording_math_verify(verify_result, raise_in):
    """呼び出しを記録する math_verify スタブモジュールを生成する。
    raise_in: "none" | "parse" | "verify" — 該当関数が呼ばれたら例外を送出する。"""
    calls = {"parse_args": [], "parse_kwargs": [], "verify_args": [], "verify_kwargs": []}

    def _parse(expr, **kwargs):
        calls["parse_args"].append(expr)
        calls["parse_kwargs"].append(kwargs)
        if raise_in == "parse":
            raise RuntimeError("boom in parse (stub)")
        return ("parsed", expr)

    def _verify(parsed_a, parsed_b, **kwargs):
        calls["verify_args"].append((parsed_a, parsed_b))
        calls["verify_kwargs"].append(kwargs)
        if raise_in == "verify":
            raise RuntimeError("boom in verify (stub)")
        return verify_result

    mod = types.ModuleType("math_verify")
    mod.parse = _parse
    mod.verify = _verify
    return mod, calls


def _run_with_math_verify_stub(verify_result, raise_in, body):
    """math_verify を記録スタブに差し替えて body(calls) を実行し、必ず元に戻す
    （L341-353 と同じ swap-and-restore パターン）。"""
    mod, calls = _make_recording_math_verify(verify_result, raise_in)
    orig = sys.modules.get("math_verify")
    sys.modules["math_verify"] = mod
    try:
        body(calls)
    finally:
        if orig is not None:
            sys.modules["math_verify"] = orig
        else:
            del sys.modules["math_verify"]


# 高速パスを迂回する入力ペア（正規化後も非空・lower()不一致・Fraction変換失敗）。
# 2026-07-24: 以前はここで \frac{1}{2} vs 0.5 を使っていたが、iteration 122 で
# normalize_answer に単純数値の \frac{a}{b} 正規化（"a/b" 化）を追加したため、
# \frac{1}{2} は "1/2" に正規化されて Fraction 高速パスに乗るようになった
# （math_verify フォールバックには到達しなくなる）。math_verify 呼び出しの実検証には
# 高速パスに絶対に乗らないペア（\sqrt は非対応マクロなので数値化されない）に差し替える。
_MV_A, _MV_B = r"\sqrt{2}", "1.41421356"


def _t_mv_verify_true(calls):
    result = f.answers_equivalent(_MV_A, _MV_B)
    check("sc: math_verifyフォールバック parseが実呼出しされる", len(calls["parse_args"]) >= 2)
    check("sc: math_verifyフォールバック verifyが実呼出しされる", len(calls["verify_args"]) == 1)
    check("sc: math_verifyフォールバック parsing_timeout=None (gotcha#6)",
          len(calls["parse_kwargs"]) >= 2
          and all(kw.get("parsing_timeout", "MISSING") is None for kw in calls["parse_kwargs"]))
    check("sc: math_verifyフォールバック timeout_seconds=None (gotcha#6)",
          len(calls["verify_kwargs"]) == 1
          and all(kw.get("timeout_seconds", "MISSING") is None for kw in calls["verify_kwargs"]))
    check("sc: math_verifyフォールバック verify=Trueを伝播", result is True)


_run_with_math_verify_stub(True, "none", _t_mv_verify_true)


def _t_mv_verify_false(calls):
    result = f.answers_equivalent(_MV_A, _MV_B)
    check("sc: math_verifyフォールバック verify=Falseを伝播", result is False)


_run_with_math_verify_stub(False, "none", _t_mv_verify_false)


def _t_mv_parse_raises(calls):
    result = f.answers_equivalent(_MV_A, _MV_B)
    check("sc: math_verifyフォールバック parse例外はFalseに握り潰す（例外は漏れない）",
          result is False)


_run_with_math_verify_stub(None, "parse", _t_mv_parse_raises)


def _t_mv_verify_raises(calls):
    result = f.answers_equivalent(_MV_A, _MV_B)
    check("sc: math_verifyフォールバック verify例外はFalseに握り潰す（例外は漏れない）",
          result is False)


_run_with_math_verify_stub(None, "verify", _t_mv_verify_raises)

# math_verify の差し替えが確実に元へ復元されていること（スタブ混入なら False に化けるはずの
# 高速パスが正常動作することで間接確認する）
check("sc: math_verifyスタブ解除後も高速パスが正常動作（sys.modules復元確認）",
      f.answers_equivalent("42", "42") and not f.answers_equivalent("41", "42"))

_top, _cnt, _cls = f.vote_answers(["42", "42", "41", "0.5", "1/2", None, ""])
check("sc: 投票 最多クラス", _top == "42" and _cnt == 2)
check("sc: 投票 同値クラス集約", any(c[1] == 2 and f.answers_equivalent(c[0], "0.5") for c in _cls))
check("sc: 投票 空リスト", f.vote_answers([]) == (None, 0, []))

# ---------- vote_answers（tie安定性・代表選定・None/空フィルタ・件数降順） ----------
# gotcha #7 の自己整合性投票（solve_verifiable/_arbitrate）が構造的に依拠する契約を
# 直接ロックする。iteration 32 (answers_equivalent) / iteration 52 (_sc_sample) の
# 姉妹カバレッジ。全て answers_equivalent の高速パス（normalize一致/Fraction/_FW_TRANS）
# のみで解決する組み合わせを使い、math_verify の実体（サブプロセス/ライブラリ有無）には
# 依存しない。

# (a) 同数タイは「先出現が勝つ」: classes.sort(key=lambda c: -c[1]) は Python の安定
#     ソートなので、件数が同じクラス同士は出現順（挿入順）が保たれる。solve_verifiable の
#     拮抗判定 classes[0][1]==classes[1][1] や、_arbitrate が None を返した際の
#     先出現フォールバックは、この安定性に構造的に依存している。Counter や非安定ソートへの
#     リファクタはここを静かに壊しうる。
_top_a1, _cnt_a1, _cls_a1 = f.vote_answers(["7", "3", "3", "7"])
check("vote: 同数タイは先出現(7)が勝つ",
      _top_a1 == "7" and _cnt_a1 == 2 and _cls_a1 == [["7", 2], ["3", 2]])

_top_a2, _cnt_a2, _cls_a2 = f.vote_answers(["3", "7", "7", "3"])
check("vote: 入力順を逆にすると先出現(3)が勝つ(値の大小ではない)",
      _top_a2 == "3" and _cnt_a2 == 2 and _cls_a2 == [["3", 2], ["7", 2]])

# (b) 統合クラスの代表文字列は「最初に見た表記」。Fraction 高速パスで 1/2 と 0.5 が
#     同一クラスに集約されても、代表（res['votes'] のキー/_arbitrate の候補ラベルに使われる）
#     は後から来た '0.5' ではなく先出の '1/2' のまま。
_top_b, _cnt_b, _cls_b = f.vote_answers(["1/2", "0.5", "0.5"])
check("vote: 統合クラスの代表は最初に見た表記(1/2)のまま(0.5にならない)",
      (_top_b, _cnt_b, _cls_b) == ("1/2", 3, [["1/2", 3]]))

# (c) normalize_answer / _FW_TRANS 高速パス経由の統合（全角マイナス U+2212 → 半角ハイフン）。
_top_c, _cnt_c, _cls_c = f.vote_answers(["-5", "−5"])
check("vote: 全角マイナス(U+2212)と半角ハイフンの'-5'は同一クラスに統合される",
      (_top_c, _cnt_c, _cls_c) == ("-5", 2, [["-5", 2]]))

# (d) None/空文字は集計から除外される。
_top_d1, _cnt_d1, _cls_d1 = f.vote_answers(["5", None, "", "5"])
check("vote: Noneと空文字は集計から除外される",
      _top_d1 == "5" and _cnt_d1 == 2 and len(_cls_d1) == 1)
check("vote: 全てNone/空文字なら(None, 0, [])",
      f.vote_answers([None, "", None]) == (None, 0, []))

# (e) classes は件数降順で返る。入力の先頭に出た低頻度クラスも件数順どおり後ろに回る。
_top_e, _cnt_e, _cls_e = f.vote_answers(["5", "9", "9", "9", "2", "2"])
check("vote: classesは件数降順(入力先頭の低頻度クラスは後ろに回る)",
      _top_e == "9" and _cnt_e == 3
      and _cls_e == [["9", 3], ["2", 2], ["5", 1]])

# ---------- solve_verifiable（ask をモックして適応サンプリングを検証） ----------
_orig_ask2 = f.ask
_orig_props2 = f.PROPOSERS
_orig_reasoning = f.REASONING_MODELS
_orig_cheap = f.SC_CHEAP_VOTES
_orig_pot = f.SC_POT
_sc_calls = []


def _fake_sc_ask(model, messages, temperature, think=None, fmt=None,
                 label=None, num_predict=None, num_ctx=None):
    _sc_calls.append(model)
    return f"reasoning...\n\\boxed{{{'42' if len(_sc_calls) % 2 else '42'}}}"


try:
    f.PROPOSERS = ["m1", "m2"]
    f.REASONING_MODELS = ["m1", "m2"]
    f.SC_CHEAP_VOTES = 0
    f.SC_POT = False
    f.ask = _fake_sc_ask
    _res = f.solve_verifiable("test question", "math")
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props2
    f.REASONING_MODELS = _orig_reasoning
    f.SC_CHEAP_VOTES = _orig_cheap
    f.SC_POT = _orig_pot
check("sc: 全会一致で早期確定", _res is not None and _res["answer"] == "42")
check("sc: 初回バッチのみで停止", len(_sc_calls) == f.SC_INITIAL)
check("sc: モデルを交互に使う", set(_sc_calls) == {"m1", "m2"})
# 2026-07-22 回帰: 拮抗/裁定が一切発生しないこのパスでは、今回の裁定後cnt/votes再計算
# 修正の影響を受けず、votes/n_samples が従来どおり返ることを明示的に確認する。
check("sc: 拮抗なし(全会一致)のvotes/n_samplesは修正の影響を受けない(不変)",
      _res is not None and _res["votes"] == {"42": f.SC_INITIAL}
      and _res["n_samples"] == f.SC_INITIAL)

# 票が割れるケース: 第1バッチで拮抗 → 追加サンプリング後に過半数で確定。
# バッチ化により第1バッチ(SC_INITIAL)は m1 まとめ→m2 まとめの順。call index で答えを固定し、
# 第1バッチを均等割り(過半数なし)にして 2 バッチ目で決着させる。
_sc_calls.clear()
_seq = (["\\boxed{1}"] * (f.SC_INITIAL // 2) + ["\\boxed{2}"] * (f.SC_INITIAL - f.SC_INITIAL // 2)
        + ["\\boxed{1}"] * 100)   # 第1バッチは均等、以降は 1 が積み上がる


def _fake_sc_ask2(model, messages, temperature, think=None, fmt=None,
                  label=None, num_predict=None, num_ctx=None):
    _sc_calls.append(model)
    idx = len(_sc_calls) - 1
    return _seq[idx] if idx < len(_seq) else "\\boxed{1}"


try:
    f.PROPOSERS = ["m1", "m2"]
    f.REASONING_MODELS = ["m1", "m2"]
    f.SC_CHEAP_VOTES = 0
    f.SC_POT = False
    f.ask = _fake_sc_ask2
    _res2 = f.solve_verifiable("test question", "math")
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props2
    f.REASONING_MODELS = _orig_reasoning
    f.SC_CHEAP_VOTES = _orig_cheap
    f.SC_POT = _orig_pot
check("sc: 割れたら追加サンプリング", len(_sc_calls) > f.SC_INITIAL)
check("sc: 追加後に過半数で確定", _res2 is not None and _res2["answer"] == "1")

check("sc: SC_MIN_VOTES 定数", f.SC_MIN_VOTES == 3)

# 疑似全会一致ガード: 第1バッチで抽出成功が1票だけ（他は thinking打ち切り/boxed無しで
# 抽出失敗）だと、旧ロジックでは cnt(1)==n(1) で「全会一致」扱いになり k=1 で確定して
# しまっていた（2026-07-21 に発見・修正）。SC_MIN_VOTES 導入後は n<3 の全会一致では
# 確定させず、add_batch(SC_STEP) で追加サンプリングされることを検証する。
_orig_min_votes = f.SC_MIN_VOTES
_sc_calls.clear()
_seq3 = (["\\boxed{42}"] + ["すみません、答えが導けませんでした。"] * 5
         + ["\\boxed{42}"] * 100)  # 第1バッチ: 1票のみ抽出成功、以降は追加分がすべて42に収束


def _fake_sc_ask3(model, messages, temperature, think=None, fmt=None,
                  label=None, num_predict=None, num_ctx=None):
    _sc_calls.append(model)
    idx = len(_sc_calls) - 1
    return _seq3[idx] if idx < len(_seq3) else "\\boxed{42}"


try:
    f.PROPOSERS = ["m1", "m2"]
    f.REASONING_MODELS = ["m1", "m2"]
    f.SC_CHEAP_VOTES = 0
    f.SC_POT = False
    f.ask = _fake_sc_ask3
    _res3 = f.solve_verifiable("test question", "math")
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props2
    f.REASONING_MODELS = _orig_reasoning
    f.SC_CHEAP_VOTES = _orig_cheap
    f.SC_POT = _orig_pot
    f.SC_MIN_VOTES = _orig_min_votes
check("sc: n<SC_MIN_VOTES の疑似全会一致では確定しない(追加サンプリング)",
      len(_sc_calls) > f.SC_INITIAL)
check("sc: 追加サンプリング後に正しく確定", _res3 is not None and _res3["answer"] == "42")

# 抽出成功が一度もない場合: 全バッチで n=0 のまま SC_MAX に到達し、無限ループせず
# None を返して MoA フォールバックへ委ねることを検証する（打ち切り自体は既存ロジック）。
_sc_calls.clear()


def _fake_sc_ask_noextract(model, messages, temperature, think=None, fmt=None,
                           label=None, num_predict=None, num_ctx=None):
    _sc_calls.append(model)
    return "考え中ですが、最終的な答えを出せませんでした。"


try:
    f.PROPOSERS = ["m1", "m2"]
    f.REASONING_MODELS = ["m1", "m2"]
    f.SC_CHEAP_VOTES = 0
    f.SC_POT = False
    f.ask = _fake_sc_ask_noextract
    _res4 = f.solve_verifiable("test question", "math")
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props2
    f.REASONING_MODELS = _orig_reasoning
    f.SC_CHEAP_VOTES = _orig_cheap
    f.SC_POT = _orig_pot
    f.SC_MIN_VOTES = _orig_min_votes
check("sc: 抽出0票が続いてもハングせず終了", len(_sc_calls) > 0)
check("sc: 抽出0票なら None を返す(MoAへフォールバック)", _res4 is None)

check("sc: 拮抗なし時は _arbitrate 未使用・勝者サンプルの本文を採用(従来通り)",
      _res is not None and _res["text"].startswith("reasoning..."))

# ---------- SC_MIN_VOTES の床を最終returnにも適用（SC_MAX消化パス） ----------
# 2026-07-21: ループ内の早期確定条件（cnt==n and n>=SC_MIN_VOTES / n>=4 and cnt*2>n）は
# while ループの break だけを守っており、SC_MAX 消化で抜けた最終 return には床が
# 掛かっていなかった。thinking打ち切りで __ERROR__、PoT失敗、\boxed{}欠落などにより
# ほとんどのサンプルが抽出失敗すると、1〜2票しか無い「勝者」がそのまま確定扱いで返る
# 疑似全会一致バグが再現する。ここでは終始 2 サンプルしか抽出成功しない（残りは全て
# 抽出不能）状況を作り、SC_MAX に到達して None（MoA フォールバック）になること、かつ
# 無限ループせず打ち切られることを検証する。
_orig_min_votes2 = f.SC_MIN_VOTES
_sc_calls.clear()
_seq5 = (["\\boxed{42}", "\\boxed{42}"]
         + ["すみません、答えが導けませんでした。"] * 60)  # 以降は一切抽出できない


def _fake_sc_ask5(model, messages, temperature, think=None, fmt=None,
                  label=None, num_predict=None, num_ctx=None):
    _sc_calls.append(model)
    idx = len(_sc_calls) - 1
    return _seq5[idx] if idx < len(_seq5) else "すみません、答えが導けませんでした。"


try:
    f.PROPOSERS = ["m1", "m2"]
    f.REASONING_MODELS = ["m1", "m2"]
    f.SC_CHEAP_VOTES = 0
    f.SC_POT = False
    f.ask = _fake_sc_ask5
    _res5 = f.solve_verifiable("test question", "math")
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props2
    f.REASONING_MODELS = _orig_reasoning
    f.SC_CHEAP_VOTES = _orig_cheap
    f.SC_POT = _orig_pot
    f.SC_MIN_VOTES = _orig_min_votes2
check("sc: SC_MAX消化まで無限ループせず打ち切られる(有限回で終了)",
      0 < len(_sc_calls) <= f.SC_MAX + 10)
check("sc: 最終return(SC_MAX消化パス)でも SC_MIN_VOTES 未満の勝者は None(床が効く)",
      _res5 is None)

# 境界値の回帰ガード: SC_MAX 消化パスでも勝者票数がちょうど SC_MIN_VOTES(3) に達して
# いれば床は発火せず、通常どおり dict を返さねばならない（floor の over-fire 防止）。
# 早期break条件（unanimous/majority）はどの中間状態でも満たさないよう票を分散させ、
# 最終的に SC_MAX 消化で抜けた時点で初めて勝者(42)が3票に達するようにしてある。
_orig_min_votes3 = f.SC_MIN_VOTES
_sc_calls.clear()
_seq6 = ["\\boxed{7}", "\\boxed{9}"]                      # batch1: idx0-1 (残り idx2-5 は抽出不能)
_seq6 += ["error"] * 4
_seq6 += ["\\boxed{42}", "\\boxed{42}", "\\boxed{7}", "error"]   # batch2: idx6-9
_seq6 += ["\\boxed{42}", "error", "error", "error"]              # batch3: idx10-13
_seq6 += ["error"] * 4                                             # batch4: idx14-17
_seq6 += ["error"] * 4                                             # batch5: idx18-21
_seq6 += ["error"] * 20                                            # 余裕分


def _fake_sc_ask6(model, messages, temperature, think=None, fmt=None,
                  label=None, num_predict=None, num_ctx=None):
    _sc_calls.append(model)
    idx = len(_sc_calls) - 1
    return _seq6[idx] if idx < len(_seq6) else "error"


try:
    f.PROPOSERS = ["m1", "m2"]
    f.REASONING_MODELS = ["m1", "m2"]
    f.SC_CHEAP_VOTES = 0
    f.SC_POT = False
    f.ask = _fake_sc_ask6
    _res6 = f.solve_verifiable("test question", "math")
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props2
    f.REASONING_MODELS = _orig_reasoning
    f.SC_CHEAP_VOTES = _orig_cheap
    f.SC_POT = _orig_pot
    f.SC_MIN_VOTES = _orig_min_votes3
check("sc: SC_MAX消化パスでも勝者票数=SC_MIN_VOTESなら床は発火しない(通常確定)",
      _res6 is not None and _res6["answer"] == "42")
check("sc: 床が過剰発火していない場合はvotes/n_samplesも通常どおり返る",
      _res6 is not None and _res6["votes"].get("42") == 3 and _res6["n_samples"] == 22)

# ---------- 拮抗時の裁定（_arbitrate）----------
# 上位2クラスが同数で並ぶと _arbitrate が呼ばれる。かつては裁定役の答えだけを採用し、
# 本文(res['text'])は SAMPLE プールから _representative_text で再選出していたため、
# 裁定で数値が変わったり第三の答えに覆ったりすると、本文が敗者側候補の主張のまま
# 残る内部矛盾があった（2026-07-21 発見・修正）。ここでは _arbitrate 自身の解答
# テキストが res['text'] として使われることを検証する。
_orig_installed = f.installed_models
_orig_arbiter_model = f.ARBITER_MODEL


def _fake_installed_m1m2():
    return ["m1", "m2"]


# ケース1: 拮抗 → 裁定役が既存候補の一方(1)を支持。本文は裁定役自身の推論であること
# （敗者候補(2)の本文であってはならない）。
_arb_calls = []


def _fake_ask_arb_pick_existing(model, messages, temperature, think=None, fmt=None,
                                label=None, num_predict=None, num_ctx=None):
    _arb_calls.append((label, model))
    if label == "arbiter":
        return ("ARBITER_REASONING: candidate B miscalculates in step 2; "
                 "re-solving from scratch gives \\boxed{1}")
    idx = len(_arb_calls) - 1
    ans = "1" if idx % 2 == 0 else "2"
    return f"sc reasoning candidate {ans}\n\\boxed{{{ans}}}"


try:
    f.PROPOSERS = ["m1", "m2"]
    f.REASONING_MODELS = ["m1", "m2"]
    f.SC_CHEAP_VOTES = 0
    f.SC_POT = False
    f.ARBITER_MODEL = None
    f.installed_models = _fake_installed_m1m2
    f.ask = _fake_ask_arb_pick_existing
    _res_arb1 = f.solve_verifiable("test question", "math")
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props2
    f.REASONING_MODELS = _orig_reasoning
    f.SC_CHEAP_VOTES = _orig_cheap
    f.SC_POT = _orig_pot
    f.ARBITER_MODEL = _orig_arbiter_model
    f.installed_models = _orig_installed
check("arb: 票が同数で拮抗 → 裁定役が呼ばれる", any(lab == "arbiter" for lab, _m in _arb_calls))
check("arb: 裁定役の答えを採用", _res_arb1 is not None and _res_arb1["answer"] == "1")
check("arb: 本文は裁定役自身の推論(敗者候補の本文ではない)",
      _res_arb1 is not None and "ARBITER_REASONING" in _res_arb1["text"]
      and "sc reasoning candidate" not in _res_arb1["text"])
# 2026-07-22: 裁定役が既存の票クラス('1')をそのまま支持したケースでは votes 辞書は
# 従来どおり(裁定前の同値クラス集計そのまま)で、answer は必ずそのキーとして存在すること。
check("arb: 裁定役が既存候補を採用した場合、votesにanswerがキーとして存在し票数は正しい",
      _res_arb1 is not None and _res_arb1["answer"] in _res_arb1["votes"]
      and _res_arb1["votes"]["1"] == _res_arb1["votes"]["2"])

# ケース2: 裁定役が両候補と異なる第三の答えを提示 → 本文が敗者側候補の主張になって
# はいけない（旧ロジックのバグ: _representative_text が第三の答えと同値のサンプルを
# 見つけられず「最長サンプル」＝どちらかの敗者の本文にフォールバックしていた）。
_arb_calls2 = []


def _fake_ask_arb_new_answer(model, messages, temperature, think=None, fmt=None,
                             label=None, num_predict=None, num_ctx=None):
    _arb_calls2.append((label, model))
    if label == "arbiter":
        return ("ARBITER_REASONING_NEW: both candidates share the same wrong "
                 "assumption; the correct value is \\boxed{3}")
    idx = len(_arb_calls2) - 1
    ans = "1" if idx % 2 == 0 else "2"
    return f"sc reasoning candidate {ans}\n\\boxed{{{ans}}}"


try:
    f.PROPOSERS = ["m1", "m2"]
    f.REASONING_MODELS = ["m1", "m2"]
    f.SC_CHEAP_VOTES = 0
    f.SC_POT = False
    f.ARBITER_MODEL = None
    f.installed_models = _fake_installed_m1m2
    f.ask = _fake_ask_arb_new_answer
    _res_arb2 = f.solve_verifiable("test question", "math")
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props2
    f.REASONING_MODELS = _orig_reasoning
    f.SC_CHEAP_VOTES = _orig_cheap
    f.SC_POT = _orig_pot
    f.ARBITER_MODEL = _orig_arbiter_model
    f.installed_models = _orig_installed
check("arb: 裁定役が出した第三の答えを採用", _res_arb2 is not None and _res_arb2["answer"] == "3")
check("arb: 本文は裁定役の推論(第三の答え。敗者候補の主張ではない)",
      _res_arb2 is not None and "ARBITER_REASONING_NEW" in _res_arb2["text"]
      and "sc reasoning candidate" not in _res_arb2["text"])
# 2026-07-22 本修正の回帰テスト: 裁定役が票の無い第三の答え('3')を採用した場合、
# res['votes'] にその答えが「真の票数(0票)」で必ずキーとして載ること。また旧トップ
# ('1'または'2'、どちらも拮抗していた同数票)の票数が、そのまま裁定結果('3')の
# 票数であるかのように誤って表示されてはならない。
check("arb: 第三の答え採用時、votesにanswerが0票としてキーで存在する",
      _res_arb2 is not None and _res_arb2["votes"].get("3") == 0)
check("arb: 第三の答え採用時、敗者候補('1'/'2')の票数が勝者の票数として流用されていない",
      _res_arb2 is not None and _res_arb2["votes"]["1"] == _res_arb2["votes"]["2"]
      and _res_arb2["votes"]["1"] > 0
      and _res_arb2["votes"]["3"] != _res_arb2["votes"]["1"])

# ケース3 (2026-07-22 本修正): 裁定役が既存の拮抗クラスの一つと数学的に同値だが、
# 書き方だけ異なる文字列を返す（例: クラス代表 '1/2' に対し裁定役は小数表記 '0.5' を
# 提示。分数⇄小数の書き直しは裁定役がよくやる）。旧コードは
# 「match is None or match[0] != top」を一括りに扱っていたため、match が見つかって
# いても新規クラス [top, cnt] を無条件追加してしまい、同じ票が '1/2' と '0.5' の
# 二つのキーに二重計上されていた（sum(votes.values()) が実際の有効票数を超える）。
_arb_calls3b = []


def _fake_ask_arb_equivalent_diff_string(model, messages, temperature, think=None, fmt=None,
                                          label=None, num_predict=None, num_ctx=None):
    _arb_calls3b.append((label, model))
    if label == "arbiter":
        return ("ARBITER_REASONING_EQUIV: both are the same value; the correct "
                 "final answer is \\boxed{0.5}")
    idx = len(_arb_calls3b) - 1
    ans = "1/2" if idx % 2 == 0 else "2"
    return f"sc reasoning candidate {ans}\n\\boxed{{{ans}}}"


try:
    f.PROPOSERS = ["m1", "m2"]
    f.REASONING_MODELS = ["m1", "m2"]
    f.SC_CHEAP_VOTES = 0
    f.SC_POT = False
    f.ARBITER_MODEL = None
    f.installed_models = _fake_installed_m1m2
    f.ask = _fake_ask_arb_equivalent_diff_string
    _res_arb_eq = f.solve_verifiable("test question", "math")
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props2
    f.REASONING_MODELS = _orig_reasoning
    f.SC_CHEAP_VOTES = _orig_cheap
    f.SC_POT = _orig_pot
    f.ARBITER_MODEL = _orig_arbiter_model
    f.installed_models = _orig_installed

_valid_votes_eq = sum(1 for lab, _m in _arb_calls3b if lab != "arbiter")
check("arb-eq: 裁定役が同値・別表記('0.5')を返す → answer はその文字列そのもの",
      _res_arb_eq is not None and _res_arb_eq["answer"] == "0.5")
check("arb-eq: votes に '0.5' が一度だけ載り、旧代表 '1/2' キーは残らない(二重計上なし)",
      _res_arb_eq is not None and _res_arb_eq["votes"].get("0.5") is not None
      and "1/2" not in _res_arb_eq["votes"])
check("arb-eq: '0.5' の票数は '2' クラスと同じ(旧'1/2'の真の票数)で0ではない",
      _res_arb_eq is not None and _res_arb_eq["votes"]["0.5"] == _res_arb_eq["votes"]["2"]
      and _res_arb_eq["votes"]["0.5"] > 0)
check("arb-eq: 票の合計が実際の有効票数と一致する(水増しなし)",
      _res_arb_eq is not None and sum(_res_arb_eq["votes"].values()) == _valid_votes_eq)

# ---------- 2026-07-25 (iteration 169): _arbitrate のrep選出にCoT>PoT優先を適用 ----------
# 従来の _arbitrate は、タイ候補ごとに samples を先頭から走査して canon と
# answers_equivalent な最初のサンプルをそのまま裁定役へ見せていた(単純な最初の一致)。
# add_batch(~L3731-3738)は各バッチの末尾にそのバッチのPoTサンプルを1件追加するため、
# 「バッチ1のPoTサンプル」が「バッチ2以降のCoTサンプル」より samples 内で先に並ぶ。
# その結果、同じ答えに後から自然言語のCoT推論が存在していても、先に並んだPoT実行結果
# (コード＋'[PoT execution output]')の方が「代表解答」として裁定役に渡ってしまい、
# 「各候補の誤りを指摘せよ」という裁定役への指示に対して弱い入力になっていた。
# _representative_text(iteration 2/55, L3594)はユーザー向け代表解答の選出で全く同じ
# 状況をCoT優先(リスト順序に非依存)で既に解決済みで、L3471-3472のコメントは
# 「_arbitrate のrep選出はそれとは別contract」と明記していた。ここではその契約が
# _arbitrate側にも揃ったことを直接ロックする。_arbitrate を直接叩き、f.ask をモックして
# label=='arbiter'呼び出しに渡されたプロンプト文字列を捕捉する(cap テストと同じ手法)。
def _arbrep_sample(answer, text, pot=False, model="m1"):
    return {"answer": answer, "text": text, "pot": pot, "model": model}


_orig_installed_rep = f.installed_models
_orig_arbiter_model_rep = f.ARBITER_MODEL
_orig_reasoning_rep = f.REASONING_MODELS
_orig_props_rep = f.PROPOSERS


def _fake_installed_rep():
    return ["m1"]


def _make_fake_ask_arb_rep(prompts_out, boxed="1"):
    def _fake(model, messages, temperature, think=None, fmt=None,
              label=None, num_predict=None, num_ctx=None):
        if label == "arbiter":
            prompts_out.append(messages[0]["content"])
            return f"ARBITER_REP_TEST \\boxed{{{boxed}}}"
        return f"\\boxed{{{boxed}}}"
    return _fake


# (1) 順序非依存の主眼のテスト: タイ候補'1'について、samples内で先に並ぶのはPoT一致
#     サンプル(コード+'[PoT execution output]')、後から並ぶのがCoT一致サンプル。
#     裁定役に渡るのは後続のCoTテキストであり、先行PoTサンプルの本文ではないこと。
#     もう一方のタイ候補'2'は最初の一致がCoTのままなので、そちらのrepは変化しない。
_rep1_samples = [
    _arbrep_sample("1", "def solve():\n    return 21*2\n\n[PoT execution output]\n42", pot=True),
    _arbrep_sample("2", "COT_REASONING_FOR_2", pot=False),
    _arbrep_sample("1", "COT_REASONING_FOR_1_LATER", pot=False),
]
_rep1_classes = [["1", 2], ["2", 2]]
_rep1_prompts = []
try:
    f.PROPOSERS = ["m1"]
    f.REASONING_MODELS = ["m1"]
    f.ARBITER_MODEL = None
    f.installed_models = _fake_installed_rep
    f.ask = _make_fake_ask_arb_rep(_rep1_prompts)
    _rep1_result = f._arbitrate("test question", "math", _rep1_samples, _rep1_classes)
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props_rep
    f.REASONING_MODELS = _orig_reasoning_rep
    f.ARBITER_MODEL = _orig_arbiter_model_rep
    f.installed_models = _orig_installed_rep

_rep1_prompt = _rep1_prompts[0] if _rep1_prompts else ""
check("arb-rep: 後続CoT一致サンプルのテキストが裁定役プロンプトに含まれる(順序非依存)",
      "COT_REASONING_FOR_1_LATER" in _rep1_prompt)
check("arb-rep: 先行PoT一致サンプルの本文([PoT execution output])はそのタイ候補のrepとして渡らない",
      "[PoT execution output]" not in _rep1_prompt
      and "def solve():" not in _rep1_prompt)
check("arb-rep: 別タイ候補('2')のrepは影響を受けず従来通りCoTテキストのまま",
      "COT_REASONING_FOR_2" in _rep1_prompt)
check("arb-rep: 有効な(answer, text)タプルを返す(既存契約維持)",
      _rep1_result is not None and _rep1_result[0] == "1")

# (2) 回帰テスト: 各タイ候補の最初の一致サンプルが既にCoTである通常ケースでは、
#     裁定役プロンプトに埋め込まれるテキストが変更前(最初の一致をそのまま採用)と
#     バイト単位で同一であること。
_rep2_samples = [
    _arbrep_sample("1", "COT_TEXT_1", pot=False),
    _arbrep_sample("2", "COT_TEXT_2", pot=False),
]
_rep2_classes = [["1", 2], ["2", 2]]
_rep2_prompts = []
try:
    f.PROPOSERS = ["m1"]
    f.REASONING_MODELS = ["m1"]
    f.ARBITER_MODEL = None
    f.installed_models = _fake_installed_rep
    f.ask = _make_fake_ask_arb_rep(_rep2_prompts)
    _rep2_result = f._arbitrate("test question", "math", _rep2_samples, _rep2_classes)
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props_rep
    f.REASONING_MODELS = _orig_reasoning_rep
    f.ARBITER_MODEL = _orig_arbiter_model_rep
    f.installed_models = _orig_installed_rep

_rep2_prompt = _rep2_prompts[0] if _rep2_prompts else ""
_rep2_expected_listing = (
    "### Candidate A (final answer: 1)\nCOT_TEXT_1\n\n"
    "### Candidate B (final answer: 2)\nCOT_TEXT_2")
check("arb-rep: 各候補の最初の一致が既にCoTの通常ケースは変更前と同一のlistingになる(回帰ピン)",
      _rep2_expected_listing in _rep2_prompt)
check("arb-rep: 通常ケースでも有効な(answer, text)タプルを返す",
      _rep2_result is not None and _rep2_result[0] == "1")

# (3) PoT-onlyフォールバック回帰: タイ候補の一致サンプルがPoTしか無い場合、
#     そのPoTサンプルのテキストがそのまま提示され、候補が脱落したり空文字列に
#     なったりしないこと。
_rep3_samples = [
    _arbrep_sample("1", "POT_ONLY_TEXT_FOR_1", pot=True),
    _arbrep_sample("2", "COT_TEXT_FOR_2", pot=False),
]
_rep3_classes = [["1", 2], ["2", 2]]
_rep3_prompts = []
try:
    f.PROPOSERS = ["m1"]
    f.REASONING_MODELS = ["m1"]
    f.ARBITER_MODEL = None
    f.installed_models = _fake_installed_rep
    f.ask = _make_fake_ask_arb_rep(_rep3_prompts)
    _rep3_result = f._arbitrate("test question", "math", _rep3_samples, _rep3_classes)
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props_rep
    f.REASONING_MODELS = _orig_reasoning_rep
    f.ARBITER_MODEL = _orig_arbiter_model_rep
    f.installed_models = _orig_installed_rep

_rep3_prompt = _rep3_prompts[0] if _rep3_prompts else ""
check("arb-rep: CoT一致が無い候補はPoT一致サンプルのテキストがそのまま提示される(脱落しない)",
      "POT_ONLY_TEXT_FOR_1" in _rep3_prompt and "final answer: 1" in _rep3_prompt)
check("arb-rep: PoT-onlyフォールバック時も有効な(answer, text)タプルを返す",
      _rep3_result is not None)

# (4) 同値判定ガード: タイ候補'0.5'に一致するのは文字列としては異なる'1/2'の
#     サンプルのみ(answers_equivalent経由でFraction高速パス一致)。無関係な
#     敗者候補('7')のテキストが紛れ込まないこと、別タイ候補('3')は完全一致のまま
#     であることも合わせて確認する。
_rep4_samples = [
    _arbrep_sample("7", "LOSER_TEXT_7", pot=False),
    _arbrep_sample("1/2", "EQUIV_MATCH_TEXT_FOR_HALF", pot=False),
    _arbrep_sample("3", "OTHER_TIED_TEXT_3", pot=False),
]
_rep4_classes = [["0.5", 2], ["3", 2]]
_rep4_prompts = []
try:
    f.PROPOSERS = ["m1"]
    f.REASONING_MODELS = ["m1"]
    f.ARBITER_MODEL = None
    f.installed_models = _fake_installed_rep
    f.ask = _make_fake_ask_arb_rep(_rep4_prompts)
    _rep4_result = f._arbitrate("test question", "math", _rep4_samples, _rep4_classes)
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props_rep
    f.REASONING_MODELS = _orig_reasoning_rep
    f.ARBITER_MODEL = _orig_arbiter_model_rep
    f.installed_models = _orig_installed_rep

_rep4_prompt = _rep4_prompts[0] if _rep4_prompts else ""
check("arb-rep: 候補'0.5'のrepは文字列が異なる同値サンプル('1/2')から選ばれる(answers_equivalent経由)",
      "EQUIV_MATCH_TEXT_FOR_HALF" in _rep4_prompt and "final answer: 0.5" in _rep4_prompt)
check("arb-rep: 無関係な敗者候補('7')のテキストは紛れ込まない",
      "LOSER_TEXT_7" not in _rep4_prompt)
check("arb-rep: 別タイ候補('3')は完全一致のテキストのまま",
      "OTHER_TIED_TEXT_3" in _rep4_prompt)

# ---------- 2026-07-22: N択拮抗で _arbitrate が classes[:2] に打ち切らないこと ----------
# 上位2クラスが同数(classes[0][1]==classes[1][1])で拮抗判定が発火する状況は、実際には
# 3クラス以上が同数タイになるケース(例: 票数 [k,k,k])も含みうる。従来の _arbitrate は
# classes[:2] で常に先頭2クラスしか裁定役に見せておらず、3番目以降の同数クラス
# （それが正解かもしれない）が黙って握りつぶされていた。ここでは 3モデル×均等分配で
# 常に3クラスが同数のまま SC_MAX まで積み上がる状況を作り、裁定役に3候補全てが
# 提示されること、かつ裁定役が(既存候補の一つである)第3クラスの答えを採用した場合に
# 事後の票数再集計(recount)が正しく合成されることを検証する。
_orig_installed3 = f.installed_models
_orig_arbiter_model3 = f.ARBITER_MODEL
_arb3_prompts = []
_sc3_idx = [0]


def _fake_installed_m1m2m3():
    return ["m1", "m2", "m3"]


def _fake_ask_arb_3way(model, messages, temperature, think=None, fmt=None,
                       label=None, num_predict=None, num_ctx=None):
    if label == "arbiter":
        _arb3_prompts.append(messages[0]["content"])
        # 裁定役は「票の無い新答え」ではなく、拮抗していた3クラスのうち3番目(='3')を支持する。
        return "ARBITER_REASONING_3WAY: candidate C is correct after re-derivation \\boxed{3}"
    idx = _sc3_idx[0]
    _sc3_idx[0] += 1
    ans = str((idx % 3) + 1)   # '1','2','3' を均等に繰り返す → 常に3クラス同数タイ
    return f"sc reasoning candidate {ans}\n\\boxed{{{ans}}}"


try:
    f.PROPOSERS = ["m1", "m2", "m3"]
    f.REASONING_MODELS = ["m1", "m2", "m3"]
    f.SC_CHEAP_VOTES = 0
    f.SC_POT = False
    f.ARBITER_MODEL = None
    f.installed_models = _fake_installed_m1m2m3
    f.ask = _fake_ask_arb_3way
    _res_arb3 = f.solve_verifiable("test question", "math")
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props2
    f.REASONING_MODELS = _orig_reasoning
    f.SC_CHEAP_VOTES = _orig_cheap
    f.SC_POT = _orig_pot
    f.ARBITER_MODEL = _orig_arbiter_model3
    f.installed_models = _orig_installed3

check("arb3: 3択で拮抗 → 裁定役が呼ばれる", len(_arb3_prompts) > 0)
_arb3_prompt = _arb3_prompts[0] if _arb3_prompts else ""
check("arb3: 裁定役プロンプトに候補A(final answer: 1)が含まれる",
      "Candidate A" in _arb3_prompt and "final answer: 1" in _arb3_prompt)
check("arb3: 裁定役プロンプトに候補B(final answer: 2)が含まれる",
      "Candidate B" in _arb3_prompt and "final answer: 2" in _arb3_prompt)
check("arb3: 裁定役プロンプトに候補C(final answer: 3)が含まれる(3番目のタイ候補が握りつぶされていない)",
      "Candidate C" in _arb3_prompt and "final answer: 3" in _arb3_prompt)
check("arb3: プロンプト文言が候補数に応じている(count-agnostic)",
      "3 candidate solutions disagree" in _arb3_prompt
      and "Two candidate solutions disagree" not in _arb3_prompt)
check("arb3: 裁定役が採用した第3候補('3')がresの答えになる",
      _res_arb3 is not None and _res_arb3["answer"] == "3")
check("arb3: votes再集計により'3'の真の票数(既存タイ候補の実票数)が反映される",
      _res_arb3 is not None and _res_arb3["votes"].get("3") is not None
      and _res_arb3["votes"]["3"] == _res_arb3["votes"]["1"] == _res_arb3["votes"]["2"]
      and _res_arb3["votes"]["3"] > 0)

# ---------- 2026-07-22: ARBITRATE_MAX_CANDIDATES による上限保護(病的な多択タイ) ----------
# num_ctx(8192/16384に固定)を溢れさせないよう、同数タイが上限を超える場合は上限件数
# のみ提示し、超過分は黙って捨てずログに出す。_arbitrate を直接叩いて検証する
# （solve_verifiable 経由でここまで多くの均等クラスを作るのは非現実的なため）。
_cap_samples = [{"answer": str(i), "text": f"reasoning for {i}", "model": "m1", "pot": False}
                for i in range(1, 6)]                       # 1..5 の5クラス、全て同数(2票)
_cap_classes = [[str(i), 2] for i in range(1, 6)]
_orig_installed_cap = f.installed_models
_orig_arbiter_model_cap = f.ARBITER_MODEL
_orig_reasoning_cap = f.REASONING_MODELS
_orig_props_cap = f.PROPOSERS
_cap_prompts = []


def _fake_ask_cap(model, messages, temperature, think=None, fmt=None,
                  label=None, num_predict=None, num_ctx=None):
    if label == "arbiter":
        _cap_prompts.append(messages[0]["content"])
        return "ARBITER_REASONING_CAP \\boxed{1}"
    return "\\boxed{1}"


_cap_stdout = io.StringIO()
try:
    f.PROPOSERS = ["m1"]
    f.REASONING_MODELS = ["m1"]
    f.ARBITER_MODEL = None
    f.installed_models = lambda: ["m1"]
    f.ask = _fake_ask_cap
    with contextlib.redirect_stdout(_cap_stdout):
        _cap_result = f._arbitrate("test question", "math", _cap_samples, _cap_classes)
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props_cap
    f.REASONING_MODELS = _orig_reasoning_cap
    f.ARBITER_MODEL = _orig_arbiter_model_cap
    f.installed_models = _orig_installed_cap

_cap_prompt = _cap_prompts[0] if _cap_prompts else ""
check("cap: 5択タイでも上限(ARBITRATE_MAX_CANDIDATES=4)件しか提示されない",
      sum(f"final answer: {i}" in _cap_prompt for i in range(1, 6)) == f.ARBITRATE_MAX_CANDIDATES)
check("cap: 上限超過分(5)は黙って捨てずログに出力される",
      "5" in _cap_stdout.getvalue() and "提示されません" in _cap_stdout.getvalue())
check("cap: 上限保護時も有効な(answer, text)タプルを返す",
      _cap_result is not None and _cap_result[0] == "1")

# ---------- 2026-07-22: _arbitrate が ask() の __ERROR__ センチネルを誤って
# 数値解答として採用しないこと ----------
# ask() は失敗時 '__ERROR__: HTTP Error 500: Internal Server Error {...}' のような
# 文字列を返す（line ~1079）。旧 _arbitrate はこれをチェックせず strip_think →
# extract_final_answer に渡していたため、math タスクの最終数値フォールバック
# （line ~2299, `nums = re.findall(...)`）がエラーメッセージ中の '500'/'429' を
# 「裁定役の最終解答」として誤採用し、拮抗投票がでっち上げの自信満々な数値に化けて
# いた（_sc_sample=iter4, ask()自体=iter9, _critic_judge/second_opinion=iter15 で
# 直した同種バグの兄弟ケース）。ここでは (a) solve_verifiable 経由で拮抗した全裁定役が
# エラーになっても最終結果にエラー文字列/誤答が漏れないこと、(b) _arbitrate を直接叩いて
# チェーンの先頭がエラーでも次の裁定役へフォールバックすること、(c) 全裁定役がエラーなら
# _arbitrate が None を返すこと、の3点を検証する。

# (a) solve_verifiable レベル: ARBITER_MODEL 無し・REASONING_MODELS=PROPOSERS=[m1,m2] の
# 従来ケース1と同じ拮抗を作り、裁定役(m1もm2も)が毎回 __ERROR__ を返す状況。
_arb_err_calls = []


def _fake_ask_arb_error_only(model, messages, temperature, think=None, fmt=None,
                              label=None, num_predict=None, num_ctx=None):
    _arb_err_calls.append((label, model))
    if label == "arbiter":
        return "__ERROR__: HTTP Error 500: Internal Server Error"
    idx = len(_arb_err_calls) - 1
    ans = "1" if idx % 2 == 0 else "2"
    return f"sc reasoning candidate {ans}\n\\boxed{{{ans}}}"


try:
    f.PROPOSERS = ["m1", "m2"]
    f.REASONING_MODELS = ["m1", "m2"]
    f.SC_CHEAP_VOTES = 0
    f.SC_POT = False
    f.ARBITER_MODEL = None
    f.installed_models = _fake_installed_m1m2
    f.ask = _fake_ask_arb_error_only
    _res_arb_err = f.solve_verifiable("test question", "math")
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props2
    f.REASONING_MODELS = _orig_reasoning
    f.SC_CHEAP_VOTES = _orig_cheap
    f.SC_POT = _orig_pot
    f.ARBITER_MODEL = _orig_arbiter_model
    f.installed_models = _orig_installed
check("arb-err: 拮抗した全裁定役が__ERROR__を返してもanswerに'500'が誤採用されない",
      _res_arb_err is None or _res_arb_err["answer"] != "500")
check("arb-err: 拮抗した全裁定役が__ERROR__を返してもtextにエラー文字列が漏れない",
      _res_arb_err is None or "__ERROR__" not in _res_arb_err["text"])

# (b) _arbitrate 直接: ARBITER_MODEL(裁定役1番手)が__ERROR__、REASONING_MODELS の
# フォールバック(2番手)が有効な \boxed 解答を返す → チェーンを進めてその有効解答を
# 採用すること（エラーで止まって None になったり、エラー文中の数値を拾ったりしない）。
_orig_installed_e2 = f.installed_models
_orig_arbiter_model_e2 = f.ARBITER_MODEL
_orig_reasoning_e2 = f.REASONING_MODELS
_orig_props_e2 = f.PROPOSERS


def _fake_ask_arb_chain_fallback(model, messages, temperature, think=None, fmt=None,
                                  label=None, num_predict=None, num_ctx=None):
    assert label == "arbiter"
    if model == "arb_big":
        return "__ERROR__: HTTP Error 500: Internal Server Error"
    return "ARBITER_REASONING_FALLBACK: re-derived correctly \\boxed{7}"


_fb_samples = [{"answer": "1", "text": "reasoning for 1", "model": "m1", "pot": False},
               {"answer": "2", "text": "reasoning for 2", "model": "m1", "pot": False}]
_fb_classes = [["1", 2], ["2", 2]]
try:
    f.PROPOSERS = ["m1"]
    f.REASONING_MODELS = ["arb_big", "m1"]
    f.ARBITER_MODEL = "arb_big"
    f.installed_models = lambda: ["arb_big", "m1"]
    f.ask = _fake_ask_arb_chain_fallback
    _e2_result = f._arbitrate("test question", "math", _fb_samples, _fb_classes)
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props_e2
    f.REASONING_MODELS = _orig_reasoning_e2
    f.ARBITER_MODEL = _orig_arbiter_model_e2
    f.installed_models = _orig_installed_e2
check("arb-err: 先頭裁定役が__ERROR__ → 次の裁定役の有効な\\boxed解答へフォールバック",
      _e2_result is not None and _e2_result[0] == "7")
check("arb-err: フォールバック採用時の本文は次裁定役自身の推論(エラー文ではない)",
      _e2_result is not None and "ARBITER_REASONING_FALLBACK" in _e2_result[1]
      and "__ERROR__" not in _e2_result[1])

# (c) _arbitrate 直接: チェーン全員が__ERROR__(しかも数字入り)を返す → None を返し、
# 誤った数値タプルをでっち上げないこと。
_orig_installed_e3 = f.installed_models
_orig_arbiter_model_e3 = f.ARBITER_MODEL
_orig_reasoning_e3 = f.REASONING_MODELS
_orig_props_e3 = f.PROPOSERS


def _fake_ask_arb_all_error(model, messages, temperature, think=None, fmt=None,
                             label=None, num_predict=None, num_ctx=None):
    assert label == "arbiter"
    if model == "arb_big":
        return "__ERROR__: HTTP Error 500: Internal Server Error"
    return "__ERROR__: HTTP Error 429: Too Many Requests"


try:
    f.PROPOSERS = ["m1"]
    f.REASONING_MODELS = ["arb_big", "m1"]
    f.ARBITER_MODEL = "arb_big"
    f.installed_models = lambda: ["arb_big", "m1"]
    f.ask = _fake_ask_arb_all_error
    _e3_result = f._arbitrate("test question", "math", _fb_samples, _fb_classes)
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props_e3
    f.REASONING_MODELS = _orig_reasoning_e3
    f.ARBITER_MODEL = _orig_arbiter_model_e3
    f.installed_models = _orig_installed_e3
check("arb-err: 全裁定役が__ERROR__ → _arbitrate は None を返す(数値をでっち上げない)",
      _e3_result is None)

# ---------- 2026-07-22: _arbitrate プロンプト本文の候補数非依存化(iteration 16の続き) ----------
# iteration 16 はヘッダー行("{len(reps)} candidate solutions disagree:")を候補数に
# 応じた表現に直したが、本文の指示文 "Carefully check both, find the flaw in the wrong
# one" は2択決め打ちのまま残っていた。3択/4択タイでは候補が3-4件出るのに「both」
# 「the wrong one」(単数)と言われ、3件目以降の精査が軽視されるリスクがある。ここでは
# _arbitrate を直接叩き、(a) 2択では従来と同じ意図(両候補を精査し\boxed{}で単一の正解を
# 出す)を保ったまま候補数非依存の文言になっていること、(b) 3択では "check both" や
# 単数形の「the wrong one」という誤解を招く表現が出ないこと、(c) いずれもヘッダー行と
# \boxed{} 指示は変更されておらず、有効な(answer, text)タプル/Noneの契約を保つこと、を検証する。
_orig_installed_body2 = f.installed_models
_orig_arbiter_model_body2 = f.ARBITER_MODEL
_orig_reasoning_body2 = f.REASONING_MODELS
_orig_props_body2 = f.PROPOSERS
_body2_prompts = []


def _fake_ask_body2(model, messages, temperature, think=None, fmt=None,
                     label=None, num_predict=None, num_ctx=None):
    if label == "arbiter":
        _body2_prompts.append(messages[0]["content"])
        return "ARBITER_REASONING_BODY2 \\boxed{1}"
    return "\\boxed{1}"


_body2_samples = [{"answer": "1", "text": "reasoning for 1", "model": "m1", "pot": False},
                  {"answer": "2", "text": "reasoning for 2", "model": "m1", "pot": False}]
_body2_classes = [["1", 2], ["2", 2]]
try:
    f.PROPOSERS = ["m1"]
    f.REASONING_MODELS = ["m1"]
    f.ARBITER_MODEL = None
    f.installed_models = lambda: ["m1"]
    f.ask = _fake_ask_body2
    _body2_result = f._arbitrate("test question", "math", _body2_samples, _body2_classes)
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props_body2
    f.REASONING_MODELS = _orig_reasoning_body2
    f.ARBITER_MODEL = _orig_arbiter_model_body2
    f.installed_models = _orig_installed_body2

_body2_prompt = _body2_prompts[0] if _body2_prompts else ""
check("arb-body2: 2択のヘッダー行は従来通り('2 candidate solutions disagree:')",
      "2 candidate solutions disagree:" in _body2_prompt)
check("arb-body2: 2択でも本文は各候補の精査を指示している(same intent as 'check both')",
      "each candidate" in _body2_prompt)
check("arb-body2: \\boxed{} による単一最終解答の指示は維持されている",
      "\\boxed{}" in _body2_prompt)
check("arb-body2: 有効な(answer, text)タプルを返す",
      _body2_result is not None and _body2_result[0] == "1")

_orig_installed_body3 = f.installed_models
_orig_arbiter_model_body3 = f.ARBITER_MODEL
_orig_reasoning_body3 = f.REASONING_MODELS
_orig_props_body3 = f.PROPOSERS
_body3_prompts = []


def _fake_ask_body3(model, messages, temperature, think=None, fmt=None,
                     label=None, num_predict=None, num_ctx=None):
    if label == "arbiter":
        _body3_prompts.append(messages[0]["content"])
        return "ARBITER_REASONING_BODY3 \\boxed{1}"
    return "\\boxed{1}"


_body3_samples = [{"answer": str(i), "text": f"reasoning for {i}", "model": "m1", "pot": False}
                  for i in (1, 2, 3)]
_body3_classes = [["1", 2], ["2", 2], ["3", 2]]
try:
    f.PROPOSERS = ["m1"]
    f.REASONING_MODELS = ["m1"]
    f.ARBITER_MODEL = None
    f.installed_models = lambda: ["m1"]
    f.ask = _fake_ask_body3
    _body3_result = f._arbitrate("test question", "math", _body3_samples, _body3_classes)
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props_body3
    f.REASONING_MODELS = _orig_reasoning_body3
    f.ARBITER_MODEL = _orig_arbiter_model_body3
    f.installed_models = _orig_installed_body3

_body3_prompt = _body3_prompts[0] if _body3_prompts else ""
check("arb-body3: 3択のヘッダー行は候補数に応じている('3 candidate solutions disagree:')",
      "3 candidate solutions disagree:" in _body3_prompt)
check("arb-body3: 3択の本文に2択決め打ちの'check both'が残っていない",
      "check both" not in _body3_prompt.lower())
check("arb-body3: 3択の本文に単数形決め打ちの'the wrong one'が残っていない",
      "the wrong one" not in _body3_prompt.lower())
check("arb-body3: 3択の本文は候補数非依存の表現('each candidate'/'incorrect one(s)')になっている",
      "each candidate" in _body3_prompt and "incorrect one" in _body3_prompt)
check("arb-body3: \\boxed{} による単一最終解答の指示は維持されている",
      "\\boxed{}" in _body3_prompt)
check("arb-body3: 有効な(answer, text)タプルを返す",
      _body3_result is not None and _body3_result[0] == "1")

# ---------- 2026-07-24: _arbitrate のプロンプト末尾指示を task_type で分岐 ----------
# _arbitrate は solve_verifiable から math/mcq 両方の拮抗解消で共用されるが、末尾の
# 出力形式指示はこれまで math 前提の "put ONLY the correct final answer in \boxed{}"
# 一本だった。extract_final_answer(text, 'mcq')（iter 3 で確立、iter 26/102 で誤爆
# 修正済み）は \boxed{} の中身が選択肢文字 A-E で始まる場合のみ採用するため、math
# 寄りの文言に引きずられた裁定役が計算値や選択肢本文（\boxed{7} や \boxed{Paris} 等）を
# 箱に入れると mcq の拮抗解決チェーンが丸ごと失敗し、MoA フォールバックへ黙って
# 劣化していた（iter 16/45 は math 側の文言のみ調整で、この mcq 側の齟齬は未対応
# だった）。ここでは (a) math/その他のプロンプトが従来と完全に同一(バイト単位)で
# あること、(b) mcq のプロンプトが選択肢文字を要求する新しい文言になっていて、かつ
# math 向けの汎用文言が残っていないこと、(c) 両 task_type とも \boxed{} からの
# 抽出・(answer, text) 契約が従来通り機能すること、を検証する。
_MATH_FINAL_INSTRUCTION = (
    "Carefully check each candidate, find the flaw(s) in the incorrect one(s), "
    "and solve the problem yourself if needed. At the very end, put ONLY the "
    "correct final answer in \\boxed{}.")

_orig_installed_tt = f.installed_models
_orig_arbiter_model_tt = f.ARBITER_MODEL
_orig_reasoning_tt = f.REASONING_MODELS
_orig_props_tt = f.PROPOSERS

# (a) math: プロンプト全文がバイト単位で不変であること（回帰ピン）。
#     裁定役の \boxed{7} は従来通り '7' を返すことも合わせて確認する。
_tt_math_prompts = []
_tt_samples = [{"answer": "1", "text": "reasoning for 1", "model": "m1", "pot": False},
               {"answer": "2", "text": "reasoning for 2", "model": "m1", "pot": False}]
_tt_classes = [["1", 2], ["2", 2]]


def _fake_ask_tt_math(model, messages, temperature, think=None, fmt=None,
                       label=None, num_predict=None, num_ctx=None):
    if label == "arbiter":
        _tt_math_prompts.append(messages[0]["content"])
        return "ARBITER_REASONING_TT_MATH: re-derived correctly \\boxed{7}"
    return "\\boxed{1}"


try:
    f.PROPOSERS = ["m1"]
    f.REASONING_MODELS = ["m1"]
    f.ARBITER_MODEL = None
    f.installed_models = lambda: ["m1"]
    f.ask = _fake_ask_tt_math
    _tt_math_result = f._arbitrate("test question", "math", _tt_samples, _tt_classes)
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props_tt
    f.REASONING_MODELS = _orig_reasoning_tt
    f.ARBITER_MODEL = _orig_arbiter_model_tt
    f.installed_models = _orig_installed_tt

_tt_math_prompt = _tt_math_prompts[0] if _tt_math_prompts else ""
check("arb-tt: mathプロンプトの末尾指示はバイト単位で従来と同一(回帰ピン)",
      _tt_math_prompt.endswith(_MATH_FINAL_INSTRUCTION))
check("arb-tt: math裁定役の\\boxed{7}は従来通り'7'を返す((answer, text)契約維持)",
      _tt_math_result is not None and _tt_math_result[0] == "7"
      and "ARBITER_REASONING_TT_MATH" in _tt_math_result[1])

# (b) mcq: プロンプトが選択肢文字を要求する新しい文言になっていること。
#     裁定役の \boxed{B} は 'B' として抽出されることも合わせて確認する。
_tt_mcq_prompts = []
_tt_mcq_samples = [{"answer": "A", "text": "reasoning for A", "model": "m1", "pot": False},
                   {"answer": "C", "text": "reasoning for C", "model": "m1", "pot": False}]
_tt_mcq_classes = [["A", 2], ["C", 2]]


def _fake_ask_tt_mcq(model, messages, temperature, think=None, fmt=None,
                      label=None, num_predict=None, num_ctx=None):
    if label == "arbiter":
        _tt_mcq_prompts.append(messages[0]["content"])
        return "ARBITER_REASONING_TT_MCQ: candidate C misreads the question \\boxed{B}"
    return "\\boxed{A}"


try:
    f.PROPOSERS = ["m1"]
    f.REASONING_MODELS = ["m1"]
    f.ARBITER_MODEL = None
    f.installed_models = lambda: ["m1"]
    f.ask = _fake_ask_tt_mcq
    _tt_mcq_result = f._arbitrate("test mcq question", "mcq", _tt_mcq_samples, _tt_mcq_classes)
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props_tt
    f.REASONING_MODELS = _orig_reasoning_tt
    f.ARBITER_MODEL = _orig_arbiter_model_tt
    f.installed_models = _orig_installed_tt

_tt_mcq_prompt = _tt_mcq_prompts[0] if _tt_mcq_prompts else ""
check("arb-tt: mcqプロンプトは選択肢文字(A-E)のみをboxedに入れるよう明示指示している",
      "single choice letter" in _tt_mcq_prompt and "A, B, C, D, or E" in _tt_mcq_prompt
      and "\\boxed{}" in _tt_mcq_prompt)
check("arb-tt: mcqプロンプトにmath向けの汎用文言('the correct final answer')が残っていない",
      "the correct final answer" not in _tt_mcq_prompt)
check("arb-tt: mcqプロンプトのヘッダー行/候補見出しは変更されていない(既存契約維持)",
      "2 candidate solutions disagree:" in _tt_mcq_prompt and "Candidate A" in _tt_mcq_prompt
      and "Candidate B" in _tt_mcq_prompt)
check("arb-tt: mcq裁定役の\\boxed{B}は従来通り'B'を返す((answer, text)契約維持)",
      _tt_mcq_result is not None and _tt_mcq_result[0] == "B"
      and "ARBITER_REASONING_TT_MCQ" in _tt_mcq_result[1])

# ---------- 2026-07-24: solve_verifiable レベルでの新挙動保証（mcqタイの裁定成功）----------
# 修正前は math 前提の文言に引きずられた裁定役が \boxed{選択肢本文/計算値} を返して
# extract_final_answer(text, 'mcq') が None を返し続け、_arbitrate 全体が None になって
# mcq の拮抗が黙って MoA フォールバックへ劣化しうる状況だった。ここでは 2-2 の mcq
# 拮抗を強制発生させ、裁定役が素の選択肢文字(\boxed{A})を返すケースで、None への
# 劣化ではなくその文字へ解決すること・votes が水増しなく truthful であることを検証する。
_orig_installed_mcqtie = f.installed_models
_orig_arbiter_model_mcqtie = f.ARBITER_MODEL
_orig_reasoning_mcqtie = f.REASONING_MODELS
_orig_props_mcqtie = f.PROPOSERS
_orig_cheap_mcqtie = f.SC_CHEAP_VOTES
_orig_pot_mcqtie = f.SC_POT
_arb_mcqtie_calls = []


def _fake_ask_mcq_tie(model, messages, temperature, think=None, fmt=None,
                       label=None, num_predict=None, num_ctx=None):
    _arb_mcqtie_calls.append((label, model))
    if label == "arbiter":
        return ("ARBITER_REASONING_MCQTIE: candidate C misreads the question; "
                 "the correct choice is \\boxed{A}")
    idx = len(_arb_mcqtie_calls) - 1
    ans = "A" if idx % 2 == 0 else "C"
    return f"sc reasoning candidate {ans}\n\\boxed{{{ans}}}"


try:
    f.PROPOSERS = ["m1", "m2"]
    f.REASONING_MODELS = ["m1", "m2"]
    f.SC_CHEAP_VOTES = 0
    f.SC_POT = False
    f.ARBITER_MODEL = None
    f.installed_models = _fake_installed_m1m2
    f.ask = _fake_ask_mcq_tie
    _res_mcqtie = f.solve_verifiable("test mcq question", "mcq")
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props_mcqtie
    f.REASONING_MODELS = _orig_reasoning_mcqtie
    f.SC_CHEAP_VOTES = _orig_cheap_mcqtie
    f.SC_POT = _orig_pot_mcqtie
    f.ARBITER_MODEL = _orig_arbiter_model_mcqtie
    f.installed_models = _orig_installed_mcqtie

_valid_votes_mcqtie = sum(1 for lab, _m in _arb_mcqtie_calls if lab != "arbiter")
check("arb-mcqtie: mcqの拮抗で裁定役が呼ばれる",
      any(lab == "arbiter" for lab, _m in _arb_mcqtie_calls))
check("arb-mcqtie: mcqの拮抗が裁定で解決しMoAフォールバック(None)へ劣化しない",
      _res_mcqtie is not None)
check("arb-mcqtie: 裁定役が採用した文字'A'がそのままanswerになる",
      _res_mcqtie is not None and _res_mcqtie["answer"] == "A")
check("arb-mcqtie: votesが truthful('A'/'C'が実票数でキー登録、水増しなし)",
      _res_mcqtie is not None and _res_mcqtie["votes"].get("A") is not None
      and _res_mcqtie["votes"].get("C") is not None
      and _res_mcqtie["votes"]["A"] == _res_mcqtie["votes"]["C"] > 0
      and sum(_res_mcqtie["votes"].values()) == _valid_votes_mcqtie)

# ---------- 2026-07-24: mcq拮抗+裁定+票再集計(iteration17)の複合経路カバレッジ ----------
# 上の arb-mcqtie は「裁定役が拮抗クラスの一つをそのまま採用」する経路のみを見ており、
# iteration 17（math向け）が固めた「裁定役の答えを既存タイクラスへ answers_equivalent で
# 同値照合してから votes へ合成する」ロジックが mcq の選択肢文字（大文字/小文字/全角）
# 表記ゆれと組み合わさったときに二重計上や欠落を起こさないかは未検証だった。
# extract_final_answer(text, 'mcq') は boxed 内の全角/小文字文字も正規化してASCII大文字
# 1文字にしてから返す（iter 13の_FW_TRANS + .upper()）ため、裁定役が全角小文字'ａ'を
# \boxed{}に入れても最終的な top は既存タイクラス'A'と文字列として完全一致するはずだが、
# 「一致はするが二重キーにならない」契約を solve_verifiable レベルで直接ロックする
# （_arbitrate単体テストでは votes 合成後の姿までは見えないため）。
_orig_installed_mcqmerge = f.installed_models
_orig_arbiter_model_mcqmerge = f.ARBITER_MODEL
_orig_reasoning_mcqmerge = f.REASONING_MODELS
_orig_props_mcqmerge = f.PROPOSERS
_orig_cheap_mcqmerge = f.SC_CHEAP_VOTES
_orig_pot_mcqmerge = f.SC_POT
_arb_mcqmerge_calls = []


def _fake_ask_mcq_merge(model, messages, temperature, think=None, fmt=None,
                         label=None, num_predict=None, num_ctx=None):
    _arb_mcqmerge_calls.append((label, model))
    if label == "arbiter":
        # 全角小文字'ａ'(U+FF41) → normalize_answer で 'a' → .upper() で 'A' に正規化される。
        # 既存のタイクラス'A'と文字列として完全一致するかどうかを検証する。
        return "ARBITER_REASONING_MCQMERGE: candidate C misreads it; correct is \\boxed{ａ}"
    idx = len(_arb_mcqmerge_calls) - 1
    ans = "A" if idx % 2 == 0 else "C"
    return f"sc reasoning candidate {ans}\n\\boxed{{{ans}}}"


try:
    f.PROPOSERS = ["m1", "m2"]
    f.REASONING_MODELS = ["m1", "m2"]
    f.SC_CHEAP_VOTES = 0
    f.SC_POT = False
    f.ARBITER_MODEL = None
    f.installed_models = _fake_installed_m1m2
    f.ask = _fake_ask_mcq_merge
    _res_mcqmerge = f.solve_verifiable("test mcq question", "mcq")
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props_mcqmerge
    f.REASONING_MODELS = _orig_reasoning_mcqmerge
    f.SC_CHEAP_VOTES = _orig_cheap_mcqmerge
    f.SC_POT = _orig_pot_mcqmerge
    f.ARBITER_MODEL = _orig_arbiter_model_mcqmerge
    f.installed_models = _orig_installed_mcqmerge

_valid_votes_mcqmerge = sum(1 for lab, _m in _arb_mcqmerge_calls if lab != "arbiter")
check("arb-mcqmerge: 裁定役の全角小文字'ａ'が正規化されタイクラス'A'に解決される",
      _res_mcqmerge is not None and _res_mcqmerge["answer"] == "A")
check("arb-mcqmerge: votesに'a'/'ａ'等の重複キーが作られない(正規化キー'A'のみ)",
      _res_mcqmerge is not None
      and set(_res_mcqmerge["votes"].keys()) == {"A", "C"})
check("arb-mcqmerge: votesが truthful('A'/'C'が実票数でキー登録、水増しなし)",
      _res_mcqmerge is not None and _res_mcqmerge["votes"].get("A") is not None
      and _res_mcqmerge["votes"].get("C") is not None
      and _res_mcqmerge["votes"]["A"] == _res_mcqmerge["votes"]["C"] > 0
      and sum(_res_mcqmerge["votes"].values()) == _valid_votes_mcqmerge)

# ---------- 2026-07-24: mcq拮抗+裁定役が既存タイクラスに無い新規文字を採用するケース ----------
# iteration 22(math向け)が固めた「裁定役の答えが既存タイクラスと無関係な第三の答えなら
# 0票の新規クラスとしてvotesに追加する（旧トップの票数を誤流用しない）」契約を、mcqの
# 選択肢文字でも直接ロックする。タイは'A'/'C'（2択）だが裁定役は両方とも誤りとして
# 第三の選択肢'E'を採用する。
_orig_installed_mcqnew = f.installed_models
_orig_arbiter_model_mcqnew = f.ARBITER_MODEL
_orig_reasoning_mcqnew = f.REASONING_MODELS
_orig_props_mcqnew = f.PROPOSERS
_orig_cheap_mcqnew = f.SC_CHEAP_VOTES
_orig_pot_mcqnew = f.SC_POT
_arb_mcqnew_calls = []


def _fake_ask_mcq_new(model, messages, temperature, think=None, fmt=None,
                       label=None, num_predict=None, num_ctx=None):
    _arb_mcqnew_calls.append((label, model))
    if label == "arbiter":
        return ("ARBITER_REASONING_MCQNEW: both candidates A and C misread the question; "
                 "the correct choice is \\boxed{E}")
    idx = len(_arb_mcqnew_calls) - 1
    ans = "A" if idx % 2 == 0 else "C"
    return f"sc reasoning candidate {ans}\n\\boxed{{{ans}}}"


try:
    f.PROPOSERS = ["m1", "m2"]
    f.REASONING_MODELS = ["m1", "m2"]
    f.SC_CHEAP_VOTES = 0
    f.SC_POT = False
    f.ARBITER_MODEL = None
    f.installed_models = _fake_installed_m1m2
    f.ask = _fake_ask_mcq_new
    _res_mcqnew = f.solve_verifiable("test mcq question", "mcq")
finally:
    f.ask = _orig_ask2
    f.PROPOSERS = _orig_props_mcqnew
    f.REASONING_MODELS = _orig_reasoning_mcqnew
    f.SC_CHEAP_VOTES = _orig_cheap_mcqnew
    f.SC_POT = _orig_pot_mcqnew
    f.ARBITER_MODEL = _orig_arbiter_model_mcqnew
    f.installed_models = _orig_installed_mcqnew

check("arb-mcqnew: 裁定役が採用した新規文字'E'がそのままanswerになる",
      _res_mcqnew is not None and _res_mcqnew["answer"] == "E")
check("arb-mcqnew: 'E'は既存タイクラスに無い第三候補として0票クラスで追加される(旧トップ票数の誤流用なし)",
      _res_mcqnew is not None and _res_mcqnew["votes"].get("E") == 0)
check("arb-mcqnew: 'A'/'C'の実票数は裁定後も書き換わらず据え置かれる(truthful)",
      _res_mcqnew is not None and _res_mcqnew["votes"].get("A") is not None
      and _res_mcqnew["votes"].get("C") is not None
      and _res_mcqnew["votes"]["A"] == _res_mcqnew["votes"]["C"] > 0)

# ==================================================
# ---------- solve_verifiable: SC_POT(PoT票混入)統合テスト (2026-07-23) ----------
# ==================================================
# gotcha #7: solve_verifiable は精度最優先の自己一貫性投票パス。SC_POT は本番既定値が
# True (fugu_local.py L2352) で、real な math ベンチでは常に PoT 票が混入する。しかし
# 上のブロック(L800〜)の solve_verifiable テストは全て SC_POT=False / SC_CHEAP_VOTES=0
# を強制しており(一つは L829 相当で len(_sc_calls)==SC_INITIAL を厳密一致で断言し、
# これは SC_POT=False でのみ成立する)、add_batch の PoT 分岐
# (`if SC_POT and task_type == 'math': add(models[0], pot=True)`)・PoT サンプルの
# 投票への計上・main_cot_count() が PoT/安価票を意図的に SC_MAX 上限計算から除外する
# 挙動は、これまで統合レベルで一切検証されていなかった。
#
# ここでは f.ask ではなく f._sc_sample そのものを丸ごとモックする。_sc_sample は
# add() 内でモジュール名前空間経由で解決されるため（`add` は solve_verifiable の
# ネスト関数で、bare name `_sc_sample` を毎回モジュールグローバルから引く）、
# f._sc_sample を差し替えれば add_batch の実ループ・cheap_ok ゲート・
# main_cot_count() の除外判定・投票集計・SC_MIN_VOTES 下限・終了判定は本物のまま
# 実行される。extract_code/run_python のサブプロセス生成には一切触れない
# （PoT の _sc_sample 内部自体は iteration 4/52/61 で既にモック ask 経由で検証済み）。
#
# 既存の solve_verifiable テスト（L829 の SC_POT=False 前提の
# len(_sc_calls)==SC_INITIAL 断言を含む）は一切変更しない。ここは追加ブロックのみ。

_orig_scpot_sc_sample = f._sc_sample
_orig_scpot_pot = f.SC_POT
_orig_scpot_cheap_votes = f.SC_CHEAP_VOTES
_orig_scpot_cheap_model = f.SC_CHEAP_MODEL
_orig_scpot_props = f.PROPOSERS
_orig_scpot_reasoning = f.REASONING_MODELS
_orig_scpot_installed = f.installed_models
_orig_scpot_initial = f.SC_INITIAL
_orig_scpot_step = f.SC_STEP
_orig_scpot_max = f.SC_MAX
_orig_scpot_min_votes = f.SC_MIN_VOTES


def _make_fake_sc_sample(answer_map, calls_log):
    """(model, pot) をキーに、あらかじめ用意した答えのリストを呼び出し順に払い出す
    f._sc_sample の代替品。リストを使い切ったら最後の値を繰り返す。
    calls_log には実際に呼ばれた (model, pot) を呼び出し順そのまま記録する
    （ask/run_python は一切呼ばない）。"""
    _idx = {}

    def _fake(model, question, task_type, pot=False, history=None):
        key = (model, pot)
        calls_log.append(key)
        i = _idx.get(key, 0)
        _idx[key] = i + 1
        lst = answer_map.get(key, [None])
        ans = lst[i] if i < len(lst) else lst[-1]
        kind = "pot" if pot else "cot"
        return ans, f"[{kind}:{model}:{i}] answer={ans}"

    return _fake


# ---- (A) math + SC_POT=True: PoTサンプルはadd_batch毎に1件、models[0]のみで実行され、
#      その答えは投票(n/cnt/votes)にちゃんと参加する。全会一致で即確定する単純ケース ----
_scpot_a_calls = []
_scpot_a_map = {
    ("m1", False): ["42", "42"],
    ("m2", False): ["42", "42"],
    ("m1", True): ["42"],
}
try:
    f.PROPOSERS = ["m1", "m2"]
    f.REASONING_MODELS = ["m1", "m2"]
    f.SC_CHEAP_VOTES = 0
    f.SC_POT = True
    f.SC_INITIAL = 4
    f._sc_sample = _make_fake_sc_sample(_scpot_a_map, _scpot_a_calls)
    _res_scpot_a = f.solve_verifiable("test question", "math")
finally:
    f._sc_sample = _orig_scpot_sc_sample
    f.PROPOSERS = _orig_scpot_props
    f.REASONING_MODELS = _orig_scpot_reasoning
    f.SC_CHEAP_VOTES = _orig_scpot_cheap_votes
    f.SC_POT = _orig_scpot_pot
    f.SC_INITIAL = _orig_scpot_initial

_scpot_a_pot_calls = [c for c in _scpot_a_calls if c[1]]
check("sc-pot: math+SC_POT=TrueでPoTサンプルがadd_batch毎に1件だけ追加される",
      len(_scpot_a_pot_calls) == 1)
check("sc-pot: PoTサンプルはmodels[0](m1)でのみ実行される",
      _scpot_a_pot_calls == [("m1", True)])
check("sc-pot: CoT+PoT計5サンプル全会一致で確定(n_samplesにPoT分も含む)",
      _res_scpot_a is not None and _res_scpot_a["answer"] == "42"
      and _res_scpot_a["n_samples"] == 5)
check("sc-pot: PoTサンプルの答えがvotes/cntに計上されている(5票すべて'42')",
      _res_scpot_a is not None and _res_scpot_a["votes"] == {"42": 5})

# ---- (B) mcq + SC_POT=True: add_batchの `and task_type == 'math'` ガードにより
#      PoTサンプルは一切追加されない ----
_scpot_b_calls = []
_scpot_b_map = {
    ("m1", False): ["A", "A"],
    ("m2", False): ["A", "A"],
}
try:
    f.PROPOSERS = ["m1", "m2"]
    f.REASONING_MODELS = ["m1", "m2"]
    f.SC_CHEAP_VOTES = 0
    f.SC_POT = True
    f.SC_INITIAL = 4
    f._sc_sample = _make_fake_sc_sample(_scpot_b_map, _scpot_b_calls)
    _res_scpot_b = f.solve_verifiable("test question", "mcq")
finally:
    f._sc_sample = _orig_scpot_sc_sample
    f.PROPOSERS = _orig_scpot_props
    f.REASONING_MODELS = _orig_scpot_reasoning
    f.SC_CHEAP_VOTES = _orig_scpot_cheap_votes
    f.SC_POT = _orig_scpot_pot
    f.SC_INITIAL = _orig_scpot_initial

check("sc-pot: task_type='mcq'ではSC_POT=TrueでもPoTサンプルは追加されない(samples中にpot=Trueが無い)",
      not any(pot for _m, pot in _scpot_b_calls))
check("sc-pot: mcqはSC_INITIAL(4件のCoTのみ)で全会一致確定する",
      _res_scpot_b is not None and _res_scpot_b["answer"] == "A"
      and _res_scpot_b["n_samples"] == 4 and len(_scpot_b_calls) == 4)

# ---- (C) math + SC_POT=True: 1バッチ目のCoTだけなら2-2で決着しないが、PoT票が
#      一方に加わることで過半数(3/5)に押し上げ、そのPoT票自体が確定に寄与すること
#      を示す(全会一致ではない、n<全票のケース) ----
_scpot_c_calls = []
_scpot_c_map = {
    ("m1", False): ["1", "1"],
    ("m2", False): ["2", "2"],
    ("m1", True): ["1"],
}
try:
    f.PROPOSERS = ["m1", "m2"]
    f.REASONING_MODELS = ["m1", "m2"]
    f.SC_CHEAP_VOTES = 0
    f.SC_POT = True
    f.SC_INITIAL = 4
    f._sc_sample = _make_fake_sc_sample(_scpot_c_map, _scpot_c_calls)
    _res_scpot_c = f.solve_verifiable("test question", "math")
finally:
    f._sc_sample = _orig_scpot_sc_sample
    f.PROPOSERS = _orig_scpot_props
    f.REASONING_MODELS = _orig_scpot_reasoning
    f.SC_CHEAP_VOTES = _orig_scpot_cheap_votes
    f.SC_POT = _orig_scpot_pot
    f.SC_INITIAL = _orig_scpot_initial

check("sc-pot: CoTのみなら2-2で拮抗するところ、PoT票が'1'側に加わり3-2の過半数で確定",
      _res_scpot_c is not None and _res_scpot_c["answer"] == "1"
      and _res_scpot_c["n_samples"] == 5)
check("sc-pot: 過半数側の勝者票(3)にPoT票が含まれる(CoT'1'は2票のみ、PoTの+1で3票)",
      _res_scpot_c is not None and _res_scpot_c["votes"] == {"1": 3, "2": 2})

# ---- (D) math + SC_POT=True: PoTはmain_cot_count()（=SC_MAXの上限判定）から
#      意図的に除外される。SC_MAXを「PoTを含めた総サンプル数」で見た場合と
#      「PoT除外の主力CoT数」で見た場合とで、ループが継続するかどうかが分岐する
#      値(SC_MAX=5)を選び、実際には除外仕様どおり2バッチ目まで継続する
#      （もしPoTがカウントされていたら1バッチ目=総数5で即打ち切られてしまうはず）
#      ことを検証する。かつ無限ループにならないこと。 ----
_scpot_d_calls = []
_scpot_d_map = {
    ("m1", False): ["1", "1", "1", "1"],
    ("m2", False): ["2", "2", "2", "2"],
    ("m1", True): ["3", "1"],
}
try:
    f.PROPOSERS = ["m1", "m2"]
    f.REASONING_MODELS = ["m1", "m2"]
    f.SC_CHEAP_VOTES = 0
    f.SC_POT = True
    f.SC_INITIAL = 4
    f.SC_STEP = 4
    f.SC_MAX = 5
    f._sc_sample = _make_fake_sc_sample(_scpot_d_map, _scpot_d_calls)
    _res_scpot_d = f.solve_verifiable("test question", "math")
finally:
    f._sc_sample = _orig_scpot_sc_sample
    f.PROPOSERS = _orig_scpot_props
    f.REASONING_MODELS = _orig_scpot_reasoning
    f.SC_CHEAP_VOTES = _orig_scpot_cheap_votes
    f.SC_POT = _orig_scpot_pot
    f.SC_INITIAL = _orig_scpot_initial
    f.SC_STEP = _orig_scpot_step
    f.SC_MAX = _orig_scpot_max

_scpot_d_pot_calls = [c for c in _scpot_d_calls if c[1]]
_scpot_d_main_calls = [c for c in _scpot_d_calls if not c[1]]
check("sc-pot: main_cot_count()除外によりadd_batchが2回走る"
      "(PoTを含めた総数だとSC_MAX=5に1バッチ目で到達し1回で打ち切られるはずだが、"
      "実際は主力CoT数だけで判定するため2バッチ目まで継続する)",
      len(_scpot_d_pot_calls) == 2 and len(_scpot_d_calls) == 10)
check("sc-pot: PoTはadd_batch毎(2回)に1件だけ、常にmodels[0]で実行される",
      _scpot_d_pot_calls == [("m1", True), ("m1", True)])
check("sc-pot: 主力CoTサンプル数がSC_MAX(5)以上に達した時点で打ち切られる(有限回で終了)",
      len(_scpot_d_main_calls) == 8 and len(_scpot_d_main_calls) >= 5)
check("sc-pot: PoT票を含む総サンプル数(10)は主力CoT数(8)より多い"
      "(SC_MAXの判定が総数ではなく主力CoT数であることの直接証拠)",
      _res_scpot_d is not None and _res_scpot_d["n_samples"] == 10)
check("sc-pot: PoT票も含めた最終投票結果は真実(votesにPoTの寄与が正しく反映)",
      _res_scpot_d is not None and _res_scpot_d["answer"] == "1"
      and _res_scpot_d["votes"] == {"1": 5, "2": 4, "3": 1})

# ---- (E, optional) 安価票(SC_CHEAP_VOTES>0)ケース: installed_modelsをモックして
#      cheap_ok=Trueにし、初回バッチ直後に安価票がSC_CHEAP_VOTES件だけ一度だけ追加され、
#      投票には参加するがmain_cot_count()（SC_MAXの上限判定）からは除外されること。
#      SC_MAXをmain-onlyとtotal-with-cheapの間に置き、除外の有無で分岐する形で検証する
#      （もし安価票がカウントされていたら安価票追加直後=総数6でSC_MAX=5に到達し
#      即打ち切られるはずだが、実際は主力CoT数(4)だけで判定するため2バッチ目まで
#      継続する）----
_scpot_e_calls = []
_scpot_e_map = {
    ("m1", False): ["1", "1", "1", "1"],
    ("m2", False): ["2", "2", "4", "4"],
    ("cheapM", False): ["3", "3"],
}
_orig_scpot_e_installed = f.installed_models
try:
    f.PROPOSERS = ["m1", "m2"]
    f.REASONING_MODELS = ["m1", "m2"]
    f.SC_POT = False
    f.SC_CHEAP_MODEL = "cheapM"
    f.SC_CHEAP_VOTES = 2
    f.installed_models = lambda: ["m1", "m2", "cheapM"]
    f.SC_INITIAL = 4
    f.SC_STEP = 4
    f.SC_MAX = 5
    f._sc_sample = _make_fake_sc_sample(_scpot_e_map, _scpot_e_calls)
    _res_scpot_e = f.solve_verifiable("test question", "math")
finally:
    f._sc_sample = _orig_scpot_sc_sample
    f.PROPOSERS = _orig_scpot_props
    f.REASONING_MODELS = _orig_scpot_reasoning
    f.SC_POT = _orig_scpot_pot
    f.SC_CHEAP_MODEL = _orig_scpot_cheap_model
    f.SC_CHEAP_VOTES = _orig_scpot_cheap_votes
    f.installed_models = _orig_scpot_e_installed
    f.SC_INITIAL = _orig_scpot_initial
    f.SC_STEP = _orig_scpot_step
    f.SC_MAX = _orig_scpot_max

_scpot_e_cheap_calls = [c for c in _scpot_e_calls if c[0] == "cheapM"]
check("sc-cheap: 安価票はSC_CHEAP_VOTES(2)件だけ、初回バッチ直後に一度だけ追加される"
      "(先頭4件がCoT、続く2件が安価票)",
      _scpot_e_calls[:4] == [("m1", False), ("m1", False), ("m2", False), ("m2", False)]
      and _scpot_e_calls[4:6] == [("cheapM", False), ("cheapM", False)]
      and len(_scpot_e_cheap_calls) == 2)
check("sc-cheap: main_cot_count()除外により安価票追加後も打ち切られず2バッチ目まで継続する"
      "(総数基準ならSC_MAX=5に安価票追加直後の総数6で到達し打ち切られるはずだが、"
      "実際は主力CoT数だけで判定するため継続する)",
      len(_scpot_e_calls) == 10)
check("sc-cheap: 安価票は投票結果(votes/n_samples)に正しく計上される",
      _res_scpot_e is not None and _res_scpot_e["n_samples"] == 10
      and _res_scpot_e["votes"].get("3") == 2)

check("sc-pot/cheap: テスト後にf._sc_sampleが元へ復元されている",
      f._sc_sample == _orig_scpot_sc_sample)
check("sc-pot/cheap: テスト後にSC_POT/SC_CHEAP_VOTES/SC_CHEAP_MODELが元へ復元されている",
      f.SC_POT == _orig_scpot_pot and f.SC_CHEAP_VOTES == _orig_scpot_cheap_votes
      and f.SC_CHEAP_MODEL == _orig_scpot_cheap_model)
check("sc-pot/cheap: テスト後にPROPOSERS/REASONING_MODELS/installed_modelsが元へ復元されている",
      f.PROPOSERS == _orig_scpot_props and f.REASONING_MODELS == _orig_scpot_reasoning
      and f.installed_models == _orig_scpot_installed)
check("sc-pot/cheap: テスト後にSC_INITIAL/SC_STEP/SC_MAX/SC_MIN_VOTESが元へ復元されている",
      f.SC_INITIAL == _orig_scpot_initial and f.SC_STEP == _orig_scpot_step
      and f.SC_MAX == _orig_scpot_max and f.SC_MIN_VOTES == _orig_scpot_min_votes)

# ---------- task_type ガードレール ----------
def _tt(q, declared=""):
    return f._apply_tasktype_guardrails(q, {"task_type": declared})["task_type"]


check("tt: AIME風は math", _tt("Find the number of ordered pairs...") == "math")
check("tt: 日本語計算は math", _tt("1000円の3割引の支払額を求めよ") == "math")
check("tt: 選択肢列挙は mcq", _tt("正しいものを選べ\nA) foo\nB) bar") == "mcq")
check("tt: which of the following は mcq", _tt("Which of the following is true?") == "mcq")
check("tt: コードは code", _tt("フィボナッチ関数を実装して") == "code")
check("tt: 証明は math にしない", _tt("3連続整数の積が6の倍数であることを証明して求めよ") != "math")
check("tt: Conductor申告を尊重", _tt("こんにちは", "chat") == "chat")
check("tt: 不明シグナルは chat", _tt("よろしくね", "") == "chat")

# --- 誤申告レスキュー: _apply_tasktype_guardrails は solve_verifiable(自己一貫性投票、
# gotcha #7) への入口ゲート。小型 Conductor が math/mcq 問題を誤って chat/knowledge 等に
# 分類しても、確実なキーワードシグナルがあれば強制的に正しい task_type へ補正し、投票
# 経路を失わないことを検証する（Conductor申告のみを尊重する既存チェックの逆側）。
check("tt: 誤申告(chat)でも強い math シグナルで math へレスキュー",
      _tt("Find the remainder when 7^100 is divided by 13.", "chat") == "math")
check("tt: 誤申告(knowledge)でも選択肢列挙シグナルで mcq へレスキュー",
      _tt("正しいものを選べ\nA) foo\nB) bar", "knowledge") == "mcq")

# --- シグナル優先順位: 実装コードは if/elif mcq -> code -> math の順で判定されるため、
# mcq シグナルが最優先、次いで code、最後に math。複数シグナルが同居する問題での
# 優先順位を固定する。
check("tt: code と math シグナルが同居 -> elif順で code が勝つ",
      _tt("フィボナッチ関数を実装して。その関数の余りを求めよ。") == "code")
check("tt: mcq と code シグナルが同居 -> mcq が最優先",
      _tt("次のうち、コードの実装として正しいものを選べ\nA) foo\nB) bar") == "mcq")

# --- 自由記述デモーションのスコープ: t == "math" のときのみ証明/説明系シグナルで
# knowledge へ格下げされる（＝投票に回さない）。declared="mcq" など math 以外に確定した
# 場合はこのデモーション条件に入らず、証明系ワードが含まれていても格下げされないことを
# 確認する。
check("tt: 申告math + 自由記述シグナル(シグナル一致なし) -> knowledge へ格下げ",
      _tt("なぜ空は青いのか", "math") == "knowledge")
check("tt: 申告mcq + 自由記述ワードが含まれていても mcq のまま(格下げされない)",
      _tt("なぜ空は青いか説明して", "mcq") == "mcq")

# --- 2026-07-26: _FREEFORM_SIGNALS の英語対応漏れ修正。_MATH_TASK_SIGNALS には
# \bhow many\b / \bcompute\b 等の英語トリガーが揃っているのに、降格側(_FREEFORM_SIGNALS)
# は日本語の説明系(証明/説明して/解説して/なぜ)しか見ておらず、英語の explain/describe/why
# が欠けていた。結果、"How many SOLID principles are there? Explain each." のように
# \bhow many\b で math 判定されつつ説明を求める問いが投票経路(solve_verifiable、gotcha #7)
# に流れ、要求された説明が黙って落ちる回帰があった。explain/describe/why の3語追加で
# 日本語側と対称になったことを検証する。
check("tt: 英語 how many+explain は knowledge へ格下げ(新規)",
      _tt("How many SOLID principles are there? Explain each.") == "knowledge")
check("tt: 英語 how many+describe は knowledge へ格下げ(新規)",
      _tt("How many design patterns exist? Describe when to use each.") == "knowledge")
check("tt: 英語 how many+why は knowledge へ格下げ(新規)",
      _tt("How many degrees does a triangle's interior angles sum to, and why?")
      == "knowledge")
# --- 回帰: 説明系ワードを伴わない正当な英語計数/計算問題は math のまま
# (bare \bhow\b を追加していないことの確認。\bhow many\b との衝突を避けるのが本修正の
# 必須制約 -- gotcha #7 の自己一貫性投票ゲートを誤って外さないため)。
check("tt: how many のみ(説明要求なし)は math のまま(回帰)",
      _tt("How many prime numbers are less than 100?") == "math")
check("tt: how many のみ(説明要求なし・別問題)は math のまま(回帰)",
      _tt("How many divisors does 60 have?") == "math")
check("tt: compute のみ(説明要求なし)は math のまま(回帰)",
      _tt("Compute 17*23.") == "math")
check("tt: find the value のみ(説明要求なし)は math のまま(回帰)",
      _tt("Find the value of x if 2x=10.") == "math")
# --- 回帰: 日本語側の既存デモーション/確定は byte-for-byte 不変
check("tt: 日本語math専用(を計算し)は引き続き math のまま(回帰)",
      _tt("42の階乗を計算しなさい") == "math")
check("tt: 日本語証明系は引き続き knowledge へ格下げ(回帰)",
      _tt("3連続整数の積が6の倍数であることを証明して求めよ", "math") == "knowledge")
# --- 回帰: 新規English語彙があっても mcq/code 分類は不変(mcqはmath分岐より先に確定、
# 格下げは t=="math" のときのみ発火するため構造的に影響を受けない)
check("tt: explain を含む mcq 問題は mcq のまま(回帰)",
      _tt("Which of the following is correct? Explain your reasoning.\nA) foo\nB) bar")
      == "mcq")
check("tt: explain を含む code 問題は code のまま(回帰)",
      _tt("Implement a function to reverse a string, then explain your approach.")
      == "code")

check("tt: validate が task_type を保持",
      f.validate_plan({"mode": "moa", "selected_proposers": [],
                       "task_type": "math"})["task_type"] == "math")
check("tt: validate が不正 task_type を空へ",
      f.validate_plan({"mode": "moa", "selected_proposers": [],
                       "task_type": "quiz"})["task_type"] == "")

# ---------- MODEL_CONFIG 解決 ----------
check("cfg: 既知モデルの num_ctx", f.model_cfg("gpt-oss:20b", "num_ctx") == 16384)
check("cfg: 未知モデルは default", f.model_cfg("nonexistent", "num_ctx", 8192) == 8192)
check("cfg: think 段階指定", f.model_cfg("gpt-oss:20b", "think") == "high")

# ---------- is_installed（インストール済み判定：厳密タグ一致） ----------
# resolve_models() が DESIRED_PROPOSERS を採否判定し、_arbitrate() が ARBITER_MODEL の
# 起用可否を判定する土台。docstring の通り、旧 startswith 実装は 'qwen3:4b' が
# 'qwen3:4b-instruct' に誤ヒットするバグを持っていたため厳密一致へ変更され、
# タグ無し指定のときだけ ':latest' を許容する例外が残された。この判定を誤ると
# 未導入モデルを誤って起用/裁定役に据えたり、導入済みの正規プロポーザーを
# 黙って除外したりして、精度優先のアンサンブル構成が静かに壊れる。
check("inst: 厳密タグ一致で導入判定", f.is_installed("qwen3:4b", ["qwen3:4b"]) is True)
check("inst: 旧startswithの誤検知(タグ違い)は拒否",
      f.is_installed("qwen3:4b", ["qwen3:4b-instruct"]) is False)
check("inst: タグ無し指定は :latest 導入を許容",
      f.is_installed("qwen3", ["qwen3:latest"]) is True)
check("inst: タグ無し指定はタグ無し導入も許容",
      f.is_installed("qwen3", ["qwen3"]) is True)
check("inst: タグ無し指定は任意のタグ付き導入には一致しない",
      f.is_installed("qwen3", ["qwen3:4b"]) is False)
check("inst: タグ付き指定はタグ無し導入では満たされない(非対称)",
      f.is_installed("qwen3:latest", ["qwen3"]) is False)
check("inst: 空リストは未導入", f.is_installed("gpt-oss:20b", []) is False)
check("inst: 無関係な導入リストのみでは未導入",
      f.is_installed("gpt-oss:20b", ["phi4", "gemma4:26b"]) is False)

# ---------- 大VRAMプロファイル ----------
_hv_saved = (dict(f.MODEL_CONFIG), f.PARALLEL_PROPOSERS, f.SC_INITIAL, f.SC_MAX,
             f.SC_CHEAP_VOTES, f.MODEL_NUM_CTX)
try:
    _cfg_snapshot = {m: dict(c) for m, c in f.MODEL_CONFIG.items()}
    f.MODEL_CONFIG = {m: dict(c) for m, c in _cfg_snapshot.items()}
    f.apply_high_vram_profile()
    check("hv: 並列ON", f.PARALLEL_PROPOSERS is True)
    check("hv: SC上限を引き上げ", f.SC_MAX >= 40)
    check("hv: num_ctx拡大", f.model_cfg("gpt-oss:20b", "num_ctx") == 65536)
    check("hv: 安価票を有効化", f.SC_CHEAP_VOTES >= 8)
finally:
    (f.MODEL_CONFIG, f.PARALLEL_PROPOSERS, f.SC_INITIAL, f.SC_MAX,
     f.SC_CHEAP_VOTES, f.MODEL_NUM_CTX) = _hv_saved
_orig_pt = f.PROPOSER_THINK
try:
    f.PROPOSER_THINK = None
    check("cfg: think解決 グローバルNoneは設定値", f.proposer_think_for("gpt-oss:20b") == "high")
    f.PROPOSER_THINK = False
    check("cfg: think解決 グローバル優先", f.proposer_think_for("gpt-oss:20b") is False)
finally:
    f.PROPOSER_THINK = _orig_pt

# ---------- use_jp_aggregator ----------
check("jp: ひらがな", f.use_jp_aggregator("これはテストです"))
check("jp: カタカナ", f.use_jp_aggregator("テスト"))
check("jp: 漢字のみ(旧版の取りこぼし)", f.use_jp_aggregator("東京都の人口密度?"))
check("jp: 英語はFalse", not f.use_jp_aggregator("What is the capital of France?"))
check("jp: 空/None耐性", not f.use_jp_aggregator("") and not f.use_jp_aggregator(None))

# ---------- pick_aggregator（統合役の選定ルーティング） ----------
# has_code / 日本語 / それ以外 の各分岐と優先順位（コード > 日本語 > 推論型既定）を
# module globals をセンチネル値へ差し替えて検証する。JP_AGGREGATOR_STRONG / JP_AGGREGATOR /
# AGGREGATOR_REASONING / AGGREGATOR / CONDUCTOR / PROPOSERS はすべて try/finally で
# 復元し、後続チェック（aggregate 経由で pick_aggregator を間接的に叩く箇所を含む）に
# 汚染したグローバルを見せないようにする。
# 判定の曖昧さを避けるため、日本語入力はひらがな/カタカナのみ、非日本語入力はASCIIのみを使う
# （use_jp_aggregator の docstring が指摘する「漢字だけでは日中判別不能」問題を踏まない）。
_JP_TEXT = "これはテストです"   # ひらがな -> use_jp_aggregator は True
_EN_TEXT = "What is this?"    # ASCIIのみ -> use_jp_aggregator は False

_pa_orig_jp_strong = f.JP_AGGREGATOR_STRONG
_pa_orig_jp = f.JP_AGGREGATOR
_pa_orig_agg_reasoning = f.AGGREGATOR_REASONING
_pa_orig_agg = f.AGGREGATOR
_pa_orig_conductor = f.CONDUCTOR
_pa_orig_proposers = f.PROPOSERS
try:
    f.AGGREGATOR = "SENTINEL-AGGREGATOR-CODE"
    f.AGGREGATOR_REASONING = "SENTINEL-AGGREGATOR-REASONING"
    f.JP_AGGREGATOR_STRONG = "SENTINEL-JP-STRONG"
    f.JP_AGGREGATOR = "SENTINEL-JP-BASE"
    f.CONDUCTOR = "SENTINEL-CONDUCTOR"

    # (1) has_code=True は言語に関係なく常に AGGREGATOR
    f.PROPOSERS = []
    check("pick_agg: has_code優先(日本語入力でもAGGREGATOR)",
          f.pick_aggregator(_JP_TEXT, has_code=True) == f.AGGREGATOR)
    check("pick_agg: has_code優先(英語入力でもAGGREGATOR)",
          f.pick_aggregator(_EN_TEXT, has_code=True) == f.AGGREGATOR)

    # (2) 日本語 + JP_AGGREGATOR_STRONG が導入済み(PROPOSERSに存在) -> STRONGを返す
    f.PROPOSERS = [f.JP_AGGREGATOR_STRONG]
    check("pick_agg: 日本語+強JPモデル導入済みはSTRONGを返す",
          f.pick_aggregator(_JP_TEXT, has_code=False) == f.JP_AGGREGATOR_STRONG)

    # (3a) 日本語 + STRONGがfalsy(未設定) + JP_AGGREGATORがPROPOSERSに存在 -> JP_AGGREGATORを返す
    f.JP_AGGREGATOR_STRONG = None
    f.PROPOSERS = [f.JP_AGGREGATOR]
    check("pick_agg: 日本語+STRONG未設定はJP_AGGREGATORへ(導入済み)",
          f.pick_aggregator(_JP_TEXT, has_code=False) == f.JP_AGGREGATOR)

    # (3b) STRONGが値を持っていてもPROPOSERS未導入なら同様にJP_AGGREGATORへ落ちる
    f.JP_AGGREGATOR_STRONG = "SENTINEL-JP-STRONG"
    f.PROPOSERS = [f.JP_AGGREGATOR]  # STRONGはPROPOSERSに含まれない=未導入扱い
    check("pick_agg: 強JPモデルがPROPOSERS未導入ならJP_AGGREGATORへ",
          f.pick_aggregator(_JP_TEXT, has_code=False) == f.JP_AGGREGATOR)

    # (4) 日本語 + JP_AGGREGATORはPROPOSERS未導入だがCONDUCTORと一致 -> JP_AGGREGATORを返す
    f.JP_AGGREGATOR_STRONG = None
    f.PROPOSERS = ["SENTINEL-OTHER-MODEL"]
    f.CONDUCTOR = f.JP_AGGREGATOR
    check("pick_agg: JP_AGGREGATOR未導入でもCONDUCTORと一致すれば採用",
          f.pick_aggregator(_JP_TEXT, has_code=False) == f.JP_AGGREGATOR)

    # (5) 日本語 + JP系条件を何も満たさない -> 最終フォールバックのAGGREGATOR
    #     （現行docstring通りの意図的フォールバックであり、qwen3系のAGGREGATORも
    #     deepseek-r1の言語混入問題を踏まない選択なので、これは特性確認であってバグ報告ではない）
    f.CONDUCTOR = "SENTINEL-CONDUCTOR-OTHER"
    check("pick_agg: 日本語でJP系条件すべて不成立ならAGGREGATOR(特性確認)",
          f.pick_aggregator(_JP_TEXT, has_code=False) == f.AGGREGATOR)

    # (6) 非日本語 + AGGREGATOR_REASONINGが導入済み -> AGGREGATOR_REASONINGを返す
    f.PROPOSERS = [f.AGGREGATOR_REASONING]
    check("pick_agg: 非日本語は導入済みのAGGREGATOR_REASONINGへ",
          f.pick_aggregator(_EN_TEXT, has_code=False) == f.AGGREGATOR_REASONING)

    # (7) 非日本語 + AGGREGATOR_REASONINGが未導入 -> AGGREGATORへフォールバック
    f.PROPOSERS = ["SENTINEL-OTHER-MODEL"]
    check("pick_agg: 非日本語でAGGREGATOR_REASONING未導入はAGGREGATORへ",
          f.pick_aggregator(_EN_TEXT, has_code=False) == f.AGGREGATOR)

    # (8) 相互作用: 日本語だがhas_code=True -> コードがJP系より優先されAGGREGATORを返す
    f.JP_AGGREGATOR_STRONG = "SENTINEL-JP-STRONG"
    f.PROPOSERS = [f.JP_AGGREGATOR_STRONG, f.JP_AGGREGATOR]  # JP条件は満たせる状態にしておく
    f.CONDUCTOR = f.JP_AGGREGATOR
    check("pick_agg: 日本語でもhas_code=TrueならAGGREGATOR(コード優先)",
          f.pick_aggregator(_JP_TEXT, has_code=True) == f.AGGREGATOR)
finally:
    f.JP_AGGREGATOR_STRONG = _pa_orig_jp_strong
    f.JP_AGGREGATOR = _pa_orig_jp
    f.AGGREGATOR_REASONING = _pa_orig_agg_reasoning
    f.AGGREGATOR = _pa_orig_agg
    f.CONDUCTOR = _pa_orig_conductor
    f.PROPOSERS = _pa_orig_proposers

# ---------- aggregate のフォールバック（ask をモンキーパッチ） ----------
_orig_ask = f.ask
_ask_log = []


def _fake_ask_empty(model, messages, temperature, think=None, fmt=None,
                    label=None, num_predict=None):
    """アグリゲータ/再統合/Criticすべて空返答 → 保険2の最終分岐(最長の提案)まで落ちる。
    ※Critic は extract_json 失敗時 ok=True 既定なので、実際は最初の提案が返る。"""
    _ask_log.append((label, model, think))
    return ""


f.ask = _fake_ask_empty
try:
    out = f.aggregate("Q?", [("m1", "short"), ("m2", "much longer answer")])
finally:
    f.ask = _orig_ask
_agg_calls = [(m, th) for lab, m, th in _ask_log if lab == "aggregator"]
check("agg: 全滅時も提案のどれかを返す(空にしない)", out in ("short", "much longer answer"))
check("agg: 再統合(保険1)が試行されている", len(_agg_calls) == 2)
check("agg: 保険1は JP_AGGREGATOR + think=False で再統合",
      _agg_calls[1] == (f.JP_AGGREGATOR, False))

_ask_log.clear()


def _fake_ask_ok(model, messages, temperature, think=None, fmt=None,
                 label=None, num_predict=None):
    _ask_log.append((label, model, think))
    return "aggregated!"


f.ask = _fake_ask_ok
try:
    out = f.aggregate("Q?", [("m1", "a"), ("m2", "b")])
finally:
    f.ask = _orig_ask
check("agg: 正常時は統合結果を返す", out == "aggregated!")
check("agg: 正常時は1回だけ呼ぶ", len(_ask_log) == 1)

# エラー提案しかない場合
check("agg: 全プロポーザー失敗は__ERROR__",
      f.aggregate("Q?", [("m1", "__ERROR__: x")]).startswith("__ERROR__"))

# ---------- agg: 保険2(insurance-2)フォールバックが [Execution check: ...] を漏らさない ----------
# 2026-07-22: aggregate() は以前、コード付き提案に実行結果タグを付ける際に `good` 自体を
# タグ付き版で上書きしていた。保険2(統合失敗時に good から直接返す)経路は「主アグリゲータ」と
# 「JP_AGGREGATOR(think=False)再統合」の両方が空/エラーを返す場合に到達する
# (2026-07-04 の空返答実測に基づく既知の実運用経路)。このタグ/生トレースバックが
# ユーザー向け回答にそのまま漏れていたバグの回帰テスト。


def _fake_ask_always_empty(model, messages, temperature, think=None, fmt=None,
                            label=None, num_predict=None):
    """主アグリゲータ・再統合(保険1)ともに空/エラーを返し、保険2まで必ず落とす。"""
    return "" if label != "force_error" else "__ERROR__"


_orig_code_execution = f.CODE_EXECUTION
_orig_critique = f.critique
f.CODE_EXECUTION = True

_code_ans = "Here you go:\n\n```python\nprint(2 + 2)\n```\n"
_prose_ans = "This is a plain prose answer with no code block at all, just text."

try:
    # --- critique が最初の合格案をそのまま採用するケース ---
    f.ask = _fake_ask_always_empty
    f.critique = lambda question, answer: (True, "")
    out = f.aggregate("Q?", [("m1", _code_ans), ("m2", _prose_ans)])
    check("agg: 保険2(critique採用)はコード本文を含む", "print(2 + 2)" in out)
    check("agg: 保険2(critique採用)は[Execution check:]タグを漏らさない",
          "[Execution check:" not in out)

    # --- critique が全案を却下し、最長案(max fallback)まで落ちるケース ---
    # コード付き案をわざと最長にして、タグ付け前の `good`(クリーン版)から
    # 選ばれることを検証する。
    _code_ans_long = _code_ans + ("x" * 200)
    f.critique = lambda question, answer: (False, "no good")
    out2 = f.aggregate("Q?", [("m1", _prose_ans), ("m2", _code_ans_long)])
    check("agg: 保険2(最長fallback)は最長案(コード付き)を返す", "print(2 + 2)" in out2)
    check("agg: 保険2(最長fallback)は[Execution check:]タグを漏らさない",
          "[Execution check:" not in out2)
finally:
    f.ask = _orig_ask
    f.critique = _orig_critique
    f.CODE_EXECUTION = _orig_code_execution

# ---------- agg: 保険2は実行済みで FAILED と判明済みのコード提案を避ける (2026-07-24) ----------
# insurance-2 の critique() 承認スキャン/最長フォールバックとも、上の CODE_EXECUTION
# アノテーションループが既に計算した code_check() の結果（≒ [Execution check: ...] タグの
# 元ネタ）を使い、コードが FAILED と判明済みの提案を「他に正常/コード無しの候補がある限り」
# 避けるべき、という回帰テスト。AGGREGATOR_SYS ルール6「FAILED のコードを最終回答の根拠に
# しない」に対応。iteration 9 (good はクリーンな (model, strip_think(answer)) のまま保持) の
# 不変条件はここでも維持する ― 返り値に [Execution check: ...] タグや生トレースバックが
# 混じらないことも合わせて確認する。
_failing_code_ans = "This code fails:\n\n```python\n1/0\n```"

_orig_code_execution2 = f.CODE_EXECUTION
_orig_critique2 = f.critique
f.CODE_EXECUTION = True

try:
    # --- 失敗コード案 + プローズ案、critique は全案 ok → FAILED でない方(プローズ)を選ぶ ---
    f.ask = _fake_ask_always_empty
    f.critique = lambda question, answer: (True, "")
    out_fail_first = f.aggregate("Q?", [("m1", _failing_code_ans), ("m2", _prose_ans)])
    check("agg: 保険2はFAILED済みコード案でなく正常な案を返す(critique採用)",
          out_fail_first == _prose_ans)
    check("agg: 保険2(FAILED回避)は1/0を含む案を返さない",
          "1/0" not in out_fail_first)
    check("agg: 保険2(FAILED回避/critique採用)は[Execution check:]タグを漏らさない",
          "[Execution check:" not in out_fail_first)

    # --- 全提案のコードが FAILED → 空にせず全候補(good)から最長フォールバック ---
    _failing_code_ans_long = _failing_code_ans + ("y" * 200)
    f.critique = lambda question, answer: (False, "no good")
    out_all_failed = f.aggregate(
        "Q?", [("m1", _failing_code_ans), ("m2", _failing_code_ans_long)]
    )
    check("agg: 全案FAILEDでも例外にならず提案を返す",
          out_all_failed in (_failing_code_ans, _failing_code_ans_long))
    check("agg: 全案FAILED時は最長案(good全体からのフォールバック)を返す",
          out_all_failed == _failing_code_ans_long)
    check("agg: 全案FAILED時も[Execution check:]タグを漏らさない",
          "[Execution check:" not in out_all_failed)

    # --- CODE_EXECUTION=False なら failed_idxs は作られず、従来通り good の先頭順で選ぶ ---
    f.CODE_EXECUTION = False
    f.critique = lambda question, answer: (True, "")
    out_no_exec = f.aggregate("Q?", [("m1", _failing_code_ans), ("m2", _prose_ans)])
    check("agg: CODE_EXECUTION=Falseなら保険2は従来通りgoodの先頭案を返す",
          out_no_exec == _failing_code_ans)
finally:
    f.ask = _orig_ask
    f.critique = _orig_critique2
    f.CODE_EXECUTION = _orig_code_execution2

# ---------- agg: 正常時、アグリゲータへの user プロンプトにはタグが残っていること ----------
# (AGGREGATOR_SYS ルール6はこのタグを判断材料にするため、アグリゲータ自身が見る
#  プロンプトからタグを消してはいけない。good を汚さない修正がここを壊していないことの回帰確認)
_captured_user = []


def _fake_ask_capture(model, messages, temperature, think=None, fmt=None,
                       label=None, num_predict=None):
    for msg in messages:
        if msg.get("role") == "user":
            _captured_user.append(msg.get("content", ""))
    return "aggregated!"


f.CODE_EXECUTION = True
f.ask = _fake_ask_capture
try:
    out3 = f.aggregate("Q?", [("m1", _code_ans), ("m2", _prose_ans)])
finally:
    f.ask = _orig_ask
    f.CODE_EXECUTION = _orig_code_execution
check("agg: 正常系はアグリゲータ出力をそのまま返す", out3 == "aggregated!")
check("agg: アグリゲータへのプロンプトには[Execution check: PASSED]が残る",
      any("[Execution check: PASSED]" in u for u in _captured_user))

# ---------- get_single_proposal: think 解決に proposer_think_for を使う (2026-07-23) ----------
# 2026-07-23 fix: get_single_proposal は隣の num_predict こそ proposer_predict_for(model) で
# MODEL_CONFIG 対応していたが、think は生の PROPOSER_THINK グローバル(既定 None)を直渡し
# していた欠落サイト。これだと gpt-oss:20b/qwen3.6:35b が MoA 提案(get_single_proposal 経由)
# で think:"high"/True を一度も受け取れず、既に proposer_think_for 経由で正しく解決していた
# _sc_sample(SC経路) とだけ非対称になっていた。ここでは f.ask をモックして think/num_predict
# kwarg を直接キャプチャし、reference無(新規回答)/reference有(改善分岐)の両方で
# proposer_think_for(model) の解決結果がそのまま渡ることを検証する。ネットワーク/GPU呼び出し無し。
_gsp_think_calls = []


def _fake_ask_capture_think(model, messages, temperature, think=None, fmt=None,
                             label=None, num_predict=None):
    _gsp_think_calls.append({"model": model, "think": think, "num_predict": num_predict})
    return "proposal text"


_orig_ask_gsp_think = f.ask
_orig_pt_gsp = f.PROPOSER_THINK
f.ask = _fake_ask_capture_think
try:
    f.PROPOSER_THINK = None
    _gsp_think_calls.clear()
    f.get_single_proposal("gpt-oss:20b", "Q?", None)
    check("gsp: think解決 gpt-oss:20b(reference無=新規回答)はMODEL_CONFIGのhighを渡す",
          _gsp_think_calls[-1]["think"] == "high")
    check("gsp: num_predictはproposer_predict_for(model)のまま(回帰・不変)",
          _gsp_think_calls[-1]["num_predict"] == f.proposer_predict_for("gpt-oss:20b"))

    _gsp_think_calls.clear()
    f.get_single_proposal("gpt-oss:20b", "Q?", "draft answer")
    check("gsp: think解決 gpt-oss:20b(reference有=ドラフト改善分岐)もhighを渡す",
          _gsp_think_calls[-1]["think"] == "high")

    _gsp_think_calls.clear()
    f.get_single_proposal("qwen3.6:35b", "Q?", None)
    check("gsp: think解決 qwen3.6:35bはTrueを渡す", _gsp_think_calls[-1]["think"] is True)

    _gsp_think_calls.clear()
    f.get_single_proposal("qwen3-coder:30b", "Q?", None)
    check("gsp: think解決 qwen3-coder:30b(MODEL_CONFIGにthink無し)はNoneのまま(回帰・400防止)",
          _gsp_think_calls[-1]["think"] is None)

    _gsp_think_calls.clear()
    f.get_single_proposal("gemma4:26b", "Q?", None)
    check("gsp: think解決 gemma4:26b(MODEL_CONFIGにthink無し)はNoneのまま(回帰・400防止)",
          _gsp_think_calls[-1]["think"] is None)

    # override優先: PROPOSER_THINK が None 以外なら全モデルそちらが最優先(eval一括OFF等)
    f.PROPOSER_THINK = False
    _gsp_think_calls.clear()
    f.get_single_proposal("gpt-oss:20b", "Q?", None)
    check("gsp: PROPOSER_THINK override時はgpt-oss:20bもFalse(eval一括OFFとbyte一致)",
          _gsp_think_calls[-1]["think"] is False)
finally:
    f.ask = _orig_ask_gsp_think
    f.PROPOSER_THINK = _orig_pt_gsp
check("gsp: テスト後にf.askが元に復元されている", f.ask == _orig_ask_gsp_think)
check("gsp: テスト後にPROPOSER_THINKが元に復元されている", f.PROPOSER_THINK == _orig_pt_gsp)

# ---------- get_proposals の多様性（先頭はドラフト無しで新規回答） ----------
_seen_refs = []


def _fake_proposal(model, question, reference, issue=None, history=None):
    _seen_refs.append(reference)
    return model, "ans"


_orig_gsp = f.get_single_proposal
f.get_single_proposal = _fake_proposal
try:
    f.PARALLEL_PROPOSERS = False
    f.get_proposals(["m1", "m2", "m3"], "Q?", reference="draft", issue="x")
finally:
    f.get_single_proposal = _orig_gsp
check("prop: ラウンド2の先頭は新規回答(reference=None)", _seen_refs[0] is None)
check("prop: 2体目以降はドラフト改善", _seen_refs[1] == "draft" and _seen_refs[2] == "draft")

# ---------- get_proposals: PARALLEL_PROPOSERS=True でも多様性契約が保たれる ----------
# 2026-07-23: 並列分岐(get_single_proposalをThreadPoolExecutorで実行)でも、
# jobs[0](=models[0])だけがreference=Noneで呼ばれる契約が崩れていないことを確認する。
_seen_refs_parallel = {}


def _fake_proposal_parallel(model, question, reference, issue=None, history=None):
    _seen_refs_parallel[model] = reference
    return model, "ans"


_orig_gsp_par = f.get_single_proposal
_orig_parallel_flag_1 = f.PARALLEL_PROPOSERS
f.get_single_proposal = _fake_proposal_parallel
try:
    f.PARALLEL_PROPOSERS = True
    f.get_proposals(["m1", "m2", "m3"], "Q?", reference="draft", issue="x")
finally:
    f.get_single_proposal = _orig_gsp_par
    f.PARALLEL_PROPOSERS = _orig_parallel_flag_1
check("prop(parallel): 先頭(models[0])は新規回答(reference=None)",
      _seen_refs_parallel.get("m1") is None)
check("prop(parallel): 2体目以降はドラフト改善",
      _seen_refs_parallel.get("m2") == "draft" and _seen_refs_parallel.get("m3") == "draft")

# ---------- get_proposals: PARALLEL_PROPOSERS=True は completion順ではなく submission順で返す ----------
# 2026-07-23: as_completed()の完了順ではなく、futsリスト(=jobs=投入順)の順序で
# 結果を集めることを確認する。最初に投入するm1のジョブをthreading.Eventで
# 意図的にブロックし、最後に投入するm3が先に完了してイベントをセットするよう
# 強制することで、完了順(m3,m2,m1相当)とsubmission順(m1,m2,m3)が確実に食い違う
# 状況を作る。さらに、m1が実際にイベントを受信できた(timeoutしなかった)ことを
# 確認することで、「全ジョブ投入 → その後に結果収集」という並行性
# (apply_high_vram_profileが有効化する96GB構成でのwall-clock短縮の前提)が
# 壊れていないことも同時に検証する。
import threading as _threading

_order_flags = {}
_order_event = _threading.Event()


def _fake_proposal_ordered(model, question, reference, issue=None, history=None):
    if model == "m1":
        # 提出順では最初のジョブだが、m3がイベントをセットするまで完了しない
        # → 完了順ではm1が最後になる。もし「全ジョブ投入前に結果収集」に
        # 退行していればm3のジョブがまだ投入されておらずここでtimeoutする。
        _order_flags["m1_got_event"] = _order_event.wait(timeout=5)
    elif model == "m3":
        _order_event.set()
    return model, f"ans-{model}"


_orig_gsp_ord = f.get_single_proposal
_orig_parallel_flag_2 = f.PARALLEL_PROPOSERS
f.get_single_proposal = _fake_proposal_ordered
try:
    f.PARALLEL_PROPOSERS = True
    _out_ordered = f.get_proposals(["m1", "m2", "m3"], "Q?")
finally:
    f.get_single_proposal = _orig_gsp_ord
    f.PARALLEL_PROPOSERS = _orig_parallel_flag_2
check("prop(parallel): 全ジョブ投入後に収集される(m1がイベントを受信できた)",
      _order_flags.get("m1_got_event") is True)
check("prop(parallel): 完了順(m1が最後)ではなくsubmission順(m1,m2,m3)で返る",
      _out_ordered == [("m1", "ans-m1"), ("m2", "ans-m2"), ("m3", "ans-m3")])

# ---------- コード実行検証 ----------
check("code: python フェンス抽出", f.extract_code("x\n```python\nprint(1)\n```\ny") == "print(1)\n")
check("code: タグ無しフェンスも拾う", f.extract_code("```\nx = 1\n```") == "x = 1\n")
check("code: コード無しは None", f.extract_code("no code here") is None)

# 2026-07-22 回帰: 非python フェンス(```json 等)が先行しても、その閉じフェンスを
# 開始フェンスと誤認してブロック間のプロースを「コード」として誤抽出しないこと。
check(
    "code: jsonブロックの後のpythonブロックを正しく抽出",
    f.extract_code(
        "Here:\n```json\n{\"a\": 1}\n```\nNow code:\n```python\nprint(2+2)\n```"
    ) == "print(2+2)\n",
)
check(
    "code: text/outputブロックの後のpythonブロックを正しく抽出",
    f.extract_code(
        "```text\nsome output\n```\n説明\n```output\nmore output\n```\n"
        "```python\nprint(3+3)\n```"
    ) == "print(3+3)\n",
)
check(
    "code: 非pythonブロックのみなら None",
    f.extract_code("```json\n{\"a\": 1}\n```") is None,
)

# 2026-07-22 回帰: extract_code の言語タグ比較を「info string 全体」ではなく
# 「最初の空白区切りトークン」で行うよう修正（CommonMark の info string 仕様。
# iteration 7 のブロック選択修正、iteration 18 の _extract_code_for_output 修正に
# 続く、同じ誤抽出クラスの3件目）。装飾/メタデータ付きの python フェンスが誤って
# 読み飛ばされないことを確認する。
check(
    "code: python+装飾タグ({.line-numbers})も正しく抽出",
    f.extract_code("```python {.line-numbers}\nprint(1)\n```") == "print(1)\n",
)
check(
    "code: python+装飾タグ(title=...)も正しく抽出",
    f.extract_code("```python title=\"sol.py\"\nx = 1\n```") == "x = 1\n",
)
check(
    "code: py3タグは受理集合外のためNone(広げすぎていないことの確認)",
    f.extract_code("```py3\ncode\n```") is None,
)
check(
    "code: 先行する装飾非pythonフェンス(json {.foo})の後のpythonブロックを正しく抽出",
    f.extract_code(
        "```json {.foo}\n{\"a\":1}\n```\n```python\ncode\n```"
    ) == "code\n",
)
check(
    "code: ハイフン複合タグ(python-repl)は空白区切りが無いため受理されない",
    f.extract_code("```python-repl\ncode\n```") is None,
)

# 2026-07-22: _extract_code_for_output (_save_as_code が使うファイル出力用抽出) に
# iteration-7 の extract_code と同じ誤抽出クラスの修正を適用した回帰テスト。
check(
    "code_out: jsonブロックの後のbareフェンスからpythonを正しく抽出",
    f._extract_code_for_output(
        "```json\n{\"a\": 1}\n```\n```\ndef f():\n    return 42\n```", ".py"
    ) == "def f():\n    return 42\n",
)
check(
    "code_out: textブロック+ブロック間プロースの後のbareフェンスを正しく抽出",
    f._extract_code_for_output(
        "```text\nsome output\n```\n説明のプロース\n```\ndef g():\n    return 1\n```",
        ".py",
    ) == "def g():\n    return 1\n",
)
check(
    "code_out: python3タグ単体ブロックを抽出",
    f._extract_code_for_output("```python3\nprint('hi')\n```", ".py")
    == "print('hi')\n",
)
check(
    "code_out: jsonブロックに続くpython3ブロックを正しく抽出",
    f._extract_code_for_output(
        "```json\n{\"x\": 1}\n```\n```python3\nprint('hi')\n```", ".py"
    ) == "print('hi')\n",
)
check(
    "code_out: 単一pythonブロック(回帰・従来通り)",
    f._extract_code_for_output("```python\nprint(1)\n```", ".py") == "print(1)\n",
)
check(
    "code_out: 単一bareフェンス(回帰・従来通り)",
    f._extract_code_for_output("```\nx = 1\n```", ".py") == "x = 1\n",
)
check(
    "code_out: フェンス無しはマークダウン見出し除去にフォールバック(回帰・従来通り)",
    f._extract_code_for_output("# Title\nSome text\n# Another\nMore text", ".py")
    == "Some text\nMore text",
)
check(
    "code_out: 唯一のフェンスが別言語(c)の実コードなら保守的に採用(スキップリストで飲み込まない)",
    f._extract_code_for_output("```c\nint main(){return 0;}\n```", ".py")
    == "int main(){return 0;}\n",
)

# 2026-07-23 回帰: フェンス無しフォールバックが l.startswith("#") で '#' 始まりの
# 行を無条件に全削除していたため、--out file.<ext> で保存される素の(フェンス無し)
# コードに含まれる #include/#define/#pragma、シェバン行、Rust属性が「見出し」
# 扱いされて誤って削除されていたバグの修正確認 (iteration 7/18/28/29 が対処した
# フェンス選択側の兄弟分岐)。CommonMark の ATX heading (^#{1,6}(?:\s|$)) にのみ
# 一致する行を除去するよう修正した。
check(
    "code_out: フェンス無しのC #includeは見出し扱いされず保持される",
    "#include <stdio.h>"
    in f._extract_code_for_output(
        "#include <stdio.h>\nint main(void){return 0;}", ".c"
    ),
)
check(
    "code_out: フェンス無しの#define/#pragma/シェバン/Rust属性はすべて保持される",
    f._extract_code_for_output(
        "#!/usr/bin/env python\n"
        "#define X 1\n"
        "#pragma once\n"
        "#[derive(Debug)]\n"
        "code_line()",
        ".py",
    )
    == "#!/usr/bin/env python\n#define X 1\n#pragma once\n#[derive(Debug)]\ncode_line()",
)
check(
    "code_out: 7個以上の連続'#'はATX見出しとして無効なため保持される",
    "####### note"
    in f._extract_code_for_output("####### note\ncode_line()", ".py"),
)
check(
    "code_out: フェンス無しでも真のATX見出し(# Title/## Section)は従来通り除去される",
    f._extract_code_for_output(
        "# Title\ncode_a()\n## Section\ncode_b()", ".py"
    )
    == "code_a()\ncode_b()",
)

# 2026-07-22: iteration 28 の extract_code 修正 (info string は最初の空白区切り
# トークンのみを言語タグとする) を _extract_code_for_output にも追随適用した
# 回帰・新規テスト。装飾付き info string (```python title="sol.py" や
# ```python {.line-numbers}) が tier-1 の対象言語一致から漏れたり、装飾付き
# 非コードフェンス(```json {.line-numbers})が _NON_CODE_TAGS と不一致になって
# tier-3 で誤って採用されたりしないことを確認する。
check(
    "code_out: title=属性で装飾されたpythonフェンスもtier-1(対象言語一致)で抽出",
    f._extract_code_for_output(
        "```python title=\"sol.py\"\nprint(1)\n```", ".py"
    ) == "print(1)\n",
)
check(
    "code_out: {.line-numbers}属性で装飾されたpythonフェンスもtier-1で抽出",
    f._extract_code_for_output(
        "```python {.line-numbers}\nprint(2)\n```", ".py"
    ) == "print(2)\n",
)
check(
    "code_out: 装飾付き非コードフェンス(json {.foo})はNON_CODE_TAGSとして読み飛ばされ、後続pythonを抽出",
    f._extract_code_for_output(
        "```json {.foo}\n{\"a\":1}\n```\n\n```python\nprint(3)\n```", ".py"
    ) == "print(3)\n",
)

# 2026-07-23: _CODE_EXTENSIONS (25拡張子) と _extract_code_for_output 内 lang_map
# (13拡張子のみ) の同期漏れ修正の回帰・新規テスト。lang_map に無い残り12拡張子
# (.jsx .tsx .mjs .h .hpp .kt .swift .php .bat .ps1 .m .jl) は langs が空集合に
# なり tier-1 (対象言語タグ一致) が絶対に発火せず、非対象言語のブロック
# (例: 使い方説明の```bash)が先行すると tier-3 でそれを誤って採用してしまう
# (iteration 7/18/28/29/56 と同じブロック誤選択バグクラス)。実例として挙げられた
# 「```bash の npm install の後に本命の```jsx」で、修正前は 'npm install\n' が
# 返っていたことを確認しつつ、12拡張子それぞれの代表ケースで tier-1 が正しく
# 発火することを検証する。
check(
    "code_out: 先行するbashブロックの後のjsxブロックを正しく抽出(修正前はnpm installを誤抽出)",
    "export default"
    in f._extract_code_for_output(
        "```bash\nnpm install\n```\n"
        "```jsx\nexport default function App(){return null}\n```",
        ".jsx",
    )
    and "npm install"
    not in f._extract_code_for_output(
        "```bash\nnpm install\n```\n"
        "```jsx\nexport default function App(){return null}\n```",
        ".jsx",
    ),
)
check(
    "code_out: 先行するbashブロックの後のtsx(tsxタグ)ブロックを正しく抽出",
    f._extract_code_for_output(
        "```bash\nnpm install\n```\n```tsx\nconst x: number = 1;\n```", ".tsx"
    ) == "const x: number = 1;\n",
)
check(
    "code_out: 先行するbashブロックの後のtsx(typescriptタグ)ブロックを正しく抽出",
    f._extract_code_for_output(
        "```bash\nnpm install\n```\n```typescript\nconst y: number = 2;\n```", ".tsx"
    ) == "const y: number = 2;\n",
)
check(
    "code_out: 先行するshブロックの後のmjs(javascriptタグ)ブロックを正しく抽出",
    f._extract_code_for_output(
        "```sh\nnode -v\n```\n```javascript\nexport const z = 1;\n```", ".mjs"
    ) == "export const z = 1;\n",
)
check(
    "code_out: 先行するbashブロックの後のh(cタグ)ブロックを正しく抽出",
    f._extract_code_for_output(
        "```bash\ngcc --version\n```\n```c\nint add(int a, int b){return a+b;}\n```",
        ".h",
    ) == "int add(int a, int b){return a+b;}\n",
)
check(
    "code_out: 先行するbashブロックの後のhpp(cppタグ)ブロックを正しく抽出",
    f._extract_code_for_output(
        "```bash\ncmake .\n```\n```cpp\nint add(int a,int b){return a+b;}\n```",
        ".hpp",
    ) == "int add(int a,int b){return a+b;}\n",
)
check(
    "code_out: 先行するbashブロックの後のphp(phpタグ)ブロックを正しく抽出",
    f._extract_code_for_output(
        "```bash\ncomposer install\n```\n```php\necho 1;\n```", ".php"
    ) == "echo 1;\n",
)
check(
    "code_out: 先行するbashブロックの後のkt(kotlinタグ)ブロックを正しく抽出",
    f._extract_code_for_output(
        "```bash\ngradle build\n```\n```kotlin\nfun main() {}\n```", ".kt"
    ) == "fun main() {}\n",
)
check(
    "code_out: 先行するbashブロックの後のswift(swiftタグ)ブロックを正しく抽出",
    f._extract_code_for_output(
        "```bash\nswift build\n```\n```swift\nprint(\"hi\")\n```", ".swift"
    ) == "print(\"hi\")\n",
)
check(
    "code_out: 先行するshブロックの後のps1(powershellタグ)ブロックを正しく抽出",
    f._extract_code_for_output(
        "```sh\necho setup\n```\n```powershell\nGet-Process\n```", ".ps1"
    ) == "Get-Process\n",
)
check(
    "code_out: 先行するbashブロックの後のbat(batタグ)ブロックを正しく抽出",
    f._extract_code_for_output(
        "```bash\necho setup\n```\n```bat\n@echo off\n```", ".bat"
    ) == "@echo off\n",
)
check(
    "code_out: 先行するbashブロックの後のjl(juliaタグ)ブロックを正しく抽出",
    f._extract_code_for_output(
        "```bash\njulia --version\n```\n```julia\nprintln(\"hi\")\n```", ".jl"
    ) == "println(\"hi\")\n",
)
check(
    "code_out: 先行するbashブロックの後のm(matlabタグ)ブロックを正しく抽出",
    f._extract_code_for_output(
        "```bash\necho setup\n```\n```matlab\ndisp('hi')\n```", ".m"
    ) == "disp('hi')\n",
)

# 回帰: 新規マッピング拡張子でも単一ブロック/bareフェンス/フェンス無しATXフォール
# バックは従来のtier-2/tier-3/フェンス無し挙動のまま変わらないことを確認する。
check(
    "code_out: phpの単一ブロック(回帰・従来通り)",
    f._extract_code_for_output("```php\necho 1;\n```", ".php") == "echo 1;\n",
)
check(
    "code_out: jsx向けでもbareフェンス(タグ無し)はtier-2で従来通り抽出",
    f._extract_code_for_output("```\nconst a = 1;\n```", ".jsx") == "const a = 1;\n",
)
check(
    "code_out: swift向けでもフェンス無しはATX見出し除去フォールバックに従来通り落ちる",
    f._extract_code_for_output(
        "# Title\nlet a = 1\n## Section\nlet b = 2", ".swift"
    ) == "let a = 1\nlet b = 2",
)

# End-to-end: _save_as_code (--out <newly-mapped-suffix> の実書き込み経路) でも
# tier-1 が発火し、先行するbashブロックではなく対象言語ブロックが書き出される
# ことを一時ファイルへの実書き込みで確認する。Ollama/ネットワーク呼び出しは無い。
import tempfile as _cfo_tempfile
import pathlib as _cfo_pathlib
with _cfo_tempfile.TemporaryDirectory() as _cfo_dir:
    _cfo_root = _cfo_pathlib.Path(_cfo_dir)
    _cfo_out = _cfo_root / "app.tsx"
    f._save_as_code(
        _cfo_out,
        "```bash\nnpm install\n```\n```typescript\nexport const x: number = 1;\n```",
    )
    _cfo_content = _cfo_out.read_text(encoding="utf-8")
    check(
        "code_out/_save_as_code: .tsxへの実書き込みは先行bashではなく対象言語ブロックになる",
        _cfo_content == "export const x: number = 1;\n\n"
        and "npm install" not in _cfo_content,
    )

ok, out = f.run_python("print('hello_runner')")
check("code: 実行成功", ok and "hello_runner" in out)
ok, out = f.run_python("raise ValueError('boom')")
check("code: 例外を検知して traceback を返す", (not ok) and "boom" in out)
ok, out = f.run_python("while True:\n    pass", timeout=2)
check("code: 無限ループはタイムアウト", (not ok) and "TIMEOUT" in out)

# 2026-07-22 回帰: run_python は子プロセスの stdin を DEVNULL にするので、
# input() を呼ぶコードはタイムアウトまでハングせず即座に EOFError で失敗する
# （親の stdin を継承していた旧挙動では repl() の対話入力を子に奪われたり、
#   TTY/pipe/closed の違いで非決定的に振る舞ったりしていた）。
# timeout は「タイムアウトまで待っていない」ことを示すため意図的に長め(30s)にする。
ok_input, out_input = f.run_python("data = input()\nprint(data)", timeout=30)
check(
    "code: input()はDEVNULL stdinによりEOFErrorで即失敗しTIMEOUTしない",
    (not ok_input) and ("EOFError" in out_input) and ("TIMEOUT" not in out_input),
)

# regression guard: DEVNULL stdin が通常コードの成功経路を邪魔しないこと
ok_normal, out_normal = f.run_python("print('no_input_needed_here')")
check(
    "code: input()を使わない通常コードはDEVNULL化後も正常に成功",
    ok_normal and "no_input_needed_here" in out_normal,
)

# regression guard: stdout_only=True も DEVNULL化後、成功時はstdoutのみを返す(iteration 4挙動)
ok_normal_only, out_normal_only = f.run_python("print('clean_stdout_only')", stdout_only=True)
check(
    "code: stdout_only=Trueの成功経路はDEVNULL化後もstdoutのみ",
    ok_normal_only and out_normal_only.strip() == "clean_stdout_only",
)

# stdout_only: 既定(False)は stdout+stderr 結合のまま(バイトレベルで不変)。
# stdout_only=True かつ成功時は stdout のみ返し、stderr の警告文で末尾行が汚染されない。
_warn_code = (
    "import sys\n"
    "print('a warning', file=sys.stderr)\n"
    "print('42')\n"
)
ok_default, out_default = f.run_python(_warn_code)
check("code: stdout_only既定Falseはstderrも含む", ok_default and "a warning" in out_default
      and "42" in out_default)
ok_only, out_only = f.run_python(_warn_code, stdout_only=True)
check("code: stdout_only=Trueはstdoutのみ・最終行が正しい値",
      ok_only and "a warning" not in out_only and out_only.splitlines()[-1].strip() == "42")

# stdout_only=True でも失敗時(returncode!=0)は traceback 込みの結合出力を返す
# （code-repair loop がエラー内容を見えるようにするための回帰ガード）。
ok_fail_only, out_fail_only = f.run_python("raise ValueError('boom_stdout_only')", stdout_only=True)
check("code: stdout_only=Trueでも失敗時はtraceback付き結合出力",
      (not ok_fail_only) and "boom_stdout_only" in out_fail_only)

# _sc_sample の PoT 分岐: 生成コードが正しい答えを stdout に、警告を stderr に出すケースで、
# run_python(stdout_only=True) により警告文が投票を汚染しないことを確認する
# （2026-07-21 修正の回帰ガード。修正前は out.splitlines()[-1] が stderr の警告行になり得た）。
_orig_ask_pot = f.ask


def _fake_ask_pot(model, messages, temperature, think=None, fmt=None,
                   label=None, num_predict=None, num_ctx=None):
    return (
        "考え方の説明です。\n```python\n"
        "import sys\n"
        "print('RuntimeWarning: something noisy', file=sys.stderr)\n"
        "print(7)\n"
        "```\n"
    )


try:
    f.ask = _fake_ask_pot
    _pot_ans, _pot_text = f._sc_sample("m1", "1+2+4=?", "math", pot=True)
finally:
    f.ask = _orig_ask_pot
check("sc: PoT stdout_onlyで警告行ではなく印字された答えが投票になる",
      _pot_ans == f.normalize_answer("7"))

# ---------- _sc_sample: __ERROR__ センチネルは投票を汚染せずそのまま素通しする (2026-07-23) ----------
# ask() は失敗時 '__ERROR__: HTTP Error 500: Internal Server Error' のような文字列を返す
# (line ~1079)。_sc_sample の L2558 (`if raw.startswith("__ERROR__"): return None, raw`) は
# これを extract_final_answer/PoT実行より前段で弾いている。弾かなければ math の最終数値
# フォールバック(extract_final_answer, line ~2497)がエラーメッセージ中の '500' を確信ありの
# 1票として誤採用しうる。同種ガードは ask()自身(iter9)・_critic_judge/second_opinion(iter15)・
# _arbitrate(iter20)では回帰確認済みだったが、自己一貫性投票(gotcha#7)の最小単位である
# _sc_sample 自身は iteration 4 で PoT の happy path しか検証されておらず、非PoT分岐・mcq分岐・
# PoTのガード早期リターン(サブプロセス不起動)は無防備だった。ここで塞ぐ。
_orig_ask_scerr = f.ask
_orig_extract_code_scerr = f.extract_code
_orig_run_python_scerr = f.run_python

try:
    # (1) 非PoT math: __ERROR__ は数値 '500' に化けず、素の (None, raw) を返す
    f.ask = lambda *a, **k: "__ERROR__: HTTP Error 500: Internal Server Error"
    _ans_math_err, _text_math_err = f._sc_sample("m1", "1+1=?", "math", pot=False)
    check("sc-err: 非PoT mathの__ERROR__は'500'を誤採用せずanswer=None",
          _ans_math_err is None)
    check("sc-err: 非PoT mathの__ERROR__はtextに生のエラー文字列をそのまま返す",
          _text_math_err == "__ERROR__: HTTP Error 500: Internal Server Error")

    # (2) 非PoT mcq: エラー文中にA-E文字や数字が含まれていても選択肢を誤採用しない
    f.ask = lambda *a, **k: "__ERROR__: HTTP Error 400 for model A"
    _ans_mcq_err, _text_mcq_err = f._sc_sample("m1", "Which is correct?", "mcq", pot=False)
    check("sc-err: 非PoT mcqの__ERROR__は選択肢文字/数字を誤採用せずanswer=None",
          _ans_mcq_err is None)

    # (3) PoT math: ガードが extract_code/run_python より前段で早期returnし、
    # サブプロセスが一切起動されないことを証明する(到達すればAssertionErrorで即検知)。
    _pot_err_touched = []

    def _extract_code_forbidden_scerr(*a, **kw):
        _pot_err_touched.append("extract_code")
        raise AssertionError("__ERROR__ガード後にextract_codeへ到達してはならない")

    def _run_python_forbidden_scerr(*a, **kw):
        _pot_err_touched.append("run_python")
        raise AssertionError("__ERROR__ガード後にrun_pythonへ到達してはならない(subprocess起動)")

    f.extract_code = _extract_code_forbidden_scerr
    f.run_python = _run_python_forbidden_scerr
    f.ask = lambda *a, **k: "__ERROR__: HTTP Error 500: Internal Server Error"
    _ans_pot_err, _text_pot_err = f._sc_sample("m1", "1+1=?", "math", pot=True)
    check("sc-err: PoT分岐でも__ERROR__は(None, raw)を返す",
          _ans_pot_err is None
          and _text_pot_err == "__ERROR__: HTTP Error 500: Internal Server Error")
    check("sc-err: PoT分岐の__ERROR__はextract_code/run_pythonへ到達しない(subprocess不起動)",
          not _pot_err_touched)
finally:
    f.ask = _orig_ask_scerr
    f.extract_code = _orig_extract_code_scerr
    f.run_python = _orig_run_python_scerr

# regression: 正常な非PoT math応答は従来どおり\boxed{}から答えが採れる(ガードの過検知なし)
try:
    f.ask = lambda *a, **k: "計算しました。\\boxed{42}"
    _ans_ok_scerr, _text_ok_scerr = f._sc_sample("m1", "6*7=?", "math", pot=False)
finally:
    f.ask = _orig_ask_scerr
check("sc-err: 通常応答は__ERROR__ガードに過検知されず'42'を採票",
      _ans_ok_scerr == "42" and _text_ok_scerr == "計算しました。\\boxed{42}")

# ---------- _sc_sample: PoT分岐の4つの棄却ガードは無投票を返す (2026-07-23) ----------
# _sc_sample の PoT 分岐 (line ~2587-2598) には「怪しい PoT サンプルを黙って捨て、
# SC投票を汚染しない」ための4つの早期returnガードがある:
#   (a) extract_code(text) が None（```python フェンス無し）
#   (b) run_python の ok が False（実行時エラー=トレースバック）
#   (c) 実行はできたが stdout が空
#   (d) stdout 最終行が80文字を超える（自由記述の暴走を数値解答として誤採用しない）
# いずれも (None, text) を返し無投票となる契約（「無投票 > 誤投票」の精度優先方針、
# gotcha#7）。iteration 4 は PoT の happy path のみ、iteration 52 は __ERROR__ センチ
# ネルのみを検証しており（52自身のコメントが認める通り）、この4ガードはこれまで無防備
# だった。ここでは f.ask のみをモックして _sc_sample を直接呼ぶ（モデル/ネットワーク
# 呼び出しは一切発生しない）。(b)(c)(d) は run_python に本物の軽量・決定的なローカル
# subprocess（1/0 や print(...) 程度）を実際に実行させる — iteration 4/46/52 と同じ流儀。
_orig_ask_potguard = f.ask

# (a) コードフェンス無し: extract_code が None を返し、run_python は起動されない
#     （到達すれば forbidden 関数が AssertionError を送出し、即座に検知できる）。
_pot_nocode_touched = []


def _run_python_forbidden_potguard_a(*a, **kw):
    _pot_nocode_touched.append("run_python")
    raise AssertionError("コードフェンス無しなのにrun_pythonへ到達してはならない(subprocess起動)")


_orig_run_python_potguard_a = f.run_python
try:
    f.run_python = _run_python_forbidden_potguard_a
    f.ask = lambda *a, **k: "説明だけの回答です。コードブロックはありません。答えはたぶん42。"
    _ans_nocode, _text_nocode = f._sc_sample("m1", "1+1=?", "math", pot=True)
finally:
    f.ask = _orig_ask_potguard
    f.run_python = _orig_run_python_potguard_a
check("sc-pot: (a)コードフェンス無し→extract_code=Noneでanswer=None(無投票)",
      _ans_nocode is None)
check("sc-pot: (a)コードフェンス無し→run_pythonへ到達しない(subprocess不起動)",
      not _pot_nocode_touched)

# (b) 実行時エラー(ZeroDivisionError): run_python の ok=False → answer=None。
#     注意: ok=False 分岐の _sc_sample は「text + run_python出力」を連結せず、素の
#     text（LLMが書いたコード込みの生テキスト）をそのまま返す契約（連結は成功時のみ）。
#     よって「tracebackの数字が誤ってanswerに混入しない」ことは _sc_sample の戻り値
#     answer が None であること自体で保証される。ここではまず run_python を直接呼び、
#     1/0 の実行が本当に失敗し(ok=False)、その出力(traceback)に数字が実在すること
#     （テスト前提の健全性確認＝ガードが無ければ数字が採用され得た状況であること）を
#     確認したうえで、_sc_sample がそれを answer=None として無投票に倒すことを検証する。
_ok_sanity_execfail, _out_sanity_execfail = f.run_python("1 / 0", stdout_only=True)
check("sc-pot: (b)前提確認: 1/0の実行は失敗し(ok=False)traceback出力に数字が実在する",
      (not _ok_sanity_execfail) and any(ch.isdigit() for ch in _out_sanity_execfail))
try:
    f.ask = lambda *a, **k: "計算コードです。\n```python\n1 / 0\n```\n"
    _ans_execfail, _text_execfail = f._sc_sample("m1", "1/0 は何？", "math", pot=True)
finally:
    f.ask = _orig_ask_potguard
check("sc-pot: (b)実行失敗(ZeroDivisionError, ok=False)はtraceback中の数字を誤採用せずanswer=None",
      _ans_execfail is None)

# (c) 実行は成功するが stdout が空: ok=True だが out=="" → answer=None
try:
    f.ask = lambda *a, **k: "何も出力しないコードです。\n```python\nx = 1 + 1\n```\n"
    _ans_emptyout, _text_emptyout = f._sc_sample("m1", "1+1=?", "math", pot=True)
finally:
    f.ask = _orig_ask_potguard
check("sc-pot: (c)stdoutが空(print無し)はanswer=None(無投票)",
      _ans_emptyout is None)

# (d) stdout 最終行が80文字を超える: 自由記述の暴走を「解答」として誤採用しない
_long_line = "9" * 90
try:
    f.ask = lambda *a, **k: "長い出力のコードです。\n```python\nprint('" + _long_line + "')\n```\n"
    _ans_toolong, _text_toolong = f._sc_sample("m1", "何か計算して", "math", pot=True)
finally:
    f.ask = _orig_ask_potguard
check("sc-pot: (d)前提確認: 生成した最終行は80文字を超える",
      len(_long_line) > 80)
check("sc-pot: (d)stdout最終行が80文字超はanswer=None(無投票、暴走出力を誤採用しない)",
      _ans_toolong is None)

# regression: 短い答えを最終行に印字 + 無関係なstderr警告、という正常系は上記4ガード
# 追加後も従来通り採票される（過検知が無いことの確認）。stdout_only=Trueによるstderr
# 除去自体は2026-07-21の別テスト（line ~2028）で既に確認済みだが、ここでは本セクション
# 自身のモック配線内でも happy path が壊れていないことを直接ロックする。
try:
    f.ask = lambda *a, **k: (
        "考え方の説明です。\n```python\n"
        "import sys\n"
        "print('DeprecationWarning: ignore me', file=sys.stderr)\n"
        "print(99)\n"
        "```\n"
    )
    _ans_happy_potguard, _text_happy_potguard = f._sc_sample("m1", "regression", "math", pot=True)
finally:
    f.ask = _orig_ask_potguard
check("sc-pot: regression: 短い答え+無関係なstderr警告は従来通り採票される(4ガード追加後も過検知なし)",
      _ans_happy_potguard == f.normalize_answer("99"))

# ---------- _sc_sample: PoT分岐の最終行\boxed{}ラッパをunwrapして票割れを防ぐ (2026-07-25) ----------
# CoT分岐(直下のextract_final_answer)は\boxed{}を剥がしてからnormalize_answerへ渡すが、PoT分岐は
# stdout最終行(iteration 4)をnormalize_answerへ素通しするだけで\boxed{}は非対応だった。数学寄りの
# モデルがコード内でprint(f'\\boxed{{{ans}}}')のように反射的に答えを\boxed{}で包むと、最終行が
# 丸ごと"\boxed{42}"になりnormalize_answerを素通りし、素の"42"というCoT側の投票クラスとは別クラスに
# 割れてgotcha #7の自己整合性投票（PoTでCoTの計算ミスを裏取り/上書きする仕組み）でPoT票が丸ごと
# 死票化していた。ここでは f.ask のみをモックし、run_pythonには実際の軽量・決定的なローカル
# subprocess（print(...)のみ）を実行させる（iteration 4/46/52/61と同じ流儀。Ollama/ネットワーク
# 呼び出しは一切発生しない）。値を捏造しない・既存票を壊さない「票を拾えるところだけ拾う」方針は
# iteration 11/23/25/78/122/134/136/140/148の票救出系修正列と同じ。
def _make_ask_pot_print(stdout_value):
    """f.ask差し替え用: 生成コードがstdout_valueをそのまま1行printするPoTサンプルを模す。
    repr()経由で埋め込むためbackslash等の手動エスケープが不要（誤エスケープでの事故を避ける）。"""
    _code = "print(" + repr(stdout_value) + ")"
    return lambda *a, **k: "説明。\n```python\n" + _code + "\n```\n"


_orig_ask_potbox = f.ask

# (1) \boxed{42} → '42' を採票し、素のCoT'42'票とanswers_equivalent/vote_answersで併合される
try:
    f.ask = _make_ask_pot_print("\\boxed{42}")
    _ans_potbox42, _text_potbox42 = f._sc_sample("m1", "6*7=?", "math", pot=True)
finally:
    f.ask = _orig_ask_potbox
check("sc-pot-boxed: 最終行が\\boxed{42}のPoTサンプルは'42'を採票する",
      _ans_potbox42 == f.normalize_answer("42"))
check("sc-pot-boxed: \\boxed{42}のPoT票と素のCoT'42'票はanswers_equivalentで同一視される",
      f.answers_equivalent(_ans_potbox42, "42"))
_classes_potbox42 = f.vote_answers(["42", _ans_potbox42])
check("sc-pot-boxed: vote_answersでもCoT'42'票とPoTの\\boxed{42}票が単一クラス(2票)に併合される",
      len(_classes_potbox42[2]) == 1 and _classes_potbox42[1] == 2)

# (2) \boxed{1/2} → '1/2' を採票し、CoTの'0.5'/'1/2'とFractionファストパスで併合される
try:
    f.ask = _make_ask_pot_print("\\boxed{1/2}")
    _ans_potboxfrac, _text_potboxfrac = f._sc_sample("m1", "1/2 は？", "math", pot=True)
finally:
    f.ask = _orig_ask_potbox
check("sc-pot-boxed: \\boxed{1/2}のPoTサンプルは'1/2'を採票する",
      _ans_potboxfrac == f.normalize_answer("1/2"))
check("sc-pot-boxed: \\boxed{1/2}のPoT票と素のCoT'0.5'票はFractionファストパスで同一視される",
      f.answers_equivalent(_ans_potboxfrac, "0.5"))
check("sc-pot-boxed: \\boxed{1/2}のPoT票と素のCoT'1/2'票も同一視される",
      f.answers_equivalent(_ans_potboxfrac, "1/2"))

# (3) \boxed{-5} → '-5' を採票する
try:
    f.ask = _make_ask_pot_print("\\boxed{-5}")
    _ans_potboxneg, _text_potboxneg = f._sc_sample("m1", "-5 は？", "math", pot=True)
finally:
    f.ask = _orig_ask_potbox
check("sc-pot-boxed: \\boxed{-5}のPoTサンプルは'-5'を採票する",
      _ans_potboxneg == f.normalize_answer("-5"))

# regression(バイト単位で不変): \boxed{}を含まない最終行は従来のnormalize_answer(生の行)と
# 完全に同じ結果を返す（本修正が非boxedケースを一切変えないことの直接ロック）
for _rv in ("42", "1/2", "-5", "0.5", "1,234"):
    try:
        f.ask = _make_ask_pot_print(_rv)
        _ans_rv, _text_rv = f._sc_sample("m1", "regression-nonboxed", "math", pot=True)
    finally:
        f.ask = _orig_ask_potbox
    check("sc-pot-boxed: regression: 非boxedの'%s'出力はnormalize_answer(ans)と完全一致(変化なし)" % _rv,
          _ans_rv == f.normalize_answer(_rv))

# regression: 未閉鎖の\boxed{（波括弧が閉じていない）はextract_boxedがNoneを返し、
# クラッシュせず・値も捏造せず、従来通りnormalize_answer(生の行)へフォールスルーする
# （iteration 11/23で確立した「閉じていなければNone」という安全側判定はそのまま）
try:
    f.ask = _make_ask_pot_print("\\boxed{")
    _ans_potboxunterm, _text_potboxunterm = f._sc_sample("m1", "壊れたboxed", "math", pot=True)
finally:
    f.ask = _orig_ask_potbox
check("sc-pot-boxed: regression: 未閉鎖の\\boxed{はクラッシュせずnormalize_answer(生の行)にフォールスルーする",
      _ans_potboxunterm == f.normalize_answer("\\boxed{"))

# regression: 閉じてはいるが中身が空の\boxed{}もextract_boxedがNoneを返し（iteration 25で
# 確立した「空boxedは無投票」判定）、同様にnormalize_answer(生の行)へフォールスルーする
try:
    f.ask = _make_ask_pot_print("\\boxed{}")
    _ans_potboxempty, _text_potboxempty = f._sc_sample("m1", "空boxed", "math", pot=True)
finally:
    f.ask = _orig_ask_potbox
check("sc-pot-boxed: regression: 空の\\boxed{}はクラッシュせずnormalize_answer(生の行)にフォールスルーする",
      _ans_potboxempty == f.normalize_answer("\\boxed{}"))

# regression: len>80ガードはunwrap後の「投票される値」に対して適用される。\boxed{42}の前後に
# 80字を優に超えるパディングが付いた行でも、unwrap後の中身'42'が短ければ採票される
# （旧実装は生の行の長さで即棄却していたため、この場合は無投票にすり替わっていた）。
_padded_boxed = "The answer is \\boxed{42} " + "x" * 90
try:
    f.ask = _make_ask_pot_print(_padded_boxed)
    _ans_potboxpad, _text_potboxpad = f._sc_sample("m1", "パディング", "math", pot=True)
finally:
    f.ask = _orig_ask_potbox
check("sc-pot-boxed: 前提確認: パディング付き生の行は80文字を超える",
      len(_padded_boxed) > 80)
check("sc-pot-boxed: len>80ガードはunwrap後の値に適用され、\\boxed{42}前後の長いパディングでも'42'を採票する",
      _ans_potboxpad == f.normalize_answer("42"))

# regression: unwrap後の値そのものが80文字を超える場合は、従来通り無投票(None)になる
_long_boxed_val = "9" * 90
try:
    f.ask = _make_ask_pot_print("\\boxed{" + _long_boxed_val + "}")
    _ans_potboxlong, _text_potboxlong = f._sc_sample("m1", "長いboxed中身", "math", pot=True)
finally:
    f.ask = _orig_ask_potbox
check("sc-pot-boxed: unwrap後の値自体が80文字を超える場合はanswer=None(無投票、暴走出力を誤採用しない)",
      _ans_potboxlong is None)

check("sc-pot-boxed: テスト後にf.askが元へ復元されている",
      f.ask == _orig_ask_potbox)

# ---------- _representative_text: SC勝者クラスの代表テキスト選出 (2026-07-23) ----------
# solve_verifiable の最終return直前(line ~2788)でres['text']を決める、自己一貫性投票
# 経路(gotcha#7)で最後まで直接テストされていなかったヘルパー。選出順序は
# 「CoT(pot=False)優先 → 無ければ最初に一致したPoTサンプル → それも無ければ全サンプル中
# 最長」。一致判定は文字列==ではなくanswers_equivalent経由（表記違いの同値も拾う）。
# iteration 2 のchangelogで「代表テキストが敗者候補の主張になりうる」懸念が挙がって
# いた箇所（_arbitrateのrep選出とは別に、_representative_text自身の契約）を直接ロックする。
def _rt_sample(answer, text, pot=False, model="m1"):
    return {"answer": answer, "text": text, "pot": pot, "model": model}


# (1) CoT優先: リスト上で先に出てくるPoT一致サンプルより、後から出てくるCoT一致サンプルの
#     テキストを返す（"pot でない"優先はリストの並び順に依存しないことを証明する）。
_rt1 = [
    _rt_sample("42", "POT_TEXT_EARLIER", pot=True),
    _rt_sample("42", "COT_TEXT_LATER", pot=False),
]
check("rt: CoT優先(先行PoT一致より後続CoT一致を採用、順序非依存)",
      f._representative_text(_rt1, "42") == "COT_TEXT_LATER")

# (2) PoTフォールバック: CoT一致が一つも無い場合は最初に一致したPoTサンプルのテキストを返す
_rt2 = [
    _rt_sample("7", "POT_TEXT_NOMATCH", pot=True),      # 不一致(勝者は42)
    _rt_sample("42", "POT_TEXT_FIRST_MATCH", pot=True),
    _rt_sample("42", "POT_TEXT_SECOND_MATCH", pot=True),
]
check("rt: CoT一致無し→最初に一致したPoTサンプルのテキストを返す",
      f._representative_text(_rt2, "42") == "POT_TEXT_FIRST_MATCH")

# (3) 同値判定: 勝者'0.5'に対しサンプルの答えが'1/2'(Fractionによる高速パス、math_verifyは
#     不要)でも一致とみなし、文字列==ではなくanswers_equivalent経由でマッチしていることを示す
_rt3 = [_rt_sample("1/2", "FRACTION_MATCH_TEXT", pot=False)]
check("rt: 同値判定(1/2 と 0.5、Fraction高速パス)でマッチしテキストを返す",
      f._representative_text(_rt3, "0.5") == "FRACTION_MATCH_TEXT")

# (4) 敗者除外: 勝者と一致するサンプルのテキストが返り、別の敗者候補のテキストは
#     決して返らない(iteration 2 changelogの「敗者候補の主張になりうる」懸念の直接ガード)
_rt4 = [
    _rt_sample("7", "LOSER_TEXT", pot=False),
    _rt_sample("42", "WINNER_TEXT", pot=False),
]
_rt4_result = f._representative_text(_rt4, "42")
check("rt: 勝者一致サンプルのテキストを返す(敗者除外 1/2)", _rt4_result == "WINNER_TEXT")
check("rt: 敗者候補のテキストは返らない(敗者除外 2/2)", _rt4_result != "LOSER_TEXT")

# (5) 無一致フォールバック: どのサンプルも勝者と同値でない場合は全サンプル中最長の
#     テキストを返す(max-by-len分岐)。空リストは""を返す。
_rt5 = [
    _rt_sample("7", "short", pot=False),
    _rt_sample("9", "much much longer non-matching text", pot=False),
]
check("rt: 無一致→全サンプル中最長のテキストにフォールバック",
      f._representative_text(_rt5, "999") == "much much longer non-matching text")
check("rt: 空サンプルリストは空文字列を返す", f._representative_text([], "42") == "")

# (6) answer=Noneのサンプルは一致判定の対象から除外される(最長でも誤って採用されない)。
#     ただし全体が無一致で最長フォールバックに落ちた場合は、その候補にもなりうる
#     (max-by-len分岐はanswerの有無を見ないため)。
_rt6a = [
    _rt_sample(None, "NONE_ANSWER_BUT_LONGEST_TEXT_AAAAAAAAAAAAAAAA", pot=False),
    _rt_sample("42", "SHORT_MATCH_TEXT", pot=False),
]
check("rt: answer=Noneのサンプルは一致判定から除外(最長でも採用されない)",
      f._representative_text(_rt6a, "42") == "SHORT_MATCH_TEXT")

_rt6b = [
    _rt_sample(None, "NONE_ANSWER_LONGEST_FALLBACK_TEXT_AAAAAAAAAAAAAA", pot=False),
    _rt_sample("7", "short_nomatch", pot=False),
]
check("rt: 無一致時に限り、answer=Noneのサンプルも最長フォールバックの対象になりうる",
      f._representative_text(_rt6b, "999") == "NONE_ANSWER_LONGEST_FALLBACK_TEXT_AAAAAAAAAAAAAA")

check("code: code_check 正常コードは None", f.code_check("```python\nprint(1)\n```") is None)
_issue = f.code_check("```python\n1/0\n```")
check("code: code_check 失敗はエラー要約", _issue is not None and "ZeroDivision" in _issue)
check("code: コード無し回答は None", f.code_check("plain text answer") is None)

# 2026-07-22 回帰: 先行する非pythonブロック(```json)があっても code_check が
# ブロック間のプロースではなく実物の python を検証すること（見せかけの失敗を防ぐ）。
_issue_leading_json = f.code_check(
    "```json\n{\"note\": \"ignore me\"}\n```\n```python\n1/0\n```"
)
check(
    "code: 先行jsonブロックがあってもpythonの失敗を正しく検知",
    _issue_leading_json is not None and "ZeroDivision" in _issue_leading_json,
)
check(
    "code: 先行jsonブロック+正しいpythonはNone",
    f.code_check(
        "```json\n{\"note\": \"ignore me\"}\n```\n```python\nprint(1)\n```"
    ) is None,
)

# 2026-07-22 回帰: 先行する装飾付き非codeブロック(```json {.foo})があっても、
# code_check が info string の最初のトークンで正しく python ブロックを見つけ、
# 見せかけの実行失敗("code execution FAILED")を報告しないこと。
check(
    "code: 先行装飾jsonブロック+正しいpythonはNone(見せかけの失敗なし)",
    f.code_check(
        "```json {.foo}\n{\"note\": \"ignore me\"}\n```\n```python\nprint(1)\n```"
    ) is None,
)
_issue_leading_decorated_json = f.code_check(
    "```json {.foo}\n{\"note\": \"ignore me\"}\n```\n```python\n1/0\n```"
)
check(
    "code: 先行装飾jsonブロックがあってもpythonの失敗を正しく検知",
    _issue_leading_decorated_json is not None
    and "ZeroDivision" in _issue_leading_decorated_json,
)

_good_fib = ("説明します。\n```python\n"
             "def fib(n):\n"
             "    a, b = 1, 1\n"
             "    for _ in range(n - 1):\n"
             "        a, b = b, a + b\n"
             "    return a\n\n"
             "assert fib(10) == 55\n"
             "```")
_bad_fib = "```python\ndef fib(n):\n    return n\n```"
check("eval: fib 正解コード→OK", e.grade_code_fib(_good_fib) is True)
check("eval: fib 誤りコード→NG", e.grade_code_fib(_bad_fib) is False)
check("eval: コード無し回答→NG", e.grade_code_fib("fib(10)は55です") is False)

# ---------- eval_fugu の採点 ----------
check("eval: has_num 境界", e.has_num("answer is 391.", "391") and not e.has_num("3910", "391"))
check("eval: has_num 部分数字を弾く", not e.has_num("17 and 13", "7"))
check("eval: batball 0.05のみ→OK", e.grade_batball("The ball costs $0.05.") is True)
check("eval: batball 0.05なし→NG", e.grade_batball("The bat costs $1.05.") is False)

# ---------- second_opinion のバイアス対策（PROPOSERS から除外したケース） ----------
_orig_proposers = f.PROPOSERS
_orig_second_opinion_model = f.SECOND_OPINION_MODEL
_orig_disabled_flag = f._SECOND_OPINION_DISABLED
try:
    f.PROPOSERS = ["qwen3:4b"]  # phi4-mini を除外
    f.SECOND_OPINION_MODEL = "phi4-mini"
    f._SECOND_OPINION_DISABLED = False
    ok, issue = f.second_opinion("test", "test answer")
    check("so: PROPOSERS外のモデルは ok=True で即返す", ok is True and issue == "")
    check("so: 無効化フラグがセットされる", f._SECOND_OPINION_DISABLED is True)
finally:
    f.PROPOSERS = _orig_proposers
    f.SECOND_OPINION_MODEL = _orig_second_opinion_model
    f._SECOND_OPINION_DISABLED = _orig_disabled_flag

# ---------- _critic_judge / second_opinion: __ERROR__ センチネルは ok=False (2026-07-22) ----------
# ask() が通信/モデル失敗で '__ERROR__:...' を返したとき、extract_json は None になり
# 旧実装は p.get("ok", True) で黙って ok=True（審査合格）にしてしまっていた。
# critic 呼び出し自体が失敗しているだけなのに「回答は問題なし」と誤判定するのは
# 精度優先の方針に反するため、__ERROR__ センチネルだけを ok=False に反転させる。
# 一方、空文字や非JSONの地の文（gpt-oss:20b の think 予算切れ等）は既存どおり
# ok=True 既定を維持する必要があり、それも合わせて回帰確認する。
_orig_ask = f.ask
try:
    f.ask = lambda *a, **k: "__ERROR__: simulated transport failure"
    ok, issue = f._critic_judge("q", "a", think=False)
    check("critic: __ERROR__センチネル(think=False)はok=False", ok is False and bool(issue))
    ok, issue = f._critic_judge("q", "a", think=True)
    check("critic: __ERROR__センチネル(think=True)もok=False", ok is False and bool(issue))

    f.ask = lambda *a, **k: ""
    ok, issue = f._critic_judge("q", "a", think=False)
    check("critic: 空文字は既定どおりok=True(gpt-oss think予算切れ対策を維持)", ok is True)

    f.ask = lambda *a, **k: "Looks fine to me, no issues here."
    ok, issue = f._critic_judge("q", "a", think=False)
    check("critic: 非JSONの地の文も既定どおりok=True", ok is True)
finally:
    f.ask = _orig_ask

# ---------- critique(): 2段階エスカレーションの直接テスト (iter 58) ----------
# critique() は fugu_answer() の MoA 継続判定(need_more = not ok, fugu_local.py:2929)と
# aggregate() の保険2採用判定(fugu_local.py:2245)の両方を左右するが、既存テストは
# critique 自体を lambda で丸ごと差し替えて使っており(前掲の agg 保険2テスト参照)、
# 本物の2段階ロジック(1段目: think=False+スキーマで高速判定、ok ならそこで確定。
# 2段目: 1段目が NG のときだけ think=True で再検算し、その結果を最終判定とする)は
# 一度も直接実行されていなかった。この2段階構成は 2026-07-03 のフル評価で実測された
# 偽エスカレーション(think=False critic が正答 '700' を誤って NG にし 310秒浪費)への
# 対策であり、以下の不変条件を直接ロックする:
#   (1) think=False が ok なら think=True は一切呼ばずに (True, "") で確定(高速パス短絡)
#   (2) think=False が NG のときだけ think=True 再検算が発火し、その結果(issue含む)が
#       そのまま最終判定になる(権威は think=True 側であって think=False の issue ではない)
#   (3) think=True 側が _critic_judge の __ERROR__ センチネル形(ok=False, "critic call
#       failed: ...")を返した場合でも、critique() は黙って ok=True にせず reject する
# f._critic_judge をモックし、think 引数(キーワード呼び出し)で分岐・記録することで検証する。
# ask() 自体は一切呼ばない(_critic_judge の一段上を差し替えるため、実運用の think=False/True
# ラベル分岐が critique() 側で正しく起きているかだけを見る)。
_orig_critic_judge = f._critic_judge


def _make_critique_judge(ok1, issue1, ok2=False, issue2="", calls=None):
    """think 引数(bool化)で呼び出しを分岐・記録するモック。calls に think 値を記録する。"""
    if calls is None:
        calls = []

    def _judge(question, answer, think=False):
        calls.append(bool(think))
        if think:
            return ok2, issue2
        return ok1, issue1
    _judge.calls = calls
    return _judge


try:
    # (1) think=False が ok → 即 (True, "") で確定し、think=True は呼ばれない(高速パス短絡)。
    _calls = []
    f._critic_judge = _make_critique_judge(ok1=True, issue1="unused fast-path issue",
                                            calls=_calls)
    ok, issue = f.critique("2+2?", "4")
    check("critique: think=False合格は(True,'')で確定", ok is True and issue == "")
    check("critique: think=False合格時はthink=Trueを呼ばない(高速パス短絡)",
          _calls == [False])

    # (2a) think=False NG → think=True 再検算が発火し、think=True 側が ok なら採用。
    _calls = []
    f._critic_judge = _make_critique_judge(ok1=False, issue1="fast doubt",
                                            ok2=True, issue2="", calls=_calls)
    ok, issue = f.critique("700は正しいか?", "700")
    check("critique: think=False NGでthink=True再検算が発火する", _calls == [False, True])
    check("critique: think=True再検算がokなら(True,'')で採用", ok is True and issue == "")

    # (2b) think=False NG → think=True も NG。最終issueはthink=True側(issue3)がそのまま
    #      使われ、think=False側のissue1は使われない(権威は再検算側であることの確認)。
    _calls = []
    f._critic_judge = _make_critique_judge(ok1=False, issue1="fast doubt",
                                            ok2=False, issue2="authoritative issue",
                                            calls=_calls)
    ok, issue = f.critique("q", "a")
    check("critique: think=True再検算もNGならreject(False)", ok is False)
    check("critique: 最終issueはthink=True側(権威)がそのまま使われる",
          issue == "authoritative issue")

    # (3) think=True 側が _critic_judge の __ERROR__ センチネル形(ok=False, "critic call
    #     failed: ...")を返す場合。critic 呼び出し自体の失敗であって「回答に問題なし」
    #     ではないため、critique() はこれを ok=True に握りつぶさず reject のまま返す必要がある。
    _calls = []
    f._critic_judge = _make_critique_judge(
        ok1=False, issue1="fast doubt",
        ok2=False, issue2="critic call failed: __ERROR__: simulated transport failure",
        calls=_calls)
    ok, issue = f.critique("q", "a")
    check("critique: think=True側がcritic呼び出し失敗センチネルでもok=Trueに握りつぶさない",
          ok is False)
finally:
    f._critic_judge = _orig_critic_judge

_orig_proposers = f.PROPOSERS
_orig_second_opinion_model = f.SECOND_OPINION_MODEL
_orig_disabled_flag = f._SECOND_OPINION_DISABLED
_orig_ask = f.ask
try:
    # SECOND_OPINION_MODEL を PROPOSERS に含めて「有効」経路を通す。
    f.SECOND_OPINION_MODEL = "phi4-mini"
    f.PROPOSERS = ["phi4-mini", "qwen3:4b"]
    f._SECOND_OPINION_DISABLED = False

    f.ask = lambda *a, **k: "__ERROR__: simulated transport failure"
    ok, issue = f.second_opinion("q", "a")
    check("so: __ERROR__センチネルはok=False", ok is False and bool(issue))

    f.ask = lambda *a, **k: ""
    ok, issue = f.second_opinion("q", "a")
    check("so: 空文字は既定どおりok=True", ok is True)

    f.ask = lambda *a, **k: "Looks fine to me, no issues here."
    ok, issue = f.second_opinion("q", "a")
    check("so: 非JSONの地の文も既定どおりok=True", ok is True)

    # 無効化パス(SECOND_OPINION_MODEL not in PROPOSERS)は ask を一切呼ばずに (True, "") を返す。
    _calls = []
    f.PROPOSERS = ["qwen3:4b"]  # phi4-mini を除外
    f.ask = lambda *a, **k: _calls.append(1) or "__ERROR__: should not be reached"
    ok, issue = f.second_opinion("q", "a")
    check("so: 無効化パスはaskを呼ばずok=True", ok is True and issue == "" and not _calls)
finally:
    f.PROPOSERS = _orig_proposers
    f.SECOND_OPINION_MODEL = _orig_second_opinion_model
    f._SECOND_OPINION_DISABLED = _orig_disabled_flag
    f.ask = _orig_ask

# ---------- verify_single: think=True 最終審判の __ERROR__ は MoA へエスカレーション (2026-07-22) ----------
# verify_single は高速チェックのどちらかが疑義を出したときだけ think=True 再検算を
# 最終審判にする。その think=True 呼び出し自体が __ERROR__ で失敗した場合、_critic_judge の
# 修正により ok=False になるはずで、verify_single はそれを受けて False を返し MoA パネルへの
# 格上げを引き起こす（黙って True を返して壊れた回答を採用しない）ことを確認する。
_orig_proposers = f.PROPOSERS
_orig_second_opinion_model = f.SECOND_OPINION_MODEL
_orig_disabled_flag = f._SECOND_OPINION_DISABLED
_orig_ask = f.ask
try:
    # second_opinion を無効化パスに固定し、think=True 最終審判の挙動だけを見る。
    f.PROPOSERS = ["qwen3:4b"]
    f.SECOND_OPINION_MODEL = "phi4-mini"
    f._SECOND_OPINION_DISABLED = False

    def _fake_ask_escalate(model, messages, temperature, think=None, fmt=None,
                            label=None, num_predict=None, num_ctx=None):
        if think:
            return "__ERROR__: simulated transport failure"
        return json.dumps({"ok": False, "issue": "fast check flagged"})
    f.ask = _fake_ask_escalate
    ok, issue = f.verify_single("2+2?", "4")
    check("verify_single: think=True最終審判の__ERROR__はMoAへ格上げ(ok=False)",
          ok is False and bool(issue))

    def _fake_ask_control(model, messages, temperature, think=None, fmt=None,
                           label=None, num_predict=None, num_ctx=None):
        if think:
            return json.dumps({"ok": True, "issue": ""})
        return json.dumps({"ok": False, "issue": "fast check flagged"})
    f.ask = _fake_ask_control
    ok, issue = f.verify_single("2+2?", "4")
    check("verify_single: think=True最終審判が正常なJSONなら採用(control)",
          ok is True and issue == "")
finally:
    f.PROPOSERS = _orig_proposers
    f.SECOND_OPINION_MODEL = _orig_second_opinion_model
    f._SECOND_OPINION_DISABLED = _orig_disabled_flag
    f.ask = _orig_ask

# ---------- verify_single: second_opinion 有効時（自己評価バイアス対策 seam, 2026-07-23）----------
# 上の 2 ブロック（second_opinion のバイアス対策 / verify_single の think=True 最終審判）は
# どちらも PROPOSERS=["qwen3:4b"] のみで SECOND_OPINION_MODEL(既定 gpt-oss:20b、ここでは
# phi4-mini) を含めていない。そのため second_opinion() は毎回 (True, "") を即返して
# _SECOND_OPINION_DISABLED=True に自己ラッチするだけで、ask を一度も呼ばない
# 「独立チェック無効」分岐しか通っていなかった。verify_single の目玉である
# 「2系統独立チェック（qwen3 think=False の自己批評 と 別モデル phi4-mini の独立チェック）を
# 両方揃えたときの高速一致即採用」と「片方だけが疑義を出したら think=True 最終審判に
# 格上げする自己評価バイアス対策そのもの」は一度も検証されていなかった。
# ここでは PROPOSERS に SECOND_OPINION_MODEL を含めて「有効」経路を通し、label
# ("critic"/"critic2") と think の組み合わせで 3 種類の呼び出し
# （fast self-critic=critic/think=False, 独立second opinion=critic2, think=True最終審判=
# critic/think=True）を区別してモックする。fast self-critic と second_opinion はどちらも
# think を渡さない（falsy）ため、think だけでは区別できず label 分岐が必須になる点に注意。
_orig_proposers = f.PROPOSERS
_orig_second_opinion_model = f.SECOND_OPINION_MODEL
_orig_disabled_flag = f._SECOND_OPINION_DISABLED
_orig_ask = f.ask


def _make_vs_ask(calls, ok1, issue1, ok2, issue2, ok3=False, issue3=""):
    """label と think(bool化) で 3種の呼び出しを切り分けて記録するモック。
    calls には (label, think) の実呼び出しログを積む（呼ばれた事実自体をテストで検証する）。"""
    def _ask(model, messages, temperature, think=None, fmt=None, label=None,
             num_predict=None, num_ctx=None):
        calls.append((label, bool(think)))
        if label == "critic" and think:
            return json.dumps({"ok": ok3, "issue": issue3})
        if label == "critic2":
            return json.dumps({"ok": ok2, "issue": issue2})
        return json.dumps({"ok": ok1, "issue": issue1})
    return _ask


try:
    f.PROPOSERS = ["phi4-mini", "qwen3:4b"]  # SECOND_OPINION_MODEL を含める→有効経路
    f.SECOND_OPINION_MODEL = "phi4-mini"

    # (A) ok1=True かつ ok2=True: 高速2系統が一致→即採用。think=True 最終審判は
    #     呼ばれないはず（高速パスの短絡が壊れていないことをロックする）。
    f._SECOND_OPINION_DISABLED = False
    _calls = []
    f.ask = _make_vs_ask(_calls, ok1=True, issue1="", ok2=True, issue2="")
    ok, issue = f.verify_single("2+2?", "4")
    check("vs(有効) A: ok1=True/ok2=True→即採用", ok is True and issue == "")
    check("vs(有効) A: second_opinionが実際に呼ばれた(critic2が無効化せず発火)",
          any(lbl == "critic2" for lbl, _th in _calls))
    check("vs(有効) A: 高速一致時はthink=True最終審判を呼ばない(短絡ロック)",
          not any(lbl == "critic" and th for lbl, th in _calls))

    # (B) ok1=True(fast self-critic は合格) だが ok2=False(独立モデルが疑義)。
    #     自己評価バイアス対策の本丸：片方だけの疑義でも think=True 最終審判へ格上げされる。
    #   (B1) 最終審判が ok→採用(True)。
    f._SECOND_OPINION_DISABLED = False
    _calls = []
    f.ask = _make_vs_ask(_calls, ok1=True, issue1="", ok2=False, issue2="second opinion doubt",
                          ok3=True, issue3="")
    ok, issue = f.verify_single("2+2?", "4")
    check("vs(有効) B1: ok1=True/ok2=False→think=True最終審判が呼ばれる",
          any(lbl == "critic" and th for lbl, th in _calls))
    check("vs(有効) B1: 最終審判okなら採用(True)", ok is True and issue == "")

    #   (B2) 最終審判もNG→不採用(False)。理由は最終審判(issue3)が権威として使われる。
    f._SECOND_OPINION_DISABLED = False
    _calls = []
    f.ask = _make_vs_ask(_calls, ok1=True, issue1="", ok2=False, issue2="second opinion doubt",
                          ok3=False, issue3="final check: real issue found")
    ok, issue = f.verify_single("2+2?", "4")
    check("vs(有効) B2: 最終審判もNGなら不採用(False)でissue非空",
          ok is False and bool(issue))
    check("vs(有効) B2: 不採用理由はthink=True最終審判が権威(issue3を採用)",
          issue == "final check: real issue found")

    # (C) ok1=False(fast self-critic が疑義) だが ok2=True(独立モデルは合格)。
    #     こちらの片方だけの疑義でも think=True 最終審判へ格上げされることを確認する。
    #     最終審判がNGでissue3が空文字のときは doubt(=issue1) にフォールバックする。
    f._SECOND_OPINION_DISABLED = False
    _calls = []
    f.ask = _make_vs_ask(_calls, ok1=False, issue1="fast critic doubt", ok2=True, issue2="",
                          ok3=False, issue3="")
    ok, issue = f.verify_single("2+2?", "4")
    check("vs(有効) C: ok1=False/ok2=True→think=True最終審判が呼ばれる(doubtはissue1側)",
          any(lbl == "critic" and th for lbl, th in _calls))
    check("vs(有効) C: 最終審判issue3が空ならdoubt(issue1)にフォールバック",
          ok is False and issue == "fast critic doubt")
finally:
    f.PROPOSERS = _orig_proposers
    f.SECOND_OPINION_MODEL = _orig_second_opinion_model
    f._SECOND_OPINION_DISABLED = _orig_disabled_flag
    f.ask = _orig_ask

# ---------- bench_queue: 異常終了コード分類（gotcha 8 再発防止, 2026-07-21）----------
# job 4 (math500/sc+pot) が rc=1073807364 で落ちた際、旧実装は成功/失敗/クラッシュを
# 区別せず、以降のジョブが rc=3221226091 で連鎖即死してもキューは気づかず
# 「正常終了」の top-level status を書いていた。classify_exit_code はその判定を
# 担う純粋関数（I/O無し・import 時に副作用も無い）で、ここではモデル呼び出しも
# subprocess も一切使わずオフラインで検証する。
check("bq: rc=0 は ok", bq.classify_exit_code(0) == "ok")
check("bq: 通常の Python 失敗(rc=1)は error(crashではない)", bq.classify_exit_code(1) == "error")
check("bq: 通常の Python 失敗(rc=2)は error", bq.classify_exit_code(2) == "error")
check("bq: 実際に遭遇した job4 の異常終了コードは crash",
      bq.classify_exit_code(1073807364) == "crash")
check("bq: 連鎖即死した後続ジョブの異常終了コードも crash",
      bq.classify_exit_code(3221226091) == "crash")
check("bq: 負のシグナル終了コードも crash", bq.classify_exit_code(-1) == "crash")
check("bq: 閾値未満の巨大値は crash 扱いしない(境界)", bq.classify_exit_code(0x40000000 - 1) == "error")
check("bq: 閾値ちょうどは crash(境界)", bq.classify_exit_code(0x40000000) == "crash")

# main() のループを直接叩かず、同じ分類ロジックを模擬ジョブ列に適用して、
# crash が発生したジョブの status が 'ok' でなく、かつ全体の総合結果
# (status['ok'] 相当) が「クリーンな成功」を主張しないことを確認する。
_sim_rcs = [0, 0, 1073807364, 3221226091]
_sim_jobs = [{"n": i + 1, "rc": rc, "status": bq.classify_exit_code(rc)}
             for i, rc in enumerate(_sim_rcs)]
check("bq: crash ジョブの status は 'ok' ではない", _sim_jobs[2]["status"] != "ok")
_sim_crashed = [j for j in _sim_jobs if j["status"] == "crash"]
_sim_failed = [j for j in _sim_jobs if j["status"] == "error"]
_sim_overall_ok = not _sim_crashed and not _sim_failed
check("bq: gotcha8同型シナリオでは総合結果が失敗を示す(クリーン成功を主張しない)",
      _sim_overall_ok is False and len(_sim_crashed) >= 1)

# bench_queue の import 自体が副作用を持たないこと（main() は __main__ ガード下）の
# 簡易確認: モジュールに main はあるが、import 時点でジョブが実行されていないこと
# （_sim_jobs はテスト側のダミーであり QUEUE の長さと無関係）。
check("bq: import時点でQUEUEは未実行のジョブ一覧のまま(副作用なし)",
      hasattr(bq, "QUEUE") and hasattr(bq, "main") and hasattr(bq, "classify_exit_code"))


# ---------- ask(): think-strip リトライがループ末尾で握り潰されない (gotcha 3 関連バグ修正) ----------
# 旧実装は「thinking非対応」400 を検知したら payload から think を pop して
# for attempt in (1, 2) の次のイテレーションに continue するだけだった。
# attempt=1 が一過性の500（ロード直後によくある）で、attempt=2 で初めて
# 「thinking非対応」400 が出ると、continue しても for ループは尽きており、
# 組み直したリクエストは一度も送信されずに __ERROR__: think_stripped_retry が
# そのまま最終戻り値になっていた（SC投票/提案が黙って1票失われる = 精度低下）。
# ここでは urllib.request.urlopen と time.sleep のみをモックし、実際の
# Ollama/ネットワーク呼び出しは一切発生させずに検証する。


class _FakeHTTPResponse:
    """urllib.request.urlopen が返す `with ... as r:` 用の最小モック。"""

    def __init__(self, body_bytes):
        self._body = body_bytes

    def __enter__(self):
        return self

    def __exit__(self, *exc_info):
        return False

    def read(self):
        return self._body


def _http_error(code, body_text):
    return urllib.error.HTTPError(
        f"{f.OLLAMA_URL}/api/chat", code, "mock error", {},
        io.BytesIO(body_text.encode("utf-8")),
    )


def _make_fake_urlopen(steps, calls_log):
    """steps: [("error", code, body_text), ...] または [("ok", content_text), ...] のリスト。
    定義された手数を超えて呼ばれたら AssertionError にする(無限ループを検知するため)。
    呼び出しごとに送信 payload(dict) を calls_log に積む。"""
    state = {"i": 0}

    def _fake_urlopen(req, timeout=None):
        i = state["i"]
        state["i"] += 1
        if i >= len(steps):
            raise AssertionError(f"unexpected extra urlopen call #{i + 1} (bounded-loop violation)")
        calls_log.append(json.loads(req.data.decode("utf-8")))
        step = steps[i]
        if step[0] == "error":
            raise _http_error(step[1], step[2])
        return _FakeHTTPResponse(json.dumps({"message": {"content": step[1]}}).encode("utf-8"))

    return _fake_urlopen


_orig_urlopen = urllib.request.urlopen
_orig_sleep = f.time.sleep

# --- シナリオ1: attempt1=一過性500, attempt2=thinking非対応400 → 組み直しリクエストが
#     独自に送信され、最終応答が失われない(バグ再現シナリオそのもの) ---
_calls1 = []
try:
    f.time.sleep = lambda s: None  # 一過性リトライの sleep(2) を待たない
    urllib.request.urlopen = _make_fake_urlopen(
        [("error", 500, "internal error, model loading"),
         ("error", 400, "this model does not support thinking"),
         ("ok", "the real final answer")],
        _calls1,
    )
    _r1 = f.ask("m1", [{"role": "user", "content": "hi"}], 0.7, think=True)
    check("ask: 500→thinking400→success で最終応答が失われない",
          _r1 == "the real final answer")
    check("ask: __ERROR__を返さない(think_stripped_retryで確定終了しない)",
          not str(_r1).startswith("__ERROR__"))
    check("ask: 成功リクエストのpayloadにthinkキーが残っていない",
          "think" not in _calls1[-1])
finally:
    urllib.request.urlopen = _orig_urlopen
    f.time.sleep = _orig_sleep

# --- シナリオ2: attempt1で即thinking非対応400(既存の主要ケース)→引き続き成功する回帰確認 ---
_calls2 = []
try:
    f.time.sleep = lambda s: None
    urllib.request.urlopen = _make_fake_urlopen(
        [("error", 400, "this model does not support thinking"),
         ("ok", "stripped retry answer")],
        _calls2,
    )
    _r2 = f.ask("m1", [{"role": "user", "content": "hi"}], 0.7, think=True)
    check("ask: 初回でthinking400→即座のstrip再送で成功(既存ケースの回帰なし)",
          _r2 == "stripped retry answer")
    check("ask: シナリオ2でも成功payloadにthinkキーが残っていない",
          "think" not in _calls2[-1])
finally:
    urllib.request.urlopen = _orig_urlopen
    f.time.sleep = _orig_sleep

# --- シナリオ3: 毎回thinking非対応400 → think はpop済みなので分岐は高々1回のみ発火し、
#     有限回のurlopen呼び出しで(無限ループせず)__ERROR__を返す ---
_calls3 = []
try:
    f.time.sleep = lambda s: None
    urllib.request.urlopen = _make_fake_urlopen(
        [("error", 400, "this model does not support thinking"),
         ("error", 400, "this model does not support thinking")],
        _calls3,
    )
    _r3 = f.ask("m1", [{"role": "user", "content": "hi"}], 0.7, think=True)
    check("ask: 毎回thinking400でも有限回(<=2回)のurlopen呼び出しで打ち切り(無限ループなし)",
          len(_calls3) <= 2)
    check("ask: 毎回thinking400なら最終的に__ERROR__を返す", str(_r3).startswith("__ERROR__"))
finally:
    urllib.request.urlopen = _orig_urlopen
    f.time.sleep = _orig_sleep

# --- シナリオ4: 通常の一過性失敗(thinking非対応ではない500が2連続)は従来通り
#     ちょうど2回試行してsleep(2)を1回だけ挟み__ERROR__を返す(一過性リトライ予算は不変) ---
_calls4 = []
try:
    f.time.sleep = lambda s: None
    urllib.request.urlopen = _make_fake_urlopen(
        [("error", 500, "internal error"),
         ("error", 500, "internal error")],
        _calls4,
    )
    _r4 = f.ask("m1", [{"role": "user", "content": "hi"}], 0.7, think=True)
    check("ask: 通常の一過性失敗(500,500)はちょうど2回試行して__ERROR__",
          len(_calls4) == 2 and str(_r4).startswith("__ERROR__"))
finally:
    urllib.request.urlopen = _orig_urlopen
    f.time.sleep = _orig_sleep

# --- gotcha #1 / #2 回帰: /api/chat 固定 & options.num_ctx 常時pin ---
# 既存の _make_fake_urlopen は payload(dict) だけを calls_log に積む設計で、既存の
# _calls1..4 のインデックス方法(payload dict として直接参照)を変えると回帰するため
# ここでは触らない。URL も検証したいこのセクション専用に別のフェイク urlopen を用意する
# (calls_log の各要素は {"full_url":.., "payload":..} の dict)。


def _make_fake_urlopen_url(steps, calls_log):
    """_make_fake_urlopen と同じ挙動だが、送信 payload に加えて req.full_url も記録する。"""
    state = {"i": 0}

    def _fake_urlopen(req, timeout=None):
        i = state["i"]
        state["i"] += 1
        if i >= len(steps):
            raise AssertionError(f"unexpected extra urlopen call #{i + 1} (bounded-loop violation)")
        calls_log.append({
            "full_url": req.full_url,
            "payload": json.loads(req.data.decode("utf-8")),
        })
        step = steps[i]
        if step[0] == "error":
            raise _http_error(step[1], step[2])
        return _FakeHTTPResponse(json.dumps({"message": {"content": step[1]}}).encode("utf-8"))

    return _fake_urlopen


# シナリオ5: 通常成功呼び出し(think/num_predict/fmt すべて未指定・未知モデル) →
# native /api/chat を叩き(/v1 は使わない)、options.num_ctx が既定値で必ず pin される。
_calls5 = []
try:
    f.time.sleep = lambda s: None
    urllib.request.urlopen = _make_fake_urlopen_url(
        [("ok", "plain answer")],
        _calls5,
    )
    _r5 = f.ask("m-unknown-nonthinking", [{"role": "user", "content": "hi"}], 0.7)
    _url5 = _calls5[-1]["full_url"]
    _opts5 = _calls5[-1]["payload"].get("options", {})
    check("ask: 通常呼び出しは /api/chat を叩く(gotcha#1)", _url5.endswith("/api/chat"))
    check("ask: 通常呼び出しで /v1 エンドポイントは使わない(gotcha#1)", "/v1" not in _url5)
    check("ask: think/num_predict/fmt が全てNoneでもoptions.num_ctxは省略されない(gotcha#2)",
          "num_ctx" in _opts5 and _opts5["num_ctx"])
    check("ask: 未知モデルはMODEL_NUM_CTXが既定値になる",
          _opts5["num_ctx"] == f.MODEL_NUM_CTX)
finally:
    urllib.request.urlopen = _orig_urlopen
    f.time.sleep = _orig_sleep

# シナリオ6: MODEL_CONFIG に登録された思考モデル(gpt-oss:20b)は num_ctx=16384 が
# model_cfg 由来で pin される(8192 のままでは思考が truncate される既知不具合の回帰防止)。
_calls6 = []
try:
    f.time.sleep = lambda s: None
    urllib.request.urlopen = _make_fake_urlopen_url(
        [("ok", "thinking model answer")],
        _calls6,
    )
    _r6 = f.ask("gpt-oss:20b", [{"role": "user", "content": "hi"}], 0.7)
    _opts6 = _calls6[-1]["payload"].get("options", {})
    _expected_ctx6 = f.model_cfg("gpt-oss:20b", "num_ctx", f.MODEL_NUM_CTX)
    check("ask: gpt-oss:20b(思考モデル)はMODEL_CONFIG由来のnum_ctxになる",
          _opts6["num_ctx"] == _expected_ctx6 == 16384)
finally:
    urllib.request.urlopen = _orig_urlopen
    f.time.sleep = _orig_sleep

# シナリオ7: 明示的な num_ctx=... 引数はモデル既定値より優先される。
_calls7 = []
try:
    f.time.sleep = lambda s: None
    urllib.request.urlopen = _make_fake_urlopen_url(
        [("ok", "explicit ctx answer")],
        _calls7,
    )
    _r7 = f.ask("gpt-oss:20b", [{"role": "user", "content": "hi"}], 0.7, num_ctx=12345)
    _opts7 = _calls7[-1]["payload"].get("options", {})
    check("ask: 明示的なnum_ctx引数はモデル既定値より優先される",
          _opts7["num_ctx"] == 12345)
finally:
    urllib.request.urlopen = _orig_urlopen
    f.time.sleep = _orig_sleep

# シナリオ8(重要): think-strip再送パス(500→thinking非対応400→success)でも、
# 最終的に再送されるリクエストが num_ctx pin と /api/chat エンドポイントの両方を
# 維持していること(L1054のpayload再構築で options.num_ctx や URL が失われていないか)。
_calls8 = []
try:
    f.time.sleep = lambda s: None
    urllib.request.urlopen = _make_fake_urlopen_url(
        [("error", 500, "internal error, model loading"),
         ("error", 400, "this model does not support thinking"),
         ("ok", "final answer after think strip")],
        _calls8,
    )
    _r8 = f.ask("gpt-oss:20b", [{"role": "user", "content": "hi"}], 0.7, think=True)
    _expected_ctx8 = f.model_cfg("gpt-oss:20b", "num_ctx", f.MODEL_NUM_CTX)
    _final_call8 = _calls8[-1]
    _final_opts8 = _final_call8["payload"].get("options", {})
    check("ask: think-strip再送でも最終応答が失われない",
          _r8 == "final answer after think strip")
    check("ask: think-strip再送の最終リクエストもoptions.num_ctxを維持する(gotcha#2)",
          "num_ctx" in _final_opts8 and _final_opts8["num_ctx"] == _expected_ctx8)
    check("ask: think-strip再送の最終リクエストも/api/chatを維持する(gotcha#1)",
          _final_call8["full_url"].endswith("/api/chat") and "/v1" not in _final_call8["full_url"])
    check("ask: think-strip再送の最終リクエストはthinkキーが除去されている",
          "think" not in _final_call8["payload"])
finally:
    urllib.request.urlopen = _orig_urlopen
    f.time.sleep = _orig_sleep

# ---------- fugu_answer: SC結果のユーザー提示 ----------
# fugu_answer() の自己一貫性投票(SC)結果 → ユーザー提示の接合点（行 2602-2615 付近）の回帰テスト。
# solve_verifiable の戻り値と、本文から実際に抽出した答え(extract_final_answer/answers_equivalent
# は本物をそのまま使う)がずれた場合にのみ「(自己一貫性投票による最終解答: X)」を付記し、
# 一致すれば本文をそのまま返す。None ならMoA(合議)へフォールバックする、という分岐を検証する。
# validate_plan() 済みの明示プランを渡して conduct() を経由させない。

_orig_sc_enabled = f.SC_ENABLED
_orig_solve_verifiable = f.solve_verifiable
_orig_get_proposals = f.get_proposals
_orig_aggregate = f.aggregate


def _validated_plan(task_type, mode="moa", rounds=1):
    # default_plan()/validate_plan() が生成する形と同じキー構成の明示プラン。
    return {
        "mode": mode,
        "task_type": task_type,
        "selected_proposers": ["m1", "m2", "m3"],
        "rounds": rounds,
        "use_image_generation": False,
        "image_only": False,
        "make_pptx": False,
        "search_required": False,
        "reason": "test",
        "_fallback": False,
    }


def _make_moa_forbidden(touched_list):
    """SCが成功した経路ではget_proposals/aggregateへ絶対に到達してはならないことを
    検証するための番人。呼ばれたら記録した上で必ず例外を送出する。"""
    def _get_proposals_forbidden(*a, **kw):
        touched_list.append(True)
        raise AssertionError("SC成功時はMoA(get_proposals)へ到達してはならない")

    def _aggregate_forbidden(*a, **kw):
        touched_list.append(True)
        raise AssertionError("SC成功時はMoA(aggregate)へ到達してはならない")

    return _get_proposals_forbidden, _aggregate_forbidden


# --- Case A: 本文の結論(裁定等で差し替わった\boxed)が投票結果と食い違う → 明示注記が付く ---
_moa_touched_a = []
_get_proposals_never_a, _aggregate_never_a = _make_moa_forbidden(_moa_touched_a)

try:
    f.SC_ENABLED = True
    f.solve_verifiable = lambda question, task_type, history=None: {
        "answer": "5",
        "text": "途中式…裁定により \\boxed{7} に差し替え。",
        "votes": {"7": 2, "5": 1},
        "n_samples": 3,
    }
    f.get_proposals = _get_proposals_never_a
    f.aggregate = _aggregate_never_a
    with contextlib.redirect_stdout(io.StringIO()):
        _ans_a = f.fugu_answer("2+3は?", plan=_validated_plan("math"))
finally:
    f.SC_ENABLED = _orig_sc_enabled
    f.solve_verifiable = _orig_solve_verifiable
    f.get_proposals = _orig_get_proposals
    f.aggregate = _orig_aggregate

check("fugu_answer: 本文の結論と投票結果が食い違う場合は明示注記を付す",
      "(自己一貫性投票による最終解答: 5)" in _ans_a)
check("fugu_answer: 食い違いケースでも元の本文はそのまま含まれる",
      "裁定により \\boxed{7} に差し替え。" in _ans_a)
check("fugu_answer: 食い違いケースはSC経路で返りMoAへ到達しない", not _moa_touched_a)

# --- Case B: 本文の結論(\boxedの答え)が投票結果と一致 → 注記なしでそのまま返す（math） ---
_moa_touched_b1 = []
_get_proposals_never_b1, _aggregate_never_b1 = _make_moa_forbidden(_moa_touched_b1)
try:
    f.SC_ENABLED = True
    f.solve_verifiable = lambda question, task_type, history=None: {
        "answer": "42",
        "text": "計算の結果、\\boxed{42} である。",
        "votes": {"42": 3},
        "n_samples": 3,
    }
    f.get_proposals = _get_proposals_never_b1
    f.aggregate = _aggregate_never_b1
    with contextlib.redirect_stdout(io.StringIO()):
        _ans_b1 = f.fugu_answer("6*7は?", plan=_validated_plan("math"))
finally:
    f.SC_ENABLED = _orig_sc_enabled
    f.solve_verifiable = _orig_solve_verifiable
    f.get_proposals = _orig_get_proposals
    f.aggregate = _orig_aggregate

check("fugu_answer: 本文とSC結果(math)が一致すれば注記なし",
      "自己一貫性投票による最終解答" not in _ans_b1)
check("fugu_answer: 一致ケース(math)は本文をそのまま返す",
      _ans_b1 == "計算の結果、\\boxed{42} である。")
check("fugu_answer: 一致ケース(math)もMoAへ到達しない", not _moa_touched_b1)

# --- Case B': mcq版（選択肢文字が一致） ---
_moa_touched_b2 = []
_get_proposals_never_b2, _aggregate_never_b2 = _make_moa_forbidden(_moa_touched_b2)
try:
    f.SC_ENABLED = True
    f.solve_verifiable = lambda question, task_type, history=None: {
        "answer": "C",
        "text": "検討の結果、\\boxed{C} が正解。",
        "votes": {"C": 3},
        "n_samples": 3,
    }
    f.get_proposals = _get_proposals_never_b2
    f.aggregate = _aggregate_never_b2
    with contextlib.redirect_stdout(io.StringIO()):
        _ans_b2 = f.fugu_answer("次のうち正しいものは?", plan=_validated_plan("mcq"))
finally:
    f.SC_ENABLED = _orig_sc_enabled
    f.solve_verifiable = _orig_solve_verifiable
    f.get_proposals = _orig_get_proposals
    f.aggregate = _orig_aggregate

check("fugu_answer: 本文とSC結果(mcq選択肢)が一致すれば注記なし",
      "自己一貫性投票による最終解答" not in _ans_b2)
check("fugu_answer: 一致ケース(mcq)は本文をそのまま返す",
      _ans_b2 == "検討の結果、\\boxed{C} が正解。")
check("fugu_answer: 一致ケース(mcq)もMoAへ到達しない", not _moa_touched_b2)

# --- Case C: solve_verifiable が None(投票不成立) → MoAへフォールスルーする ---
# plan["rounds"]=MAX_ROUNDS にして、r>=limit のブレークが「計画分残っているか」判定より
# 先に効くようにし、critique()（本物のask呼び出しを要する）へ到達せずに済ませる
# （このテストが検証したいのはSC→MoAへの委譲そのものであり、MoAの反復打ち切りロジック自体は
# 既存の他テストが担保している）。
_MOA_SENTINEL = "MOA_FALLBACK_SENTINEL: 単体/合議側で生成された最終回答"
_get_proposals_calls_c = []
_aggregate_calls_c = []


def _fake_get_proposals_c(models, question, reference=None, issue=None, history=None):
    _get_proposals_calls_c.append((tuple(models), reference, issue))
    return [(m, "dummy proposal (SCフォールバック検証用ダミー)") for m in models]


def _fake_aggregate_c(question, proposals):
    _aggregate_calls_c.append(len(proposals))
    return _MOA_SENTINEL


try:
    f.SC_ENABLED = True
    f.solve_verifiable = lambda question, task_type, history=None: None
    f.get_proposals = _fake_get_proposals_c
    f.aggregate = _fake_aggregate_c
    with contextlib.redirect_stdout(io.StringIO()):
        _ans_c = f.fugu_answer(
            "解けない問題?", plan=_validated_plan("math", mode="moa", rounds=f.MAX_ROUNDS))
finally:
    f.SC_ENABLED = _orig_sc_enabled
    f.solve_verifiable = _orig_solve_verifiable
    f.get_proposals = _orig_get_proposals
    f.aggregate = _orig_aggregate

check("fugu_answer: SCがNoneならMoA(get_proposals/aggregate)へフォールスルーする",
      _ans_c == _MOA_SENTINEL)
check("fugu_answer: フォールバック時は実際にget_proposals/aggregateが呼ばれる",
      len(_get_proposals_calls_c) >= 1 and len(_aggregate_calls_c) >= 1)

# ---------- fugu_answer MoAループ: 次ラウンドの reference は aggregate 出力の think を
# 持ち越さない (2026-07-24) ----------
# fugu_local.py ~L3052-3062: 従来は `reference = final`（aggregate の生出力）をそのまま
# 次ラウンドの get_proposals へ渡していた。aggregate/ask_fugu の他の受け渡し箇所（aggregate
# L2321、ask_fugu L3196、_sc_sample/_arbitrate）はすべて strip_think 済みの値を使っている
# のに、この1箇所だけ think を残したまま次ラウンドの get_single_proposal に
# 'A draft answer from the panel:\n{reference}' として渡していた。これは(1) 内部思考を
# 「ドラフト回答」としてプロポーザーに誤呈示し改善を誤誘導し、(2) 8192/16384に固定された
# num_ctx（gotcha #2）を think ブロックが圧迫し本来のドラフト/質問が切り詰められかねない
# （精度優先=gotcha #7）という2つの問題を持つ。修正: aggregate 直後に1回だけ
# fin=strip_think(final) を計算し、reference にはその fin を使う。戻り値の final 自体は
# 変更しない（ask_fugu 側で最終的に strip_think される）。
_orig_code_check_ref = f.code_check
_orig_critique_ref = f.critique
_orig_ask_ref = f.ask
_orig_verify_single_ref = f.verify_single
_orig_allow_recursion_ref = f.ALLOW_RECURSION
_orig_adaptive_escalation_ref = f.ADAPTIVE_ESCALATION

_ROUND1_THINK_REF = "<think>これは内部思考であり最終回答ではない、長々とした考察が続く</think>"
_ROUND1_BODY_REF = "ラウンド1の統合結果の本文です。"
_ROUND1_FINAL_REF = _ROUND1_THINK_REF + "\n" + _ROUND1_BODY_REF
_ROUND2_THINK_REF = "<think>ラウンド2の内部思考</think>"
_ROUND2_BODY_REF = "ラウンド2の最終統合結果です。"
_ROUND2_FINAL_REF = _ROUND2_THINK_REF + "\n" + _ROUND2_BODY_REF

_get_proposals_calls_ref = []
_aggregate_call_n_ref = [0]


def _fake_get_proposals_ref(models, question, reference=None, issue=None, history=None):
    _get_proposals_calls_ref.append({"reference": reference, "issue": issue})
    return [(m, "dummy proposal") for m in models]


def _fake_aggregate_ref(question, proposals):
    _aggregate_call_n_ref[0] += 1
    return _ROUND1_FINAL_REF if _aggregate_call_n_ref[0] == 1 else _ROUND2_FINAL_REF


try:
    f.get_proposals = _fake_get_proposals_ref
    f.aggregate = _fake_aggregate_ref
    f.code_check = lambda answer: None  # コード実行検証は本テストの対象外
    # r>=planned後にcritique(本物のask呼び出し)へ到達させないための最短経路
    f.ALLOW_RECURSION = False
    with contextlib.redirect_stdout(io.StringIO()):
        _ans_ref = f.fugu_answer(
            "テスト質問(reference持ち越し検証用)?",
            plan=_validated_plan(None, mode="moa", rounds=2))
finally:
    f.get_proposals = _orig_get_proposals
    f.aggregate = _orig_aggregate
    f.code_check = _orig_code_check_ref
    f.ALLOW_RECURSION = _orig_allow_recursion_ref

check("fugu_answer MoA reference: 2ラウンド強制でget_proposalsが2回呼ばれる",
      len(_get_proposals_calls_ref) == 2)
check("fugu_answer MoA reference: round2のreferenceは<think>タグを含まない",
      "<think>" not in (_get_proposals_calls_ref[1]["reference"] or ""))
check("fugu_answer MoA reference: round2のreferenceは内部思考の本文を含まない",
      "内部思考であり最終回答ではない" not in (_get_proposals_calls_ref[1]["reference"] or ""))
check("fugu_answer MoA reference: round2のreferenceはround1の本文を含む",
      _ROUND1_BODY_REF in (_get_proposals_calls_ref[1]["reference"] or ""))
check("fugu_answer MoA reference: round2のreferenceはstrip_think(round1 final)と厳密一致",
      _get_proposals_calls_ref[1]["reference"] == f.strip_think(_ROUND1_FINAL_REF))
check("fugu_answer MoA reference(回帰): 戻り値は最終ラウンドの生出力のまま"
      "(returnはask_fugu側でstrip_thinkされる前提で変更していない)",
      _ans_ref == _ROUND2_FINAL_REF)
check("fugu_answer MoA reference(回帰): 戻り値には<think>タグが残っている",
      "<think>" in _ans_ref)
check("fugu_answer MoA reference: テスト後にf.get_proposalsが元に復元されている",
      f.get_proposals == _orig_get_proposals)
check("fugu_answer MoA reference: テスト後にf.aggregateが元に復元されている",
      f.aggregate == _orig_aggregate)

# ---------- fugu_answer MoAループ: コード修復ラウンドでも reference はコードフェンスを
# 保持したまま think だけ除去される (2026-07-24 / 回帰) ----------
# strip_think は <think>/<thinking> のみを除去し、```python コードフェンスには触れない。
# code_check がラウンド1の回答にエラーを見つけて追加ラウンドを要求するケースでも、
# ラウンド2への reference にコードフェンスがそのまま残ることを確認する
# （code-repair ループが失敗コードを見失わないことの確認）。
_ROUND1_THINK_CODE = "<think>コードを検討中の内部思考</think>"
_ROUND1_CODE_FENCE = "```python\nprint(1/0)\n```"
_ROUND1_FINAL_CODE = (_ROUND1_THINK_CODE + "\n"
                      + "ゼロ除算を試すコードです。\n" + _ROUND1_CODE_FENCE)
_ROUND2_FINAL_CODE = "<think>修正を検討</think>\n修正済みの回答（コード除去済み）。"
_CODE_ISSUE_TEXT = "code execution FAILED:\nZeroDivisionError: division by zero"

_get_proposals_calls_code = []
_code_check_calls_code = []


def _fake_get_proposals_code(models, question, reference=None, issue=None, history=None):
    _get_proposals_calls_code.append({"reference": reference, "issue": issue})
    return [(m, "dummy proposal") for m in models]


def _fake_aggregate_code(question, proposals):
    return _ROUND1_FINAL_CODE if len(_get_proposals_calls_code) <= 1 else _ROUND2_FINAL_CODE


def _fake_code_check_code(answer):
    _code_check_calls_code.append(answer)
    if len(_code_check_calls_code) == 1:
        return _CODE_ISSUE_TEXT
    return None


try:
    f.get_proposals = _fake_get_proposals_code
    f.aggregate = _fake_aggregate_code
    f.code_check = _fake_code_check_code
    f.critique = lambda question, answer: (True, None)  # 到達しても即終了、本物のaskは呼ばない
    with contextlib.redirect_stdout(io.StringIO()):
        _ans_code = f.fugu_answer(
            "コードを書いて(修復ラウンド検証用)?",
            plan=_validated_plan(None, mode="moa", rounds=1))
finally:
    f.get_proposals = _orig_get_proposals
    f.aggregate = _orig_aggregate
    f.code_check = _orig_code_check_ref
    f.critique = _orig_critique_ref

check("fugu_answer コード修復(回帰): code_checkの指摘でget_proposalsが2回呼ばれる",
      len(_get_proposals_calls_code) == 2)
check("fugu_answer コード修復(回帰): round2のissueはcode_checkの指摘そのもの",
      _get_proposals_calls_code[1]["issue"] == _CODE_ISSUE_TEXT)
check("fugu_answer コード修復(回帰): round2のreferenceはコードフェンスを保持している",
      "```python" in (_get_proposals_calls_code[1]["reference"] or "")
      and "print(1/0)" in (_get_proposals_calls_code[1]["reference"] or ""))
check("fugu_answer コード修復(回帰): round2のreferenceから<think>タグは除去されている",
      "<think>" not in (_get_proposals_calls_code[1]["reference"] or ""))
check("fugu_answer コード修復(回帰): round2のreferenceはstrip_think(round1 final)と厳密一致",
      _get_proposals_calls_code[1]["reference"] == f.strip_think(_ROUND1_FINAL_CODE))

# ---------- fugu_answer 単体→MoA エスカレーション: seed_answer は今回の変更の影響を
# 受けない (2026-07-24 / 回帰) ----------
# seed_answer は単体モードの ask() 直後に既に strip_think 済み(fugu_local.py L3006)。
# 今回の変更は「aggregate() の出力を次ラウンドの reference にする箇所」のみが対象で、
# エスカレーション直後の最初の reference(=seed_answer)には触れていない。
_SEED_ANSWER_ESC = "単体モードの回答本文（think除去済み）"
_get_proposals_calls_esc = []


def _fake_get_proposals_esc(models, question, reference=None, issue=None, history=None):
    _get_proposals_calls_esc.append({"reference": reference, "issue": issue})
    return [(m, "dummy proposal") for m in models]


def _fake_aggregate_esc(question, proposals):
    return "<think>集約時の内部思考</think>\nエスカレーション後の統合結果。"


try:
    f.ADAPTIVE_ESCALATION = True
    f.ALLOW_RECURSION = False  # 1ラウンドで打ち切り、本物のcritique/askへ到達させない
    f.ask = lambda *a, **kw: _SEED_ANSWER_ESC  # think タグ無し = strip_think後も不変
    f.verify_single = lambda question, answer: (False, "単体回答に疑義あり(テスト用)")
    f.get_proposals = _fake_get_proposals_esc
    f.aggregate = _fake_aggregate_esc
    _esc_plan = _validated_plan(None, mode="single", rounds=1)
    with contextlib.redirect_stdout(io.StringIO()):
        f.fugu_answer("テスト質問(エスカレーション用)?", plan=_esc_plan, history=[])
finally:
    f.ask = _orig_ask_ref
    f.verify_single = _orig_verify_single_ref
    f.get_proposals = _orig_get_proposals
    f.aggregate = _orig_aggregate
    f.ADAPTIVE_ESCALATION = _orig_adaptive_escalation_ref
    f.ALLOW_RECURSION = _orig_allow_recursion_ref

check("fugu_answer エスカレーション(回帰): 単体回答失敗でMoAへ切替りget_proposalsが呼ばれる",
      len(_get_proposals_calls_esc) >= 1)
check("fugu_answer エスカレーション(回帰): 最初のroundのreferenceは単体回答そのまま(変更なし)",
      _get_proposals_calls_esc[0]["reference"] == _SEED_ANSWER_ESC)
check("fugu_answer エスカレーション(回帰): 最初のroundのissueはverify_singleの指摘そのまま",
      _get_proposals_calls_esc[0]["issue"] == "単体回答に疑義あり(テスト用)")
check("fugu_answer エスカレーション: テスト後にf.askが元に復元されている", f.ask == _orig_ask_ref)
check("fugu_answer エスカレーション: テスト後にf.verify_singleが元に復元されている",
      f.verify_single == _orig_verify_single_ref)
check("fugu_answer エスカレーション: テスト後にf.ALLOW_RECURSIONが元に復元されている",
      f.ALLOW_RECURSION == _orig_allow_recursion_ref)

# ---------- fugu_answer 単体モード: think 解決に proposer_think_for を使う (2026-07-23) ----------
# 2026-07-23 fix の適用対象その2（get_single_proposal と全く同じ欠落パターン）: 単体モードの
# ask 呼び出しも think=PROPOSER_THINK の生グローバルを直渡ししており、隣の
# num_predict=proposer_predict_for(model) だけがMODEL_CONFIG対応済みという非対称があった。
# f.ask をモックしてthink/num_predict kwargを捕捉し、f.verify_singleをok=Trueに固定して
# 単体モードが合議へフォールバックせず即returnする経路（ADAPTIVE_ESCALATION=True既定）を通す。
# task_type をmath/mcq以外にしてSC経路(solve_verifiable)には触れさせない。
_fa_think_calls = []


def _fake_ask_capture_single(model, messages, temperature, think=None, fmt=None,
                              label=None, num_predict=None):
    _fa_think_calls.append({"model": model, "think": think, "num_predict": num_predict})
    return "single-mode answer"


_orig_ask_fa = f.ask
_orig_verify_single_fa = f.verify_single
_orig_pt_fa = f.PROPOSER_THINK
f.ask = _fake_ask_capture_single
f.verify_single = lambda question, answer: (True, None)
try:
    f.PROPOSER_THINK = None
    _fa_think_calls.clear()
    _fa_plan = _validated_plan(None, mode="single")
    _fa_plan["selected_proposers"] = ["gpt-oss:20b"]
    with contextlib.redirect_stdout(io.StringIO()):
        _fa_out = f.fugu_answer("Q?", plan=_fa_plan, history=[])
    check("fugu_answer単体: think解決 gpt-oss:20bはMODEL_CONFIGのhighを渡す(旧: 生PROPOSER_THINKで欠落)",
          _fa_think_calls and _fa_think_calls[-1]["think"] == "high")
    check("fugu_answer単体: num_predictはproposer_predict_for(model)のまま(回帰・不変)",
          _fa_think_calls[-1]["num_predict"] == f.proposer_predict_for("gpt-oss:20b"))
    check("fugu_answer単体: verify_single ok=Trueで即returnし単体回答をそのまま返す",
          _fa_out == "single-mode answer")

    # PROPOSER_THINK override優先の確認（eval の一括OFF等とbyte一致でなければならない）
    f.PROPOSER_THINK = False
    _fa_think_calls.clear()
    _fa_plan2 = _validated_plan(None, mode="single")
    _fa_plan2["selected_proposers"] = ["gpt-oss:20b"]
    with contextlib.redirect_stdout(io.StringIO()):
        f.fugu_answer("Q?", plan=_fa_plan2, history=[])
    check("fugu_answer単体: PROPOSER_THINK override時はgpt-oss:20bもFalse(eval一括OFFとbyte一致)",
          _fa_think_calls[-1]["think"] is False)
finally:
    f.ask = _orig_ask_fa
    f.verify_single = _orig_verify_single_fa
    f.PROPOSER_THINK = _orig_pt_fa
check("fugu_answer単体: テスト後にf.askが元に復元されている", f.ask == _orig_ask_fa)
check("fugu_answer単体: テスト後にf.verify_singleが元に復元されている",
      f.verify_single == _orig_verify_single_fa)
check("fugu_answer単体: テスト後にPROPOSER_THINKが元に復元されている",
      f.PROPOSER_THINK == _orig_pt_fa)

# ---------- _load_rag_chunks: '[' 始まりの過剰フィルタ回帰防止 (2026-07-22) ----------
# _read_excel/_read_pptx の成功時出力（"[Sheet: ...]" / "[Slide 1]"）が
# text.startswith("[") だけで誤スキップされ、RAGから丸ごと欠落していたバグの回帰テスト。
# ライブラリ未インストール通知（1行・pip install を含む）だけが正しくスキップされることも検証。
# ローカル一時ファイルのみを使用し、Ollama/ネットワーク/bench呼び出しは一切行わない。
import tempfile as _tempfile
import os as _os

check("_is_lib_missing_notice: PDF未インストール通知はTrue",
      f._is_lib_missing_notice(
          "[PDF: foo.pdf — テキスト抽出には pdfplumber or pypdf が必要: pip install pdfplumber]"))
check("_is_lib_missing_notice: DOCX未インストール通知はTrue",
      f._is_lib_missing_notice("[DOCX: foo.docx — python-docx が必要: pip install python-docx]"))
check("_is_lib_missing_notice: Excel未インストール通知はTrue",
      f._is_lib_missing_notice("[Excel: foo.xlsx — openpyxl or pandas が必要: pip install openpyxl]"))
check("_is_lib_missing_notice: PPTX未インストール通知はTrue",
      f._is_lib_missing_notice("[PPTX: foo.pptx — python-pptx が必要: pip install python-pptx]"))
check("_is_lib_missing_notice: Excel成功時の'[Sheet: ...]'はFalse",
      not f._is_lib_missing_notice("[Sheet: Sheet1]\nA\tB\n1\t2"))
check("_is_lib_missing_notice: PPTX成功時の'[Slide 1]'はFalse",
      not f._is_lib_missing_notice("[Slide 1]\nこんにちは"))
check("_is_lib_missing_notice: JSON配列先頭'[1, 2, 3]'はFalse",
      not f._is_lib_missing_notice("[1, 2, 3]"))
check("_is_lib_missing_notice: Markdownリンク'[link](url)'はFalse",
      not f._is_lib_missing_notice("[link](url)\n本文がここに続く"))
check("_is_lib_missing_notice: 空文字はFalse", not f._is_lib_missing_notice(""))

with _tempfile.TemporaryDirectory() as _rag_dir:
    import pathlib as _pathlib
    _rag_root = _pathlib.Path(_rag_dir)

    # 成功したExcel抽出を模したテキストファイル（本物のopenpyxl/pandas呼び出しは不要、
    # read_file_text の出力形状だけを .txt として直接再現して検証する）
    (_rag_root / "sheet_like.txt").write_text(
        "[Sheet: Sheet1]\n" + ("data\t" * 5 + "\n") * 50, encoding="utf-8")
    # 成功したPPTX抽出を模したテキストファイル
    (_rag_root / "slide_like.txt").write_text(
        "[Slide 1]\n" + ("スライド本文の内容です。" * 20 + "\n") * 20, encoding="utf-8")
    # JSON配列（トップレベルが '[' で始まる正当なドキュメント）
    (_rag_root / "array.json").write_text(
        json.dumps(list(range(200))), encoding="utf-8")
    # ブラケットリンクで始まる正当なMarkdown
    (_rag_root / "note.md").write_text(
        "[link](https://example.com)\n" + ("本文テキストです。" * 20 + "\n") * 20,
        encoding="utf-8")
    # 本物のライブラリ未インストール通知そのもの（1行）をそのまま模したファイル
    (_rag_root / "notice_like.txt").write_text(
        "[Excel: dummy.xlsx — openpyxl or pandas が必要: pip install openpyxl]",
        encoding="utf-8")
    # 通常のブラケット無しテキスト（既存挙動のバイト単位不変性チェック用）
    _plain_text = ("普通の本文テキストです。" * 30 + "\n") * 10
    (_rag_root / "plain.txt").write_text(_plain_text, encoding="utf-8")

    _rag_chunks = f._load_rag_chunks([str(_rag_root)])
    _rag_by_file = {}
    for _fp, _chunk in _rag_chunks:
        _rag_by_file.setdefault(_os.path.basename(_fp), []).append(_chunk)

    check("RAG: 成功Excel様('[Sheet: ...]')出力がチャンク化される",
          len(_rag_by_file.get("sheet_like.txt", [])) >= 1)
    check("RAG: 成功PPTX様('[Slide 1]')出力がチャンク化される",
          len(_rag_by_file.get("slide_like.txt", [])) >= 1)
    check("RAG: JSON配列('[1, 2, ...']')がチャンク化される",
          len(_rag_by_file.get("array.json", [])) >= 1)
    check("RAG: ブラケットリンクMarkdownがチャンク化される",
          len(_rag_by_file.get("note.md", [])) >= 1)
    check("RAG: ライブラリ未インストール通知そのものはチャンク化されない",
          "notice_like.txt" not in _rag_by_file)

    # チャンク分割/オーバーラップ計算が非ブラケットテキストで従来通りバイト一致すること
    _expected_plain_chunks = []
    _start = 0
    while _start < len(_plain_text):
        _end = _start + f.RAG_CHUNK_CHARS
        _expected_plain_chunks.append(_plain_text[_start:_end])
        _start += f.RAG_CHUNK_CHARS - f.RAG_CHUNK_OVERLAP
    check("RAG: 通常テキストのチャンク分割はバイト単位で従来通り",
          _rag_by_file.get("plain.txt", []) == _expected_plain_chunks)

# ---------- _load_rag_chunks: 1ファイルの読み込み例外でRAG全体が落ちない (2026-07-22 / iter42) ----------
# read_file_text(fp) を裸で呼んでいたため、破損/未対応ファイル1件がImportError以外の例外を
# 送出すると _load_rag_chunks -> _get_rag_chunks -> rag_search -> build_context まで伝播し、
# 質問のたびにRAGコンテキストが丸ごと失われていた。ここでは1ファイル単位に例外を隔離し、
# 他ファイルは正常にチャンク化されることを検証する。iter41のgraceful-degradation方針を踏襲。

# (1) read_file_textをmonkeypatchし、片方のファイルパスだけ例外を送出させる。
#     実ファイルは一時ディレクトリに置き、read_file_textをすり替えるだけで
#     Ollama/ネットワーク呼び出しは一切行わない。
_orig_read_file_text = f.read_file_text
with _tempfile.TemporaryDirectory() as _rag_dir2:
    _rag_root2 = _pathlib.Path(_rag_dir2)
    _bad_fp = _rag_root2 / "corrupt.xlsx"
    _good_fp = _rag_root2 / "good.txt"
    _bad_fp.write_bytes(b"not a real xlsx file, just garbage bytes")
    _good_text = ("これは正常に読めるファイルの本文です。" * 10 + "\n") * 5
    _good_fp.write_text(_good_text, encoding="utf-8")

    def _fake_read_file_text(path):
        if _pathlib.Path(path).name == "corrupt.xlsx":
            raise ValueError("simulated corrupt file read failure")
        return _orig_read_file_text(path)

    try:
        f.read_file_text = _fake_read_file_text
        _rag_chunks2 = f._load_rag_chunks([str(_rag_root2)])
    finally:
        f.read_file_text = _orig_read_file_text

    check("_load_rag_chunks: 1ファイルの読み込み例外で全体が例外送出しない(到達できていること自体が検証)",
          True)
    _rag_by_file2 = {}
    for _fp2, _chunk2 in _rag_chunks2:
        _rag_by_file2.setdefault(_os.path.basename(_fp2), []).append((_fp2, _chunk2))
    check("_load_rag_chunks: 読み込み失敗したファイルのチャンクは一切含まれない",
          "corrupt.xlsx" not in _rag_by_file2)
    check("_load_rag_chunks: 正常ファイルのチャンクは含まれる",
          len(_rag_by_file2.get("good.txt", [])) >= 1)
    _expected_good_chunks2 = []
    _start2 = 0
    while _start2 < len(_good_text):
        _end2 = _start2 + f.RAG_CHUNK_CHARS
        _expected_good_chunks2.append((str(_good_fp), _good_text[_start2:_end2]))
        _start2 += f.RAG_CHUNK_CHARS - f.RAG_CHUNK_OVERLAP
    check("_load_rag_chunks: 正常ファイルの(パス,チャンク)タプルが正しい",
          _rag_by_file2.get("good.txt", []) == _expected_good_chunks2)

# (2) ディレクトリ内が「読み込みに失敗するファイルのみ」の場合、空リストを返し例外を送出しない。
with _tempfile.TemporaryDirectory() as _rag_dir3:
    _rag_root3 = _pathlib.Path(_rag_dir3)
    (_rag_root3 / "onlybad.xlsx").write_bytes(b"garbage garbage garbage")

    def _always_fail_read_file_text(path):
        raise RuntimeError("simulated total read failure")

    try:
        f.read_file_text = _always_fail_read_file_text
        _rag_chunks3 = f._load_rag_chunks([str(_rag_root3)])
    finally:
        f.read_file_text = _orig_read_file_text

    check("_load_rag_chunks: 全ファイルが読み込み失敗するディレクトリでは空リストを返す",
          _rag_chunks3 == [])

# (3) 回帰: 正常に読めるファイルのみのディレクトリでは、変更前と完全に同一のチャンク出力
#     (境界・オーバーラップ・順序含め)になること。
with _tempfile.TemporaryDirectory() as _rag_dir4:
    _rag_root4 = _pathlib.Path(_rag_dir4)
    _text_a4 = ("ファイルAの本文テキストです。" * 15 + "\n") * 8
    _text_b4 = ("File B plain ascii content line. " * 20 + "\n") * 6
    (_rag_root4 / "a_file.txt").write_text(_text_a4, encoding="utf-8")
    (_rag_root4 / "b_file.md").write_text(_text_b4, encoding="utf-8")

    _rag_chunks4 = f._load_rag_chunks([str(_rag_root4)])

    _expected_chunks4 = []
    for _fname4, _text4 in sorted([("a_file.txt", _text_a4), ("b_file.md", _text_b4)]):
        _fp4 = str(_rag_root4 / _fname4)
        _start4 = 0
        while _start4 < len(_text4):
            _end4 = _start4 + f.RAG_CHUNK_CHARS
            _expected_chunks4.append((_fp4, _text4[_start4:_end4]))
            _start4 += f.RAG_CHUNK_CHARS - f.RAG_CHUNK_OVERLAP
    check("_load_rag_chunks: 正常ファイルのみの場合はチャンク出力が変更前とバイト単位で完全一致(境界/オーバーラップ/順序)",
          _rag_chunks4 == _expected_chunks4)

# (4) 任意: 実際のリーダー例外経路の検証（monkeypatchではなく、本物の壊れた.xlsxを
#     _read_excel に読ませて例外を発生させる）。openpyxlが利用可能な環境でのみ実施。
try:
    import openpyxl as _openpyxl_probe2  # noqa: F401
    _HAS_OPENPYXL2 = True
except Exception:
    _HAS_OPENPYXL2 = False

if _HAS_OPENPYXL2:
    with _tempfile.TemporaryDirectory() as _rag_dir5:
        _rag_root5 = _pathlib.Path(_rag_dir5)
        # 本物の壊れた.xlsx（ZIP/XMLとして無効なゴミバイト列）
        (_rag_root5 / "broken.xlsx").write_bytes(b"\x00\x01\x02not a zip or xlsx file at all\xff\xfe")
        _good_text5 = "This is a genuinely readable plain text file for RAG.\n" * 30
        (_rag_root5 / "readable.txt").write_text(_good_text5, encoding="utf-8")

        _rag_chunks5 = f._load_rag_chunks([str(_rag_root5)])
        _rag_by_file5 = {}
        for _fp5, _chunk5 in _rag_chunks5:
            _rag_by_file5.setdefault(_os.path.basename(_fp5), []).append(_chunk5)

        check("_load_rag_chunks: 本物の破損.xlsx(実リーダー例外)はスキップされる",
              "broken.xlsx" not in _rag_by_file5)
        check("_load_rag_chunks: 破損.xlsxと同居する正常な.txtは読み込まれる",
              len(_rag_by_file5.get("readable.txt", [])) >= 1)
else:
    print("   [SKIP] openpyxl未インストールのため実.xlsx破損読み込みテストをスキップ")

# ---------- _load_rag_chunks: 非正値step(RAG_CHUNK_OVERLAP>=RAG_CHUNK_CHARS)による
# 無限ループガード (2026-07-25) ----------
# RAG_CHUNK_CHARS/RAG_CHUNK_OVERLAP は本ファイル上部でチューニング可能なモジュール
# 定数として明記されている。RAG_CHUNK_OVERLAP >= RAG_CHUNK_CHARS に設定されると、
# 従来のコード（start += RAG_CHUNK_CHARS - RAG_CHUNK_OVERLAP）は非正値のstepを
# 生み、while start < len(text) のstartが二度と前進せず無限ループ（ハング）に
# なっていた。fugu_local.py側の修正は
# `step = max(1, RAG_CHUNK_CHARS - RAG_CHUNK_OVERLAP)` によるクランプ。
# ここではRAG_CHUNK_CHARS/RAG_CHUNK_OVERLAPを一時的に書き換えて検証するため、
# 既存の慣習に倣いtry/finallyで必ず既定値へ復元する（復元しないと後続の
# iter37/42 RAG回帰テストが既定値を前提に壊れる）。
_orig_rag_chunk_chars = f.RAG_CHUNK_CHARS
_orig_rag_chunk_overlap = f.RAG_CHUNK_OVERLAP

# (a) 既定値でのクランプが厳密にno-opであること（RAG_CHUNK_OVERLAP < RAG_CHUNK_CHARS
#     のとき max(1, positive) == positive）。
check("_load_rag_chunks: 既定値(600/100)ではクランプ後のstepが従来通りRAG_CHUNK_CHARS-RAG_CHUNK_OVERLAP(=500)と一致",
      max(1, f.RAG_CHUNK_CHARS - f.RAG_CHUNK_OVERLAP) == f.RAG_CHUNK_CHARS - f.RAG_CHUNK_OVERLAP == 500)

# (b) 有効な設定(オーバーラップ<チャンクサイズ)での回帰: 手計算した期待チャンクと
#     バイト単位で一致すること。既存のiter37/42テストとは独立に、生成式を使わず
#     文字列リテラルで期待値を書き下す。RAG_CHUNK_CHARS=10, RAG_CHUNK_OVERLAP=4
#     (step=6)、text="ABCDEFGHIJKLMNOP"(16文字)の場合:
#       start=0  -> text[0:10]  = "ABCDEFGHIJ"
#       start=6  -> text[6:16]  = "GHIJKLMNOP"
#       start=12 -> text[12:22] = "MNOP"（len=16のため実際は[12:16]相当）
#       start=18 -> 18 < 16 は偽なので終了
try:
    f.RAG_CHUNK_CHARS = 10
    f.RAG_CHUNK_OVERLAP = 4
    with _tempfile.TemporaryDirectory() as _rag_dir6:
        _rag_root6 = _pathlib.Path(_rag_dir6)
        _text6 = "ABCDEFGHIJKLMNOP"
        (_rag_root6 / "small.txt").write_text(_text6, encoding="utf-8")
        _rag_chunks6 = f._load_rag_chunks([str(_rag_root6)])
    _chunk_texts6 = [c for (_p6, c) in _rag_chunks6]
    check("_load_rag_chunks: 有効config(overlap<chars)での手計算チャンクとバイト単位一致",
          _chunk_texts6 == ["ABCDEFGHIJ", "GHIJKLMNOP", "MNOP"])
finally:
    f.RAG_CHUNK_CHARS = _orig_rag_chunk_chars
    f.RAG_CHUNK_OVERLAP = _orig_rag_chunk_overlap

# (c) 誤設定(RAG_CHUNK_OVERLAP == RAG_CHUNK_CHARS、stepが0になるケース)でも
#     有限かつ完結すること（クランプが無ければここでテストプロセスごと無限ループに
#     なり、以下には絶対に到達しない＝到達できていること自体が検証）。
try:
    f.RAG_CHUNK_CHARS = 5
    f.RAG_CHUNK_OVERLAP = 5
    with _tempfile.TemporaryDirectory() as _rag_dir7:
        _rag_root7 = _pathlib.Path(_rag_dir7)
        _text7 = "ABCDEFGHIJKL"  # 12文字
        (_rag_root7 / "eq.txt").write_text(_text7, encoding="utf-8")
        _rag_chunks7 = f._load_rag_chunks([str(_rag_root7)])
    check("_load_rag_chunks: RAG_CHUNK_OVERLAP==RAG_CHUNK_CHARS(step=0would-be)でもハングせず到達",
          True)
    check("_load_rag_chunks: RAG_CHUNK_OVERLAP==RAG_CHUNK_CHARSでもチャンク数は有限かつlen(text)以下",
          0 < len(_rag_chunks7) <= len(_text7))
finally:
    f.RAG_CHUNK_CHARS = _orig_rag_chunk_chars
    f.RAG_CHUNK_OVERLAP = _orig_rag_chunk_overlap

# (d) 誤設定(RAG_CHUNK_OVERLAP > RAG_CHUNK_CHARS、stepが負になるケース)でも
#     同様に有限かつ完結すること。
try:
    f.RAG_CHUNK_CHARS = 600
    f.RAG_CHUNK_OVERLAP = 700
    with _tempfile.TemporaryDirectory() as _rag_dir8:
        _rag_root8 = _pathlib.Path(_rag_dir8)
        _text8 = "x" * 50
        (_rag_root8 / "over.txt").write_text(_text8, encoding="utf-8")
        _rag_chunks8 = f._load_rag_chunks([str(_rag_root8)])
    check("_load_rag_chunks: RAG_CHUNK_OVERLAP>RAG_CHUNK_CHARS(step負would-be)でもハングせず到達",
          True)
    check("_load_rag_chunks: RAG_CHUNK_OVERLAP>RAG_CHUNK_CHARSでもチャンク数は有限かつlen(text)以下",
          0 < len(_rag_chunks8) <= len(_text8))
finally:
    f.RAG_CHUNK_CHARS = _orig_rag_chunk_chars
    f.RAG_CHUNK_OVERLAP = _orig_rag_chunk_overlap

# (e) 空ファイル(0バイト)は従来通りチャンク0件（`while start < len(text)` は
#     text=""のとき即偽になるため、このクランプ変更でも挙動は変わらない）。
with _tempfile.TemporaryDirectory() as _rag_dir9:
    _rag_root9 = _pathlib.Path(_rag_dir9)
    (_rag_root9 / "empty.txt").write_text("", encoding="utf-8")
    _rag_chunks9 = f._load_rag_chunks([str(_rag_root9)])
check("_load_rag_chunks: 空ファイル(0バイト)はチャンク0件のまま(クランプの影響を受けない)",
      _rag_chunks9 == [])

# (f) 設定復元の確認: 既定値(600/100)に戻っていること（後続テストへの汚染防止）。
check("_load_rag_chunks: RAG_CHUNK_CHARS/RAG_CHUNK_OVERLAPが既定値に復元されている",
      f.RAG_CHUNK_CHARS == _orig_rag_chunk_chars and f.RAG_CHUNK_OVERLAP == _orig_rag_chunk_overlap)

# ---------- read_file_text: ディスパッチ例外を握りつぶして""を返す (2026-07-23 / iter53) ----------
# _read_pdf/_read_docx/_read_excel/_read_pptx/_read_html は import 文と実パース処理を
# 同じ try/except ImportError で包んでいるため、ライブラリ自体は入っているがファイルが
# 壊れている場合の非ImportError例外（zipfile.BadZipFile, PDFSyntaxError等）が
# read_file_text の外まで素通しで伝播していた。main()の--file呼び出し
# （`read_file_text(fp).strip()`）は_load_rag_chunksと違って無防備だったため、
# 壊れたOffice/PDFファイルを渡すとCLI全体がクラッシュしていた（iter42はRAG経路のみ保護、
# 全リーダー書き換えを試みたiter51は行き詰まった）。ここではread_file_text自身が
# 例外を握りつぶして""を返すことを、個々のリーダーをmonkeypatchして直接検証する
# （実ライブラリ・実ファイル・Ollama呼び出しは一切不要）。
import zipfile as _zipfile

_orig_read_pdf2 = f._read_pdf
_orig_read_excel2 = f._read_excel


def _rft_raise_value_error(path):
    raise ValueError("simulated corrupt PDF parse failure")


def _rft_raise_bad_zip(path):
    raise _zipfile.BadZipFile("simulated corrupt xlsx (not a zip file)")


try:
    f._read_pdf = _rft_raise_value_error
    with contextlib.redirect_stdout(io.StringIO()) as _rft_out1:
        _rft_result1 = f.read_file_text(_pathlib.Path("dummy_corrupt.pdf"))
finally:
    f._read_pdf = _orig_read_pdf2

check("read_file_text: _read_pdfが非ImportError例外(ValueError)を送出しても伝播せず\"\"を返す",
      _rft_result1 == "")
check("read_file_text: PDF読み込み失敗時に警告(ファイル名+例外型)を表示する",
      "dummy_corrupt.pdf" in _rft_out1.getvalue() and "ValueError" in _rft_out1.getvalue())

try:
    f._read_excel = _rft_raise_bad_zip
    with contextlib.redirect_stdout(io.StringIO()) as _rft_out2:
        _rft_result2 = f.read_file_text(_pathlib.Path("dummy_corrupt.xlsx"))
finally:
    f._read_excel = _orig_read_excel2

check("read_file_text: _read_excelがBadZipFileを送出しても伝播せず\"\"を返す(ディスパッチ全体の保護を確認)",
      _rft_result2 == "")
check("read_file_text: Excel読み込み失敗時に警告(ファイル名+例外型)を表示する",
      "dummy_corrupt.xlsx" in _rft_out2.getvalue() and "BadZipFile" in _rft_out2.getvalue())

# 回帰: 正常なテキストファイル/_BINARY_SKIP拡張子は従来通りの挙動を維持
with _tempfile.TemporaryDirectory() as _rft_dir:
    _rft_root = _pathlib.Path(_rft_dir)
    _rft_txt_content = "これは通常のテキストファイルです。\nsecond line.\n"
    (_rft_root / "plain.txt").write_text(_rft_txt_content, encoding="utf-8")
    (_rft_root / "note.md").write_text(_rft_txt_content, encoding="utf-8")
    check("read_file_text: 通常の.txtファイルはバイト単位で内容が一致(回帰)",
          f.read_file_text(_rft_root / "plain.txt") == _rft_txt_content)
    check("read_file_text: 通常の.mdファイルはバイト単位で内容が一致(回帰)",
          f.read_file_text(_rft_root / "note.md") == _rft_txt_content)

    (_rft_root / "image.png").write_bytes(b"\x89PNG\r\n\x1a\nnot a real png")
    check("read_file_text: _BINARY_SKIP拡張子(.png)は従来通り\"\"を返す(回帰)",
          f.read_file_text(_rft_root / "image.png") == "")

# ライブラリ未インストール通知文字列（例外ではなくリーダーが正常returnした場合）は
# ""に変換されず、そのままバイト単位で素通しされること（_is_lib_missing_notice判定は
# 呼び出し側=_load_rag_chunksの仕事であり、read_file_text自体は加工しない）。
_rft_notice_text = "[Excel: dummy.xlsx — openpyxl or pandas が必要: pip install openpyxl]"


def _rft_return_notice(path):
    return _rft_notice_text


try:
    f._read_excel = _rft_return_notice
    _rft_result3 = f.read_file_text(_pathlib.Path("dummy_notice.xlsx"))
finally:
    f._read_excel = _orig_read_excel2

check("read_file_text: リーダーがpip install通知文字列を正常returnした場合はそのまま素通し(\"\"化されない)",
      _rft_result3 == _rft_notice_text)

# ---------- read_file_text: 汎用テキストフォールバック分岐の例外も警告を出す (2026-07-24 / iter125) ----------
# 上のiter53テストはsuffix-dispatch分岐(.pdf/.docx/.xlsx/.pptx/.html/.ipynb)のみをカバーする。
# 未知拡張子(.txt等)は最後のelse=汎用テキストフォールバック分岐(Path.read_bytes +
# _decode_text_bytes)を通るが、従来この分岐は except Exception: return "" のみで、
# 権限エラー(PermissionError)やread_bytes自体の失敗が無警告のまま握りつぶされていた。
# ここではPath.read_bytesをmonkeypatchしてPermissionErrorを送出させ、
# read_file_textが(1)従来通り""を返す((2)警告print(ファイル名+例外型)を出すことを確認する。
_orig_path_read_bytes = _pathlib.Path.read_bytes


def _rft_raise_permission_error(self):
    raise PermissionError("simulated permission denied")


try:
    _pathlib.Path.read_bytes = _rft_raise_permission_error
    with contextlib.redirect_stdout(io.StringIO()) as _rft_out4:
        _rft_result4 = f.read_file_text(_pathlib.Path("dummy_protected.txt"))
finally:
    _pathlib.Path.read_bytes = _orig_path_read_bytes

check("read_file_text: 汎用テキスト分岐でPath.read_bytesがPermissionErrorを送出しても伝播せず\"\"を返す(契約不変)",
      _rft_result4 == "")
check("read_file_text: 汎用テキスト分岐の読み込み失敗時にも警告(ファイル名+例外型)を表示する(従来は無警告だった)",
      "dummy_protected.txt" in _rft_out4.getvalue() and "PermissionError" in _rft_out4.getvalue())

# ---------- _read_html / _read_ipynb: 直接カバレッジ (2026-07-23 / iter71) ----------
# _read_html/_read_ipynb はstdlibのみ(html.parser/json)で書かれた汎用ファイルリーダーで、
# RAG(_load_rag_chunks -> rag_search -> build_context)や --file 経由でproposerの
# プロンプトに直接混入する（= 精度優先の回答経路に効くコンテキスト源）。これまでは
# iter53のread_file_textディスパッチ例外テストがコメント中で言及するのみで、
# HTML/Notebookパース処理そのものを直接呼び出す検証が一切なかった。
# ここでは f._read_html(path) / f._read_ipynb(path) を一時ファイルに対して直接呼び、
# f.ask/urlopen/subprocessは一切mockしない（本物のパースロジックのみを検証）。
with _tempfile.TemporaryDirectory() as _rh_dir:
    _rh_root = _pathlib.Path(_rh_dir)

    # (a) 通常のタグ除去: 可視テキストが残る
    _rh_a = _rh_root / "a.html"
    _rh_a.write_text(
        "<html><body><div><h1>Title</h1><p>Paragraph one.</p>"
        "<p>Paragraph two.</p></div></body></html>",
        encoding="utf-8")
    _rh_out_a = f._read_html(_rh_a)
    check("_read_html: 見出し/段落の可視テキストが残る",
          "Title" in _rh_out_a and "Paragraph one." in _rh_out_a
          and "Paragraph two." in _rh_out_a)
    check("_read_html: タグ自体は出力に残らない(山括弧なし)",
          "<" not in _rh_out_a and ">" not in _rh_out_a)

    # (b) <script>/<style>の中身は除外される
    _rh_b = _rh_root / "b.html"
    _rh_b.write_text(
        "<html><head><style>body{color:red}</style></head><body>"
        "<script>alert('hi'); var secret = 42;</script>"
        "<p>Visible text here</p></body></html>",
        encoding="utf-8")
    _rh_out_b = f._read_html(_rh_b)
    check("_read_html: <script>本文は出力に含まれない",
          "alert" not in _rh_out_b and "secret" not in _rh_out_b)
    check("_read_html: <style>本文は出力に含まれない",
          "color" not in _rh_out_b and "red" not in _rh_out_b)
    check("_read_html: <script>/<style>と同居する通常テキストは残る",
          "Visible text here" in _rh_out_b)

    # (c) HTML文字参照(&amp; &lt; &gt;)がデコードされる(HTMLParserのconvert_charrefs既定)
    _rh_c = _rh_root / "c.html"
    _rh_c.write_text(
        "<p>Fish &amp; Chips &lt;tasty&gt;</p>", encoding="utf-8")
    _rh_out_c = f._read_html(_rh_c)
    check("_read_html: &amp;が'&'にデコードされる",
          "Fish & Chips" in _rh_out_c)
    check("_read_html: &lt;/&gt;が'<'/'>'にデコードされる",
          "<tasty>" in _rh_out_c)
    check("_read_html: デコード後もエンティティの生表記は残らない",
          "&amp;" not in _rh_out_c and "&lt;" not in _rh_out_c)

    # (d) テキストを持たない文書(空白のみ)は空文字を返す
    _rh_d = _rh_root / "d.html"
    _rh_d.write_text(
        "<html><body>   <div>\n\n   </div>\t</body></html>", encoding="utf-8")
    check("_read_html: テキストなし文書は空文字を返す",
          f._read_html(_rh_d) == "")

    # ---------- _read_html: ブロック対応セパレータ (2026-07-26) ----------
    # 旧実装はhandle_dataのテキストノードごとにdata.strip()して"\n".join()して
    # いたため、<b>/<strong>/<a>/<sub>/<sup>/<span>/<code>のようなインライン
    # 要素が子孫テキストを別ノードに分割するだけで、文/フレーズの断片化と
    # ノード間空白の消失を引き起こしていた。特に日本語では、この断片化が
    # 直上(iter179)で追加したCJKバイグラムトークナイザ(_tokenize)を無力化し
    # ていた（改行を挟むと境界を跨ぐバイグラムが二度と生成されない）。以下は
    # 修正（ブロックレベルタグの開始/終了時のみ区切りを挿入し、インライン/
    # 未知タグでは何も挿入しない・テキストノードはstripせず生のまま連結）の
    # 直接検証。iter71の直接カバレッジ・iter94のcp932デコードラダーと合わせて
    # 参照。

    # (e) インライン要素をまたぐ英文が断片化せず1つに連結される
    # (旧実装は'The\nquick\nbrown fox'に断片化していた)
    _rh_e = _rh_root / "e.html"
    _rh_e.write_text("<p>The <b>quick</b> brown fox</p>", encoding="utf-8")
    _rh_out_e = f._read_html(_rh_e)
    check("_read_html(2026-07-26): インライン要素をまたぐ文が断片化せず連結される",
          "The quick brown fox" in _rh_out_e)
    check("_read_html(2026-07-26): インライン結合結果にもタグは残らない(山括弧なし)",
          "<" not in _rh_out_e and ">" not in _rh_out_e)

    # (f) 日本語: インライン要素をまたぐ断片化はiter179のCJKバイグラム
    # トークナイザを無力化していた('機械\n学習'では境界バイグラム'械学'が
    # 生成されない)。修正後は1行に連結され、境界バイグラムが復活することを
    # _tokenize経由で直接確認する。
    _rh_f = _rh_root / "f.html"
    _rh_f.write_text(
        "<p>機械<strong>学習</strong>技術の解説</p>", encoding="utf-8")
    _rh_out_f = f._read_html(_rh_f)
    check("_read_html(2026-07-26): 日本語がインライン要素をまたいでも断片化せず連結される",
          "機械学習技術の解説" in _rh_out_f)
    check("_read_html(2026-07-26/iter179再現率修正): 修正後の出力から境界バイグラム'械学'が_tokenizeで得られる",
          "械学" in f._tokenize(_rh_out_f))
    check("_read_html(2026-07-26/対照): 断片化された'機械\\n学習技術の解説'相当の文字列には'械学'が含まれない"
          "(=修正前の挙動ではiter179の再現率が回復しないことの確認)",
          "械学" not in f._tokenize("機械\n学習技術の解説"))

    # (g) '10<sup>3</sup>'のような上付き表記は'103'に連結される(旧実装の
    # '10\n3'分断を回避)。真の上付き/下付きの意味論はプレーンテキスト抽出
    # からは原理的に復元不能で、本修正が保証するのは「分断されない」ことまで。
    _rh_g = _rh_root / "g.html"
    _rh_g.write_text("<p>10<sup>3</sup></p>", encoding="utf-8")
    check("_read_html(2026-07-26): '10<sup>3</sup>'が'103'に連結され'10/3'に分断されない",
          "103" in f._read_html(_rh_g))

    # (h) 回帰: ブロックレベルの区切り(table/td/th/tr, ul/li, br)は引き続き
    # 行として分離されること。インライン修正が誤ってセル/項目/改行を1行に
    # 潰してしまわないことの確認。
    _rh_h1 = _rh_root / "h1.html"
    _rh_h1.write_text(
        "<table><tr><td>A</td><td>B</td></tr></table>", encoding="utf-8")
    check("_read_html(2026-07-26回帰): <td>セルはインライン修正後も別行のまま(A/Bが結合しない)",
          f._read_html(_rh_h1) == "A\nB")

    _rh_h2 = _rh_root / "h2.html"
    _rh_h2.write_text("<ul><li>x</li><li>y</li></ul>", encoding="utf-8")
    check("_read_html(2026-07-26回帰): <li>項目はインライン修正後も別行のまま(x/yが結合しない)",
          f._read_html(_rh_h2) == "x\ny")

    _rh_h3 = _rh_root / "h3.html"
    _rh_h3.write_text("A<br>B", encoding="utf-8")
    check("_read_html(2026-07-26回帰): <br>で区切られたA/Bも別行のまま(結合しない)",
          f._read_html(_rh_h3) == "A\nB")

    # (i) 回帰再確認: 単一段落の完全一致(exact-equality)はブロック区切り
    # 導入後も崩れない(既存の~L5550の完全一致回帰とは別の新規フィクスチャ)。
    _rh_i = _rh_root / "i.html"
    _rh_i.write_text(
        "<html><body><p>Hello again world. こんにちは、世界。</p></body></html>",
        encoding="utf-8")
    check("_read_html(2026-07-26回帰再確認): 単一段落の完全一致は崩れない",
          f._read_html(_rh_i) == "Hello again world. こんにちは、世界。")

    # (j) 回帰再確認: cp932保存HTMLの日本語復元(iter94のutf-8→cp932→replace
    # ラダー)はブロック区切り導入後も無傷であること。
    _rh_j = _rh_root / "j.html"
    _rh_j.write_bytes(
        "<html><body><h2>速報</h2><p>台風が接近しています。</p></body></html>"
        .encode("cp932"))
    _rh_out_j = f._read_html(_rh_j)
    check("_read_html(2026-07-26回帰再確認): cp932保存HTMLの日本語復元は無傷",
          "速報" in _rh_out_j and "台風が接近しています。" in _rh_out_j
          and "�" not in _rh_out_j and "<" not in _rh_out_j)

    # (k) 回帰再確認: HTML文字参照のデコードはブロック区切り導入後も無傷。
    _rh_k = _rh_root / "k.html"
    _rh_k.write_text("<p>Rock &amp; Roll &lt;loud&gt;</p>", encoding="utf-8")
    _rh_out_k = f._read_html(_rh_k)
    check("_read_html(2026-07-26回帰再確認): エンティティデコード('&amp;'->'&'等)は無傷",
          "Rock & Roll" in _rh_out_k and "<loud>" in _rh_out_k)

    # (l) 回帰再確認: <script>/<style>本文除外はブロック区切り導入後も無傷。
    _rh_l = _rh_root / "l.html"
    _rh_l.write_text(
        "<html><head><style>p{color:blue}</style></head><body>"
        "<script>doEvil(); var token = 99;</script>"
        "<p>Still visible</p></body></html>",
        encoding="utf-8")
    _rh_out_l = f._read_html(_rh_l)
    check("_read_html(2026-07-26回帰再確認): <script>/<style>本文除外は無傷",
          "doEvil" not in _rh_out_l and "token" not in _rh_out_l
          and "color" not in _rh_out_l and "blue" not in _rh_out_l
          and "Still visible" in _rh_out_l)

with _tempfile.TemporaryDirectory() as _ri_dir:
    _ri_root = _pathlib.Path(_ri_dir)

    # (a) 整形されたNotebook: コードセルは```python フェンス、markdownセルは平文、
    #     空白のみ/空文字セルはスキップ、単一文字列source(リストでなくとも正当なnbformat)も扱える
    _ri_nb_a = {
        "cells": [
            {"cell_type": "code", "source": ["import os\n", "print(os.getcwd())"]},
            {"cell_type": "markdown", "source": ["# Title\n", "some text"]},
            {"cell_type": "code", "source": ["   \n", "  "]},  # 空白のみ -> skip
            {"cell_type": "markdown", "source": ""},  # 空文字 -> skip
            {"cell_type": "code", "source": "x = 1\nprint(x)"},  # 単一文字列source
        ]
    }
    _ri_a = _ri_root / "a.ipynb"
    _ri_a.write_text(json.dumps(_ri_nb_a), encoding="utf-8")
    _ri_out_a = f._read_ipynb(_ri_a)
    _ri_expected_a = (
        "```python\nimport os\nprint(os.getcwd())\n```"
        "\n\n# Title\nsome text"
        "\n\n```python\nx = 1\nprint(x)\n```"
    )
    check("_read_ipynb: コード/markdownセルの整形・空白セルのスキップ・単一文字列sourceが期待通り",
          _ri_out_a == _ri_expected_a)
    check("_read_ipynb: 空白のみセルの痕跡(空フェンス等)が出力に残らない",
          "```python\n\n```" not in _ri_out_a and "```python\n   \n```" not in _ri_out_a)

    # (b) 壊れた(非JSON)notebookはexcept分岐でraw読み込みにフォールバックし、例外を送出しない
    _ri_b = _ri_root / "b.ipynb"
    _ri_bad_text = "this is { not : valid json at all"
    _ri_b.write_text(_ri_bad_text, encoding="utf-8")
    _ri_out_b = f._read_ipynb(_ri_b)
    check("_read_ipynb: 非JSONファイルは例外を送出せずraw textを返す",
          _ri_out_b == _ri_bad_text)

    # (c) 修正後挙動(iter72): JSONとしては妥当だが cellの'source'がjoin可能な
    #     list[str]ではない「壊れているがパース可能」なnotebook。かつては"".join()が
    #     TypeErrorを送出し、_read_ipynbの外側try/exceptで捕捉されて問題のセルだけで
    #     なくnotebook全体が生JSON全文としてRAG/--fileコンテキストへ丸ごと注入されて
    #     いた（iteration 71がtest-onlyで発見し特性検証テストとしてピン留め、将来の
    #     iterationでの修正候補と明示的にフラグしていたもの）。iter72で該当セルのみを
    #     安全にスキップするよう修正したため、ここでは修正後の正しい挙動を検証する。
    _ri_nb_c = {"cells": [{"cell_type": "code", "source": None}]}
    _ri_c = _ri_root / "c.ipynb"
    _ri_raw_c = json.dumps(_ri_nb_c)
    _ri_c.write_text(_ri_raw_c, encoding="utf-8")
    _ri_out_c = f._read_ipynb(_ri_c)
    check("_read_ipynb: source=Noneのセルはskipされ空文字を返す(生JSONダンプではない)",
          _ri_out_c == "" and _ri_out_c != _ri_raw_c)

    _ri_nb_d = {"cells": [{"cell_type": "code", "source": ["ok\n", 42]}]}
    _ri_d = _ri_root / "d.ipynb"
    _ri_raw_d = json.dumps(_ri_nb_d)
    _ri_d.write_text(_ri_raw_d, encoding="utf-8")
    _ri_out_d = f._read_ipynb(_ri_d)
    check("_read_ipynb: source内に非文字列混入のセルは文字列要素のみjoinされ非str要素は無視される",
          _ri_out_d == "```python\nok\n\n```")
    check("_read_ipynb: 非文字列混入セルの結果は生JSONダンプではない",
          _ri_out_d != _ri_raw_d and "cell_type" not in _ri_out_d and "source" not in _ri_out_d)

    # (e) 良好なセルと壊れたセル(source=None)が混在する場合: 壊れたセルのみskipされ、
    #     良好なセルの構造化抽出は維持される（1セルの破損でnotebook全体が道連れにならない）。
    _ri_nb_e = {
        "cells": [
            {"cell_type": "code", "source": ["import os\n", "print(1)"]},
            {"cell_type": "code", "source": None},
        ]
    }
    _ri_e = _ri_root / "e.ipynb"
    _ri_raw_e = json.dumps(_ri_nb_e)
    _ri_e.write_text(_ri_raw_e, encoding="utf-8")
    _ri_out_e = f._read_ipynb(_ri_e)
    check("_read_ipynb: 良好なセル+壊れたセル混在時は良好なセルのみ構造化抽出される",
          _ri_out_e == "```python\nimport os\nprint(1)\n```")
    check("_read_ipynb: 良好+壊れたセル混在時も生JSONマーカーが出力に含まれない",
          "cell_type" not in _ri_out_e and "source" not in _ri_out_e)

    # (f) cellsリスト内に非dictの要素が混じっている場合: そのエントリだけskipし、
    #     以降の正当なセルは正常に構造化抽出される(cell.get(...)でAttributeErrorを
    #     起こして生JSONダンプへ道連れにならない)。
    _ri_nb_f = {
        "cells": [
            "junk",
            {"cell_type": "markdown", "source": ["# Title"]},
        ]
    }
    _ri_f = _ri_root / "f.ipynb"
    _ri_raw_f = json.dumps(_ri_nb_f)
    _ri_f.write_text(_ri_raw_f, encoding="utf-8")
    _ri_out_f = f._read_ipynb(_ri_f)
    check("_read_ipynb: cells内の非dictエントリはskipされ後続の正当なセルは抽出される",
          _ri_out_f == "# Title")

    # (g) json.loadsは成功するがトップレベルがdictでない(配列)場合: iter72はセル単位の
    #     壊れたsourceのみ対処しており、トップレベル構造自体が非dictだと
    #     nb.get("cells", [])がAttributeErrorを送出して外側exceptに落ち、notebook
    #     全体の生JSON(この場合は配列全文)がRAG/--fileコンテキストへ丸ごと注入
    #     されていた(iter71がフラグ、iter72未対応部分)。本iterationでの修正後は
    #     例外を送出せず空文字列を返す。
    _ri_g = _ri_root / "g.ipynb"
    _ri_raw_g = json.dumps([{"cell_type": "code", "source": ["x = 1"]}])
    _ri_g.write_text(_ri_raw_g, encoding="utf-8")
    _ri_out_g = f._read_ipynb(_ri_g)
    check("_read_ipynb: トップレベルが配列(非dict)の場合は例外を送出せず空文字を返す",
          _ri_out_g == "")
    check("_read_ipynb: トップレベル配列ケースの結果に生JSONマーカーが含まれない",
          "cell_type" not in _ri_out_g and "source" not in _ri_out_g and "[" not in _ri_out_g)

    # (g2) json.loadsは成功するがトップレベルが裸のスカラー(数値)の場合: (g)と同じ
    #     AttributeErrorの経路だが、int.get自体が存在しないという別の失敗形。
    _ri_g2 = _ri_root / "g2.ipynb"
    _ri_raw_g2 = json.dumps(12345)
    _ri_g2.write_text(_ri_raw_g2, encoding="utf-8")
    _ri_out_g2 = f._read_ipynb(_ri_g2)
    check("_read_ipynb: トップレベルが裸の数値(非dict)の場合は例外を送出せず空文字を返す",
          _ri_out_g2 == "")

    # (g3) json.loadsは成功するがトップレベルが裸の文字列の場合: 同様にAttributeError経路。
    _ri_g3 = _ri_root / "g3.ipynb"
    _ri_raw_g3 = json.dumps("not a notebook")
    _ri_g3.write_text(_ri_raw_g3, encoding="utf-8")
    _ri_out_g3 = f._read_ipynb(_ri_g3)
    check("_read_ipynb: トップレベルが裸の文字列(非dict)の場合は例外を送出せず空文字を返す",
          _ri_out_g3 == "")

    # (h) トップレベルはdictだが nb["cells"] がtruthyな非list(整数)の場合: 旧実装は
    #     for cell in <int> でTypeErrorを送出し外側exceptで生JSON全文がRAG/--file
    #     コンテキストへ丸ごと注入されていた。修正後はcellsを[]へ強制変換して空文字を返す。
    _ri_nb_h = {"cells": 42}
    _ri_h = _ri_root / "h.ipynb"
    _ri_raw_h = json.dumps(_ri_nb_h)
    _ri_h.write_text(_ri_raw_h, encoding="utf-8")
    _ri_out_h = f._read_ipynb(_ri_h)
    check("_read_ipynb: cellsがtruthyな非list(整数)の場合は例外を送出せず空文字を返す",
          _ri_out_h == "")
    check("_read_ipynb: cells=整数ケースの結果に生JSONマーカーが含まれない",
          "cells" not in _ri_out_h and "42" not in _ri_out_h)

    # (i) nb["cells"] がdict(非list)の場合: (h)と同じ非list強制変換で救済される。
    _ri_nb_i = {"cells": {"0": {"cell_type": "code", "source": ["x = 1"]}}}
    _ri_i = _ri_root / "i.ipynb"
    _ri_raw_i = json.dumps(_ri_nb_i)
    _ri_i.write_text(_ri_raw_i, encoding="utf-8")
    _ri_out_i = f._read_ipynb(_ri_i)
    check("_read_ipynb: cellsがdict(非list)の場合は例外を送出せず空文字を返す",
          _ri_out_i == "")
    check("_read_ipynb: cells=dictケースの結果に生JSONマーカーが含まれない",
          "cell_type" not in _ri_out_i and "source" not in _ri_out_i)

# ---------- _read_ipynb: コードセルのstream出力(stdout/stderr)抽出 (2026-07-26) ----------
# _read_ipynbは従来コードセルの'source'(入力コード)のみを抽出し、'outputs'配列
# (実行結果)は完全に無視していた。データ分析notebookではprint(...)が出力する実際の
# 数値・結果こそが質問への回答に直結する事実であることが多く、精度criticalな
# RAG/--fileコンテキストからそれが黙って欠落していた。iteration 188はstream/
# execute_result/display_data(text/plain)の全種類を一度に扱おうとして3回試みたが
# 行き詰まり断念しており、本セクションはあえて"stream"(stdout/stderr)出力1種類のみに
# 絞った縮小スコープの直接検証。'source'正規化(iteration 72)・cellの非dictスキップ
# (iteration 72)・トップレベル構造ガード(iteration 113)は上のセクションで既に検証済みで
# 本変更でも一切変更されていないため、ここではstream出力抽出そのものに焦点を絞る。
# f.ask/urlopen/subprocessは一切呼ばない(すべてtempfile上のオフラインI/O)。
with _tempfile.TemporaryDirectory() as _rio_dir:
    _rio_root = _pathlib.Path(_rio_dir)

    # (j) stream stdout出力(text: list[str])が```python フェンス直後にラベル付きで追加される
    _rio_nb_j = {
        "cells": [
            {"cell_type": "code", "source": ["print('hello')"],
             "outputs": [
                 {"output_type": "stream", "name": "stdout", "text": ["hello\n"]},
             ]},
        ]
    }
    _rio_j = _rio_root / "j.ipynb"
    _rio_j.write_text(json.dumps(_rio_nb_j), encoding="utf-8")
    _rio_out_j = f._read_ipynb(_rio_j)
    check("_read_ipynb: stream stdout(text=list[str])がフェンス直後にラベル付きで追加される",
          _rio_out_j == "```python\nprint('hello')\n```"
          "\n\n[Notebook stdout/stderr output]\nhello\n")

    # (k) stream stdout出力(text: 単一str)も同様に抽出される
    _rio_nb_k = {
        "cells": [
            {"cell_type": "code", "source": ["print('world')"],
             "outputs": [
                 {"output_type": "stream", "name": "stdout", "text": "world\n"},
             ]},
        ]
    }
    _rio_k = _rio_root / "k.ipynb"
    _rio_k.write_text(json.dumps(_rio_nb_k), encoding="utf-8")
    _rio_out_k = f._read_ipynb(_rio_k)
    check("_read_ipynb: stream stdout(text=単一str)も同様に抽出される",
          _rio_out_k == "```python\nprint('world')\n```"
          "\n\n[Notebook stdout/stderr output]\nworld\n")

    # (l) stream stderr出力も同様に抽出される(stdout限定ではない)
    _rio_nb_l = {
        "cells": [
            {"cell_type": "code", "source": ["import sys; sys.stderr.write('oops')"],
             "outputs": [
                 {"output_type": "stream", "name": "stderr", "text": ["oops\n"]},
             ]},
        ]
    }
    _rio_l = _rio_root / "l.ipynb"
    _rio_l.write_text(json.dumps(_rio_nb_l), encoding="utf-8")
    _rio_out_l = f._read_ipynb(_rio_l)
    check("_read_ipynb: stream stderr出力も同様に抽出される",
          _rio_out_l == "```python\nimport sys; sys.stderr.write('oops')\n```"
          "\n\n[Notebook stdout/stderr output]\noops\n")

    # (m) outputsがtruthyな非list(整数)の場合: iteration 113のcells非list強制変換と
    #     同じ作法(`or []`のtruthinessトリックには頼らない)で空へ倒し、例外を送出せず
    #     コード自体の抽出は無傷のまま出力ブロックのみ省略される。
    _rio_nb_m = {"cells": [{"cell_type": "code", "source": ["print(1)"], "outputs": 42}]}
    _rio_m = _rio_root / "m.ipynb"
    _rio_m.write_text(json.dumps(_rio_nb_m), encoding="utf-8")
    _rio_out_m = f._read_ipynb(_rio_m)
    check("_read_ipynb: outputsが非list(整数)でも例外を送出せずコード抽出のみ維持される",
          _rio_out_m == "```python\nprint(1)\n```")

    # (m2) outputsがtruthyな非list(dict)の場合も同様に空へ倒れる
    _rio_nb_m2 = {"cells": [{"cell_type": "code", "source": ["print(1)"],
                              "outputs": {"0": {"output_type": "stream", "text": ["x"]}}}]}
    _rio_m2 = _rio_root / "m2.ipynb"
    _rio_m2.write_text(json.dumps(_rio_nb_m2), encoding="utf-8")
    _rio_out_m2 = f._read_ipynb(_rio_m2)
    check("_read_ipynb: outputsが非list(dict)でも例外を送出せずコード抽出のみ維持される",
          _rio_out_m2 == "```python\nprint(1)\n```")

    # (m3) outputs=None(キーはあるが値がnull)の場合も同様に空へ倒れる
    _rio_nb_m3 = {"cells": [{"cell_type": "code", "source": ["print(1)"], "outputs": None}]}
    _rio_m3 = _rio_root / "m3.ipynb"
    _rio_m3.write_text(json.dumps(_rio_nb_m3), encoding="utf-8")
    _rio_out_m3 = f._read_ipynb(_rio_m3)
    check("_read_ipynb: outputs=Noneでも例外を送出せずコード抽出のみ維持される",
          _rio_out_m3 == "```python\nprint(1)\n```")

    # (n) outputsリスト内に非dictエントリが混在する場合: そのエントリだけskipされ、
    #     後続の正当なstream出力エントリは抽出される(iteration 72の非dictセルskipと
    #     同じ「1件の破損で全体を道連れにしない」パターン)。
    _rio_nb_n = {
        "cells": [
            {"cell_type": "code", "source": ["print(1)"],
             "outputs": ["junk", {"output_type": "stream", "name": "stdout", "text": ["ok\n"]}]},
        ]
    }
    _rio_n = _rio_root / "n.ipynb"
    _rio_n.write_text(json.dumps(_rio_nb_n), encoding="utf-8")
    _rio_out_n = f._read_ipynb(_rio_n)
    check("_read_ipynb: outputs内の非dictエントリはskipされ後続の正当なstream出力は抽出される",
          _rio_out_n == "```python\nprint(1)\n```"
          "\n\n[Notebook stdout/stderr output]\nok\n")

    # (o) stream 'text'のlistに非str要素が混入する場合: source正規化(iteration 72)と
    #     全く同じパターンでstr要素のみjoinされる(42は無視されて'a'+'b'='ab')。
    _rio_nb_o = {
        "cells": [
            {"cell_type": "code", "source": ["print(1)"],
             "outputs": [{"output_type": "stream", "name": "stdout", "text": ["a", 42, "b"]}]},
        ]
    }
    _rio_o = _rio_root / "o.ipynb"
    _rio_o.write_text(json.dumps(_rio_nb_o), encoding="utf-8")
    _rio_out_o = f._read_ipynb(_rio_o)
    check("_read_ipynb: stream textのlist内の非str要素は無視されstr要素のみjoinされる",
          _rio_out_o == "```python\nprint(1)\n```\n\n[Notebook stdout/stderr output]\nab")

    # (o2) stream 'text'がstrでもlistでもない(整数)場合: 空文字へ正規化されblank後
    #     strip判定でそのoutputエントリ自体がskipされる(例外は送出しない)。
    _rio_nb_o2 = {
        "cells": [
            {"cell_type": "code", "source": ["print(1)"],
             "outputs": [{"output_type": "stream", "name": "stdout", "text": 42}]},
        ]
    }
    _rio_o2 = _rio_root / "o2.ipynb"
    _rio_o2.write_text(json.dumps(_rio_nb_o2), encoding="utf-8")
    _rio_out_o2 = f._read_ipynb(_rio_o2)
    check("_read_ipynb: stream textが非str/非list(整数)でも例外を送出せずそのoutputはskipされる",
          _rio_out_o2 == "```python\nprint(1)\n```")

    # (o3) stream 'text'のlistが非str要素のみで構成される場合: join結果が空文字になり
    #     blank後stripでそのoutputエントリがskipされる。
    _rio_nb_o3 = {
        "cells": [
            {"cell_type": "code", "source": ["print(1)"],
             "outputs": [{"output_type": "stream", "name": "stdout", "text": [42, None]}]},
        ]
    }
    _rio_o3 = _rio_root / "o3.ipynb"
    _rio_o3.write_text(json.dumps(_rio_nb_o3), encoding="utf-8")
    _rio_out_o3 = f._read_ipynb(_rio_o3)
    check("_read_ipynb: stream textのlistが非str要素のみの場合もそのoutputはskipされる",
          _rio_out_o3 == "```python\nprint(1)\n```")

    # (p) execute_result/display_data/errorの混在ケース。iter194時点ではこの3種すべて
    #     対象外だったが、iter195がexecute_result/display_dataの'text/plain'のみを、
    #     iter196(本iteration)がerrorの'ename'/'evalue'のみをそれぞれ追加で対象に
    #     したため、期待値を更新する(下の新セクションでより網羅的に検証):
    #     execute_resultのtext/plainは抽出され、display_data側はimage/pngのみ
    #     (text/plainなし)なので抽出されずbase64も混入しない。errorはename/evalueが
    #     [Notebook error]ブロックとして抽出されるが、tracebackは引き続き対象外の
    #     まま混入しないことを、判別可能なマーカー文字列で確認する。
    _rio_nb_p = {
        "cells": [
            {"cell_type": "code", "source": ["1 + 1"],
             "outputs": [
                 {"output_type": "execute_result", "execution_count": 1,
                  "data": {"text/plain": ["EXECRESULT_MARKER_42"]}, "metadata": {}},
                 {"output_type": "display_data",
                  "data": {"image/png": "QkFTRTY0X01BUktFUg=="}, "metadata": {}},
                 {"output_type": "error", "ename": "ValueError", "evalue": "boom",
                  "traceback": ["ERROR_TRACEBACK_MARKER"]},
             ]},
        ]
    }
    _rio_p = _rio_root / "p.ipynb"
    _rio_p.write_text(json.dumps(_rio_nb_p), encoding="utf-8")
    _rio_out_p = f._read_ipynb(_rio_p)
    check("_read_ipynb(iter195): execute_resultのtext/plainは[Notebook result output]として抽出される",
          _rio_out_p == "```python\n1 + 1\n```"
          "\n\n[Notebook result output]\nEXECRESULT_MARKER_42"
          "\n\n[Notebook error]\nValueError: boom")
    check("_read_ipynb: display_dataがimage/pngのみ(text/plainなし)の場合はbase64が出力に混入しない",
          "QkFTRTY0X01BUktFUg==" not in _rio_out_p and "image/png" not in _rio_out_p)
    check("_read_ipynb(iter196更新): errorはename/evalueが[Notebook error]として抽出されるが、tracebackは対象外のまま混入しない",
          "ValueError: boom" in _rio_out_p and "ERROR_TRACEBACK_MARKER" not in _rio_out_p)

    # (q) outputsが空listの場合: 出力ブロックは追加されずコード抽出のみ(回帰: 既存の
    #     outputsキー自体が無いケースは上のセクション(a)/(e)等で既に検証済み)。
    _rio_nb_q = {"cells": [{"cell_type": "code", "source": ["print(1)"], "outputs": []}]}
    _rio_q = _rio_root / "q.ipynb"
    _rio_q.write_text(json.dumps(_rio_nb_q), encoding="utf-8")
    _rio_out_q = f._read_ipynb(_rio_q)
    check("_read_ipynb: outputs=[]の場合は出力ブロックが追加されずコード抽出のみ",
          _rio_out_q == "```python\nprint(1)\n```")

    # (r) 暴走したprintループ相当の長大なstdoutは上限(4000文字)で切り詰められ、
    #     切り詰めマーカーが付与される(num_ctx保護)。
    _rio_long_text = "A" * 5000
    _rio_nb_r = {
        "cells": [
            {"cell_type": "code", "source": ["print('x' * 5000)"],
             "outputs": [{"output_type": "stream", "name": "stdout", "text": [_rio_long_text]}]},
        ]
    }
    _rio_r = _rio_root / "r.ipynb"
    _rio_r.write_text(json.dumps(_rio_nb_r), encoding="utf-8")
    _rio_out_r = f._read_ipynb(_rio_r)
    _rio_expected_r = (
        "```python\nprint('x' * 5000)\n```"
        "\n\n[Notebook stdout/stderr output]\n" + "A" * 4000 + "\n...(出力が長いため切り詰め)"
    )
    check("_read_ipynb: 長大なstdoutは上限4000文字+切り詰めマーカーで打ち切られる",
          _rio_out_r == _rio_expected_r)
    check("_read_ipynb: 切り詰め後の出力に元の全文(5000文字連続のA)は残らない",
          _rio_long_text not in _rio_out_r)

# ---------- _read_ipynb: execute_result/display_dataの計算結果(text/plain)抽出 (2026-07-26 / iter195) ----------
# 発端: iter194はstream(stdout/stderr)出力のみをRAG/--fileコンテキストへ抽出する
# ように縮小スコープで対応したが、その際の自己コメントでexecute_result/
# display_dataの'data'辞書(text/plainを含む)は明示的に対象外のまま据え置いていた。
# データ分析notebookでは、print(...)によるstream出力ではなく、セル末尾の式評価
# 結果や変数のrepr(戻り値、DataFrameのテキスト表現など)が質問への回答に直結する
# 事実であることが多く、それが精度criticalなRAG/--fileコンテキストから黙って
# 欠落していた。iter188はstream/execute_result/display_data(text/plain)の全種類を
# 一度に扱おうとして行き詰まったため、本iterationはiter194が確立した縮小
# スコープの作法(str->そのまま、list->str要素のみjoin、その他->空、空白のみ
# skip、cap超過時は切り詰めマーカー付与)をそのまま流用し、'data'辞書のうち
# 'text/plain'キー1つだけを追加対象とする。text/html・image/png等の非'text/plain'
# MIME、およびerrorのtracebackは引き続き意図的に対象外(オーバーサイトではない)。
# 上のセクション(iter71/72/113/159構造ガード、iter194 stream抽出)は本セクションの
# 変更で一切変更されていない。f.ask/urlopen/subprocessは一切呼ばない(すべて
# tempfile上のオフラインI/O)。
with _tempfile.TemporaryDirectory() as _rid_dir:
    _rid_root = _pathlib.Path(_rid_dir)

    # (s) execute_resultのtext/plain(list[str])が[Notebook result output]ラベル付きで
    #     フェンス直後に追加される。
    _rid_nb_s = {
        "cells": [
            {"cell_type": "code", "source": ["1 + 1"],
             "outputs": [
                 {"output_type": "execute_result", "execution_count": 1,
                  "data": {"text/plain": ["2"]}, "metadata": {}},
             ]},
        ]
    }
    _rid_s = _rid_root / "s.ipynb"
    _rid_s.write_text(json.dumps(_rid_nb_s), encoding="utf-8")
    _rid_out_s = f._read_ipynb(_rid_s)
    check("_read_ipynb: execute_resultのtext/plain(list[str])が[Notebook result output]で抽出される",
          _rid_out_s == "```python\n1 + 1\n```\n\n[Notebook result output]\n2")

    # (t) execute_resultのtext/plainが単一strでも同様に抽出される。
    _rid_nb_t = {
        "cells": [
            {"cell_type": "code", "source": ["1 + 1"],
             "outputs": [
                 {"output_type": "execute_result", "data": {"text/plain": "2"}},
             ]},
        ]
    }
    _rid_t = _rid_root / "t.ipynb"
    _rid_t.write_text(json.dumps(_rid_nb_t), encoding="utf-8")
    _rid_out_t = f._read_ipynb(_rid_t)
    check("_read_ipynb: execute_resultのtext/plain(単一str)も同様に抽出される",
          _rid_out_t == "```python\n1 + 1\n```\n\n[Notebook result output]\n2")

    # (u) display_dataのtext/plain(DataFrameのテキスト表現相当)も同じラベルで抽出される。
    _rid_nb_u = {
        "cells": [
            {"cell_type": "code", "source": ["df"],
             "outputs": [
                 {"output_type": "display_data",
                  "data": {"text/plain": ["   a  b\n0  1  2"]}, "metadata": {}},
             ]},
        ]
    }
    _rid_u = _rid_root / "u.ipynb"
    _rid_u.write_text(json.dumps(_rid_nb_u), encoding="utf-8")
    _rid_out_u = f._read_ipynb(_rid_u)
    check("_read_ipynb: display_dataのtext/plainも[Notebook result output]で抽出される",
          _rid_out_u == "```python\ndf\n```\n\n[Notebook result output]\n   a  b\n0  1  2")

    # (v) execute_result/display_dataがimage/pngのみ(text/plainなし)の場合は結果ブロックが
    #     追加されず、base64データも一切出力に混入しない。
    _rid_nb_v = {
        "cells": [
            {"cell_type": "code", "source": ["plot(df)"],
             "outputs": [
                 {"output_type": "execute_result",
                  "data": {"image/png": "QkFTRTY0X0VYRUNfSU1H"}, "metadata": {}},
                 {"output_type": "display_data",
                  "data": {"image/png": "QkFTRTY0X0RJU1BfSU1H"}, "metadata": {}},
             ]},
        ]
    }
    _rid_v = _rid_root / "v.ipynb"
    _rid_v.write_text(json.dumps(_rid_nb_v), encoding="utf-8")
    _rid_out_v = f._read_ipynb(_rid_v)
    check("_read_ipynb: execute_result/display_dataがimage/pngのみの場合は結果ブロックが追加されない",
          _rid_out_v == "```python\nplot(df)\n```")
    check("_read_ipynb: image/pngのみのケースでbase64文字列が出力に混入しない",
          "QkFTRTY0X0VYRUNfSU1H" not in _rid_out_v
          and "QkFTRTY0X0RJU1BfSU1H" not in _rid_out_v
          and "image/png" not in _rid_out_v)

    # (w) text/plainとimage/pngが同一dataに同居する場合、text/plainのみが抽出され
    #     base64は一切出力に混入しない。
    _rid_nb_w = {
        "cells": [
            {"cell_type": "code", "source": ["plot(df)"],
             "outputs": [
                 {"output_type": "display_data",
                  "data": {"text/plain": ["<Figure size 640x480>"],
                            "image/png": "QkFTRTY0X01JWEVEX0lNRw=="},
                  "metadata": {}},
             ]},
        ]
    }
    _rid_w = _rid_root / "w.ipynb"
    _rid_w.write_text(json.dumps(_rid_nb_w), encoding="utf-8")
    _rid_out_w = f._read_ipynb(_rid_w)
    check("_read_ipynb: text/plain+image/png同居時はtext/plainのみ抽出される",
          _rid_out_w == "```python\nplot(df)\n```\n\n[Notebook result output]\n<Figure size 640x480>")
    check("_read_ipynb: text/plain+image/png同居時もbase64は出力に混入しない",
          "QkFTRTY0X01JWEVEX0lNRw==" not in _rid_out_w and "image/png" not in _rid_out_w)

    # (x) 'data'キー自体が無い場合: 例外を送出せず結果ブロックも追加されない。
    _rid_nb_x = {
        "cells": [
            {"cell_type": "code", "source": ["1 + 1"],
             "outputs": [{"output_type": "execute_result", "execution_count": 1}]},
        ]
    }
    _rid_x = _rid_root / "x.ipynb"
    _rid_x.write_text(json.dumps(_rid_nb_x), encoding="utf-8")
    _rid_out_x = f._read_ipynb(_rid_x)
    check("_read_ipynb: 'data'キーが無いexecute_resultは例外を送出せず結果ブロックも追加されない",
          _rid_out_x == "```python\n1 + 1\n```")

    # (y) 'data'がtruthyな非dict(整数/文字列/None)の場合も、iteration 113のcells非list
    #     強制変換と同じ作法(`or []`のtruthinessトリックに頼らない)でそのoutputをskipし、
    #     例外を送出しない。
    for _rid_bad_data, _rid_label in (
        (42, "整数"), ("notadict", "文字列"), (None, "None"),
    ):
        _rid_nb_y = {
            "cells": [
                {"cell_type": "code", "source": ["1 + 1"],
                 "outputs": [{"output_type": "execute_result", "data": _rid_bad_data}]},
            ]
        }
        _rid_y = _rid_root / f"y_{_rid_label}.ipynb"
        _rid_y.write_text(json.dumps(_rid_nb_y), encoding="utf-8")
        _rid_out_y = f._read_ipynb(_rid_y)
        check(f"_read_ipynb: 'data'が非dict({_rid_label})の場合は例外を送出せずそのoutputはskipされる",
              _rid_out_y == "```python\n1 + 1\n```")

    # (z) 'text/plain'の値が非str/非list(整数/None)の場合、空文字へ正規化されblank後
    #     strip判定でそのoutputはskipされる(例外は送出しない)。
    for _rid_bad_tp, _rid_label2 in ((123, "整数"), (None, "None")):
        _rid_nb_z = {
            "cells": [
                {"cell_type": "code", "source": ["1 + 1"],
                 "outputs": [{"output_type": "execute_result",
                              "data": {"text/plain": _rid_bad_tp}}]},
            ]
        }
        _rid_z = _rid_root / f"z_{_rid_label2}.ipynb"
        _rid_z.write_text(json.dumps(_rid_nb_z), encoding="utf-8")
        _rid_out_z = f._read_ipynb(_rid_z)
        check(f"_read_ipynb: text/plainが非str/非list({_rid_label2})でも例外を送出せずskipされる",
              _rid_out_z == "```python\n1 + 1\n```")

    # (z2) text/plainが空白のみ(str/list双方)の場合、strip後blankとしてskipされる。
    _rid_nb_z2 = {
        "cells": [
            {"cell_type": "code", "source": ["1 + 1"],
             "outputs": [
                 {"output_type": "execute_result", "data": {"text/plain": "   \n  "}},
                 {"output_type": "display_data", "data": {"text/plain": ["   ", "\n"]}},
             ]},
        ]
    }
    _rid_z2 = _rid_root / "z2.ipynb"
    _rid_z2.write_text(json.dumps(_rid_nb_z2), encoding="utf-8")
    _rid_out_z2 = f._read_ipynb(_rid_z2)
    check("_read_ipynb: text/plainが空白のみ(str/list双方)の場合は結果ブロックが追加されない",
          _rid_out_z2 == "```python\n1 + 1\n```")

    # (aa) text/plainのlist内に非str要素が混入する場合、source正規化(iteration 72)と
    #      全く同じパターンでstr要素のみjoinされる(42は無視されて'a'+'b'='ab')。
    _rid_nb_aa = {
        "cells": [
            {"cell_type": "code", "source": ["1 + 1"],
             "outputs": [
                 {"output_type": "execute_result", "data": {"text/plain": ["a", 42, "b"]}},
             ]},
        ]
    }
    _rid_aa = _rid_root / "aa.ipynb"
    _rid_aa.write_text(json.dumps(_rid_nb_aa), encoding="utf-8")
    _rid_out_aa = f._read_ipynb(_rid_aa)
    check("_read_ipynb: text/plainのlist内の非str要素は無視されstr要素のみjoinされる",
          _rid_out_aa == "```python\n1 + 1\n```\n\n[Notebook result output]\nab")

    # (bb) 複数のexecute_result/display_dataが存在する場合、outputsのlist順を保った
    #      決定的な順序で結合される(streamのcombined_outと同じ結合パターン)。
    _rid_nb_bb = {
        "cells": [
            {"cell_type": "code", "source": ["1 + 1"],
             "outputs": [
                 {"output_type": "execute_result", "data": {"text/plain": ["first "]}},
                 {"output_type": "display_data", "data": {"text/plain": ["second"]}},
             ]},
        ]
    }
    _rid_bb = _rid_root / "bb.ipynb"
    _rid_bb.write_text(json.dumps(_rid_nb_bb), encoding="utf-8")
    _rid_out_bb = f._read_ipynb(_rid_bb)
    check("_read_ipynb: 複数のexecute_result/display_dataはoutputsのlist順で決定的に結合される",
          _rid_out_bb == "```python\n1 + 1\n```\n\n[Notebook result output]\nfirst second")

    # (cc) stream出力とexecute_result出力が同一セルに同居する場合、両方のブロックが
    #      それぞれ[Notebook stdout/stderr output]/[Notebook result output]の別ラベルで
    #      インターリーブせず(混ざらず)フェンス直後に順番通り追加される。
    _rid_nb_cc = {
        "cells": [
            {"cell_type": "code", "source": ["print('out1'); 2 + 2"],
             "outputs": [
                 {"output_type": "stream", "name": "stdout", "text": ["out1\n"]},
                 {"output_type": "execute_result", "data": {"text/plain": ["4"]}},
             ]},
        ]
    }
    _rid_cc = _rid_root / "cc.ipynb"
    _rid_cc.write_text(json.dumps(_rid_nb_cc), encoding="utf-8")
    _rid_out_cc = f._read_ipynb(_rid_cc)
    check("_read_ipynb: streamとexecute_resultが同居する場合、両ブロックが別ラベルで非インターリーブに追加される",
          _rid_out_cc == "```python\nprint('out1'); 2 + 2\n```"
          "\n\n[Notebook stdout/stderr output]\nout1\n"
          "\n\n[Notebook result output]\n4")

    # (dd) 巨大なtext/plain(暴走したreprを想定)は、streamと同じ上限(4000文字)+
    #      切り詰めマーカーで打ち切られる(num_ctx保護、_IPYNB_STREAM_OUTPUT_CAP流用)。
    _rid_long_text = "B" * 5000
    _rid_nb_dd = {
        "cells": [
            {"cell_type": "code", "source": ["huge_repr"],
             "outputs": [
                 {"output_type": "execute_result", "data": {"text/plain": [_rid_long_text]}},
             ]},
        ]
    }
    _rid_dd = _rid_root / "dd.ipynb"
    _rid_dd.write_text(json.dumps(_rid_nb_dd), encoding="utf-8")
    _rid_out_dd = f._read_ipynb(_rid_dd)
    _rid_expected_dd = (
        "```python\nhuge_repr\n```"
        "\n\n[Notebook result output]\n" + "B" * 4000 + "\n...(出力が長いため切り詰め)"
    )
    check("_read_ipynb: 巨大なtext/plainは上限4000文字+切り詰めマーカーで打ち切られる",
          _rid_out_dd == _rid_expected_dd)
    check("_read_ipynb: 切り詰め後の出力に元の全文(5000文字連続のB)は残らない",
          _rid_long_text not in _rid_out_dd)

    # (ee) iter194/195時点ではerrorのoutput_typeは全面的に対象外だったが、iter196
    #      (本iteration)がename/evalueのみを[Notebook error]として追加抽出するように
    #      なった。execute_result/display_data向けの[Notebook result output]ブロックは
    #      (output_typeがerrorのため)引き続き追加されず、tracebackも引き続き出力に
    #      混入しない(下の専用セクションでより網羅的に検証)。
    _rid_nb_ee = {
        "cells": [
            {"cell_type": "code", "source": ["1 / 0"],
             "outputs": [
                 {"output_type": "error", "ename": "ZeroDivisionError",
                  "evalue": "division by zero", "traceback": ["ERROR_TRACEBACK_MARKER_EE"]},
             ]},
        ]
    }
    _rid_ee = _rid_root / "ee.ipynb"
    _rid_ee.write_text(json.dumps(_rid_nb_ee), encoding="utf-8")
    _rid_out_ee = f._read_ipynb(_rid_ee)
    check("_read_ipynb(iter196更新): errorは[Notebook result output]としては追加されず[Notebook error]として抽出される",
          _rid_out_ee == "```python\n1 / 0\n```"
          "\n\n[Notebook error]\nZeroDivisionError: division by zero")
    check("_read_ipynb: errorのtracebackは出力に混入しない",
          "ERROR_TRACEBACK_MARKER_EE" not in _rid_out_ee)

# ---------- _read_ipynb: コードセルのerror出力(ename/evalue)抽出 (2026-07-26 / iter196) ----------
# 発端: iter194はstream(stdout/stderr)出力を、iter195はexecute_result/display_dataの
# 'text/plain'をそれぞれ追加抽出したが、output_type=='error'は「iter188がstream/
# execute_result/display_data/errorの全種類を一度に扱おうとして行き詰まり断念した」
# ことを受けて、iter194/195とも意図的に対象外のまま据え置いてきた(iter195のコメント:
# 「'error'のtracebackも同様に対象外のまま据え置く」)。この据え置きは単なる網羅性の
# 欠落ではなく精度上のハザードであり、保存済み出力がerrorのコードセルでも従来は
# ```pythonフェンスのみが出力され、セルが失敗した事実が一切示されないまま提案者に
# 渡っていた(提案者がコードは成功したものと誤って推論しうる)。本iterationは
# iter194/195の縮小スコープの作法をそのまま踏襲し、'error'の'ename'/'evalue'という
# 単一の最小形状のみを追加抽出する。'traceback'(list of str、ANSIカラーエスケープ
# シーケンスを含む)は意図的に対象外のまま据え置く(iter195からのdeferralの継続。
# ANSIエスケープの除去・正規化には別途の慎重な設計が必要であり、iter188の轍を
# 踏まないため)。上のセクション(iter71/72/113/159構造ガード、iter194 stream抽出、
# iter195 result抽出)は本セクションの変更で一切変更されていない。f.ask/urlopen/
# subprocessは一切呼ばない(すべてtempfile上のオフラインI/O)。
with _tempfile.TemporaryDirectory() as _rie_dir:
    _rie_root = _pathlib.Path(_rie_dir)

    # (a) 基本ケース: ename='ZeroDivisionError'/evalue='division by zero'が
    #     [Notebook error]ブロックとしてフェンス直後に抽出される。
    _rie_nb_a = {
        "cells": [
            {"cell_type": "code", "source": ["1 / 0"],
             "outputs": [
                 {"output_type": "error", "ename": "ZeroDivisionError",
                  "evalue": "division by zero", "traceback": ["tb line 1", "tb line 2"]},
             ]},
        ]
    }
    _rie_a = _rie_root / "a.ipynb"
    _rie_a.write_text(json.dumps(_rie_nb_a), encoding="utf-8")
    _rie_out_a = f._read_ipynb(_rie_a)
    check("_read_ipynb: error出力のename/evalueが[Notebook error]としてフェンス直後に抽出される",
          _rie_out_a == "```python\n1 / 0\n```"
          "\n\n[Notebook error]\nZeroDivisionError: division by zero")

    # (b) tracebackにANSIカラーエスケープ(\x1b[...m)が含まれていても、ename/evalueは
    #     正しく抽出されつつ、ANSIエスケープ文字自体は出力のどこにも混入しない
    #     (traceback自体を一切読まないことの直接的な証明)。
    _rie_nb_b = {
        "cells": [
            {"cell_type": "code", "source": ["1 / 0"],
             "outputs": [
                 {"output_type": "error", "ename": "ZeroDivisionError",
                  "evalue": "division by zero",
                  "traceback": [
                      "[0;31m----- Traceback -----[0m",
                      "[0;31mZeroDivisionError[0m: division by zero",
                  ]},
             ]},
        ]
    }
    _rie_b = _rie_root / "b.ipynb"
    _rie_b.write_text(json.dumps(_rie_nb_b), encoding="utf-8")
    _rie_out_b = f._read_ipynb(_rie_b)
    check("_read_ipynb: tracebackのANSIエスケープ文字(\\x1b)は出力のどこにも混入しない",
          "" not in _rie_out_b)
    check("_read_ipynb: ANSI混入tracebackが同居してもename/evalueは正しく抽出される",
          "ZeroDivisionError: division by zero" in _rie_out_b)

    # (c) evalueが欠落('evalue'キー自体が無い)場合、ename単体が抽出され末尾に
    #     余分な': 'は付かない。
    _rie_nb_c = {
        "cells": [
            {"cell_type": "code", "source": ["raise RuntimeError"],
             "outputs": [{"output_type": "error", "ename": "RuntimeError"}]},
        ]
    }
    _rie_c = _rie_root / "c.ipynb"
    _rie_c.write_text(json.dumps(_rie_nb_c), encoding="utf-8")
    _rie_out_c = f._read_ipynb(_rie_c)
    check("_read_ipynb: evalueキー自体が無い場合もenameのみ抽出され末尾に': 'が付かない",
          _rie_out_c == "```python\nraise RuntimeError\n```\n\n[Notebook error]\nRuntimeError")

    # (d) evalueが空白のみの場合も(c)と同様、ename単体が抽出され余分な': 'は付かない。
    _rie_nb_d = {
        "cells": [
            {"cell_type": "code", "source": ["raise RuntimeError"],
             "outputs": [{"output_type": "error", "ename": "RuntimeError", "evalue": "   "}]},
        ]
    }
    _rie_d = _rie_root / "d.ipynb"
    _rie_d.write_text(json.dumps(_rie_nb_d), encoding="utf-8")
    _rie_out_d = f._read_ipynb(_rie_d)
    check("_read_ipynb: evalueが空白のみの場合もenameのみ抽出され末尾に': 'が付かない",
          _rie_out_d == "```python\nraise RuntimeError\n```\n\n[Notebook error]\nRuntimeError")

    # (e) enameが欠落/空白でevalueのみ存在する場合は対称的にevalue単体が抽出される
    #     (先頭に余分な': 'が付かない)。
    _rie_nb_e = {
        "cells": [
            {"cell_type": "code", "source": ["assert False"],
             "outputs": [{"output_type": "error", "evalue": "assertion failed"}]},
        ]
    }
    _rie_e = _rie_root / "e.ipynb"
    _rie_e.write_text(json.dumps(_rie_nb_e), encoding="utf-8")
    _rie_out_e = f._read_ipynb(_rie_e)
    check("_read_ipynb: enameが無くevalueのみの場合はevalue単体が抽出される",
          _rie_out_e == "```python\nassert False\n```\n\n[Notebook error]\nassertion failed")

    # (f) enameとevalueの両方が欠落/空白の場合はerrorブロック自体が追加されない
    #     (キー無し・明示的な空文字・空白のみの3パターン)。
    for _rie_bad_out, _rie_label_f in (
        ({"output_type": "error"}, "両方キー無し"),
        ({"output_type": "error", "ename": "", "evalue": ""}, "両方空文字"),
        ({"output_type": "error", "ename": "  ", "evalue": "\n"}, "両方空白のみ"),
    ):
        _rie_nb_f = {"cells": [{"cell_type": "code", "source": ["pass"], "outputs": [_rie_bad_out]}]}
        _rie_f = _rie_root / f"f_{_rie_label_f}.ipynb"
        _rie_f.write_text(json.dumps(_rie_nb_f), encoding="utf-8")
        _rie_out_f = f._read_ipynb(_rie_f)
        check(f"_read_ipynb: ename/evalueが{_rie_label_f}の場合はerrorブロックが追加されない",
              _rie_out_f == "```python\npass\n```")

    # (g) ename/evalueが非str/非list(int/dict/None)の場合、例外を送出せずstr()での
    #     強制変換(dict/int等のrepr混入)もせず空文字へ正規化される。evalueは固定の
    #     正常値にして、enameの異常値だけを切り替える。
    for _rie_bad_ename, _rie_label_g in ((42, "整数"), ({"x": 1}, "dict"), (None, "None")):
        _rie_nb_g = {
            "cells": [
                {"cell_type": "code", "source": ["boom()"],
                 "outputs": [{"output_type": "error", "ename": _rie_bad_ename,
                              "evalue": "GDEF_MARKER"}]},
            ]
        }
        _rie_g = _rie_root / f"g_{_rie_label_g}.ipynb"
        _rie_g.write_text(json.dumps(_rie_nb_g), encoding="utf-8")
        _rie_out_g = f._read_ipynb(_rie_g)
        check(f"_read_ipynb: enameが非str/非list({_rie_label_g})でも例外を送出せず空文字へ正規化される",
              _rie_out_g == "```python\nboom()\n```\n\n[Notebook error]\nGDEF_MARKER")

    # (h) enameがlist内に非str要素を含む場合、source正規化(iteration 72)と全く同じ
    #     パターンでstr要素のみjoinされる(42は無視されて'a'+'b'='ab')。
    _rie_nb_h = {
        "cells": [
            {"cell_type": "code", "source": ["boom()"],
             "outputs": [{"output_type": "error", "ename": ["a", 42, "b"], "evalue": "m"}]},
        ]
    }
    _rie_h = _rie_root / "h.ipynb"
    _rie_h.write_text(json.dumps(_rie_nb_h), encoding="utf-8")
    _rie_out_h = f._read_ipynb(_rie_h)
    check("_read_ipynb: enameのlist内の非str要素は無視されstr要素のみjoinされる",
          _rie_out_h == "```python\nboom()\n```\n\n[Notebook error]\nab: m")

    # (i) outputs内に非dictエントリが混在しても例外を送出せずskipされ、後続の正当な
    #     error出力は引き続き抽出される(iteration 72の非dictセルskipと同じ作法)。
    _rie_nb_i = {
        "cells": [
            {"cell_type": "code", "source": ["boom()"],
             "outputs": ["not_a_dict_output",
                         {"output_type": "error", "ename": "KeyError", "evalue": "boom2"}]},
        ]
    }
    _rie_i = _rie_root / "i.ipynb"
    _rie_i.write_text(json.dumps(_rie_nb_i), encoding="utf-8")
    _rie_out_i = f._read_ipynb(_rie_i)
    check("_read_ipynb: outputs内の非dictエントリはskipされ後続の正当なerror出力は抽出される",
          _rie_out_i == "```python\nboom()\n```\n\n[Notebook error]\nKeyError: boom2")

    # (j) 'outputs'自体がtruthyな非list(整数)の場合も、iteration 113/72の構造ガードを
    #     再利用しているため例外を送出せずコード抽出のみが残る。
    _rie_nb_j = {"cells": [{"cell_type": "code", "source": ["boom()"], "outputs": 42}]}
    _rie_j = _rie_root / "j.ipynb"
    _rie_j.write_text(json.dumps(_rie_nb_j), encoding="utf-8")
    _rie_out_j = f._read_ipynb(_rie_j)
    check("_read_ipynb: outputsが非list(整数)でも例外を送出せずコード抽出のみ維持される",
          _rie_out_j == "```python\nboom()\n```")

    # (k) stream + execute_result(text/plain) + errorが同一セルに同居する場合、
    #     3ブロックが[Notebook stdout/stderr output] -> [Notebook result output] ->
    #     [Notebook error]の順で、インターリーブせず・欠落/重複なく追加される。
    _rie_nb_k = {
        "cells": [
            {"cell_type": "code", "source": ["print('before'); 1/0"],
             "outputs": [
                 {"output_type": "stream", "name": "stdout", "text": ["before\n"]},
                 {"output_type": "execute_result", "data": {"text/plain": ["<partial>"]}},
                 {"output_type": "error", "ename": "ZeroDivisionError", "evalue": "division by zero"},
             ]},
        ]
    }
    _rie_k = _rie_root / "k.ipynb"
    _rie_k.write_text(json.dumps(_rie_nb_k), encoding="utf-8")
    _rie_out_k = f._read_ipynb(_rie_k)
    check("_read_ipynb: stream+result+errorが同居する場合、3ブロックが決定的な順序で1回ずつ追加される",
          _rie_out_k == "```python\nprint('before'); 1/0\n```"
          "\n\n[Notebook stdout/stderr output]\nbefore\n"
          "\n\n[Notebook result output]\n<partial>"
          "\n\n[Notebook error]\nZeroDivisionError: division by zero")

    # (l) 巨大なevalue(暴走したtraceback相当)は、stream/result出力と同じ上限
    #     (4000文字、_IPYNB_STREAM_OUTPUT_CAP流用)+切り詰めマーカーで打ち切られる
    #     (num_ctx保護)。enameは空にして、切り詰め位置の計算をevalue単体に単純化する。
    _rie_long_evalue = "E" * 5000
    _rie_nb_l = {
        "cells": [
            {"cell_type": "code", "source": ["boom()"],
             "outputs": [{"output_type": "error", "ename": "", "evalue": _rie_long_evalue}]},
        ]
    }
    _rie_l = _rie_root / "l.ipynb"
    _rie_l.write_text(json.dumps(_rie_nb_l), encoding="utf-8")
    _rie_out_l = f._read_ipynb(_rie_l)
    _rie_expected_l = (
        "```python\nboom()\n```"
        "\n\n[Notebook error]\n" + "E" * 4000 + "\n...(出力が長いため切り詰め)"
    )
    check("_read_ipynb: 巨大なevalueは上限4000文字+切り詰めマーカーで打ち切られる",
          _rie_out_l == _rie_expected_l)
    check("_read_ipynb: 切り詰め後の出力に元の全文(5000文字連続のE)は残らない",
          _rie_long_evalue not in _rie_out_l)

# ---------- _read_ipynb: cp932(Shift-JIS)デコードラダー (2026-07-25 / iter159) ----------
# 発端: _read_ipynbは、iteration 94が_read_html/read_file_text汎用テキスト分岐に
# 導入した_decode_text_bytes()のutf-8→cp932→replaceラダー(上のセクション、
# iteration 70のerrors="replace"はfugu自身がUTF-8で書いたファイルの読み戻し用の
# 保険であり他所由来ファイルの元エンコーディングは救済しないという整理も同じ)を
# 使わない最後のサードパーティファイルリーダーだった。notebookのJSON構造
# (波括弧・引用符・キー名)はASCIIのみで構成されるため、cp932(Shift-JIS)保存の
# .ipynb(既知の落とし穴#4のcp932コンソールと同根の環境要因でこのマシンには
# 普通に存在する)でもjson.loads自体は成功し、iteration 72/113(上のセクション)の
# セル単位/トップレベル構造ガードも問題なく通過してしまうが、従来の
# path.read_text(encoding="utf-8", errors="replace")がデコード時点で各セルの
# 'source'内の日本語をU+FFFDへ全て潰していたため、精度criticalなRAG/--file
# コンテキストへ文字化けがそのまま注入されていた。ここでは(1) cp932保存
# notebookのコード/markdown両セルで日本語sourceが正しく復元されること、
# (2) 通常のUTF-8 notebook(単一文字列source含む)は変更前とバイト単位で完全
# 一致すること(回帰)、(3) 真に非JSONなnotebookのexceptフォールバックも
# 再読み込みではなく同一の_decode_text_bytes結果を使うため、そのフォールバック
# テキスト自体もcp932復元されること、の3点を検証する。iteration 71/72/113の
# 既存フィクスチャ(直上のセクション、source=None・非str要素混入・非dictセル・
# 非dictトップレベル・非listなcells)は本変更後も全てutf-8保存のまま変更前と
# 同じ結果を返すことを直上のテストで既に回帰確認済み。Ollama/ネットワーク
# 呼び出しは一切不要(すべてtempfile上のオフラインI/O)。
with _tempfile.TemporaryDirectory() as _rin_dir:
    _rin_root = _pathlib.Path(_rin_dir)

    # (1) cp932保存のnotebook: コード/markdown両セルの日本語sourceが正しく復元される。
    #     ensure_ascii=Falseでdumpすることで、JSON文字列中に実際の日本語文字を
    #     含めた上でcp932へエンコードする(ensure_ascii既定のままだと\uXXXXエスケープ
    #     のみのASCII文字列になり、cp932マルチバイト列を全く含まないためテストとして
    #     無意味になる)。
    _rin_nb_jp = {
        "cells": [
            {"cell_type": "code",
             "source": ["import os\n", "print('こんにちは')  # 日本語コメント"]},
            {"cell_type": "markdown",
             "source": ["# 見出し\n", "これは日本語の説明文です。株式会社。"]},
        ]
    }
    _rin_jp_json = json.dumps(_rin_nb_jp, ensure_ascii=False)
    _rin_jp_path = _rin_root / "sjis.ipynb"
    _rin_jp_path.write_bytes(_rin_jp_json.encode("cp932"))
    _rin_jp_out = f._read_ipynb(_rin_jp_path)
    _rin_jp_expected = (
        "```python\nimport os\nprint('こんにちは')  # 日本語コメント\n```"
        "\n\n# 見出し\nこれは日本語の説明文です。株式会社。"
    )
    check("_read_ipynb: cp932保存notebookのコード/markdownセルの日本語sourceが正しく復元される",
          _rin_jp_out == _rin_jp_expected)
    check("_read_ipynb: cp932保存notebookの復元結果にU+FFFD(文字化け)が含まれない",
          "�" not in _rin_jp_out)

    # (2) 回帰: 通常のUTF-8 notebook(複数セル・単一文字列source含む)は
    #     変更前とバイト単位で完全一致する。
    _rin_nb_utf8 = {
        "cells": [
            {"cell_type": "code", "source": ["import sys\n", "print(sys.version)"]},
            {"cell_type": "markdown", "source": "Plain single-string markdown source."},
            {"cell_type": "code", "source": "x = 1\ny = 2\nprint(x + y)"},
        ]
    }
    _rin_utf8_path = _rin_root / "utf8.ipynb"
    _rin_utf8_path.write_text(json.dumps(_rin_nb_utf8), encoding="utf-8")
    _rin_utf8_out = f._read_ipynb(_rin_utf8_path)
    _rin_utf8_expected = (
        "```python\nimport sys\nprint(sys.version)\n```"
        "\n\nPlain single-string markdown source."
        "\n\n```python\nx = 1\ny = 2\nprint(x + y)\n```"
    )
    check("_read_ipynb: 通常のUTF-8 notebook(単一文字列source含む)は従来通りバイト単位で完全一致(回帰)",
          _rin_utf8_out == _rin_utf8_expected)

    # (3) 真に非JSONなnotebook(cp932保存の平文): exceptフォールバックが
    #     path.read_text()での再読み込みではなく、try節冒頭で一度だけ
    #     _decode_text_bytes()した結果を再利用するため、フォールバックテキスト
    #     自体もcp932復元される(旧実装なら再読み込み時もutf-8+replaceのままで
    #     日本語がU+FFFDに化けていた)。
    _rin_bad_jp_text = "これは有効なJSONではないメモです。{ 壊れています"
    _rin_bad_path = _rin_root / "sjis_broken.ipynb"
    _rin_bad_path.write_bytes(_rin_bad_jp_text.encode("cp932"))
    _rin_bad_out = f._read_ipynb(_rin_bad_path)
    check("_read_ipynb: 非JSONかつcp932保存のnotebookはフォールバックでも日本語が正しく復元される",
          _rin_bad_out == _rin_bad_jp_text)
    check("_read_ipynb: 非JSON cp932フォールバック結果にU+FFFD(文字化け)が含まれない",
          "�" not in _rin_bad_out)

# ---------- _decode_text_bytes / read_file_text・_read_html cp932フォールバック (2026-07-24) ----------
# このマシンのコンソールが cp932(Shift-JIS) である既知の落とし穴 #4 と同根の環境要因で、
# ローカルに保存された非UTF-8(Shift-JIS/cp932)の .txt/.csv/.html 等が普通に存在する。
# 従来 read_file_text の汎用テキスト分岐と _read_html は encoding="utf-8",
# errors="replace" を無条件に使っており、cp932ファイルを読むと日本語部分が「全て」
# U+FFFDに化けて精度critical な RAG/--file コンテキストへそのまま注入されていた。
# iter47(_save_as_markdown/_save_as_text/_save_as_htmlの読み戻し)・iter70(会話履歴JSON
# 読み込み)のerrors="replace"適用はfugu自身がUTF-8で書いたファイルの読み戻し用の保険で
# あり、他所由来ファイルの元エンコーディングは救済しない。本iterationはutf-8 厳密 ->
# cp932 厳密 -> utf-8 errors="replace" のデコードラダー(_decode_text_bytes)を追加し、
# 他所由来ファイルの元エンコーディングを損失なく復元する。ここでは
# _decode_text_bytes単体・read_file_text経由・_read_html経由の3レイヤーで検証する。
# Ollama/ネットワーク呼び出しは一切不要（すべてtempfile上のオフラインI/O）。

# (1) _decode_text_bytes単体: utf-8成功 / utf-8失敗+cp932成功 / 両方失敗の3分岐
_dtb_ascii_bytes = "plain ascii text, no surprises.".encode("utf-8")
check("_decode_text_bytes: 純粋ASCIIはutf-8厳密デコードでそのまま返る",
      f._decode_text_bytes(_dtb_ascii_bytes) == "plain ascii text, no surprises.")

_dtb_jp_text = "日本語のテスト文章です。utf-8で保存。"
check("_decode_text_bytes: utf-8厳密デコード成功時は日本語もそのまま復元される",
      f._decode_text_bytes(_dtb_jp_text.encode("utf-8")) == _dtb_jp_text)

_dtb_cp932_text = "これはcp932(Shift-JIS)で保存された日本語テキストです。株式会社。"
_dtb_cp932_bytes = _dtb_cp932_text.encode("cp932")
_dtb_cp932_decoded = f._decode_text_bytes(_dtb_cp932_bytes)
check("_decode_text_bytes: utf-8失敗・cp932厳密デコード成功時はcp932デコード結果を返す",
      _dtb_cp932_decoded == _dtb_cp932_text)
check("_decode_text_bytes: cp932復元結果にU+FFFD(文字化け)が含まれない",
      "�" not in _dtb_cp932_decoded)

_dtb_undecodable = b"\x80\x81\xfe\xff"  # utf-8としてもcp932としても不正なバイト列
_dtb_expected_fallback = _dtb_undecodable.decode("utf-8", errors="replace")
check("_decode_text_bytes: utf-8/cp932とも失敗した場合は例外を送出せずutf-8+replaceにフォールバック",
      f._decode_text_bytes(_dtb_undecodable) == _dtb_expected_fallback)
check("_decode_text_bytes: 完全劣化フォールバック結果にはU+FFFDが含まれる(想定通りの置換・退行なし)",
      "�" in f._decode_text_bytes(_dtb_undecodable))

# (2) read_file_text経由: cp932保存の.txt/.csvが正しい日本語として復元される(従来は全滅)
with _tempfile.TemporaryDirectory() as _dtb_dir:
    _dtb_root = _pathlib.Path(_dtb_dir)

    _dtb_txt_jp = "議事録: 本日の会議では新製品の売上について議論した。担当者は田中さんです。"
    (_dtb_root / "sjis.txt").write_bytes(_dtb_txt_jp.encode("cp932"))
    _dtb_txt_out = f.read_file_text(_dtb_root / "sjis.txt")
    check("read_file_text: cp932保存の.txtが正しい日本語として復元される(U+FFFD化けなし)",
          _dtb_txt_out == _dtb_txt_jp and "�" not in _dtb_txt_out)

    _dtb_csv_jp = "名前,部署,売上\n田中太郎,営業部,1200000\n鈴木花子,開発部,980000\n"
    (_dtb_root / "sjis.csv").write_bytes(_dtb_csv_jp.encode("cp932"))
    _dtb_csv_out = f.read_file_text(_dtb_root / "sjis.csv")
    check("read_file_text: cp932保存の.csvが正しい日本語として復元される(U+FFFD化けなし)",
          _dtb_csv_out == _dtb_csv_jp and "�" not in _dtb_csv_out)

    # 回帰: 通常のUTF-8 .txtは従来(utf-8厳密読み)とバイト単位で完全一致
    _dtb_utf8_txt = "通常のUTF-8保存テキストです。\n2行目。\n"
    (_dtb_root / "utf8.txt").write_text(_dtb_utf8_txt, encoding="utf-8")
    check("read_file_text: 通常のUTF-8 .txtは従来の読み込み結果とバイト単位で完全一致(回帰)",
          f.read_file_text(_dtb_root / "utf8.txt") == _dtb_utf8_txt)

    # 回帰: 純粋ASCIIファイルは無変化
    _dtb_ascii_txt = "pure ascii content, nothing exotic.\nsecond line.\n"
    (_dtb_root / "ascii.txt").write_text(_dtb_ascii_txt, encoding="utf-8")
    check("read_file_text: 純粋ASCIIファイルは無変化(回帰)",
          f.read_file_text(_dtb_root / "ascii.txt") == _dtb_ascii_txt)

    # 劣化耐性: read_bytes自体が失敗する(存在しないファイル)場合も従来通り""を返し
    # 例外を送出しない(read_file_textのgraceful-degradation契約は不変)
    check("read_file_text: 存在しないファイルはread_bytes失敗を握りつぶし\"\"を返す(劣化耐性)",
          f.read_file_text(_dtb_root / "does_not_exist.txt") == "")

# (3) _read_html経由: cp932保存の.htmlがタグ除去済みの正しい日本語として復元される
with _tempfile.TemporaryDirectory() as _dtb_html_dir:
    _dtb_html_root = _pathlib.Path(_dtb_html_dir)

    _dtb_html_src = (
        "<html><body><h1>お知らせ</h1>"
        "<p>本日は晴天なり。会社の業績は好調です。</p></body></html>"
    )
    (_dtb_html_root / "sjis.html").write_bytes(_dtb_html_src.encode("cp932"))
    _dtb_html_out = f._read_html(_dtb_html_root / "sjis.html")
    check("_read_html: cp932保存の.htmlが正しい日本語として復元される(タグ除去済み)",
          "お知らせ" in _dtb_html_out and "本日は晴天なり。会社の業績は好調です。" in _dtb_html_out)
    check("_read_html: cp932復元結果にタグが残らない(山括弧なし)",
          "<" not in _dtb_html_out and ">" not in _dtb_html_out)
    check("_read_html: cp932復元結果にU+FFFD(文字化け)が含まれない",
          "�" not in _dtb_html_out)

    # 回帰: 通常のUTF-8 .htmlは従来の読み込み結果(utf-8厳密読み)と完全一致
    _dtb_html_utf8_src = "<html><body><p>Hello UTF-8 world. こんにちは。</p></body></html>"
    (_dtb_html_root / "utf8.html").write_text(_dtb_html_utf8_src, encoding="utf-8")
    check("_read_html: 通常のUTF-8 .htmlは'Hello UTF-8 world. こんにちは。'を正しく抽出する(回帰)",
          f._read_html(_dtb_html_root / "utf8.html") == "Hello UTF-8 world. こんにちは。")

# ---------- read_file_text: 汎用テキスト分岐のNULバイトバイナリ判定 (2026-07-25) ----------
# read_file_text()のdocstringは「バイナリはスキップして空文字を返す」と約束するが、
# 実際にこれを支えるのは_BINARY_SKIP（約30拡張子のみのdenylist）だけだった。
# .npy/.h5/.parquet/.safetensors/.gguf/.sqlite/.db/.woff/.ttf/.class/.wasm/.pyc や
# 拡張子なしバイナリなど未収載の拡張子は汎用テキスト分岐まで落ちてきて
# _decode_text_bytes()（iter94: utf-8→cp932→replaceラダー、例外を送出しない設計）に
# 通され、文字化けした「ゴミテキスト」がそのまま返っていた。このゴミは--file経路
# ではmain()で質問全文そのものになり（下流フィルタなし）、RAG経路では
# _load_rag_chunks（iter42のファイル単位隔離）を通じて精度criticalなコンテキストへ
# チャンク注入される。ここではNUL(0x00)バイトの有無で真のバイナリを検出し""へ
# 落とす追加ガード（read_file_textの汎用テキスト分岐のみに限定）を検証する。
# iter94のcp932救済ラダー・iter53/125のgraceful-degradation契約・iter42のRAG単位
# 隔離のいずれも変更しないことも併せて回帰確認する。Ollama/ネットワーク呼び出しは
# 一切不要（すべてtempfile上のオフラインI/O）。
with _tempfile.TemporaryDirectory() as _nulb_dir:
    _nulb_root = _pathlib.Path(_nulb_dir)

    # (1) _BINARY_SKIP未収載の拡張子でもNULバイトを含めば""を返し警告を表示する
    _nulb_npy = _nulb_root / "weights.npy"
    _nulb_npy.write_bytes(b"\x93NUMPY\x01\x00\x00\x00\x00\x00")
    _nulb_cap1 = io.StringIO()
    with contextlib.redirect_stdout(_nulb_cap1):
        _nulb_result1 = f.read_file_text(_nulb_npy)
    check("read_file_text: _BINARY_SKIP未収載拡張子(.npy)でもNULバイト含有なら\"\"を返す",
          _nulb_result1 == "")
    check("read_file_text: NULバイト検出時にスキップ警告を表示する(ファイル名を含む)",
          "スキップ" in _nulb_cap1.getvalue() and "weights.npy" in _nulb_cap1.getvalue())

    # でっち上げ拡張子(.xyz)や別の未収載拡張子(.dat2)でも同様
    _nulb_xyz = _nulb_root / "blob.xyz"
    _nulb_xyz.write_bytes(b"garbage\x00moregarbage\x00\x01\x02")
    check("read_file_text: でっち上げ拡張子(.xyz)でもNULバイト含有なら\"\"を返す",
          f.read_file_text(_nulb_xyz) == "")

    _nulb_dat2 = _nulb_root / "raw.dat2"
    _nulb_dat2.write_bytes(b"\x00\x00\x00\x01binarydata")
    check("read_file_text: 未収載拡張子(.dat2)でもNULバイト含有なら\"\"を返す",
          f.read_file_text(_nulb_dat2) == "")

    # (2) 回帰(iter94): NULバイトを含まないcp932(Shift-JIS)日本語は従来通り正しく復元される
    _nulb_jp = "これはNULバイトを含まないcp932の日本語テキストです。会議は明日です。"
    (_nulb_root / "sjis_no_nul.txt").write_bytes(_nulb_jp.encode("cp932"))
    _nulb_jp_out = f.read_file_text(_nulb_root / "sjis_no_nul.txt")
    check("read_file_text: NUL非含有のcp932日本語は引き続き正しく復元される(iter94回帰なし)",
          _nulb_jp_out == _nulb_jp)
    check("read_file_text: NUL非含有のcp932復元結果にU+FFFD(文字化け)が含まれない(iter94回帰なし)",
          "�" not in _nulb_jp_out)

    # (3) 回帰(iter94): NUL非含有だがutf-8/cp932とも不正なバイト列はガードで弾かれず
    #     従来通りreplaceラダーを経由する("" にならない、NULガードのfalse positive防止)
    _nulb_undecodable = b"\x80\x81\xfe\xff"
    (_nulb_root / "undecodable.xyz").write_bytes(_nulb_undecodable)
    _nulb_undecodable_out = f.read_file_text(_nulb_root / "undecodable.xyz")
    check("read_file_text: NUL非含有の不正バイト列はガードで弾かれず従来通りreplaceラダーを経由する",
          _nulb_undecodable_out == _nulb_undecodable.decode("utf-8", errors="replace"))
    check("read_file_text: 上記replaceラダー結果は空文字にならない(NULガードのfalse positiveなし)",
          _nulb_undecodable_out != "")

    # (4) 回帰: NUL非含有のASCII/UTF-8テキストは未知拡張子(iter120)・_CODE_EXTENSIONS(iter121)問わず不変
    _nulb_ascii = "plain ascii content, nothing binary here at all."
    (_nulb_root / "note.xyz123").write_text(_nulb_ascii, encoding="utf-8")
    check("read_file_text: NUL非含有ASCIIの未知拡張子(.xyz123)は従来通り読み込める(iter120回帰なし)",
          f.read_file_text(_nulb_root / "note.xyz123") == _nulb_ascii)

    _nulb_code = "def f():\n    return 1\n"
    (_nulb_root / "script.py").write_text(_nulb_code, encoding="utf-8")
    check("read_file_text: NUL非含有の_CODE_EXTENSIONS(.py)は従来通り読み込める(iter121回帰なし)",
          f.read_file_text(_nulb_root / "script.py") == _nulb_code)

    # (5) 回帰: 既存の_BINARY_SKIP拡張子(.png/.exe)は従来通り早期リターンで""を返す(NULガードより前段)
    (_nulb_root / "image.png").write_bytes(b"\x89PNG\r\n\x1a\n\x00\x00\x00")
    check("read_file_text: _BINARY_SKIP拡張子(.png)は従来通り早期リターンで\"\"を返す(回帰)",
          f.read_file_text(_nulb_root / "image.png") == "")
    (_nulb_root / "app.exe").write_bytes(b"MZ\x90\x00\x03\x00\x00\x00")
    check("read_file_text: _BINARY_SKIP拡張子(.exe)は従来通り早期リターンで\"\"を返す(回帰)",
          f.read_file_text(_nulb_root / "app.exe") == "")

# (6) _load_rag_chunks経由: NULバイト含有バイナリ(未収載拡張子)は同居する正常な.txtを
#     巻き添えにせずスキップされ、RAGコンテキストへゴミが注入されない
with _tempfile.TemporaryDirectory() as _nulb_rag_dir:
    _nulb_rag_root = _pathlib.Path(_nulb_rag_dir)
    (_nulb_rag_root / "model.safetensors").write_bytes(
        b"header\x00\x00binarypayload\x00\x01\x02"
    )
    _nulb_rag_txt = "これは正常なテキストファイルの内容です。RAGに載るべき本文。"
    (_nulb_rag_root / "notes.txt").write_text(_nulb_rag_txt, encoding="utf-8")

    _nulb_rag_chunks = f._load_rag_chunks([str(_nulb_rag_root)])
    check("_load_rag_chunks: NULバイト含有バイナリ(未収載拡張子.safetensors)はチャンクに含まれない",
          all("model.safetensors" not in _p for _p, _c in _nulb_rag_chunks))
    check("_load_rag_chunks: 同居する正常な.txtのチャンクのみが含まれる",
          _nulb_rag_chunks == [(str(_nulb_rag_root / "notes.txt"), _nulb_rag_txt)])

# ---------- _read_excel/_read_pptx/_read_docx: 成功時の構造化抽出カバレッジ (2026-07-23 / iter82) ----------
# iteration 71 は stdlib のみで書かれた _read_html/_read_ipynb の「成功時」抽出を直接
# 検証したが、ライブラリ依存の _read_docx/_read_excel/_read_pptx はこれまで失敗/劣化
# 経路（iter42: 破損.xlsxがRAGから静かにスキップされる、iter53: read_file_textの
# ディスパッチ例外を握りつぶす）しかテストされておらず、「正しく抽出できた場合」の
# シート/スライド/段落マーカー・表の行・順序・空行スキップは一切検証されていなかった。
# これらの抽出結果はRAG(_load_rag_chunks)や--fileを経由してproposer/aggregatorの
# プロンプトへそのまま混入するため、抽出のサイレントな退行（シート脱落・行の
# 入れ替わり・内容切り詰め）を検出するテストがこれまで存在しなかった。
# ここでは各ライブラリ自身でtempdirへ実ファイルを書き出し(openpyxl.Workbook().save()/
# python-docx Document().save()/python-pptx Presentation().save())、
# f._read_excel/_read_pptx/_read_docx を直接呼んで成功時の出力を検証する。
# Excel/PPTXについては、iter37が修正した「出力が'['で始まるだけでライブラリ未
# インストール通知と誤判定し、RAGから丸ごと欠落させていた」過剰フィルタの回帰を
# _is_lib_missing_notice で併せて確認する(3661行目のRAG経由の回帰テストと対をなす、
# リーダー自体に対する直接の回帰ガード)。
# ライブラリ未インストール環境でもスイートが全体PASSするよう、各ライブラリの
# import可否をtry/exceptで検出し、成功時抽出テストはそのif分岐でのみ実行する
# (iteration 41/43/68のスタイルを踏襲)。ライブラリが実際に利用可能な環境でも
# 「未インストール」フォールバック行自体のカバレッジを取りこぼさないよう、
# sys.modulesを一時的に書き換えてImportErrorを模擬する経路も併せて検証する
# (iteration 43/44で確立済みの手法: sys.modules[name]=Noneでその名前のimport文に
# ImportErrorを送出させる)。テスト用の実ファイルはローカル一時ディレクトリのみに
# 書き込み、Ollama/ネットワーク/subprocess呼び出しは一切行わない。
# 表明的注記: 3ライブラリすべてで、実際に生成・読み戻した結果は仕様通り
# （マーカー順序維持・タブ結合・Noneの空文字化・完全空行/空白シェイプ/空白段落の
# スキップ）であり、抽出処理自体に新たな欠陥は見つからなかった(surface-don't-fix:
# 欠陥が見つかった場合のみ特性テストとしてピン留めしフラグする方針/iter48/71踏襲)。

# ----- _read_excel -----
try:
    import openpyxl as _openpyxl_probe_rd
    _HAS_OPENPYXL_RD = True
except ImportError:
    _HAS_OPENPYXL_RD = False

with _tempfile.TemporaryDirectory() as _rxl_dir:
    _rxl_root = _pathlib.Path(_rxl_dir)
    _rxl_path = _rxl_root / "book.xlsx"

    if _HAS_OPENPYXL_RD:
        _rxl_wb = _openpyxl_probe_rd.Workbook()
        _rxl_ws1 = _rxl_wb.active
        _rxl_ws1.title = "Sheet1"
        _rxl_ws1.append(["name", "age"])
        _rxl_ws1.append(["Alice", 30])
        _rxl_ws1.append([None, None])   # 完全に空の行 -> 出力から省略される
        _rxl_ws1.append(["Bob", None])  # Noneセルを含むが行自体は非空 -> 保持される
        _rxl_ws2 = _rxl_wb.create_sheet("Sheet2")
        _rxl_ws2.append(["x", "y"])
        _rxl_ws2.append([1, 2])
        _rxl_wb.save(str(_rxl_path))

        _rxl_out = f._read_excel(_rxl_path)
        # openpyxlのread_only=Trueワークブックはワークシートとの間に参照循環を持ち、
        # 単純な参照カウントだけでは_read_excel関数を抜けても即座に解放されず、
        # 内部のzipファイルハンドルがWindows上でロックされたままになりうる
        # (openpyxl公式にwb.close()を推奨する既知の挙動)。fugu_local.py側は
        # 変更できないため、テスト側でgc.collect()を呼び循環GCを強制発火させ、
        # 一時ディレクトリのクリーンアップ時に"プロセスはファイルにアクセスできません"
        # (WinError 32)が起きないようにする(テスト専用の後始末、抽出ロジックとは無関係)。
        import gc as _gc_rd
        _gc_rd.collect()
        _rxl_expected = "\n".join([
            "[Sheet: Sheet1]",
            "name\tage",
            "Alice\t30",
            "Bob\t",
            "[Sheet: Sheet2]",
            "x\ty",
            "1\t2",
        ])
        check("_read_excel: 複数シートの抽出結果がマーカー順序/タブ結合/空行スキップ/"
              "Noneの空文字化を含めbyte-for-byte一致",
              _rxl_out == _rxl_expected)
        check("_read_excel: '[Sheet: Sheet1]'が'[Sheet: Sheet2]'より前に出現する(ワークブック順)",
              _rxl_out.index("[Sheet: Sheet1]") < _rxl_out.index("[Sheet: Sheet2]"))
        check("_is_lib_missing_notice: _read_excel成功時の出力全体はFalse(iter37の過剰フィルタ回帰ガード)",
              not f._is_lib_missing_notice(_rxl_out))

        # sys.modules差し替えでopenpyxl/pandas双方が未インストールの状態を模擬し、
        # 実行環境に関わらず「ライブラリ未インストール」通知フォールバック行の
        # カバレッジを確保する。_read_excelはopenpyxl -> pandasの順にフォールバック
        # するため、通知文字列に到達させるには両方をImportError化する必要がある。
        _orig_openpyxl_mod_rd = sys.modules.get("openpyxl")
        _orig_pandas_mod_rd = sys.modules.get("pandas")
        sys.modules["openpyxl"] = None
        sys.modules["pandas"] = None
        try:
            _rxl_notice = f._read_excel(_rxl_path)
        finally:
            if _orig_openpyxl_mod_rd is not None:
                sys.modules["openpyxl"] = _orig_openpyxl_mod_rd
            else:
                del sys.modules["openpyxl"]
            if _orig_pandas_mod_rd is not None:
                sys.modules["pandas"] = _orig_pandas_mod_rd
            else:
                del sys.modules["pandas"]
        check("_read_excel: openpyxl/pandas双方未インストール模擬時は'pip install'通知文字列を返す",
              f._is_lib_missing_notice(_rxl_notice))
        check("_read_excel: 未インストール通知にファイル名が含まれる",
              _rxl_path.name in _rxl_notice)
    else:
        # openpyxlが無いとタスク制約(本物のopenpyxl.Workbook().save()で実.xlsxを
        # 生成すること)を満たすテスト用フィクスチャを作れない。このホストでは
        # openpyxlが常にインストール済みのため通常は通らない分岐だが、万一欠けている
        # 環境でもスイート全体を落とさないためのガードとしてスキップする
        # (3850行目の既存の実.xlsx破損読み込みテストと同じ理由・同じスタイル)。
        print("   [SKIP] openpyxl未インストールのため_read_excel成功時抽出テストをスキップ")

check("_read_excel: テスト後にsys.modulesの'openpyxl'エントリが元通り解決可能(復元確認)",
      ("openpyxl" not in sys.modules) or (sys.modules["openpyxl"] is not None))

# ----- _read_pptx -----
try:
    import pptx as _pptx_probe_rd
    _HAS_PPTX_RD = True
except ImportError:
    _HAS_PPTX_RD = False

with _tempfile.TemporaryDirectory() as _rpx_dir:
    _rpx_root = _pathlib.Path(_rpx_dir)
    _rpx_path = _rpx_root / "deck.pptx"

    if _HAS_PPTX_RD:
        from pptx.util import Inches as _Inches_rd

        _rpx_prs = _pptx_probe_rd.Presentation()
        _rpx_blank = _rpx_prs.slide_layouts[6]
        _rpx_slide1 = _rpx_prs.slides.add_slide(_rpx_blank)
        _rpx_tb1 = _rpx_slide1.shapes.add_textbox(
            _Inches_rd(1), _Inches_rd(1), _Inches_rd(4), _Inches_rd(1))
        _rpx_tb1.text_frame.text = "Slide One Text"
        _rpx_slide2 = _rpx_prs.slides.add_slide(_rpx_blank)
        _rpx_tb2 = _rpx_slide2.shapes.add_textbox(
            _Inches_rd(1), _Inches_rd(1), _Inches_rd(4), _Inches_rd(1))
        _rpx_tb2.text_frame.text = "Slide Two Text"
        _rpx_tb2b = _rpx_slide2.shapes.add_textbox(
            _Inches_rd(1), _Inches_rd(2), _Inches_rd(4), _Inches_rd(1))
        _rpx_tb2b.text_frame.text = "   "  # 空白のみのシェイプ -> スキップされる
        _rpx_prs.save(str(_rpx_path))

        _rpx_out = f._read_pptx(_rpx_path)
        _rpx_expected = "\n\n".join([
            "[Slide 1]\nSlide One Text",
            "[Slide 2]\nSlide Two Text",
        ])
        check("_read_pptx: 2スライド分の抽出結果がスライド順マーカー/本文/空白シェイプ"
              "スキップを含めbyte-for-byte一致",
              _rpx_out == _rpx_expected)
        check("_read_pptx: 空白のみのテキストボックスは出力に含まれない",
              "   " not in _rpx_out)
        check("_is_lib_missing_notice: _read_pptx成功時の出力全体はFalse(iter37の過剰フィルタ回帰ガード)",
              not f._is_lib_missing_notice(_rpx_out))

        _orig_pptx_mod_rd = sys.modules.get("pptx")
        sys.modules["pptx"] = None
        try:
            _rpx_notice = f._read_pptx(_rpx_path)
        finally:
            if _orig_pptx_mod_rd is not None:
                sys.modules["pptx"] = _orig_pptx_mod_rd
            else:
                del sys.modules["pptx"]
        check("_read_pptx: python-pptx未インストール模擬時は'pip install'通知文字列を返す",
              f._is_lib_missing_notice(_rpx_notice) and _rpx_path.name in _rpx_notice)
    else:
        # from pptx import Presentation はpathへ触れる前にImportErrorを送出するため、
        # 実ファイルを用意できなくても通知文字列フォールバック自体は安全に検証できる。
        _rpx_notice_real = f._read_pptx(_rpx_path)
        check("_read_pptx: python-pptx未インストール環境では通知文字列を返す",
              f._is_lib_missing_notice(_rpx_notice_real) and _rpx_path.name in _rpx_notice_real)

check("_read_pptx: テスト後にsys.modulesの'pptx'エントリが元通り解決可能(復元確認)",
      ("pptx" not in sys.modules) or (sys.modules["pptx"] is not None))

# ----- _read_pptx: 表シェイプ/グループシェイプの抽出 (2026-07-24 / iter87) -----
# 背景: python-pptxの表シェイプ(GraphicFrame)とグループシェイプ(GroupShape)には
# `.text` 属性そのものが存在しない。そのため上のiter82テストが固定した従来の
# _read_pptx実装 `[sh.text for sh in slide.shapes if hasattr(sh, "text") ...]` は
# 表・グループ内の文字列を静かに全て読み飛ばしていた。実データのPPTXでは表が
# 情報の大半を占めることが多く、これがRAG/--fileコンテキストへ一切届かない
# 精度事故だった(iter83の_read_pdf/iter84の_read_excelと同系統の「呼び出し元が
# 静かに失っていたコンテンツを救済する」修正)。ここでは(1)非結合2x2表、
# (2)ネストしたグループを含むグループシェイプ、(3)テキストボックス+表の混在
# スライドの3フィクスチャを実際にpython-pptxで生成し、_read_pptxが表セルを
# タブ結合/行を改行結合で拾い、グループ内テキストも再帰的に拾い、既存の
# プレーンテキストボックス経路(iter82)と併存しても重複しないことを検証する。
try:
    import pptx as _pptx_probe_tbl
    _HAS_PPTX_TBL = True
except ImportError:
    _HAS_PPTX_TBL = False

with _tempfile.TemporaryDirectory() as _rpt_dir:
    _rpt_root = _pathlib.Path(_rpt_dir)

    if _HAS_PPTX_TBL:
        from pptx.util import Inches as _Inches_tbl

        # (1) 非結合2x2表単独のスライド
        _rpt_path1 = _rpt_root / "table.pptx"
        _rpt_prs1 = _pptx_probe_tbl.Presentation()
        _rpt_blank1 = _rpt_prs1.slide_layouts[6]
        _rpt_slide1 = _rpt_prs1.slides.add_slide(_rpt_blank1)
        _rpt_gf1 = _rpt_slide1.shapes.add_table(
            2, 2, _Inches_tbl(1), _Inches_tbl(1), _Inches_tbl(4), _Inches_tbl(2))
        _rpt_tbl1 = _rpt_gf1.table
        _rpt_tbl1.cell(0, 0).text = "A1"
        _rpt_tbl1.cell(0, 1).text = "B1"
        _rpt_tbl1.cell(1, 0).text = "A2"
        _rpt_tbl1.cell(1, 1).text = "B2"
        _rpt_prs1.save(str(_rpt_path1))

        _rpt_out1 = f._read_pptx(_rpt_path1)
        _rpt_expected1 = "[Slide 1]\n" + "\n".join(["A1\tB1", "A2\tB2"])
        check("_read_pptx: 2x2表(非結合)の全セルがタブ結合/行が改行結合でbyte-for-byte一致",
              _rpt_out1 == _rpt_expected1)

        # (2) グループシェイプ(テキストボックス)+ネストしたグループのスライド
        _rpt_path2 = _rpt_root / "group.pptx"
        _rpt_prs2 = _pptx_probe_tbl.Presentation()
        _rpt_blank2 = _rpt_prs2.slide_layouts[6]
        _rpt_slide2 = _rpt_prs2.slides.add_slide(_rpt_blank2)
        _rpt_grp2 = _rpt_slide2.shapes.add_group_shape()
        _rpt_tb2 = _rpt_grp2.shapes.add_textbox(
            _Inches_tbl(1), _Inches_tbl(1), _Inches_tbl(2), _Inches_tbl(1))
        _rpt_tb2.text_frame.text = "Grouped Text"
        _rpt_nested2 = _rpt_grp2.shapes.add_group_shape()
        _rpt_tb2n = _rpt_nested2.shapes.add_textbox(
            _Inches_tbl(1), _Inches_tbl(2), _Inches_tbl(2), _Inches_tbl(1))
        _rpt_tb2n.text_frame.text = "Nested Group Text"
        _rpt_prs2.save(str(_rpt_path2))

        _rpt_out2 = f._read_pptx(_rpt_path2)
        check("_read_pptx: グループシェイプ直下のテキストが抽出される(旧コードでは脱落)",
              "Grouped Text" in _rpt_out2)
        check("_read_pptx: ネストしたグループ内のテキストも再帰的に抽出される",
              "Nested Group Text" in _rpt_out2)
        check("_read_pptx: グループ抽出結果は'[Slide 1]'マーカー配下に収まる",
              _rpt_out2 == "[Slide 1]\n" + "\n".join(["Grouped Text", "Nested Group Text"]))

        # (3) テキストボックス+表が同一スライドに混在(シェイプ追加順=出力順)
        _rpt_path3 = _rpt_root / "mixed.pptx"
        _rpt_prs3 = _pptx_probe_tbl.Presentation()
        _rpt_blank3 = _rpt_prs3.slide_layouts[6]
        _rpt_slide3 = _rpt_prs3.slides.add_slide(_rpt_blank3)
        _rpt_tb3 = _rpt_slide3.shapes.add_textbox(
            _Inches_tbl(1), _Inches_tbl(1), _Inches_tbl(4), _Inches_tbl(1))
        _rpt_tb3.text_frame.text = "Mixed Slide Text"
        _rpt_gf3 = _rpt_slide3.shapes.add_table(
            2, 2, _Inches_tbl(1), _Inches_tbl(2), _Inches_tbl(4), _Inches_tbl(2))
        _rpt_tbl3 = _rpt_gf3.table
        _rpt_tbl3.cell(0, 0).text = "C1"
        _rpt_tbl3.cell(0, 1).text = "D1"
        _rpt_tbl3.cell(1, 0).text = "C2"
        _rpt_tbl3.cell(1, 1).text = "D2"
        _rpt_prs3.save(str(_rpt_path3))

        _rpt_out3 = f._read_pptx(_rpt_path3)
        _rpt_expected3 = "[Slide 1]\n" + "\n".join(
            ["Mixed Slide Text", "C1\tD1", "C2\tD2"])
        check("_read_pptx: テキストボックス+表混在スライドは両方をシェイプ追加順で"
              "byte-for-byte抽出(重複無し)",
              _rpt_out3 == _rpt_expected3)
        check("_read_pptx: 混在スライドでテキストボックス文字列が重複挿入されない",
              _rpt_out3.count("Mixed Slide Text") == 1)

        check("_is_lib_missing_notice: 表抽出成功時の出力全体はFalse(iter37の過剰フィルタ回帰ガード)",
              not f._is_lib_missing_notice(_rpt_out1))
        check("_is_lib_missing_notice: グループ抽出成功時の出力全体はFalse(iter37の過剰フィルタ回帰ガード)",
              not f._is_lib_missing_notice(_rpt_out2))
    else:
        # python-pptxが無いとタスク制約(本物のpython-pptxで実.pptxを生成すること)を
        # 満たすテスト用フィクスチャを作れないため、上のiter82ブロックと同じ理由・
        # 同じスタイルでスキップする(このホストでは通常通らない分岐)。
        print("   [SKIP] python-pptx未インストールのため_read_pptx表/グループ抽出テストをスキップ")

check("_read_pptx: 表/グループテスト後にsys.modulesの'pptx'エントリが元通り解決可能(復元確認)",
      ("pptx" not in sys.modules) or (sys.modules["pptx"] is not None))

# ----- _read_pptx: スピーカーノート抽出 (2026-07-24 / iter98) -----
# 背景: _read_pptxは従来slide.shapesしか走査せず、slide.notes_slideには一切触れて
# いなかった。ノートには箇条書き本文が要約している詳細な説明が書かれていることが
# 多く、これがRAG/--fileコンテキストへ届かないとモデルが古い学習知識で答えてしまう
# 精度事故になる(iter87の表/グループ救済と同系統の「静かに落としているコンテンツを
# 拾う」修正)。ここでは(1)ノート付きスライド、(2)ノート無しスライド(iter82/87の
# 既存フィクスチャがbyte-for-byte不変であることの再確認)、(3)空白のみのノート、
# (4)ノート読み取りが例外を投げる壊れたケース、の4パターンを検証する。
try:
    import pptx as _pptx_probe_notes
    _HAS_PPTX_NOTES = True
except ImportError:
    _HAS_PPTX_NOTES = False

with _tempfile.TemporaryDirectory() as _rpn_dir:
    _rpn_root = _pathlib.Path(_rpn_dir)

    if _HAS_PPTX_NOTES:
        from pptx.util import Inches as _Inches_notes

        # (1)+(3) 1枚目はノートあり、2枚目は空白のみノート(→ノートブロック無し)、
        # 3枚目はノート無し(has_notes_slideがFalseのまま)のマルチスライド構成。
        _rpn_path = _rpn_root / "notes.pptx"
        _rpn_prs = _pptx_probe_notes.Presentation()
        _rpn_blank = _rpn_prs.slide_layouts[6]

        _rpn_slide1 = _rpn_prs.slides.add_slide(_rpn_blank)
        _rpn_tb1 = _rpn_slide1.shapes.add_textbox(
            _Inches_notes(1), _Inches_notes(1), _Inches_notes(4), _Inches_notes(1))
        _rpn_tb1.text_frame.text = "Slide One Body"
        _rpn_slide1.notes_slide.notes_text_frame.text = "Detailed narration for slide one."

        _rpn_slide2 = _rpn_prs.slides.add_slide(_rpn_blank)
        _rpn_tb2 = _rpn_slide2.shapes.add_textbox(
            _Inches_notes(1), _Inches_notes(1), _Inches_notes(4), _Inches_notes(1))
        _rpn_tb2.text_frame.text = "Slide Two Body"
        _rpn_slide2.notes_slide.notes_text_frame.text = "   "  # 空白のみ -> ノートブロック無し

        _rpn_slide3 = _rpn_prs.slides.add_slide(_rpn_blank)
        _rpn_tb3 = _rpn_slide3.shapes.add_textbox(
            _Inches_notes(1), _Inches_notes(1), _Inches_notes(4), _Inches_notes(1))
        _rpn_tb3.text_frame.text = "Slide Three Body"
        # slide3は notes_slide に一切触れない(has_notes_slideはFalseのまま)。

        _rpn_prs.save(str(_rpn_path))

        _rpn_out = f._read_pptx(_rpn_path)
        _rpn_expected = "\n\n".join([
            "[Slide 1]\nSlide One Body",
            "[Slide 1 Notes]\nDetailed narration for slide one.",
            "[Slide 2]\nSlide Two Body",
            "[Slide 3]\nSlide Three Body",
        ])
        check("_read_pptx: ノート付きスライドは'[Slide N Notes]'ブロックが本文の直後・"
              "正しいスライド番号でbyte-for-byte一致",
              _rpn_out == _rpn_expected)
        check("_read_pptx: 空白のみのノートは'[Slide 2 Notes]'マーカーを一切出力しない",
              "[Slide 2 Notes]" not in _rpn_out)
        check("_read_pptx: ノート未使用スライドは'[Slide 3 Notes]'マーカーを一切出力しない",
              "[Slide 3 Notes]" not in _rpn_out)
        check("_is_lib_missing_notice: ノート込み出力全体はFalse(iter37の過剰フィルタ回帰ガード)",
              not f._is_lib_missing_notice(_rpn_out))

        # (2) ノートに一切触れていないiter82の平文フィクスチャ相当の構成が、
        # 本変更後もbyte-for-byte不変であることを再確認する(有害な回帰ガード)。
        _rpn_path_plain = _rpn_root / "plain_no_notes.pptx"
        _rpn_prs_plain = _pptx_probe_notes.Presentation()
        _rpn_blank_plain = _rpn_prs_plain.slide_layouts[6]
        _rpn_slide_plain = _rpn_prs_plain.slides.add_slide(_rpn_blank_plain)
        _rpn_tb_plain = _rpn_slide_plain.shapes.add_textbox(
            _Inches_notes(1), _Inches_notes(1), _Inches_notes(4), _Inches_notes(1))
        _rpn_tb_plain.text_frame.text = "Plain Slide Text"
        _rpn_prs_plain.save(str(_rpn_path_plain))

        _rpn_out_plain = f._read_pptx(_rpn_path_plain)
        check("_read_pptx: ノート無しフィクスチャの出力は'[Slide 1]\\nPlain Slide Text'と"
              "byte-for-byte一致(iter82互換の回帰ガード)",
              _rpn_out_plain == "[Slide 1]\nPlain Slide Text")
        check("_read_pptx: ノート無しフィクスチャに'Notes'マーカーは一切現れない",
              "Notes]" not in _rpn_out_plain)

        # (4) notes_text_frameへのアクセスが例外を投げる「壊れたノートスライド」を
        # unittest.mock.patch.objectでNotesSlide.notes_text_frameプロパティに対して
        # 一時的にシミュレートし、ノートだけがスキップされてスライド本文
        # (Slide One Body)は失われず、かつ_read_pptx自体は例外を送出しないことを
        # 検証する(iter72のskip-bad-part-keep-the-rest方針)。
        import unittest.mock as _mock_notes
        import pptx.slide as _pptx_slide_mod

        def _rpn_boom(self):
            raise RuntimeError("simulated malformed notes slide")

        with _mock_notes.patch.object(
                _pptx_slide_mod.NotesSlide, "notes_text_frame",
                property(_rpn_boom)):
            _rpn_out_broken = f._read_pptx(_rpn_path)

        check("_read_pptx: notes_text_frameが例外を投げても_read_pptxはクラッシュしない"
              "(壊れたノートはスキップされるだけ)",
              "Slide One Body" in _rpn_out_broken and "Slide Two Body" in _rpn_out_broken
              and "Slide Three Body" in _rpn_out_broken)
        check("_read_pptx: 壊れたノートアクセス時は該当スライドの'[Slide N Notes]'ブロックが"
              "一切出力されない",
              "Notes]" not in _rpn_out_broken)
    else:
        # python-pptxが無いとタスク制約(本物のpython-pptxで実.pptxを生成すること)を
        # 満たすテスト用フィクスチャを作れないため、上のiter82/87ブロックと同じ理由・
        # 同じスタイルでスキップする(このホストでは通常通らない分岐)。
        print("   [SKIP] python-pptx未インストールのため_read_pptxノート抽出テストをスキップ")

check("_read_pptx: ノートテスト後にsys.modulesの'pptx'エントリが元通り解決可能(復元確認)",
      ("pptx" not in sys.modules) or (sys.modules["pptx"] is not None))

# ----- _pptx_shape_texts: 直接単体テスト (2026-07-25) -----
# 背景: カバレッジ調査により_pptx_shape_texts(fugu_local.py:997-1026)への直接の
# テスト参照が皆無であることが判明した。上のiter87テストは_read_pptx経由の間接
# テストで非結合2x2表・グループ・ネストしたグループのみをカバーしており、表の
# 水平方向マージ(gridSpan相当)セルの経路は未検証のまま残っていた。_read_docxの
# 表処理はiter91で水平マージセルの重複排除(id(getattr(c,"_tc",None))による
# トラッキング)が入ったが、本関数のdocstringにある「_read_docxの表取り扱い
# (L708-710)と同じ規約」という一文は現在の行番号(iter93以降の本文順序抽出化で
# _table_rows_textはL806-839へ移動済み)とズレて久しく、かつiter91の重複排除が
# 本関数側にも必要かは当時から未検証だった。python-docxは水平マージされた全ての
# グリッド座標で同一の_Cellオブジェクトを共有するためc.textが重複するが、
# python-pptxの仕様はこれと異なる: マージ起点セル(is_merge_origin=True)のみが
# テキストを保持し、被マージセル(is_spanned=True)の.textは空文字列を返す。これは
# 本タスク実装時に実際にpython-pptx 1.0.2を動かして実測したものであり、推測では
# ない(下のexpected値はすべて実行結果から採取した)。そのためタブ結合してもマージ
# 起点セルのテキストが重複することはなく、_read_docxのiter91のような重複排除
# ロジックはpython-pptx側には不要と判明した。よってfugu_local.py自体は変更せず、
# 以下は確認できた挙動を特性テストとして固定するものである
# (iter48/66/71と同じ「バグではなく仕様を実測確認した上でテストに固定する」方針)。
# _pptx_shape_texts(shape)を_read_pptx経由ではなく直接呼び出し、(a)平文テキスト
# ボックス、(b)空白のみ/空文字のテキストボックス、(c)非結合2x2表(空白行スキップ・
# セル順序保持)、(d)水平マージされたヘッダーセルを持つ表(マージ領域のテキストが
# 重複しないことを直接呼び出し・_read_pptx経由のデッキ出力の両方で確認)、
# (e)ネストしたグループを含むグループシェイプ(再帰・シェイプ追加順の保持)、
# (f)テキスト/表/グループのいずれでもないシェイプ(画像)、の6パターンを検証する。
try:
    import pptx as _pst_probe
    _HAS_PPTX_PST = True
except ImportError:
    _HAS_PPTX_PST = False

if _HAS_PPTX_PST:
    from pptx.util import Inches as _Inches_pst

    _pst_prs = _pst_probe.Presentation()
    _pst_blank = _pst_prs.slide_layouts[6]
    _pst_slide = _pst_prs.slides.add_slide(_pst_blank)

    # (a) 平文テキストボックス(非空)
    _pst_tb_a = _pst_slide.shapes.add_textbox(
        _Inches_pst(1), _Inches_pst(1), _Inches_pst(2), _Inches_pst(1))
    _pst_tb_a.text_frame.text = "Hello Text"
    check("_pptx_shape_texts: 平文テキストボックス(非空)はテキスト1件のリストを返す",
          f._pptx_shape_texts(_pst_tb_a) == ["Hello Text"])

    # (b) 空白のみ/空文字のテキストボックスは空リスト
    _pst_tb_b1 = _pst_slide.shapes.add_textbox(
        _Inches_pst(1), _Inches_pst(2), _Inches_pst(2), _Inches_pst(1))
    _pst_tb_b1.text_frame.text = "   "
    check("_pptx_shape_texts: 空白のみのテキストボックスは空リスト",
          f._pptx_shape_texts(_pst_tb_b1) == [])

    _pst_tb_b2 = _pst_slide.shapes.add_textbox(
        _Inches_pst(1), _Inches_pst(3), _Inches_pst(2), _Inches_pst(1))
    _pst_tb_b2.text_frame.text = ""
    check("_pptx_shape_texts: 空文字のテキストボックスは空リスト",
          f._pptx_shape_texts(_pst_tb_b2) == [])

    # (c) 非結合3x2表: 空白行(2行目)はスキップされ、セル順序はそのまま保持される
    _pst_gf_c = _pst_slide.shapes.add_table(
        3, 2, _Inches_pst(1), _Inches_pst(4), _Inches_pst(4), _Inches_pst(2))
    _pst_tbl_c = _pst_gf_c.table
    _pst_tbl_c.cell(0, 0).text = "A1"
    _pst_tbl_c.cell(0, 1).text = "B1"
    _pst_tbl_c.cell(1, 0).text = " "   # 空白のみ
    _pst_tbl_c.cell(1, 1).text = ""    # 空文字
    _pst_tbl_c.cell(2, 0).text = "C1"
    _pst_tbl_c.cell(2, 1).text = "D1"
    check("_pptx_shape_texts: 非結合3x2表は空白行(2行目)をスキップし、"
          "タブ結合された行のみをセル順序を保ったまま返す",
          f._pptx_shape_texts(_pst_gf_c) == ["A1\tB1", "C1\tD1"])

    # (d) 水平マージ(2列にまたがるヘッダーセル)を持つ表
    # python-pptxのTable.cell(r,c).merge()は起点セル(is_merge_origin=True)にのみ
    # テキストを保持し、被マージセル(is_spanned=True)の.textは空文字列を返す
    # (python-docxが両座標で同一_Cellを共有しc.textが重複するのとは異なる挙動。
    # 実測値をそのままexpectedにしている)。そのためタブ結合してもマージ起点セルの
    # テキストが重複することはない。
    _pst_gf_d = _pst_slide.shapes.add_table(
        2, 2, _Inches_pst(1), _Inches_pst(7), _Inches_pst(4), _Inches_pst(2))
    _pst_tbl_d = _pst_gf_d.table
    _pst_tbl_d.cell(0, 0).merge(_pst_tbl_d.cell(0, 1))
    _pst_tbl_d.cell(0, 0).text = "Header"
    _pst_tbl_d.cell(1, 0).text = "R1C1"
    _pst_tbl_d.cell(1, 1).text = "R1C2"
    _pst_out_d = f._pptx_shape_texts(_pst_gf_d)
    check("_pptx_shape_texts: 水平マージされたヘッダーセルを含む表からの直接呼び出し"
          "結果(実測: マージ起点セル'Header'+空文字の被マージセルがタブ結合される)",
          _pst_out_d == ["Header\t", "R1C1\tR1C2"])
    check("_pptx_shape_texts: 水平マージ領域のテキスト'Header'は直接呼び出し結果"
          "全体でちょうど1回しか出現しない(python-docx側iter91が対処した重複は"
          "python-pptxでは発生しない=特性テストとして固定)",
          "".join(_pst_out_d).count("Header") == 1)
    check("_pptx_shape_texts: マージ行に続く非マージ行(R1C1/R1C2)はタブ結合のまま"
          "影響を受けない",
          _pst_out_d[1] == "R1C1\tR1C2")

    # (d-2) 同じマージ表を実ファイルに保存し、_read_pptx経由のデッキ全体出力でも
    # 'Header'がちょうど1回しか出現しないことを確認する(受け入れ基準(d)の
    # 「デッキ出力全体で1回のみ」を直接呼び出しだけでなくエンドツーエンドでも保証)。
    with _tempfile.TemporaryDirectory() as _pstd_dir:
        _pstd_root = _pathlib.Path(_pstd_dir)
        _pstd_path = _pstd_root / "merged_table.pptx"
        _pstd_prs = _pst_probe.Presentation()
        _pstd_blank = _pstd_prs.slide_layouts[6]
        _pstd_slide = _pstd_prs.slides.add_slide(_pstd_blank)
        _pstd_gf = _pstd_slide.shapes.add_table(
            2, 2, _Inches_pst(1), _Inches_pst(1), _Inches_pst(4), _Inches_pst(2))
        _pstd_tbl = _pstd_gf.table
        _pstd_tbl.cell(0, 0).merge(_pstd_tbl.cell(0, 1))
        _pstd_tbl.cell(0, 0).text = "Header"
        _pstd_tbl.cell(1, 0).text = "R1C1"
        _pstd_tbl.cell(1, 1).text = "R1C2"
        _pstd_prs.save(str(_pstd_path))

        _pstd_out = f._read_pptx(_pstd_path)
        check("_read_pptx: 水平マージ表を含むデッキ全体の出力で'Header'はちょうど1回"
              "しか出現しない(重複なし、直接呼び出しテスト(d)と整合)",
              _pstd_out.count("Header") == 1)
        check("_read_pptx: 水平マージ表を含むデッキ全体の出力がbyte-for-byte一致"
              "(マージ起点セル+空文字の被マージセルをタブ結合)",
              _pstd_out == "[Slide 1]\n" + "\n".join(["Header\t", "R1C1\tR1C2"]))

    # (e) グループシェイプ + ネストしたグループ: 再帰してもシェイプ追加順が保たれる
    _pst_grp_e = _pst_slide.shapes.add_group_shape()
    _pst_tb_e1 = _pst_grp_e.shapes.add_textbox(
        _Inches_pst(1), _Inches_pst(1), _Inches_pst(2), _Inches_pst(1))
    _pst_tb_e1.text_frame.text = "First"
    _pst_nested_e = _pst_grp_e.shapes.add_group_shape()
    _pst_tb_e2 = _pst_nested_e.shapes.add_textbox(
        _Inches_pst(1), _Inches_pst(2), _Inches_pst(2), _Inches_pst(1))
    _pst_tb_e2.text_frame.text = "Nested"
    _pst_tb_e3 = _pst_grp_e.shapes.add_textbox(
        _Inches_pst(1), _Inches_pst(3), _Inches_pst(2), _Inches_pst(1))
    _pst_tb_e3.text_frame.text = "Third"
    check("_pptx_shape_texts: グループ+ネストしたグループはシェイプ追加順"
          "(First, ネスト内Nested, Third)を保ったまま再帰的に抽出する",
          f._pptx_shape_texts(_pst_grp_e) == ["First", "Nested", "Third"])

    # (f) テキスト/表/グループのいずれでもないシェイプ(画像)は空リスト
    import base64 as _base64_pst

    _pst_png_f = _base64_pst.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+A8AAQUBAScY"
        "42YAAAAASUVORK5CYII="
    )
    _pst_pic_f = _pst_slide.shapes.add_picture(
        io.BytesIO(_pst_png_f), _Inches_pst(5), _Inches_pst(5))
    check("_pptx_shape_texts: 画像シェイプ(text/table/groupのいずれも持たない)は"
          "空リスト", f._pptx_shape_texts(_pst_pic_f) == [])
else:
    print("   [SKIP] python-pptx未インストールのため_pptx_shape_texts直接テストをスキップ")

check("_pptx_shape_texts: テスト後にsys.modulesの'pptx'エントリが元通り解決可能(復元確認)",
      ("pptx" not in sys.modules) or (sys.modules["pptx"] is not None))

# ----- _pptx_shape_texts: ネストグループ再帰の深さ上限 (2026-07-26) -----
# 背景: _pptx_shape_texts(iter87, 上のブロック)はグループシェイプへ
# `elif hasattr(sh, "shapes"): for sub in sh.shapes: out.extend(_pptx_shape_texts(sub))`
# で無条件に再帰しており、病的/破損した<p:grpSp>の異常に深いネスト連鎖を与えると
# sys.recursionlimitを超えてRecursionErrorを送出しうる。_read_pptx側のtry節は
# import pptx失敗によるImportErrorしか捕捉していないため、このRecursionErrorは
# そのまま外へ伝播し、read_file_text(iter53)・_load_rag_chunks(iter42)の広い
# except Exceptionガードに握りつぶされて、PowerPointファイル全体がRAG/--file
# コンテキストから無言で丸ごと脱落する(iter87が救済した「表/グループ内容が静かに
# 読み飛ばされる」精度事故と同じ害のクラス)。本タスクはiter157(_read_docxの
# ネスト表再帰への深さ上限_DOCX_NESTED_TABLE_MAX_DEPTH)と同じ方針で、
# _pptx_shape_textsに_depth引数と module-level の _PPTX_GROUP_MAX_DEPTH を追加した。
# 上のiter87/iter25(2026-07-25)テストは実python-pptxオブジェクトに依存しており
# python-pptx未インストール環境ではスキップされてしまうため、ここではダック
# タイピングの疑似シェイプ(.textを持つ葉/.shapesを持つグループ/has_table+
# .table.rowsを持つ表)だけで構成し、python-pptxの有無に関わらず必ず実行される。


class _FakePptxTextShape:
    """.textだけを持つ葉シェイプ(平文テキストボックス相当)のダックタイピング版。"""

    def __init__(self, text):
        self.text = text


class _FakePptxGroupShape:
    """.shapesだけを持つグループシェイプ相当のダックタイピング版
    (.textも.has_tableも持たない=python-pptxのGroupShapeの実挙動どおり)。"""

    def __init__(self, shapes):
        self.shapes = shapes


class _FakePptxCell:
    def __init__(self, text):
        self.text = text


class _FakePptxRow:
    def __init__(self, cell_texts):
        self.cells = [_FakePptxCell(t) for t in cell_texts]


class _FakePptxTable:
    def __init__(self, rows):
        self.rows = [_FakePptxRow(r) for r in rows]


class _FakePptxTableShape:
    """has_table=True + .table.rows/.cells だけを持つ表シェイプ相当
    (.textは持たない=python-pptxのGraphicFrameの実挙動どおり)。"""

    has_table = True

    def __init__(self, rows):
        self.table = _FakePptxTable(rows)


# (1) 単一の非グループシェイプ(平文.text)/空白のみ・空文字の.textは従来どおり。
check("_pptx_shape_texts: ダックタイピング版の平文テキストシェイプはテキスト1件のリスト",
      f._pptx_shape_texts(_FakePptxTextShape("Hello Fake")) == ["Hello Fake"])
check("_pptx_shape_texts: ダックタイピング版の空白のみ.textは空リスト",
      f._pptx_shape_texts(_FakePptxTextShape("   ")) == [])
check("_pptx_shape_texts: ダックタイピング版の空文字.textは空リスト",
      f._pptx_shape_texts(_FakePptxTextShape("")) == [])

# (2) cap以下(深さ3)のネストグループ+テキストボックス混在は変更前の実装と
# byte-for-byte同一の出力(シェイプ追加順)を返す。
_ppwc_inner = _FakePptxGroupShape([_FakePptxTextShape("Nested2")])
_ppwc_mid = _FakePptxGroupShape([_FakePptxTextShape("Nested1"), _ppwc_inner])
_ppwc_top = _FakePptxGroupShape(
    [_FakePptxTextShape("First"), _ppwc_mid, _FakePptxTextShape("Third")])
check("_pptx_shape_texts: cap以下(深さ3)のネストグループはシェイプ追加順で"
      "byte-for-byte一致(First, Nested1, Nested2, Third)",
      f._pptx_shape_texts(_ppwc_top) == ["First", "Nested1", "Nested2", "Third"])

# (3) cap以下のグループ内にネストした表シェイプはhas_table分岐がそのまま働き、
# タブ結合された行が抽出される(_depth引数の追加はhas_table分岐を一切変更しない)。
_ppwc_table_shape = _FakePptxTableShape([["A1", "B1"], ["C1", "D1"]])
_ppwc_group_with_table = _FakePptxGroupShape(
    [_FakePptxTextShape("Before"), _ppwc_table_shape])
check("_pptx_shape_texts: cap以下のグループ内にネストした表もタブ結合行として抽出される",
      f._pptx_shape_texts(_ppwc_group_with_table) == ["Before", "A1\tB1", "C1\tD1"])

# (4) capを超える病的ネスト(cap+10段)でもRecursionErrorを送出せず、
# 深さ<= capのテキストは救済され、それより深いテキストは打ち切られる。
_ppgd_cap = f._PPTX_GROUP_MAX_DEPTH
_ppgd_chain_depth = _ppgd_cap + 10
_ppgd_node = _FakePptxTextShape(f"T{_ppgd_chain_depth}")
for _ppgd_i in range(_ppgd_chain_depth - 1, -1, -1):
    _ppgd_node = _FakePptxGroupShape(
        [_FakePptxTextShape(f"T{_ppgd_i}"), _ppgd_node])

_ppgd_exc = None
try:
    _ppgd_out = f._pptx_shape_texts(_ppgd_node)
except Exception as _ppgd_e:
    _ppgd_exc = _ppgd_e
    _ppgd_out = None

check(f"_pptx_shape_texts: cap({_ppgd_cap})+10段の病的ネストグループでも"
      "RecursionErrorを送出しない",
      _ppgd_exc is None)
check("_pptx_shape_texts: 病的ネストでも浅い階層(T0)のテキストは救済される",
      _ppgd_out is not None and "T0" in _ppgd_out)
check("_pptx_shape_texts: capを超える深さの最深部テキストは打ち切られ含まれない",
      _ppgd_out is not None and f"T{_ppgd_chain_depth}" not in _ppgd_out)
check("_pptx_shape_texts: 病的ネストでも深さ<=cap-1のテキストはT0,T1,...の"
      "追加順(昇順)のままcap件だけ救済される(それ以深は静かに打ち切り)",
      _ppgd_out is not None and _ppgd_out == [f"T{_i}" for _i in range(_ppgd_cap)])

# ----- _read_pptx: 1シェイプ単位の抽出単離 (2026-07-26) -----
# 背景: fugu_local.py側の_read_pptxコメントの通り、_pptx_shape_texts内部の
# hasattr(sh, "text")/getattr(sh, "has_table", False)によるダックタイピング分岐は
# AttributeErrorしか吸収しない(Python 3の仕様でhasattrはAttributeError以外の
# 例外をそのまま伝播させる)。そのため壊れた1シェイプの.text/.table参照が別の
# 例外を送出すると、従来は_read_pptxの
# `for sh in slide.shapes: texts.extend(_pptx_shape_texts(sh))` ループの外、
# ひいては_read_pptx自体の外まで例外が伝播していた。_read_pptxには
# _read_pdf/_read_excelと違って代替ライブラリへのフォールスルーが無いため、
# 1シェイプの異常だけでデッキ全体(他の全スライド・同スライドの他の全シェイプ・
# 既にiter98で個別保護済みのスピーカーノート)がread_file_text(iter53)の広い
# except Exceptionに握りつぶされて""へ丸ごと脱落していた。本タスクは
# `texts.extend(_pptx_shape_texts(sh))` の呼び出し1回分だけをtry/exceptで
# 保護し、iter98(ノート個別保護)/iter93(_read_docx本文走査保護)と同じ
# skip-bad-part-keep-the-rest方針を本文シェイプ抽出ループへ適用する。
# python-pptxの有無に関わらず必ず実行されるよう、実python-pptxではなく
# sys.modulesへ注入するフェイクの'pptx'モジュール(iter43/44で確立済みの
# swap-restoreパターンをこのブロック内だけで再現)と、上のブロックで定義済みの
# ダックタイピング疑似シェイプ(_FakePptxTextShape)+新設の「.textアクセスで
# 例外を送出する壊れたシェイプ」だけで構成する。


class _FakePptxRaisingShape:
    """.textへのアクセスが例外を送出する壊れたシェイプ(破損したXML片や非対応の
    グラフィック要素を想定)。hasattr(sh, "text")はPython 3の仕様上
    AttributeErrorしか吸収しないため、この例外は_pptx_shape_texts内の
    hasattr呼び出し自体からそのまま伝播する。"""

    @property
    def text(self):
        raise RuntimeError("simulated malformed shape (corrupt XML)")


class _FakePptxSlideForIsolation:
    """.shapes(リスト)と.has_notes_slide(常にFalse)だけを持つスライド相当の
    ダックタイピング版。ノート抽出自体は本テストの対象外(iter98で別途検証済み)。"""

    def __init__(self, shapes):
        self.shapes = shapes
        self.has_notes_slide = False


def _rpi_make_fake_pptx_module(slides):
    """Presentation(path)が`slides`をそのまま`.slides`として返すフェイク'pptx'
    モジュールを組み立てる。"""
    mod = types.ModuleType("pptx")

    class _FakePresentation:
        def __init__(self, path=None):
            self.slides = slides

    mod.Presentation = _FakePresentation
    return mod


def _rpi_swap_pptx_module(fake_mod, body):
    """sys.modules['pptx']をfake_modへ一時的に差し替えてbody()を実行し、
    元の状態(存在した/しなかった)へ必ず復元する(iter43/44のswap-restore
    パターンをこのブロック専用に再現したもの)。"""
    _orig = sys.modules.get("pptx")
    sys.modules["pptx"] = fake_mod
    try:
        return body()
    finally:
        if _orig is not None:
            sys.modules["pptx"] = _orig
        else:
            del sys.modules["pptx"]


def _rpi_is_cp932_safe(s):
    """gotcha#4: Windowsコンソール(cp932)でエンコード不能な文字(絵文字等)を
    含んでいないかを確認する。"""
    try:
        s.encode("cp932")
        return True
    except UnicodeEncodeError:
        return False


with _tempfile.TemporaryDirectory() as _rpi_dir:
    _rpi_path = _pathlib.Path(_rpi_dir) / "isolate.pptx"
    # フェイクPresentationは中身を一切パースしないため、実体は空でよい
    # (存在確認や拡張子ベースの分岐が万一入っても壊れないよう実ファイルにする)。
    _rpi_path.write_bytes(b"dummy pptx bytes, not parsed by the fake Presentation")

    # (1) 1枚のスライドに正常シェイプ1件+壊れたシェイプ1件。壊れたシェイプの
    #     例外は伝播せず、正常シェイプのテキストが救済され、cp932安全な警告が
    #     path.nameと例外型名(RuntimeError)を伴って出力される。
    _rpi_slide1 = _FakePptxSlideForIsolation(
        [_FakePptxTextShape("Good Shape Text"), _FakePptxRaisingShape()])
    _rpi_cap1 = io.StringIO()
    with contextlib.redirect_stdout(_rpi_cap1):
        _rpi_out1 = _rpi_swap_pptx_module(
            _rpi_make_fake_pptx_module([_rpi_slide1]),
            lambda: f._read_pptx(_rpi_path),
        )
    check("_read_pptx: 壊れたシェイプが例外を送出してもクラッシュせず正常シェイプの"
          "テキストが救済される",
          _rpi_out1 == "[Slide 1]\nGood Shape Text")
    check("_read_pptx: 壊れたシェイプの警告にファイル名(path.name)が出力される",
          _rpi_path.name in _rpi_cap1.getvalue())
    check("_read_pptx: 壊れたシェイプの警告に例外型名(RuntimeError)が出力される",
          "RuntimeError" in _rpi_cap1.getvalue())
    check("_read_pptx: 壊れたシェイプの警告メッセージはcp932でエンコード可能"
          "(絵文字等の非cp932文字を含まない、gotcha#4)",
          _rpi_is_cp932_safe(_rpi_cap1.getvalue()))

    # (2) 正常シェイプ2件の間に壊れたシェイプを挟んでも、両方の正常シェイプが
    #     元の追加順のまま残り、壊れたシェイプの寄与だけが欠落する。さらに
    #     全シェイプ正常な2枚目のスライドは一切影響を受けない。
    _rpi_slide2a = _FakePptxSlideForIsolation([
        _FakePptxTextShape("First Good"),
        _FakePptxRaisingShape(),
        _FakePptxTextShape("Second Good"),
    ])
    _rpi_slide2b = _FakePptxSlideForIsolation([_FakePptxTextShape("Slide Two Text")])
    _rpi_out2 = _rpi_swap_pptx_module(
        _rpi_make_fake_pptx_module([_rpi_slide2a, _rpi_slide2b]),
        lambda: f._read_pptx(_rpi_path),
    )
    check("_read_pptx: 正常シェイプ2件に挟まれた壊れたシェイプは両側の正常シェイプを"
          "元の順序のまま残す(byte-for-byte一致)",
          _rpi_out2 == "[Slide 1]\nFirst Good\nSecond Good\n\n[Slide 2]\nSlide Two Text")
    check("_read_pptx: 壊れたシェイプ自身の抽出結果は出力に一切現れない",
          "Raising" not in _rpi_out2)
    check("_read_pptx: 1枚目の壊れたシェイプは2枚目(全シェイプ正常)に影響しない",
          "Slide Two Text" in _rpi_out2)

    # (3) 空デッキ境界: 全スライドの全シェイプが例外を送出する場合でも
    #     python-pptxは正常にimportできている(=ImportError分岐ではない)ため、
    #     pip installの未インストール通知ではなく空文字列を返す。
    _rpi_slide3 = _FakePptxSlideForIsolation(
        [_FakePptxRaisingShape(), _FakePptxRaisingShape()])
    _rpi_out3 = _rpi_swap_pptx_module(
        _rpi_make_fake_pptx_module([_rpi_slide3]),
        lambda: f._read_pptx(_rpi_path),
    )
    check("_read_pptx: 全シェイプが例外を送出する空デッキは''を返す"
          "(pip install通知ではない)",
          _rpi_out3 == "")
    check("_is_lib_missing_notice: 全滅デッキの''はFalse(未インストール通知と"
          "誤判定されない)",
          not f._is_lib_missing_notice(_rpi_out3))

check("_read_pptx: シェイプ単離テスト後にsys.modulesの'pptx'エントリが元通り"
      "解決可能(復元確認)",
      ("pptx" not in sys.modules) or (sys.modules["pptx"] is not None))

# ----- _read_docx -----
try:
    import docx as _docx_probe_rd
    _HAS_DOCX_RD = True
except ImportError:
    _HAS_DOCX_RD = False

with _tempfile.TemporaryDirectory() as _rdx_dir:
    _rdx_root = _pathlib.Path(_rdx_dir)
    _rdx_path = _rdx_root / "doc.docx"

    if _HAS_DOCX_RD:
        _rdx_doc = _docx_probe_rd.Document()
        _rdx_doc.add_paragraph("First paragraph.")
        _rdx_doc.add_paragraph("")  # 空白段落 -> "if para.text.strip()" でスキップされる
        _rdx_doc.add_paragraph("Second paragraph.")
        _rdx_table = _rdx_doc.add_table(rows=2, cols=2)  # マージセル無しの単純な表
        _rdx_table.cell(0, 0).text = "H1"
        _rdx_table.cell(0, 1).text = "H2"
        _rdx_table.cell(1, 0).text = "R1C1"
        _rdx_table.cell(1, 1).text = "R1C2"
        _rdx_doc.save(str(_rdx_path))

        _rdx_out = f._read_docx(_rdx_path)
        _rdx_expected = "\n".join([
            "First paragraph.",
            "Second paragraph.",
            "H1\tH2",
            "R1C1\tR1C2",
        ])
        check("_read_docx: 段落2件(空白段落はスキップ)+表2行の抽出結果が"
              "byte-for-byte一致(段落->表の順序含む)",
              _rdx_out == _rdx_expected)
        check("_read_docx: 空白段落が出力中に空行として残らない",
              "" not in _rdx_out.split("\n"))

        _orig_docx_mod_rd = sys.modules.get("docx")
        sys.modules["docx"] = None
        try:
            _rdx_notice = f._read_docx(_rdx_path)
        finally:
            if _orig_docx_mod_rd is not None:
                sys.modules["docx"] = _orig_docx_mod_rd
            else:
                del sys.modules["docx"]
        _rdx_expected_notice = f"[DOCX: {_rdx_path.name} — python-docx が必要: pip install python-docx]"
        check("_read_docx: python-docx未インストール模擬時は'pip install'通知文字列を返す",
              _rdx_notice == _rdx_expected_notice)
        check("_is_lib_missing_notice: _read_docxの未インストール通知はTrue(併せて確認)",
              f._is_lib_missing_notice(_rdx_notice))
    else:
        # docx.Document(...)へ到達する前にImportErrorが送出されるため、
        # 実ファイルを用意できなくても通知文字列フォールバック自体は安全に検証できる。
        _rdx_notice_real = f._read_docx(_rdx_path)
        check("_read_docx: python-docx未インストール環境では通知文字列を返す",
              _rdx_notice_real == f"[DOCX: {_rdx_path.name} — python-docx が必要: pip install python-docx]")

check("_read_docx: テスト後にsys.modulesの'docx'エントリが元通り解決可能(復元確認)",
      ("docx" not in sys.modules) or (sys.modules["docx"] is not None))

# ---------- _read_docx: 水平マージセル(gridSpan)の重複排除 (2026-07-24 / iter91) ----------
# python-docxのrow.cellsは表のグリッド座標の数だけ_Cellを返すが、水平方向(gridSpan)に
# マージされたセルは全ての被マージ座標で同一の<w:tc>要素を共有するため、c.textは
# マージされたテキストを座標数だけ重複して返す(例: 2列にまたがるヘッダーは
# 'Header\tHeader')。iter82のテストは「マージセル無しの単純な表」しかカバーしておらず、
# この重複抽出の穴は未検証だった。本テストはiter87(_read_pptxの表/グループ抽出)の
# 直系である_read_docx側の同種の抽出漏れ/重複修正を検証する。
with _tempfile.TemporaryDirectory() as _rdxm_dir:
    _rdxm_root = _pathlib.Path(_rdxm_dir)
    _rdxm_path = _rdxm_root / "merged.docx"

    if _HAS_DOCX_RD:
        _rdxm_doc = _docx_probe_rd.Document()
        _rdxm_table = _rdxm_doc.add_table(rows=2, cols=2)
        # 1行目: 0列目と1列目を水平マージして1つのヘッダーセルにする
        _rdxm_merged_cell = _rdxm_table.cell(0, 0).merge(_rdxm_table.cell(0, 1))
        _rdxm_merged_cell.text = "Header"
        # 2行目: マージ無しの通常セル(従来通りの挙動が保たれることの対照群)
        _rdxm_table.cell(1, 0).text = "R1C1"
        _rdxm_table.cell(1, 1).text = "R1C2"
        _rdxm_doc.save(str(_rdxm_path))

        _rdxm_out = f._read_docx(_rdxm_path)
        _rdxm_expected = "\n".join([
            "Header",
            "R1C1\tR1C2",
        ])
        check("_read_docx: 水平マージされたヘッダーセルは行につき1回だけ抽出される"
              "('Header\\tHeader'のように重複しない)",
              _rdxm_out == _rdxm_expected)
        check("_read_docx: マージ行に続く非マージ行(R1C1/R1C2)はタブ結合のまま影響を受けない",
              _rdxm_out.split("\n")[1] == "R1C1\tR1C2")
    else:
        # python-docx未インストール環境ではマージ表フィクスチャ自体を構築できないため、
        # 通知文字列フォールバックの確認は既存の"_read_docx未インストール"テストに委譲する。
        pass

# _read_docx: gridSpanの重複排除で参照する getattr(c, "_tc", None) が None を返す場合
# (将来のpython-docx実装変化などで内部属性 _tc 自体が取得できなくなるケースの模擬)でも
# 例外を送出せず、重複排除を行わない従来通りの挙動へ安全にフォールバックすることを検証する。
# python-docxが未インストールの環境でも成立するよう、docx.Document自体をダックタイピングの
# 偽オブジェクトに差し替えて検証する(python-docxの実クラス/内部実装には一切触れない)。
_rdxtc_fake_mod = types.ModuleType("docx")


class _RdxTcFakeCell:
    """_tc属性を意図的に持たないセルの模擬(python-docxの_Cellではない)。"""

    def __init__(self, text):
        self.text = text


class _RdxTcFakeRow:
    def __init__(self, cells):
        self.cells = cells


class _RdxTcFakeTable:
    def __init__(self, rows):
        self.rows = rows


class _RdxTcFakeDoc:
    def __init__(self):
        self.paragraphs = []
        # 同一テキスト"A"を持つ2セル: _tcが無いため重複排除されず、
        # 従来通り"A\tA"として出力されるはず。
        self.tables = [_RdxTcFakeTable([_RdxTcFakeRow([_RdxTcFakeCell("A"), _RdxTcFakeCell("A")])])]


_rdxtc_fake_mod.Document = lambda _path: _RdxTcFakeDoc()

_orig_docx_mod_tc = sys.modules.get("docx")
sys.modules["docx"] = _rdxtc_fake_mod
try:
    _rdxtc_out = f._read_docx(_pathlib.Path("dummy_for_tc_fallback_test.docx"))
finally:
    if _orig_docx_mod_tc is not None:
        sys.modules["docx"] = _orig_docx_mod_tc
    else:
        del sys.modules["docx"]

check("_read_docx: getattr(c,'_tc',None)がNoneの場合は例外を送出せず"
      "重複排除なしの従来通りの挙動にフォールバックする(id依存しない安全側動作)",
      _rdxtc_out == "A\tA")

check("_read_docx: テスト後にsys.modulesの'docx'エントリが元通り解決可能(_tcフォールバック検証後の復元確認)",
      ("docx" not in sys.modules) or (sys.modules["docx"] is not None))

# ---------- _read_docx: 本文の読み取り順序を保持する(段落<->表の交互出現) (2026-07-24 / iter93) ----------
# 従来の_read_docxはdoc.paragraphsを全件処理してからdoc.tablesを全件処理する2パス
# 構成だった。導入文 -> データ表 -> 結論文、のように本文中に表が挟まる文書では、
# 実際の読み順(導入文・表・結論文)ではなく「導入文・結論文」の後に「表」が
# まとめて出力される形に入れ替わり、表とそれを説明する地の文が同じRAGチャンク/
# --file文脈に絶対に同居できなくなる精度事故だった。iter82(非マージ表テスト)は
# 表が段落の後ろに追加されるフィクスチャしか使っておらず、iter91(水平マージ
# 重複排除)もこの順序入れ替わり自体には触れていなかった。ここではdoc.element.body
# を直接歩くPython実装が実際に本文出現順で段落・表を交互抽出することと、
# iter82/iter91の非交互フィクスチャ(表が段落群の後ろにあるだけ)が退行しないことの
# 両方を検証する。
with _tempfile.TemporaryDirectory() as _rdxo_dir:
    _rdxo_root = _pathlib.Path(_rdxo_dir)
    _rdxo_path = _rdxo_root / "order.docx"

    if _HAS_DOCX_RD:
        _rdxo_doc = _docx_probe_rd.Document()
        _rdxo_doc.add_paragraph("Intro")
        _rdxo_doc.add_paragraph("")     # 空白段落 -> 表の直前でもスキップされる
        _rdxo_t1 = _rdxo_doc.add_table(rows=2, cols=2)
        _rdxo_t1.cell(0, 0).text = "A1"
        _rdxo_t1.cell(0, 1).text = "B1"
        _rdxo_t1.cell(1, 0).text = "A2"
        _rdxo_t1.cell(1, 1).text = "B2"
        _rdxo_doc.add_paragraph("   ")  # 空白のみの段落 -> 表の直後でもスキップされる
        _rdxo_t2 = _rdxo_doc.add_table(rows=1, cols=2)  # 本文中の2つ目の表
        _rdxo_t2.cell(0, 0).text = "C1"
        _rdxo_t2.cell(0, 1).text = "C2"
        _rdxo_doc.add_paragraph("Conclusion")
        _rdxo_doc.save(str(_rdxo_path))

        _rdxo_out = f._read_docx(_rdxo_path)
        _rdxo_expected = "\n".join([
            "Intro",
            "A1\tB1",
            "A2\tB2",
            "C1\tC2",
            "Conclusion",
        ])
        check("_read_docx (iter93): Intro->表1(2行)->表2(1行)->Conclusionが本文出現順の"
              "まま交互にbyte-for-byte一致で抽出される",
              _rdxo_out == _rdxo_expected)
        check("_read_docx (iter93): コアアサーション— 表(A1)は両方の段落の'間'に出現する"
              "(前後にまとめられない、旧2パス実装では失敗するはずの検証)",
              _rdxo_out.index("Intro") < _rdxo_out.index("A1\tB1") < _rdxo_out.index("Conclusion"))
        check("_read_docx (iter93): 表の直前/直後の空白のみ段落は交互抽出中もスキップされる",
              "" not in _rdxo_out.split("\n"))
        _rdxo_old_two_pass_order = "\n".join([
            "Intro", "Conclusion", "A1\tB1", "A2\tB2", "C1\tC2",
        ])
        check("_read_docx (iter93): 修正前の段落→表2パス構成が生成していたはずの順序"
              "(退行ガード。旧実装ならこのcheckはFalseになる)",
              _rdxo_out != _rdxo_old_two_pass_order)

        # iter82/iter91のフィクスチャは表を段落群の後ろに追加するだけで本文順序が
        # 交互にならないため、この順序修正では出力が変化しないはずの回帰確認。
        # (実体は_rdx_out/_rdxm_out として上でbyte-for-byte一致を既に検証済みで、
        # ここでは同じ2フィクスチャがiter93の変更後も再現することを明示的に確認する)
        check("_read_docx (iter93回帰): iter82の非マージ表フィクスチャ(表が段落群の後ろ)"
              "は順序修正後も従来通りbyte-for-byte一致のまま",
              _rdx_out == "\n".join(["First paragraph.", "Second paragraph.", "H1\tH2", "R1C1\tR1C2"]))
        check("_read_docx (iter93回帰): iter91の水平マージ表フィクスチャ(表のみ、段落無し)"
              "は順序修正後も従来通りbyte-for-byte一致のまま(マージ行につき1回)",
              _rdxm_out == "\n".join(["Header", "R1C1\tR1C2"]))
    else:
        print("   [SKIP] python-docx未インストールのため_read_docx本文順序テストをスキップ")

# ---------- _read_docx: セル内にネストした表へ再帰し、内容の欠落を解消する (2026-07-25) ----------
# 背景: python-docxの_Cell.textはそのセル直下の<w:p>段落だけを連結し、セル内に
# ネストされた<w:tbl>は一切含めない。そのため「表の中の表」構成の文書では、
# ネスト表の内容が_read_docxの出力から――ひいては_load_rag_chunks(RAG)や
# --fileの全文コンテキストからも――無言で欠落していた。この欠落自体は当時から
# 認識されていたが、iter91(水平マージ重複排除の導入)・iter93(本文読み取り順序の
# 修正)はどちらも「ネストした表の再帰抽出は意図的にスコープ外(将来の
# フォローアップ課題)」と明記して先送りしており、直前のこのテストブロックは
# その「スコープ外である('InnerCell'が出力に現れない)」ことそのものを固定する
# 内容だった。本コミットでiter91/93のフォローアップを完了させたため、ここでは
# 従来の「スコープ外」アサーションを反転させ、ネスト表の内容が実際に出力される
# ことと、水平マージ重複排除(iter91)・本文順序抽出(iter93)の両方と正しく共存する
# ことを検証する。
with _tempfile.TemporaryDirectory() as _rdxn_dir:
    _rdxn_root = _pathlib.Path(_rdxn_dir)
    _rdxn_path = _rdxn_root / "nested.docx"

    if _HAS_DOCX_RD:
        _rdxn_doc = _docx_probe_rd.Document()
        _rdxn_doc.add_paragraph("Before")
        _rdxn_outer = _rdxn_doc.add_table(rows=1, cols=1)
        _rdxn_cell = _rdxn_outer.cell(0, 0)
        _rdxn_cell.text = "OuterCell"
        _rdxn_inner = _rdxn_cell.add_table(rows=1, cols=1)  # セル内にネストした表
        _rdxn_inner.cell(0, 0).text = "InnerCell"
        _rdxn_doc.add_paragraph("After")
        _rdxn_doc.save(str(_rdxn_path))

        _rdxn_exc = None
        try:
            _rdxn_out = f._read_docx(_rdxn_path)
        except Exception as _rdxn_e:
            _rdxn_exc = _rdxn_e
        check("_read_docx: セル内にネストした表があっても例外を送出しない",
              _rdxn_exc is None)
        if _rdxn_exc is None:
            check("_read_docx: ネストした表の内容('InnerCell')が出力に現れる"
                  "(修正前は'InnerCell'が出力に現れずスコープ外として欠落していた)",
                  "InnerCell" in _rdxn_out)
            check("_read_docx: 自身の段落テキストとネスト表を両方持つセル"
                  "('OuterCell')は二重出力されない",
                  _rdxn_out.count("OuterCell") == 1)
            check("_read_docx: ネスト表の内容('InnerCell')も二重出力されない",
                  _rdxn_out.count("InnerCell") == 1)
            _rdxn_lines = [_ln for _ln in _rdxn_out.split("\n") if _ln]
            check("_read_docx: ネスト表を含む文書でも本文順(Before->OuterCell->InnerCell->After)"
                  "が保たれ、ネスト行は自身を含む行の直後に追加される(行の途中に混在しない)",
                  _rdxn_lines == ["Before", "OuterCell", "InnerCell", "After"])
    else:
        print("   [SKIP] python-docx未インストールのため_read_docxネスト表テストをスキップ")

# ---------- _read_docx: ネスト表 x 水平マージ(gridSpan)の相互作用 (2026-07-25) ----------
# iter91の水平マージ重複排除は、被マージ座標を同一<w:tc>のidで検知してスキップする。
# ネスト表への再帰はそのdedupループの「初出セルのみ処理する」分岐の内側で行うため、
# マージされたセルにネスト表がある場合でも、被マージ座標側でネスト表が再度
# 処理される(内容が重複する)ことがないことを明示的に検証する。
with _tempfile.TemporaryDirectory() as _rdxhm_dir:
    _rdxhm_root = _pathlib.Path(_rdxhm_dir)
    _rdxhm_path = _rdxhm_root / "merged_nested.docx"

    if _HAS_DOCX_RD:
        _rdxhm_doc = _docx_probe_rd.Document()
        _rdxhm_table = _rdxhm_doc.add_table(rows=1, cols=2)
        _rdxhm_merged = _rdxhm_table.cell(0, 0).merge(_rdxhm_table.cell(0, 1))
        _rdxhm_merged.text = "MergedHeader"
        _rdxhm_inner = _rdxhm_merged.add_table(rows=1, cols=1)
        _rdxhm_inner.cell(0, 0).text = "MergedInner"
        _rdxhm_doc.save(str(_rdxhm_path))

        _rdxhm_out = f._read_docx(_rdxhm_path)
        _rdxhm_lines = [_ln for _ln in _rdxhm_out.split("\n") if _ln]
        check("_read_docx: 水平マージされたセルのネスト表内容('MergedInner')は"
              "被マージ座標の数だけ重複せず1回だけ出力される",
              _rdxhm_out.count("MergedInner") == 1)
        check("_read_docx: 水平マージ+ネスト表の出力は"
              "['MergedHeader', 'MergedInner']のまま(マージヘッダー直後にネスト行)",
              _rdxhm_lines == ["MergedHeader", "MergedInner"])
    else:
        print("   [SKIP] python-docx未インストールのため水平マージ+ネスト表テストをスキップ")

# ---------- _read_docx: ネスト表の行は「自身を含む行」の直後にまとめて追加される (2026-07-25) ----------
# 複数セルを持つ行の一部のセルだけにネスト表がある場合でも、ネスト表の行が
# 行の途中(セルとセルの間)に混在せず、行全体のタブ結合済みテキストの直後に
# まとめて追加されることを検証する(非ネスト表の出力がbyte-for-byte不変である
# ための配置規約)。
with _tempfile.TemporaryDirectory() as _rdxp_dir:
    _rdxp_root = _pathlib.Path(_rdxp_dir)
    _rdxp_path = _rdxp_root / "placement.docx"

    if _HAS_DOCX_RD:
        _rdxp_doc = _docx_probe_rd.Document()
        _rdxp_table = _rdxp_doc.add_table(rows=1, cols=2)
        _rdxp_table.cell(0, 0).text = "A"
        _rdxp_cell_b = _rdxp_table.cell(0, 1)
        _rdxp_cell_b.text = "B"
        _rdxp_inner = _rdxp_cell_b.add_table(rows=1, cols=1)
        _rdxp_inner.cell(0, 0).text = "Nested"
        _rdxp_doc.save(str(_rdxp_path))

        _rdxp_out = f._read_docx(_rdxp_path)
        _rdxp_lines = [_ln for _ln in _rdxp_out.split("\n") if _ln]
        check("_read_docx: 行の1セルのみにネスト表があっても行結合('A\\tB')は"
              "タブ結合のまま不変(セル間に割り込まない)",
              _rdxp_lines[0] == "A\tB")
        check("_read_docx: ネスト表の行はその行の直後に追加される(混在しない配置規約)",
              _rdxp_lines == ["A\tB", "Nested"])
    else:
        print("   [SKIP] python-docx未インストールのため行内配置テストをスキップ")

# ---------- _read_docx: ネスト表の再帰深さに上限があり病的な深いネストでも終了する (2026-07-25) ----------
# 自己参照的/異常に深いネスト構造でも再帰が終わらなくならないよう、_table_rows_text
# には固定の深さ上限(_DOCX_NESTED_TABLE_MAX_DEPTH)を設けている。深さ上限を超える
# 階層にあるネスト表は(安全側に倒し)無視されることを、上限を超える深さのチェーンを
# 実際に構築して検証する。テストがハングせず完走すること自体が上限の効果の証拠。
with _tempfile.TemporaryDirectory() as _rdxd_dir:
    _rdxd_root = _pathlib.Path(_rdxd_dir)
    _rdxd_path = _rdxd_root / "deep_nested.docx"

    if _HAS_DOCX_RD:
        _rdxd_doc = _docx_probe_rd.Document()
        _rdxd_table = _rdxd_doc.add_table(rows=1, cols=1)
        _rdxd_cell = _rdxd_table.cell(0, 0)
        _rdxd_cell.text = "L0"
        _RDXD_LEVELS = 8  # 深さ上限(6)より深いL7まで作り、打ち切りを検証する
        for _rdxd_lvl in range(1, _RDXD_LEVELS):
            _rdxd_inner = _rdxd_cell.add_table(rows=1, cols=1)
            _rdxd_cell = _rdxd_inner.cell(0, 0)
            _rdxd_cell.text = f"L{_rdxd_lvl}"
        _rdxd_doc.save(str(_rdxd_path))

        _rdxd_exc = None
        try:
            _rdxd_out = f._read_docx(_rdxd_path)
        except Exception as _rdxd_e:
            _rdxd_exc = _rdxd_e
        check("_read_docx: 深さ上限を超えるネスト表チェーンでも例外を送出せず完走する"
              "(再帰が終了しないことへの防御)",
              _rdxd_exc is None)
        if _rdxd_exc is None:
            for _rdxd_lvl in range(0, 7):  # L0..L6(深さ0..6)は上限内なので出力に現れる
                check(f"_read_docx: 深さ上限内のネスト表(L{_rdxd_lvl})は出力に現れる",
                      f"L{_rdxd_lvl}" in _rdxd_out)
            check("_read_docx: 深さ上限を超えた階層(L7、深さ7)は打ち切られ出力に現れない",
                  "L7" not in _rdxd_out)
    else:
        print("   [SKIP] python-docx未インストールのため深さ上限テストをスキップ")

check("_read_docx: ネスト表テスト群の後にsys.modulesの'docx'エントリが元通り解決可能(復元確認)",
      ("docx" not in sys.modules) or (sys.modules["docx"] is not None))

# ---------- _read_pdf: ImportError以外の実行時例外でもpypdf/PyPDF2へフォールスルー (2026-07-23 / iter83) ----------
# _read_pdf は pdfplumber -> pypdf -> PyPDF2 の順で試行するが、各ブロックは従来
# except ImportError のみで、下位ライブラリへのフォールスルーは「上位ライブラリが
# 未インストール」の場合にしか起きなかった。pdfplumber等が実際にはインストール
# 済みでも、暗号化/破損/パーサ固有のエッジケースで実行時に例外を送出するPDFに
# 対しては例外がそのまま _read_pdf の外へ伝播し、read_file_text(iter53)の
# 呼び出し側ガードがそれを握りつぶして""を返してしまい、pypdf/PyPDF2 なら
# 救えたはずのテキストがPDF丸ごとRAG/--fileコンテキストから失われていた
# （精度優先の方針に反する）。全リーダー関数の書き換えを試みたiter51は行き詰まった
# スタック案件だが、これは「呼び出し側でクラッシュさせない」話ではなく「下位
# ライブラリへフォールスルーしてテキストを救済する」話であり別角度の問題
# （iter41-44のgraceful degradation方針の延長）。
# ここではpdfplumber/pypdf/PyPDF2という実ライブラリを一切必要とせず、sys.modules
# にフェイクモジュールを注入することで(iter43/44/82と同じ手法)、インストール状態に
# 関わらず決定的に検証する。変更したsys.modulesエントリ・キャプチャしたstdoutは
# すべてtry/finallyで確実に復元する。
import pathlib as _rpdf_pathlib
import tempfile as _rpdf_tempfile


def _rpdf_make_reader_module(mod_name, page_texts):
    """pypdf/PyPDF2 と同一インタフェース(PdfReader(f).pages[i].extract_text())の
    フェイクモジュールを作る(実PDFパースは一切行わない、常に成功する版)。"""
    mod = types.ModuleType(mod_name)

    class _FakePage:
        def __init__(self, text):
            self._text = text

        def extract_text(self):
            return self._text

    class _FakePdfReader:
        def __init__(self, fileobj):
            self.pages = [_FakePage(t) for t in page_texts]

    mod.PdfReader = _FakePdfReader
    return mod


def _rpdf_make_raising_reader_module(mod_name, exc):
    """PdfReader構築時に指定した例外を送出するフェイクモジュール
    (importには成功するが実行時に失敗するケースを模擬)。"""
    mod = types.ModuleType(mod_name)

    class _FakePdfReader:
        def __init__(self, fileobj):
            raise exc

    mod.PdfReader = _FakePdfReader
    return mod


def _rpdf_swap_modules(entries, body):
    """entries: {module_name: fake_module_or_None} をsys.modulesへ差し替えてbody()を
    実行し、必ず元の状態(存在した/しなかった)へ復元する
    (iter43/44/82のswap-restoreパターンを複数モジュール分まとめて適用)。
    値がNoneのエントリはsys.modules[name]=Noneとなり、その名前のimport文に
    ImportErrorを送出させる(iter43/44で確立済みの標準手法)。"""
    originals = {}
    for name, fake in entries.items():
        originals[name] = sys.modules.get(name)
        sys.modules[name] = fake
    try:
        return body()
    finally:
        for name, orig in originals.items():
            if orig is not None:
                sys.modules[name] = orig
            else:
                del sys.modules[name]


def _rpdf_is_cp932_safe(s):
    """gotcha#4: Windowsコンソール(cp932)でエンコード不能な文字(絵文字等)を
    含んでいないかを確認する。"""
    try:
        s.encode("cp932")
        return True
    except UnicodeEncodeError:
        return False


with _rpdf_tempfile.TemporaryDirectory() as _rpdf_dir:
    _rpdf_path = _rpdf_pathlib.Path(_rpdf_dir) / "doc.pdf"
    # pypdf/PyPDF2ブロックは open(path, "rb") で実際にファイルを開くため
    # (中身はフェイクPdfReaderが解釈するので問われない)、ファイル自体は実在させる。
    _rpdf_path.write_bytes(b"%PDF-1.4 dummy content, not a real PDF")

    # (1) pdfplumberはimport可能だが.open()が実行時にRuntimeErrorを送出、
    #     pypdfはimport可能かつ正常にテキストを返す
    #     -> フォールスルーが働き、pypdfのテキストが失われずに返ってくること。
    def _rpdf_pdfplumber_open_raises(path):
        raise RuntimeError("simulated pdfplumber runtime failure (corrupt/encrypted PDF)")

    _fake_pdfplumber_raising = types.ModuleType("pdfplumber")
    _fake_pdfplumber_raising.open = _rpdf_pdfplumber_open_raises
    _fake_pypdf_ok = _rpdf_make_reader_module("pypdf", ["pypdf extracted text page1"])

    _rpdf_cap1 = io.StringIO()
    with contextlib.redirect_stdout(_rpdf_cap1):
        _rpdf_result1 = _rpdf_swap_modules(
            {"pdfplumber": _fake_pdfplumber_raising, "pypdf": _fake_pypdf_ok},
            lambda: f._read_pdf(_rpdf_path),
        )
    check("_read_pdf: pdfplumberが実行時RuntimeErrorを送出してもpypdfへフォールスルーしテキストを失わない",
          _rpdf_result1 == "pypdf extracted text page1")
    check("_read_pdf: pdfplumber実行時失敗の警告にファイル名(path.name)が出力される",
          _rpdf_path.name in _rpdf_cap1.getvalue())
    check("_read_pdf: pdfplumber実行時失敗の警告に例外型名(RuntimeError)が出力される",
          "RuntimeError" in _rpdf_cap1.getvalue())
    check("_read_pdf: 実行時失敗の警告メッセージはcp932でエンコード可能(絵文字等の非cp932文字を含まない、gotcha#4)",
          _rpdf_is_cp932_safe(_rpdf_cap1.getvalue()))

    # (2) 回帰: pdfplumberが未インストール(ImportError)でも、pypdfが正常にテキストを
    #     返せば従来通りpypdfの結果が返ること(フォールスルー自体は既存挙動)。
    _fake_pypdf_ok2 = _rpdf_make_reader_module("pypdf", ["pypdf text (pdfplumber absent)"])
    _rpdf_result2 = _rpdf_swap_modules(
        {"pdfplumber": None, "pypdf": _fake_pypdf_ok2},
        lambda: f._read_pdf(_rpdf_path),
    )
    check("_read_pdf(回帰): pdfplumber未インストール(ImportError)でもpypdfへフォールスルーする(既存挙動維持)",
          _rpdf_result2 == "pypdf text (pdfplumber absent)")

    # (3) 3層すべて失敗(ImportErrorと実行時例外が混在) -> 既存の
    #     '[PDF: {name} ... pip install pdfplumber]' 通知文字列をそのまま返すこと、
    #     かつ _is_lib_missing_notice がTrueと判定すること(RAGスキップ対象のまま)。
    def _rpdf_pdfplumber_open_raises3(path):
        raise RuntimeError("simulated pdfplumber failure (all-tiers-fail case)")

    _fake_pdfplumber_raising3 = types.ModuleType("pdfplumber")
    _fake_pdfplumber_raising3.open = _rpdf_pdfplumber_open_raises3
    _fake_pypdf2_raising = _rpdf_make_raising_reader_module(
        "PyPDF2", ValueError("simulated PyPDF2 failure (all-tiers-fail case)"))

    _rpdf_cap3 = io.StringIO()
    with contextlib.redirect_stdout(_rpdf_cap3):
        _rpdf_result3 = _rpdf_swap_modules(
            {"pdfplumber": _fake_pdfplumber_raising3, "pypdf": None, "PyPDF2": _fake_pypdf2_raising},
            lambda: f._read_pdf(_rpdf_path),
        )
    _rpdf_expected_notice = (
        f"[PDF: {_rpdf_path.name} — テキスト抽出には pdfplumber or pypdf が必要: pip install pdfplumber]"
    )
    check("_read_pdf: 3層すべて失敗(ImportError+実行時例外混在)なら既存の通知文字列をbyte-for-byteで返す",
          _rpdf_result3 == _rpdf_expected_notice)
    check("_is_lib_missing_notice: 3層すべて失敗時の通知文字列はTrueと判定される(RAGスキップ対象のまま)",
          f._is_lib_missing_notice(_rpdf_result3))
    check("_read_pdf: 3層すべて失敗時、pdfplumber(RuntimeError)とPyPDF2(ValueError)双方の"
          "実行時失敗警告が出力される",
          "RuntimeError" in _rpdf_cap3.getvalue() and "ValueError" in _rpdf_cap3.getvalue())

check("_read_pdf: テスト後にsys.modulesの'pdfplumber'エントリが元通り解決可能(復元確認)",
      ("pdfplumber" not in sys.modules) or (sys.modules["pdfplumber"] is not None))
check("_read_pdf: テスト後にsys.modulesの'pypdf'エントリが元通り解決可能(復元確認)",
      ("pypdf" not in sys.modules) or (sys.modules["pypdf"] is not None))
check("_read_pdf: テスト後にsys.modulesの'PyPDF2'エントリが元通り解決可能(復元確認)",
      ("PyPDF2" not in sys.modules) or (sys.modules["PyPDF2"] is not None))

# ---------- _read_pdf: 上位ライブラリが例外なしで空文字列を返した場合も下位へ
# フォールスルー (2026-07-25) ----------
# 直上のiter83は「実行時例外」に対してのみpypdf/PyPDF2へのフォールスルーを追加したが、
# iter83自身のコメントが明示する通りその修正は例外ケースに限定されており、
# pdfplumberが例外を送出せず単に空文字列/空白のみを返すケース（フォントエンコー
# ディング等に起因する既知の仕様上の癖で、実際には読める有効なPDFでも起こりうる）は
# 未対応のまま残っていた。fugu_local.py:783-784(pdfplumber.open()が成功し、結合した
# ページテキストが空)で即座に""を返して以降の層を一切試さないため、pypdf/PyPDF2なら
# 救えたはずのテキストがPDF丸ごとRAG/--fileコンテキストから静かに失われる
# （精度優先の方針に反する、iter83が閉じ損ねた「例外ではなく結果が空」という穴）。
# ここではiter41-44のgraceful degradation方針とiter83の例外フォールスルーの系譜を
# 継ぎ、「そのライブラリで抽出処理自体は完走したが結果が空だった」場合も次候補を
# 試す挙動を検証する。iter83と同じくpdfplumber/pypdf/PyPDF2という実ライブラリを
# 一切必要とせず、sys.modulesにフェイクモジュールを注入する(_rpdf_swap_modules/
# _rpdf_make_reader_moduleをそのまま再利用)。変更したsys.modulesエントリは
# すべてtry/finallyで確実に復元する。


def _rpdf_make_pdfplumber_module(page_texts):
    """pdfplumber.open(path) と同一インタフェース(with文 + .pages[i].extract_text())の
    フェイクモジュールを作る(実PDFパースは一切行わない、常に成功する版)。
    _rpdf_make_reader_module(pypdf/PyPDF2用)のpdfplumber版。"""
    mod = types.ModuleType("pdfplumber")

    class _FakePlumberPage:
        def __init__(self, text):
            self._text = text

        def extract_text(self):
            return self._text

    class _FakePlumberPdf:
        def __init__(self, texts):
            self.pages = [_FakePlumberPage(t) for t in texts]

        def __enter__(self):
            return self

        def __exit__(self, exc_type, exc_val, exc_tb):
            return False

    def _fake_open(path):
        return _FakePlumberPdf(page_texts)

    mod.open = _fake_open
    return mod


def _rpdf_make_never_called_module(mod_name):
    """importはできるが、PdfReader()が呼び出されたら即AssertionErrorになる
    フェイクモジュール。「上位層が成功したら下位層には一切触れない」ことを
    検証するためのカナリア(_read_pdfはexcept Exceptionで例外を握りつぶすため、
    このAssertionError自体は外へは伝播しないが、握りつぶされた結果として次候補へ
    フォールスルーしてしまい、期待した上位層のテキストとは異なる結果が返るため
    check()で検出できる)。"""
    mod = types.ModuleType(mod_name)

    class _NeverCalledPdfReader:
        def __init__(self, fileobj):
            raise AssertionError(
                f"{mod_name}.PdfReader が呼び出された(上位層成功時は下位層に触れないはず)")

    mod.PdfReader = _NeverCalledPdfReader
    return mod


with _rpdf_tempfile.TemporaryDirectory() as _rpdf2_dir:
    _rpdf2_path = _rpdf_pathlib.Path(_rpdf2_dir) / "doc.pdf"
    _rpdf2_path.write_bytes(b"%PDF-1.4 dummy content, not a real PDF")

    # (1) pdfplumberはimport可能・例外も出さないが抽出結果が空白のみ、
    #     pypdfはimport可能かつ正常にテキストを返す
    #     -> フォールスルーが働き、pypdfのテキストが失われずに返ってくること
    #     (これが今回追加した「結果が空」経路のリカバリケース)。
    _fake_pdfplumber_empty1 = _rpdf_make_pdfplumber_module(["", "   \n  "])
    _fake_pypdf_recovers1 = _rpdf_make_reader_module("pypdf", ["pypdf rescued text"])
    _rpdf2_result1 = _rpdf_swap_modules(
        {"pdfplumber": _fake_pdfplumber_empty1, "pypdf": _fake_pypdf_recovers1},
        lambda: f._read_pdf(_rpdf2_path),
    )
    check("_read_pdf: pdfplumberが例外なしで空文字列を返してもpypdfへフォールスルーしテキストを失わない",
          _rpdf2_result1 == "pypdf rescued text")

    # (2) pdfplumberは空、pypdfもimport可能だが空、PyPDF2はimport可能かつ正常
    #     -> 3層をまたいだフォールスルーチェーンが機能すること。
    _fake_pdfplumber_empty2 = _rpdf_make_pdfplumber_module([""])
    _fake_pypdf_empty2 = _rpdf_make_reader_module("pypdf", [None, ""])
    _fake_pypdf2_recovers2 = _rpdf_make_reader_module("PyPDF2", ["PyPDF2 rescued text"])
    _rpdf2_result2 = _rpdf_swap_modules(
        {"pdfplumber": _fake_pdfplumber_empty2, "pypdf": _fake_pypdf_empty2, "PyPDF2": _fake_pypdf2_recovers2},
        lambda: f._read_pdf(_rpdf2_path),
    )
    check("_read_pdf: pdfplumber/pypdfが空でもPyPDF2へフォールスルーしテキストを失わない(3層チェーン)",
          _rpdf2_result2 == "PyPDF2 rescued text")

    # (3) 3層すべてimport可能だが全て空(スキャンPDF等、本当に中身が無いケース)
    #     -> pip installの通知ではなく "" (空文字列)を返すこと。
    _fake_pdfplumber_empty3 = _rpdf_make_pdfplumber_module(["", None])
    _fake_pypdf_empty3 = _rpdf_make_reader_module("pypdf", [""])
    _fake_pypdf2_empty3 = _rpdf_make_reader_module("PyPDF2", [None, "  "])
    _rpdf2_result3 = _rpdf_swap_modules(
        {"pdfplumber": _fake_pdfplumber_empty3, "pypdf": _fake_pypdf_empty3, "PyPDF2": _fake_pypdf2_empty3},
        lambda: f._read_pdf(_rpdf2_path),
    )
    check("_read_pdf: 3層すべてimport可能だが全て空文字列/空白のみなら\"\"(空文字列)を返す(通知文字列ではない)",
          _rpdf2_result3 == "")
    check("_is_lib_missing_notice: 3層すべて空の\"\"はFalse(スキャン済PDF等を未インストール扱いしない)",
          not f._is_lib_missing_notice(_rpdf2_result3))

    # (4) 3層ともimport不可(全てImportError) -> 従来通りpip install通知を返すこと
    #     (「importできない」と「importできたが空」を混同していないことの回帰確認)。
    _rpdf2_result4 = _rpdf_swap_modules(
        {"pdfplumber": None, "pypdf": None, "PyPDF2": None},
        lambda: f._read_pdf(_rpdf2_path),
    )
    _rpdf2_expected_notice = (
        f"[PDF: {_rpdf2_path.name} — テキスト抽出には pdfplumber or pypdf が必要: pip install pdfplumber]"
    )
    check("_read_pdf: 3層ともImportErrorならpip install通知をbyte-for-byteで返す(空との混同なし)",
          _rpdf2_result4 == _rpdf2_expected_notice)
    check("_is_lib_missing_notice: 3層ともImportError時の通知文字列はTrue",
          f._is_lib_missing_notice(_rpdf2_result4))

    # (5) 回帰: pdfplumberが非空テキストを返した場合、pypdf/PyPDF2には一切触れず
    #     (importされたら即AssertionErrorになるフェイクで検証)、即座に
    #     pdfplumberのテキストをそのまま(未加工で)返すこと(共通経路は不変)。
    _fake_pdfplumber_ok5 = _rpdf_make_pdfplumber_module(["  non-empty pdfplumber text  "])
    _fake_pypdf_guard5 = _rpdf_make_never_called_module("pypdf")
    _fake_pypdf2_guard5 = _rpdf_make_never_called_module("PyPDF2")
    _rpdf2_result5 = _rpdf_swap_modules(
        {"pdfplumber": _fake_pdfplumber_ok5, "pypdf": _fake_pypdf_guard5, "PyPDF2": _fake_pypdf2_guard5},
        lambda: f._read_pdf(_rpdf2_path),
    )
    check("_read_pdf(回帰): pdfplumberが非空テキストを返せば即座にそれを未加工で返す(strip()されない、下位層は未使用)",
          _rpdf2_result5 == "  non-empty pdfplumber text  ")

check("_read_pdf: 空結果フォールスルーのテスト後にsys.modulesの'pdfplumber'エントリが元通り(復元確認)",
      ("pdfplumber" not in sys.modules) or (sys.modules["pdfplumber"] is not None))
check("_read_pdf: 空結果フォールスルーのテスト後にsys.modulesの'pypdf'エントリが元通り(復元確認)",
      ("pypdf" not in sys.modules) or (sys.modules["pypdf"] is not None))
check("_read_pdf: 空結果フォールスルーのテスト後にsys.modulesの'PyPDF2'エントリが元通り(復元確認)",
      ("PyPDF2" not in sys.modules) or (sys.modules["PyPDF2"] is not None))

# ---------- _read_excel: ImportError以外の実行時例外でもpandas/xlrdへフォールスルー (2026-07-24 / iter84) ----------
# _read_excel は openpyxl -> pandas/xlrd の順で試行するが、各ブロックは従来
# except ImportError のみで、下位ライブラリへのフォールスルーは「上位ライブラリが
# 未インストール」の場合にしか起きなかった。openpyxl は legacy な .xls
# （バイナリ形式）を一切読めず、openpyxl.load_workbook() は
# openpyxl.utils.exceptions.InvalidFileException という実行時例外
# （ImportError ではない）を送出する。read_file_text（L869付近）は .xlsx と .xls
# の両方をこの _read_excel へディスパッチするため、.xls は公式にサポートされる
# 入力形式でありながら、pandas+xlrd がインストール済みで読めるはずでもフォール
# スルーが起きず例外がそのまま外へ伝播し、read_file_text（iter53）/
# _load_rag_chunks（iter42）の呼び出し側ガードがそれを握りつぶしてスプレッド
# シート丸ごとがRAG/--fileコンテキストから静かに失われていた（精度優先の方針に
# 反する）。破損/openpyxl未対応の.xlsxがpandas側で再試行されない点も同じ隙間に
# 起因する。これは iter83 が _read_pdf に対して行った修正（下位ライブラリへ
# フォールスルーしてテキストを救済する）の直系の姉妹修正であり、iter41-44 の
# graceful degradation 方針の延長でもある。
# ここではopenpyxl/pandasという実ライブラリを一切必要とせず、sys.modules に
# フェイクモジュールを注入することで(iter43/44/82/83と同じ手法)、インストール
# 状態に関わらず決定的に検証する。多モジュールのswap-restoreはiter83で確立済みの
# _rpdf_swap_modules/_rpdf_is_cp932_safeをそのまま再利用する。変更した
# sys.modulesエントリ・キャプチャしたstdoutはすべてtry/finallyで確実に復元する。
import pathlib as _rxl2_pathlib
import tempfile as _rxl2_tempfile


class _Rxl2FakeInvalidFileException(Exception):
    """openpyxl.utils.exceptions.InvalidFileException を模した実行時例外
    (.xls等openpyxlが読めない形式で実際に送出される例外の代替。実際の例外クラスを
    使わずテスト側で定義することで、openpyxl未インストール環境でも決定的に
    テストできる)。"""
    pass


def _rxl2_make_openpyxl_raising_module(mod_name, exc):
    """openpyxl.load_workbook(...)呼び出し時に指定した例外を送出するフェイク
    モジュール(importには成功するが実行時に失敗するケースを模擬)。"""
    mod = types.ModuleType(mod_name)

    def _raising_load_workbook(*args, **kwargs):
        raise exc

    mod.load_workbook = _raising_load_workbook
    return mod


def _rxl2_make_pandas_module(mod_name, sheets, excelfile_raises=None):
    """pandas.ExcelFile(path).sheet_names / .parse(sheet).to_csv(index=False) と
    同一インタフェースのフェイクモジュールを作る(実Excelパースは一切行わない)。
    sheets: {シート名: csvテキスト} の辞書。excelfile_raisesを指定すると
    ExcelFile(...)コンストラクタ自体がその例外を送出する(pandas/xlrd層も実行時に
    失敗するケース、すなわち全層失敗ケースを模擬)。"""
    mod = types.ModuleType(mod_name)

    class _FakeDataFrame:
        def __init__(self, csv_text):
            self._csv_text = csv_text

        def to_csv(self, index=False):
            return self._csv_text

    class _FakeExcelFile:
        def __init__(self, path):
            if excelfile_raises is not None:
                raise excelfile_raises
            self.sheet_names = list(sheets.keys())

        def parse(self, sheet):
            return _FakeDataFrame(sheets[sheet])

    mod.ExcelFile = _FakeExcelFile
    return mod


with _rxl2_tempfile.TemporaryDirectory() as _rxl2_dir:
    _rxl2_path = _rxl2_pathlib.Path(_rxl2_dir) / "legacy.xls"
    # pandasブロックは pd.ExcelFile(str(path)) をフェイクモジュール経由で呼ぶだけで
    # 実際にバイト列を解釈しないため、中身は問われないがファイル自体は実在させる。
    _rxl2_path.write_bytes(b"\xd0\xcf\x11\xe0 dummy legacy .xls binary content, not real")

    # (1) openpyxlはimport可能だがload_workbook()が実行時にInvalidFileException相当を
    #     送出(.xls等openpyxlが読めない形式を想定)、pandasはimport可能かつ正常に
    #     テキストを返す -> フォールスルーが働き、pandas由来のテキストが失われずに
    #     返ってくること(受け入れ基準: fall-through/recovery テスト)。
    _rxl2_exc1 = _Rxl2FakeInvalidFileException(
        "simulated openpyxl InvalidFileException (legacy .xls / unsupported format)")
    _fake_openpyxl_raising1 = _rxl2_make_openpyxl_raising_module("openpyxl", _rxl2_exc1)
    _fake_pandas_ok1 = _rxl2_make_pandas_module("pandas", {"Sheet1": "name,age\nAlice,30\n"})

    _rxl2_cap1 = io.StringIO()
    with contextlib.redirect_stdout(_rxl2_cap1):
        _rxl2_result1 = _rpdf_swap_modules(
            {"openpyxl": _fake_openpyxl_raising1, "pandas": _fake_pandas_ok1},
            lambda: f._read_excel(_rxl2_path),
        )
    _rxl2_expected1 = "[Sheet: Sheet1]\nname,age\nAlice,30\n"
    check("_read_excel: openpyxlが実行時InvalidFileException相当を送出してもpandasへ"
          "フォールスルーしテキストを失わない(.xls想定)",
          _rxl2_result1 == _rxl2_expected1)
    check("_read_excel: openpyxl実行時失敗の警告にファイル名(path.name)が出力される",
          _rxl2_path.name in _rxl2_cap1.getvalue())
    check("_read_excel: openpyxl実行時失敗の警告に例外型名が出力される",
          "_Rxl2FakeInvalidFileException" in _rxl2_cap1.getvalue())
    check("_read_excel: 実行時失敗の警告メッセージはcp932でエンコード可能"
          "(絵文字等の非cp932文字を含まない、gotcha#4)",
          _rpdf_is_cp932_safe(_rxl2_cap1.getvalue()))

    # (2) openpyxl/pandas両方が実行時失敗 -> 既存の
    #     '[Excel: {name} ... pip install openpyxl]' 通知文字列をそのまま返すこと、
    #     かつ _is_lib_missing_notice がTrueと判定すること(RAGスキップ対象のまま、
    #     受け入れ基準: all tiers fail テスト)。例外が_read_excelの外へ伝播しない
    #     こと自体も、ここでcheck()が例外送出なく完了する形で確認される。
    _rxl2_exc2a = _Rxl2FakeInvalidFileException("simulated openpyxl failure (all-tiers-fail case)")
    _fake_openpyxl_raising2 = _rxl2_make_openpyxl_raising_module("openpyxl", _rxl2_exc2a)
    _fake_pandas_raising2 = _rxl2_make_pandas_module(
        "pandas", {}, excelfile_raises=ValueError("simulated pandas/xlrd failure (all-tiers-fail case)"))

    _rxl2_cap2 = io.StringIO()
    with contextlib.redirect_stdout(_rxl2_cap2):
        _rxl2_result2 = _rpdf_swap_modules(
            {"openpyxl": _fake_openpyxl_raising2, "pandas": _fake_pandas_raising2},
            lambda: f._read_excel(_rxl2_path),
        )
    _rxl2_expected_notice = f"[Excel: {_rxl2_path.name} — openpyxl or pandas が必要: pip install openpyxl]"
    check("_read_excel: openpyxl/pandas両方が実行時失敗なら既存の通知文字列をbyte-for-byteで返す",
          _rxl2_result2 == _rxl2_expected_notice)
    check("_is_lib_missing_notice: 両方実行時失敗時の通知文字列はTrueと判定される(RAGスキップ対象のまま)",
          f._is_lib_missing_notice(_rxl2_result2))
    check("_read_excel: 両方実行時失敗時、openpyxlとpandas双方の実行時失敗警告が出力される",
          "_Rxl2FakeInvalidFileException" in _rxl2_cap2.getvalue() and "ValueError" in _rxl2_cap2.getvalue())

    # (3) 回帰: openpyxlが未インストール(ImportError)でも、pandasが正常にテキストを
    #     返せば従来通りpandasの結果が返ること(フォールスルー自体は既存挙動、
    #     受け入れ基準: openpyxl-absent regression テスト)。
    _fake_pandas_ok3 = _rxl2_make_pandas_module("pandas", {"Data": "x,y\n1,2\n"})
    _rxl2_result3 = _rpdf_swap_modules(
        {"openpyxl": None, "pandas": _fake_pandas_ok3},
        lambda: f._read_excel(_rxl2_path),
    )
    check("_read_excel(回帰): openpyxl未インストール(ImportError)でもpandasへフォールスルーする"
          "(既存挙動維持)",
          _rxl2_result3 == "[Sheet: Data]\nx,y\n1,2\n")

check("_read_excel: テスト後にsys.modulesの'openpyxl'エントリが元通り解決可能(復元確認、iter84)",
      ("openpyxl" not in sys.modules) or (sys.modules["openpyxl"] is not None))
check("_read_excel: テスト後にsys.modulesの'pandas'エントリが元通り解決可能(復元確認、iter84)",
      ("pandas" not in sys.modules) or (sys.modules["pandas"] is not None))

# ---------- _read_excel: 変更後も成功パス(openpyxl実ワークブック)の出力が不変であることの確認 (iter84) ----------
# 上記の except Exception 追加はopenpyxlブロックの成功パス(try本体のreturn文)には
# 一切触れていないが、実際のopenpyxlの実行結果でも回帰していないことを直接確認する
# (受け入れ基準: success-path regression テスト。iter82のフィクスチャ手法を踏襲し、
# gc.collect()によるWindowsファイルハンドルロック回避も同様に行う)。
try:
    import openpyxl as _rxl2_openpyxl_probe
    _RXL2_HAS_OPENPYXL = True
except ImportError:
    _RXL2_HAS_OPENPYXL = False

if _RXL2_HAS_OPENPYXL:
    with _rxl2_tempfile.TemporaryDirectory() as _rxl2_ok_dir:
        _rxl2_ok_path = _rxl2_pathlib.Path(_rxl2_ok_dir) / "book_after_fix.xlsx"
        _rxl2_ok_wb = _rxl2_openpyxl_probe.Workbook()
        _rxl2_ok_ws = _rxl2_ok_wb.active
        _rxl2_ok_ws.title = "Data"
        _rxl2_ok_ws.append(["col1", "col2"])
        _rxl2_ok_ws.append(["v1", 1])
        _rxl2_ok_wb.save(str(_rxl2_ok_path))

        _rxl2_ok_out = f._read_excel(_rxl2_ok_path)
        # iter82と同じ理由(read_only=Trueワークブックの参照循環によるWindowsでの
        # 遅延ファイルハンドル解放)でgc.collect()を挟む(テスト専用の後始末)。
        import gc as _rxl2_gc
        _rxl2_gc.collect()
        _rxl2_ok_expected = "\n".join([
            "[Sheet: Data]",
            "col1\tcol2",
            "v1\t1",
        ])
        check("_read_excel: 変更後もopenpyxl成功パスの出力はbyte-for-byte従来通り(iter84回帰ガード)",
              _rxl2_ok_out == _rxl2_ok_expected)
else:
    print("   [SKIP] openpyxl未インストールのため_read_excel成功パス回帰テスト(iter84)をスキップ")

# ---------- _read_excel: 読み取り後にワークブック/ExcelFileをclose()する (2026-07-24 / iter88) ----------
# openpyxlブロック(load_workbook(read_only=True, data_only=True))とpandasブロック
# (pd.ExcelFile)はどちらも抽出後にハンドルを解放していなかった。read_only=True の
# ワークブックはワークシートとの間に参照循環を持つため単純な参照カウントだけでは
# 解放されず、iter82のテストはこれを直接踏んで、tempディレクトリのクリーンアップ時に
# gc.collect()を挟まないとWinError 32(他のプロセスがファイルを使用中)で失敗する
# ことを確認していた(=本番コードが実際にハンドルをリークしている直接証拠)。
# pandas.ExcelFileも内部でzipハンドルを開いたままにする。close()を呼ばないと
# Windowsでは同じファイルの後続の読み書き/移動をブロックし、_load_rag_chunksが
# 多数の.xlsxを読む場面ではハンドルリークにもなる。iter84はこの関数に
# 「ImportError以外の実行時例外でも下位ライブラリへフォールスルーする」修正を
# 入れた直系の姉妹修正であり、ここではその副作用(finally節の追加)がフォール
# スルー・警告文言・成功時テキストのいずれも変えていないことを併せて検証する。
# 実ライブラリのimport可否に関わらず決定的に検証するため、sys.modulesへの
# フェイクモジュール注入(iter43/44/82/83/84で確立済みの_rpdf_swap_modules)と、
# 実openpyxlが使える環境でのみ実施する成功パス/ハンドル解放の直接検証を併用する。
try:
    import openpyxl as _rxl3_real_openpyxl
    _RXL3_HAS_OPENPYXL = True
except ImportError:
    _RXL3_HAS_OPENPYXL = False

import pathlib as _rxl3_pathlib
import tempfile as _rxl3_tempfile
import os as _rxl3_os


class _Rxl3CloseCountingProxy:
    """実オブジェクト(openpyxl Workbook / pandas ExcelFile)をラップし、close()の
    呼び出し回数を記録するプロキシ。close以外の属性アクセスは__getattr__経由で
    全て実オブジェクトへ委譲するため、抽出処理そのものへの影響はない
    (close()呼び出し回数の検証専用、iter88)。"""

    def __init__(self, real_obj):
        self._rxl3_real = real_obj
        self.close_calls = 0

    def __getattr__(self, name):
        return getattr(self._rxl3_real, name)

    def close(self):
        self.close_calls += 1
        return self._rxl3_real.close()


def _rxl3_make_pandas_module_tracked(mod_name, sheets, tracker):
    """iter84の_rxl2_make_pandas_moduleにclose()呼び出し回数の記録を追加したもの
    (trackerは{"count": int}の辞書。呼び出し側で呼び出し後の値を検証する)。"""
    mod = types.ModuleType(mod_name)

    class _Rxl3FakeDataFrame:
        def __init__(self, csv_text):
            self._csv_text = csv_text

        def to_csv(self, index=False):
            return self._csv_text

    class _Rxl3FakeExcelFile:
        def __init__(self, path):
            self.sheet_names = list(sheets.keys())

        def parse(self, sheet):
            return _Rxl3FakeDataFrame(sheets[sheet])

        def close(self):
            tracker["count"] += 1

    mod.ExcelFile = _Rxl3FakeExcelFile
    return mod


# (A) openpyxl分岐: wb.close()が正確に1回呼ばれ、抽出テキストは不変であること
#     (受け入れ基準: openpyxl branch close-count テスト)。
if _RXL3_HAS_OPENPYXL:
    with _rxl3_tempfile.TemporaryDirectory() as _rxl3_dir_a:
        _rxl3_path_a = _rxl3_pathlib.Path(_rxl3_dir_a) / "close_check.xlsx"
        _rxl3_wb_a = _rxl3_real_openpyxl.Workbook()
        _rxl3_ws_a = _rxl3_wb_a.active
        _rxl3_ws_a.title = "Sheet1"
        _rxl3_ws_a.append(["name", "age"])
        _rxl3_ws_a.append(["Alice", 30])
        _rxl3_wb_a.save(str(_rxl3_path_a))

        _rxl3_captured_a = []

        def _rxl3_fake_load_workbook_a(*args, **kwargs):
            real_wb = _rxl3_real_openpyxl.load_workbook(*args, **kwargs)
            proxy = _Rxl3CloseCountingProxy(real_wb)
            _rxl3_captured_a.append(proxy)
            return proxy

        _rxl3_fake_openpyxl_mod_a = types.ModuleType("openpyxl")
        _rxl3_fake_openpyxl_mod_a.load_workbook = _rxl3_fake_load_workbook_a

        _rxl3_out_a = _rpdf_swap_modules(
            {"openpyxl": _rxl3_fake_openpyxl_mod_a},
            lambda: f._read_excel(_rxl3_path_a),
        )
        _rxl3_expected_a = "\n".join([
            "[Sheet: Sheet1]",
            "name\tage",
            "Alice\t30",
        ])
        check("_read_excel: openpyxl分岐でwb.close()が正確に1回呼ばれる(iter88)",
              len(_rxl3_captured_a) == 1 and _rxl3_captured_a[0].close_calls == 1)
        check("_read_excel: close()追加後もopenpyxl分岐の抽出テキストは従来通り(iter88)",
              _rxl3_out_a == _rxl3_expected_a)
else:
    print("   [SKIP] openpyxl未インストールのためwb.close()呼び出し検証(iter88)をスキップ")

# (B) pandas分岐: xl.close()が正確に1回呼ばれ、抽出テキストは正しいこと
#     (openpyxlは利用不可に見せかけてpandas分岐を強制する。受け入れ基準:
#     pandas branch close-count テスト)。
_rxl3_pandas_tracker_b = {"count": 0}
_rxl3_fake_pandas_mod_b = _rxl3_make_pandas_module_tracked(
    "pandas", {"Sheet1": "x,y\n1,2\n"}, _rxl3_pandas_tracker_b)

with _rxl3_tempfile.TemporaryDirectory() as _rxl3_dir_b:
    _rxl3_path_b = _rxl3_pathlib.Path(_rxl3_dir_b) / "pandas_close_check.xlsx"
    # pandas分岐はフェイクモジュール経由で呼ぶだけで実際にバイト列を解釈しないため
    # 中身は問われないが、ファイル自体は実在させる(iter84の手法を踏襲)。
    _rxl3_path_b.write_bytes(b"dummy xlsx bytes for pandas-branch close() test, not parsed for real")

    _rxl3_out_b = _rpdf_swap_modules(
        {"openpyxl": None, "pandas": _rxl3_fake_pandas_mod_b},
        lambda: f._read_excel(_rxl3_path_b),
    )
    check("_read_excel: pandas分岐でxl.close()が正確に1回呼ばれる(openpyxl不在, iter88)",
          _rxl3_pandas_tracker_b["count"] == 1)
    check("_read_excel: close()追加後もpandas分岐の抽出テキストは従来通り(iter88)",
          _rxl3_out_b == "[Sheet: Sheet1]\nx,y\n1,2\n")

# (C) 成功パスの回帰確認: 実openpyxlで複数シート/空行スキップ/None空文字化を含む
#     出力がbyte-for-byte従来通りであること、かつ_is_lib_missing_notice()がFalseに
#     なること(受け入れ基準: success-path regression テスト)。
if _RXL3_HAS_OPENPYXL:
    with _rxl3_tempfile.TemporaryDirectory() as _rxl3_dir_c:
        _rxl3_path_c = _rxl3_pathlib.Path(_rxl3_dir_c) / "regression.xlsx"
        _rxl3_wb_c = _rxl3_real_openpyxl.Workbook()
        _rxl3_ws1_c = _rxl3_wb_c.active
        _rxl3_ws1_c.title = "Sheet1"
        _rxl3_ws1_c.append(["name", "age"])
        _rxl3_ws1_c.append(["Alice", 30])
        _rxl3_ws1_c.append([None, None])   # 完全に空の行 -> 出力から省略される
        _rxl3_ws1_c.append(["Bob", None])  # Noneセルを含むが行自体は非空 -> 保持される
        _rxl3_ws2_c = _rxl3_wb_c.create_sheet("Sheet2")
        _rxl3_ws2_c.append(["x", "y"])
        _rxl3_ws2_c.append([1, 2])
        _rxl3_wb_c.save(str(_rxl3_path_c))

        _rxl3_out_c = f._read_excel(_rxl3_path_c)
        _rxl3_expected_c = "\n".join([
            "[Sheet: Sheet1]",
            "name\tage",
            "Alice\t30",
            "Bob\t",
            "[Sheet: Sheet2]",
            "x\ty",
            "1\t2",
        ])
        check("_read_excel: close()追加後もbyte-for-byte従来通りの抽出結果"
              "(複数シート/空行スキップ/None空文字化, iter88)",
              _rxl3_out_c == _rxl3_expected_c)
        check("_is_lib_missing_notice: close()追加後の成功時出力はFalse(iter37の過剰フィルタ回帰ガード, iter88)",
              not f._is_lib_missing_notice(_rxl3_out_c))
else:
    print("   [SKIP] openpyxl未インストールのため成功パス回帰検証(iter88)をスキップ")

# (D) ハンドル解放: _read_excel()から戻った直後、gc.collect()無しでも読み取った
#     .xlsxを即座にunlinkできること(受け入れ基準: handle-release テスト)。
#     Windows上でのみ意味のある検証のため、os.nameで判定してWindows以外では
#     スキップし、(A)/(B)のclose()呼び出し回数アサーションで代替する
#     (POSIXではread_onlyハンドルが残っていてもunlink自体は成功してしまうため、
#     unlinkの成否だけではハンドル解放の証明にならない)。
if _RXL3_HAS_OPENPYXL:
    with _rxl3_tempfile.TemporaryDirectory() as _rxl3_dir_d:
        _rxl3_path_d = _rxl3_pathlib.Path(_rxl3_dir_d) / "handle_release.xlsx"
        _rxl3_wb_d = _rxl3_real_openpyxl.Workbook()
        _rxl3_ws_d = _rxl3_wb_d.active
        _rxl3_ws_d.title = "S"
        _rxl3_ws_d.append(["k", "v"])
        _rxl3_wb_d.save(str(_rxl3_path_d))

        _rxl3_out_d = f._read_excel(_rxl3_path_d)
        check("_read_excel: ハンドル解放検証用フィクスチャの抽出結果も正しい(iter88)",
              _rxl3_out_d == "[Sheet: S]\nk\tv")

        if _rxl3_os.name == "nt":
            # 意図的にgc.collect()を呼ばない: close()単独でWindows上のファイル
            # ハンドルが即座に解放されることそのものを検証する対象のテスト
            # (iter82が必要としていたgc.collect()回避策がclose()追加により
            # 不要になったことの直接確認、iter88)。
            try:
                _rxl3_path_d.unlink()
                _rxl3_unlink_ok = True
            except OSError:
                _rxl3_unlink_ok = False
                # unlinkに失敗した場合、このwithブロックを抜ける際のTemporaryDirectory
                # クリーンアップがWinError 32で丸ごとクラッシュしてスイート全体を
                # 落とさないよう、iter82と同じ後始末を保険として行う(check()自体は
                # 既にFalseとして記録済みなので、この後始末は結果に影響しない)。
                import gc as _rxl3_gc_d
                _rxl3_gc_d.collect()
            check("_read_excel: Windows上でgc.collect()無しでも読み取り直後に"
                  "ファイルをunlinkできる(wb.close()によるハンドル即時解放, iter88)",
                  _rxl3_unlink_ok)
        else:
            print("   [SKIP] POSIX環境のためgc.collect()無しunlink検証(iter88)をスキップ"
                  "(close()呼び出し確認で代替)")
else:
    print("   [SKIP] openpyxl未インストールのためハンドル解放検証(iter88)をスキップ")

# (E) 抽出処理中(worksheets反復中)に例外が起きるケース: ワークブックは開けたが
#     parts構築中に失敗した場合でも、wb.close()は呼ばれ、かつiter84同様pandas側へ
#     フォールスルーしてテキストを失わないこと(受け入れ基準: mid-iteration failure
#     テスト)。openpyxl/pandasどちらも完全にフェイクなので実ライブラリの有無に
#     関わらず決定的に検証できる。
_rxl3_captured_e = []


class _Rxl3FakeWsRaising:
    title = "Sheet1"

    def iter_rows(self, values_only=True):
        raise RuntimeError("simulated mid-extraction failure after workbook opened (iter88)")


class _Rxl3FakeWbRaising:
    def __init__(self):
        self.worksheets = [_Rxl3FakeWsRaising()]
        self.close_calls = 0

    def close(self):
        self.close_calls += 1


def _rxl3_fake_load_workbook_e(*args, **kwargs):
    wb = _Rxl3FakeWbRaising()
    _rxl3_captured_e.append(wb)
    return wb


_rxl3_fake_openpyxl_mod_e = types.ModuleType("openpyxl")
_rxl3_fake_openpyxl_mod_e.load_workbook = _rxl3_fake_load_workbook_e

_rxl3_fake_pandas_mod_e = _rxl3_make_pandas_module_tracked(
    "pandas", {"Fallback": "a,b\n1,2\n"}, {"count": 0})

with _rxl3_tempfile.TemporaryDirectory() as _rxl3_dir_e:
    _rxl3_path_e = _rxl3_pathlib.Path(_rxl3_dir_e) / "mid_fail.xlsx"
    _rxl3_path_e.write_bytes(b"dummy content, not parsed for real (fully mocked)")

    _rxl3_cap_e = io.StringIO()
    with contextlib.redirect_stdout(_rxl3_cap_e):
        _rxl3_out_e = _rpdf_swap_modules(
            {"openpyxl": _rxl3_fake_openpyxl_mod_e, "pandas": _rxl3_fake_pandas_mod_e},
            lambda: f._read_excel(_rxl3_path_e),
        )
    check("_read_excel: 抽出処理中(worksheets反復中)の例外でもwb.close()は1回呼ばれる(iter88)",
          len(_rxl3_captured_e) == 1 and _rxl3_captured_e[0].close_calls == 1)
    check("_read_excel: 抽出処理中の例外でもiter84同様pandas側へフォールスルーしテキストを失わない(iter88)",
          _rxl3_out_e == "[Sheet: Fallback]\na,b\n1,2\n")
    check("_read_excel: 抽出処理中の例外の警告にファイル名と例外型名が出力される(iter88)",
          _rxl3_path_e.name in _rxl3_cap_e.getvalue() and "RuntimeError" in _rxl3_cap_e.getvalue())
    check("_read_excel: 抽出処理中例外の警告メッセージはcp932でエンコード可能(gotcha#4, iter88)",
          _rpdf_is_cp932_safe(_rxl3_cap_e.getvalue()))

check("_read_excel: テスト後にsys.modulesの'openpyxl'エントリが元通り解決可能(復元確認、iter88)",
      ("openpyxl" not in sys.modules) or (sys.modules["openpyxl"] is not None))
check("_read_excel: テスト後にsys.modulesの'pandas'エントリが元通り解決可能(復元確認、iter88)",
      ("pandas" not in sys.modules) or (sys.modules["pandas"] is not None))

# ---------- _read_excel: openpyxl.load_workbook が data_only=True/read_only=True で
# 呼ばれることの直接固定 (2026-07-26) ----------
# _read_excel(L1038)は openpyxl.load_workbook(str(path), read_only=True, data_only=True)
# を呼ぶ。data_only=True は精度に直結する既知の挙動固定であり、これは
# 数式セル(例 "=SUM(A1:A10)")を持つワークブックに対して、openpyxlが
# data_only=False(デフォルト)だと数式の"文字列そのもの"を返すのに対し、
# data_only=True だとExcelが最後に保存した時点でキャッシュした"計算済みの値"
# (例 12345)を返す、という違いを言う。_read_excelの抽出結果はread_file_text
# 経由で--fileのプロポーザーコンテキストに、また_load_rag_chunks経由でRAG
# コンテキストにそのまま混入するため、将来のリファクタでdata_only=Trueが
# 落とされる(openpyxlのデフォルトはFalse)と、数式を含む財務/データ系
# スプレッドシートで、計算済みの値の代わりに生の数式文字列がサイレントに
# LLMへ渡ってしまう(精度優先・時間は気にしないの方針に反する重大な劣化)。
# read_only=True も(iter82/88のコメント通り)意図した指定でありこれも併せて
# 固定する。
# grep上、data_only/read_only という語はこれまでこのテストファイル中の
# コメントと、フェイクload_workbookラッパー(例: 直前のiter88の
# _rxl3_fake_load_workbook_a/_e)にのみ出現しており、それらのラッパーは
# 実処理へフォワードするだけでkwargs自体は一切assertしていなかった。
# ここでは実openpyxlモジュールのload_workbook属性そのものを、
# 「本物のload_workbookを呼びつつ(args, kwargs)を記録して結果をそのまま
# 返すラッパー」でmonkeypatchし、_read_excelから実際に渡されたkwargsを
# 直接検証する(iter35がgotcha#1/#2の/api/chat・num_ctxピン留めを呼び出し
# パラメータのassertで固定したのと同じ手法)。openpyxlが利用できない環境
# でも全体PASSが崩れないよう、importの可否をtry/exceptで検出しガードする
# (iter82/84/88と同じスタイル)。
try:
    import openpyxl as _rxldo_openpyxl
    _RXLDO_HAS_OPENPYXL = True
except ImportError:
    _RXLDO_HAS_OPENPYXL = False

if _RXLDO_HAS_OPENPYXL:
    import pathlib as _rxldo_pathlib
    import tempfile as _rxldo_tempfile

    with _rxldo_tempfile.TemporaryDirectory() as _rxldo_dir:
        _rxldo_path = _rxldo_pathlib.Path(_rxldo_dir) / "data_only_check.xlsx"
        _rxldo_wb = _rxldo_openpyxl.Workbook()
        _rxldo_ws = _rxldo_wb.active
        _rxldo_ws.title = "Sheet1"
        _rxldo_ws.append(["a", "b"])
        _rxldo_ws.append([1, 2])
        _rxldo_wb.save(str(_rxldo_path))

        _rxldo_calls = []
        _rxldo_orig_load_workbook = _rxldo_openpyxl.load_workbook

        def _rxldo_recording_load_workbook(*args, **kwargs):
            """本物のopenpyxl.load_workbookを呼びつつ、_read_excelが渡した
            (args, kwargs)をそのまま記録するラッパー(呼び出しの転送のみで
            戻り値・副作用は本物と同一、抽出処理は一切変えない)。"""
            _rxldo_calls.append((args, kwargs))
            return _rxldo_orig_load_workbook(*args, **kwargs)

        _rxldo_openpyxl.load_workbook = _rxldo_recording_load_workbook
        try:
            _rxldo_out = f._read_excel(_rxldo_path)
        finally:
            _rxldo_openpyxl.load_workbook = _rxldo_orig_load_workbook

        check("_read_excel: openpyxl.load_workbookが正確に1回呼ばれる"
              "(data_only/read_only kwargs検証用フィクスチャ)",
              len(_rxldo_calls) == 1)
        _rxldo_kwargs = _rxldo_calls[0][1] if _rxldo_calls else {}
        check("_read_excel: openpyxl.load_workbook呼び出し時にdata_only=Trueが渡される"
              "(数式セルの計算済み値をRAG/--fileコンテキストへ渡すために必須。"
              "data_only=Falseだと数式文字列そのものが漏れる)",
              _rxldo_kwargs.get("data_only") is True)
        check("_read_excel: openpyxl.load_workbook呼び出し時にread_only=Trueが渡される",
              _rxldo_kwargs.get("read_only") is True)
        check("_read_excel: monkeypatch後もopenpyxl成功パスの抽出結果は従来通り"
              "(kwargs検証は抽出処理に影響しない)",
              _rxldo_out == "[Sheet: Sheet1]\na\tb\n1\t2")
        check("_read_excel: テスト後にopenpyxl.load_workbookが元の関数へ復元されている"
              "(post-restore identity確認)",
              _rxldo_openpyxl.load_workbook is _rxldo_orig_load_workbook)
else:
    print("   [SKIP] openpyxl未インストールのため_read_excel data_only/read_only"
          " kwargs検証テストをスキップ")

check("_read_excel: テスト後にsys.modulesの'openpyxl'エントリが元通り解決可能"
      "(復元確認、data_only/read_only kwargs検証)",
      ("openpyxl" not in sys.modules) or (sys.modules["openpyxl"] is not None))

# ---------- _tokenize / _score_chunk: 現行挙動の直接検証 ----------
# 2026-07-26: 以前は非ASCII連続列を丸ごと1トークンにしていたため、iter38では
# _tokenize('PINNについて') == {'pinn','について'} として「1トークンに固定化
# される」挙動そのものをロックしていた。しかしこれは、クエリの部分フレーズが
# チャンク側のより長い非ASCII連続列と完全一致しない限り _score_chunk の集合
# オーバーラップが常に0になる、という日本語RAG再現率ほぼゼロのバグを固定化
# していたに過ぎない。本プロジェクトの主要言語である日本語での RAG 再現率を
# 精度優先（時間は気にしない）の方針で改善するため、非ASCII連続列を隣接2文字
# の文字バイグラムへ分解するよう _tokenize を変更した（MeCab等の形態素解析器
# を使わないCJK情報検索の標準的対処、Lucene/ElasticsearchのCJKBigramFilterと
# 同様）。このassertionはその新挙動へ更新する（iter38の意図＝ASCIIとCJKを
# 混同しない、は維持したまま、CJK側の粒度だけを1連続列→文字バイグラムへ）。
check("_tokenize: ASCII+CJK混在を別トークンに分割し、CJK側は文字バイグラムに分解される(iter38改訂, 2026-07-26)",
      f._tokenize("PINNについて") == {"pinn", "につ", "つい", "いて"})
check("_tokenize: ASCII側の挙動は完全に不変(英語RAG回帰なし)",
      f._tokenize("apple pie recipe") == {"apple", "pie", "recipe"})
check("_tokenize: 日本語の部分フレーズがより長い連続列とバイグラムで重なる(再現率修正の核)",
      bool(f._tokenize("機械学習") & f._tokenize("本稿では機械学習の手法について述べる")))
check("_score_chunk: 空チャンクは0.0",
      f._score_chunk({"apple"}, "") == 0.0)
check("_score_chunk: 重複トークンありは正のスコア",
      f._score_chunk({"apple"}, "apple pie recipe") > 0.0)
check("_score_chunk: 重複トークンなしは0.0",
      f._score_chunk({"apple"}, "zebra mountain train") == 0.0)
check("_score_chunk: 日本語の部分フレーズがチャンク中のより長い文へ埋め込まれていても正のスコア(修正前は0.0)",
      f._score_chunk(f._tokenize("東京の人口"), "東京都の人口は約1400万人です") > 0.0)
check("_score_chunk: 内容を共有しない無関係な日本語チャンクは引き続き0.0(精度サニティ)",
      f._score_chunk(f._tokenize("東京の人口"), "桜が満開の京都で紅葉を楽しんだ") == 0.0)

# ---------- _get_rag_chunks: モジュールグローバルキャッシュの直接検証（2026-07-26） ----------
# _get_rag_chunks()(fugu_local.py ~L1531)は_load_rag_chunks()の前段にあるモジュール
# グローバルキャッシュで、rag_search -> _get_rag_chunks -> build_context という経路を
# 通じて全proposerに渡るRAGコンテキストの唯一の入口を握っている。挙動は
# 「dirs != _RAG_DIRS_LOADED を条件にした再読み込み」「list(dirs)によるコピー保持」
# 「キャッシュそのものを参照で返す」という純粋にステートフルなものだが、この
# ファイル中のrag_search関連テストは全て_get_rag_chunks自体をlambdaで丸ごと
# 差し替えて迂回しており(下のrag_searchセクション参照)、キャッシュ層そのものへの
# 直接テストがこれまで一件も存在しなかった。再読み込み条件がここで壊れると、
# (a) RAG_DIRSが変わっても古いチャンクを返し続ける(誤った文書がproposerに渡る
# 精度事故)、(b) 逆に毎ターン再読み込みしてしまい無駄なI/Oに加えiter42/156で
# ハードニングした1ファイル単位の隔離パスを毎回再発火させる、のどちらかの
# サイレントな精度劣化に直結する。ここではf._load_rag_chunksを呼び出し回数・
# 引数を記録するだけの純粋なin-processスタブに差し替え、f._get_rag_chunksを
# 直接呼び出してキャッシュ契約(再利用/無効化/コピーキー/空リストのキャッシュ)を
# 固定化する。ディスク・RAG_DIRSには一切触れない。printの文言は
# contextlib.redirect_stdoutで抑制し、文言自体はアサート対象にしない。
# 静的レビューでは欠陥は見つからなかったため、これは既存挙動の特性固定化
# (characterization)であり、意外な挙動が見つかってもここでは修正せず
# iteration 48/66/71のsurface-don't-fix方針を踏襲して明示するに留める
# (fugu_local.py自体はこの変更で一切変更していない)。

_orig_ggrc_chunks = f._RAG_CHUNKS
_orig_ggrc_dirs_loaded = f._RAG_DIRS_LOADED
_orig_ggrc_load_rag_chunks = f._load_rag_chunks

try:
    f._RAG_CHUNKS = []
    f._RAG_DIRS_LOADED = []

    _ggrc_calls = []      # 呼び出しごとに渡されたdirsの内容(コピー)を記録
    _ggrc_fixtures = []   # 呼び出しごとに返す固定チャンクリストのキュー(FIFO)

    def _ggrc_stub(dirs):
        _ggrc_calls.append(list(dirs))
        return _ggrc_fixtures.pop(0)

    f._load_rag_chunks = _ggrc_stub

    # --- 1) 初回ロード: キャッシュ未設定でdirsを渡すとスタブがちょうど1回呼ばれる ---
    _ggrc_dirs1 = ["dirA", "dirB"]
    _ggrc_fixture1 = [("dirA/x.txt", "chunk1"), ("dirB/y.txt", "chunk2")]
    _ggrc_fixtures.append(_ggrc_fixture1)
    with contextlib.redirect_stdout(io.StringIO()):
        _ggrc_result1 = f._get_rag_chunks(_ggrc_dirs1)
    check("_get_rag_chunks: 初回ロードでスタブがちょうど1回呼ばれる",
          len(_ggrc_calls) == 1)
    check("_get_rag_chunks: 初回ロードの戻り値はスタブの戻り値そのもの(同一オブジェクト)",
          _ggrc_result1 is _ggrc_fixture1)
    check("_get_rag_chunks: 初回ロード後_RAG_DIRS_LOADEDは渡したdirsの内容と一致",
          f._RAG_DIRS_LOADED == _ggrc_dirs1)

    # --- 2) 再利用: 内容が同じでも別オブジェクトのリストを渡すと再ロードされない ---
    _ggrc_dirs2 = list(_ggrc_dirs1)  # 別オブジェクト、内容は同一
    check("_get_rag_chunks(自己サニティ): dirs2はdirs1と別オブジェクト",
          _ggrc_dirs2 is not _ggrc_dirs1)
    with contextlib.redirect_stdout(io.StringIO()):
        _ggrc_result2 = f._get_rag_chunks(_ggrc_dirs2)
    check("_get_rag_chunks: 内容が同じ別オブジェクトのdirsでは再ロードされない(呼び出し回数1のまま)",
          len(_ggrc_calls) == 1)
    check("_get_rag_chunks: 再利用時は同一キャッシュオブジェクトを参照で返す(値一致ではなくis)",
          _ggrc_result2 is _ggrc_result1)

    # --- 3) 無効化(要素追加): dirsの内容が変わると再ロードされ、キャッシュが更新される ---
    _ggrc_dirs3 = ["dirA", "dirB", "dirC"]
    _ggrc_fixture3 = [("dirC/z.txt", "chunk3")]
    _ggrc_fixtures.append(_ggrc_fixture3)
    with contextlib.redirect_stdout(io.StringIO()):
        _ggrc_result3 = f._get_rag_chunks(_ggrc_dirs3)
    check("_get_rag_chunks: dirsに要素が増えると再ロードされる(呼び出し回数2)",
          len(_ggrc_calls) == 2)
    check("_get_rag_chunks: 無効化後の戻り値は新しいスタブの戻り値そのもの",
          _ggrc_result3 is _ggrc_fixture3)
    check("_get_rag_chunks: 無効化後_RAG_CHUNKSも新しいフィクスチャに更新される",
          f._RAG_CHUNKS is _ggrc_fixture3)
    check("_get_rag_chunks: 無効化後_RAG_DIRS_LOADEDも新しいdirsの内容に更新される",
          f._RAG_DIRS_LOADED == _ggrc_dirs3)

    # --- 3b) 無効化(順序違い): 同じ要素でも順序が異なれば別内容として再ロードされる ---
    _ggrc_dirs3b = ["dirB", "dirA", "dirC"]  # dirs3と要素は同じだが順序が異なる
    _ggrc_fixture3b = [("dirB/w.txt", "chunk3b")]
    _ggrc_fixtures.append(_ggrc_fixture3b)
    with contextlib.redirect_stdout(io.StringIO()):
        _ggrc_result3b = f._get_rag_chunks(_ggrc_dirs3b)
    check("_get_rag_chunks: 同じ要素でも順序が異なれば再ロードされる(呼び出し回数3)",
          len(_ggrc_calls) == 3)
    check("_get_rag_chunks: 順序違いによる無効化後の戻り値も新しいフィクスチャそのもの",
          _ggrc_result3b is _ggrc_fixture3b)

    # --- 4) コピーキー契約: ロード後に呼び出し元のdirsリストを変更しても、
    #        以後の「変更前と同内容の別オブジェクト」呼び出しは誤って再ロードを
    #        誘発しない(`_RAG_DIRS_LOADED = list(dirs)`でコピーを保持している
    #        契約のロック。もし単に`= dirs`なら下のcheckは呼び出し回数5で落ちる) ---
    _ggrc_dirs4 = ["dirX"]
    _ggrc_fixture4 = [("dirX/a.txt", "chunk4")]
    _ggrc_fixtures.append(_ggrc_fixture4)
    with contextlib.redirect_stdout(io.StringIO()):
        _ggrc_result4 = f._get_rag_chunks(_ggrc_dirs4)
    check("_get_rag_chunks(コピーキー契約, 準備): ロードされる(呼び出し回数4)",
          len(_ggrc_calls) == 4)
    _ggrc_dirs4.append("dirY_mutated_after_load")  # 呼び出し元のリストを事後に変更
    _ggrc_dirs4_fresh = ["dirX"]  # 変更前と同内容の、別の新しいリストオブジェクト
    with contextlib.redirect_stdout(io.StringIO()):
        _ggrc_result4b = f._get_rag_chunks(_ggrc_dirs4_fresh)
    check("_get_rag_chunks: ロード後に元のdirsリストを変更しても、"
          "同内容の別リストでの呼び出しは再ロードを誘発しない(コピー保持の証拠)",
          len(_ggrc_calls) == 4)
    check("_get_rag_chunks: コピーキー契約下でも再利用時は同一キャッシュオブジェクトを返す",
          _ggrc_result4b is _ggrc_result4)

    # --- 5) 空ロードのキャッシュ: スタブが空リストを返すディレクトリでも、以後の
    #        同内容呼び出しで再ロードされない(空/0ファイルディレクトリを毎ターン
    #        再スキャンしない契約のロック) ---
    _ggrc_dirs5 = ["dirEmpty"]
    _ggrc_fixtures.append([])
    with contextlib.redirect_stdout(io.StringIO()):
        _ggrc_result5 = f._get_rag_chunks(_ggrc_dirs5)
    check("_get_rag_chunks(空ロード, 準備): 空リストで初回ロードされる(呼び出し回数5)",
          len(_ggrc_calls) == 5)
    check("_get_rag_chunks: 空ディレクトリのロード結果は空リスト", _ggrc_result5 == [])
    _ggrc_dirs5_again = ["dirEmpty"]  # 別オブジェクト、内容は同一
    with contextlib.redirect_stdout(io.StringIO()):
        _ggrc_result5b = f._get_rag_chunks(_ggrc_dirs5_again)
    check("_get_rag_chunks: 空リストがキャッシュされ、同内容の再呼び出しでは再ロードされない",
          len(_ggrc_calls) == 5)
    check("_get_rag_chunks: 空リストキャッシュの再利用時も同一オブジェクトを返す(is)",
          _ggrc_result5b is _ggrc_result5)

finally:
    f._RAG_CHUNKS = _orig_ggrc_chunks
    f._RAG_DIRS_LOADED = _orig_ggrc_dirs_loaded
    f._load_rag_chunks = _orig_ggrc_load_rag_chunks

check("_get_rag_chunks: テスト後に_RAG_CHUNKSが元の状態(同一オブジェクト)へ復元されている",
      f._RAG_CHUNKS is _orig_ggrc_chunks)
check("_get_rag_chunks: テスト後に_RAG_DIRS_LOADEDが元の状態(同一オブジェクト)へ復元されている",
      f._RAG_DIRS_LOADED is _orig_ggrc_dirs_loaded)
check("_get_rag_chunks: テスト後にf._load_rag_chunksが元の状態へ復元されている",
      f._load_rag_chunks == _orig_ggrc_load_rag_chunks)

# ---------- rag_search: score>0のみ抽出（2026-07-22） ----------
# 以前は top = scored[:top_k] のうち先頭(best)が0でなければ丸ごと返しており、
# top_k内にキーワード的に無関係(score==0)なチャンクが混ざっていても
# そのままプロンプトへ注入されていた（rag_search ~L855-864）。
# score>0のみへの絞り込みが「best以外は素通し」の回帰を起こしていないか、
# また「関連チャンクを取りこぼしていない」かを、Ollama/ネットワーク一切なしで検証する。


_orig_get_rag_chunks = f._get_rag_chunks
try:
    # 1) queryは1件だけにマッチ、他はスコア0。top_k=2(>=2件)でも
    #    無関係チャンクの本文が混入しないこと。
    _chunks_partial = [
        ("dir/apple.txt", "This chunk talks about apple pie and recipe details."),
        ("dir/zebra.txt", "Completely unrelated content about zebras and mountains."),
        ("dir/cars.txt", "Another unrelated text about cars and trains."),
    ]
    f._get_rag_chunks = lambda dirs: _chunks_partial
    _res_partial = f.rag_search("apple recipe", dirs=["dummy"], top_k=2)
    check("rag_search: マッチしたチャンクのSourceが含まれる",
          "[Source: apple.txt]" in _res_partial)
    check("rag_search: スコア0チャンクの本文(zebra)は混入しない",
          "zebra" not in _res_partial and "mountains" not in _res_partial)
    check("rag_search: スコア0チャンクの本文(cars)は混入しない",
          "cars" not in _res_partial and "trains" not in _res_partial)

    # 2) 複数チャンクが全てスコア>0 -> 従来通りtop_k件まで降順・書式そのまま。
    # 注意: _score_chunk は _tokenize(chunk) を「集合」として重複除去してから
    # overlap/sqrt(len)+1 を計算するため、同じ語の反復回数はスコアに影響しない。
    # トークン集合の重複数と集合サイズを変えて意図的にスコアを分ける
    # （one: overlap2/len2≈82.8 > three: overlap1/len1=50.0 > two: overlap1/len2≈41.4）。
    # わざと入力順とスコア降順を食い違わせ、sort が実際に効いていることを検証する。
    _chunks_all_match = [
        ("dir/two.txt", "apple only"),      # overlap=1({apple}), len=2 -> ~41.4
        ("dir/three.txt", "recipe"),        # overlap=1({recipe}), len=1 -> 50.0
        ("dir/one.txt", "apple recipe"),    # overlap=2({apple,recipe}), len=2 -> ~82.8
    ]
    f._get_rag_chunks = lambda dirs: _chunks_all_match
    _res_all = f.rag_search("apple recipe", dirs=["dummy"], top_k=3)
    _expected_all = (
        "## Relevant Document Context (RAG)\n\n"
        "[Source: one.txt]\napple recipe"
        "\n\n---\n\n"
        "[Source: three.txt]\nrecipe"
        "\n\n---\n\n"
        "[Source: two.txt]\napple only"
    )
    check("rag_search: 全チャンクscore>0ならtop_k件・降順・書式が従来通り",
          _res_all == _expected_all)

finally:
    f._get_rag_chunks = _orig_get_rag_chunks

# best(=全件)が0スコアになるクエリで再検証(明示的に独立したtry/finallyで実施)
try:
    f._get_rag_chunks = lambda dirs: [
        ("dir/zebra.txt", "Completely unrelated content about zebras and mountains."),
        ("dir/cars.txt", "Another unrelated text about cars and trains."),
    ]
    check("rag_search: best(=全件)がscore0なら空文字（既存契約を維持）",
          f.rag_search("apple recipe", dirs=["dummy"], top_k=2) == "")
finally:
    f._get_rag_chunks = _orig_get_rag_chunks

# 4) 空dirs / 空チャンクリストは従来通り空文字
check("rag_search: dirsが空なら空文字",
      f.rag_search("apple", dirs=[]) == "")
try:
    f._get_rag_chunks = lambda dirs: []
    check("rag_search: チャンクリストが空なら空文字",
          f.rag_search("apple", dirs=["dummy"]) == "")
finally:
    f._get_rag_chunks = _orig_get_rag_chunks

# ---------- rag_search: 日本語クエリのEnd-to-Endバイグラム再現率検証（2026-07-26） ----------
# 修正前は非ASCII連続列が丸ごと1トークンだったため、クエリのキーフレーズが
# チャンク内のより長い連続文に埋め込まれているだけで一致しなくなり、rag_search は
# 現実的な日本語クエリ（'機械学習の応用'のような部分フレーズ）に対してほぼ常に
# ''（該当チャンクなし）を返していた。Ollama/ネットワーク一切なしで
# _get_rag_chunks のみモックし、「クエリのキーフレーズがチャンク内の長い文に
# 埋め込まれているケース」で実際に該当チャンクのSourceが返る（かつ内容を
# 共有しない無関係チャンクは混入しない）ことを検証する。
try:
    _chunks_ja = [
        ("dir/ml.txt", "本稿では機械学習の手法について詳しく説明する。"),
        ("dir/weather.txt", "今日の天気は晴れ時々曇りでした。"),
    ]
    f._get_rag_chunks = lambda dirs: _chunks_ja
    _res_ja = f.rag_search("機械学習の応用", dirs=["dummy"], top_k=2)
    check("rag_search: 日本語クエリの部分フレーズがチャンク内の長い文に埋め込まれていてもSourceが返る(修正前は'')",
          "[Source: ml.txt]" in _res_ja)
    check("rag_search: 内容を共有しない無関係な日本語チャンク(天気の話題)は混入しない",
          "天気" not in _res_ja and "曇り" not in _res_ja)
finally:
    f._get_rag_chunks = _orig_get_rag_chunks

# ---------- build_context / _with_context: Web検索+RAGのコンテキスト組み立てと注入 ----------
# build_context()(~L1104)と_with_context()(~L1125)は、Web検索(research_search)とRAG
# (rag_search)の結果を組み立てて全proposer/aggregatorへ、そしてmath/mcqでは
# solve_verifiableの自己一貫性投票(gotcha #7の精度クリティカル経路)にまで伝わる質問へ
# 注入する唯一の合流点だが、grep上コメントからの言及しかなく直接テストが皆無だった。
# ここではf.research_searchとf.rag_searchのみをモックし、パート順序・空ハンドリング
# (空検索がRAGを消さない)・use_search=Falseでの検索スキップ・結合セパレータの契約を
# 固定する。print文言そのものはアサート対象にしない(将来変わりうるため)。
# 万一モックが外れて実ネットワーク/subprocessに落ちないよう、urllib.request.urlopenと
# f.subprocess.runも「呼ばれたら即座にAssertionError」の番人に差し替える
# (gotcha #8 / iter 38・39・76のtripwire流儀を踏襲)。

_orig_research_search_bc = f.research_search
_orig_rag_search_bc = f.rag_search
_orig_urlopen_bc = urllib.request.urlopen
_orig_subprocess_run_bc = f.subprocess.run


def _bc_no_network_urlopen(*a, **k):
    raise AssertionError("build_context: モック漏れで実urlopen(ネットワーク)が呼ばれた")


def _bc_no_subprocess_run(*a, **k):
    raise AssertionError("build_context: モック漏れで実subprocess.runが呼ばれた")


try:
    urllib.request.urlopen = _bc_no_network_urlopen
    f.subprocess.run = _bc_no_subprocess_run

    # --- (a) use_search=False -> research_searchは一切呼ばれず、rag_search部分のみが返る ---
    _bc_a_search_calls = []

    def _bc_a_research_search(q):
        _bc_a_search_calls.append(q)
        return "SHOULD_NOT_APPEAR"

    try:
        f.research_search = _bc_a_research_search
        f.rag_search = lambda q, dirs=None, top_k=None: (
            "## Relevant Document Context (RAG)\n\nRAG_ONLY_BODY")
        with contextlib.redirect_stdout(io.StringIO()):
            _bc_a = f.build_context("Q_A", use_search=False, rag_dirs=["dummy"])
        check("build_context: use_search=Falseならresearch_searchは呼ばれない(呼び出し数0)",
              len(_bc_a_search_calls) == 0)
        check("build_context: use_search=Falseの戻り値はrag_search部分のみ",
              _bc_a == "## Relevant Document Context (RAG)\n\nRAG_ONLY_BODY")
    finally:
        f.research_search = _orig_research_search_bc
        f.rag_search = _orig_rag_search_bc

    # --- (b) use_search=Trueかつresearch_searchが非空 -> 検索ブロックを含み、RAGより前に来る ---
    try:
        f.research_search = lambda q: "## Web Search Results\n\nSEARCH_BODY"
        f.rag_search = lambda q, dirs=None, top_k=None: (
            "## Relevant Document Context (RAG)\n\nRAG_BODY")
        with contextlib.redirect_stdout(io.StringIO()):
            _bc_b = f.build_context("Q_B", use_search=True, rag_dirs=["dummy"])
        check("build_context: use_search=Trueで検索結果非空ならSEARCH_BODYを含む",
              "SEARCH_BODY" in _bc_b)
        check("build_context: use_search=Trueで検索結果非空ならRAG_BODYも含む",
              "RAG_BODY" in _bc_b)
        check("build_context: 検索ブロックはRAGブロックより前に来る(検索->RAGの順)",
              _bc_b.index("SEARCH_BODY") < _bc_b.index("RAG_BODY"))
    finally:
        f.research_search = _orig_research_search_bc
        f.rag_search = _orig_rag_search_bc

    # --- (c) use_search=Trueかつresearch_searchが''を返す -> 検索ブロックはないがRAGは残る ---
    try:
        f.research_search = lambda q: ""
        f.rag_search = lambda q, dirs=None, top_k=None: (
            "## Relevant Document Context (RAG)\n\nRAG_SURVIVES")
        with contextlib.redirect_stdout(io.StringIO()):
            _bc_c = f.build_context("Q_C", use_search=True, rag_dirs=["dummy"])
        check("build_context: 検索結果が空文字でもRAG結果は落とされない",
              _bc_c == "## Relevant Document Context (RAG)\n\nRAG_SURVIVES")
    finally:
        f.research_search = _orig_research_search_bc
        f.rag_search = _orig_rag_search_bc

    # --- (d) research_search・rag_searchとも'' -> build_contextはちょうど''を返す ---
    try:
        f.research_search = lambda q: ""
        f.rag_search = lambda q, dirs=None, top_k=None: ""
        with contextlib.redirect_stdout(io.StringIO()):
            _bc_d = f.build_context("Q_D", use_search=True, rag_dirs=["dummy"])
        check("build_context: 検索・RAGとも空文字ならちょうど空文字を返す", _bc_d == "")
    finally:
        f.research_search = _orig_research_search_bc
        f.rag_search = _orig_rag_search_bc

    # --- (e) 両パートがある場合、ちょうど'\n\n'で結合され検索->RAGの順である ---
    try:
        f.research_search = lambda q: "SEARCH_PART"
        f.rag_search = lambda q, dirs=None, top_k=None: "RAG_PART"
        with contextlib.redirect_stdout(io.StringIO()):
            _bc_e = f.build_context("Q_E", use_search=True, rag_dirs=["dummy"])
        check("build_context: 両パートはちょうど'\\n\\n'で結合され検索->RAG順",
              _bc_e == "SEARCH_PART\n\nRAG_PART")
    finally:
        f.research_search = _orig_research_search_bc
        f.rag_search = _orig_rag_search_bc

finally:
    f.research_search = _orig_research_search_bc
    f.rag_search = _orig_rag_search_bc
    urllib.request.urlopen = _orig_urlopen_bc
    f.subprocess.run = _orig_subprocess_run_bc

check("build_context: テスト後にresearch_searchが元の状態へ復元されている",
      f.research_search == _orig_research_search_bc)
check("build_context: テスト後にrag_searchが元の状態へ復元されている",
      f.rag_search == _orig_rag_search_bc)
check("build_context: テスト後にurllib.request.urlopenが元の状態へ復元されている",
      urllib.request.urlopen == _orig_urlopen_bc)
check("build_context: テスト後にsubprocess.runが元の状態へ復元されている",
      f.subprocess.run == _orig_subprocess_run_bc)

# --- _with_context: コンテキストの前置ロジック ---
_wc_question = "元の質問そのまま"
check("_with_context: 空文字コンテキストは質問をbyte-for-byteそのまま返す(同一オブジェクト)",
      f._with_context(_wc_question, "") is _wc_question)
check("_with_context: Noneコンテキスト(falsey)も質問をそのまま返す",
      f._with_context(_wc_question, None) is _wc_question)
check("_with_context: 非空コンテキストはcontext+'\\n\\n---\\n\\n'+questionを返す",
      f._with_context("質問本体", "CTX_BODY") == "CTX_BODY\n\n---\n\n質問本体")

# ---------- _save_as_html: コードフェンスの開始/終了タグ整合性 (2026-07-22) ----------
# 旧実装は開始/終了 ``` フェンスを両方とも "<pre><code>" にマップし、
# "</code></pre>" を一度も出力しないため <pre><code>...<pre><code> という
# 入れ子・未クローズの不整合HTMLになり、さらにコード本文が <br> 付きの
# 通常行として扱われて整形が崩れていた（L3135-3140）。ローカル一時ファイル
# への書き込み・読み戻しのみで検証し、Ollama/ネットワーク呼び出しは一切ない。
import pathlib as _html_pathlib

with _tempfile.TemporaryDirectory() as _html_dir:
    _html_root = _html_pathlib.Path(_html_dir)

    # (a) 単一のpythonフェンス付きコードブロック
    _html_out_a = _html_root / "a.html"
    _answer_a = "before\n```python\nx = 1\nprint(x)\n```\nafter"
    f._save_as_html(_html_out_a, "q1", _answer_a, 1.23)
    _content_a = _html_out_a.read_text(encoding="utf-8")
    check("_save_as_html: <pre><code>は正確に1回出現",
          _content_a.count("<pre><code>") == 1)
    check("_save_as_html: </code></pre>は正確に1回出現",
          _content_a.count("</code></pre>") == 1)
    check("_save_as_html: コード本文はescapeされている",
          "x = 1" in _content_a and "print(x)" in _content_a)
    _code_body_a = _content_a.split("<pre><code>", 1)[1].split("</code></pre>", 1)[0]
    check("_save_as_html: コード本文内に<br>が混入しない",
          "<br>" not in _code_body_a)

    # (b) プレーンテキストのみの回答: <br>維持・<pre>は出現しない
    _html_out_b = _html_root / "b.html"
    f._save_as_html(_html_out_b, "q2", "line1\nline2\nline3", 0.5)
    _content_b = _html_out_b.read_text(encoding="utf-8")
    check("_save_as_html: プレーンテキストは<br>で改行が維持される",
          _content_b.count("<br>") == 3)
    check("_save_as_html: プレーンテキストのみでは<pre>が出現しない",
          "<pre>" not in _content_b)

    # (c) 未終端フェンス（```が奇数個）でもバランスの取れた閉じタグになる
    _html_out_c = _html_root / "c.html"
    _answer_c = "intro\n```python\nx = 1\ny = 2\n"  # 閉じフェンスなし
    f._save_as_html(_html_out_c, "q3", _answer_c, 0.1)
    _content_c = _html_out_c.read_text(encoding="utf-8")
    check("_save_as_html: 未終端フェンスでも<pre><code>と</code></pre>の個数が一致",
          _content_c.count("<pre><code>") == _content_c.count("</code></pre>") == 1)

    # (d) 同一パスへ2回保存 -> 単一<body>へマージされ、両方のコードブロックがバランス
    _html_out_d = _html_root / "d.html"
    f._save_as_html(_html_out_d, "q4a", "```\ncode block one\n```", 0.1)
    f._save_as_html(_html_out_d, "q4b", "```\ncode block two\n```", 0.2)
    _content_d = _html_out_d.read_text(encoding="utf-8")
    check("_save_as_html: 2回保存しても<body>は1つにマージされる",
          _content_d.count("<body>") == 1 and _content_d.count("</body>") == 1)
    check("_save_as_html: 2回保存後も<pre><code>/</code></pre>の個数が一致(各2件)",
          _content_d.count("<pre><code>") == 2 and _content_d.count("</code></pre>") == 2)
    check("_save_as_html: 2回保存後も両方の回答本文が含まれる",
          "code block one" in _content_d and "code block two" in _content_d)

    # 回帰: コード本文中の '<' '&' がHTMLエスケープされ、生の '<b' 等が漏れない
    _html_out_e = _html_root / "e.html"
    _answer_e = "```\na < b && c\n```"
    f._save_as_html(_html_out_e, "q5", _answer_e, 0.1)
    _content_e = _html_out_e.read_text(encoding="utf-8")
    check("_save_as_html: コード本文の'<'/'&'がescapeされている",
          "a &lt; b &amp;&amp; c" in _content_e)
    check("_save_as_html: 生の(未escape)コード本文は出力に含まれない",
          "a < b && c" not in _content_e)

    # ---- (f) インデントされたフェンス (2026-07-23) ----
    # 旧実装は line.startswith("```") で列0固定だったため、番号付きリスト内に
    # 3スペースインデントで置かれた```pythonブロック（LLMがよく出す形）を
    # フェンスとして検出できず、```python/```自体がプレーン行としてescapeされ
    # コード本文も<br>付きの通常行として崩れて出力されていた。extract_boxed
    # （iteration 11）・strip_think（iteration 16）・_save_as_htmlのタグ整合性
    # （iteration 37）・_parse_slides（iteration 50）と同じフェンス未検出系
    # バグクラス。line.strip().startswith("```")への統一で解消したことを検証。
    _html_out_f = _html_root / "f.html"
    _answer_f = "1. Do this:\n   ```python\n   x = 1\n   print(x)\n   ```\n2. Done"
    f._save_as_html(_html_out_f, "q6", _answer_f, 0.2)
    _content_f = _html_out_f.read_text(encoding="utf-8")
    check("_save_as_html: インデントされたフェンスでも<pre><code>が正確に1回出現",
          _content_f.count("<pre><code>") == 1)
    check("_save_as_html: インデントされたフェンスでも</code></pre>が正確に1回出現",
          _content_f.count("</code></pre>") == 1)
    check("_save_as_html: インデントされたフェンス記号自体がエスケープされた文字列として残らない",
          "```python" not in _content_f and "```" not in _content_f)
    _code_body_f = _content_f.split("<pre><code>", 1)[1].split("</code></pre>", 1)[0]
    check("_save_as_html: インデントされたフェンスのコード本文に<br>が混入しない",
          "<br>" not in _code_body_f)
    check("_save_as_html: インデントされたフェンスのコード本文はそのまま含まれる",
          "x = 1" in _content_f and "print(x)" in _content_f)

    # ---- (g) インデントされた未終端フェンス（```が奇数個）でもバランスが取れる ----
    _html_out_g = _html_root / "g.html"
    _answer_g = "  - step:\n    ```python\n    x = 1\n"  # 閉じフェンスなし、インデント付き
    f._save_as_html(_html_out_g, "q7", _answer_g, 0.1)
    _content_g = _html_out_g.read_text(encoding="utf-8")
    check("_save_as_html: インデントされた未終端フェンスでも<pre><code>と</code></pre>の個数が一致",
          _content_g.count("<pre><code>") == _content_g.count("</code></pre>") == 1)

# ---------- _save_as_markdown/_save_as_text/_save_as_html: 既存ファイルが非UTF-8でも
# クラッシュしない (2026-07-23) ----------
# 3関数とも既存 --out ファイルの読み戻しに encoding="utf-8" のみを渡しており
# errors ハンドラが無かった（L3158/L3167/L3212）。このマシンのコンソールが
# cp932 であること(既知の落とし穴 #4)に起因し、既存ファイルが cp932/Shift_JIS
# 等の非UTF-8バイト列を含む場合に read_text が UnicodeDecodeError を送出し、
# _save_answer_to_file 経由で保存ステップ全体が異常終了して、計算し終えた
# 回答が失われていた。これは iteration 41-44 の _save_as_excel/_docx/_pdf
# 修正（保存を絶対にクラッシュさせず回答を失わない）と同じバグクラス。
# errors="replace" を既存内容の読み戻しにのみ追加し、書き込み側の
# encoding="utf-8" とUTF-8正常系の追記/マージ挙動は不変であることを検証する。
# ローカル一時ファイルの読み書きのみで検証し、Ollama/ネットワーク/bench呼び出しは一切ない。
with _tempfile.TemporaryDirectory() as _sv_dir:
    _sv_root = _html_pathlib.Path(_sv_dir)

    # --- _save_as_markdown: 既存ファイルがcp932(日本語)でも例外を送出しない ---
    _md_bad = _sv_root / "bad.md"
    _md_bad.write_bytes("日本語の既存メモ".encode("cp932"))
    try:
        f._save_as_markdown(_md_bad, "q_bad_md", "新しい回答md", 1.0, "")
        _md_bad_exc = None
    except Exception as _exc:
        _md_bad_exc = _exc
    check("_save_as_markdown: 既存ファイルがcp932でも例外を送出しない",
          _md_bad_exc is None)
    if _md_bad_exc is None:
        _md_bad_content = _md_bad.read_text(encoding="utf-8")
        check("_save_as_markdown: 非UTF-8既存ファイルでも新規回答が保存される",
              "新しい回答md" in _md_bad_content)
        try:
            _md_bad.read_text(encoding="utf-8")
            _md_bad_reread_ok = True
        except UnicodeDecodeError:
            _md_bad_reread_ok = False
        check("_save_as_markdown: 保存後ファイルはutf-8として問題なく再読込できる(残留未デコードバイトなし)",
              _md_bad_reread_ok)

    # --- _save_as_text: 既存ファイルが不正単独バイト列でも例外を送出しない ---
    _txt_bad = _sv_root / "bad.txt"
    _txt_bad.write_bytes(b"\x93\xfa\xff")
    try:
        f._save_as_text(_txt_bad, "q_bad_txt", "新しい回答txt", 0.5)
        _txt_bad_exc = None
    except Exception as _exc:
        _txt_bad_exc = _exc
    check("_save_as_text: 既存ファイルが不正単独バイトでも例外を送出しない",
          _txt_bad_exc is None)
    if _txt_bad_exc is None:
        _txt_bad_content = _txt_bad.read_text(encoding="utf-8")
        check("_save_as_text: 非UTF-8既存ファイルでも新規回答が保存される",
              "新しい回答txt" in _txt_bad_content)
        try:
            _txt_bad.read_text(encoding="utf-8")
            _txt_bad_reread_ok = True
        except UnicodeDecodeError:
            _txt_bad_reread_ok = False
        check("_save_as_text: 保存後ファイルはutf-8として問題なく再読込できる(残留未デコードバイトなし)",
              _txt_bad_reread_ok)

    # --- _save_as_html: 既存ファイルが非UTF-8(cp932)でも例外を送出しない ---
    _html_bad = _sv_root / "bad.html"
    _html_bad.write_bytes("<body>日本語</body>".encode("cp932"))
    try:
        f._save_as_html(_html_bad, "q_bad_html", "新しい回答html", 0.7)
        _html_bad_exc = None
    except Exception as _exc:
        _html_bad_exc = _exc
    check("_save_as_html: 既存ファイルがcp932でも例外を送出しない",
          _html_bad_exc is None)
    if _html_bad_exc is None:
        _html_bad_content = _html_bad.read_text(encoding="utf-8")
        check("_save_as_html: 非UTF-8既存ファイルでも新規回答が保存される",
              "新しい回答html" in _html_bad_content)
        try:
            _html_bad.read_text(encoding="utf-8")
            _html_bad_reread_ok = True
        except UnicodeDecodeError:
            _html_bad_reread_ok = False
        check("_save_as_html: 保存後ファイルはutf-8として問題なく再読込できる(残留未デコードバイトなし)",
              _html_bad_reread_ok)

    # --- 回帰: 既存ファイルが無い/有効UTF-8の場合は従来通りbyte-for-byte一致 ---

    # _save_as_markdown: 有効UTF-8の既存ファイルへの追記が従来通り
    _md_good = _sv_root / "good.md"
    _md_existing = "# 既存メモ\n\n有効なUTF-8の既存内容 täüst 日本語。\n\n"
    _md_good.write_text(_md_existing, encoding="utf-8")
    f._save_as_markdown(_md_good, "q_good_md", "有効な既存回答md", 2.5, "")
    _md_good_content = _md_good.read_text(encoding="utf-8")
    _md_ts_m = f.re.search(r"## Q \(([^)]+)\)", _md_good_content)
    check("_save_as_markdown: 追記後にtsが検出できる", _md_ts_m is not None)
    if _md_ts_m:
        _md_ts = _md_ts_m.group(1)
        _md_expected_block = (f"## Q ({_md_ts})\n\nq_good_md\n\n"
                              f"## A\n\n有効な既存回答md\n\n*所要: 2.5s*\n\n---\n\n")
        check("_save_as_markdown: 有効UTF-8既存ファイルへの追記はbyte-for-byte従来通り",
              _md_good_content == _md_existing + _md_expected_block)

    # _save_as_text: 有効UTF-8の既存ファイルへの追記が従来通り
    _txt_good = _sv_root / "good.txt"
    _txt_existing = "有効なUTF-8の既存内容 täüst 日本語。\n\n"
    _txt_good.write_text(_txt_existing, encoding="utf-8")
    f._save_as_text(_txt_good, "q_good_txt", "有効な既存回答txt", 1.5)
    _txt_good_content = _txt_good.read_text(encoding="utf-8")
    _txt_ts_m = f.re.search(r"^\[([^\]]+)\]", _txt_good_content[len(_txt_existing):])
    check("_save_as_text: 追記後にtsが検出できる", _txt_ts_m is not None)
    if _txt_ts_m:
        _txt_ts = _txt_ts_m.group(1)
        _txt_expected_block = (f"[{_txt_ts}]\nQ: q_good_txt\n\nA:\n有効な既存回答txt\n\n"
                               f"(所要 1.5s)\n{'='*60}\n\n")
        check("_save_as_text: 有効UTF-8既存ファイルへの追記はbyte-for-byte従来通り",
              _txt_good_content == _txt_existing + _txt_expected_block)

    # _save_as_html: 有効UTF-8の既存ファイル(2回保存の<body>マージ)が従来通り
    _html_good = _sv_root / "good.html"
    f._save_as_html(_html_good, "q_good_html1", "有効な既存回答html1", 0.3)
    _html_good_first = _html_good.read_text(encoding="utf-8")
    f._save_as_html(_html_good, "q_good_html2", "有効な既存回答html2", 0.4)
    _html_good_content = _html_good.read_text(encoding="utf-8")
    _html_first_body_m = f.re.search(r"<body>\n(.*)</body>", _html_good_first, f.re.DOTALL)
    _html_ts2_m = f.re.search(r"<h2>Q <small>\(([^)]+)\)</small></h2>\n<p>q_good_html2</p>",
                              _html_good_content)
    check("_save_as_html: 1回目保存のbodyが抽出できる", _html_first_body_m is not None)
    check("_save_as_html: 2回目保存のtsが検出できる", _html_ts2_m is not None)
    if _html_first_body_m and _html_ts2_m:
        _html_first_body = _html_first_body_m.group(1)
        _html_ts2 = _html_ts2_m.group(1)
        _html_second_body = (f"<h2>Q <small>({_html_ts2})</small></h2>\n<p>q_good_html2</p>\n"
                             f"<h2>A</h2>\n<p>有効な既存回答html2<br></p>\n"
                             f"<hr><p><small>所要: 0.4s</small></p>\n")
        # 実装(L3210-3218)は既存<body>...</body>間の中身をそのまま次回の
        # existing_body として引き継ぐ。1回目保存時点の内容は
        # "<body>\n" + body1 + "</body></html>" であり、<body>と</body>の
        # 間は "\n" + body1 なので、2回目保存の既存部分には <body> 直後の
        # 改行が引き継がれ "\n\n" + body1 になる（既存の仕様・quirkであり
        # 今回の修正では変更していない）。
        _html_expected = (f"<!DOCTYPE html>\n<html lang='ja'><head>"
                          f"<meta charset='UTF-8'><title>Fugu Output</title></head>\n"
                          f"<body>\n\n{_html_first_body}{_html_second_body}</body></html>")
        check("_save_as_html: 有効UTF-8既存ファイルへの2回目保存(<body>マージ)はbyte-for-byte従来通り",
              _html_good_content == _html_expected)

# ---------- _search_raw: import段階 vs 実行段階の ImportError 誤爆防止 (2026-07-23修正) ----------
# 旧実装は _ddg_full 内部の「ライブラリ import」と「DDGS().text()のクエリ実行」を
# 1つの try にまとめて except ImportError で受けていたため、ddgs/duckduckgo_search が
# 正しく import できていても、内部で遅延importされる primp/lxml 等のバックエンドが
# 実行時に ImportError を送出すると「ライブラリ未インストール」と誤認し、誤った
# pip install 警告を出した上で Instant Answer フォールバック（事実系クエリでほぼ空を
# 返し、古い知識で回答する事故の温床）に倒れていた。
# _resolve_ddgs_class() と _ddg_instant をモックし、実ネットワーク呼び出しなしで
# (1) 本当にライブラリが無い場合は従来通りInstant Answerへ縮退すること、
# (2) ライブラリ解決後の実行時ImportErrorはInstant Answerに縮退せず[]を返すこと、
# (3) 実行時の非ImportError例外も従来通り[]を返すこと、
# (4) 成功時の整形済み結果とweb_searchのヘッダ付与が従来通りであること、を検証する。

_orig_resolve_ddgs_sr = f._resolve_ddgs_class
_orig_ddg_instant_sr = f._ddg_instant


class _FakeDDGSCtx:
    """DDGS() as ... のコンテキストマネージャを模す最小限のフェイク。"""
    def __init__(self, text_fn):
        self._text_fn = text_fn

    def __enter__(self):
        return self

    def __exit__(self, exc_type, exc, tb):
        return False

    def text(self, query, max_results=None):
        return self._text_fn(query, max_results)


def _make_fake_ddgs_class(text_fn):
    def _factory(*a, **k):
        return _FakeDDGSCtx(text_fn)
    return _factory


try:
    # --- (1) 本当にライブラリ未インストール(解決段階でImportError) -> Instant Answerへ縮退 ---
    _sentinel_instant = ["INSTANT_ANSWER_SENTINEL"]

    def _resolve_raises_importerror():
        raise ImportError("ddgs / duckduckgo_search いずれも未インストール(模擬)")

    def _instant_returns_sentinel(query, max_results):
        return list(_sentinel_instant)

    f._resolve_ddgs_class = _resolve_raises_importerror
    f._ddg_instant = _instant_returns_sentinel

    _buf1 = io.StringIO()
    with contextlib.redirect_stdout(_buf1):
        _r1 = f._search_raw("ライブラリ未インストールクエリ", max_results=3)
    _out1 = _buf1.getvalue()

    check("_search_raw: ライブラリ未インストール時はInstant Answerの結果を返す(既存挙動)",
          _r1 == _sentinel_instant)
    check("_search_raw: ライブラリ未インストール時は既存のpip install警告を表示する(既存挙動)",
          "ddgs 未インストール" in _out1 and "pip install ddgs" in _out1)

    # --- (2) ライブラリ解決は成功、クエリ実行(.text())が実行時ImportErrorを送出 ---
    #     -> Instant Answerには倒れず(警告メッセージも出さず)、[]を返す。
    def _resolve_ok_for_runtime_importerror():
        return _make_fake_ddgs_class(
            lambda q, mr: (_ for _ in ()).throw(
                ImportError("primp/lxml 等バックエンドの遅延import失敗(模擬)")))

    def _instant_must_not_be_called(query, max_results):
        raise AssertionError(
            "_search_raw: 実行時ImportErrorなのにInstant Answerへ縮退した(誤爆再発)")

    f._resolve_ddgs_class = _resolve_ok_for_runtime_importerror
    f._ddg_instant = _instant_must_not_be_called

    _buf2 = io.StringIO()
    with contextlib.redirect_stdout(_buf2):
        _r2 = f._search_raw("実行時ImportErrorクエリ", max_results=3)
    _out2 = _buf2.getvalue()

    check("_search_raw: 実行時ImportErrorは[]を返す(Instant Answerへ誤爆しない)",
          _r2 == [])
    check("_search_raw: 実行時ImportErrorではpip install警告を表示しない",
          "ddgs 未インストール" not in _out2 and "pip install ddgs" not in _out2)
    check("_search_raw: 実行時ImportErrorは他の実行時エラーと同じ[Web検索エラー]表記",
          "[Web検索エラー:" in _out2)

    # --- (3) クエリ実行が非ImportErrorの通常例外を送出 -> 従来通り[]を返す ---
    def _resolve_ok_for_runtime_generic_error():
        return _make_fake_ddgs_class(
            lambda q, mr: (_ for _ in ()).throw(RuntimeError("ネットワーク断(模擬)")))

    f._resolve_ddgs_class = _resolve_ok_for_runtime_generic_error
    f._ddg_instant = _instant_must_not_be_called

    _buf3 = io.StringIO()
    with contextlib.redirect_stdout(_buf3):
        _r3 = f._search_raw("実行時RuntimeErrorクエリ", max_results=3)
    _out3 = _buf3.getvalue()

    check("_search_raw: 実行時の非ImportError例外も従来通り[]を返す(既存挙動)",
          _r3 == [])
    check("_search_raw: 実行時の非ImportError例外でもInstant Answerへ縮退しない(既存挙動)",
          "ddgs 未インストール" not in _out3)

    # --- (4) 正常系: 整形済み結果 + web_search のヘッダ付与が従来通り ---
    def _canned_rows(query, max_results):
        return [
            {"title": "結果1", "body": "本文1本文1本文1", "href": "https://example.com/1"},
            {"title": "結果2", "body": "本文2本文2本文2", "href": "https://example.com/2"},
        ]

    f._resolve_ddgs_class = lambda: _make_fake_ddgs_class(_canned_rows)
    f._ddg_instant = _instant_must_not_be_called

    _r4 = f._search_raw("正常系クエリ", max_results=5)
    _expected_r4 = [
        "[結果1]\n本文1本文1本文1\nSource: https://example.com/1",
        "[結果2]\n本文2本文2本文2\nSource: https://example.com/2",
    ]
    check("_search_raw: 正常系は_ddg_fullの整形形式(タイトル/本文/Source)通りの結果を返す",
          _r4 == _expected_r4)

    _ws4 = f.web_search("正常系クエリ", max_results=5)
    check("web_search: 正常系はDuckDuckGoヘッダを付与して結果を連結する(既存挙動)",
          _ws4 == "## Web Search Results (DuckDuckGo)\n" + "\n\n".join(_expected_r4))

    # --- (5) 解決段階(_resolve_ddgs_class())が非ImportError例外を送出 (2026-07-26追加) ---
    #     -> importが本当に失敗したわけではないので、Instant Answerへは誤爆させず、
    #     従来のクエリ実行時例外と同じく[]を返すだけにする(iter83の判断を踏襲)。
    def _resolve_raises_runtimeerror():
        raise RuntimeError("ddgs パッケージ__init__内でのクラッシュ(模擬・非ImportError)")

    f._resolve_ddgs_class = _resolve_raises_runtimeerror
    f._ddg_instant = _instant_must_not_be_called

    _buf5 = io.StringIO()
    with contextlib.redirect_stdout(_buf5):
        _r5 = f._search_raw("解決段階非ImportErrorクエリ", max_results=3)
    _out5 = _buf5.getvalue()

    check("_search_raw: 解決段階の非ImportErrorは例外を伝播させず[]を返す",
          _r5 == [])
    check("_search_raw: 解決段階の非ImportErrorはInstant Answerへ誤爆しない"
          "(pip install警告を出さない)",
          "ddgs 未インストール" not in _out5 and "pip install ddgs" not in _out5)
finally:
    f._resolve_ddgs_class = _orig_resolve_ddgs_sr
    f._ddg_instant = _orig_ddg_instant_sr

check("_search_raw: テスト後に_resolve_ddgs_classが元の状態へ復元されている",
      f._resolve_ddgs_class == _orig_resolve_ddgs_sr)
check("_search_raw: テスト後に_ddg_instantが元の状態へ復元されている",
      f._ddg_instant == _orig_ddg_instant_sr)

# ---------- _ddg_full: 壊れた行(非dict)・非文字列title/body/hrefの防御 (2026-07-25追加) ----------
# _ddg_full はプライマリ検索経路（ddgs/duckduckgo_search インストール済みの通常運用時に
# 必ず通る）でありながら、フォールバック側の _ddg_instant（イテレーション103/111/112/113/
# 138/139で「壊れた外部ペイロードの1件が全体を握り潰さない」よう isinstance ガードで段階的
# に固められてきた）とは非対称に、ddgs.text() が yield する各行 r に対して
# r.get("body")/r.get("title")/r.get("href") を型チェックなしで直接呼んでいた。行が dict
# でなければ即 AttributeError となり、_search_raw の `except Exception: return []` がクエリ
# 全体を空リストに丸め込むため、同じクエリ内の他の正常な行まで道連れにしていた
# （iter113/138と同じ「1件の破損が全体を握り潰す」問題）。ここでは
# _make_fake_ddgs_class/_FakeDDGSCtx（上のセクションで定義済み）を再利用し、
# f._resolve_ddgs_class をモック(try/finallyで復元)、さらに f._ddg_instant を
# 「呼ばれたらAssertionError」というトリップワイヤに差し替えて、フォールバック/実
# ネットワーク経路が一切使われないことも合わせて検証する。
_orig_resolve_ddgs_sr4 = f._resolve_ddgs_class
_orig_ddg_instant_sr4 = f._ddg_instant


def _ddg_instant_tripwire_sr4(query, max_results):
    raise AssertionError(
        "_ddg_full: 壊れた行の防御テスト中にフォールバック_ddg_instantが呼ばれた"
        "(プライマリ経路内で処理されるべき)")


try:
    f._ddg_instant = _ddg_instant_tripwire_sr4

    # --- (1) dict行と非dict行(str/int/None)が混在 -> 非dict行だけ読み飛ばし、
    #         有効な行は提出順のまま全て残る(1件の破損で全件を道連れにしない) ---
    def _mixed_rows_sr4(query, max_results):
        return [
            {"title": "T1", "body": "B1", "href": "https://example.com/1"},
            "こんな行はdictではない(壊れた形状)",
            {"title": "T2", "body": "B2", "href": "https://example.com/2"},
            12345,
            None,
            {"title": "T3", "body": "B3", "href": "https://example.com/3"},
        ]

    f._resolve_ddgs_class = lambda: _make_fake_ddgs_class(_mixed_rows_sr4)
    _expected_mixed_sr4 = [
        "[T1]\nB1\nSource: https://example.com/1",
        "[T2]\nB2\nSource: https://example.com/2",
        "[T3]\nB3\nSource: https://example.com/3",
    ]

    _r_mixed_sr4 = f._search_raw("非dict行混在クエリ", max_results=10)
    check("_search_raw: 非dict行(str/int/None)が混在しても例外を出さず、"
          "有効な行だけを提出順のまま全て返す(1件の破損で全件を道連れにしない)",
          _r_mixed_sr4 == _expected_mixed_sr4)

    # _ddg_full 自体を直接呼んでも(=_search_rawのexcept Exceptionに守られない状態でも)
    # 例外が伝播しないことを確認する。
    _exc_mixed_sr4 = None
    try:
        _r_mixed_direct_sr4 = f._ddg_full("非dict行混在クエリ", 10)
    except Exception as _e:
        _exc_mixed_sr4 = _e
        _r_mixed_direct_sr4 = None
    check("_ddg_full: 非dict行が混在しても関数自体が例外を送出しない(直接呼び出しで確認)",
          _exc_mixed_sr4 is None)
    check("_ddg_full: 直接呼び出しでも_search_raw経由と同じ結果(有効な行のみ提出順)",
          _r_mixed_direct_sr4 == _expected_mixed_sr4)

    # --- (2) dict行だがbodyが非文字列(list/dict/int) -> その行は残るがスニペットは
    #         ''になり、list/dictのrepr文字列はコンテキストへ混入しない ---
    for _bad_body, _label in [(["not", "a", "string"], "list"),
                               ({"nested": "dict"}, "dict"),
                               (12345, "int")]:
        def _bad_body_rows_sr4(query, max_results, _bad_body=_bad_body):
            return [{"title": "TB", "body": _bad_body, "href": "https://example.com/b"}]

        f._resolve_ddgs_class = lambda _fn=_bad_body_rows_sr4: _make_fake_ddgs_class(_fn)
        _r_bad_body_sr4 = f._search_raw(f"非文字列body({_label})クエリ", max_results=5)
        check(f"_ddg_full: bodyが非文字列({_label})でも行自体は残り、"
              f"スニペットは''になりrepr文字列は混入しない",
              _r_bad_body_sr4 == ["[TB]\n\nSource: https://example.com/b"])

    # --- (3) dict行だがtitleが非文字列(list) -> ''になり行は残る ---
    def _bad_title_rows_sr4(query, max_results):
        return [{"title": ["not", "a", "string"], "body": "Body text",
                 "href": "https://example.com/t"}]

    f._resolve_ddgs_class = lambda: _make_fake_ddgs_class(_bad_title_rows_sr4)
    _r_bad_title_sr4 = f._search_raw("非文字列titleクエリ", max_results=5)
    check("_ddg_full: titleが非文字列(list)でも行自体は残り、"
          "タイトルは''になりrepr文字列は混入しない",
          _r_bad_title_sr4 == ["[]\nBody text\nSource: https://example.com/t"])

    # --- (4) dict行だがhrefが非文字列(dict) -> ''になり行は残る ---
    def _bad_href_rows_sr4(query, max_results):
        return [{"title": "TH", "body": "Body text2", "href": {"nested": "dict"}}]

    f._resolve_ddgs_class = lambda: _make_fake_ddgs_class(_bad_href_rows_sr4)
    _r_bad_href_sr4 = f._search_raw("非文字列hrefクエリ", max_results=5)
    check("_ddg_full: hrefが非文字列(dict)でも行自体は残り、"
          "Source欄は''になりrepr文字列は混入しない",
          _r_bad_href_sr4 == ["[TH]\nBody text2\nSource: "])
finally:
    f._resolve_ddgs_class = _orig_resolve_ddgs_sr4
    f._ddg_instant = _orig_ddg_instant_sr4

check("_ddg_full: 壊れた行防御テスト後にf._resolve_ddgs_classが元の状態へ復元されている",
      f._resolve_ddgs_class == _orig_resolve_ddgs_sr4)
check("_ddg_full: 壊れた行防御テスト後にf._ddg_instantが元の状態へ復元されている",
      f._ddg_instant == _orig_ddg_instant_sr4)

# ---------- _ddg_instant: Abstract/RelatedTopics スニペットの長さ上限 (2026-07-24修正) ----------
# _ddg_full は (r.get("body") or "")[:WEB_SEARCH_SNIPPET_CHARS] で各スニペットを必ず
# 切り詰めるが、_ddg_instant (ddgs/duckduckgo_search 未インストール時のフォールバック。
# 上の 2026-07-23 セクションの通り維持対象の経路) は Abstract / RelatedTopics の Text を
# 無制限に追加していた。DuckDuckGo の Abstract は数KBに及ぶことがあり、research_search の
# 本文組み立てループで先頭アイテムが SEARCH_CONTEXT_CHARS を超えると
# `if not body: body = item[:SEARCH_CONTEXT_CHARS]` で break するため、巨大Abstract 1件が
# 他の収集済み事実を全て握り潰す事故につながる。ここでは urllib.request.urlopen のみを
# モックし、実ネットワーク呼び出しは一切発生させずに検証する。

_orig_urlopen_di = urllib.request.urlopen


def _fake_urlopen_di_payload(payload):
    def _fake(req, timeout=None):
        return _FakeHTTPResponse(json.dumps(payload).encode("utf-8"))
    return _fake


def _fake_urlopen_di_raises(exc):
    def _fake(req, timeout=None):
        raise exc
    return _fake


try:
    # --- (1) Abstractが上限を大幅に超える -> 本文のみ切り詰め、タイトル接頭辞と
    #     Source行(AbstractURL)はそのまま ---
    _long_abstract = "A" * (f.WEB_SEARCH_SNIPPET_CHARS * 3)
    urllib.request.urlopen = _fake_urlopen_di_payload({
        "Abstract": _long_abstract,
        "AbstractTitle": "巨大要約",
        "AbstractURL": "https://example.com/abstract",
        "RelatedTopics": [],
    })
    _di1 = f._ddg_instant("巨大Abstractクエリ", max_results=10)
    check("_ddg_instant: 巨大Abstractの本文はWEB_SEARCH_SNIPPET_CHARSに切り詰められ、"
          "タイトル接頭辞とSource行はそのまま",
          len(_di1) == 1 and _di1[0] == (
              f"[巨大要約]\n{_long_abstract[:f.WEB_SEARCH_SNIPPET_CHARS]}\n"
              f"Source: https://example.com/abstract"))
    check("_ddg_instant: 切り詰め後の本文長は正確にWEB_SEARCH_SNIPPET_CHARS",
          _di1[0].split("\n")[1] == "A" * f.WEB_SEARCH_SNIPPET_CHARS)

    # --- (2) RelatedTopicsのTextが上限超過 -> 切り詰められる(Source行は付与しない=既存の形状) ---
    _long_text = "B" * (f.WEB_SEARCH_SNIPPET_CHARS * 2)
    urllib.request.urlopen = _fake_urlopen_di_payload({
        "Abstract": "",
        "RelatedTopics": [{"Text": _long_text, "FirstURL": "https://example.com/rt"}],
    })
    _di2 = f._ddg_instant("巨大RelatedTopicsクエリ", max_results=10)
    check("_ddg_instant: 巨大RelatedTopics Textも上限に切り詰められる",
          _di2 == ["B" * f.WEB_SEARCH_SNIPPET_CHARS])

    # --- (3) 短いAbstract/Text(上限未満)はbyte-for-byteでそのまま(回帰なし) ---
    _short_abstract = "短い要約テキスト"
    _short_text = "短い関連トピック"
    urllib.request.urlopen = _fake_urlopen_di_payload({
        "Abstract": _short_abstract,
        "AbstractTitle": "短い要約",
        "AbstractURL": "https://example.com/short",
        "RelatedTopics": [{"Text": _short_text}],
    })
    _di3 = f._ddg_instant("短いクエリ", max_results=10)
    check("_ddg_instant: 短いAbstractはSource行も含め従来通り無変更(既存挙動)",
          _di3[0] == f"[短い要約]\n{_short_abstract}\nSource: https://example.com/short")
    check("_ddg_instant: 短いRelatedTopics Textも従来通り無変更(既存挙動)",
          _di3[1] == _short_text)

    # --- (4) max_resultsが従来通り尊重される(ループ早期break + 最終スライス、既存挙動) ---
    _many_topics = [{"Text": f"トピック{i}"} for i in range(20)]
    urllib.request.urlopen = _fake_urlopen_di_payload({
        "Abstract": "",
        "RelatedTopics": _many_topics,
    })
    _di4 = f._ddg_instant("多数RelatedTopicsクエリ", max_results=3)
    check("_ddg_instant: max_resultsで結果数が従来通り上限に収まる(既存挙動)",
          len(_di4) == 3)

    # --- (5) AbstractもRelatedTopicsも無い -> [](既存挙動) ---
    urllib.request.urlopen = _fake_urlopen_di_payload({
        "Abstract": "",
        "RelatedTopics": [],
    })
    _di5 = f._ddg_instant("空クエリ", max_results=5)
    check("_ddg_instant: AbstractもRelatedTopicsも無い場合は[]を返す(既存挙動)",
          _di5 == [])

    # --- (6) urlopen例外(JSON取得失敗相当) -> [](既存挙動) ---
    urllib.request.urlopen = _fake_urlopen_di_raises(RuntimeError("ネットワーク断(模擬)"))
    _di6 = f._ddg_instant("例外クエリ", max_results=5)
    check("_ddg_instant: urlopen例外時は[]を返す(既存挙動)",
          _di6 == [])

    # ---------- (7)-(11) RelatedTopicsのグループ化トピック1階層フラット化 (2026-07-24追加) ----------
    # DuckDuckGo Instant Answer APIのRelatedTopicsは、トップレベルにTextを持つ
    # 「直接トピック」と{"Name": "カテゴリ名", "Topics": [...]}形式の「グループ化
    # トピック」が混在する一覧である。従来はisinstance(t, dict) and t.get("Text")
    # のみを見ており、グループ化エントリ配下にネストされた事実が無条件で握り潰され
    # ていた(イテレーション85がこの関数を整備した際に残った既知の欠落)。以下は
    # その1階層フラット化(再帰はしない)を検証する。

    # --- (7) 直接トピックとグループ化トピックが混在 -> 順序を保ったまま全て展開 ---
    urllib.request.urlopen = _fake_urlopen_di_payload({
        "Abstract": "",
        "RelatedTopics": [
            {"Text": "t0"},
            {"Name": "X", "Topics": [{"Text": "t1"}, {"Text": "t2"}]},
        ],
    })
    _di7 = f._ddg_instant("混在RelatedTopicsクエリ", max_results=50)
    check("_ddg_instant: 直接トピックとグループ化トピックの混在は順序を保ったまま展開される",
          _di7 == ["t0", "t1", "t2"])

    # --- (8) ネストされたTextも上限に切り詰められる(直接トピックと同じ切り詰め) ---
    _long_nested_text = "C" * (f.WEB_SEARCH_SNIPPET_CHARS * 2)
    urllib.request.urlopen = _fake_urlopen_di_payload({
        "Abstract": "",
        "RelatedTopics": [
            {"Name": "Y", "Topics": [{"Text": _long_nested_text}]},
        ],
    })
    _di8 = f._ddg_instant("巨大ネストTextクエリ", max_results=10)
    check("_ddg_instant: グループ化トピック内のネストTextも上限に正確に切り詰められる",
          _di8 == ["C" * f.WEB_SEARCH_SNIPPET_CHARS])

    # --- (9) max_resultsがフラット化後のネストトピックにも従来通り適用される ---
    urllib.request.urlopen = _fake_urlopen_di_payload({
        "Abstract": "",
        "RelatedTopics": [
            {"Name": "Z", "Topics": [{"Text": f"nested{i}"} for i in range(20)]},
        ],
    })
    _di9 = f._ddg_instant("多数ネストトピッククエリ", max_results=3)
    check("_ddg_instant: max_resultsはフラット化後のネストトピックにも適用され、"
          "内側ループのbreakと最終スライスの両方が効く",
          _di9 == ["nested0", "nested1", "nested2"])

    # --- (10) 壊れた形状は例外を出さず読み飛ばす ---
    urllib.request.urlopen = _fake_urlopen_di_payload({
        "Abstract": "",
        "RelatedTopics": [
            {"Name": "A", "Topics": "not-a-list"},
            {"Name": "B", "Topics": []},
            {"Name": "C", "Topics": ["not-a-dict"]},
            {"Name": "D", "Topics": [{"NoText": "x"}]},
        ],
    })
    _di10 = f._ddg_instant("壊れた形状クエリ", max_results=10)
    check("_ddg_instant: 壊れた形状(Topicsがリストでない/空/非dict要素/Text欠落)は"
          "例外を出さずに読み飛ばされる",
          _di10 == [])

    # --- (11) トップレベルTextとTopicsを両方持つ場合はトップレベルTextのみ追加 ---
    urllib.request.urlopen = _fake_urlopen_di_payload({
        "Abstract": "",
        "RelatedTopics": [
            {"Text": "only-top", "Topics": [{"Text": "should-not-appear"}]},
        ],
    })
    _di11 = f._ddg_instant("Text優先クエリ", max_results=10)
    check("_ddg_instant: トップレベルTextとTopicsが両方ある場合はTextのみ追加され"
          "二重追加されない",
          _di11 == ["only-top"])

    # ---------- (12)-(15) 壊れたペイロード全体への耐性 (2026-07-24追加) ----------
    # data.get("RelatedTopics", []) は "RelatedTopics" キーが存在するが値が null の
    # 場合はdefaultが使われずNoneを返す(dict.getのdefaultは「キー不在」時のみ有効という
    # 仕様の落とし穴)ため、`for t in None` がTypeErrorになる。また DuckDuckGo Instant
    # Answer APIはクエリ次第でトップレベルがdictでないJSON(配列・文字列・数値)を返す
    # こともあり、その場合は data.get("Abstract") がAttributeErrorになる。_ddg_instant は
    # _search_raw() の except ImportError ブロックからtryに包まれず直接呼ばれる(L544相当)
    # ため、これらの例外は _search_raw の「失敗時は空リスト(呼び出し側を止めない)」契約
    # (イテレーション75)をすり抜けてターン全体を落としてしまう。以下は本イテレーションで
    # 追加した防御(dict以外のトップレベル->[]、非listなRelatedTopics->[]に丸め)を検証する。

    # --- (12) RelatedTopicsキーは存在するが値がnull -> []を返し例外なし ---
    urllib.request.urlopen = _fake_urlopen_di_payload({
        "Abstract": "",
        "RelatedTopics": None,
    })
    _di12 = f._ddg_instant("null-RelatedTopicsクエリ", max_results=5)
    check("_ddg_instant: RelatedTopicsが存在するがnullの場合は[]を返し例外を出さない",
          _di12 == [])

    # --- (13) トップレベルJSONがdictでない(配列/文字列/数値) -> []を返し例外なし ---
    for _bad_payload, _label in [
        ([{"Text": "x"}], "配列"),
        ("just-a-string", "文字列"),
        (12345, "数値"),
    ]:
        urllib.request.urlopen = _fake_urlopen_di_payload(_bad_payload)
        _di13 = f._ddg_instant(f"非dictペイロード({_label})クエリ", max_results=5)
        check(f"_ddg_instant: トップレベルJSONが{_label}(dict以外)の場合は"
              "[]を返し例外を出さない",
              _di13 == [])

    # --- (14) RelatedTopicsがdict(非list)でも、有効なAbstractは従来通り返る ---
    urllib.request.urlopen = _fake_urlopen_di_payload({
        "Abstract": "正常な要約",
        "AbstractTitle": "タイトル",
        "AbstractURL": "https://example.com/ok",
        "RelatedTopics": {"not": "a-list"},
    })
    _di14 = f._ddg_instant("非listRelatedTopics(dict)クエリ", max_results=5)
    check("_ddg_instant: RelatedTopicsがdict(非list)でもAbstractは従来通り返り例外なし",
          _di14 == ["[タイトル]\n正常な要約\nSource: https://example.com/ok"])

    # --- (15) RelatedTopicsが文字列(非list) -> トピック側は[]に丸められ例外なし ---
    urllib.request.urlopen = _fake_urlopen_di_payload({
        "Abstract": "",
        "RelatedTopics": "not-a-list-either",
    })
    _di15 = f._ddg_instant("非listRelatedTopics(文字列)クエリ", max_results=5)
    check("_ddg_instant: RelatedTopicsが文字列(非list)の場合も[]を返し例外を出さない",
          _di15 == [])
finally:
    urllib.request.urlopen = _orig_urlopen_di

check("_ddg_instant: テスト後にurllib.request.urlopenが元の状態へ復元されている",
      urllib.request.urlopen == _orig_urlopen_di)

# ---------- _search_raw: _ddg_instantフォールバック経路での壊れたペイロード伝播防止 (2026-07-24追加) ----------
# _search_raw は _resolve_ddgs_class() が ImportError を送出した場合、
# except ImportError: ブロックの中から `return _ddg_instant(...)` を直接呼ぶ(L544相当)。
# _ddg_full を包む try/except Exception のような保護がここには無いため、_ddg_instant が
# 上で追加した防御(非dictペイロード/null RelatedTopicsへの耐性)を持たないと、
# _search_raw の「失敗時は空リスト(呼び出し側を止めない)」契約(イテレーション75)が
# ここで破られる。_resolve_ddgs_class を ImportError へ、urllib.request.urlopen を
# 壊れたペイロードへそれぞれモックし(_ddg_instant自体はモックせず本物を経由させる)、
# 実ネットワーク呼び出しなしでこの伝播防止を検証する。

_orig_resolve_ddgs_sr2 = f._resolve_ddgs_class
_orig_urlopen_sr2 = urllib.request.urlopen

try:
    f._resolve_ddgs_class = _resolve_raises_importerror
    urllib.request.urlopen = _fake_urlopen_di_payload({"RelatedTopics": None})

    _buf_sr2 = io.StringIO()
    with contextlib.redirect_stdout(_buf_sr2):
        _r_sr2 = f._search_raw("壊れたペイロード経由クエリ", max_results=5)

    check("_search_raw: ライブラリ未インストール経由で_ddg_instantに渡った壊れたペイロード"
          "(RelatedTopics=null)も[]を返し例外を出さない(never-raise契約の維持)",
          _r_sr2 == [])
finally:
    f._resolve_ddgs_class = _orig_resolve_ddgs_sr2
    urllib.request.urlopen = _orig_urlopen_sr2

check("_search_raw: 壊れたペイロード伝播テスト後にurllib.request.urlopenが"
      "元の状態へ復元されている",
      urllib.request.urlopen == _orig_urlopen_sr2)
check("_search_raw: 壊れたペイロード伝播テスト後に_resolve_ddgs_classが"
      "元の状態へ復元されている",
      f._resolve_ddgs_class == _orig_resolve_ddgs_sr2)

# ---------- _ddg_instant: 末端(Abstract/Text)がtruthy非文字列の場合の防御 (2026-07-25追加) ----------
# イテレーション103はコンテナ形状(トップレベルdata・RelatedTopics)のみをisinstanceで
# 丸め、末端(Abstract文字列本体・直接トピックのText・ネストされたTopics[].Text)は
# `if data.get("Abstract"):` のような真偽値チェックの直後で無条件にスライスしていた。
# DDGは壊れたペイロードでText/Abstractにint/float/bool/list/dictを返すことがあり、
# int/float/boolはスライス不能のTypeError、dict/listは例外を出さない場合でも
# research_searchのre.search/'\n\n'.joinを壊す汚染データになる。以下はイテレーション
# 111/112/113と同じ「isinstanceで判定し例外を出さず読み飛ばす」方式の末端版を検証する。
# urllib.request.urlopenのみをモックし、実ネットワーク呼び出しは発生させない。

_orig_urlopen_di2 = urllib.request.urlopen
_BAD_LEAF_VALUES = [12345, 3.14, True, ["not", "a", "string"], {"bad": "dict"}]

try:
    # --- (16) Abstractがtruthy非文字列(int/float/bool/list/dict) -> Abstract項目のみ
    #     省略され、例外を出さず、兄弟のRelatedTopics Textは従来通り返る ---
    for _bad_abstract in _BAD_LEAF_VALUES:
        urllib.request.urlopen = _fake_urlopen_di_payload({
            "Abstract": _bad_abstract,
            "AbstractTitle": "無視されるタイトル",
            "AbstractURL": "https://example.com/bad-abstract",
            "RelatedTopics": [{"Text": "有効な兄弟トピック"}],
        })
        _di16 = f._ddg_instant(f"非文字列Abstract({type(_bad_abstract).__name__})クエリ",
                                max_results=10)
        check(f"_ddg_instant: Abstractがtruthy非文字列({type(_bad_abstract).__name__})でも"
              "例外を出さずAbstract項目のみ省略され、兄弟のTextは返る",
              _di16 == ["有効な兄弟トピック"])

    # --- (17) 直接トピックのTextがtruthy非文字列 -> そのリーフのみ省略され、
    #     兄弟の直接トピックは従来通り返る(direct-topic分岐自体は変わらず選択される) ---
    for _bad_text in _BAD_LEAF_VALUES:
        urllib.request.urlopen = _fake_urlopen_di_payload({
            "Abstract": "",
            "RelatedTopics": [
                {"Text": _bad_text},
                {"Text": "有効な直接トピック"},
            ],
        })
        _di17 = f._ddg_instant(f"非文字列直接Text({type(_bad_text).__name__})クエリ",
                                max_results=10)
        check(f"_ddg_instant: 直接トピックのTextがtruthy非文字列({type(_bad_text).__name__})"
              "でも例外を出さずそのリーフのみ省略され、兄弟の直接トピックは返る",
              _di17 == ["有効な直接トピック"])

    # --- (17b) 直接トピックのTextが非文字列で、かつ同じ要素にTopicsも併存する場合、
    #     direct-Text優先の分岐選択自体は変えない(Topicsへのフォールバックはしない)。
    #     既存のTextTopics優先(_di11相当)を壊れたTextでも保つことを確認する ---
    urllib.request.urlopen = _fake_urlopen_di_payload({
        "Abstract": "",
        "RelatedTopics": [
            {"Text": 999, "Topics": [{"Text": "should-not-appear"}]},
            {"Text": "good-sibling"},
        ],
    })
    _di17b = f._ddg_instant("非文字列Text優先クエリ", max_results=10)
    check("_ddg_instant: 直接トピックのTextが非文字列でもdirect-Text優先の分岐選択は"
          "変わらずTopicsへフォールバックしない(既存の優先順位を維持)",
          _di17b == ["good-sibling"])

    # --- (18) ネストされたTopics[].Textがtruthy非文字列 -> そのリーフのみ省略され、
    #     兄弟のネストTextは従来通り返る ---
    for _bad_nested in _BAD_LEAF_VALUES:
        urllib.request.urlopen = _fake_urlopen_di_payload({
            "Abstract": "",
            "RelatedTopics": [
                {"Name": "G", "Topics": [{"Text": _bad_nested}, {"Text": "有効なネスト"}]},
            ],
        })
        _di18 = f._ddg_instant(f"非文字列ネストText({type(_bad_nested).__name__})クエリ",
                                max_results=10)
        check(f"_ddg_instant: ネストTopics[].Textがtruthy非文字列"
              f"({type(_bad_nested).__name__})でも例外を出さずそのリーフのみ省略され、"
              "兄弟のネストTextは返る",
              _di18 == ["有効なネスト"])

    # --- (19) Abstract/直接Text/ネストTextの全てが同時に壊れていても、正常なリーフは
    #     全て保たれ、max_resultsの内側/外側breakと最終スライスも従来通り効く ---
    urllib.request.urlopen = _fake_urlopen_di_payload({
        "Abstract": ["not", "a", "string"],
        "AbstractTitle": "無視されるタイトル",
        "AbstractURL": "https://example.com/bad",
        "RelatedTopics": [
            {"Text": {"bad": "dict"}},
            {"Text": "good-direct-1"},
            {"Name": "grp", "Topics": [
                {"Text": 3.14}, {"Text": "good-nested-1"}, {"Text": "good-nested-2"},
            ]},
            {"Text": "good-direct-2"},
        ],
    })
    _di19 = f._ddg_instant("複合壊れ形状クエリ", max_results=3)
    check("_ddg_instant: Abstract/直接Text/ネストTextが同時に壊れていても、正常な"
          "リーフは順序を保ったまま返り、max_resultsの内側/外側breakも従来通り効く",
          _di19 == ["good-direct-1", "good-nested-1", "good-nested-2"])

    # --- (20) 回帰: 完全に整形済み(string Abstract/Text)なペイロードはbyte-for-byteで
    #     従来通り(Abstract+直接トピック+グループ化ネストトピックを1ペイロードに同居させ、
    #     [AbstractTitle]\nAbstract\nSource: URL の形とネスト展開順序を両方確認する) ---
    urllib.request.urlopen = _fake_urlopen_di_payload({
        "Abstract": "正常な要約文",
        "AbstractTitle": "正常タイトル",
        "AbstractURL": "https://example.com/normal",
        "RelatedTopics": [
            {"Text": "直接トピック1"},
            {"Name": "カテゴリ", "Topics": [{"Text": "ネストトピック1"}, {"Text": "ネストトピック2"}]},
        ],
    })
    _di20 = f._ddg_instant("完全整形ペイロードクエリ", max_results=10)
    check("_ddg_instant: 完全に整形済みのペイロードはisinstanceガード追加後もbyte-for-byteで"
          "従来通り(Abstract整形・直接トピック・ネストトピックの展開順序が全て一致)",
          _di20 == [
              "[正常タイトル]\n正常な要約文\nSource: https://example.com/normal",
              "直接トピック1",
              "ネストトピック1",
              "ネストトピック2",
          ])
finally:
    urllib.request.urlopen = _orig_urlopen_di2

check("_ddg_instant: 末端型ガードのテスト後にurllib.request.urlopenが"
      "元の状態へ復元されている",
      urllib.request.urlopen == _orig_urlopen_di2)

# ---------- _search_raw: _ddg_instant末端型ガード経由での壊れたペイロード伝播防止 (2026-07-25追加) ----------
# 上のセクションで_ddg_instant単体の末端型ガードを検証したのに続き、ここでは
# 2026-07-24追加分の「_search_rawはexcept ImportErrorの内側からtryに包まずに
# _ddg_instantを直接呼ぶ(L544相当)」テストと同じ構成で、末端(Abstract/Text)が
# truthy非文字列という壊れ方をする実際の経路(f._resolve_ddgs_classをImportErrorへ、
# urllib.request.urlopenを壊れたペイロードへモックし、_ddg_instant自体はモックせず
# 本物を経由させる)を通しても、_search_rawのnever-raise契約
# (「失敗時は空リスト、呼び出し側を止めない」・イテレーション75)が保たれることを
# 確認する。

_orig_resolve_ddgs_sr3 = f._resolve_ddgs_class
_orig_urlopen_sr3 = urllib.request.urlopen

try:
    f._resolve_ddgs_class = _resolve_raises_importerror
    urllib.request.urlopen = _fake_urlopen_di_payload({
        "Abstract": ["not", "a", "string"],
        "AbstractTitle": "T",
        "AbstractURL": "U",
        "RelatedTopics": [
            {"Text": {"bad": "dict"}},
            {"Text": "good-direct"},
            {"Name": "grp", "Topics": [{"Text": 12345}, {"Text": "good-nested"}]},
        ],
    })

    _buf_sr3 = io.StringIO()
    with contextlib.redirect_stdout(_buf_sr3):
        _r_sr3 = f._search_raw("末端型ガード経由クエリ", max_results=5)

    check("_search_raw: ライブラリ未インストール経由で_ddg_instantに渡った"
          "末端型が壊れたペイロード(Abstract/直接Text/ネストTextが非文字列)でも"
          "例外を出さず、正常なリーフのみのリストを返す(never-raise契約の維持)",
          _r_sr3 == ["good-direct", "good-nested"])
finally:
    f._resolve_ddgs_class = _orig_resolve_ddgs_sr3
    urllib.request.urlopen = _orig_urlopen_sr3

check("_search_raw: 末端型ガード伝播テスト後にurllib.request.urlopenが"
      "元の状態へ復元されている",
      urllib.request.urlopen == _orig_urlopen_sr3)
check("_search_raw: 末端型ガード伝播テスト後に_resolve_ddgs_classが"
      "元の状態へ復元されている",
      f._resolve_ddgs_class == _orig_resolve_ddgs_sr3)

# ---------- research_search: 反復リサーチループ (dedup / ラウンド上限 / 早期終了) ----------
# research_search は「Conductor/proposer 全員に注入される権威コンテキスト」を作る
# 精度クリティカルな経路だが、これまでオフラインテストが皆無だった。
# f._search_raw と f.ask の両方をモックし、実ネットワーク/Ollama呼び出しを一切発生させずに
# 分岐(Source URL重複排除・大小無視のクエリ重複排除・sufficient判定・空queries早期終了・
# SEARCH_MAX_ROUNDS上限)を検証する。extract_json は本物をそのまま使う(ask()の戻り値の
# 生文字列だけをモックする)。
#
# さらに、万一モックが外れて本物の _search_raw/ask 経由で実ネットワークに落ちないことを
# 保証するため、urllib.request.urlopen と f.subprocess.run も「呼ばれたら即座に
# AssertionError」の番人(センチネル)に差し替える(gotcha #8 の「bounded-loop違反は例外で
# 可視化する」流儀を踏襲)。

_orig_search_raw_rs = f._search_raw
_orig_ask_rs = f.ask
_orig_urlopen_rs = urllib.request.urlopen
_orig_subprocess_run_rs = f.subprocess.run


def _rs_no_network_urlopen(*a, **k):
    raise AssertionError("research_search: モック漏れで実urlopen(ネットワーク)が呼ばれた")


def _rs_no_subprocess_run(*a, **k):
    raise AssertionError("research_search: モック漏れで実subprocess.runが呼ばれた")


def _rs_search_factory(mapping, calls_log, max_calls):
    """mapping: {query: [item, ...]}。想定回数(max_calls)を超えたら例外(上限違反を可視化)。"""
    def _fake(query, max_results=None):
        calls_log.append(query)
        if len(calls_log) > max_calls:
            raise AssertionError(
                f"research_search: 想定回数({max_calls})を超えて_search_rawが呼ばれた"
                "(SEARCH_MAX_ROUNDS違反疑い)")
        return list(mapping.get(query, []))
    return _fake


def _rs_ask_factory(responses, calls_log):
    """responses: 十分性判定として順に返すdictのリスト。想定回数を超えたら例外。"""
    def _fake(model, messages, temperature, think=None, fmt=None, label=None,
              num_predict=None, num_ctx=None):
        calls_log.append(messages)
        idx = len(calls_log) - 1
        if idx >= len(responses):
            raise AssertionError(
                "research_search: 想定回数を超えてask()が呼ばれた(早期終了/ラウンド上限違反疑い)")
        return json.dumps(responses[idx])
    return _fake


try:
    urllib.request.urlopen = _rs_no_network_urlopen
    f.subprocess.run = _rs_no_subprocess_run

    # --- (A) Source-URL重複排除 と Sourceなし項目の先頭80文字重複排除 ---
    _itemA1 = "[T1]\nBody1 for source dedup test\nSource: http://example.com/a"
    _itemA2 = "No-source item padding text to reach eighty chars exactly xxxxxxxxxxxxxxxxxxxxxxxxx"
    _itemA3 = "[T3]\nBody3 brand new item\nSource: http://example.com/b"
    _searchA_calls = []
    _askA_calls = []
    try:
        f._search_raw = _rs_search_factory(
            {"Q_DEDUP": [_itemA1, _itemA2],
             "Q_DEDUP_R2": [_itemA1, _itemA2, _itemA3]},  # R2で同じ2件+新規1件を返す
            _searchA_calls, max_calls=2)
        f.ask = _rs_ask_factory(
            [{"sufficient": False, "missing": "x", "queries": ["Q_DEDUP_R2"]},
             {"sufficient": True, "missing": "", "queries": []}],
            _askA_calls)
        _resA = f.research_search("Q_DEDUP")
        check("research_search: Source URL重複は2巡目で再注入されない",
              _resA.count("http://example.com/a") == 1)
        check("research_search: Sourceなし項目は先頭80文字一致で重複排除される",
              _resA.count("No-source item padding text to reach eighty chars exactly") == 1)
        check("research_search: 新規Source項目はきちんと追加される",
              "http://example.com/b" in _resA)
        check("research_search: (A)_search_rawはR1・R2の2回のみ呼ばれる",
              _searchA_calls == ["Q_DEDUP", "Q_DEDUP_R2"])
    finally:
        f._search_raw = _orig_search_raw_rs
        f.ask = _orig_ask_rs

    # --- (B) 実行済みクエリの大小無視の重複排除(同一クエリの再検索防止) ---
    _itemB1 = "[B1]\nfirst round body\nSource: http://example.com/r1"
    _itemB2 = "[B2]\nsecond round body\nSource: http://example.com/r2"
    _searchB_calls = []
    _askB_calls = []
    try:
        f._search_raw = _rs_search_factory(
            {"Case Test Query": [_itemB1], "New Angle Query": [_itemB2]},
            _searchB_calls, max_calls=2)
        f.ask = _rs_ask_factory(
            [{"sufficient": False, "missing": "x",
              "queries": ["CASE TEST QUERY", "New Angle Query"]},  # 1つ目は既実行クエリの大小違い
             {"sufficient": True, "missing": "", "queries": []}],
            _askB_calls)
        _resB = f.research_search("Case Test Query")
        check("research_search: 既実行クエリは大小無視で再検索されない",
              _searchB_calls == ["Case Test Query", "New Angle Query"])
        check("research_search: 重複クエリをスキップしつつ新規クエリの結果は反映される",
              "http://example.com/r1" in _resB and "http://example.com/r2" in _resB)
    finally:
        f._search_raw = _orig_search_raw_rs
        f.ask = _orig_ask_rs

    # --- (C) sufficient=true で即座に早期終了(以降のラウンドの検索が発生しない) ---
    _itemC1 = "[C1]\nsufficient stop body\nSource: http://example.com/c1"
    _searchC_calls = []
    _askC_calls = []
    try:
        f._search_raw = _rs_search_factory({"Suff Stop Query": [_itemC1]},
                                            _searchC_calls, max_calls=1)
        f.ask = _rs_ask_factory([{"sufficient": True, "missing": "", "queries": []}],
                                 _askC_calls)
        _resC = f.research_search("Suff Stop Query")
        check("research_search: sufficient=trueで即座に早期終了(検索は1ラウンドのみ)",
              _searchC_calls == ["Suff Stop Query"])
        check("research_search: sufficient=trueなら判定は1回のみ呼ばれる",
              len(_askC_calls) == 1)
        check("research_search: 早期終了時もそのラウンドの結果は反映される",
              "http://example.com/c1" in _resC)
    finally:
        f._search_raw = _orig_search_raw_rs
        f.ask = _orig_ask_rs

    # --- (D) 空/欠落のqueriesで早期終了(以降のラウンドの検索が発生しない) ---
    _itemD1 = "[D1]\nempty queries stop body\nSource: http://example.com/d1"
    _searchD_calls = []
    _askD_calls = []
    try:
        f._search_raw = _rs_search_factory({"Empty Queries Query": [_itemD1]},
                                            _searchD_calls, max_calls=1)
        f.ask = _rs_ask_factory(
            [{"sufficient": False, "missing": "still missing", "queries": []}],
            _askD_calls)
        _resD = f.research_search("Empty Queries Query")
        check("research_search: 空queriesリストで早期終了(検索は1ラウンドのみ)",
              _searchD_calls == ["Empty Queries Query"])
        check("research_search: 空queries早期終了時もそのラウンドの結果は反映される",
              "http://example.com/d1" in _resD)
    finally:
        f._search_raw = _orig_search_raw_rs
        f.ask = _orig_ask_rs

    _itemD2 = "[D2]\nmissing queries key stop body\nSource: http://example.com/d2"
    _searchD2_calls = []
    _askD2_calls = []
    try:
        f._search_raw = _rs_search_factory({"Missing Key Query": [_itemD2]},
                                            _searchD2_calls, max_calls=1)
        # "queries" キー自体が欠落したJSON(j.get("queries") が None になるケース)
        f.ask = _rs_ask_factory([{"sufficient": False, "missing": "still missing"}],
                                 _askD2_calls)
        _resD2 = f.research_search("Missing Key Query")
        check("research_search: queriesキー欠落でも早期終了する(検索は1ラウンドのみ)",
              _searchD2_calls == ["Missing Key Query"])
    finally:
        f._search_raw = _orig_search_raw_rs
        f.ask = _orig_ask_rs

    # --- (E) SEARCH_MAX_ROUNDS上限: 常にsufficient=falseかつ新規クエリでも上限で打ち切り ---
    _searchE_calls = []
    _askE_calls = []
    try:
        f._search_raw = _rs_search_factory(
            {"Bound Round Query": ["[E1]\nr1\nSource: http://example.com/e1"],
             "Bound Q2": ["[E2]\nr2\nSource: http://example.com/e2"],
             "Bound Q3": ["[E3]\nr3\nSource: http://example.com/e3"]},
            _searchE_calls, max_calls=f.SEARCH_MAX_ROUNDS)
        f.ask = _rs_ask_factory(
            [{"sufficient": False, "missing": "x", "queries": ["Bound Q2"]},
             {"sufficient": False, "missing": "y", "queries": ["Bound Q3"]}],
            _askE_calls)  # ちょうど MAX_ROUNDS-1 回分しか用意しない(最終ラウンドは判定なし)
        _resE = f.research_search("Bound Round Query")
        check("research_search: 新規クエリが尽きなくてもSEARCH_MAX_ROUNDSでちょうど打ち切り",
              len(_searchE_calls) == f.SEARCH_MAX_ROUNDS == 3)
        check("research_search: 最終ラウンドでは十分性判定(ask)を呼ばない",
              len(_askE_calls) == f.SEARCH_MAX_ROUNDS - 1)
        check("research_search: 上限到達までの全ラウンドの結果が反映される",
              all(u in _resE for u in ("http://example.com/e1", "http://example.com/e2",
                                        "http://example.com/e3")))
    finally:
        f._search_raw = _orig_search_raw_rs
        f.ask = _orig_ask_rs

    # --- (F) 全ラウンドで結果ゼロ -> 空文字を返す ---
    _searchF_calls = []
    _askF_calls = []
    try:
        f._search_raw = _rs_search_factory({}, _searchF_calls, max_calls=f.SEARCH_MAX_ROUNDS)
        f.ask = _rs_ask_factory(
            [{"sufficient": False, "missing": "x", "queries": ["F Round2"]},
             {"sufficient": False, "missing": "y", "queries": ["F Round3"]}],
            _askF_calls)
        _resF = f.research_search("F Round1")
        check("research_search: 全ラウンドで結果ゼロなら空文字を返す", _resF == "")
        check("research_search: (F)検索はSEARCH_MAX_ROUNDS回実施された上での空文字",
              len(_searchF_calls) == f.SEARCH_MAX_ROUNDS)
    finally:
        f._search_raw = _orig_search_raw_rs
        f.ask = _orig_ask_rs

    # --- (G) 結果が1件でもあれば日付入りフレッシュネスヘッダーが付与される ---
    _searchG_calls = []
    _askG_calls = []
    try:
        f._search_raw = _rs_search_factory(
            {"Header Query": ["[G1]\nheader body\nSource: http://example.com/g1"]},
            _searchG_calls, max_calls=1)
        f.ask = _rs_ask_factory([{"sufficient": True, "missing": "", "queries": []}],
                                 _askG_calls)
        _resG = f.research_search("Header Query")
        _expected_date_g = f.time.strftime("%Y-%m-%d")
        check("research_search: 結果ありなら取得日入りヘッダーが付与される",
              _resG.startswith(f"## Web Search Results (取得日: {_expected_date_g})"))
        check("research_search: ヘッダーに本文(検索結果)が続く",
              "http://example.com/g1" in _resG)
    finally:
        f._search_raw = _orig_search_raw_rs
        f.ask = _orig_ask_rs

    # --- (H) 2026-07-22 修正済み挙動の固定化: 先頭(唯一の)結果が SEARCH_CONTEXT_CHARS を
    #     超えていても、body を空文字のまま break せず、先頭結果を上限まで切り詰めて必ず
    #     注入する(精度優先。sufficient=True と判定された唯一の具体的事実を黙って
    #     落とさない)。旧挙動(body="")はイテレーション38で特性テストとして固定されていたが、
    #     イテレーション39で修正した。
    _hugeH = "[Huge]\n" + ("Z" * (f.SEARCH_CONTEXT_CHARS + 500)) + "\nSource: http://example.com/huge"
    _searchH_calls = []
    _askH_calls = []
    try:
        f._search_raw = _rs_search_factory({"Huge Item Query": [_hugeH]},
                                            _searchH_calls, max_calls=1)
        f.ask = _rs_ask_factory([{"sufficient": True, "missing": "", "queries": []}],
                                 _askH_calls)
        _resH = f.research_search("Huge Item Query")
        _headerH = f"## Web Search Results (取得日: {f.time.strftime('%Y-%m-%d')})"
        _bodyH = _resH[len(_headerH):] if _resH.startswith(_headerH) else _resH
        check("research_search: 先頭項目がSEARCH_CONTEXT_CHARS超過でもheaderが付与される",
              _resH.startswith(_headerH))
        check("research_search: 先頭項目がSEARCH_CONTEXT_CHARS超過でもbodyが空にならない"
              "(切り詰めてでも必ず注入)",
              _bodyH.strip() != "")
        check("research_search: 切り詰められたbodyは先頭結果のプレフィックスを含む",
              _hugeH[:200] in _resH)
        check("research_search: 切り詰められたbodyの全文Sourceまでは含まれない"
              "(SEARCH_CONTEXT_CHARSで打ち切られている)",
              "http://example.com/huge" not in _resH)
    finally:
        f._search_raw = _orig_search_raw_rs
        f.ask = _orig_ask_rs

    # --- (I) 2026-07-24: "queries" が非list(int/float/bool/dict/str/None)でも
    #     TypeErrorを送出せず、追加ラウンドの検索を発行しない(=[]として扱われ即座に
    #     早期終了する)ことを固定化する。従来の `j.get("queries") or []` は
    #     "queries" が falsy(None/[]/""/0等)の場合のみ[]に丸めるtruthinessトリック
    #     で、真値だが非反復可能(non-iterable)なint/float/bool/dictはそのまま通り
    #     `for x in ...` がTypeErrorを送出していた(未捕捉のままbuild_context経由で
    #     ask_fugu呼び出し元まで伝播しターン全体を落とす)。また反復可能だが意図しない
    #     truthy str/dictは1文字ずつ/キーごとに分解された無意味なクエリを発行して
    #     いた。イテレーション103の_ddg_instant(非list RelatedTopics)・イテレーション
    #     111のplan_pptx_images(非list images)と同じisinstance判定への修正を検証する。
    _non_list_queries_cases = [
        ("int", 5),
        ("float", 3.14),
        ("bool_true", True),
        ("dict", {"a": "b"}),
        ("str", "not a list of strings"),
        ("none_explicit", None),
    ]
    for _case_name, _bad_queries in _non_list_queries_cases:
        _qname = f"NonList Query {_case_name}"
        _itemI = (f"[I-{_case_name}]\nnon-list queries body\n"
                  f"Source: http://example.com/i-{_case_name}")
        _searchI_calls = []
        _askI_calls = []
        try:
            f._search_raw = _rs_search_factory({_qname: [_itemI]}, _searchI_calls,
                                                max_calls=1)
            f.ask = _rs_ask_factory(
                [{"sufficient": False, "missing": "x", "queries": _bad_queries}],
                _askI_calls)
            _resI, _excI = None, None
            try:
                _resI = f.research_search(_qname)
            except Exception as _exc:
                _excI = _exc
            check(f"research_search: queries非list({_case_name})でも例外を送出しない",
                  _excI is None)
            check(f"research_search: queries非list({_case_name})では追加ラウンドの"
                  "検索が発生しない(1文字/1キーずつのクエリも発行されない)",
                  _searchI_calls == [_qname])
            check(f"research_search: queries非list({_case_name})でもそのラウンドの"
                  "結果は反映される",
                  _resI is not None and f"http://example.com/i-{_case_name}" in _resI)
        finally:
            f._search_raw = _orig_search_raw_rs
            f.ask = _orig_ask_rs

    # --- (I-valid) 対照実験: 有効な非空listのqueries + sufficient=falseは、
    #     修正後も従来通り追加ラウンドの検索を駆動する(既存挙動に変更がないことの確認)。
    _itemIv1 = "[Iv1]\nvalid list round1\nSource: http://example.com/iv1"
    _itemIv2 = "[Iv2]\nvalid list round2\nSource: http://example.com/iv2"
    _searchIv_calls = []
    _askIv_calls = []
    try:
        f._search_raw = _rs_search_factory(
            {"Valid List Query": [_itemIv1], "Valid Follow Up": [_itemIv2]},
            _searchIv_calls, max_calls=2)
        f.ask = _rs_ask_factory(
            [{"sufficient": False, "missing": "x", "queries": ["Valid Follow Up"]},
             {"sufficient": True, "missing": "", "queries": []}],
            _askIv_calls)
        _resIv = f.research_search("Valid List Query")
        check("research_search: 有効な非空listのqueriesはR2の追加検索を駆動する"
              "(既存挙動は不変)",
              _searchIv_calls == ["Valid List Query", "Valid Follow Up"])
        check("research_search: (I-valid)R1・R2両方の結果が反映される",
              "http://example.com/iv1" in _resIv and "http://example.com/iv2" in _resIv)
    finally:
        f._search_raw = _orig_search_raw_rs
        f.ask = _orig_ask_rs

    # --- (J) 2026-07-25: greedy first-fit パッキング — 途中の1件が SEARCH_CONTEXT_CHARS
    #     を超えていても、そこで break せず、それより後ろの小さく収まる結果は捨てずに拾う
    #     (採用順は維持)。イテレーション38の特性テスト(H)が発見した「巨大な1件が他の
    #     収集済み事実を握り潰す」問題のうち、イテレーション39では「先頭1件だけが単独で
    #     超過する」サブケースのみ対処され、「途中の1件が超過するとそれ以降が"全部"
    #     捨てられる」根本原因は未修正のまま残っていた回帰テスト。
    def _rs_pad_item(label, source_url, total_len, filler="Q"):
        """dedup用のSource URLを保ったまま、item全体をtotal_len文字ちょうどに
        パディングするテスト用ヘルパー(手計算での文字数ズレを避けるため)。"""
        head = f"[{label}]\n"
        tail = f"\nSource: {source_url}"
        pad_len = total_len - len(head) - len(tail)
        assert pad_len >= 0, "テスト用itemのtotal_lenが短すぎる"
        return head + (filler * pad_len) + tail

    _itemJ_a = "[Ja]\nsmall first result\nSource: http://example.com/ja"
    _itemJ_b = _rs_pad_item("Jb", "http://example.com/jb", f.SEARCH_CONTEXT_CHARS + 500)
    _itemJ_c = "[Jc]\nsmall third result\nSource: http://example.com/jc"
    _searchJ_calls = []
    _askJ_calls = []
    try:
        f._search_raw = _rs_search_factory(
            {"Greedy Skip Middle Query": [_itemJ_a, _itemJ_b, _itemJ_c]},
            _searchJ_calls, max_calls=1)
        f.ask = _rs_ask_factory([{"sufficient": True, "missing": "", "queries": []}],
                                 _askJ_calls)
        _resJ = f.research_search("Greedy Skip Middle Query")
        check("research_search: (J)中間の結果が上限超過でも前後の小さい結果は両方注入される",
              "http://example.com/ja" in _resJ and "http://example.com/jc" in _resJ)
        check("research_search: (J)上限超過した中間結果はbodyに含まれない",
              "http://example.com/jb" not in _resJ)
        check("research_search: (J)採用された結果の順序は元の順序のまま維持される"
              "(1番目の結果が3番目の結果より前に出現する)",
              _resJ.index("http://example.com/ja") < _resJ.index("http://example.com/jc"))
    finally:
        f._search_raw = _orig_search_raw_rs
        f.ask = _orig_ask_rs

    # --- (K) 2026-07-25: 単体では上限に収まる小さな結果でも、その前の結果群で既に
    #     予算の大半が埋まっていれば(残り予算に収まらないため)スキップされ、かつ
    #     最終bodyの長さがSEARCH_CONTEXT_CHARSを("\n\n"区切りのスラック分を除いて)
    #     超えないことを固定化する。スキップ後も走査は続くため、続く十分に小さい結果は
    #     引き続き拾われる(greedy first-fitの核: skipはbreakではない)。
    _itemK_a = _rs_pad_item("Ka", "http://example.com/ka", f.SEARCH_CONTEXT_CHARS - 100)
    _itemK_b = _rs_pad_item("Kb", "http://example.com/kb", 150)  # 単体ならCHARSに余裕で収まる
    _itemK_c = _rs_pad_item("Kc", "http://example.com/kc", 50)   # Ka採用後の残り予算には収まる
    _searchK_calls = []
    _askK_calls = []
    try:
        f._search_raw = _rs_search_factory(
            {"Greedy Budget Fill Query": [_itemK_a, _itemK_b, _itemK_c]},
            _searchK_calls, max_calls=1)
        f.ask = _rs_ask_factory([{"sufficient": True, "missing": "", "queries": []}],
                                 _askK_calls)
        _resK = f.research_search("Greedy Budget Fill Query")
        check("research_search: (K)先に予算の大半を占めた結果は採用される",
              "http://example.com/ka" in _resK)
        check("research_search: (K)単体ならCHARSに収まる結果でも残り予算不足ならスキップされる",
              "http://example.com/kb" not in _resK)
        check("research_search: (K)スキップ後も走査は続き、残り予算に収まる後続結果は採用される",
              "http://example.com/kc" in _resK)
        # ヘッダー（日付行＋固定の注意書き段落）は本文の一部ではないため、実際に注入
        # された検索結果本文（先頭は必ず itemK_a の "[Ka]"）以降だけを取り出して測る。
        _bodyK = _resK[_resK.index("[Ka]"):]
        check("research_search: (K)最終bodyの長さがSEARCH_CONTEXT_CHARSを"
              "(\"\\n\\n\"区切りのスラック分を除いて)超えない",
              len(_bodyK.strip()) <= f.SEARCH_CONTEXT_CHARS + 2)
    finally:
        f._search_raw = _orig_search_raw_rs
        f.ask = _orig_ask_rs

    # --- (L) 回帰: 全結果が小さく何も上限超過しないケースは、変更前と同じく
    #     全件を元の順序のまま連結する(body組み立てロジックの挙動は不変)。
    _itemL1 = "[L1]\nfirst small item\nSource: http://example.com/l1"
    _itemL2 = "[L2]\nsecond small item\nSource: http://example.com/l2"
    _itemL3 = "[L3]\nthird small item\nSource: http://example.com/l3"
    _searchL_calls = []
    _askL_calls = []
    try:
        f._search_raw = _rs_search_factory(
            {"All Small Query": [_itemL1, _itemL2, _itemL3]}, _searchL_calls, max_calls=1)
        f.ask = _rs_ask_factory([{"sufficient": True, "missing": "", "queries": []}],
                                 _askL_calls)
        _resL = f.research_search("All Small Query")
        _expected_body_l = _itemL1 + "\n\n" + _itemL2 + "\n\n" + _itemL3
        _headerL = f"## Web Search Results (取得日: {f.time.strftime('%Y-%m-%d')})"
        check("research_search: (L)結果ありなら取得日入りヘッダーが付与される(既存挙動)",
              _resL.startswith(_headerL))
        check("research_search: (L)全件が小さければ元の順序のまま連結される"
              "(変更前とbyte-for-byte同一の本文)",
              _resL.endswith(_expected_body_l))
    finally:
        f._search_raw = _orig_search_raw_rs
        f.ask = _orig_ask_rs

finally:
    f._search_raw = _orig_search_raw_rs
    f.ask = _orig_ask_rs
    urllib.request.urlopen = _orig_urlopen_rs
    f.subprocess.run = _orig_subprocess_run_rs

# ---------- _save_as_excel: XML不正制御文字によるIllegalCharacterError耐性 (2026-07-22) ----------
# openpyxl の ws.append() は XML 1.0 で禁止された制御文字(0x00-0x08/0x0B/0x0C/0x0E-0x1F)を
# 含むセルに対して openpyxl.utils.exceptions.IllegalCharacterError（ImportErrorではない
# 素のException）を送出する。従来コードは except ImportError しか捕捉しておらず、
# LLM回答に混入したフォームフィード(\x0c)/ANSIエスケープ(\x1b)/NUL(\x00)等が原因で
# _save_answer_to_file 全体が異常終了していた。ここでは openpyxl の有無を検出し、
# 両方の分岐（サニタイズして.xlsx生成 / 未インストール時の.csvフォールバック）を検証する。
# 本物のOllama/ネットワーク呼び出しは一切行わない。
import pathlib as _pathlib_xlsx

try:
    import openpyxl as _openpyxl_probe
    _HAS_OPENPYXL = True
except ImportError:
    _HAS_OPENPYXL = False

with _tempfile.TemporaryDirectory() as _xlsx_dir:
    _xlsx_root = _pathlib_xlsx.Path(_xlsx_dir)

    def _trim_trailing_none(_row):
        # openpyxl の iter_rows() はシート全体の最大列数まで各行を None で
        # パディングして返す(この関数の変更とは無関係の既存仕様)。比較用に
        # 末尾の None パディングだけを取り除く。
        _r = list(_row)
        while _r and _r[-1] is None:
            _r.pop()
        return _r

    if _HAS_OPENPYXL:
        # 制御文字はわざと文字列の「途中」に置く。ただし \x0b/\x0c/\x1c-\x1e は
        # Python の str.splitlines() 自体が改行境界として解釈し、answer.splitlines()
        # の時点でセル文字列に残らず消費されてしまう(この関数の既存仕様であり、
        # 今回の修正対象外)。そのため実際にセルへ到達し検証可能な、XML不正かつ
        # splitlines非対象の制御文字 \x1b(ESC)/\x00(NUL) を用いる。
        _illegal_answer = "name,age\nAli\x1bce,30\nB\x00ob,25"
        _out_illegal = _xlsx_root / "illegal.xlsx"
        _exc = None
        try:
            f._save_as_excel(_out_illegal, _illegal_answer)
        except Exception as _e:
            _exc = _e
        check("_save_as_excel: XML不正制御文字混入でも例外を送出しない(IllegalCharacterError回帰)",
              _exc is None)
        check("_save_as_excel: 制御文字混入時も.xlsxファイルが生成される", _out_illegal.exists())

        if _exc is None and _out_illegal.exists():
            _wb_illegal = _openpyxl_probe.load_workbook(str(_out_illegal))
            _rows_illegal = [_trim_trailing_none(row) for row in _wb_illegal.active.iter_rows(values_only=True)]
            check("_save_as_excel: 制御文字は除去されつつ実データ(表の中身)は保持される",
                  _rows_illegal == [["name", "age"], ["Alice", "30"], ["Bob", "25"]])

        # 制御文字を含まない通常回答は、列分割(re.split(r"[,\t|]", line))が
        # 従来とバイト単位で同一であること。
        _clean_answer = "a,b\tc|d\nx, y , z"
        _out_clean = _xlsx_root / "clean.xlsx"
        f._save_as_excel(_out_clean, _clean_answer)
        _wb_clean = _openpyxl_probe.load_workbook(str(_out_clean))
        _rows_clean = [_trim_trailing_none(row) for row in _wb_clean.active.iter_rows(values_only=True)]
        _expected_clean = [
            [c.strip() for c in f.re.split(r"[,\t|]", "a,b\tc|d")],
            [c.strip() for c in f.re.split(r"[,\t|]", "x, y , z")],
        ]
        check("_save_as_excel: 制御文字なしの通常回答は従来通り列分割される(既存挙動不変)",
              _rows_clean == _expected_clean)

        # ---------- _save_as_excel: Markdown表のパース (2026-07-26) ----------
        # 従来は全行を re.split(r"[,\t|]", line) で分割していたため、Markdown表の
        # 行 '| Name | Age |' が ['', 'Name', 'Age', ''] という前後に空列を持つ
        # 行になり、区切り線 '| --- | --- |' がそのままゴミの1データ行として
        # 書き込まれ、さらにセル内のカンマ（例: '1,234'）まで列区切りとして
        # 誤分割されていた。ここでは Markdown表を正しく読み、区切り線が
        # 除去され、セル内カンマが分割されないことを検証する。
        _md_answer = "| Name | Age |\n| --- | --- |\n| Alice | 30 |"
        _out_md = _xlsx_root / "md_table.xlsx"
        f._save_as_excel(_out_md, _md_answer)
        _wb_md = _openpyxl_probe.load_workbook(str(_out_md))
        _rows_md = [_trim_trailing_none(row) for row in _wb_md.active.iter_rows(values_only=True)]
        check("_save_as_excel: Markdown表は外側の'|'による空列なしで解釈される",
              _rows_md == [["Name", "Age"], ["Alice", "30"]])

        # コロン付きの整列指定を含む区切り線（左寄せ/右寄せ/中央寄せ）も
        # データ行として書き込まれず、丸ごとスキップされること。
        _md_align_answer = "| Name | Score | Note |\n| :--- | ---: | :---: |\n| Bob | 90 | ok |"
        _out_md_align = _xlsx_root / "md_align.xlsx"
        f._save_as_excel(_out_md_align, _md_align_answer)
        _wb_md_align = _openpyxl_probe.load_workbook(str(_out_md_align))
        _rows_md_align = [_trim_trailing_none(row) for row in _wb_md_align.active.iter_rows(values_only=True)]
        check("_save_as_excel: コロン付き整列区切り線(:---/---:/:---:)もスキップされる",
              _rows_md_align == [["Name", "Score", "Note"], ["Bob", "90", "ok"]])

        # セル内のカンマ（桁区切り数値など）が列区切りとして誤分割されないこと。
        _md_comma_answer = "| 1,234 | total |"
        _out_md_comma = _xlsx_root / "md_comma.xlsx"
        f._save_as_excel(_out_md_comma, _md_comma_answer)
        _wb_md_comma = _openpyxl_probe.load_workbook(str(_out_md_comma))
        _rows_md_comma = [_trim_trailing_none(row) for row in _wb_md_comma.active.iter_rows(values_only=True)]
        check("_save_as_excel: Markdown表のセル内カンマは列区切りとみなされない",
              _rows_md_comma == [["1,234", "total"]])

        # Markdown表の行の中に混入した制御文字も、iteration 41 のサニタイズが
        # 新しい分岐内でも適用され、除去されつつ実データは保持されること。
        _md_illegal_answer = "| Ali\x1bce | B\x00ob |"
        _out_md_illegal = _xlsx_root / "md_illegal.xlsx"
        _md_illegal_exc = None
        try:
            f._save_as_excel(_out_md_illegal, _md_illegal_answer)
        except Exception as _e:
            _md_illegal_exc = _e
        check("_save_as_excel: Markdown表内の制御文字混入でも例外を送出しない",
              _md_illegal_exc is None)
        if _md_illegal_exc is None:
            _wb_md_illegal = _openpyxl_probe.load_workbook(str(_out_md_illegal))
            _rows_md_illegal = [_trim_trailing_none(row) for row in _wb_md_illegal.active.iter_rows(values_only=True)]
            check("_save_as_excel: Markdown表内の制御文字は除去されつつ実データは保持される",
                  _rows_md_illegal == [["Alice", "Bob"]])
    else:
        # openpyxl 未インストール環境: 既存の.csvフォールバック(拡張子/メッセージ/戻り値)が
        # 従来通り機能すること。
        _out_missing = _xlsx_root / "missing.xlsx"
        _result_missing = f._save_as_excel(_out_missing, "a,b\nc,d")
        check("_save_as_excel: openpyxl未インストール時は.csvへフォールバックする",
              _result_missing == _out_missing.with_suffix(".csv"))
        check("_save_as_excel: フォールバック時の.csvファイルが実際に書かれる",
              _result_missing is not None and _result_missing.exists())

# ---------- _save_as_docx: XML不正制御文字によるValueError耐性 (2026-07-22) ----------
# python-docx (lxml) の add_paragraph/add_heading は XML 1.0 で禁止された制御文字
# (0x00-0x08/0x0B/0x0C/0x0E-0x1F) を含む文字列を渡されると ValueError を送出する。
# 従来コードは except ImportError しか捕捉しておらず、question をそのまま
# add_paragraph に渡していた事もあり、LLM回答/questionに混入したANSIエスケープ
# (\x1b)やNUL(\x00)等が原因で _save_answer_to_file 全体が異常終了していた
# （iteration 41 の _save_as_excel の IllegalCharacterError 修正と同じバグクラス）。
# ここでは python-docx の有無を検出し、サニタイズして.docx生成できること・
# 未インストール時は既存通り.mdへフォールバックすること・ビルド/保存失敗時も
# .mdへ安全に降格することを検証する。本物のOllama/ネットワーク呼び出しは一切行わない。
import pathlib as _pathlib_docx

try:
    import docx as _docx_probe
    _HAS_DOCX = True
except ImportError:
    _HAS_DOCX = False

with _tempfile.TemporaryDirectory() as _docx_dir:
    _docx_root = _pathlib_docx.Path(_docx_dir)

    if _HAS_DOCX:
        # 制御文字はわざと単語の間の空白の位置に置く。除去後も単語同士がくっつかず
        # 実データ(本文)が読み取れることを確認する。
        _illegal_question = "Hello \x1bWorld \x00Test?"
        _illegal_answer = "Answer \x1bline one.\n\n```python\nprint(\x001)\n```\n\nFinal \x00line."
        _out_illegal_docx = _docx_root / "illegal.docx"
        _exc_docx = None
        try:
            _ret_illegal = f._save_as_docx(_out_illegal_docx, _illegal_question, _illegal_answer, 0.42)
        except Exception as _e:
            _exc_docx = _e
            _ret_illegal = None
        check("_save_as_docx: XML不正制御文字混入でも例外を送出しない(ValueError回帰)",
              _exc_docx is None)
        check("_save_as_docx: 制御文字混入時も.docxファイルが生成される", _out_illegal_docx.exists())
        check("_save_as_docx: 制御文字混入でも成功時は.docxパス(None)を返す(戻り値契約維持)",
              _ret_illegal is None)

        if _exc_docx is None and _out_illegal_docx.exists():
            _doc_illegal = _docx_probe.Document(str(_out_illegal_docx))
            _texts_illegal = [p.text for p in _doc_illegal.paragraphs]
            _all_text_illegal = "\n".join(_texts_illegal)
            check("_save_as_docx: 制御文字(ESC/NUL)は本文から除去される",
                  "\x1b" not in _all_text_illegal and "\x00" not in _all_text_illegal)
            check("_save_as_docx: 制御文字除去後もquestion本文の単語は保持される",
                  "Hello World Test?" in _texts_illegal)
            check("_save_as_docx: 制御文字除去後もanswer本文の単語は保持される",
                  "Answer line one." in _texts_illegal and "Final line." in _texts_illegal)
            check("_save_as_docx: 制御文字除去後もコードブロック本文は保持される",
                  "print(1)" in _texts_illegal)

        # 制御文字を含まない通常回答は、コードフェンス解析・見出し構造・所要時間行が
        # 従来通り生成されること(既存挙動不変の回帰確認)。
        _clean_question = "Plain question"
        _clean_answer = "Intro line\n\n```python\nprint('hi')\n```\n\nOutro line"
        _out_clean_docx = _docx_root / "clean.docx"
        _ret_clean = f._save_as_docx(_out_clean_docx, _clean_question, _clean_answer, 2.5)
        check("_save_as_docx: 制御文字なしの通常回答も成功時はNoneを返す", _ret_clean is None)
        _doc_clean = _docx_probe.Document(str(_out_clean_docx))
        _paras_clean = list(_doc_clean.paragraphs)
        _texts_clean = [p.text for p in _paras_clean]
        _styles_clean = [p.style.name for p in _paras_clean]
        check("_save_as_docx: Q見出しが先頭に生成される(既存挙動不変)",
              len(_texts_clean) > 0 and _texts_clean[0].startswith("Q (") and _styles_clean[0] == "Heading 1")
        check("_save_as_docx: question本文がQ見出しの直後に生成される(既存挙動不変)",
              len(_texts_clean) > 1 and _texts_clean[1] == "Plain question")
        check("_save_as_docx: A見出しが生成される(既存挙動不変)",
              len(_texts_clean) > 2 and _texts_clean[2] == "A" and _styles_clean[2] == "Heading 1")
        check("_save_as_docx: コードフェンス本文はNo Spacingスタイルの段落になる(既存挙動不変)",
              "print('hi')" in _texts_clean and
              _styles_clean[_texts_clean.index("print('hi')")] == "No Spacing")
        check("_save_as_docx: 所要時間の行が末尾に生成される(既存挙動不変)",
              _texts_clean[-1] == "所要: 2.5s")

        # python-docx が未インストールの場合の分岐: sys.modules['docx'] を None に
        # することで `import docx` に ImportError を送出させる（実インストール状態を
        # 変更せずに未インストール環境を模擬する標準的な手法）。
        _orig_docx_mod = sys.modules.get("docx")
        sys.modules["docx"] = None
        try:
            _out_missing_docx = _docx_root / "missing.docx"
            _ret_missing = f._save_as_docx(_out_missing_docx, "Q?", "A.", 1.0)
            check("_save_as_docx: python-docx未インストール時は.mdへフォールバックする",
                  _ret_missing == _out_missing_docx.with_suffix(".md"))
            check("_save_as_docx: フォールバック時の.mdファイルが実際に書かれる",
                  _ret_missing is not None and _ret_missing.exists())
        finally:
            if _orig_docx_mod is not None:
                sys.modules["docx"] = _orig_docx_mod
            else:
                del sys.modules["docx"]

        # ビルド/保存自体が失敗するケース(IllegalXml以外の残存エラーも含む)を
        # docx.document.Document.save をモンキーパッチして模擬し、例外が
        # 外へ漏れずに.mdへ降格することを確認する。
        import docx.document as _docx_document_mod
        _orig_save_method = _docx_document_mod.Document.save

        def _boom_save(self, *_a, **_kw):
            raise RuntimeError("simulated docx save failure")

        _docx_document_mod.Document.save = _boom_save
        try:
            _out_fail_docx = _docx_root / "fail.docx"
            _exc_fail = None
            try:
                _ret_fail = f._save_as_docx(_out_fail_docx, "Q?", "A.", 1.0)
            except Exception as _e:
                _exc_fail = _e
                _ret_fail = None
            check("_save_as_docx: 保存失敗時も例外は外へ伝播しない", _exc_fail is None)
            check("_save_as_docx: 保存失敗時は.mdへフォールバックする",
                  _ret_fail == _out_fail_docx.with_suffix(".md"))
            check("_save_as_docx: 保存失敗フォールバック時の.mdファイルが実際に書かれる",
                  _ret_fail is not None and _ret_fail.exists())
        finally:
            _docx_document_mod.Document.save = _orig_save_method
    else:
        # python-docx 未インストール環境: 既存の.mdフォールバック(拡張子/戻り値)が
        # 従来通り機能すること。
        _out_missing_docx = _docx_root / "missing.docx"
        _result_missing_docx = f._save_as_docx(_out_missing_docx, "Q?", "A.", 1.0)
        check("_save_as_docx: python-docx未インストール時は.mdへフォールバックする",
              _result_missing_docx == _out_missing_docx.with_suffix(".md"))
        check("_save_as_docx: フォールバック時の.mdファイルが実際に書かれる",
              _result_missing_docx is not None and _result_missing_docx.exists())

# ---------- _save_as_pdf: fpdf2 ビルド/出力失敗時の .md フォールバック耐性 (2026-07-22) ----------
# fpdf2 には 'DejaVu' という名前で事前登録された組み込み Unicode フォントは無く、
# set_font("DejaVu") は FPDFException となり Helvetica (コアlatin-1フォント) へ
# 縮退する。既定言語である日本語などの非ASCII文字を multi_cell に渡すと
# FPDFUnicodeEncodingException（ImportErrorではない素のException）が送出される。
# 従来コードは except ImportError しか捕捉しておらず、_save_as_pdf の呼び出し元
# _save_answer_to_file 自体にもガードが無いため、回答保存ステップ全体が
# 異常終了していた（iteration 41 の _save_as_excel の IllegalCharacterError 修正、
# iteration 43 の _save_as_docx の制御文字 ValueError 修正と同じバグクラス）。
# ここでは fpdf2 の有無を検出し、(1) 日本語などUnicodeを含む通常呼び出しが例外を
# 送出せず実ファイルを生成すること、(2) ビルド/出力自体が失敗しても.mdへ安全に
# 降格すること、(3) fpdf2未インストール時の既存.mdフォールバック(メッセージ/戻り値)
# が変わらないことを検証する。本物のOllama/ネットワーク呼び出しは一切行わない。
import pathlib as _pathlib_pdf

try:
    import fpdf as _fpdf_probe
    _HAS_FPDF = True
except ImportError:
    _HAS_FPDF = False

with _tempfile.TemporaryDirectory() as _pdf_dir:
    _pdf_root = _pathlib_pdf.Path(_pdf_dir)

    if _HAS_FPDF:
        # 日本語(Unicode)を含む question/answer で呼んでも例外が伝播せず、
        # 実ファイル(.pdf または .md フォールバック)が生成されること。
        _ja_question = "日本語の質問です。テスト？"
        _ja_answer = "日本語の回答です。\n改行を含む本文。"
        _out_ja = _pdf_root / "ja.pdf"
        _exc_ja = None
        try:
            _ret_ja = f._save_as_pdf(_out_ja, _ja_question, _ja_answer, 1.23)
        except Exception as _e:
            _exc_ja = _e
            _ret_ja = None
        check("_save_as_pdf: 日本語/Unicode本文でも例外を送出しない(FPDFUnicodeEncodingException回帰)",
              _exc_ja is None)
        _produced_ja = _out_ja.exists() or _out_ja.with_suffix(".md").exists()
        check("_save_as_pdf: 日本語/Unicode本文でも実ファイル(.pdfまたは.md)が生成される",
              _produced_ja)

        # ビルド/出力自体が失敗するケース(Unicode以外の残存エラーも含む)を
        # fpdf.FPDF.output をモンキーパッチして模擬し、例外が外へ漏れずに
        # .md へ降格することを確認する。
        _orig_output_method = _fpdf_probe.FPDF.output

        def _boom_output(self, *_a, **_kw):
            raise RuntimeError("simulated pdf output failure")

        _fpdf_probe.FPDF.output = _boom_output
        try:
            _out_fail_pdf = _pdf_root / "fail.pdf"
            _exc_fail_pdf = None
            try:
                _ret_fail_pdf = f._save_as_pdf(_out_fail_pdf, "Q?", "A.", 1.0)
            except Exception as _e:
                _exc_fail_pdf = _e
                _ret_fail_pdf = None
            check("_save_as_pdf: 保存失敗時も例外は外へ伝播しない", _exc_fail_pdf is None)
            check("_save_as_pdf: 保存失敗時は.mdへフォールバックする",
                  _ret_fail_pdf == _out_fail_pdf.with_suffix(".md"))
            check("_save_as_pdf: 保存失敗フォールバック時の.mdファイルが実際に書かれる",
                  _ret_fail_pdf is not None and _ret_fail_pdf.exists())
        finally:
            _fpdf_probe.FPDF.output = _orig_output_method
    else:
        # fpdf2 が本当に存在しない環境: 既存の.mdフォールバック(拡張子/戻り値)が
        # 従来通り機能すること(iteration 41 の else分岐スタイルを踏襲)。
        _out_missing_pdf_real = _pdf_root / "missing_real.pdf"
        _result_missing_pdf_real = f._save_as_pdf(_out_missing_pdf_real, "Q?", "A.", 1.0)
        check("_save_as_pdf: fpdf2未インストール環境では.mdへフォールバックする",
              _result_missing_pdf_real == _out_missing_pdf_real.with_suffix(".md"))
        check("_save_as_pdf: フォールバック時の.mdファイルが実際に書かれる(未インストール環境)",
              _result_missing_pdf_real is not None and _result_missing_pdf_real.exists())

    # 回帰ガード: fpdf2 が実際にインストールされている環境でも sys.modules['fpdf']
    # を None にすることで `from fpdf import FPDF` に ImportError を送出させ
    # （実インストール状態を変更せずに未インストール環境を模擬する標準的な手法）、
    # 既存の ImportError フォールバック分岐(メッセージ/戻り値)が変わっていない
    # ことを検証する。
    _orig_fpdf_mod = sys.modules.get("fpdf")
    sys.modules["fpdf"] = None
    try:
        _out_missing_pdf = _pdf_root / "missing.pdf"
        _ret_missing_pdf = f._save_as_pdf(_out_missing_pdf, "Q?", "A.", 1.0)
        check("_save_as_pdf: fpdf2未インストール時は.mdへフォールバックする(既存メッセージ/戻り値不変)",
              _ret_missing_pdf == _out_missing_pdf.with_suffix(".md"))
        check("_save_as_pdf: フォールバック時の.mdファイルが実際に書かれる",
              _ret_missing_pdf is not None and _ret_missing_pdf.exists())
    finally:
        if _orig_fpdf_mod is not None:
            sys.modules["fpdf"] = _orig_fpdf_mod
        else:
            del sys.modules["fpdf"]

# ---------- build_pptx: XML不正制御文字によるValueError耐性 + 保存失敗の.md降格 (2026-07-23) ----------
# add_textbox が run.text に渡す文字列（タイトルスライド見出し・コンテンツスライド
# 見出し・箇条書き）に LLM 回答由来の制御文字 (NUL 0x00, ESC 0x1B 等) が混入すると、
# python-pptx/lxml が ValueError ('All strings must be XML compatible: no NULL bytes
# or control characters') を送出する。従来コードは except ImportError しか捕捉して
# おらず、_save_answer_to_file・さらには ask_fugu の make_pptx 経路まで例外が
# そのまま伝播し、計算済み（math/mcqではSC投票済みの）回答ごと失われていた
# （iteration 41 の _save_as_excel の IllegalCharacterError 修正・iteration 43 の
# _save_as_docx の制御文字 ValueError 修正・iteration 44 の _save_as_pdf の
# FPDFUnicodeEncodingException 修正と同じバグクラス）。ここでは python-pptx の
# 有無を検出し、(1) 制御文字混入でもサニタイズして.pptx生成できること・実データは
# 保持されること、(2) 制御文字なしの通常回答は従来通りの構造で生成されること
# （既存挙動不変）、(3) デッキ構築/保存自体が失敗しても.mdへ安全に降格すること、
# (4) python-pptx未インストール時の既存.mdフォールバックが変わらないことを検証する。
# 画像バックエンドは f.IMAGE_BACKEND="off" に固定し（試行後に必ず復元）、
# _detect_backend() がネットワーク探索なしで None を返すようにして画像生成経路を
# 完全に迂回する。本物のOllama/ネットワーク/画像バックエンド呼び出しは一切行わない。
import pathlib as _pathlib_pptx

try:
    import pptx as _pptx_probe
    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False

with _tempfile.TemporaryDirectory() as _pptx_dir:
    _pptx_root = _pathlib_pptx.Path(_pptx_dir)
    _orig_image_backend_pptx = f.IMAGE_BACKEND

    if _HAS_PPTX:
        f.IMAGE_BACKEND = "off"
        try:
            # (1) 制御文字混入(ESC/NUL、見出し行・箇条書き行それぞれの単語の途中)でも
            # 例外を送出せず.pptxが生成され、除去後の実データが読み取れること。
            _illegal_question_pptx = "PPTX 制御文字テスト"
            _illegal_answer_pptx = (
                "## Sec\x1btion One\n"
                "\n"
                "- Bul\x00let Al\x1bpha\n"
                "- Bullet Beta\n"
            )
            _out_illegal_pptx = _pptx_root / "illegal.pptx"
            _exc_pptx = None
            try:
                _ret_illegal_pptx = f.build_pptx(_illegal_question_pptx, _illegal_answer_pptx,
                                                  out_path=_out_illegal_pptx)
            except Exception as _e:
                _exc_pptx = _e
                _ret_illegal_pptx = None
            check("build_pptx: XML不正制御文字混入でも例外を送出しない(ValueError回帰)",
                  _exc_pptx is None)
            check("build_pptx: 制御文字混入時も.pptxファイルが生成される(戻り値がPath)",
                  isinstance(_ret_illegal_pptx, _pathlib_pptx.Path) and
                  _ret_illegal_pptx.suffix == ".pptx" and _ret_illegal_pptx.exists())

            if _exc_pptx is None and _ret_illegal_pptx is not None and _ret_illegal_pptx.exists():
                _prs_illegal = _pptx_probe.Presentation(str(_ret_illegal_pptx))
                _runs_illegal = [run.text
                                 for slide in _prs_illegal.slides
                                 for shape in slide.shapes if shape.has_text_frame
                                 for para in shape.text_frame.paragraphs
                                 for run in para.runs]
                _all_text_illegal_pptx = "\n".join(_runs_illegal)
                check("build_pptx: 制御文字(ESC/NUL)は本文から除去される",
                      "\x1b" not in _all_text_illegal_pptx and "\x00" not in _all_text_illegal_pptx)
                check("build_pptx: 制御文字除去後も見出しの実データは保持される(単語がくっつく)",
                      "Section One" in _runs_illegal)
                check("build_pptx: 制御文字除去後も箇条書きの実データは保持される",
                      "• Bullet Alpha" in _runs_illegal and "• Bullet Beta" in _runs_illegal)

            # (2) 制御文字を含まない通常の複数セクション回答は、従来通りタイトル
            # スライド+各見出しごとのコンテンツスライド+箇条書きが生成される
            # (既存挙動不変の回帰確認)。
            _clean_question_pptx = "Clean Question"
            _clean_answer_pptx = "## Intro\n- Point A\n- Point B\n\n## Details\n- Point C\n- Point D\n"
            _out_clean_pptx = _pptx_root / "clean.pptx"
            _ret_clean_pptx = f.build_pptx(_clean_question_pptx, _clean_answer_pptx,
                                            out_path=_out_clean_pptx)
            check("build_pptx: 制御文字なしの通常回答は成功時に.pptxのPathを返す(既存挙動不変)",
                  _ret_clean_pptx == _out_clean_pptx and _out_clean_pptx.exists())
            _prs_clean = _pptx_probe.Presentation(str(_ret_clean_pptx))
            _slides_clean = list(_prs_clean.slides)
            check("build_pptx: タイトルスライド+見出し2枚=計3枚のスライドが生成される(既存挙動不変)",
                  len(_slides_clean) == 3)

            def _slide_texts(_slide):
                return [run.text
                        for shape in _slide.shapes if shape.has_text_frame
                        for para in shape.text_frame.paragraphs
                        for run in para.runs]

            check("build_pptx: タイトルスライドに質問文がそのまま使われる(既存挙動不変)",
                  len(_slides_clean) > 0 and _clean_question_pptx in _slide_texts(_slides_clean[0]))
            check("build_pptx: 1枚目のコンテンツスライド見出しがIntroになる(既存挙動不変)",
                  len(_slides_clean) > 1 and "Intro" in _slide_texts(_slides_clean[1]))
            check("build_pptx: 1枚目のコンテンツスライドに箇条書きが両方含まれる(既存挙動不変)",
                  len(_slides_clean) > 1 and
                  "• Point A" in _slide_texts(_slides_clean[1]) and
                  "• Point B" in _slide_texts(_slides_clean[1]))
            check("build_pptx: 2枚目のコンテンツスライド見出しがDetailsになる(既存挙動不変)",
                  len(_slides_clean) > 2 and "Details" in _slide_texts(_slides_clean[2]))

            # (3) デッキ構築/保存自体が失敗するケースを pptx.presentation.Presentation.save
            # をモンキーパッチして模擬し、例外が外へ漏れずに既存の.mdフォールバックへ
            # 安全に降格することを確認する。
            import pptx.presentation as _pptx_presentation_mod
            _orig_pptx_save = _pptx_presentation_mod.Presentation.save

            def _boom_pptx_save(self, *_a, **_kw):
                raise RuntimeError("simulated pptx save failure")

            _pptx_presentation_mod.Presentation.save = _boom_pptx_save
            try:
                _out_fail_pptx = _pptx_root / "fail.pptx"
                _exc_fail_pptx = None
                try:
                    _ret_fail_pptx = f.build_pptx("Q?", "A.", out_path=_out_fail_pptx)
                except Exception as _e:
                    _exc_fail_pptx = _e
                    _ret_fail_pptx = None
                check("build_pptx: 保存失敗時も例外は外へ伝播しない", _exc_fail_pptx is None)
                check("build_pptx: 保存失敗時は.mdへフォールバックする",
                      _ret_fail_pptx == _out_fail_pptx.with_suffix(".md"))
                check("build_pptx: 保存失敗フォールバック時の.mdファイルが実際に書かれる",
                      _ret_fail_pptx is not None and _ret_fail_pptx.exists())
            finally:
                _pptx_presentation_mod.Presentation.save = _orig_pptx_save

            # (4) python-pptx が未インストールの場合の分岐: sys.modules['pptx'] を None に
            # することで `from pptx import Presentation` に ImportError を送出させる
            # （実インストール状態を変更せずに未インストール環境を模擬する標準的な手法。
            # iteration 43/44 と同じ手法）。
            _orig_pptx_mod = sys.modules.get("pptx")
            sys.modules["pptx"] = None
            try:
                _out_missing_pptx = _pptx_root / "missing.pptx"
                _ret_missing_pptx = f.build_pptx("Q?", "A.", out_path=_out_missing_pptx)
                check("build_pptx: python-pptx未インストール時は.mdへフォールバックする(既存挙動不変)",
                      _ret_missing_pptx == _out_missing_pptx.with_suffix(".md"))
                check("build_pptx: フォールバック時の.mdファイルが実際に書かれる",
                      _ret_missing_pptx is not None and _ret_missing_pptx.exists())
            finally:
                if _orig_pptx_mod is not None:
                    sys.modules["pptx"] = _orig_pptx_mod
                else:
                    del sys.modules["pptx"]
        finally:
            f.IMAGE_BACKEND = _orig_image_backend_pptx
    else:
        # python-pptx が本当に存在しない環境: 既存の.mdフォールバック(拡張子/戻り値)が
        # 従来通り機能すること(iteration 41/43/44 の else分岐スタイルを踏襲)。
        try:
            _out_missing_pptx_real = _pptx_root / "missing_real.pptx"
            _result_missing_pptx_real = f.build_pptx("Q?", "A.", out_path=_out_missing_pptx_real)
            check("build_pptx: python-pptx未インストール環境では.mdへフォールバックする",
                  _result_missing_pptx_real == _out_missing_pptx_real.with_suffix(".md"))
            check("build_pptx: フォールバック時の.mdファイルが実際に書かれる(未インストール環境)",
                  _result_missing_pptx_real is not None and _result_missing_pptx_real.exists())
        finally:
            f.IMAGE_BACKEND = _orig_image_backend_pptx

check("build_pptx: IMAGE_BACKEND はテスト後に既定値へ復元されている",
      f.IMAGE_BACKEND == _orig_image_backend_pptx)

# ---------- build_pptx: 画像プランが index 0 を省略し枠が満杯でも
# タイトルヒーロー画像が生き残ることの確認 (2026-07-23) ----------
# 背景: plan_pptx_images() は最大 PPTX_MAX_IMAGES 件の {index: prompt} を返す。
# build_pptx はその直後に plan.setdefault(0, None) で「タイトルには必ず
# ヒーロー画像」という不変条件を保証しようとするが、dict は挿入順を保持する
# ため、plan が index 0 を含まずに PPTX_MAX_IMAGES 件ちょうどで満杯だった
# 場合、setdefault は 0 を末尾に追加するだけになる。従来コードの
# list(plan.items())[:PPTX_MAX_IMAGES] スライスは直後にその末尾の 0 を
# 切り捨ててしまい、LLM が plan_pptx_images のシステムプロンプトの
# 「Include index 0」指示に従わなかった場合、タイトルスライドへヒーロー
# 画像が入らないまま不変条件が静かに破られていた。本テストは
# _detect_backend / plan_pptx_images / generate_image / author_image_prompt を
# 全てモックし、実際の Ollama・画像バックエンド・ネットワーク呼び出しを
# 一切行わずにこのバグと修正後の挙動、および既存の非バグケース(byte-for-byte
# 不変)を検証する。
import base64 as _base64_pptx_hero

_ONE_PX_PNG_HERO = _base64_pptx_hero.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+A8AAQUBAScY"
    "42YAAAAASUVORK5CYII="
)

if _HAS_PPTX:
    with _tempfile.TemporaryDirectory() as _pptx_hero_dir:
        _pptx_hero_root = _pathlib_pptx.Path(_pptx_hero_dir)
        _fake_img_path_hero = _pptx_hero_root / "fake.png"
        _fake_img_path_hero.write_bytes(_ONE_PX_PNG_HERO)

        _orig_ib_hero = f.IMAGE_BACKEND
        _orig_detect_hero = f._detect_backend
        _orig_plan_hero = f.plan_pptx_images
        _orig_genimg_hero = f.generate_image
        _orig_authimg_hero = f.author_image_prompt

        _genimg_calls_hero = []
        _authimg_calls_hero = []

        def _fake_generate_image_hero(prompt, negative=""):
            _genimg_calls_hero.append((prompt, negative))
            return str(_fake_img_path_hero)

        def _fake_author_image_prompt_hero(base_text, panel=None):
            _authimg_calls_hero.append(base_text)
            return ("TITLE_PROMPT", "TITLE_NEG")

        try:
            f.IMAGE_BACKEND = "a1111"
            f._detect_backend = lambda: "a1111"
            f.generate_image = _fake_generate_image_hero
            f.author_image_prompt = _fake_author_image_prompt_hero

            check("build_pptx: このテストはPPTX_MAX_IMAGES==4を前提にしている",
                  f.PPTX_MAX_IMAGES == 4)

            _hero_question = "Hero Bug Test"
            _hero_answer = "## S1\n- b1\n\n## S2\n- b2\n\n## S3\n- b3\n\n## S4\n- b4\n"

            # (A) バグ再現ケース: plan が PPTX_MAX_IMAGES 件ちょうど、index 0 を含まない。
            f.plan_pptx_images = lambda title, slides: {1: "p1", 2: "p2", 3: "p3", 4: "p4"}
            _genimg_calls_hero.clear()
            _authimg_calls_hero.clear()
            _out_hero_a = _pptx_hero_root / "hero_full_no_zero.pptx"
            _ret_hero_a = f.build_pptx(_hero_question, _hero_answer, out_path=_out_hero_a)
            check("build_pptx/画像0省略+満杯: 例外を送出せず.pptxを生成する",
                  isinstance(_ret_hero_a, _pathlib_pptx.Path) and _ret_hero_a.suffix == ".pptx"
                  and _ret_hero_a.exists())
            check("build_pptx/画像0省略+満杯: タイトルヒーロー画像生成が実行される"
                  "(author_image_promptがタイトルで呼ばれる)",
                  _hero_question in _authimg_calls_hero)
            check("build_pptx/画像0省略+満杯: タイトルヒーロー画像がgenerate_imageに渡る(TITLE_PROMPT)",
                  any(p == "TITLE_PROMPT" for p, _n in _genimg_calls_hero))
            check("build_pptx/画像0省略+満杯: 画像総数はPPTX_MAX_IMAGESを超えない"
                  "(タイトル用に1枠を再割当するだけで追加はしない)",
                  len(_genimg_calls_hero) <= f.PPTX_MAX_IMAGES)
            check("build_pptx/画像0省略+満杯: contentスライド側は1枠だけ犠牲になり3枚は生き残る",
                  sum(1 for p, _n in _genimg_calls_hero if p in {"p1", "p2", "p3", "p4"})
                  == f.PPTX_MAX_IMAGES - 1)

            # (B) 回帰ケース1: plan が既に index 0 を含み、かつ枠がちょうど満杯
            # (PPTX_MAX_IMAGES件)。0 が既に予算内にあるので並び替え分岐に入らず、
            # 従来通り全4件がそのまま生成対象になる(byte-for-byte不変)。
            f.plan_pptx_images = lambda title, slides: {0: "p0", 1: "p1", 2: "p2", 3: "p3"}
            _genimg_calls_hero.clear()
            _authimg_calls_hero.clear()
            _out_hero_b = _pptx_hero_root / "hero_full_with_zero.pptx"
            f.build_pptx(_hero_question, _hero_answer, out_path=_out_hero_b)
            check("build_pptx/画像0含む+満杯(回帰): author_image_promptは一切呼ばれない"
                  "(全エントリにプロンプトがあるため)",
                  _authimg_calls_hero == [])
            check("build_pptx/画像0含む+満杯(回帰): 生成対象の集合は従来通り{p0,p1,p2,p3}のまま",
                  {p for p, _n in _genimg_calls_hero} == {"p0", "p1", "p2", "p3"}
                  and len(_genimg_calls_hero) == 4)

            # (C) 回帰ケース2: plan が PPTX_MAX_IMAGES 未満で index 0 を含まない。
            # setdefaultで0が追加されても件数は定員未満のままなのでスライスは
            # 何も切り捨てず、0を含め全エントリがそのまま生き残る(従来通り)。
            f.plan_pptx_images = lambda title, slides: {2: "p2"}
            _genimg_calls_hero.clear()
            _authimg_calls_hero.clear()
            _out_hero_c = _pptx_hero_root / "hero_partial_no_zero.pptx"
            f.build_pptx(_hero_question, _hero_answer, out_path=_out_hero_c)
            check("build_pptx/画像0省略+定員未満(回帰): タイトル用にauthor_image_promptが呼ばれる",
                  _hero_question in _authimg_calls_hero)
            check("build_pptx/画像0省略+定員未満(回帰): 元々あったp2も生き残る",
                  any(p == "p2" for p, _n in _genimg_calls_hero))
            check("build_pptx/画像0省略+定員未満(回帰): タイトル用のTITLE_PROMPTも生成される",
                  any(p == "TITLE_PROMPT" for p, _n in _genimg_calls_hero))
            check("build_pptx/画像0省略+定員未満(回帰): 生成呼び出しは2件のみ(p2 + タイトル)",
                  len(_genimg_calls_hero) == 2)
        finally:
            f.IMAGE_BACKEND = _orig_ib_hero
            f._detect_backend = _orig_detect_hero
            f.plan_pptx_images = _orig_plan_hero
            f.generate_image = _orig_genimg_hero
            f.author_image_prompt = _orig_authimg_hero

    check("build_pptx: 画像ヒーローテスト後にIMAGE_BACKENDが既定値へ復元されている",
          f.IMAGE_BACKEND == _orig_image_backend_pptx)

# ---------- _trim_history: 直近ペアを消さない (2026-07-23回帰) ----------
# 背景: ask_fugu は直近の [user, assistant] ペアを _HISTORY に追記した「直後」に
# _trim_history を呼ぶ。旧ガード `len(history) >= 2` だと、履歴が直近ペア1組
# (長さ2) まで削られた状態でもまだ MAX_HISTORY_CHARS 超過ならループに入り、
# pop(0) を2回叩いて直近ペアごと消してしまい、履歴が完全に空になっていた
# (直後に "[会話履歴: 0 往復保持中]" と出る既知の再現手順)。
# `len(history) > 2` に直したことで、削除対象が最後の1ペアだけになった時点で
# 必ず停止し、多ターン対話の文脈(精度に直結)を失わないことをここで検証する。
_orig_max_hist_chars = f.MAX_HISTORY_CHARS

# (1) 予算内なら無加工でそのまま(同一オブジェクト・同一順序・同一長)
_hist_small = [{"role": "user", "content": "hello"},
               {"role": "assistant", "content": "world"}]
_hist_small_ids_before = [id(m) for m in _hist_small]
f._trim_history(_hist_small)
check("trim_history: 予算内は要素数不変(無加工)", len(_hist_small) == 2)
check("trim_history: 予算内は同一オブジェクト・順序も不変(無加工)",
      [id(m) for m in _hist_small] == _hist_small_ids_before)

# (2) バグ再現(主目的): 直近1ペアだけの状態で単独で予算超過しても消えない
try:
    f.MAX_HISTORY_CHARS = 4000
    _hist_one_pair = [{"role": "user", "content": "q"},
                       {"role": "assistant", "content": "X" * 6000}]
    f._trim_history(_hist_one_pair)
    check("trim_history: 直近1ペアが単独で予算超過しても消えない(バグ再現)",
          len(_hist_one_pair) == 2)
    check("trim_history: 残るのは直近のuser/assistantペアそのもの(内容不変)",
          _hist_one_pair[0]["role"] == "user"
          and _hist_one_pair[0]["content"] == "q"
          and _hist_one_pair[1]["role"] == "assistant"
          and _hist_one_pair[1]["content"] == "X" * 6000)
finally:
    f.MAX_HISTORY_CHARS = _orig_max_hist_chars

# (3) 回帰: 複数ペアが予算超過なら古いペアから先頭ごと削除され、
#     生き残りは元リストの末尾(tail)と一致する(順序維持・ペア単位で削除)
try:
    f.MAX_HISTORY_CHARS = 250
    _pair1 = [{"role": "user", "content": "u" * 50},
              {"role": "assistant", "content": "a" * 50}]
    _pair2 = [{"role": "user", "content": "u" * 50},
              {"role": "assistant", "content": "a" * 50}]
    _pair3 = [{"role": "user", "content": "u" * 50},
              {"role": "assistant", "content": "a" * 50}]
    _hist_multi = _pair1 + _pair2 + _pair3
    _hist_multi_expected_tail = _pair2 + _pair3
    f._trim_history(_hist_multi)
    check("trim_history: 予算超過時は古いペアが先頭からペア単位で削除される",
          _hist_multi == _hist_multi_expected_tail)
    check("trim_history: 生き残りは元リストの末尾と一致する(順序維持)",
          len(_hist_multi) == 4 and _hist_multi[0]["content"] == "u" * 50)
finally:
    f.MAX_HISTORY_CHARS = _orig_max_hist_chars

# (4) 全ペアが予算超過でも直近の1ペアだけは必ず残る(先頭ペアのみ削除)
try:
    f.MAX_HISTORY_CHARS = 4000
    _big_pair1 = [{"role": "user", "content": "u" * 3000},
                  {"role": "assistant", "content": "a" * 3000}]
    _big_pair2 = [{"role": "user", "content": "u" * 3000},
                  {"role": "assistant", "content": "a" * 3000}]
    _hist_two_big = _big_pair1 + _big_pair2
    f._trim_history(_hist_two_big)
    check("trim_history: 2ペアとも予算超過でも直近ペアのみ残る(先頭ペア削除)",
          _hist_two_big == _big_pair2)
finally:
    f.MAX_HISTORY_CHARS = _orig_max_hist_chars

# (5) 一般化: 極小予算・様々なペア数でも空にならず、例外を送出せず完走する(ハングしない)
try:
    f.MAX_HISTORY_CHARS = 1
    _no_hang_exc = None
    try:
        for _n_pairs in (1, 2, 3, 5):
            _h = []
            for _i in range(_n_pairs):
                _h.append({"role": "user", "content": "u" * 500})
                _h.append({"role": "assistant", "content": "a" * 500})
            f._trim_history(_h)
            check(f"trim_history: 予算1でも空にならない(ペア数={_n_pairs})",
                  len(_h) >= 2)
    except Exception as _e:
        _no_hang_exc = _e
    check("trim_history: 極小予算でも例外を送出せず完走する(ハング/クラッシュなし)",
          _no_hang_exc is None)
finally:
    f.MAX_HISTORY_CHARS = _orig_max_hist_chars

# ---------- 会話履歴の永続化（load/save_history_file） ----------
# load_history_file (L353) / save_history_file (L369) はセッションをまたいだ
# 複数ターン会話のメモリを担う。iteration 59 の _trim_history 修正と同じ
# 「文脈を失わない」精度クリティカルな経路だが、これまでテストが皆無だった
# (grep しても関数本体・SESSION_SAVE・MAX_HISTORY_TURNS_SAVED のいずれも
# test_fugu_offline.py に一度も出現しない)。
# 重要: 以下は全呼び出しで path= に一時ファイルを明示的に渡す。デフォルト
# 引数は Path.home() / ".fugu_history.json"（ユーザーの実際の会話履歴）を
# 指すため、path= を省略すると本物の履歴を読み込み/破壊してしまう危険がある。
# ここは純粋な一時ファイルI/Oのみで、Ollama/ネットワーク/subprocess呼び出しは
# 一切ない。gotcha 該当箇所（/api/chat・num_ctx固定・think除去リトライ・
# cp932 reconfigure・OLLAMA_MAX_LOADED_MODELS・math_verifyタイムアウト・
# SC投票/solve_verifiable）はいずれも本セクションでは触れていない。
_orig_session_save = f.SESSION_SAVE
_orig_max_hist_turns_saved = f.MAX_HISTORY_TURNS_SAVED

with _tempfile.TemporaryDirectory() as _hist_dir:
    _hist_root = f.Path(_hist_dir)

    # (1) ファイルが存在しない -> 空リスト
    _hp_missing = _hist_root / "missing.json"
    check("history: 存在しないファイルは空リストを返す",
          f.load_history_file(path=_hp_missing) == [])

    # (2) 保存->読込のラウンドトリップは等価なメッセージ列を返す
    _hp_roundtrip = _hist_root / "roundtrip.json"
    _hist_valid = [{"role": "user", "content": "こんにちは"},
                   {"role": "assistant", "content": "こんにちは、ご用件は？"},
                   {"role": "user", "content": "元気です"},
                   {"role": "assistant", "content": "それは良かったです"}]
    f.save_history_file(_hist_valid, path=_hp_roundtrip)
    check("history: save->loadのラウンドトリップは等価なリストを返す",
          f.load_history_file(path=_hp_roundtrip) == _hist_valid)

    # (3) 壊れたJSON -> 空リスト(広い except で例外を送出しない)
    _hp_corrupt = _hist_root / "corrupt.json"
    _hp_corrupt.write_text("{not valid json!!!", encoding="utf-8")
    check("history: 壊れたJSONファイルは空リストを返す",
          f.load_history_file(path=_hp_corrupt) == [])

    # (4) トップレベルが list ではない(dict) -> 空リスト
    _hp_notlist = _hist_root / "notlist.json"
    _hp_notlist.write_text(json.dumps({"role": "user", "content": "x"}),
                            encoding="utf-8")
    check("history: トップレベルが list以外(dict)なら空リストを返す",
          f.load_history_file(path=_hp_notlist) == [])

    # (5) 不正エントリのフィルタリング: dict以外・role欠落・content欠落を除外し、
    #     整形済みメッセージのみを元の順序のまま返す
    _hp_mixed = _hist_root / "mixed.json"
    _mixed_raw = [
        {"role": "user", "content": "1問目"},
        "not a dict",
        {"role": "assistant"},          # content 欠落 -> 除外
        {"content": "role欠落"},         # role 欠落 -> 除外
        {"role": "assistant", "content": "1答目"},
        123,
        None,
        {"role": "user", "content": "2問目"},
    ]
    _hp_mixed.write_text(json.dumps(_mixed_raw, ensure_ascii=False), encoding="utf-8")
    check("history: 不正エントリを除外し整形済みメッセージのみ順序維持で返す",
          f.load_history_file(path=_hp_mixed) == [
              {"role": "user", "content": "1問目"},
              {"role": "assistant", "content": "1答目"},
              {"role": "user", "content": "2問目"},
          ])

    # (5b) 2026-07-23: iteration 66 で指摘・未修正だった値型バグの回帰テスト。
    #      role/content の「キーの有無」だけでなく「値が str かどうか」も
    #      検証しなければならない。content が None/int/list/dict、あるいは
    #      role が str でないエントリを混入させ、いずれも除外され
    #      文字列値の整形済みメッセージのみが元の順序で残ることを確認する。
    _hp_badval = _hist_root / "badval.json"
    _badval_raw = [
        {"role": "user", "content": "1問目"},
        {"role": "user", "content": None},       # content が None -> 除外
        {"role": "user", "content": 123},        # content が int -> 除外
        {"role": "user", "content": ["x"]},      # content が list -> 除外
        {"role": "user", "content": {"x": 1}},   # content が dict -> 除外
        {"role": "assistant", "content": "1答目"},
        {"role": 123, "content": "x"},           # role が int(非str) -> 除外
        {"role": "user", "content": "2問目"},
    ]
    _hp_badval.write_text(json.dumps(_badval_raw, ensure_ascii=False), encoding="utf-8")
    _loaded_badval = f.load_history_file(path=_hp_badval)
    check("history: contentが非str(None/int/list/dict)のエントリを除外する",
          _loaded_badval == [
              {"role": "user", "content": "1問目"},
              {"role": "assistant", "content": "1答目"},
              {"role": "user", "content": "2問目"},
          ])
    check("history: roleが非str(int)のエントリを除外する",
          all(isinstance(m["role"], str) for m in _loaded_badval))

    # (5c) エンドツーエンドのクラッシュ経路防止確認: content:null を含む壊れた
    #      ファイルを load_history_file で読み込んだ結果を _trim_history に
    #      そのまま渡しても、以前クラッシュしていた
    #      `sum(len(m["content"]) for m in history)` の TypeError(len(None))
    #      経路が発生しないこと(=事前に非strエントリが除外されていること)を確認する。
    _hp_crash = _hist_root / "crash.json"
    _hp_crash.write_text(
        json.dumps([{"role": "user", "content": "生き残る"},
                    {"role": "user", "content": None}], ensure_ascii=False),
        encoding="utf-8")
    _loaded_crash = f.load_history_file(path=_hp_crash)
    _trim_crash_exc = None
    try:
        f._trim_history(_loaded_crash)
    except Exception as _e:
        _trim_crash_exc = _e
    check("history: content:nullを含む壊れたファイルをload後、_trim_historyが例外を送出しない",
          _trim_crash_exc is None)

    # (6) 末尾スライスの上限(MAX_HISTORY_TURNS_SAVED*2)を load/save 双方で検証
    try:
        f.MAX_HISTORY_TURNS_SAVED = 2  # 上限を一時的に縮小(*2 = 4件まで保持)
        _hist_long = [{"role": ("user" if i % 2 == 0 else "assistant"),
                       "content": f"msg{i}"} for i in range(10)]
        _expected_tail = _hist_long[-4:]

        # load側: 上限より多いエントリを持つファイルを直接書いて確認
        _hp_cap_load = _hist_root / "cap_load.json"
        _hp_cap_load.write_text(json.dumps(_hist_long, ensure_ascii=False),
                                 encoding="utf-8")
        check("history: loadは直近 MAX_HISTORY_TURNS_SAVED*2 件のみ返す",
              f.load_history_file(path=_hp_cap_load) == _expected_tail)

        # save側: 書き込み時点で切り詰められ、ファイルの生JSONも上限内に収まる
        _hp_cap_save = _hist_root / "cap_save.json"
        f.save_history_file(_hist_long, path=_hp_cap_save)
        check("history: saveは書き込み時点で MAX_HISTORY_TURNS_SAVED*2 件に切り詰める",
              f.load_history_file(path=_hp_cap_save) == _expected_tail)
        _raw_cap_save = json.loads(_hp_cap_save.read_text(encoding="utf-8"))
        check("history: 切り詰め後にファイルへ書かれる生JSONも上限件数のみ",
              len(_raw_cap_save) == 4)
    finally:
        f.MAX_HISTORY_TURNS_SAVED = _orig_max_hist_turns_saved

    # (7) SESSION_SAVE=False は save を no-op にする(ファイル未作成/既存ファイル不変)
    try:
        f.SESSION_SAVE = False
        _hp_nosave = _hist_root / "nosave.json"
        f.save_history_file([{"role": "user", "content": "saved?"}], path=_hp_nosave)
        check("history: SESSION_SAVE=Falseなら新規ファイルを作成しない",
              not _hp_nosave.exists())

        _hp_nosave_existing = _hist_root / "nosave_existing.json"
        _existing_content = json.dumps([{"role": "user", "content": "keep"}],
                                        ensure_ascii=False)
        _hp_nosave_existing.write_text(_existing_content, encoding="utf-8")
        f.save_history_file([{"role": "user", "content": "should not overwrite"}],
                             path=_hp_nosave_existing)
        check("history: SESSION_SAVE=Falseなら既存ファイルも変更しない",
              _hp_nosave_existing.read_text(encoding="utf-8") == _existing_content)
    finally:
        f.SESSION_SAVE = _orig_session_save

    # (8) 非ASCII(日本語)コンテンツがラウンドトリップで無傷(ensure_ascii=False)
    _hp_ja = _hist_root / "ja.json"
    _hist_ja = [{"role": "user",
                 "content": "日本語の質問です。特殊文字：あいうえお、漢字も含む"},
                {"role": "assistant", "content": "はい、日本語で回答します。"}]
    f.save_history_file(_hist_ja, path=_hp_ja)
    check("history: 日本語コンテンツはラウンドトリップで無傷",
          f.load_history_file(path=_hp_ja) == _hist_ja)
    _raw_ja_text = _hp_ja.read_text(encoding="utf-8")
    check("history: 保存ファイルの生テキストに日本語がそのまま含まれる(\\uXXXXエスケープでない)",
          "日本語の質問です" in _raw_ja_text)

    # (9) 2026-07-23: load_history_file の読み込みに errors="replace" を追加した
    #     修正の回帰テスト(このタスク本体)。落とし穴 #4(このマシンのコンソールが
    #     cp932)と同種の環境要因で、セッションファイルが Shift-JIS エディタでの
    #     開き直し等により非UTF-8バイト列を部分的に含むことがある。JSON構造自体は
    #     健全 (dictのlist、role/contentともstr) でも、修正前の
    #     path.read_text(encoding="utf-8") はバイト列全体のデコードに失敗して
    #     UnicodeDecodeError を送出し、広い except Exception: pass に捕まって
    #     空リスト [] を返していた＝読み取れる他のエントリまで含め会話履歴を
    #     丸ごと喪失していた。iteration 47 で _save_as_markdown/_save_as_text/
    #     _save_as_html の読み戻しに適用した errors="replace" パターンと同じ
    #     一引数の変更で、読めないバイトだけを置換文字(U+FFFD)に落として
    #     JSONパースを成功させ、破損エントリの前後の文字列も無傷の兄弟エントリも
    #     回収できることを確認する。
    _hp_cp932 = _hist_root / "cp932_corrupt.json"
    _hist_cp932_raw = [
        {"role": "user", "content": "壊れる前@@BADBYTES@@壊れる後"},
        {"role": "assistant", "content": "無傷であるはずの兄弟エントリ：日本語も含む"},
    ]
    # 構造は正当なJSONとして組み立ててからUTF-8バイト列化し、マーカー部分だけを
    # UTF-8として不正なバイト列(0xFF, 0xFE は単独では常に不正)へ差し替える。
    # JSON区切り文字(引用符・波括弧・カンマ)自体はASCIIのまま保つため、
    # errors="replace" 適用後のテキストは引き続き妥当なJSONとしてパースできる。
    _hist_cp932_bytes = (
        json.dumps(_hist_cp932_raw, ensure_ascii=False).encode("utf-8")
        .replace(b"@@BADBYTES@@", b"\xff\xfe\xff\xfe")
    )
    _hp_cp932.write_bytes(_hist_cp932_bytes)

    _loaded_cp932 = f.load_history_file(path=_hp_cp932)
    check("history: 非UTF-8バイトを含む(構造は健全な)セッションファイルでも"
          "空リストにならず読み込める(修正前はUnicodeDecodeErrorで[]だった)",
          len(_loaded_cp932) == 2)
    check("history: 復元されたエントリのrole/contentは常にstr型",
          all(isinstance(m.get("role"), str) and isinstance(m.get("content"), str)
              for m in _loaded_cp932))
    check("history: 壊れたバイトはU+FFFDに置換されつつ前後の文字列は保持される",
          "壊れる前" in _loaded_cp932[0]["content"]
          and "壊れる後" in _loaded_cp932[0]["content"]
          and "�" in _loaded_cp932[0]["content"])
    check("history: 破損していない兄弟エントリはバイト単位で無傷のまま読み込まれる",
          _loaded_cp932[1]["content"] == "無傷であるはずの兄弟エントリ：日本語も含む")

    # (10) 2026-07-26: save_history_file に force フラグ + 真偽値の戻り値を追加した
    #      修正の回帰テスト。repl() の 'save <path>' コマンドはユーザーが明示的に
    #      指示した履歴エクスポートだが、修正前は save_history_file が
    #      SESSION_SAVE=False (--no-history) 時に無条件で早期returnし、かつ常に
    #      None を返していたため、repl 側は書き込みの成否に関わらず
    #      「[履歴を保存しました: ...]」と表示していた（実際にはファイルが
    #      一切書かれていない/エラーで失敗していてもユーザーには成功したように
    #      見える、という無言のデータロス）。ここでは save_history_file 単体の
    #      契約を直接検証する: force=False（既定）は iteration 66/67 と同じ
    #      SESSION_SAVE=False no-op 契約を保つ（バイト単位で不変）ことと、
    #      force=True は SESSION_SAVE の値に関わらず書き込みを行うこと、
    #      戻り値が「実際に書けたか」を正確に反映すること、書き込み失敗時に
    #      例外を伝播させず既存の '[履歴保存エラー: ...]' 表示を保つことを
    #      確認する。SESSION_SAVE はこのブロック内で try/finally により復元する。
    try:
        # (10a) force=True + SESSION_SAVE=False -> 書き込みが行われ True を返す
        f.SESSION_SAVE = False
        _hp_force = _hist_root / "force_write.json"
        _hist_force = [{"role": "user", "content": "force保存されるはず"}]
        _ret_force = f.save_history_file(_hist_force, path=_hp_force, force=True)
        check("history: force=Trueなら SESSION_SAVE=False でもファイルを書き込む",
              _hp_force.exists())
        check("history: force=True かつ書き込み成功時は True を返す",
              _ret_force is True)
        check("history: force=True で書き込んだ内容はload/jsonでラウンドトリップする",
              f.load_history_file(path=_hp_force) == _hist_force)

        # (10b) force=False(既定) + SESSION_SAVE=False -> 依然として no-op かつ False
        _hp_default_nosave = _hist_root / "default_nosave.json"
        _ret_default_nosave = f.save_history_file(
            [{"role": "user", "content": "saved?"}], path=_hp_default_nosave)
        check("history: 既定(force=False)は従来通りSESSION_SAVE=Falseでファイルを作成しない"
              "(iteration66/67のno-op契約の回帰防止)",
              not _hp_default_nosave.exists())
        check("history: 既定(force=False)でSESSION_SAVE=Falseならスキップを示すFalseを返す",
              _ret_default_nosave is False)

        # (10c) force=False(既定) + SESSION_SAVE=True -> 従来通り書き込み、True を返す
        f.SESSION_SAVE = True
        _hp_default_save = _hist_root / "default_save.json"
        _hist_default_save = [{"role": "assistant", "content": "自動保存されるはず"}]
        _ret_default_save = f.save_history_file(_hist_default_save, path=_hp_default_save)
        check("history: 既定(force=False)でSESSION_SAVE=Trueなら従来通り書き込む",
              f.load_history_file(path=_hp_default_save) == _hist_default_save)
        check("history: 既定(force=False)でSESSION_SAVE=Trueなら書き込み成功でTrueを返す",
              _ret_default_save is True)
    finally:
        f.SESSION_SAVE = _orig_session_save

    # (10d) 強制書き込みが失敗する場合(存在しないディレクトリ配下への書き込み)は
    #       例外を伝播させず、既存の '[履歴保存エラー: ...]' 表示を維持しつつ
    #       False を返す(force=True/SESSION_SAVE いずれの値でも同じ契約)。
    _hp_bad_dir = _hist_root / "no_such_subdir" / "unwritable.json"
    with contextlib.redirect_stdout(io.StringIO()) as _hbw_out:
        _ret_bad_dir = f.save_history_file(
            [{"role": "user", "content": "x"}], path=_hp_bad_dir, force=True)
    check("history: 存在しないディレクトリへのforce書き込みは例外を送出せずFalseを返す",
          _ret_bad_dir is False)
    check("history: 書き込み失敗時に既存の[履歴保存エラー: ...]表示を維持する",
          "履歴保存エラー" in _hbw_out.getvalue())
    check("history: 書き込みに失敗したパスにファイルは作成されない",
          not _hp_bad_dir.exists())

# ---------- ask_fugu: _HISTORY には成果物付き final ではなく text_answer を保存 ----------
# 2026-07-23: fugu_local.py L3098 のコメント「履歴にはテキスト本文のみ保存する」が
# 実装(L3129 で _HISTORY に final を追記)と乖離していたバグの修正を検証する。
# make_pptx / イラスト付き回答の各経路では、final は text_answer に
# '## 生成した PowerPoint / 保存先: <deck>' や '## 生成画像 / <img markup-or-status>'
# という「その場限りの成果物ノート」を追記した文字列であり、これをそのまま履歴に
# 積むと次ターンの Conductor/proposers がファイルパスや画像生成ステータスを
# 「前回回答の実質的内容」として誤読しうる(iteration 59/67/70 で対処した複数ターン
# 間の忠実性劣化と同種の回帰)上、MAX_HISTORY_CHARS の予算も無駄になる。
# 修正は _HISTORY への assistant 追記だけを text_answer に変更するもので、
# 戻り値・コンソール出力・notify_slack・_save_answer_to_file は従来通り成果物付き
# final を使い続ける。setup/conduct/fugu_answer/build_pptx/handle_image_generation/
# notify_slack/save_history_file をすべてモックしており、実際の Ollama・
# ネットワーク・サブプロセス呼び出しは一切発生しない。
_orig_af_hist = list(f._HISTORY)
_orig_af_setup = f.setup
_orig_af_conduct = f.conduct
_orig_af_fugu_answer = f.fugu_answer
_orig_af_build_pptx = f.build_pptx
_orig_af_handle_image = f.handle_image_generation
_orig_af_notify = f.notify_slack
_orig_af_save_hist = f.save_history_file


def _af_base_plan(**overrides):
    plan = {
        "mode": "moa",
        "rounds": 1,
        "selected_proposers": [],
        "search_required": False,
        "use_image_generation": False,
        "image_only": False,
        "make_pptx": False,
        "task_type": "creative",
        "reason": "offline test",
    }
    plan.update(overrides)
    return plan


try:
    f.setup = lambda: True
    f.notify_slack = lambda *a, **k: None
    f.save_history_file = lambda *a, **k: None

    # --- (1) PPTX 経路: _HISTORY にはクリーンな本文のみ、戻り値には成果物ノートも含む ---
    f._HISTORY = []
    _clean_body_pptx = "これはクリーンな本文です。PowerPoint化されます。"
    f.conduct = lambda question, history=None, office_attached=False: (
        _af_base_plan(make_pptx=True), {})
    f.fugu_answer = lambda question, plan=None, history=None: _clean_body_pptx
    f.build_pptx = lambda question, answer, out_path=None: "C:/fake/deck.pptx"

    _ret_pptx = f.ask_fugu("PPTXにして", baseline=False)

    check("ask_fugu/pptx: _HISTORY[-1]は成果物ノートを含まないクリーンな本文と一致",
          f._HISTORY[-1]["content"] == _clean_body_pptx)
    check("ask_fugu/pptx: _HISTORY[-1]に'生成した PowerPoint'が混入しない",
          "生成した PowerPoint" not in f._HISTORY[-1]["content"])
    check("ask_fugu/pptx: _HISTORY[-1]に'保存先'が混入しない",
          "保存先" not in f._HISTORY[-1]["content"])
    check("ask_fugu/pptx: 戻り値には'生成した PowerPoint'/'保存先'ノートが含まれる(従来通り)",
          "生成した PowerPoint" in _ret_pptx and "保存先" in _ret_pptx)
    check("ask_fugu/pptx: _HISTORY[0]は元の質問文をそのまま保存(user側は既存挙動)",
          f._HISTORY[0]["content"] == "PPTXにして")

    # --- (2) イラスト付き回答経路: _HISTORY にはクリーンな本文のみ ---
    f._HISTORY = []
    _clean_body_img = "これはクリーンな本文です。イラストが付きます。"
    f.conduct = lambda question, history=None, office_attached=False: (
        _af_base_plan(use_image_generation=True, image_only=False), {})
    f.fugu_answer = lambda question, plan=None, history=None: _clean_body_img
    f.handle_image_generation = lambda user_request, **k: "![img](fake.png)"

    _ret_img = f.ask_fugu("イラスト付きで説明して", baseline=False)

    check("ask_fugu/image: _HISTORY[-1]は成果物ノート('生成画像')を含まないクリーンな本文と一致",
          f._HISTORY[-1]["content"] == _clean_body_img
          and "生成画像" not in f._HISTORY[-1]["content"])
    check("ask_fugu/image: 戻り値には'生成画像'ノートが含まれる(従来通り)",
          "生成画像" in _ret_img)

    # --- (3) 通常経路(画像/PPTXなし): text_answer と final が同一 -> 既存挙動の回帰なし ---
    f._HISTORY = []
    _plain_body = "画像もPPTXも使わない通常のMoA回答本文です。"
    f.conduct = lambda question, history=None, office_attached=False: (
        _af_base_plan(), {})
    f.fugu_answer = lambda question, plan=None, history=None: _plain_body

    _ret_plain = f.ask_fugu("普通の質問", baseline=False)

    check("ask_fugu/plain: 画像/PPTXなし経路では_HISTORY[-1]・戻り値・本文がバイト一致",
          f._HISTORY[-1]["content"] == _ret_plain == _plain_body)
finally:
    f._HISTORY = _orig_af_hist
    f.setup = _orig_af_setup
    f.conduct = _orig_af_conduct
    f.fugu_answer = _orig_af_fugu_answer
    f.build_pptx = _orig_af_build_pptx
    f.handle_image_generation = _orig_af_handle_image
    f.notify_slack = _orig_af_notify
    f.save_history_file = _orig_af_save_hist

check("ask_fugu: テスト後に_HISTORYが元の状態へ復元されている",
      f._HISTORY == _orig_af_hist)

# 念のため: 本セクションはグローバル状態を try/finally で復元済みであることを確認
check("history: SESSION_SAVE はテスト後に既定値へ復元されている",
      f.SESSION_SAVE == _orig_session_save)
check("history: MAX_HISTORY_TURNS_SAVED はテスト後に既定値へ復元されている",
      f.MAX_HISTORY_TURNS_SAVED == _orig_max_hist_turns_saved)

# ---------- ask_fugu 経路2(PPTX): out_file が .pptx/.ppt 以外の場合に汎用保存へ
# フォールバックすることを検証 (2026-07-26) ----------
# 2026-07-26: 経路2(make_pptx)は build_pptx 呼び出し後に out_file を無条件で
# None化していた(旧実装)。--out が .pptx/.ppt 以外(例: notes.md)の場合でも
# pptx_out は最初から None しか build_pptx に渡らず(=既定のPPTX_OUT_DIRへ保存)、
# ユーザーが明示指定した --out 先には何も保存されない一方、下の汎用
# _save_answer_to_file も out_file=None化により実行されず、コストの高いMoA回答が
# どこにも保存されずに消えていた。iteration 186 で表面化した repl 経路の --out
# 取りこぼしバグと同種であり、「保存ステップで計算済みの回答を失わない」という
# iteration 41-47/80 の原則にも反する。修正は pptx_out が実際に消費された
# (=out_file が .pptx/.ppt だった)場合のみ out_file を None化するように変更し、
# それ以外は out_file を残して下の汎用保存に委ねるのみ。pptx_out の計算・
# build_pptx 呼び出し・final の構築(デッキノート付与)・iteration 73 の
# _HISTORY 分離(クリーンな text_answer のみ積む)はいずれも不変。
# setup/conduct/fugu_answer/build_pptx/notify_slack/save_history_file/
# _save_answer_to_file をすべてモックしており、実際の Ollama・ネットワーク・
# python-pptx バックエンド呼び出しは一切発生しない。
_orig_af4_hist = list(f._HISTORY)
_orig_af4_setup = f.setup
_orig_af4_conduct = f.conduct
_orig_af4_fugu_answer = f.fugu_answer
_orig_af4_build_pptx = f.build_pptx
_orig_af4_notify = f.notify_slack
_orig_af4_save_hist = f.save_history_file
_orig_af4_save_file = f._save_answer_to_file

try:
    f.setup = lambda: True
    f.notify_slack = lambda *a, **k: None
    f.save_history_file = lambda *a, **k: None

    _clean_body_pptx2 = "これはクリーンな本文です(経路2 out_fileテスト)。"

    def _af4_pptx_conduct(question, history=None, office_attached=False):
        return _af_base_plan(make_pptx=True), {}

    # --- (1) out_file が .pptx/.ppt 以外(notes.md): 汎用保存が発動する(新規挙動) ---
    f._HISTORY = []
    f.conduct = _af4_pptx_conduct
    f.fugu_answer = lambda question, plan=None, history=None: _clean_body_pptx2
    _build_pptx_calls_1 = []
    f.build_pptx = lambda question, answer, out_path=None: (
        _build_pptx_calls_1.append(out_path) or "C:/fake/deck1.pptx")
    _save_calls_1 = []
    f._save_answer_to_file = lambda *a, **k: _save_calls_1.append((a, k))

    _ret_1 = f.ask_fugu("PPTXにして(notes.md指定)", baseline=False,
                        out_file="notes.md")

    check("ask_fugu/経路2 out_file=notes.md: build_pptxの第3引数はNone",
          _build_pptx_calls_1 == [None])
    check("ask_fugu/経路2 out_file=notes.md: _save_answer_to_fileが1回呼ばれる",
          len(_save_calls_1) == 1)
    check("ask_fugu/経路2 out_file=notes.md: _save_answer_to_fileにpath='notes.md'が渡る",
          len(_save_calls_1) == 1 and _save_calls_1[0][0][3] == "notes.md")
    check("ask_fugu/経路2 out_file=notes.md: 保存される回答にクリーンな本文が含まれる",
          len(_save_calls_1) == 1 and _clean_body_pptx2 in _save_calls_1[0][0][1])
    check("ask_fugu/経路2 out_file=notes.md: 保存される回答に'生成した PowerPoint'ノートが含まれる",
          len(_save_calls_1) == 1 and "生成した PowerPoint" in _save_calls_1[0][0][1])
    check("ask_fugu/経路2 out_file=notes.md: _HISTORY[-1]はクリーンな本文のみ(デッキノート混入なし)",
          f._HISTORY[-1]["content"] == _clean_body_pptx2
          and "生成した PowerPoint" not in f._HISTORY[-1]["content"])
    check("ask_fugu/経路2 out_file=notes.md: 戻り値にはデッキノートが含まれる",
          "生成した PowerPoint" in _ret_1)

    # --- (2) out_file が .pptx: build_pptxが受け取り、汎用保存は呼ばれない(回帰) ---
    f._HISTORY = []
    f.conduct = _af4_pptx_conduct
    _build_pptx_calls_2 = []
    f.build_pptx = lambda question, answer, out_path=None: (
        _build_pptx_calls_2.append(out_path) or "C:/fake/deck2.pptx")
    _save_calls_2 = []
    f._save_answer_to_file = lambda *a, **k: _save_calls_2.append((a, k))

    _ret_2 = f.ask_fugu("PPTXにして(deck.pptx指定)", baseline=False,
                        out_file="deck.pptx")

    check("ask_fugu/経路2 out_file=deck.pptx(回帰): build_pptxの第3引数はout_fileそのもの",
          _build_pptx_calls_2 == ["deck.pptx"])
    check("ask_fugu/経路2 out_file=deck.pptx(回帰): _save_answer_to_fileは呼ばれない(二重保存なし)",
          len(_save_calls_2) == 0)
    check("ask_fugu/経路2 out_file=deck.pptx(回帰): _HISTORY[-1]はクリーンな本文のみ",
          f._HISTORY[-1]["content"] == _clean_body_pptx2)
    check("ask_fugu/経路2 out_file=deck.pptx(回帰): 戻り値にはデッキノートが含まれる",
          "生成した PowerPoint" in _ret_2)

    # --- (2b) out_file が大文字 .PPTX 拡張子(既存の str(...).lower() 判定の回帰) ---
    f._HISTORY = []
    f.conduct = _af4_pptx_conduct
    _build_pptx_calls_2b = []
    f.build_pptx = lambda question, answer, out_path=None: (
        _build_pptx_calls_2b.append(out_path) or "C:/fake/deck2b.pptx")
    _save_calls_2b = []
    f._save_answer_to_file = lambda *a, **k: _save_calls_2b.append((a, k))

    _ret_2b = f.ask_fugu("PPTXにして(DECK.PPTX指定)", baseline=False,
                         out_file="DECK.PPTX")

    check("ask_fugu/経路2 out_file=DECK.PPTX(大文字回帰): build_pptxの第3引数はout_fileそのもの",
          _build_pptx_calls_2b == ["DECK.PPTX"])
    check("ask_fugu/経路2 out_file=DECK.PPTX(大文字回帰): _save_answer_to_fileは呼ばれない",
          len(_save_calls_2b) == 0)

    # --- (3) out_file=None: build_pptxはNoneを受け取り、汎用保存も呼ばれない(既存挙動不変) ---
    f._HISTORY = []
    f.conduct = _af4_pptx_conduct
    _build_pptx_calls_3 = []
    f.build_pptx = lambda question, answer, out_path=None: (
        _build_pptx_calls_3.append(out_path) or "C:/fake/deck3.pptx")
    _save_calls_3 = []
    f._save_answer_to_file = lambda *a, **k: _save_calls_3.append((a, k))

    _ret_3 = f.ask_fugu("PPTXにして(out_file無指定)", baseline=False,
                        out_file=None)

    check("ask_fugu/経路2 out_file=None(既存挙動不変): build_pptxの第3引数はNone",
          _build_pptx_calls_3 == [None])
    check("ask_fugu/経路2 out_file=None(既存挙動不変): _save_answer_to_fileは呼ばれない",
          len(_save_calls_3) == 0)
    check("ask_fugu/経路2 out_file=None(既存挙動不変): _HISTORY[-1]はクリーンな本文のみ",
          f._HISTORY[-1]["content"] == _clean_body_pptx2)
    check("ask_fugu/経路2 out_file=None(既存挙動不変): 戻り値にはデッキノートが含まれる",
          "生成した PowerPoint" in _ret_3)

    # --- (4) 非PPTX経路のサニティ: make_pptx=Falseなら従来通り_save_answer_to_fileが呼ばれる ---
    f._HISTORY = []
    f.conduct = lambda question, history=None, office_attached=False: (
        _af_base_plan(), {})
    _plain_body_af4 = "画像もPPTXも使わない通常のMoA回答本文です(経路2隣接サニティ)。"
    f.fugu_answer = lambda question, plan=None, history=None: _plain_body_af4
    _save_calls_4 = []
    f._save_answer_to_file = lambda *a, **k: _save_calls_4.append((a, k))

    _ret_4 = f.ask_fugu("普通の質問(out_file指定)", baseline=False,
                        out_file="plain.md")

    check("ask_fugu/非PPTX経路サニティ: _save_answer_to_fileが1回呼ばれる(既存挙動不変)",
          len(_save_calls_4) == 1)
    check("ask_fugu/非PPTX経路サニティ: _save_answer_to_fileにpath='plain.md'が渡る",
          len(_save_calls_4) == 1 and _save_calls_4[0][0][3] == "plain.md")
    check("ask_fugu/非PPTX経路サニティ: 戻り値はクリーン本文と一致(成果物ノートなし)",
          _ret_4 == _plain_body_af4)
finally:
    f._HISTORY = _orig_af4_hist
    f.setup = _orig_af4_setup
    f.conduct = _orig_af4_conduct
    f.fugu_answer = _orig_af4_fugu_answer
    f.build_pptx = _orig_af4_build_pptx
    f.notify_slack = _orig_af4_notify
    f.save_history_file = _orig_af4_save_hist
    f._save_answer_to_file = _orig_af4_save_file

check("ask_fugu/経路2テスト: テスト後に_HISTORYが元の状態へ復元されている",
      f._HISTORY == _orig_af4_hist)

# ---------- ask_fugu 経路3(イラスト付き回答): 画像生成失敗時に __ERROR__ センチネルが
# 最終回答/保存ファイルへ漏出しないことを検証 ----------
# 2026-07-24: handle_image_generation は失敗時に内部センチネル '__ERROR__: ...' を
# 返す(L1945/1947/1955)。旧実装はこれを素通しで
# final = text_answer + '\n\n---\n## 生成画像\n' + img と連結していたため、final は
# text_answer から始まり "final.startswith('__ERROR__')" では常に False になる
# (=正常応答として扱われる)。結果としてコンソール表示・notify_slack・
# _HISTORY への追記可否判定・_save_answer_to_file 呼び出し可否判定は正しく
# 「成功」経路を通る一方、final 文字列の中に生の '__ERROR__: ...' センチネルが
# そのまま残り、画面表示/Slack通知/--out 保存ファイルへ内部マーカーが漏出して
# いた。aggregate()(iteration 9)、_critic_judge/second_opinion(iteration 15)、
# _arbitrate(iteration 20)で対処した「内部センチネル/タグをユーザ向け出力に
# 漏らさない」バグと同種。修正は route3 内で img が '__ERROR__' で始まる場合に
# プレフィックスを剥がした人間可読な失敗ノートへ置き換えるのみで、テキスト
# 本文自体を失敗扱いにはしない(text_answer は従来通り保存・保存・通知される)。
# iteration 73 で確立した「_HISTORY にはクリーンな text_answer のみを積む」
# 分離も変更しない。setup/conduct/fugu_answer/handle_image_generation/
# notify_slack/save_history_file/_save_answer_to_file をすべてモックしており、
# 実際の Ollama・ネットワーク・画像バックエンド・サブプロセス呼び出しは
# 一切発生しない。
_orig_af2_hist = list(f._HISTORY)
_orig_af2_setup = f.setup
_orig_af2_conduct = f.conduct
_orig_af2_fugu_answer = f.fugu_answer
_orig_af2_handle_image = f.handle_image_generation
_orig_af2_notify = f.notify_slack
_orig_af2_save_hist = f.save_history_file
_orig_af2_save_file = f._save_answer_to_file

try:
    f.setup = lambda: True
    f.notify_slack = lambda *a, **k: None
    f.save_history_file = lambda *a, **k: None
    _save_file_calls = []
    f._save_answer_to_file = lambda *a, **k: _save_file_calls.append((a, k))

    # --- (4) イラスト付き回答経路で画像生成が失敗: センチネルが漏出しない ---
    f._HISTORY = []
    _clean_body_imgfail = "これはクリーンな本文です。画像生成は失敗します。"
    _img_sentinel = "__ERROR__: 画像生成に失敗しました（バックエンドが画像を返しませんでした）。"
    f.conduct = lambda question, history=None, office_attached=False: (
        _af_base_plan(use_image_generation=True, image_only=False), {})
    f.fugu_answer = lambda question, plan=None, history=None: _clean_body_imgfail
    f.handle_image_generation = lambda user_request, **k: _img_sentinel

    _ret_imgfail = f.ask_fugu("イラスト付きで説明して(画像失敗)", baseline=False,
                              out_file="C:/fake/out_imgfail.md")

    check("ask_fugu/image失敗: 戻り値に内部センチネル'__ERROR__'が漏出しない",
          "__ERROR__" not in _ret_imgfail)
    check("ask_fugu/image失敗: 戻り値にクリーンな本文が保持されている",
          _clean_body_imgfail in _ret_imgfail)
    check("ask_fugu/image失敗: 戻り値の'## 生成画像'配下に人間可読な失敗ノートがある",
          "## 生成画像" in _ret_imgfail and "画像生成に失敗しました" in _ret_imgfail)
    check("ask_fugu/image失敗: _HISTORY[-1]はクリーンなtext_answerと一致(センチネル/ノート混入なし)",
          f._HISTORY[-1]["content"] == _clean_body_imgfail)
    check("ask_fugu/image失敗: _HISTORY[-1]に'__ERROR__'/'生成画像'ノートが混入しない",
          "__ERROR__" not in f._HISTORY[-1]["content"]
          and "生成画像" not in f._HISTORY[-1]["content"])
    check("ask_fugu/image失敗: テキスト本文は失敗扱いされず_save_answer_to_fileが呼ばれる",
          len(_save_file_calls) == 1)

    # --- (5) 回帰: 画像生成成功時は従来通りセンチネル処理を経由しない(iteration73ケース2と同型) ---
    f._HISTORY = []
    _clean_body_imgok = "これはクリーンな本文です。画像生成は成功します。"
    f.fugu_answer = lambda question, plan=None, history=None: _clean_body_imgok
    f.handle_image_generation = lambda user_request, **k: "![img](fake.png)"

    _ret_imgok = f.ask_fugu("イラスト付きで説明して(画像成功)", baseline=False)

    check("ask_fugu/image成功(回帰): 戻り値に'## 生成画像'配下の成功マークアップがバイト一致で含まれる",
          "## 生成画像\n![img](fake.png)" in _ret_imgok)
    check("ask_fugu/image成功(回帰): 戻り値に'__ERROR__'が含まれない",
          "__ERROR__" not in _ret_imgok)
finally:
    f._HISTORY = _orig_af2_hist
    f.setup = _orig_af2_setup
    f.conduct = _orig_af2_conduct
    f.fugu_answer = _orig_af2_fugu_answer
    f.handle_image_generation = _orig_af2_handle_image
    f.notify_slack = _orig_af2_notify
    f.save_history_file = _orig_af2_save_hist
    f._save_answer_to_file = _orig_af2_save_file

check("ask_fugu/image失敗テスト: テスト後に_HISTORYが元の状態へ復元されている",
      f._HISTORY == _orig_af2_hist)

# ---------- ask_fugu 経路1(画像のみ): 画像生成失敗時に __ERROR__ センチネルが
# コンソール表示へ漏出しないことを検証 ----------
# 2026-07-24: 経路1(use_image_generation かつ image_only)は handle_image_generation
# が失敗時に返す内部センチネル '__ERROR__: ...' を print(result) でそのまま
# '===== 画像生成結果 =====' の下に出力していた。プレーンパス(最終回答表示、
# final.startswith('__ERROR__') 時は「生成に失敗しました:」と整形して表示)や、
# 兄弟経路である経路3のイラスト付き回答(iteration 99 で対処済み)には同種の
# 人間可読化があったが、経路1だけ未対応だった。aggregate()(iteration 9)、
# _critic_judge/second_opinion(iteration 15)、_arbitrate(iteration 20)で対処した
# 「内部センチネル/タグをユーザ向け出力に漏らさない」バグと同種。修正はコンソール
# 表示のみで、result が '__ERROR__' で始まる場合にプレフィックスを剥がした
# 人間可読な失敗ノートへ置き換える。notify_slack(iteration 119 の失敗アイコン
# 判定用)には生の result を渡し続け、_save_answer_to_file のゲート
# (iteration 80 のエラー時未保存)・関数の戻り値(プレーンパスの
# '__ERROR__'-on-failure 契約と同型)も生の result のまま変更しない。
# setup/conduct/handle_image_generation/notify_slack/_save_answer_to_file を
# すべてモックしており、実際の Ollama・ネットワーク・画像バックエンド呼び出しは
# 一切発生しない。
_orig_af3_hist = list(f._HISTORY)
_orig_af3_setup = f.setup
_orig_af3_conduct = f.conduct
_orig_af3_handle_image = f.handle_image_generation
_orig_af3_notify = f.notify_slack
_orig_af3_save_file = f._save_answer_to_file

try:
    f.setup = lambda: True

    # --- (1) 画像のみ経路で画像生成が失敗: コンソールにセンチネルが漏出しない ---
    f._HISTORY = []
    _notify_calls_r1 = []
    f.notify_slack = lambda *a, **k: _notify_calls_r1.append(a)
    _save_calls_r1 = []
    f._save_answer_to_file = lambda *a, **k: _save_calls_r1.append((a, k))
    _img_only_sentinel = "__ERROR__: 画像生成は無効化されています（IMAGE_BACKEND=off）。"
    f.conduct = lambda question, history=None, office_attached=False: (
        _af_base_plan(use_image_generation=True, image_only=True), {})
    f.handle_image_generation = lambda user_request, **k: _img_only_sentinel

    _out_r1 = io.StringIO()
    with contextlib.redirect_stdout(_out_r1):
        _ret_r1 = f.ask_fugu("猫の絵だけ描いて(失敗)", baseline=False,
                             out_file="C:/fake/out_imgonly_fail.md")
    _printed_r1 = _out_r1.getvalue()

    check("ask_fugu/経路1失敗: コンソール出力に生の'__ERROR__'トークンが含まれない",
          "__ERROR__" not in _printed_r1)
    check("ask_fugu/経路1失敗: コンソール出力に人間可読な失敗文言が含まれる",
          "画像生成に失敗しました" in _printed_r1)
    check("ask_fugu/経路1失敗: notify_slackには生のresult('__ERROR__'始まり)が渡る",
          len(_notify_calls_r1) == 1 and _notify_calls_r1[0][1] == _img_only_sentinel
          and _notify_calls_r1[0][1].startswith("__ERROR__"))
    check("ask_fugu/経路1失敗: out_file指定済みでも_save_answer_to_fileは呼ばれない",
          len(_save_calls_r1) == 0)
    check("ask_fugu/経路1失敗: 戻り値は生のresultのまま('__ERROR__'始まり)",
          _ret_r1 == _img_only_sentinel and _ret_r1.startswith("__ERROR__"))

    # --- (2) 回帰: 画像のみ経路で画像生成が成功時はメッセージがそのまま表示される ---
    f._HISTORY = []
    _notify_calls_r2 = []
    f.notify_slack = lambda *a, **k: _notify_calls_r2.append(a)
    _save_calls_r2 = []
    f._save_answer_to_file = lambda *a, **k: _save_calls_r2.append((a, k))
    _img_only_ok = "画像を生成しました。保存先: C:/fake/cat.png"
    f.handle_image_generation = lambda user_request, **k: _img_only_ok

    _out_r2 = io.StringIO()
    with contextlib.redirect_stdout(_out_r2):
        _ret_r2 = f.ask_fugu("猫の絵だけ描いて(成功)", baseline=False,
                             out_file="C:/fake/out_imgonly_ok.md")
    _printed_r2 = _out_r2.getvalue()

    check("ask_fugu/経路1成功(回帰): コンソール出力に成功メッセージがそのまま含まれる",
          _img_only_ok in _printed_r2)
    check("ask_fugu/経路1成功(回帰): コンソール出力に'画像生成に失敗しました'が混入しない",
          "画像生成に失敗しました" not in _printed_r2)
    check("ask_fugu/経路1成功(回帰): notify_slackには成功メッセージが渡る",
          len(_notify_calls_r2) == 1 and _notify_calls_r2[0][1] == _img_only_ok)
    check("ask_fugu/経路1成功(回帰): out_file指定時は_save_answer_to_fileが呼ばれる",
          len(_save_calls_r2) == 1)
    check("ask_fugu/経路1成功(回帰): 戻り値は成功メッセージと一致",
          _ret_r2 == _img_only_ok)
finally:
    f._HISTORY = _orig_af3_hist
    f.setup = _orig_af3_setup
    f.conduct = _orig_af3_conduct
    f.handle_image_generation = _orig_af3_handle_image
    f.notify_slack = _orig_af3_notify
    f._save_answer_to_file = _orig_af3_save_file

check("ask_fugu/経路1テスト: テスト後に_HISTORYが元の状態へ復元されている",
      f._HISTORY == _orig_af3_hist)

# ---------- 画像プロンプト起草チェーン: _sd_prompt_from_request / moa_image_prompt /
# author_image_prompt (2026-07-24) ----------
# この3関数はSDXL画像プロンプト起草パイプラインを構成し、ask_fugu 経路1(画像のみ)・
# 経路3(イラスト付き回答、直上のテスト)・build_pptx のヒーロー/スライド画像
# (handle_image_generation L1970経由でauthor_image_promptを呼ぶ)の全てから使われる
# が、grep上は一貫してモックされるだけで直接のテストが皆無だった。
# moa_image_prompt/_sd_prompt_from_request はどちらも、ask() が壊れたJSONや
# '__ERROR__: ...' センチネルを返した場合に extract_json(raw) が None を返す
# (または j.get('prompt') が falsy になる)ことに暗黙的に依存してフォールバック
# している。他の全ての類似箇所 ―― iteration 9 の aggregate、iteration 15 の
# _critic_judge/second_opinion、iteration 20 の _arbitrate、iteration 52 の
# _sc_sample、直上(iteration 99)の ask_fugu 経路3 ―― はセンチネル漏出を明示的に
# テスト済みだったのに対し、ここだけ空白地帯だった。
# f.ask のみをモックし、extract_json は本物を実行してセンチネル拒否経路を実際に
# 通す(extract_jsonそのものはモックしない)。urllib.request.urlopen と
# f.subprocess.run にも「呼ばれたら即AssertionError」の番人を仕込み
# (gotcha #8 / iteration 38・39・76の流儀)、モック漏れで実ネットワーク/
# サブプロセスへ落ちないことを保証する。実際の Ollama・画像バックエンド・
# ネットワーク・subprocess 呼び出しは一切発生しない。
_orig_ip_ask = f.ask
_orig_ip_proposers = f.PROPOSERS
_orig_ip_conductor = f.CONDUCTOR
_orig_ip_moa_flag = f.IMAGE_PROMPT_MOA
_orig_ip_translate = f.IMAGE_TRANSLATE_PROMPT
_orig_ip_panel = f.IMAGE_PROMPT_PANEL
_orig_ip_urlopen = urllib.request.urlopen
_orig_ip_subprocess_run = f.subprocess.run


def _ip_no_network_urlopen(*a, **k):
    raise AssertionError("image-prompt起草チェーン: モック漏れで実urlopen(ネットワーク)が呼ばれた")


def _ip_no_subprocess_run(*a, **k):
    raise AssertionError("image-prompt起草チェーン: モック漏れで実subprocess.runが呼ばれた")


_IP_SENTINEL = "__ERROR__: HTTP Error 500 (模擬)"

try:
    urllib.request.urlopen = _ip_no_network_urlopen
    f.subprocess.run = _ip_no_subprocess_run
    f.PROPOSERS = ["ip_p1", "ip_p2", "ip_p3"]
    f.CONDUCTOR = "ip_conductor"
    f.IMAGE_PROMPT_PANEL = 2

    _ip_calls = []

    def _make_ip_ask(moa_responses=None, merge_response=None, sd_response=None):
        moa_responses = moa_responses or {}

        def _fake_ask(model, messages, temperature, think=None, fmt=None,
                      label=None, num_predict=None, num_ctx=None):
            _ip_calls.append((model, label))
            if label == "img-prompt":
                return sd_response
            if label == "img-merge":
                return merge_response
            if label == "img-moa":
                return moa_responses.get(model, "__ERROR__: no response configured (test)")
            raise AssertionError(
                f"image-prompt起草チェーン: 想定外のlabel={label!r} model={model!r}")
        return _fake_ask

    # === _sd_prompt_from_request ===
    # (a) IMAGE_TRANSLATE_PROMPT=False: askを一切呼ばず(user_request, "")をそのまま返す
    f.IMAGE_TRANSLATE_PROMPT = False
    _ip_calls.clear()
    f.ask = _make_ip_ask()
    _sd_a = f._sd_prompt_from_request("猫の絵")
    check("_sd_prompt_from_request: IMAGE_TRANSLATE_PROMPT=Falseはaskを呼ばない",
          _ip_calls == [])
    check("_sd_prompt_from_request: IMAGE_TRANSLATE_PROMPT=Falseは(user_request,'')をそのまま返す",
          _sd_a == ("猫の絵", ""))

    # (b) True + 正当なJSON -> パース結果の(prompt,negative)を返す
    f.IMAGE_TRANSLATE_PROMPT = True
    _ip_calls.clear()
    f.ask = _make_ip_ask(sd_response='{"prompt": "a cat, masterpiece", "negative": "blurry"}')
    _sd_b = f._sd_prompt_from_request("猫の絵")
    check("_sd_prompt_from_request: 正当なJSONはパースした(prompt,negative)を返す",
          _sd_b == ("a cat, masterpiece", "blurry"))
    check("_sd_prompt_from_request: 正当なJSON経路ではaskが1回だけ呼ばれる(label=img-prompt)",
          _ip_calls == [("ip_conductor", "img-prompt")])

    # (c) 壊れたJSON(パース不能) -> (user_request, "")へフォールバック
    _ip_calls.clear()
    f.ask = _make_ip_ask(sd_response="not json at all, just prose")
    _sd_c = f._sd_prompt_from_request("犬の絵")
    check("_sd_prompt_from_request: パース不能出力は(user_request,'')へフォールバックする",
          _sd_c == ("犬の絵", ""))
    check("_sd_prompt_from_request: パース不能フォールバックの戻り値に'__ERROR__'を含まない",
          "__ERROR__" not in _sd_c[0] and "__ERROR__" not in _sd_c[1])

    # (d) '__ERROR__'センチネル -> (user_request, "")へフォールバックし生センチネルを漏らさない
    _ip_calls.clear()
    f.ask = _make_ip_ask(sd_response=_IP_SENTINEL)
    _sd_d = f._sd_prompt_from_request("鳥の絵")
    check("_sd_prompt_from_request: '__ERROR__'センチネルは(user_request,'')へフォールバックする",
          _sd_d == ("鳥の絵", ""))
    check("_sd_prompt_from_request: '__ERROR__'センチネルは戻り値に漏出しない",
          "__ERROR__" not in _sd_d[0] and "__ERROR__" not in _sd_d[1])

    # === moa_image_prompt ===
    _valid_ip_p1 = {"prompt": "cand from p1", "negative": "neg p1"}
    _valid_ip_p2 = {"prompt": "cand from p2", "negative": "neg p2"}
    _valid_ip_p3 = {"prompt": "cand from p3", "negative": "neg p3"}

    # (e) 使える候補が1件のみ -> マージaskを一切呼ばずその候補をそのまま返す(短絡)
    _ip_calls.clear()
    f.ask = _make_ip_ask(moa_responses={
        "ip_p1": json.dumps(_valid_ip_p1),
        "ip_p2": "__ERROR__: proposer failed (test)",
    })
    _moa_e = f.moa_image_prompt("猫を描いて")
    check("moa_image_prompt: 使える候補1件は短絡しその(prompt,negative)をそのまま返す",
          _moa_e == ("cand from p1", "neg p1"))
    check("moa_image_prompt: 候補1件のときマージask(img-merge)は呼ばれない",
          all(lbl != "img-merge" for _m, lbl in _ip_calls))
    check("moa_image_prompt: 候補1件のときも起草askはpanel全員(既定2件)に対して行われる",
          [m for m, lbl in _ip_calls if lbl == "img-moa"] == ["ip_p1", "ip_p2"])

    # (f) 使える候補が2件以上 -> マージaskが発行されその結果を返す
    _ip_calls.clear()
    _merge_result = {"prompt": "merged prompt", "negative": "merged neg"}
    f.ask = _make_ip_ask(
        moa_responses={"ip_p1": json.dumps(_valid_ip_p1), "ip_p2": json.dumps(_valid_ip_p2)},
        merge_response=json.dumps(_merge_result))
    _moa_f = f.moa_image_prompt("猫を描いて")
    check("moa_image_prompt: 候補2件以上はマージaskの結果を返す",
          _moa_f == ("merged prompt", "merged neg"))
    check("moa_image_prompt: 候補2件以上のときマージaskがCONDUCTORに対して発行される",
          ("ip_conductor", "img-merge") in _ip_calls)

    # (g) マージaskがパース不能/センチネル -> 最初の候補の(prompt,negative)へフォールバック
    _ip_calls.clear()
    f.ask = _make_ip_ask(
        moa_responses={"ip_p1": json.dumps(_valid_ip_p1), "ip_p2": json.dumps(_valid_ip_p2)},
        merge_response="__ERROR__: merge failed (test)")
    _moa_g = f.moa_image_prompt("猫を描いて")
    check("moa_image_prompt: マージask失敗時は最初の候補の(prompt,negative)へフォールバックする",
          _moa_g == ("cand from p1", "neg p1"))
    check("moa_image_prompt: マージask失敗フォールバックは戻り値に'__ERROR__'を含まない",
          "__ERROR__" not in _moa_g[0] and "__ERROR__" not in _moa_g[1])

    # (h) 全proposerが失敗('__ERROR__'/非JSON、使える候補0件) -> Noneを返す
    _ip_calls.clear()
    f.ask = _make_ip_ask(moa_responses={
        "ip_p1": "__ERROR__: fail1 (test)",
        "ip_p2": "not json either",
    })
    _moa_h = f.moa_image_prompt("猫を描いて")
    check("moa_image_prompt: 全proposer失敗(候補0件)はNoneを返す", _moa_h is None)
    check("moa_image_prompt: 全proposer失敗時はマージask(img-merge)を呼ばない",
          all(lbl != "img-merge" for _m, lbl in _ip_calls))

    # (i) panelはPROPOSERSの構成員へフィルタされ、IMAGE_PROMPT_PANEL件で打ち切られる
    _ip_calls.clear()
    f.ask = _make_ip_ask(moa_responses={
        "ip_p1": json.dumps(_valid_ip_p1),
        "ip_p2": json.dumps(_valid_ip_p2),
        "ip_p3": json.dumps(_valid_ip_p3),
    })
    f.moa_image_prompt("猫を描いて", panel=["ip_p2", "not_a_proposer", "ip_p1", "ip_p3"])
    check("moa_image_prompt: panelはPROPOSERS外('not_a_proposer')を除外し"
          "IMAGE_PROMPT_PANEL(=2)件で打ち切って実際にaskする",
          [m for m, lbl in _ip_calls if lbl == "img-moa"] == ["ip_p2", "ip_p1"])

    # === author_image_prompt ===
    # moa_image_prompt/_sd_prompt_from_request 自体をモックしてルーティングのみを検証する
    # (ask経由のセンチネル拒否は上のmoa_image_prompt/_sd_prompt_from_requestテストで
    # 既に直接カバー済みのため)。
    _orig_ip_moa_fn = f.moa_image_prompt
    _orig_ip_sd_fn = f._sd_prompt_from_request
    try:
        _sd_calls_author = []

        def _fake_sd_author(user_request):
            _sd_calls_author.append(user_request)
            return ("SD_FALLBACK_PROMPT", "SD_FALLBACK_NEG")

        f._sd_prompt_from_request = _fake_sd_author

        # (j) IMAGE_PROMPT_MOA=True かつ moa_image_prompt がタプルを返す
        # -> それをそのまま返し、単独翻訳フォールバックは一切呼ばれない
        f.IMAGE_PROMPT_MOA = True
        _sd_calls_author.clear()
        f.moa_image_prompt = lambda base_text, panel=None: ("MOA_PROMPT", "MOA_NEG")
        _auth_j = f.author_image_prompt("何か描いて")
        check("author_image_prompt: MOA=Trueでmoa_image_promptがタプルを返せばそれを返す",
              _auth_j == ("MOA_PROMPT", "MOA_NEG"))
        check("author_image_prompt: moa_image_prompt成功時は_sd_prompt_from_requestを呼ばない",
              _sd_calls_author == [])

        # (k) moa_image_promptがNoneを返す -> _sd_prompt_from_requestへフォールバックする
        _sd_calls_author.clear()
        f.moa_image_prompt = lambda base_text, panel=None: None
        _auth_k = f.author_image_prompt("何か描いて")
        check("author_image_prompt: moa_image_promptがNoneなら_sd_prompt_from_requestへ"
              "フォールバックする",
              _auth_k == ("SD_FALLBACK_PROMPT", "SD_FALLBACK_NEG"))
        check("author_image_prompt: Noneフォールバック時に_sd_prompt_from_requestが実際に呼ばれる",
              _sd_calls_author == ["何か描いて"])

        # (l) IMAGE_PROMPT_MOA=False -> MoA起草を一切経由せず直接_sd_prompt_from_requestへ行く
        f.IMAGE_PROMPT_MOA = False
        _sd_calls_author.clear()

        def _moa_forbidden_author(base_text, panel=None):
            raise AssertionError(
                "author_image_prompt: IMAGE_PROMPT_MOA=Falseなのにmoa_image_promptが呼ばれた")

        f.moa_image_prompt = _moa_forbidden_author
        _auth_l = f.author_image_prompt("何か描いて")
        check("author_image_prompt: MOA=Falseは直接_sd_prompt_from_requestへ行く"
              "(MoA起草askを経由しない)",
              _auth_l == ("SD_FALLBACK_PROMPT", "SD_FALLBACK_NEG"))
        check("author_image_prompt: MOA=Falseのとき_sd_prompt_from_requestが呼ばれる",
              _sd_calls_author == ["何か描いて"])
    finally:
        f.moa_image_prompt = _orig_ip_moa_fn
        f._sd_prompt_from_request = _orig_ip_sd_fn
finally:
    f.ask = _orig_ip_ask
    f.PROPOSERS = _orig_ip_proposers
    f.CONDUCTOR = _orig_ip_conductor
    f.IMAGE_PROMPT_MOA = _orig_ip_moa_flag
    f.IMAGE_TRANSLATE_PROMPT = _orig_ip_translate
    f.IMAGE_PROMPT_PANEL = _orig_ip_panel
    urllib.request.urlopen = _orig_ip_urlopen
    f.subprocess.run = _orig_ip_subprocess_run

check("image-prompt起草チェーン: テスト後にaskが元へ復元されている", f.ask == _orig_ip_ask)
check("image-prompt起草チェーン: テスト後にPROPOSERSが元へ復元されている",
      f.PROPOSERS == _orig_ip_proposers)
check("image-prompt起草チェーン: テスト後にCONDUCTORが元へ復元されている",
      f.CONDUCTOR == _orig_ip_conductor)
check("image-prompt起草チェーン: テスト後にIMAGE_PROMPT_MOAが元へ復元されている",
      f.IMAGE_PROMPT_MOA == _orig_ip_moa_flag)
check("image-prompt起草チェーン: テスト後にIMAGE_TRANSLATE_PROMPTが元へ復元されている",
      f.IMAGE_TRANSLATE_PROMPT == _orig_ip_translate)
check("image-prompt起草チェーン: テスト後にIMAGE_PROMPT_PANELが元へ復元されている",
      f.IMAGE_PROMPT_PANEL == _orig_ip_panel)
check("image-prompt起草チェーン: テスト後にurllib.request.urlopenが元へ復元されている",
      urllib.request.urlopen == _orig_ip_urlopen)
check("image-prompt起草チェーン: テスト後にf.subprocess.runが元へ復元されている",
      f.subprocess.run == _orig_ip_subprocess_run)

# ---------- resolve_models: PROPOSERS/AGGREGATOR/CONDUCTOR プール構成 ----------
# resolve_models() (fugu_local.py L1084) は SC投票パネル(solve_verifiableが
# `[m for m in REASONING_MODELS if m in PROPOSERS]` で組む)と _arbitrate のアービター
# チェーン(PROPOSERSから構成)の土台になるが、これまで直接のテストが皆無だった。
# installed_models()(ネットワーク越しに /api/tags を叩く)と pull()(実際に
# `ollama pull` を subprocess.run で起動する)の両方を必ずモックする。さらに
# モック漏れを即座に可視化するため、f.subprocess.run と urllib.request.urlopen も
# 「呼ばれたら即AssertionError」の番人に差し替える(iteration 38/39 と同じ流儀)。
# 各ケースの期待(pool, agg, cond) はソースを手でトレースして導出したもの(推測ではない)。
# 静的レビューでは欠陥は見つからなかったため、ここは既存挙動の特性固定化(characterization)
# であり、意外な挙動があってもここでは修正せず iteration 48/66/71 の
# surface-don't-fix 方針を踏襲して明示するに留める。

_orig_rm_desired_proposers = f.DESIRED_PROPOSERS
_orig_rm_desired_aggregator = f.DESIRED_AGGREGATOR
_orig_rm_desired_conductor = f.DESIRED_CONDUCTOR
_orig_rm_fallback_model = f.FALLBACK_MODEL
_orig_rm_installed_models = f.installed_models
_orig_rm_pull = f.pull
_orig_rm_subprocess_run = f.subprocess.run
_orig_rm_urlopen = urllib.request.urlopen


def _rm_no_subprocess_run(*a, **k):
    raise AssertionError("resolve_models: モック漏れで実subprocess.run(ollama pull起動)が呼ばれた")


def _rm_no_urlopen(*a, **k):
    raise AssertionError("resolve_models: モック漏れで実urlopen(ネットワーク)が呼ばれた")


def _rm_installed_factory(fixed_list):
    """installed_models() は resolve_models() 内で複数回呼ばれるため、
    呼び出し回数に関わらず常に同じリストを返す固定スタブにする。"""
    def _fake():
        return list(fixed_list)
    return _fake


def _rm_pull_factory(succeed_set, calls_log):
    """pull() の代役。実subprocessは一切呼ばない。呼び出されたモデル名を記録し、
    succeed_set に含まれていれば成功(True)、それ以外は失敗(False)を返す。"""
    def _fake(model):
        calls_log.append(model)
        return model in succeed_set
    return _fake


try:
    f.subprocess.run = _rm_no_subprocess_run
    urllib.request.urlopen = _rm_no_urlopen

    # --- (1) 全DESIRED_PROPOSERSが導入済み -> pool==DESIRED_PROPOSERS、pullは一度も呼ばれない ---
    f.DESIRED_PROPOSERS = ["fake-p1:1b", "fake-p2:1b", "fake-p3:1b"]
    f.DESIRED_AGGREGATOR = "fake-p2:1b"   # 既にpool内 -> agg用のpullは発生しない
    f.DESIRED_CONDUCTOR = "fake-p1:1b"    # 既にpool内 -> cond用のpullは発生しない
    f.FALLBACK_MODEL = "fake-fallback:1b"
    _calls1 = []
    try:
        f.installed_models = _rm_installed_factory(
            ["fake-p1:1b", "fake-p2:1b", "fake-p3:1b"])
        f.pull = _rm_pull_factory(set(), _calls1)  # 何が来ても失敗扱い(呼ばれない想定)
        _pool1, _agg1, _cond1 = f.resolve_models()
        check("resolve_models(1): 全proposer導入済みならpoolはDESIRED_PROPOSERSと一致",
              _pool1 == ["fake-p1:1b", "fake-p2:1b", "fake-p3:1b"])
        check("resolve_models(1): aggはpool内既存モデル(fake-p2)のまま",
              _agg1 == "fake-p2:1b")
        check("resolve_models(1): condはpool内既存モデル(fake-p1)のまま",
              _cond1 == "fake-p1:1b")
        check("resolve_models(1): pull()は一度も呼ばれない", _calls1 == [])
    finally:
        f.installed_models = _orig_rm_installed_models
        f.pull = _orig_rm_pull

    # --- (2) proposerが1つ未導入だがpull成功 -> poolに追加される ---
    f.DESIRED_PROPOSERS = ["fake-p1:1b", "fake-p2:1b", "fake-p3:1b"]
    f.DESIRED_AGGREGATOR = "fake-p1:1b"   # 既にpool内
    f.DESIRED_CONDUCTOR = "fake-p3:1b"    # 既にpool内
    _calls2 = []
    try:
        f.installed_models = _rm_installed_factory(["fake-p1:1b", "fake-p3:1b"])  # p2欠落
        f.pull = _rm_pull_factory({"fake-p2:1b"}, _calls2)
        _pool2, _agg2, _cond2 = f.resolve_models()
        check("resolve_models(2): pull成功したproposerはpoolに追加される(順序も維持)",
              _pool2 == ["fake-p1:1b", "fake-p2:1b", "fake-p3:1b"])
        check("resolve_models(2): pullは欠落していたp2に対してのみ呼ばれる",
              _calls2 == ["fake-p2:1b"])
        check("resolve_models(2): agg/condはpool内既存モデルのまま(追加pull無し)",
              (_agg2, _cond2) == ("fake-p1:1b", "fake-p3:1b"))
    finally:
        f.installed_models = _orig_rm_installed_models
        f.pull = _orig_rm_pull

    # --- (3) proposerが1つ未導入かつpull失敗 -> poolから除外され、他は無傷 ---
    _calls3 = []
    try:
        f.installed_models = _rm_installed_factory(["fake-p1:1b", "fake-p3:1b"])  # p2欠落
        f.pull = _rm_pull_factory(set(), _calls3)  # p2のpullを失敗させる
        _pool3, _agg3, _cond3 = f.resolve_models()
        check("resolve_models(3): pull失敗したproposerはpoolから除外される",
              _pool3 == ["fake-p1:1b", "fake-p3:1b"])
        check("resolve_models(3): 他のproposer(p1,p3)は無傷のまま2件だけ残る",
              "fake-p1:1b" in _pool3 and "fake-p3:1b" in _pool3 and len(_pool3) == 2)
        check("resolve_models(3): pullは欠落p2に対してのみ呼ばれる(1回)",
              _calls3 == ["fake-p2:1b"])
    finally:
        f.installed_models = _orig_rm_installed_models
        f.pull = _orig_rm_pull

    # --- (4) aggregatorが未導入かつpool外、pull失敗 -> agg=pool[0]で代用 ---
    f.DESIRED_PROPOSERS = ["fake-p1:1b", "fake-p2:1b"]
    f.DESIRED_AGGREGATOR = "fake-agg:1b"   # 未導入かつpool外
    f.DESIRED_CONDUCTOR = "fake-p2:1b"     # 既にpool内(cond側の分岐に影響させない)
    _calls4 = []
    try:
        f.installed_models = _rm_installed_factory(["fake-p1:1b", "fake-p2:1b"])
        f.pull = _rm_pull_factory(set(), _calls4)  # aggregatorのpullを失敗させる
        _pool4, _agg4, _cond4 = f.resolve_models()
        check("resolve_models(4): proposer側は無傷", _pool4 == ["fake-p1:1b", "fake-p2:1b"])
        check("resolve_models(4): aggregatorのpull失敗時はpool[0]で代用される",
              _agg4 == "fake-p1:1b")
        check("resolve_models(4): condはpool内既存モデルのまま", _cond4 == "fake-p2:1b")
        check("resolve_models(4): pullはaggregatorに対してのみ呼ばれる(1回)",
              _calls4 == ["fake-agg:1b"])
    finally:
        f.installed_models = _orig_rm_installed_models
        f.pull = _orig_rm_pull

    # --- (5) conductorがaggregatorと重複(cond == agg) -> conductor用のpullは発生しない ---
    #     ("cond in pool" 分岐は上の(1)〜(4)でcond=既存proposerとして既に踏んでいるため、
    #      ここでは判定式のもう一方の分岐 cond == agg を明示的に踏む)
    f.DESIRED_PROPOSERS = ["fake-p1:1b", "fake-p2:1b"]
    f.DESIRED_AGGREGATOR = "fake-shared:1b"  # 未導入かつpool外 -> pullが必要
    f.DESIRED_CONDUCTOR = "fake-shared:1b"   # aggと同一モデル
    _calls5 = []
    try:
        f.installed_models = _rm_installed_factory(["fake-p1:1b", "fake-p2:1b"])
        f.pull = _rm_pull_factory({"fake-shared:1b"}, _calls5)  # aggregatorのpullは成功させる
        _pool5, _agg5, _cond5 = f.resolve_models()
        check("resolve_models(5): pull成功したaggregatorはpoolには追加されない(仕様通り)",
              _pool5 == ["fake-p1:1b", "fake-p2:1b"])
        check("resolve_models(5): aggはpull成功した共有モデルになる", _agg5 == "fake-shared:1b")
        check("resolve_models(5): cond==aggの場合はcondもその値に一致する",
              _cond5 == "fake-shared:1b")
        check("resolve_models(5): pullはaggregator用に1回だけ呼ばれる(conductor用の追加pullは無い)",
              _calls5 == ["fake-shared:1b"])
    finally:
        f.installed_models = _orig_rm_installed_models
        f.pull = _orig_rm_pull

    # --- (6) 全滅(installed_models==[]・proposer/agg/condの全pull失敗)後、
    #     FALLBACK_MODELのpullのみ成功 -> pool==[FALLBACK_MODEL]、agg/condも後埋めされる ---
    f.DESIRED_PROPOSERS = ["fake-p1:1b", "fake-p2:1b"]
    f.DESIRED_AGGREGATOR = "fake-agg:1b"
    f.DESIRED_CONDUCTOR = "fake-cond:1b"
    f.FALLBACK_MODEL = "fake-fallback:1b"
    _calls6 = []
    try:
        f.installed_models = _rm_installed_factory([])  # 何も導入されていない
        f.pull = _rm_pull_factory({"fake-fallback:1b"}, _calls6)  # FALLBACKのみ成功
        _pool6, _agg6, _cond6 = f.resolve_models()
        check("resolve_models(6): 全滅時はFALLBACK_MODEL単体のpoolになる",
              _pool6 == ["fake-fallback:1b"])
        check("resolve_models(6): aggもFALLBACK_MODELへ後埋めされる", _agg6 == "fake-fallback:1b")
        check("resolve_models(6): condもFALLBACK_MODELへ後埋めされる", _cond6 == "fake-fallback:1b")
        check("resolve_models(6): pullはp1,p2,agg,cond,FALLBACKの順で全て試行される",
              _calls6 == ["fake-p1:1b", "fake-p2:1b", "fake-agg:1b", "fake-cond:1b",
                          "fake-fallback:1b"])
    finally:
        f.installed_models = _orig_rm_installed_models
        f.pull = _orig_rm_pull
finally:
    f.DESIRED_PROPOSERS = _orig_rm_desired_proposers
    f.DESIRED_AGGREGATOR = _orig_rm_desired_aggregator
    f.DESIRED_CONDUCTOR = _orig_rm_desired_conductor
    f.FALLBACK_MODEL = _orig_rm_fallback_model
    f.installed_models = _orig_rm_installed_models
    f.pull = _orig_rm_pull
    f.subprocess.run = _orig_rm_subprocess_run
    urllib.request.urlopen = _orig_rm_urlopen

check("resolve_models: テスト後にDESIRED_PROPOSERSが元の状態へ復元されている",
      f.DESIRED_PROPOSERS == _orig_rm_desired_proposers)
check("resolve_models: テスト後にDESIRED_AGGREGATORが元の状態へ復元されている",
      f.DESIRED_AGGREGATOR == _orig_rm_desired_aggregator)
check("resolve_models: テスト後にDESIRED_CONDUCTORが元の状態へ復元されている",
      f.DESIRED_CONDUCTOR == _orig_rm_desired_conductor)
check("resolve_models: テスト後にFALLBACK_MODELが元の状態へ復元されている",
      f.FALLBACK_MODEL == _orig_rm_fallback_model)
check("resolve_models: テスト後にinstalled_modelsが元の状態へ復元されている",
      f.installed_models == _orig_rm_installed_models)
check("resolve_models: テスト後にpullが元の状態へ復元されている",
      f.pull == _orig_rm_pull)
check("resolve_models: テスト後にsubprocess.runが元の状態へ復元されている",
      f.subprocess.run == _orig_rm_subprocess_run)
check("resolve_models: テスト後にurllib.request.urlopenが元の状態へ復元されている",
      urllib.request.urlopen == _orig_rm_urlopen)

# ---------- installed_models(): /api/tags 応答の直接テスト (2026-07-25, iter152) ----------
# 旧実装は [m["name"] for m in data.get("models", [])] という素のリスト内包表記で、
# 応答中の1件でも壊れていると(dict以外の要素・"name"キー欠落・"name"が非文字列・
# "name"が空文字列)その場で例外を出し、外側の except Exception: return [] が正常な
# 要素も含めて導入済みモデル一覧を丸ごと空にしていた。これはiter103/111/112/113/139と
# 同じ「1件の壊れた要素が全件を道連れにする」失敗形。installed_models()はpull()で
# 自己修復しない2箇所の判定に直結しており、空リストへの縮退は静かに精度を落とす:
# _arbitrate の is_installed(ARBITER_MODEL, installed_models()) は最上位知性モデル
# gpt-oss:120b を裁定チェーンへ加えるか否かを決めており、空リストだと裁定が弱い
# フォールバックモデルへ静かに格下げされる(gotcha #7: SC投票のtie-break劣化)。
# solve_verifiable の is_installed(SC_CHEAP_MODEL, installed_models()) も同様に
# 安価な追加投票の有無を決めており、空リストだと投票パネルが静かに薄くなる。
# installed_models() はこれまで resolve_models() のテスト(上のセクション、iter76)で
# 常に丸ごとモックされるだけで、この関数自体を直接検証するテストが無かった。ここでは
# urllib.request.urlopen のみをモックし、実ネットワーク/Ollama呼び出しは一切発生させない。


class _ImFakeResponse:
    """installed_models() 用の `with urlopen(...) as r:` 最小モック(bytesを返す)。"""

    def __init__(self, body_bytes):
        self._body = body_bytes

    def __enter__(self):
        return self

    def __exit__(self, *exc_info):
        return False

    def read(self):
        return self._body


def _im_urlopen_body(raw_text):
    """rawテキスト(JSON文字列、または壊れた非JSON文字列)をそのままbodyとして
    返すurlopen代役を作る。"""
    body_bytes = raw_text.encode("utf-8")

    def _fake(url, timeout=None):
        return _ImFakeResponse(body_bytes)
    return _fake


def _im_urlopen_payload(payload):
    """Pythonオブジェクトをjson.dumpsしてbodyとして返すurlopen代役を作る。"""
    return _im_urlopen_body(json.dumps(payload))


def _im_urlopen_raises(exc):
    def _fake(url, timeout=None):
        raise exc
    return _fake


_orig_urlopen_im = urllib.request.urlopen

try:
    # --- (1) 正常な要素の間に壊れた要素が混在 -> 壊れた要素だけ読み飛ばし、
    #     正常な要素は元の順序で回収する(例外は伝播せず一部も失わない) ---
    urllib.request.urlopen = _im_urlopen_payload({
        "models": [
            {"name": "good-model-1:1b"},
            "not-a-dict-entry",
            {"size": 123},
            {"name": ["not", "a", "string"]},
            {"name": ""},
            {"name": "good-model-2:7b"},
        ]
    })
    check("installed_models: 混在応答は壊れた要素だけ読み飛ばし正常分を順序維持で返す",
          f.installed_models() == ["good-model-1:1b", "good-model-2:7b"])

    # --- (2) "models" が非list -> [] (isinstance判定でリスト内包表記の例外を防ぐ) ---
    for _bad_models_val, _label in [
        ({"a": 1}, "dict"), ("oops", "str"), (123, "int"),
        (True, "bool"), (None, "None"),
    ]:
        urllib.request.urlopen = _im_urlopen_payload({"models": _bad_models_val})
        check(f"installed_models: models が非list({_label})なら例外を出さず[]",
              f.installed_models() == [])

    # --- (3) "models" キー自体が無い -> [] ---
    urllib.request.urlopen = _im_urlopen_payload({"other_key": "x"})
    check("installed_models: modelsキー自体が無いなら[]", f.installed_models() == [])

    # --- (4) 応答トップレベルが非dict -> [] ---
    for _bad_top, _label in [([1, 2, 3], "list"), ("plain string", "str"), (42, "number")]:
        urllib.request.urlopen = _im_urlopen_payload(_bad_top)
        check(f"installed_models: 応答トップレベルが非dict({_label})なら[]",
              f.installed_models() == [])

    # --- (5) 不正/非JSON本文 -> [] ---
    urllib.request.urlopen = _im_urlopen_body("{not valid json!!")
    check("installed_models: 不正な(非JSON)本文でも例外を出さず[]",
          f.installed_models() == [])

    # --- (6) urlopen自体が例外を投げる(ネットワーク断等) -> [] ---
    urllib.request.urlopen = _im_urlopen_raises(RuntimeError("ネットワーク断(模擬)"))
    check("installed_models: urlopenが例外を投げても[]", f.installed_models() == [])

    # --- (7) 正常系のリグレッション: 壊れた要素が無いN件の応答は旧実装のリスト
    #     内包表記と同じ名前・同じ順序を返す(byte-for-byte) ---
    _im_happy_models = [{"name": "m1:1b"}, {"name": "m2:7b"}, {"name": "m3:30b"}]
    urllib.request.urlopen = _im_urlopen_payload({"models": _im_happy_models})
    _im_expected_happy = [m["name"] for m in _im_happy_models]  # 旧実装と同じ式
    check("installed_models: 正常系はN件を旧実装と同一・同順序で返す(回帰)",
          f.installed_models() == _im_expected_happy
          and f.installed_models() == ["m1:1b", "m2:7b", "m3:30b"])
finally:
    urllib.request.urlopen = _orig_urlopen_im

check("installed_models: テスト後にurllib.request.urlopenが元の状態へ復元されている",
      urllib.request.urlopen == _orig_urlopen_im)

# ---------- _save_answer_to_file: --out の親ディレクトリが存在しない場合でも
# 計算済み回答を失わない (2026-07-23修正) ----------
# build_pptx は既に自前で out_path.parent.mkdir を呼んでいる(L3729、python-pptx
# 不在時のフォールバック分岐にも別途mkdirあり)が、他のsuffix分岐(.md/.txt/.py等)は
# out.write_text / 各lib.save を親ディレクトリの存在チェック無しに直接呼んでおり、
# --out に存在しないサブディレクトリ(例: reports/answer.md)を指定すると
# FileNotFoundError が ask_fugu 側(L3185-3187)無捕捉のまま伝播し、MoA/SC投票まで
# 完了した高コストな計算済み回答がトレースバックと共に失われていた。さらに
# office系フォールバック(.docx/.xlsx失敗時の.md/.csv代替書き込み)も同じ存在しない
# ディレクトリに書こうとして二重に失敗し、出力が完全に消える。これは
# iteration 41-47・68(_save_as_excelのIllegalCharacterError・_save_as_docxの
# ValueError・_save_as_pdfのFPDFUnicodeEncodingException等、保存段で計算済み
# 回答を失わないための一連の修正)と同じバグクラス。_save_answer_to_file の
# out = Path(path) 直後に out.parent.mkdir(parents=True, exist_ok=True) を1箇所
# 追加し、全suffix分岐とその場フォールバックをまとめて保護する(個別saverへの
# mkdir散乱はしない)。ローカル一時ディレクトリのみで検証し、
# Ollama/ネットワーク/bench呼び出しは一切発生しない。
import tempfile as _pd_tempfile
import pathlib as _pd_pathlib

with _pd_tempfile.TemporaryDirectory() as _pd_dir:
    _pd_root = _pd_pathlib.Path(_pd_dir)

    # (1) .md を存在しない1階層サブディレクトリへ
    #     -> ディレクトリが作られ、例外なくanswerが書き込まれる
    _pd_md = _pd_root / "missing_subdir" / "answer.md"
    _pd_md_exc = None
    try:
        f._save_answer_to_file("q_pd_md", "回答pd_md", 1.0, str(_pd_md))
    except Exception as _exc:
        _pd_md_exc = _exc
    check("_save_answer_to_file: 存在しないサブディレクトリへの.md保存は例外を送出しない",
          _pd_md_exc is None)
    check("_save_answer_to_file: 存在しないサブディレクトリが作成される(.md)",
          _pd_md.parent.is_dir())
    if _pd_md_exc is None:
        check("_save_answer_to_file: .mdファイルが作成され回答を含む",
              _pd_md.is_file() and "回答pd_md" in _pd_md.read_text(encoding="utf-8"))

    # (2) コード拡張子(.py)を存在しないサブディレクトリへ
    #     -> ディレクトリが作られ、抽出されたコードが書き込まれる
    _pd_py = _pd_root / "missing_subdir_code" / "answer.py"
    _pd_py_answer = "説明文\n```python\nprint('hello')\n```\n"
    _pd_py_exc = None
    try:
        f._save_answer_to_file("q_pd_py", _pd_py_answer, 0.5, str(_pd_py))
    except Exception as _exc:
        _pd_py_exc = _exc
    check("_save_answer_to_file: 存在しないサブディレクトリへの.py保存は例外を送出しない",
          _pd_py_exc is None)
    if _pd_py_exc is None:
        _pd_py_content = _pd_py.read_text(encoding="utf-8")
        check("_save_answer_to_file: .pyファイルが作成され抽出コードを含む",
              "print('hello')" in _pd_py_content)

    # (3) .txt / .html を2階層以上ネストした存在しないサブディレクトリへ
    _pd_txt = _pd_root / "a" / "b" / "c" / "answer.txt"
    _pd_txt_exc = None
    try:
        f._save_answer_to_file("q_pd_txt", "回答pd_txt", 0.3, str(_pd_txt))
    except Exception as _exc:
        _pd_txt_exc = _exc
    check("_save_answer_to_file: 2階層以上ネストした存在しないサブディレクトリへの.txt保存は例外なし",
          _pd_txt_exc is None)
    if _pd_txt_exc is None:
        check("_save_answer_to_file: ネスト.txtファイルが作成され回答を含む",
              _pd_txt.is_file() and "回答pd_txt" in _pd_txt.read_text(encoding="utf-8"))

    _pd_html = _pd_root / "x" / "y" / "z" / "answer.html"
    _pd_html_exc = None
    try:
        f._save_answer_to_file("q_pd_html", "回答pd_html", 0.4, str(_pd_html))
    except Exception as _exc:
        _pd_html_exc = _exc
    check("_save_answer_to_file: 2階層以上ネストした存在しないサブディレクトリへの.html保存は例外なし",
          _pd_html_exc is None)
    if _pd_html_exc is None:
        check("_save_answer_to_file: ネスト.htmlファイルが作成され回答を含む",
              _pd_html.is_file() and "回答pd_html" in _pd_html.read_text(encoding="utf-8"))

    # (4) 回帰: 既に存在するディレクトリ(tempdirルート直下)への保存は
    #     byte-for-byte従来通り(.mdの既存内容への追記マージが不変であることを確認)
    _pd_regress = _pd_root / "regress.md"
    _pd_existing = "# 既存メモ\n\n既存内容 täüst 日本語。\n\n"
    _pd_regress.write_text(_pd_existing, encoding="utf-8")
    f._save_answer_to_file("q_pd_regress", "回答pd_regress", 2.0, str(_pd_regress))
    _pd_regress_content = _pd_regress.read_text(encoding="utf-8")
    _pd_ts_m = f.re.search(r"## Q \(([^)]+)\)", _pd_regress_content)
    check("_save_answer_to_file: 既存ディレクトリへの追記後にtsが検出できる", _pd_ts_m is not None)
    if _pd_ts_m:
        _pd_ts = _pd_ts_m.group(1)
        _pd_expected_block = (f"## Q ({_pd_ts})\n\nq_pd_regress\n\n"
                              f"## A\n\n回答pd_regress\n\n*所要: 2.0s*\n\n---\n\n")
        check("_save_answer_to_file: 既存ディレクトリへの.md追記はbyte-for-byte従来通り",
              _pd_regress_content == _pd_existing + _pd_expected_block)

# ==================================================
# ---------- MCQ 選択肢文字(A-E)の投票クラス整合性 (2026-07-24, gotcha #7) ----------
# ==================================================
# solve_verifiable の mcq 分岐は vote_answers -> answers_equivalent で選択肢文字(A-E)を
# 集計するが、既存の SC/投票の直接テスト（iteration 32/54/55/79/105）は全て math の
# 数値/分数の組み合わせのみを使っており、この「カテゴリカルな文字比較」経路には直接の
# カバレッジが無かった。2つの異なる文字については answers_equivalent の高速パスが
# 両方失敗する（na.lower() != nb.lower(); Fraction('A') は例外）ため、異なる文字の票が
# 別クラスのまま保たれるかどうかは実質的に math_verify（数式の同値判定エンジンであり、
# 選択肢文字のようなカテゴリカルなラベルの比較に使うことは想定されていない）の判断に
# 委ねられてしまっている。これは gotcha #7（自己整合性投票の完全性）に直結する未検証の
# 脆弱な依存であり、iteration 32/54/55/79 が固定したのと同種の特性テストとしてここに
# 直接ロックする。同一文字の大小/全角違いは併合される（normalize_answer の lower()/
# _FW_TRANS 高速パス）ことも合わせて確認し、非対称（違う文字は割れる/同じ文字は
# 統合される）であることを保証する。

# (a) 異なる選択肢文字(A/B/C)は絶対に併合されない。classesは件数降順で返る。
_top_mcqv1, _cnt_mcqv1, _cls_mcqv1 = f.vote_answers(["A", "B", "A", "C", "A"])
check("mcq-vote: 異なる選択肢文字は併合されず3クラスのまま",
      len(_cls_mcqv1) == 3)
check("mcq-vote: 最多票はA(3票)、票数はA(3)/B(1)/C(1)で件数降順",
      _top_mcqv1 == "A" and _cnt_mcqv1 == 3
      and _cls_mcqv1 == [["A", 3], ["B", 1], ["C", 1]])

# (b) tie安定性（iteration 54の数値版と同型）: 同数タイは先出現の文字が勝つ。
_top_mcqv2, _cnt_mcqv2, _cls_mcqv2 = f.vote_answers(["A", "B", "B", "A"])
check("mcq-vote: 同数タイは先出現(A)の文字が勝つ",
      _top_mcqv2 == "A" and _cnt_mcqv2 == 2 and _cls_mcqv2 == [["A", 2], ["B", 2]])
_top_mcqv3, _cnt_mcqv3, _cls_mcqv3 = f.vote_answers(["B", "A", "A", "B"])
check("mcq-vote: 入力順を逆にすると先出現(B)の文字が勝つ(値の大小ではない)",
      _top_mcqv3 == "B" and _cnt_mcqv3 == 2 and _cls_mcqv3 == [["B", 2], ["A", 2]])

# (c)+(d) math_verify を「呼ばれたら必ず例外」なスタブに差し替え(iteration 13/32と同じ
#     swap-and-restoreパターン。_fake_math_verify は上で定義済みのものを再利用)、
#     (c) 同一文字の大小/全角違いは高速パス(lower()/_FW_TRANS)だけで併合されること、
#     (d) 異なる文字は math_verify が例外を投げても(=利用不能でも)絶対に併合されない
#     ことの両方を、同一のスタブ下で確認する。
_orig_mv_mod_mcqv = sys.modules.get("math_verify")
sys.modules["math_verify"] = _fake_math_verify
try:
    check("mcq: answers_equivalent 大小文字違い(A/a)は高速パスでTrue(math_verify不要)",
          f.answers_equivalent("A", "a") is True)
    check("mcq: answers_equivalent 全角/半角文字(Ａ/A)は高速パスでTrue(math_verify不要)",
          f.answers_equivalent("Ａ", "A") is True)
    check("mcq-vote: 大小文字違い(A/a)は高速パスで単一クラス(cnt=2)へ併合",
          f.vote_answers(["A", "a"]) == ("A", 2, [["A", 2]]))
    check("mcq-vote: 全角文字(Ａ)は高速パスで単一クラス(cnt=2)へ併合",
          f.vote_answers(["A", "Ａ"]) == ("A", 2, [["A", 2]]))

    check("mcq: answers_equivalent 異なる文字(A/B)はmath_verify例外時もFalse",
          f.answers_equivalent("A", "B") is False)
    _top_mcqv4, _cnt_mcqv4, _cls_mcqv4 = f.vote_answers(["A", "B"])
    check("mcq-vote: 異なる文字(A/B)はmath_verifyが例外/利用不能でも2クラスのまま"
          "(併合されない整合性はmath_verifyの判断に依存しない)",
          len(_cls_mcqv4) == 2 and _cnt_mcqv4 == 1
          and _cls_mcqv4[0][0] != _cls_mcqv4[1][0])
finally:
    if _orig_mv_mod_mcqv is not None:
        sys.modules["math_verify"] = _orig_mv_mod_mcqv
    else:
        del sys.modules["math_verify"]
# 補足: この開発環境には math_verify が実際にはインストールされておらず(import時点で
# ModuleNotFoundError)、上のスタブ無しでも同じ結果になる。つまり「math_verifyが
# 利用不能」なケースは実質的に常時カバーされている。一方「math_verifyが実在し、かつ
# A/Bのような異なる選択肢文字を誤って同値だと判定してしまう」ケースは、ライブラリ本体を
# 実際にインストールしない限りここでは再現・検証できない。もし将来 math_verify が
# 導入され、かつ実際に異なる選択肢文字を同値だと判定する事例が見つかった場合は、
# 静かに fugu_local.py 側を直さず、票の併合バグとして別イテレーションで報告すること
# (surface-don't-fix、iteration 66/71/48 の方針)。

# (e) solve_verifiable 全体(mcq)を f._sc_sample のみモックしてE2Eで確認する。
#     f.ask には一切触れず、B,B,B,A,C のプルラリティ(B)で確定し、votesが
#     A/B/C を別キーとして正しく持ち、水増し/併合が無いことを検証する。
#     SC_POT=True のまま渡し、mcqではPoTサンプルが一切要求されないこと
#     （add_batchの `if SC_POT and task_type == 'math'` ガード）も併せて回帰確認する。
_orig_mcqv_sc_sample = f._sc_sample
_orig_mcqv_pot = f.SC_POT
_orig_mcqv_cheap_votes = f.SC_CHEAP_VOTES
_orig_mcqv_props = f.PROPOSERS
_orig_mcqv_reasoning = f.REASONING_MODELS
_orig_mcqv_initial = f.SC_INITIAL
_mcqv_calls = []
_mcqv_map = {("m1", False): ["B", "B", "B", "A", "C"]}
try:
    f.PROPOSERS = ["m1"]
    f.REASONING_MODELS = ["m1"]
    f.SC_CHEAP_VOTES = 0
    f.SC_POT = True
    f.SC_INITIAL = 5
    f._sc_sample = _make_fake_sc_sample(_mcqv_map, _mcqv_calls)
    with contextlib.redirect_stdout(io.StringIO()):
        _res_mcqv = f.solve_verifiable("次のうち正しいものはどれか", "mcq")
finally:
    f._sc_sample = _orig_mcqv_sc_sample
    f.PROPOSERS = _orig_mcqv_props
    f.REASONING_MODELS = _orig_mcqv_reasoning
    f.SC_CHEAP_VOTES = _orig_mcqv_cheap_votes
    f.SC_POT = _orig_mcqv_pot
    f.SC_INITIAL = _orig_mcqv_initial

check("mcq-e2e: B,B,B,A,Cのプルラリティ(B, 3/5)で確定",
      _res_mcqv is not None and _res_mcqv["answer"] == "B")
check("mcq-e2e: votesがtruthful(A/B/Cが別キーの実票数、水増し/併合なし)",
      _res_mcqv is not None and _res_mcqv["votes"] == {"B": 3, "A": 1, "C": 1})
check("mcq-e2e: n_samplesは5(全てCoT、PoTは含まれない)",
      _res_mcqv is not None and _res_mcqv["n_samples"] == 5)
check("mcq-e2e: SC_POT=TrueでもmcqではPoTサンプルは要求されない(全呼び出しがpot=False)",
      len(_mcqv_calls) == 5 and not any(pot for _m, pot in _mcqv_calls))
check("mcq-e2e: テスト後にf._sc_sample/SC_POT/SC_INITIAL等のグローバルが復元されている",
      f._sc_sample == _orig_mcqv_sc_sample and f.SC_POT == _orig_mcqv_pot
      and f.SC_CHEAP_VOTES == _orig_mcqv_cheap_votes and f.PROPOSERS == _orig_mcqv_props
      and f.REASONING_MODELS == _orig_mcqv_reasoning and f.SC_INITIAL == _orig_mcqv_initial)

# ==================================================
# ---------- answers_equivalent: MCQ選択肢文字(A-E)ガード (2026-07-25, iteration 129) ----------
# ==================================================
# 上のMCQ投票クラス整合性テスト(iteration 107)は「異なる選択肢文字がmath_verifyによって
# 誤って同値と判定される」リスクを指摘したが、その場ではfugu_local.py側を修正せず据え置いた
# （math_verifyがこの開発環境に無く、実際にそう誤判定する事例を再現できなかったため）。
# fugu_local.py の answers_equivalent には現在、na.lower()==nb.lower() の直後に「両方とも
# 単一の選択肢文字A-Eの形をしている場合はFraction/math_verifyへ一切渡さず、
# case-insensitive文字列比較で確定的に返す」ガードが追加されている(gotcha #6/#7参照)。
# ここでは math_verify.verify が常にTrueを返す（=iteration 107が懸念した「誤って同値と
# 判定してしまう」最悪ケースを意図的に再現する）記録スタブに差し替え、それでも異なる
# 選択肢文字が併合されないこと、かつ選択肢文字ペアではparse/verifyが一切呼ばれないこと
# （ガードがFraction/math_verifyより手前で確定的にreturnしている直接証拠）を検証する。
# 対比として、選択肢文字の形をしていない真にmath_verify依存のペア(iteration 122で選定
# された \sqrt{2} vs 1.41421356、Fraction高速パスを迂回する)は引き続きスタブへ到達し、
# ガードが過剰発火して本来のmath比較を握り潰していないことも同じスタブの下で確認する。


def _t_letter_guard_always_true(calls):
    check("letter-guard: 異なる選択肢文字(A/B)はverify常にTrueのスタブでも非同値",
          f.answers_equivalent("A", "B") is False)
    check("letter-guard: 異なる選択肢文字(E/D)はverify常にTrueのスタブでも非同値"
          "(sympyはEをEuler数として解釈しうる境界ケース)",
          f.answers_equivalent("E", "D") is False)
    check("letter-guard: 異なる選択肢文字(C/A)はverify常にTrueのスタブでも非同値",
          f.answers_equivalent("C", "A") is False)
    check("letter-guard: 選択肢文字ペア(A/B, E/D, C/A)ではmath_verify.parseが一切呼ばれない"
          "(ガードがFraction/math_verifyより手前で確定的にreturnしている証拠)",
          len(calls["parse_args"]) == 0)
    check("letter-guard: 選択肢文字ペアではmath_verify.verifyが一切呼ばれない",
          len(calls["verify_args"]) == 0)

    # 同一文字(大小違い)は従来通りna.lower()の高速パスで解決し、こちらもmath_verifyの
    # スタブに一切到達しない(直前の異なる文字ペアの検証を経てもなお呼び出し件数は0のまま)。
    check("letter-guard: 同一選択肢文字(A/A)はna.lower()高速パスでTrue",
          f.answers_equivalent("A", "A") is True)
    check("letter-guard: 大小文字違い(a/A)はna.lower()高速パスでTrue",
          f.answers_equivalent("a", "A") is True)
    check("letter-guard: 同一文字ペア判定後もmath_verifyスタブは一度も呼ばれていない",
          len(calls["parse_args"]) == 0 and len(calls["verify_args"]) == 0)

    # 対比: 選択肢文字の形をしていない真にmath_verify依存のペアは、同じ「verifyが常に
    # Trueを返す」スタブの下で引き続きparse/verifyへ到達しTrueを受け取る(ガードが過剰
    # 発火してmath比較経路そのものを塞いでいないことの確認。iteration 122のペアを再利用)。
    _result_math = f.answers_equivalent(_MV_A, _MV_B)
    check("letter-guard: 選択肢文字でない真のmath_verify依存ペアはガードを迂回されず"
          "引き続きスタブへ到達しTrueが伝播する(過剰発火していない)",
          _result_math is True and len(calls["parse_args"]) >= 2 and len(calls["verify_args"]) == 1)


_run_with_math_verify_stub(True, "none", _t_letter_guard_always_true)

# ==================================================
# ---------- plan_pptx_images: index範囲/重複/予算上限のパース契約 (2026-07-24) ----------
# ==================================================
# 背景: plan_pptx_images() (fugu_local.py L4095) は Conductor の JSON 画像プランを
# {index: prompt} の dict へパースする。この契約 ―
#   index 0 = タイトルヒーロー画像、1..len(slides) は build_pptx 側で
#   slides[idx-1] に対応、非整数/None index・辞書でないエントリ・'index'欠落は
#   try/except-continue で無視、空/空白のみの prompt はスキップ、重複 index は
#   先勝ち(idx not in out)、範囲は 0 <= idx <= len(slides) の閉区間、最大
#   PPTX_MAX_IMAGES 件で打ち切り ― は build_pptx の slides[idx-1] アクセスと
# iter77 のタイトルヒーロー保証(plan.setdefault(0, None)/L4164)が前提とする
# 不変条件そのものである。従来 test_fugu_offline.py 内では plan_pptx_images は
# build_pptx テスト(iter68/77、L7505/7533等)の中で常に丸ごとモックされるだけで、
# この関数自体への直接テストは皆無だった。将来の変更で範囲外 index や
# PPTX_MAX_IMAGES 超過件数を返すよう壊れた場合、タイトルヒーロー保証が静かに
# 崩れるか、build_pptx が IndexError で落ちる。
# ここでは plan_pptx_images の唯一の外部呼び出しである f.ask のみをモックし、
# extract_json は実物をそのまま通す(純粋ロジックのため安全)。モック漏れの
# 即時検知のため urllib.request.urlopen と f.subprocess.run も「呼ばれたら
# 即AssertionError」の番人に差し替える(iteration 38/39/76/104と同じ流儀)。
# Ollama/ネットワーク/bench呼び出しは一切発生しない。

check("plan_pptx_images: このテストはPPTX_MAX_IMAGES==4を前提にしている",
      f.PPTX_MAX_IMAGES == 4)

_orig_pi_ask = f.ask
_orig_pi_urlopen = urllib.request.urlopen
_orig_pi_subprocess_run = f.subprocess.run


def _pi_no_network_urlopen(*a, **k):
    raise AssertionError("plan_pptx_images: モック漏れで実urlopen(ネットワーク)が呼ばれた")


def _pi_no_subprocess_run(*a, **k):
    raise AssertionError("plan_pptx_images: モック漏れで実subprocess.runが呼ばれた")


# len(slides) == PPTX_MAX_IMAGES + 2。境界値(index==len(slides))と
# 予算超過(有効エントリ > PPTX_MAX_IMAGES)の両方を無理なく作れる件数にする。
_PI_N_SLIDES = f.PPTX_MAX_IMAGES + 2
_pi_slides = [{"title": f"Slide {i + 1}",
               "bullets": [f"b{i + 1}-1", f"b{i + 1}-2", f"b{i + 1}-3", f"b{i + 1}-4"]}
              for i in range(_PI_N_SLIDES)]

_pi_calls = []


def _make_pi_ask(canned):
    # 呼び出し形状は実装(fugu_local.py L4109-4112)通り: 位置引数
    # model/messages/temperature + キーワード think=/fmt=/num_predict=/label=。
    # (*args, **kwargs) で受け止め、その形をそのまま記録して後で検証する。
    def _fake_pi_ask(*args, **kwargs):
        _pi_calls.append((args, kwargs))
        return canned
    return _fake_pi_ask


try:
    urllib.request.urlopen = _pi_no_network_urlopen
    f.subprocess.run = _pi_no_subprocess_run

    # (0) 呼び出し形状の確認: askは位置引数3個(model,messages,temperature)+
    #     think=/fmt=/num_predict=/label=キーワードで呼ばれる(実装の唯一の外部呼び出し)。
    _pi_calls.clear()
    f.ask = _make_pi_ask('{"images":[{"index":0,"prompt":"Hero"}]}')
    f.plan_pptx_images("Title", _pi_slides)
    check("plan_pptx_images(0): askは位置引数3個(model,messages,temperature)で呼ばれる",
          len(_pi_calls) == 1 and len(_pi_calls[0][0]) == 3)
    check("plan_pptx_images(0): askはthink=/fmt=/num_predict=/label=キーワードで呼ばれる"
          "(label='pptx-img-plan')",
          _pi_calls[0][1].keys() == {"think", "fmt", "num_predict", "label"}
          and _pi_calls[0][1]["label"] == "pptx-img-plan")

    # (1) 正常系: index 0 + 範囲内の複数indexを含むプラン -> そのまま{index:prompt}で返る。
    f.ask = _make_pi_ask(json.dumps({"images": [
        {"index": 0, "prompt": "Hero image"},
        {"index": 2, "prompt": "Slide2 image"},
        {"index": 4, "prompt": "Slide4 image"},
    ]}))
    _pi_r1 = f.plan_pptx_images("Title", _pi_slides)
    check("plan_pptx_images(1): 正常系は指定した{index:prompt}をそのまま返す",
          _pi_r1 == {0: "Hero image", 2: "Slide2 image", 4: "Slide4 image"})
    check("plan_pptx_images(1): index 0(タイトルヒーロー)が含まれる", 0 in _pi_r1)

    # (2) 不正エントリ(非整数index/None index/辞書でないエントリ/index欠落)は
    #     例外を送出せずスキップされる(int(it.get("index"))のtry/except-continue)。
    f.ask = _make_pi_ask(json.dumps({"images": [
        {"index": "not-an-int", "prompt": "bad1"},
        {"index": None, "prompt": "bad2"},
        "just-a-string-not-a-dict",
        {"prompt": "missing index key"},
        {"index": 1, "prompt": "good"},
    ]}))
    _pi_r2, _pi_r2_exc = None, None
    try:
        _pi_r2 = f.plan_pptx_images("Title", _pi_slides)
    except Exception as _exc:
        _pi_r2_exc = _exc
    check("plan_pptx_images(2): 不正エントリ混在でも例外を送出しない", _pi_r2_exc is None)
    check("plan_pptx_images(2): 不正エントリは無視され有効な1件のみ残る",
          _pi_r2 == {1: "good"})

    # (3) 空/空白のみのpromptはスキップされる(if p and ...)。
    f.ask = _make_pi_ask(json.dumps({"images": [
        {"index": 1, "prompt": ""},
        {"index": 2, "prompt": "   "},
        {"index": 3, "prompt": "kept"},
    ]}))
    _pi_r3 = f.plan_pptx_images("Title", _pi_slides)
    check("plan_pptx_images(3): 空/空白のみのpromptは除外される", _pi_r3 == {3: "kept"})

    # (4) 重複indexは先勝ち(idx not in out)。
    f.ask = _make_pi_ask(json.dumps({"images": [
        {"index": 1, "prompt": "first"},
        {"index": 1, "prompt": "second"},
    ]}))
    _pi_r4 = f.plan_pptx_images("Title", _pi_slides)
    check("plan_pptx_images(4): 重複indexは最初のpromptが勝つ", _pi_r4 == {1: "first"})

    # (5) 範囲外indexは除外され、index==len(slides)の境界値は保持される(閉区間)。
    f.ask = _make_pi_ask(json.dumps({"images": [
        {"index": -1, "prompt": "negative"},
        {"index": _PI_N_SLIDES + 1, "prompt": "too-high"},
        {"index": _PI_N_SLIDES, "prompt": "boundary"},
    ]}))
    _pi_r5 = f.plan_pptx_images("Title", _pi_slides)
    check("plan_pptx_images(5): 負のindexは除外される", -1 not in _pi_r5)
    check("plan_pptx_images(5): len(slides)を超えるindexは除外される",
          (_PI_N_SLIDES + 1) not in _pi_r5)
    check("plan_pptx_images(5): index==len(slides)の境界値は保持される"
          "(閉区間 0<=idx<=len(slides)、切り捨てではない)",
          _pi_r5 == {_PI_N_SLIDES: "boundary"})

    # (6) 有効エントリがPPTX_MAX_IMAGESを超えても、その上限で打ち切られる。
    f.ask = _make_pi_ask(json.dumps({"images": [
        {"index": i, "prompt": f"p{i}"} for i in range(_PI_N_SLIDES + 1)
    ]}))
    _pi_r6 = f.plan_pptx_images("Title", _pi_slides)
    check("plan_pptx_images(6): 有効件数がPPTX_MAX_IMAGESを超えても上限で打ち切られる",
          len(_pi_r6) == f.PPTX_MAX_IMAGES)
    check("plan_pptx_images(6): 打ち切りは出現順(先頭からPPTX_MAX_IMAGES件が残る)",
          _pi_r6 == {i: f"p{i}" for i in range(f.PPTX_MAX_IMAGES)})

    # (7) askが解析不能な文字列/'__ERROR__'センチネルを返す
    #     -> extract_jsonがNoneになりj={}に落ち、{}を返す(例外は送出しない)。
    f.ask = _make_pi_ask("__ERROR__: HTTP Error 500 (模擬)")
    _pi_r7a = f.plan_pptx_images("Title", _pi_slides)
    check("plan_pptx_images(7a): '__ERROR__'センチネルは例外なく{}を返す", _pi_r7a == {})
    f.ask = _make_pi_ask("this is not json at all, just prose")
    _pi_r7b = f.plan_pptx_images("Title", _pi_slides)
    check("plan_pptx_images(7b): 解析不能な出力も例外なく{}を返す", _pi_r7b == {})

    # (8) 'images'がリストでない値(文字列)でもクラッシュせず{}を返す。
    #     2026-07-24: 従来の j.get("images") or [] は falsy値のみを[]化するトリック
    #     だったため、"images"が真になり得る非リスト値(文字列/int/float/bool等)では
    #     for文自体が例外の温床になり得た。文字列の場合は反復可能なため
    #     (1文字ずつのstrがitに入りit.get(...)がAttributeErrorになるがtry/exceptで
    #     握りつぶされる)従来から{}に縮退していたが、int/float/boolのような
    #     非反復可能な真値ではfor文自体がTypeErrorを送出し未捕捉のまま伝播していた
    #     (iteration110がこのケースをテストコメントとして発見・記録していたが、
    #     当時のタスクはテストのみの変更に限定されており修正は見送られていた)。
    #     本イテレーションで plan_pptx_images 側を
    #     `imgs = j.get("images"); if not isinstance(imgs, list): imgs = []` に
    #     修正し、真偽・型に関わらず非list値を確実に[]へ倒すようにした
    #     (iteration103の_ddg_instant非list RelatedTopics補正と同じ方式)。
    #     この(8)は文字列ケースの回帰確認、直後の(8b)で新たに修正された
    #     int/float/bool(非反復可能な真値)のケースを検証する。
    f.ask = _make_pi_ask(json.dumps({"images": "oops-not-a-list"}))
    _pi_r8, _pi_r8_exc = None, None
    try:
        _pi_r8 = f.plan_pptx_images("Title", _pi_slides)
    except Exception as _exc:
        _pi_r8_exc = _exc
    check("plan_pptx_images(8): 'images'が非リスト(文字列)でも例外を送出しない",
          _pi_r8_exc is None)
    check("plan_pptx_images(8): 'images'が非リスト(文字列)の場合は{}を返す", _pi_r8 == {})

    # (8b) 'images'が非反復可能な真値(int/float/bool)でも例外を送出せず{}を返す。
    #      修正前はここで素の `for it in (j.get("images") or [])` がTypeErrorを
    #      送出しており、build_pptxのXML安全化try/except(iteration68)の外・
    #      ask_fuguの無防備な呼び出し元まで伝播して計算済み回答を丸ごと失っていた。
    for _pi_bad_images in (5, 3.14, True):
        f.ask = _make_pi_ask(json.dumps({"images": _pi_bad_images}))
        _pi_r8b, _pi_r8b_exc = None, None
        try:
            _pi_r8b = f.plan_pptx_images("Title", _pi_slides)
        except Exception as _exc:
            _pi_r8b_exc = _exc
        check(f"plan_pptx_images(8b): 'images'={_pi_bad_images!r}(非反復可能な真値)"
              "でも例外を送出しない", _pi_r8b_exc is None)
        check(f"plan_pptx_images(8b): 'images'={_pi_bad_images!r}の場合は{{}}を返す",
              _pi_r8b == {})

    # (contract) build_pptxが依存する不変条件: 返る全キーが0<=k<=len(slides)を
    #     満たし、件数はPPTX_MAX_IMAGES以下である。(1)(5)(6)の結果で確認する。
    for _pi_label, _pi_r in (("(1)", _pi_r1), ("(5)", _pi_r5), ("(6)", _pi_r6)):
        check(f"plan_pptx_images contract{_pi_label}: 全キーが0<=k<=len(slides)を満たす",
              all(0 <= k <= _PI_N_SLIDES for k in _pi_r.keys()))
        check(f"plan_pptx_images contract{_pi_label}: 件数はPPTX_MAX_IMAGES以下",
              len(_pi_r) <= f.PPTX_MAX_IMAGES)
finally:
    f.ask = _orig_pi_ask
    urllib.request.urlopen = _orig_pi_urlopen
    f.subprocess.run = _orig_pi_subprocess_run

check("plan_pptx_images: テスト後にf.askが元の状態へ復元されている", f.ask == _orig_pi_ask)
check("plan_pptx_images: テスト後にurllib.request.urlopenが元の状態へ復元されている",
      urllib.request.urlopen == _orig_pi_urlopen)
check("plan_pptx_images: テスト後にf.subprocess.runが元の状態へ復元されている",
      f.subprocess.run == _orig_pi_subprocess_run)

# ==================================================
# handle_image_generation の直接テスト (iteration 118)
# ==================================================
# _detect_backend/author_image_prompt/generate_image の3つの外部呼び出しのみを
# monkeypatchし、実urlopen/subprocessには一切触れない。モック漏れを即座に
# 可視化するため、urllib.request.urlopenとf.subprocess.runにも「呼ばれたら
# 即AssertionError」の番人を仕込む(iteration 38/76/103のtripwire流儀を踏襲)。

_orig_hig_urlopen = urllib.request.urlopen
_orig_hig_subprocess_run = f.subprocess.run
_orig_hig_image_backend = f.IMAGE_BACKEND
_orig_hig_detect_backend = f._detect_backend
_orig_hig_author_image_prompt = f.author_image_prompt
_orig_hig_generate_image = f.generate_image


def _hig_no_network_urlopen(*a, **k):
    raise AssertionError("handle_image_generation: モック漏れで実urlopen(ネットワーク)が呼ばれた")


def _hig_no_subprocess_run(*a, **k):
    raise AssertionError("handle_image_generation: モック漏れで実subprocess.runが呼ばれた")


try:
    urllib.request.urlopen = _hig_no_network_urlopen
    f.subprocess.run = _hig_no_subprocess_run

    # --- (1) IMAGE_BACKEND=='off' -> __ERROR__即返し、他の3関数は一切呼ばれない ---
    f.IMAGE_BACKEND = "off"
    _hig_calls = {"detect": 0, "author": 0, "generate": 0}

    def _hig_detect_should_not_be_called():
        _hig_calls["detect"] += 1
        return "a1111"

    def _hig_author_should_not_be_called(user_request, panel=None):
        _hig_calls["author"] += 1
        return ("dummy prompt", "dummy negative")

    def _hig_generate_should_not_be_called(prompt, negative):
        _hig_calls["generate"] += 1
        return "dummy.png"

    f._detect_backend = _hig_detect_should_not_be_called
    f.author_image_prompt = _hig_author_should_not_be_called
    f.generate_image = _hig_generate_should_not_be_called

    _hig_r1 = f.handle_image_generation("猫の絵を描いて", prompt="既存プロンプト")
    check("handle_image_generation(1): IMAGE_BACKEND='off'で__ERROR__を返す",
          isinstance(_hig_r1, str) and _hig_r1.startswith("__ERROR__"))
    check("handle_image_generation(1): _detect_backendは呼ばれない", _hig_calls["detect"] == 0)
    check("handle_image_generation(1): author_image_promptは呼ばれない", _hig_calls["author"] == 0)
    check("handle_image_generation(1): generate_imageは呼ばれない", _hig_calls["generate"] == 0)

    # --- (2) _detect_backend()がNone -> __ERROR__ + A1111_URL/COMFYUI_URLの案内文、
    #     generate_imageは呼ばれない ---
    f.IMAGE_BACKEND = "auto"
    _hig_calls2 = {"author": 0, "generate": 0}
    f._detect_backend = lambda: None
    f.author_image_prompt = lambda user_request, panel=None: (
        _hig_calls2.__setitem__("author", _hig_calls2["author"] + 1) or ("p", "n"))
    f.generate_image = lambda prompt, negative: (
        _hig_calls2.__setitem__("generate", _hig_calls2["generate"] + 1) or "out.png")

    _hig_r2 = f.handle_image_generation("猫の絵を描いて")
    check("handle_image_generation(2): backend未検出で__ERROR__を返す",
          isinstance(_hig_r2, str) and _hig_r2.startswith("__ERROR__"))
    check("handle_image_generation(2): A1111_URLの案内文を含む", f.A1111_URL in _hig_r2)
    check("handle_image_generation(2): COMFYUI_URLの案内文を含む", f.COMFYUI_URL in _hig_r2)
    check("handle_image_generation(2): generate_imageは呼ばれない", _hig_calls2["generate"] == 0)
    # backend検出失敗の時点でauthor_image_promptに到達する前にreturnするはず
    check("handle_image_generation(2): author_image_promptも呼ばれない", _hig_calls2["author"] == 0)

    # --- (3) prompt=None -> author_image_prompt(user_request, panel=panel)が呼ばれ、
    #     その戻り値(prompt, negative)がgenerate_imageへ渡る ---
    f.IMAGE_BACKEND = "auto"
    f._detect_backend = lambda: "a1111"
    _hig_author_seen = {}
    _hig_generate_seen = {}
    _hig_panel_sentinel = object()

    def _hig_author3(user_request, panel=None):
        _hig_author_seen["user_request"] = user_request
        _hig_author_seen["panel"] = panel
        return ("起草されたプロンプト", "起草されたネガティブ")

    def _hig_generate3(prompt, negative):
        _hig_generate_seen["prompt"] = prompt
        _hig_generate_seen["negative"] = negative
        return "saved/path.png"

    f.author_image_prompt = _hig_author3
    f.generate_image = _hig_generate3

    _hig_r3 = f.handle_image_generation("犬の絵を描いて", panel=_hig_panel_sentinel, prompt=None)
    check("handle_image_generation(3): prompt=Noneでauthor_image_promptが呼ばれる",
          _hig_author_seen.get("user_request") == "犬の絵を描いて")
    check("handle_image_generation(3): author_image_promptにpanelがそのまま渡る",
          _hig_author_seen.get("panel") is _hig_panel_sentinel)
    check("handle_image_generation(3): author_image_promptの戻り値のpromptがgenerate_imageへ渡る",
          _hig_generate_seen.get("prompt") == "起草されたプロンプト")
    check("handle_image_generation(3): author_image_promptの戻り値のnegativeがgenerate_imageへ渡る",
          _hig_generate_seen.get("negative") == "起草されたネガティブ")
    check("handle_image_generation(3): 成功時に保存先パスを含む結果を返す",
          isinstance(_hig_r3, str) and "saved/path.png" in _hig_r3)

    # --- (4) promptを明示指定 -> author_image_promptは一切呼ばれない ---
    f.IMAGE_BACKEND = "auto"
    f._detect_backend = lambda: "a1111"
    _hig_calls4 = {"author": 0}
    f.author_image_prompt = lambda user_request, panel=None: (
        _hig_calls4.__setitem__("author", _hig_calls4["author"] + 1) or ("should-not-be-used", ""))
    _hig_generate_seen4 = {}

    def _hig_generate4(prompt, negative):
        _hig_generate_seen4["prompt"] = prompt
        _hig_generate_seen4["negative"] = negative
        return "explicit.png"

    f.generate_image = _hig_generate4

    _hig_r4 = f.handle_image_generation("猫の絵を描いて", prompt="明示プロンプト", negative="明示ネガティブ")
    check("handle_image_generation(4): prompt明示指定時はauthor_image_promptが呼ばれない",
          _hig_calls4["author"] == 0)
    check("handle_image_generation(4): 明示指定したpromptがそのままgenerate_imageへ渡る",
          _hig_generate_seen4.get("prompt") == "明示プロンプト")
    check("handle_image_generation(4): 明示指定したnegativeがそのままgenerate_imageへ渡る",
          _hig_generate_seen4.get("negative") == "明示ネガティブ")

    # --- (5) generate_imageがfalsy(None/'')を返す -> __ERROR__を返す(例外を送出しない) ---
    f.IMAGE_BACKEND = "auto"
    f._detect_backend = lambda: "a1111"
    f.author_image_prompt = lambda user_request, panel=None: ("p", "")
    for _hig_falsy in (None, ""):
        f.generate_image = lambda prompt, negative, _v=_hig_falsy: _v
        _hig_r5, _hig_r5_exc = None, None
        try:
            _hig_r5 = f.handle_image_generation("何か描いて", prompt="p")
        except Exception as _exc:
            _hig_r5_exc = _exc
        check(f"handle_image_generation(5): generate_imageが{_hig_falsy!r}を返しても例外を送出しない",
              _hig_r5_exc is None)
        check(f"handle_image_generation(5): generate_imageが{_hig_falsy!r}を返すと__ERROR__を返す",
              isinstance(_hig_r5, str) and _hig_r5.startswith("__ERROR__"))

    # --- (6) 成功時のメッセージ構築: negativeがtruthyなら'- negative:'行を含み、
    #     falsy(空文字/None)なら含まない ---
    f.IMAGE_BACKEND = "auto"
    f._detect_backend = lambda: "a1111"
    f.author_image_prompt = lambda user_request, panel=None: ("使われないはず", "使われないはず")
    f.generate_image = lambda prompt, negative: "out/dir/result.png"

    _hig_r6a = f.handle_image_generation("風景画を描いて", prompt="山と湖", negative="低品質")
    check("handle_image_generation(6a): 成功メッセージに保存先パスを含む", "out/dir/result.png" in _hig_r6a)
    check("handle_image_generation(6a): 成功メッセージにpromptを含む", "山と湖" in _hig_r6a)
    check("handle_image_generation(6a): negativeがtruthyなら'- negative:'行を含む",
          "- negative: 低品質" in _hig_r6a)

    for _hig_empty_neg in ("", None):
        _hig_r6b = f.handle_image_generation("風景画を描いて", prompt="山と湖", negative=_hig_empty_neg)
        check(f"handle_image_generation(6b): negative={_hig_empty_neg!r}では'- negative:'行を含まない",
              "- negative:" not in _hig_r6b)
        check(f"handle_image_generation(6b): negative={_hig_empty_neg!r}でも保存先パスは含む",
              "out/dir/result.png" in _hig_r6b)
finally:
    urllib.request.urlopen = _orig_hig_urlopen
    f.subprocess.run = _orig_hig_subprocess_run
    f.IMAGE_BACKEND = _orig_hig_image_backend
    f._detect_backend = _orig_hig_detect_backend
    f.author_image_prompt = _orig_hig_author_image_prompt
    f.generate_image = _orig_hig_generate_image

check("handle_image_generation: テスト後にurllib.request.urlopenが元の状態へ復元されている",
      urllib.request.urlopen == _orig_hig_urlopen)
check("handle_image_generation: テスト後にf.subprocess.runが元の状態へ復元されている",
      f.subprocess.run == _orig_hig_subprocess_run)
check("handle_image_generation: テスト後にIMAGE_BACKENDが元の状態へ復元されている",
      f.IMAGE_BACKEND == _orig_hig_image_backend)
check("handle_image_generation: テスト後に_detect_backendが元の状態へ復元されている",
      f._detect_backend == _orig_hig_detect_backend)
check("handle_image_generation: テスト後にauthor_image_promptが元の状態へ復元されている",
      f.author_image_prompt == _orig_hig_author_image_prompt)
check("handle_image_generation: テスト後にgenerate_imageが元の状態へ復元されている",
      f.generate_image == _orig_hig_generate_image)

# ============================================================
# _detect_backend / generate_image の直接テスト (iteration 176)
# ============================================================
# 2026-07-26: _detect_backend (fugu_local.py ~L2308) と generate_image
# (~L2320) はこれまで直接呼び出すテストが無く、build_pptx/handle_image_generation
# 経由のテストは常にこの2関数を丸ごとlambdaで差し替えていた(grep確認: f._detect_
# backend/f.generate_imageへの参照は全てlambda代入かsave/restoreのみで、直接呼んで
# 戻り値を検証するテストは無かった)。_detect_backendの分岐(明示指定'a1111'/
# 'comfyui'はプローブなしでそのままpass-through、'auto'はA1111をComfyUIより先に
# プローブしA1111優先で短絡、'off'/未知の値はNone)と、generate_imageのtry/except
# Exception->print->return None(バックエンド呼び出し失敗時に例外を外へ伝播させない
# 契約。高価なMoA/SC計算が完了した後でこの契約が壊れると、計算済みの回答ごと
# ターンをクラッシュさせてしまう。iter41-47/68/80と同種の「計算済みの結果を失わ
# ない」保証)を直接ロックする。_detect_backendのテストではf._backend_upのみを
# モックし(実urlopenには一切触れない)、generate_imageのテストではf._detect_backend/
# f.generate_image_a1111/f.generate_image_comfyuiをモックする。モック漏れを即座に
# 可視化するため、urllib.request.urlopenとf.subprocess.runにも「呼ばれたら即
# AssertionError」の番人を仕込む(iteration 38/76/104の流儀)。f.A1111_URL/
# f.COMFYUI_URLは読むだけで書き換えない。

_orig_db_urlopen = urllib.request.urlopen
_orig_db_subprocess_run = f.subprocess.run
_orig_db_image_backend = f.IMAGE_BACKEND
_orig_db_backend_up = f._backend_up
_orig_db_detect_backend = f._detect_backend
_orig_db_generate_image_a1111 = f.generate_image_a1111
_orig_db_generate_image_comfyui = f.generate_image_comfyui


def _db_no_network_urlopen(*a, **k):
    raise AssertionError("_detect_backend/generate_image: モック漏れで実urlopen(ネットワーク)が呼ばれた")


def _db_no_subprocess_run(*a, **k):
    raise AssertionError("_detect_backend/generate_image: モック漏れで実subprocess.runが呼ばれた")


try:
    urllib.request.urlopen = _db_no_network_urlopen
    f.subprocess.run = _db_no_subprocess_run

    _db_a1111_probe_url = f"{f.A1111_URL}/sdapi/v1/sd-models"
    _db_comfyui_probe_url = f"{f.COMFYUI_URL}/system_stats"
    _db_backend_up_calls = []

    def _db_make_backend_up(url_map):
        def _mock(url):
            _db_backend_up_calls.append(url)
            return url_map.get(url, False)
        return _mock

    # --- (1) 明示指定 'a1111'/'comfyui' -> プローブなしでそのままpass-through ---
    for _db_explicit in ("a1111", "comfyui"):
        f.IMAGE_BACKEND = _db_explicit
        _db_backend_up_calls.clear()
        f._backend_up = _db_make_backend_up({})
        _db_r1 = f._detect_backend()
        check(f"_detect_backend: IMAGE_BACKEND={_db_explicit!r}はそのまま返る(pass-through)",
              _db_r1 == _db_explicit)
        check(f"_detect_backend: IMAGE_BACKEND={_db_explicit!r}では_backend_upが一切呼ばれない",
              len(_db_backend_up_calls) == 0)

    # --- (2) 'auto' + A1111稼働中 -> 'a1111'を返す。A1111側URLのみプローブされ、
    #     ComfyUI側は短絡でプローブされない(A1111優先)。ComfyUI側もupにしておき、
    #     短絡が「たまたまComfyUIがdownだったから」ではないことを確認する ---
    f.IMAGE_BACKEND = "auto"
    _db_backend_up_calls.clear()
    f._backend_up = _db_make_backend_up({_db_a1111_probe_url: True, _db_comfyui_probe_url: True})
    _db_r2 = f._detect_backend()
    check("_detect_backend: auto+A1111稼働中は'a1111'を返す", _db_r2 == "a1111")
    check("_detect_backend: auto+A1111稼働中はA1111側URLのみをプローブする(短絡)",
          _db_backend_up_calls == [_db_a1111_probe_url])

    # --- (3) 'auto' + A1111停止/ComfyUI稼働中 -> 'comfyui'を返す。
    #     A1111が先、ComfyUIが後の順で両方プローブされる ---
    _db_backend_up_calls.clear()
    f._backend_up = _db_make_backend_up({_db_a1111_probe_url: False, _db_comfyui_probe_url: True})
    _db_r3 = f._detect_backend()
    check("_detect_backend: auto+A1111停止/ComfyUI稼働中は'comfyui'を返す", _db_r3 == "comfyui")
    check("_detect_backend: auto+A1111停止時はA1111->ComfyUIの順で両方プローブされる",
          _db_backend_up_calls == [_db_a1111_probe_url, _db_comfyui_probe_url])

    # --- (4) 'auto' + 両方停止 -> None。両方プローブされる ---
    _db_backend_up_calls.clear()
    f._backend_up = _db_make_backend_up({})
    _db_r4 = f._detect_backend()
    check("_detect_backend: auto+両方停止中はNoneを返す", _db_r4 is None)
    check("_detect_backend: auto+両方停止中でも両方プローブされる(A1111->ComfyUIの順)",
          _db_backend_up_calls == [_db_a1111_probe_url, _db_comfyui_probe_url])

    # --- (5) 'off' / 未知の値 -> None。_backend_upは一切呼ばれない ---
    for _db_none_case in ("off", "xyz"):
        f.IMAGE_BACKEND = _db_none_case
        _db_backend_up_calls.clear()
        f._backend_up = _db_make_backend_up({})
        _db_r5 = f._detect_backend()
        check(f"_detect_backend: IMAGE_BACKEND={_db_none_case!r}はNoneを返す", _db_r5 is None)
        check(f"_detect_backend: IMAGE_BACKEND={_db_none_case!r}では_backend_upが一切呼ばれない",
              len(_db_backend_up_calls) == 0)

    # --- (6) generate_image: バックエンド未検出 -> Noneを返し、a1111/comfyui
    #     いずれの低レベル関数も呼ばれない ---
    f._detect_backend = lambda: None
    _gi_calls6 = {"a1111": 0, "comfyui": 0}
    f.generate_image_a1111 = lambda prompt, negative="": (
        _gi_calls6.__setitem__("a1111", _gi_calls6["a1111"] + 1) or "should-not-be-used")
    f.generate_image_comfyui = lambda prompt, negative="": (
        _gi_calls6.__setitem__("comfyui", _gi_calls6["comfyui"] + 1) or "should-not-be-used")
    _gi_r6 = f.generate_image("a cat")
    check("generate_image: バックエンド未検出ではNoneを返す", _gi_r6 is None)
    check("generate_image: バックエンド未検出ではgenerate_image_a1111が呼ばれない",
          _gi_calls6["a1111"] == 0)
    check("generate_image: バックエンド未検出ではgenerate_image_comfyuiが呼ばれない",
          _gi_calls6["comfyui"] == 0)

    # --- (7) generate_image: 'a1111'にdispatch -> generate_image_a1111が
    #     (prompt, negative)で1回だけ呼ばれ、その戻り値をそのまま返す。
    #     generate_image_comfyuiは呼ばれない ---
    f._detect_backend = lambda: "a1111"
    _gi_seen7 = {"calls": 0}

    def _gi_a1111_7(prompt, negative=""):
        _gi_seen7["calls"] += 1
        _gi_seen7["prompt"] = prompt
        _gi_seen7["negative"] = negative
        return "a1111-result.png"

    _gi_comfyui_calls7 = {"n": 0}
    f.generate_image_a1111 = _gi_a1111_7
    f.generate_image_comfyui = lambda prompt, negative="": (
        _gi_comfyui_calls7.__setitem__("n", _gi_comfyui_calls7["n"] + 1) or "wrong-backend")

    _gi_r7 = f.generate_image("a cat", "blurry")
    check("generate_image: 'a1111'検出時はgenerate_image_a1111の戻り値をそのまま返す",
          _gi_r7 == "a1111-result.png")
    check("generate_image: 'a1111'検出時はgenerate_image_a1111がちょうど1回呼ばれる",
          _gi_seen7["calls"] == 1)
    check("generate_image: 'a1111'検出時はpromptとnegativeがそのまま渡る",
          _gi_seen7.get("prompt") == "a cat" and _gi_seen7.get("negative") == "blurry")
    check("generate_image: 'a1111'検出時はgenerate_image_comfyuiが呼ばれない",
          _gi_comfyui_calls7["n"] == 0)

    # --- (8) generate_image: 'comfyui'にdispatch (7の対称) ---
    f._detect_backend = lambda: "comfyui"
    _gi_seen8 = {"calls": 0}

    def _gi_comfyui_8(prompt, negative=""):
        _gi_seen8["calls"] += 1
        _gi_seen8["prompt"] = prompt
        _gi_seen8["negative"] = negative
        return "comfyui-result.png"

    _gi_a1111_calls8 = {"n": 0}
    f.generate_image_comfyui = _gi_comfyui_8
    f.generate_image_a1111 = lambda prompt, negative="": (
        _gi_a1111_calls8.__setitem__("n", _gi_a1111_calls8["n"] + 1) or "wrong-backend")

    _gi_r8 = f.generate_image("a dog", "ugly")
    check("generate_image: 'comfyui'検出時はgenerate_image_comfyuiの戻り値をそのまま返す",
          _gi_r8 == "comfyui-result.png")
    check("generate_image: 'comfyui'検出時はgenerate_image_comfyuiがちょうど1回呼ばれる",
          _gi_seen8["calls"] == 1)
    check("generate_image: 'comfyui'検出時はpromptとnegativeがそのまま渡る",
          _gi_seen8.get("prompt") == "a dog" and _gi_seen8.get("negative") == "ugly")
    check("generate_image: 'comfyui'検出時はgenerate_image_a1111が呼ばれない",
          _gi_a1111_calls8["n"] == 0)

    # --- (9) generate_image: never-crash契約の回帰確認。低レベル関数が例外を
    #     送出しても外へ伝播させず、Noneを返す(iter41-47/68/80と同種の
    #     「計算済みの結果を失わない」保証がこの1関数内にも効いている) ---
    f._detect_backend = lambda: "a1111"

    def _gi_a1111_raises(prompt, negative=""):
        raise RuntimeError("simulated backend failure")

    f.generate_image_a1111 = _gi_a1111_raises
    _gi_r9, _gi_exc9 = "unset", None
    try:
        _gi_r9 = f.generate_image("prompt")
    except Exception as _exc:
        _gi_exc9 = _exc
    check("generate_image: 低レベル関数が例外を送出しても外へ伝播しない",
          _gi_exc9 is None)
    check("generate_image: 低レベル関数が例外を送出した場合はNoneを返す",
          _gi_r9 is None)

    # --- (10) generate_image: negative省略時は既定の''がそのままdispatch先に渡る ---
    f._detect_backend = lambda: "a1111"
    _gi_seen10 = {}

    def _gi_a1111_10(prompt, negative=""):
        _gi_seen10["negative"] = negative
        return "ok.png"

    f.generate_image_a1111 = _gi_a1111_10
    f.generate_image("prompt only")
    check("generate_image: negative省略時は''がdispatch先に渡る",
          _gi_seen10.get("negative") == "")
finally:
    urllib.request.urlopen = _orig_db_urlopen
    f.subprocess.run = _orig_db_subprocess_run
    f.IMAGE_BACKEND = _orig_db_image_backend
    f._backend_up = _orig_db_backend_up
    f._detect_backend = _orig_db_detect_backend
    f.generate_image_a1111 = _orig_db_generate_image_a1111
    f.generate_image_comfyui = _orig_db_generate_image_comfyui

check("_detect_backend/generate_image: テスト後にurllib.request.urlopenが元の状態へ復元されている",
      urllib.request.urlopen == _orig_db_urlopen)
check("_detect_backend/generate_image: テスト後にf.subprocess.runが元の状態へ復元されている",
      f.subprocess.run == _orig_db_subprocess_run)
check("_detect_backend/generate_image: テスト後にIMAGE_BACKENDが元の状態へ復元されている",
      f.IMAGE_BACKEND == _orig_db_image_backend)
check("_detect_backend/generate_image: テスト後にf._backend_upが元の状態へ復元されている",
      f._backend_up == _orig_db_backend_up)
check("_detect_backend/generate_image: テスト後にf._detect_backendが元の状態へ復元されている",
      f._detect_backend == _orig_db_detect_backend)
check("_detect_backend/generate_image: テスト後にf.generate_image_a1111が元の状態へ復元されている",
      f.generate_image_a1111 == _orig_db_generate_image_a1111)
check("_detect_backend/generate_image: テスト後にf.generate_image_comfyuiが元の状態へ復元されている",
      f.generate_image_comfyui == _orig_db_generate_image_comfyui)

# ============================================================
# _slack_truncate / notify_slack のテスト
# ============================================================
# SLACK_WEBHOOK_URL / SLACK_Q_PREVIEW / SLACK_A_PREVIEW / SLACK_NOTIFY_TIMEOUT を
# 一時的に書き換えるため、finallyで必ず元へ戻す。urllib.request.urlopenも同様。

check("_slack_truncate: 上限以下のテキストはそのまま", f._slack_truncate("hello", 10) == "hello")
check("_slack_truncate: ちょうど上限のテキストはそのまま", f._slack_truncate("abcde", 5) == "abcde")
_st_long = "a" * 20
check("_slack_truncate: 上限超過テキストはlimit文字+省略記号に切り詰め",
      f._slack_truncate(_st_long, 10) == ("a" * 10) + "…")
check("_slack_truncate: 切り詰め結果の長さはlimit+1(省略記号込み)",
      len(f._slack_truncate(_st_long, 10)) == 11)
check("_slack_truncate: Noneは空文字扱いで例外を送出しない", f._slack_truncate(None, 10) == "")
check("_slack_truncate: 空文字はそのまま(空文字)", f._slack_truncate("", 10) == "")
check("_slack_truncate: 前後空白はstripされる", f._slack_truncate("  hi  ", 10) == "hi")

_orig_ns_webhook = f.SLACK_WEBHOOK_URL
_orig_ns_q_preview = f.SLACK_Q_PREVIEW
_orig_ns_a_preview = f.SLACK_A_PREVIEW
_orig_ns_timeout = f.SLACK_NOTIFY_TIMEOUT
_orig_ns_urlopen = urllib.request.urlopen
try:
    # --- (1) SLACK_WEBHOOK_URLがfalsy(None/'') -> urlopenは一切呼ばれず例外も出ない ---
    for _ns_falsy in (None, ""):
        f.SLACK_WEBHOOK_URL = _ns_falsy

        def _ns_no_call_urlopen(*a, **k):
            raise AssertionError("notify_slack: SLACK_WEBHOOK_URLがfalsyなのにurlopenが呼ばれた")

        urllib.request.urlopen = _ns_no_call_urlopen
        _ns_exc = None
        try:
            f.notify_slack("質問です", "回答です", 1.23)
        except Exception as _exc:
            _ns_exc = _exc
        check(f"notify_slack: SLACK_WEBHOOK_URL={_ns_falsy!r}ではurlopenを呼ばず例外も出ない",
              _ns_exc is None)

    # --- (2) 成功時: 成功アイコンを含み、失敗アイコンは含まない。長文は切り詰められる ---
    f.SLACK_WEBHOOK_URL = "https://hooks.slack.example/T000/B000/xxxx"
    f.SLACK_Q_PREVIEW = 20
    f.SLACK_A_PREVIEW = 30
    _ns_calls = []

    class _NsOkResp:
        def __enter__(self):
            return self

        def __exit__(self, *a):
            return False

        def read(self):
            return b"ok"

    def _ns_ok_urlopen(req, timeout=None):
        _ns_calls.append((req, timeout))
        return _NsOkResp()

    urllib.request.urlopen = _ns_ok_urlopen
    _ns_long_q = "質問" * 30
    _ns_long_a = "回答" * 30
    f.notify_slack(_ns_long_q, _ns_long_a, 4.56)
    check("notify_slack: 成功時にurlopenが1回だけ呼ばれる", len(_ns_calls) == 1)
    _ns_req = _ns_calls[0][0]
    _ns_body = json.loads(_ns_req.data.decode("utf-8"))["text"]
    check("notify_slack: 成功時は成功アイコンを含む", ":white_check_mark:" in _ns_body)
    check("notify_slack: 成功時は失敗アイコンを含まない", ":x:" not in _ns_body)
    check("notify_slack: 質問はSLACK_Q_PREVIEW文字に切り詰められる",
          f._slack_truncate(_ns_long_q, f.SLACK_Q_PREVIEW) in _ns_body)
    check("notify_slack: 回答はSLACK_A_PREVIEW文字に切り詰められる",
          f._slack_truncate(_ns_long_a, f.SLACK_A_PREVIEW) in _ns_body)
    check("notify_slack: timeout引数にSLACK_NOTIFY_TIMEOUTを渡す",
          _ns_calls[0][1] == f.SLACK_NOTIFY_TIMEOUT)

    # --- (3) 失敗時(__ERROR__プレフィックス): 失敗アイコンを含み、成功アイコンは含まない ---
    _ns_calls.clear()
    urllib.request.urlopen = _ns_ok_urlopen
    f.notify_slack("質問です", "__ERROR__: something broke", 7.89)
    check("notify_slack: 失敗時にurlopenが1回だけ呼ばれる", len(_ns_calls) == 1)
    _ns_req2 = _ns_calls[0][0]
    _ns_body2 = json.loads(_ns_req2.data.decode("utf-8"))["text"]
    check("notify_slack: __ERROR__プレフィックスでは失敗アイコンを含む", ":x:" in _ns_body2)
    check("notify_slack: __ERROR__プレフィックスでは成功アイコンを含まない",
          ":white_check_mark:" not in _ns_body2)

    # --- (4) urlopenが例外を送出しても呼び出し元へ伝播しない(exceptで握り潰される) ---
    def _ns_raise_urlopen(req, timeout=None):
        raise RuntimeError("ネットワーク断(模擬)")

    urllib.request.urlopen = _ns_raise_urlopen
    _ns_exc2 = None
    try:
        f.notify_slack("質問です", "回答です", 0.1)
    except Exception as _exc:
        _ns_exc2 = _exc
    check("notify_slack: urlopenが例外を送出しても呼び出し元へは伝播しない", _ns_exc2 is None)
finally:
    f.SLACK_WEBHOOK_URL = _orig_ns_webhook
    f.SLACK_Q_PREVIEW = _orig_ns_q_preview
    f.SLACK_A_PREVIEW = _orig_ns_a_preview
    f.SLACK_NOTIFY_TIMEOUT = _orig_ns_timeout
    urllib.request.urlopen = _orig_ns_urlopen

check("notify_slack: テスト後にSLACK_WEBHOOK_URLが元の状態へ復元されている",
      f.SLACK_WEBHOOK_URL == _orig_ns_webhook)
check("notify_slack: テスト後にSLACK_Q_PREVIEWが元の状態へ復元されている",
      f.SLACK_Q_PREVIEW == _orig_ns_q_preview)
check("notify_slack: テスト後にSLACK_A_PREVIEWが元の状態へ復元されている",
      f.SLACK_A_PREVIEW == _orig_ns_a_preview)
check("notify_slack: テスト後にurllib.request.urlopenが元の状態へ復元されている",
      urllib.request.urlopen == _orig_ns_urlopen)

# ---------- read_file_text: 拡張子ディスパッチの直接検証 (2026-07-24) ----------
# 個々の _read_pdf/_read_docx/_read_excel/_read_pptx の中身は既存テストで検証済みだが、
# read_file_text がどの拡張子をどのリーダーへ振り分けるか（大文字小文字の扱い、
# 旧形式 .doc/.ppt/.xls のルーティング、_BINARY_SKIP境界）自体は未検証だった。
# --file で渡されるパスはユーザー入力由来で大文字拡張子(.PDF等)もあり得るため、
# ディスパッチ自体にバグがあると RAG/--file コンテキストが静かに劣化する
# (精度critical パスへの入力なので影響が大きい)。
# ここでは各 _read_* をmonkeypatchして「どのリーダーが呼ばれたか」だけを見る
# ことで、リーダー内部から独立してディスパッチ経路だけを検証する。
import tempfile as _rfd_tempfile
import pathlib as _rfd_pathlib

_orig_read_pdf = f._read_pdf
_orig_read_docx = f._read_docx
_orig_read_excel = f._read_excel
_orig_read_pptx = f._read_pptx
_orig_read_html = f._read_html
_orig_read_ipynb = f._read_ipynb

_rfd_calls = []


def _mk_fake_reader(tag):
    def _fake(path):
        _rfd_calls.append(tag)
        return f"CALLED:{tag}"
    return _fake


try:
    f._read_pdf = _mk_fake_reader("pdf")
    f._read_docx = _mk_fake_reader("docx")
    f._read_excel = _mk_fake_reader("excel")
    f._read_pptx = _mk_fake_reader("pptx")
    f._read_html = _mk_fake_reader("html")
    f._read_ipynb = _mk_fake_reader("ipynb")

    with _rfd_tempfile.TemporaryDirectory() as _rfd_td:
        _rfd_root = _rfd_pathlib.Path(_rfd_td)

        # (1) 小文字拡張子は期待通りのリーダーへ
        _rfd_calls.clear()
        check("read_file_text: .pdf(小文字)は_read_pdfへ振り分け",
              f.read_file_text(_rfd_root / "a.pdf") == "CALLED:pdf" and _rfd_calls == ["pdf"])
        _rfd_calls.clear()
        check("read_file_text: .docx(小文字)は_read_docxへ振り分け",
              f.read_file_text(_rfd_root / "a.docx") == "CALLED:docx" and _rfd_calls == ["docx"])
        _rfd_calls.clear()
        check("read_file_text: .xlsx(小文字)は_read_excelへ振り分け",
              f.read_file_text(_rfd_root / "a.xlsx") == "CALLED:excel" and _rfd_calls == ["excel"])
        _rfd_calls.clear()
        check("read_file_text: .pptx(小文字)は_read_pptxへ振り分け",
              f.read_file_text(_rfd_root / "a.pptx") == "CALLED:pptx" and _rfd_calls == ["pptx"])

        # (2) 大文字・混在大小文字の拡張子も同じリーダーへ(ユーザー入力--fileパス対策)
        _rfd_calls.clear()
        check("read_file_text: .PDF(大文字)も_read_pdfへ振り分け(小文字と同一)",
              f.read_file_text(_rfd_root / "a.PDF") == "CALLED:pdf" and _rfd_calls == ["pdf"])
        _rfd_calls.clear()
        check("read_file_text: .Docx(混在大小文字)も_read_docxへ振り分け(小文字と同一)",
              f.read_file_text(_rfd_root / "a.Docx") == "CALLED:docx" and _rfd_calls == ["docx"])
        _rfd_calls.clear()
        check("read_file_text: .XLSX(大文字)も_read_excelへ振り分け(小文字と同一)",
              f.read_file_text(_rfd_root / "a.XLSX") == "CALLED:excel" and _rfd_calls == ["excel"])
        _rfd_calls.clear()
        check("read_file_text: .PPTX(大文字)も_read_pptxへ振り分け(小文字と同一)",
              f.read_file_text(_rfd_root / "a.PPTX") == "CALLED:pptx" and _rfd_calls == ["pptx"])

        # (3) 旧形式(.doc/.ppt/.xls)は現行形式と同じリーダー関数へルーティングされる
        _rfd_calls.clear()
        check("read_file_text: 旧形式.docは.docxと同じ_read_docxへ振り分け",
              f.read_file_text(_rfd_root / "a.doc") == "CALLED:docx" and _rfd_calls == ["docx"])
        _rfd_calls.clear()
        check("read_file_text: 旧形式.xlsは.xlsxと同じ_read_excelへ振り分け",
              f.read_file_text(_rfd_root / "a.xls") == "CALLED:excel" and _rfd_calls == ["excel"])
        _rfd_calls.clear()
        check("read_file_text: 旧形式.pptは.pptxと同じ_read_pptxへ振り分け",
              f.read_file_text(_rfd_root / "a.ppt") == "CALLED:pptx" and _rfd_calls == ["pptx"])
        # 旧形式の大文字混在も同様(.DOC等)
        _rfd_calls.clear()
        check("read_file_text: 旧形式.DOC(大文字)も_read_docxへ振り分け",
              f.read_file_text(_rfd_root / "a.DOC") == "CALLED:docx" and _rfd_calls == ["docx"])

        # (4) html/ipynbも念のため確認(既存の分岐だが、ディスパッチ経路として一括網羅)
        _rfd_calls.clear()
        check("read_file_text: .htmは_read_htmlへ振り分け",
              f.read_file_text(_rfd_root / "a.htm") == "CALLED:html" and _rfd_calls == ["html"])
        _rfd_calls.clear()
        check("read_file_text: .ipynbは_read_ipynbへ振り分け",
              f.read_file_text(_rfd_root / "a.ipynb") == "CALLED:ipynb" and _rfd_calls == ["ipynb"])

        # (5) _BINARY_SKIP境界: スキップ対象拡張子はどのリーダーも呼ばずに""を返す
        _rfd_calls.clear()
        _rfd_png = _rfd_root / "a.png"
        _rfd_png.write_bytes(b"\x89PNG fake binary content")
        check("read_file_text: _BINARY_SKIP対象(.png)は\"\"を返しどのリーダーも呼ばない",
              f.read_file_text(_rfd_png) == "" and _rfd_calls == [])
        _rfd_calls.clear()
        _rfd_exe = _rfd_root / "a.exe"
        _rfd_exe.write_bytes(b"MZ fake binary content")
        check("read_file_text: _BINARY_SKIP対象(.exe)は\"\"を返しどのリーダーも呼ばない",
              f.read_file_text(_rfd_exe) == "" and _rfd_calls == [])

        # (6) _BINARY_SKIPにも既知リーダーにも属さない未知拡張子は汎用テキスト読み込みへ
        #     フォールスルーする(どのfake readerも呼ばれず、ファイル内容がそのまま返る)
        _rfd_calls.clear()
        _rfd_unknown = _rfd_root / "a.xyz123"
        _rfd_unknown_content = "unknown suffix generic text fallback content"
        _rfd_unknown.write_text(_rfd_unknown_content, encoding="utf-8")
        check("read_file_text: 未知拡張子(.xyz123)はリーダーを介さず汎用テキスト読み込みへフォールスルー",
              f.read_file_text(_rfd_unknown) == _rfd_unknown_content and _rfd_calls == [])
finally:
    f._read_pdf = _orig_read_pdf
    f._read_docx = _orig_read_docx
    f._read_excel = _orig_read_excel
    f._read_pptx = _orig_read_pptx
    f._read_html = _orig_read_html
    f._read_ipynb = _orig_read_ipynb

check("read_file_text: テスト後に_read_pdfが元の関数へ復元されている", f._read_pdf == _orig_read_pdf)
check("read_file_text: テスト後に_read_docxが元の関数へ復元されている", f._read_docx == _orig_read_docx)
check("read_file_text: テスト後に_read_excelが元の関数へ復元されている", f._read_excel == _orig_read_excel)
check("read_file_text: テスト後に_read_pptxが元の関数へ復元されている", f._read_pptx == _orig_read_pptx)
check("read_file_text: テスト後に_read_htmlが元の関数へ復元されている", f._read_html == _orig_read_html)
check("read_file_text: テスト後に_read_ipynbが元の関数へ復元されている", f._read_ipynb == _orig_read_ipynb)

# ---------- read_file_text: .htm確認 + _CODE_EXTENSIONS/_BINARY_SKIPの相互排他性 (2026-07-24 / iter121) ----------
# 発端: 「iter120は.htm/.ipynbを'covers'とだけ記載し、.htmが本当に_read_htmlへ
# ルーティングされるか未確認では」という懸念が上がった。実際にはコード上(L1188:
# `if suffix in {".html", ".htm"}: return _read_html(path)`)でも、直上のテスト
# ブロック(2026-07-24, 「(4) html/ipynbも念のため確認」)でも .htm→_read_html は
# 既に振り分け確認済みであり、本物のディスパッチ漏れではなかった
# (静的レビューの結果、コード変更は不要と判断。強引な「修正」はしない)。
# ついでに依頼された第二の観点も点検する: _CODE_EXTENSIONS (iter60, 25拡張子)の
# 各拡張子が read_file_text の専用ディスパッチ(_BINARY_SKIP・PDF/Word/Excel/
# PowerPoint/HTML/ipynb)のどれとも重ならず、意図通り汎用テキスト分岐へ
# フォールスルーしてソースコードがそのまま読めることを、全25拡張子について
# 直接 read_file_text() を呼んで確認する(lang_mapとの同期は既存テストで別途
# カバー済みなので、ここではread_file_text側のディスパッチのみを見る)。
_rft_special_dispatch_suffixes = {
    ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt",
    ".html", ".htm", ".ipynb",
}
check("read_file_text: _CODE_EXTENSIONSと_BINARY_SKIPは重複しない(コード拡張子がバイナリ扱いされない)",
      not (f._CODE_EXTENSIONS & f._BINARY_SKIP))
check("read_file_text: _CODE_EXTENSIONSは専用リーダー(pdf/docx/xlsx/pptx/html/ipynb)とも重複しない",
      not (f._CODE_EXTENSIONS & _rft_special_dispatch_suffixes))

with _rfd_tempfile.TemporaryDirectory() as _rce_td:
    _rce_root = _rfd_pathlib.Path(_rce_td)
    _rce_all_fallthrough = True
    for _rce_suffix in sorted(f._CODE_EXTENSIONS):
        _rce_fp = _rce_root / f"sample{_rce_suffix}"
        _rce_content = f"sample code content for {_rce_suffix}"
        _rce_fp.write_text(_rce_content, encoding="utf-8")
        if f.read_file_text(_rce_fp) != _rce_content:
            _rce_all_fallthrough = False
    check("read_file_text: _CODE_EXTENSIONS全25拡張子が汎用テキスト読み込みへフォールスルーし内容がそのまま返る",
          _rce_all_fallthrough)

# 2026-07-24: normalize_answer の \frac/\dfrac/\tfrac 数値正規化（iteration 122、
# iteration 108 で3回スタックしたまま未着手だった修正のリトライ）。
check("normalize_answer: \\frac{1}{2} -> 1/2",
      f.normalize_answer(r"\frac{1}{2}") == "1/2")
check("normalize_answer: \\dfrac{3}{4} -> 3/4",
      f.normalize_answer(r"\dfrac{3}{4}") == "3/4")
check("normalize_answer: \\tfrac{-1}{2} -> -1/2",
      f.normalize_answer(r"\tfrac{-1}{2}") == "-1/2")
check("normalize_answer: 外側マイナス + 分子マイナス -> 相殺",
      f.normalize_answer(r"-\frac{-1}{2}") == "1/2")
check("normalize_answer: ネストした \\frac は素通り(クラッシュせず未変更)",
      f.normalize_answer(r"\frac{\frac{1}{2}}{3}") == r"\frac{\frac{1}{2}}{3}")
check("normalize_answer: 変数を含む \\frac は素通り(クラッシュせず未変更)",
      f.normalize_answer(r"\frac{x}{2}") == r"\frac{x}{2}")


def _t_frac_fastpath_equiv(calls):
    result = f.answers_equivalent(r"\frac{1}{2}", "0.5")
    check("answers_equivalent: \\frac{1}{2} と 0.5 がFractionファストパスで一致(math_verify不使用)",
          result is True)
    check("answers_equivalent: \\frac{1}{2}/0.5 一致判定はmath_verify.parseを一切呼ばない(高速パス経由)",
          len(calls["parse_args"]) == 0)


_run_with_math_verify_stub(None, "parse", _t_frac_fastpath_equiv)

# 2026-07-25: iteration 122 の \frac{a}{b} 正規化（すぐ上のブロック）は分子・先頭の
# 符号だけを neg へ XOR しており（コメントに明記の通りスコープを numerator/leading
# sign に限定）、分母の符号は無視して den をそのまま埋め込んでいた。だが _frac_re の
# 分母グループ (-?\d+(?:\.\d+)?) は先頭マイナスを許容するため、\frac{1}{-2}（=-1/2）は
# "1/-2" に、\frac{-1}{-2}（=+1/2）は "-1/-2" になってしまい、Fraction() は分母に符号が
# 付いた文字列を拒否するのでどちらも下の Fraction 高速パスに乗れずすり抜け、-1/2・
# \frac{-1}{2}・-0.5 という「素直な」表記の票とは別の投票クラスに分裂して自己整合性投票
# （gotcha #7）の票を薄めていた。iteration 13/22/24/30/78/122/134/136/140 と同系統の
# 「ファストパスで拾えない票を拾う」修正として、分母の符号も neg へ折り込んでから den から
# 取り除く（正の分母では den.startswith("-") が False で lstrip("-") はno-opなので、
# 上の iteration 122/134/136/140 の既存挙動はバイト単位で不変のまま）。
check("normalize_answer: \\frac{1}{-2} -> -1/2 (分母の符号をnegへ折り込む)",
      f.normalize_answer(r"\frac{1}{-2}") == "-1/2")
check("normalize_answer: \\frac{-1}{-2} -> 1/2 (分子・分母の符号が相殺)",
      f.normalize_answer(r"\frac{-1}{-2}") == "1/2")
check("normalize_answer: \\frac{3}{-4} -> -3/4",
      f.normalize_answer(r"\frac{3}{-4}") == "-3/4")
check("normalize_answer: \\dfrac{1}{-2} -> -1/2 ([dt]?fracバリアント)",
      f.normalize_answer(r"\dfrac{1}{-2}") == "-1/2")
check("normalize_answer: \\tfrac{5}{-6} -> -5/6 ([dt]?fracバリアント)",
      f.normalize_answer(r"\tfrac{5}{-6}") == "-5/6")
# 回帰: 正の分母の既存挙動はバイト単位で不変（den.startswith("-") が False のため
# den.lstrip("-") はno-op）。iteration 122/134/136/140 のテストはこのブロックとは別に
# そのまま残しており、ここでは変更前と同じ入出力を重ねて確認する。
check("normalize_answer: 分母符号折り込み後も正の分母は不変(1)",
      f.normalize_answer(r"\frac{1}{2}") == "1/2")
check("normalize_answer: 分母符号折り込み後も正の分母は不変(2)",
      f.normalize_answer(r"-\frac{1}{2}") == "-1/2")
check("normalize_answer: 分母符号折り込み後も正の分母は不変(3)",
      f.normalize_answer(r"\frac{-1}{2}") == "-1/2")


def _t_neg_den_frac_fastpath_equiv(calls):
    result_a = f.answers_equivalent(r"\frac{1}{-2}", "-1/2")
    result_b = f.answers_equivalent(r"\frac{1}{-2}", "-0.5")
    check("answers_equivalent: \\frac{1}{-2} と -1/2 がFractionファストパスで一致(math_verify不使用)",
          result_a is True)
    check("answers_equivalent: \\frac{1}{-2} と -0.5 がFractionファストパスで一致(math_verify不使用)",
          result_b is True)
    check("answers_equivalent: 負の分母frac一致判定はmath_verify.parseを一切呼ばない(高速パス経由)",
          len(calls["parse_args"]) == 0)


_run_with_math_verify_stub(None, "parse", _t_neg_den_frac_fastpath_equiv)

_neg_den_votes = [r"\frac{1}{-2}", "-1/2", "-0.5"]
_neg_den_top, _neg_den_count, _neg_den_classes = f.vote_answers(_neg_den_votes)
check("vote_answers: \\frac{1}{-2}/-1/2/-0.5 の3票が単一クラスに集約される(票割れしない)",
      _neg_den_count == 3 and len(_neg_den_classes) == 1)

# ---------- normalize_answer: 数字に挟まれたカンマの直後の空白除去 (2026-07-25) ----------
# (3, 4) と (3,4) は同じ順序対/区間/座標値だが、上の桁区切り正規表現
# (?<=\d),\s*(?=\d{3}\b) はカンマ直後にちょうど3桁が続く場合しか吸収しないため
# ("4)" は \d{3} に一致しない)、このままでは2つの投票クラスに分裂する
# (na.lower()/Fraction高速パス不通過 -> gotcha #6 のmath_verifyフォールバックに依存、
# gotcha #7 の自己整合性投票が最も嫌う票割れ)。iteration 13/22/24/30/78/122/134/136/
# 140/148 と同系統の「ファストパスで拾えない票を拾う」姉妹修正。
check("normalize_answer: (3, 4) -> (3,4) (順序対の空白除去)",
      f.normalize_answer("(3, 4)") == "(3,4)")
check("normalize_answer: [2, 5) -> [2,5) (区間の空白除去)",
      f.normalize_answer("[2, 5)") == "[2,5)")
check("normalize_answer: 1, 2 -> 1,2",
      f.normalize_answer("1, 2") == "1,2")
check("normalize_answer: 3.14, 2.71 -> 3.14,2.71 (小数点混じりでも数字境界で判定)",
      f.normalize_answer("3.14, 2.71") == "3.14,2.71")
check("normalize_answer: 既に空白除去済みの(3,4)は不変(冪等性)",
      f.normalize_answer("(3,4)") == "(3,4)")
# 桁区切り回帰: 新ルールが上の桁区切り除去 (12, 345 -> 12345) を妨げないこと
check("normalize_answer: 桁区切り回帰 1,234 -> 1234 (新ルールと非干渉)",
      f.normalize_answer("1,234") == "1234")
check("normalize_answer: 桁区切り回帰 12, 345 -> 12345 (新ルールと非干渉)",
      f.normalize_answer("12, 345") == "12345")
# 散文/記号的カンマの安全性: 数字に挟まれていないカンマは対象外のまま
check("normalize_answer: 散文カンマ'yes, it is'は不変(数字に挟まれていない)",
      f.normalize_answer("yes, it is") == "yes, it is")
check("normalize_answer: 記号的タプル(x, y)は不変(数字に挟まれていない)",
      f.normalize_answer("(x, y)") == "(x, y)")
# 末尾カンマ回帰(iteration 22): 末尾の','は引き続き除去される
check("normalize_answer: 末尾カンマ回帰 42, -> 42 (新ルールと非干渉)",
      f.normalize_answer("42,") == "42")


def _t_tuple_comma_fastpath_equiv(calls):
    result = f.answers_equivalent("(3, 4)", "(3,4)")
    check("answers_equivalent: (3, 4) と (3,4) が正規化高速パスで一致(math_verify不使用)",
          result is True)
    check("answers_equivalent: 順序対の一致判定はmath_verify.parseを一切呼ばない(高速パス経由)",
          len(calls["parse_args"]) == 0)


_run_with_math_verify_stub(None, "parse", _t_tuple_comma_fastpath_equiv)

_tuple_comma_votes = ["(3, 4)", "(3,4)", "(3,4)"]
_tuple_comma_top, _tuple_comma_count, _tuple_comma_classes = f.vote_answers(_tuple_comma_votes)
check("vote_answers: (3, 4)/(3,4)/(3,4) の3票が単一クラスに集約される(票割れしない)",
      _tuple_comma_count == 3 and len(_tuple_comma_classes) == 1)

# ---------- _read_docx/_read_pptx: Document()/Presentation()の実行時例外もクラッシュせず
# notice文字列に劣化させる (2026-07-24 / iter123) ----------
# iter83(_read_pdf)/iter84(_read_excel)が確立した「except ImportErrorだけでは
# 実際にライブラリがインストールされていて実行時に失敗するケース(破損ファイル・
# レガシーバイナリを現代拡張子のまま開いた場合等)を捕捉できない」という穴を
# _read_docx/_read_pptxにも同じ方針で塞ぐ。ここでは(1)本物のpython-docx/
# python-pptxが利用可能な環境でも決定的に検証できるよう、乱数バイト列(zipにも
# レガシーバイナリにもならない不正データ)を実際のtempファイルとして書き、
# Document()/Presentation()に本当に例外を送出させるケースと、(2)iter83/84の
# sys.modules差し替えスタイルを踏襲し、docx.Document/pptx.Presentationを
# PackageNotFoundError等を送出するフェイクに差し替えるケースの両方を確認する。
import random as _rdx_random

_rdx_orig_docx_mod = sys.modules.get("docx")
_rdx_orig_pptx_mod = sys.modules.get("pptx")

with _rfd_tempfile.TemporaryDirectory() as _rdx_td:
    _rdx_root = _rfd_pathlib.Path(_rdx_td)

    # (1) 破損/ガーベジバイト列の.docxを実際のtempファイルとして用意し、
    #     python-docxが利用可能ならDocument()に本当に例外を投げさせて確認する。
    _rdx_corrupt_docx = _rdx_root / "corrupt.docx"
    _rdx_random.seed(123)
    _rdx_corrupt_docx.write_bytes(bytes(_rdx_random.getrandbits(8) for _ in range(512)))

    try:
        import docx as _rdx_probe  # noqa: F401
        _rdx_docx_available = True
    except ImportError:
        _rdx_docx_available = False

    if _rdx_docx_available:
        _rdx_cap1 = io.StringIO()
        with contextlib.redirect_stdout(_rdx_cap1):
            _rdx_result1 = f._read_docx(_rdx_corrupt_docx)
        check("_read_docx: 破損/ガーベジバイト列の.docxを実python-docxに渡しても例外を送出せず文字列を返す",
              isinstance(_rdx_result1, str))
        check("_read_docx: 破損.docxの読み込み失敗はnotice文字列に劣化する(先頭が'[DOCX:')",
              _rdx_result1.startswith("[DOCX:") and "読み込みエラー" in _rdx_result1)
        check("_read_docx: 破損.docxのnotice文字列はファイル名を含む",
              _rdx_corrupt_docx.name in _rdx_result1)
        check("_read_docx: 破損.docxの読み込み失敗の警告メッセージはcp932でエンコード可能(gotcha#4)",
              _rpdf_is_cp932_safe(_rdx_cap1.getvalue()))
        check("_is_lib_missing_notice: _read_docxの実行時エラーnotice(pip installを含まない)はFalseと判定される"
              "(未インストールと誤判定させない意図的な区別)",
              f._is_lib_missing_notice(_rdx_result1) is False)
    else:
        print("   [SKIP] python-docx未インストールのため破損.docx実ファイルテストをスキップ")

    # (2) 破損/ガーベジバイト列の.pptxも同様に確認する。
    _rdx_corrupt_pptx = _rdx_root / "corrupt.pptx"
    _rdx_random.seed(456)
    _rdx_corrupt_pptx.write_bytes(bytes(_rdx_random.getrandbits(8) for _ in range(512)))

    try:
        import pptx as _rpx_probe  # noqa: F401
        _rdx_pptx_available = True
    except ImportError:
        _rdx_pptx_available = False

    if _rdx_pptx_available:
        _rdx_cap2 = io.StringIO()
        with contextlib.redirect_stdout(_rdx_cap2):
            _rdx_result2 = f._read_pptx(_rdx_corrupt_pptx)
        check("_read_pptx: 破損/ガーベジバイト列の.pptxを実python-pptxに渡しても例外を送出せず文字列を返す",
              isinstance(_rdx_result2, str))
        check("_read_pptx: 破損.pptxの読み込み失敗はnotice文字列に劣化する(先頭が'[PPTX:')",
              _rdx_result2.startswith("[PPTX:") and "読み込みエラー" in _rdx_result2)
        check("_read_pptx: 破損.pptxのnotice文字列はファイル名を含む",
              _rdx_corrupt_pptx.name in _rdx_result2)
        check("_read_pptx: 破損.pptxの読み込み失敗の警告メッセージはcp932でエンコード可能(gotcha#4)",
              _rpdf_is_cp932_safe(_rdx_cap2.getvalue()))
        check("_is_lib_missing_notice: _read_pptxの実行時エラーnotice(pip installを含まない)はFalseと判定される"
              "(未インストールと誤判定させない意図的な区別)",
              f._is_lib_missing_notice(_rdx_result2) is False)
    else:
        print("   [SKIP] python-pptx未インストールのため破損.pptx実ファイルテストをスキップ")

    # (3) iter83/84のsys.modules差し替えスタイル: docx.Document/pptx.Presentationを
    #     実ライブラリの有無に関わらず決定的にPackageNotFoundError相当で失敗させる
    #     フェイクモジュールに差し替え、環境非依存で同じ挙動を確認する。
    _rdx_legacy_doc = _rdx_root / "legacy.doc"
    _rdx_legacy_doc.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1 legacy binary ole marker")

    class _RdxFakePackageNotFoundError(Exception):
        pass

    def _rdx_fake_document_raises(path):
        raise _RdxFakePackageNotFoundError(f"Package not found at '{path}'")

    _rdx_fake_docx_mod = types.ModuleType("docx")
    _rdx_fake_docx_mod.Document = _rdx_fake_document_raises
    # oxml/table/text サブモジュールへの遅延importは本経路では到達しないため
    # フェイクモジュールに用意する必要はない(Document()の時点で例外が飛ぶ)。

    _rdx_cap3 = io.StringIO()
    with contextlib.redirect_stdout(_rdx_cap3):
        _rdx_result3 = _rpdf_swap_modules(
            {"docx": _rdx_fake_docx_mod},
            lambda: f._read_docx(_rdx_legacy_doc),
        )
    check("_read_docx: docx.Documentがモンキーパッチでruntime例外(PackageNotFoundError相当)を"
          "送出しても伝播せず文字列を返す(iter83/84スタイル)",
          isinstance(_rdx_result3, str) and _rdx_result3.startswith("[DOCX:") and "読み込みエラー" in _rdx_result3)
    check("_read_docx: モンキーパッチ経路の警告に例外型名が出力される",
          "_RdxFakePackageNotFoundError" in _rdx_cap3.getvalue())

    _rdx_legacy_ppt = _rdx_root / "legacy.ppt"
    _rdx_legacy_ppt.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1 legacy binary ole marker")

    def _rdx_fake_presentation_raises(path):
        raise _RdxFakePackageNotFoundError(f"Package not found at '{path}'")

    _rdx_fake_pptx_mod = types.ModuleType("pptx")
    _rdx_fake_pptx_mod.Presentation = _rdx_fake_presentation_raises

    _rdx_cap4 = io.StringIO()
    with contextlib.redirect_stdout(_rdx_cap4):
        _rdx_result4 = _rpdf_swap_modules(
            {"pptx": _rdx_fake_pptx_mod},
            lambda: f._read_pptx(_rdx_legacy_ppt),
        )
    check("_read_pptx: pptx.Presentationがモンキーパッチでruntime例外(PackageNotFoundError相当)を"
          "送出しても伝播せず文字列を返す(iter83/84スタイル)",
          isinstance(_rdx_result4, str) and _rdx_result4.startswith("[PPTX:") and "読み込みエラー" in _rdx_result4)
    check("_read_pptx: モンキーパッチ経路の警告に例外型名が出力される",
          "_RdxFakePackageNotFoundError" in _rdx_cap4.getvalue())

check("_read_docx: テスト後にsys.modulesの'docx'エントリが元通り解決可能(復元確認、iter123)",
      (sys.modules.get("docx") is _rdx_orig_docx_mod) if _rdx_orig_docx_mod is not None
      else ("docx" not in sys.modules or sys.modules["docx"] is not None))
check("_read_pptx: テスト後にsys.modulesの'pptx'エントリが元通り解決可能(復元確認、iter123)",
      (sys.modules.get("pptx") is _rdx_orig_pptx_mod) if _rdx_orig_pptx_mod is not None
      else ("pptx" not in sys.modules or sys.modules["pptx"] is not None))

# ---------- _read_docx/_read_pptx: 成功パス(正常なファイル)の出力は変更前と不変 (2026-07-24 / iter123) ----------
# 開放/パース周りにtry/exceptを追加しただけで、正常系の抽出ロジック(iter82/87/91/93/98の
# テーブル/グループ/ノート/マージセル/読み順テスト)には一切触れていないため、成功時の
# 出力が変わらないことを軽く再確認する(詳細な回帰ガードは既存のiter82等のテストが担う)。
if _rdx_docx_available:
    with _rfd_tempfile.TemporaryDirectory() as _rdx_ok_td:
        _rdx_ok_path = _rfd_pathlib.Path(_rdx_ok_td) / "ok.docx"
        _rdx_ok_doc = _rdx_probe.Document()
        _rdx_ok_doc.add_paragraph("iter123 success path paragraph")
        _rdx_ok_doc.save(str(_rdx_ok_path))
        _rdx_ok_result = f._read_docx(_rdx_ok_path)
        check("_read_docx: 正常な.docxの成功パス出力は変更後も期待通り(notice文字列に劣化しない)",
              _rdx_ok_result == "iter123 success path paragraph")

if _rdx_pptx_available:
    with _rfd_tempfile.TemporaryDirectory() as _rdx_ok_td2:
        _rdx_ok_path2 = _rfd_pathlib.Path(_rdx_ok_td2) / "ok.pptx"
        _rdx_ok_prs = _rpx_probe.Presentation()
        _rdx_ok_slide_layout = _rdx_ok_prs.slide_layouts[0]
        _rdx_ok_slide = _rdx_ok_prs.slides.add_slide(_rdx_ok_slide_layout)
        _rdx_ok_slide.shapes.title.text = "iter123 success path title"
        _rdx_ok_prs.save(str(_rdx_ok_path2))
        _rdx_ok_result2 = f._read_pptx(_rdx_ok_path2)
        check("_read_pptx: 正常な.pptxの成功パス出力は変更後も期待通り(notice文字列に劣化しない)",
              "iter123 success path title" in _rdx_ok_result2)

# ---------- apply_high_vram_profile(): FUGU_HIGH_VRAM=1 プロファイル切替 (iter124) ----------
# 7つのモジュールグローバルを一括変更する関数で、setup()冒頭(L3540-3541)からのみ呼ばれる。
# 巻き戻し機構が関数自身には無いため、テスト側でswap-and-restoreする(iter43/44/82/83と同じ流儀)。
_hvp_saved_model_config = {m: dict(c) for m, c in f.MODEL_CONFIG.items()}
_hvp_saved = dict(
    PARALLEL_PROPOSERS=f.PARALLEL_PROPOSERS,
    MODEL_NUM_CTX=f.MODEL_NUM_CTX,
    MODEL_KEEP_ALIVE=f.MODEL_KEEP_ALIVE,
    SC_INITIAL=f.SC_INITIAL,
    SC_STEP=f.SC_STEP,
    SC_MAX=f.SC_MAX,
    SC_CHEAP_VOTES=f.SC_CHEAP_VOTES,
    ARBITER_MODEL=f.ARBITER_MODEL,
)
try:
    _hvp_cap = io.StringIO()
    with contextlib.redirect_stdout(_hvp_cap):
        f.apply_high_vram_profile()

    check("apply_high_vram_profile: 全MODEL_CONFIGエントリのnum_ctxが65536になる",
          all(cfg.get("num_ctx") == 65536 for cfg in f.MODEL_CONFIG.values()))
    check("apply_high_vram_profile: 全MODEL_CONFIGエントリのnum_predictが32768になる",
          all(cfg.get("num_predict") == 32768 for cfg in f.MODEL_CONFIG.values()))
    check("apply_high_vram_profile: PARALLEL_PROPOSERSがTrueになる(高VRAM前提でのみ有効、既定の逐次動作(gotcha#5)は変えない)",
          f.PARALLEL_PROPOSERS is True)
    check("apply_high_vram_profile: MODEL_NUM_CTXが32768になる",
          f.MODEL_NUM_CTX == 32768)
    check("apply_high_vram_profile: MODEL_KEEP_ALIVEが'30m'になる",
          f.MODEL_KEEP_ALIVE == "30m")
    check("apply_high_vram_profile: ARBITER_MODELが'gpt-oss:120b'になる",
          f.ARBITER_MODEL == "gpt-oss:120b")
    check("apply_high_vram_profile: SC_INITIAL/SC_STEP/SC_MAXが12/8/48になる(自己一貫性サンプル数の主レバー)",
          (f.SC_INITIAL, f.SC_STEP, f.SC_MAX) == (12, 8, 48))
    check("apply_high_vram_profile: SC_CHEAP_VOTESが16になる",
          f.SC_CHEAP_VOTES == 16)
    check("apply_high_vram_profile: 表示サマリの文言が実際に適用した値と一致する(num_ctx/SC/cheap_votes/arbiter)",
          "num_ctx=65536" in _hvp_cap.getvalue()
          and "SC(init=12,max=48)" in _hvp_cap.getvalue()
          and "cheap_votes=16" in _hvp_cap.getvalue()
          and "arbiter=gpt-oss:120b" in _hvp_cap.getvalue())
finally:
    f.MODEL_CONFIG = _hvp_saved_model_config
    f.PARALLEL_PROPOSERS = _hvp_saved["PARALLEL_PROPOSERS"]
    f.MODEL_NUM_CTX = _hvp_saved["MODEL_NUM_CTX"]
    f.MODEL_KEEP_ALIVE = _hvp_saved["MODEL_KEEP_ALIVE"]
    f.SC_INITIAL = _hvp_saved["SC_INITIAL"]
    f.SC_STEP = _hvp_saved["SC_STEP"]
    f.SC_MAX = _hvp_saved["SC_MAX"]
    f.SC_CHEAP_VOTES = _hvp_saved["SC_CHEAP_VOTES"]
    f.ARBITER_MODEL = _hvp_saved["ARBITER_MODEL"]

check("apply_high_vram_profile: テスト後にMODEL_CONFIGが既定プロファイルの値へ復元されている(後続テストへの汚染防止)",
      f.MODEL_CONFIG == _hvp_saved_model_config)
check("apply_high_vram_profile: テスト後にPARALLEL_PROPOSERS/MODEL_NUM_CTX等7グローバルが既定値へ復元されている",
      f.PARALLEL_PROPOSERS == _hvp_saved["PARALLEL_PROPOSERS"]
      and f.MODEL_NUM_CTX == _hvp_saved["MODEL_NUM_CTX"]
      and f.MODEL_KEEP_ALIVE == _hvp_saved["MODEL_KEEP_ALIVE"]
      and f.SC_INITIAL == _hvp_saved["SC_INITIAL"]
      and f.SC_STEP == _hvp_saved["SC_STEP"]
      and f.SC_MAX == _hvp_saved["SC_MAX"]
      and f.SC_CHEAP_VOTES == _hvp_saved["SC_CHEAP_VOTES"]
      and f.ARBITER_MODEL == _hvp_saved["ARBITER_MODEL"])

# ---------- generate_image_comfyui: 壊れたSaveImage出力エントリのskip-and-recover (2026-07-25 / iter139) ----------
# generate_image_comfyui (fugu_local.py ~L1983-1997) のoutputs/imageエントリ走査は、
# 隣接するsubfolder/typeがimg.get()で守られている一方、filenameだけがimg["filename"]の
# 直接dictアクセスとして2箇所に残っていた。ComfyUI /history の壊れたエントリ
# （filenameキー欠落・null・空文字・非str、またはimages配列内の非dict要素）は
# KeyError/TypeErrorを送出し、本関数を丸ごと巻き込んで呼び出し元generate_imageの外側
# except Exceptionまで伝播、後続のノード/エントリに有効な画像が残っていても生成結果
# ごと握り潰してNoneを返していた。iter77（良い方を回収する）・iter103/111/112
# （非list/非dictの強制truthy変換に頼らない既定値フォールバック）・iter113/iter138
# （外部由来ペイロードの1件の破損で全体を道連れにしないentry単位skip）と同じ作法の
# 回帰防止テスト。generate_image_comfyui/generate_image_a1111/generate_image/
# _detect_backend/_backend_up/_http_post_jsonはこれまで直接のオフラインカバレッジが
# 0件（既存テストは全て丸ごとモック）だったため、ここで初めてComfyUIフォールバック
# 経路そのものに直接テストを追加し、カバレッジの穴も塞ぐ。
# urllib.request.urlopen(/history, /view のGET)と f._http_post_json(/prompt のPOST)の
# みをモックし、実Ollama/ComfyUI/A1111/GPU/ネットワーク呼び出しは一切発生させない。
# f.COMFYUI_CKPTを非空にしてobject_infoチェックポイント自動取得を、f.IMAGE_OUT_DIRを
# tempfile.TemporaryDirectoryにして実ファイルシステム/ホームディレクトリへの書き込みを、
# それぞれ回避する。ワークフロー(wf)辞書・チェックポイント自動検出・/prompt投入・
# /history ポーリングループ・IMAGE_TIMEOUT/deadlineロジック・generate_image_a1111/
# generate_image/_detect_backend/_backend_up/_http_post_json自体には一切触れない
# （このテストが検証するのはoutputs/imageエントリ走査ループのみ）。


class _CjFakeResponse:
    """urllib.request.urlopen が返す `with ... as r:` 用の最小モック(生バイト列を保持)。"""

    def __init__(self, body_bytes):
        self._body = body_bytes

    def __enter__(self):
        return self

    def __exit__(self, *exc_info):
        return False

    def read(self):
        return self._body


def _cj_no_network_urlopen(*a, **k):
    raise AssertionError("generate_image_comfyui: モック漏れで実urlopen(ネットワーク)が呼ばれた"
                          "(object_info取得はCOMFYUI_CKPTプリセットで回避されるはず)")


def _make_comfy_urlopen(history_payload, view_bytes_by_filename, view_calls_log):
    """/history/<pid> と /view?... のGETのみをモックする最小フェイクurlopen。
    generate_image_comfyui はどちらも `urlopen(url_str, timeout=...)` という
    プレーンなURL文字列呼び出しなので、Requestオブジェクトのdata属性は見ない。
    想定外URL(例: object_info。COMFYUI_CKPTプリセットで回避されるはず)は
    AssertionErrorで即座に可視化する。"""

    def _fake(url, timeout=None):
        if "/history/" in url:
            return _CjFakeResponse(json.dumps(history_payload).encode("utf-8"))
        if "/view?" in url:
            qs = urllib.parse.parse_qs(urllib.parse.urlparse(url).query)
            fn = (qs.get("filename") or [None])[0]
            view_calls_log.append(fn)
            if fn in view_bytes_by_filename:
                return _CjFakeResponse(view_bytes_by_filename[fn])
            raise AssertionError(f"generate_image_comfyui test: 未定義の/view filename={fn!r}")
        raise AssertionError(f"generate_image_comfyui test: 想定外のurlopen呼び出し url={url!r}")

    return _fake


_orig_cj_urlopen = urllib.request.urlopen
_orig_cj_post_json = f._http_post_json
_orig_cj_ckpt = f.COMFYUI_CKPT
_orig_cj_out_dir = f.IMAGE_OUT_DIR

# --- (1) 最初のimageエントリにfilenameが無い(欠落)が、同ノード内の後続エントリに
#     有効なfilenameがある -> Noneで諦めず後続の有効画像を回収してPathを返す ---
_cj_td1 = None
try:
    import tempfile as _cj_tempfile
    from pathlib import Path as _cj_Path

    _cj_td1 = _cj_tempfile.TemporaryDirectory()
    f.IMAGE_OUT_DIR = _cj_Path(_cj_td1.name)
    f.COMFYUI_CKPT = "dummy-ckpt-preset"  # 非空 -> object_info自動取得は発火しない
    f._http_post_json = lambda url, payload, timeout: {"prompt_id": "pid-1"}
    _cj_view_calls1 = []
    urllib.request.urlopen = _make_comfy_urlopen(
        {"pid-1": {"outputs": {
            "9": {"images": [
                {"subfolder": "", "type": "output"},         # filenameキー欠落
                {"filename": "good_first.png", "subfolder": "", "type": "output"},
            ]}
        }}},
        {"good_first.png": b"PNGDATA-first"},
        _cj_view_calls1,
    )
    _cj_r1 = f.generate_image_comfyui("a cat", "blurry")
    check("generate_image_comfyui: 先頭エントリのfilename欠落を飛ばして後続の有効画像を回収",
          _cj_r1 is not None)
    check("generate_image_comfyui: 回収したPathのファイル名に有効エントリのfilenameが埋め込まれる",
          _cj_r1 is not None and _cj_r1.name.endswith("_good_first.png"))
    check("generate_image_comfyui: 回収したPathへ/viewの応答バイト列が書き込まれている",
          _cj_r1 is not None and _cj_r1.read_bytes() == b"PNGDATA-first")
    check("generate_image_comfyui: 壊れた先頭エントリでは/viewを呼ばない(2件目のみ1回)",
          _cj_view_calls1 == ["good_first.png"])
finally:
    urllib.request.urlopen = _orig_cj_urlopen
    f._http_post_json = _orig_cj_post_json
    f.COMFYUI_CKPT = _orig_cj_ckpt
    f.IMAGE_OUT_DIR = _orig_cj_out_dir
    if _cj_td1 is not None:
        _cj_td1.cleanup()

# --- (1b) filenameがNoneの壊れたエントリを含むノードの後、別ノードに有効なfilenameが
#     ある -> ノードをまたいでも後続の有効画像を回収する ---
_cj_td1b = None
try:
    _cj_td1b = _cj_tempfile.TemporaryDirectory()
    f.IMAGE_OUT_DIR = _cj_Path(_cj_td1b.name)
    f.COMFYUI_CKPT = "dummy-ckpt-preset"
    f._http_post_json = lambda url, payload, timeout: {"prompt_id": "pid-1b"}
    _cj_view_calls1b = []
    urllib.request.urlopen = _make_comfy_urlopen(
        {"pid-1b": {"outputs": {
            "8": {"images": [{"filename": None, "subfolder": "", "type": "output"}]},
            "9": {"images": [{"filename": "good_other_node.png", "type": "output"}]},
        }}},
        {"good_other_node.png": b"PNGDATA-other-node"},
        _cj_view_calls1b,
    )
    _cj_r1b = f.generate_image_comfyui("a dog", "")
    check("generate_image_comfyui: filename=Noneの壊れたエントリを飛ばし、別ノードの有効画像を回収",
          _cj_r1b is not None and _cj_r1b.name.endswith("_good_other_node.png"))
finally:
    urllib.request.urlopen = _orig_cj_urlopen
    f._http_post_json = _orig_cj_post_json
    f.COMFYUI_CKPT = _orig_cj_ckpt
    f.IMAGE_OUT_DIR = _orig_cj_out_dir
    if _cj_td1b is not None:
        _cj_td1b.cleanup()

# --- (2) imageエントリ自体が非dict(文字列/数値) -> 例外を出さずskipし、有効な
#     エントリが1件も残らなければNoneを返す ---
_cj_td2 = None
try:
    _cj_td2 = _cj_tempfile.TemporaryDirectory()
    f.IMAGE_OUT_DIR = _cj_Path(_cj_td2.name)
    f.COMFYUI_CKPT = "dummy-ckpt-preset"
    f._http_post_json = lambda url, payload, timeout: {"prompt_id": "pid-2"}
    _cj_view_calls2 = []
    urllib.request.urlopen = _make_comfy_urlopen(
        {"pid-2": {"outputs": {
            "9": {"images": ["not-a-dict-entry", 12345, {"filename": None}, {"subfolder": "x"}]}
        }}},
        {},
        _cj_view_calls2,
    )
    _cj_exc2 = None
    try:
        _cj_r2 = f.generate_image_comfyui("prompt", "")
    except Exception as _e:
        _cj_exc2 = _e
        _cj_r2 = "__RAISED__"
    check("generate_image_comfyui: 非dictのimageエントリ(文字列/数値)で例外を送出しない",
          _cj_exc2 is None)
    check("generate_image_comfyui: 有効なfilenameが1件も残らなければNoneを返す",
          _cj_r2 is None)
    check("generate_image_comfyui: 有効エントリが無いので/viewは一度も呼ばれない",
          _cj_view_calls2 == [])
finally:
    urllib.request.urlopen = _orig_cj_urlopen
    f._http_post_json = _orig_cj_post_json
    f.COMFYUI_CKPT = _orig_cj_ckpt
    f.IMAGE_OUT_DIR = _orig_cj_out_dir
    if _cj_td2 is not None:
        _cj_td2.cleanup()

# --- (3) 回帰: 単一の正常なimageエントリでは従来通りPathを返し、/viewは1回だけ
#     フェッチされる(ハッピーパスはバイト単位で不変) ---
_cj_td3 = None
try:
    _cj_td3 = _cj_tempfile.TemporaryDirectory()
    f.IMAGE_OUT_DIR = _cj_Path(_cj_td3.name)
    f.COMFYUI_CKPT = "dummy-ckpt-preset"
    f._http_post_json = lambda url, payload, timeout: {"prompt_id": "pid-3"}
    _cj_view_calls3 = []
    urllib.request.urlopen = _make_comfy_urlopen(
        {"pid-3": {"outputs": {
            "9": {"images": [{"filename": "solo.png", "subfolder": "", "type": "output"}]}
        }}},
        {"solo.png": b"PNGDATA-solo"},
        _cj_view_calls3,
    )
    _cj_r3 = f.generate_image_comfyui("solo prompt", "neg")
    check("generate_image_comfyui: 単一正常エントリの回帰確認(Pathを返す)",
          _cj_r3 is not None and isinstance(_cj_r3, _cj_Path))
    check("generate_image_comfyui: 単一正常エントリのファイル名にfilenameが埋め込まれる(回帰)",
          _cj_r3 is not None and _cj_r3.name.endswith("_solo.png"))
    check("generate_image_comfyui: 単一正常エントリで/viewはちょうど1回だけ呼ばれる(回帰)",
          _cj_view_calls3 == ["solo.png"])
    check("generate_image_comfyui: 単一正常エントリの保存先ディレクトリがIMAGE_OUT_DIR配下(回帰)",
          _cj_r3 is not None and _cj_r3.parent == f.IMAGE_OUT_DIR)
finally:
    urllib.request.urlopen = _orig_cj_urlopen
    f._http_post_json = _orig_cj_post_json
    f.COMFYUI_CKPT = _orig_cj_ckpt
    f.IMAGE_OUT_DIR = _orig_cj_out_dir
    if _cj_td3 is not None:
        _cj_td3.cleanup()

# --- (4) 回帰: 全imageエントリが壊れている/使えるfilenameが無い -> 例外を伝播させずNoneを返す ---
_cj_td4 = None
try:
    _cj_td4 = _cj_tempfile.TemporaryDirectory()
    f.IMAGE_OUT_DIR = _cj_Path(_cj_td4.name)
    f.COMFYUI_CKPT = "dummy-ckpt-preset"
    f._http_post_json = lambda url, payload, timeout: {"prompt_id": "pid-4"}
    _cj_view_calls4 = []
    urllib.request.urlopen = _make_comfy_urlopen(
        {"pid-4": {"outputs": {
            "8": {"images": [{}, {"filename": ""}, {"filename": 123}]},
            "9": {"images": [{"filename": None}, "bogus", 42, None]},
        }}},
        {},
        _cj_view_calls4,
    )
    _cj_exc4 = None
    try:
        _cj_r4 = f.generate_image_comfyui("prompt", "")
    except Exception as _e:
        _cj_exc4 = _e
        _cj_r4 = "__RAISED__"
    check("generate_image_comfyui: 全エントリ壊れていても例外を送出しない(回帰)",
          _cj_exc4 is None)
    check("generate_image_comfyui: 全エントリ壊れていればNoneを返す(回帰)",
          _cj_r4 is None)
finally:
    urllib.request.urlopen = _orig_cj_urlopen
    f._http_post_json = _orig_cj_post_json
    f.COMFYUI_CKPT = _orig_cj_ckpt
    f.IMAGE_OUT_DIR = _orig_cj_out_dir
    if _cj_td4 is not None:
        _cj_td4.cleanup()

check("generate_image_comfyui: テスト後にurllib.request.urlopenが元の状態へ復元されている",
      urllib.request.urlopen == _orig_cj_urlopen)
check("generate_image_comfyui: テスト後にf._http_post_jsonが元の状態へ復元されている",
      f._http_post_json == _orig_cj_post_json)
check("generate_image_comfyui: テスト後にf.COMFYUI_CKPTが元の状態へ復元されている",
      f.COMFYUI_CKPT == _orig_cj_ckpt)
check("generate_image_comfyui: テスト後にf.IMAGE_OUT_DIRが元の状態へ復元されている",
      f.IMAGE_OUT_DIR == _orig_cj_out_dir)

# ---------- generate_image_a1111: 壊れたtxt2imgレスポンスのskip-and-recover (2026-07-25 / iter144) ----------
# generate_image_a1111 (fugu_local.py ~L1941-1976) の txt2img JSONレスポンス処理は、
# iter139でComfyUI側(generate_image_comfyui)のoutputs/imageエントリ走査に施した
# skip-and-recoverと対称の穴がこちら(A1111側)にも残っていた。data.get("images")は
# dataがdictである前提、images[0]はimagesがnon-emptyなlistである前提、
# images[0].split(",", 1)は先頭要素がstrである前提で、いずれも無保証だった。
# 非dictなdata・truthyだが非listなimages・非str/空/デコード不能な先頭エントリは
# AttributeError/TypeError/binascii.Errorを送出し、呼び出し元generate_imageの外側
# except Exceptionまで伝播、後続に有効な画像が残っていても生成結果ごと握り潰して
# Noneになっていた。iter103/111/112/113/138（非list/非dictの強制truthy変換に
# 頼らない既定値フォールバックとentry単位skip）・iter139(ComfyUI側の同型修正)と
# 同じ作法の回帰防止テスト。generate_image_a1111はこれまで直接のオフラインカバ
# レッジが0件（iter139のテストノート通り、iter139はComfyUI側のみに意図的に限定
# していた）だったため、ここで初めて直接テストを追加しカバレッジの穴を塞ぐ。
# f._http_post_jsonのみをモックし、実A1111/ネットワーク呼び出しは一切発生させ
# ない。f.IMAGE_OUT_DIRをtempfile.TemporaryDirectoryにして実ファイルシステムへの
# 書き込みを回避する。_http_post_json自体・generate_image/generate_image_comfyui/
# _detect_backend/_backend_upには一切触れない。

import base64 as _a1_base64
import tempfile as _a1_tempfile
from pathlib import Path as _a1_Path

_orig_a1_post_json = f._http_post_json
_orig_a1_out_dir = f.IMAGE_OUT_DIR

# --- (1) dataが非dict(list/文字列/数値/None) -> 例外を出さずNoneを返す ---
_a1_td1 = None
try:
    _a1_td1 = _a1_tempfile.TemporaryDirectory()
    f.IMAGE_OUT_DIR = _a1_Path(_a1_td1.name)
    for _bad_data in ([1, 2, 3], "not-a-dict", 12345, None):
        f._http_post_json = lambda url, payload, timeout, _v=_bad_data: _v
        _a1_exc1 = None
        try:
            _a1_r1 = f.generate_image_a1111("a cat", "blurry")
        except Exception as _e:
            _a1_exc1 = _e
            _a1_r1 = "__RAISED__"
        check(f"generate_image_a1111: dataが非dict({type(_bad_data).__name__})で例外を送出せずNone",
              _a1_exc1 is None and _a1_r1 is None)
finally:
    f._http_post_json = _orig_a1_post_json
    f.IMAGE_OUT_DIR = _orig_a1_out_dir
    if _a1_td1 is not None:
        _a1_td1.cleanup()

# --- (2) imagesがtruthyだが非list(dict/文字列/数値) -> 例外を出さずNoneを返す ---
_a1_td2 = None
try:
    _a1_td2 = _a1_tempfile.TemporaryDirectory()
    f.IMAGE_OUT_DIR = _a1_Path(_a1_td2.name)
    for _bad_images in ({"0": "x"}, "just-a-string", 999):
        f._http_post_json = lambda url, payload, timeout, _v=_bad_images: {"images": _v}
        _a1_exc2 = None
        try:
            _a1_r2 = f.generate_image_a1111("prompt", "")
        except Exception as _e:
            _a1_exc2 = _e
            _a1_r2 = "__RAISED__"
        check(f"generate_image_a1111: imagesがtruthyな非list({type(_bad_images).__name__})で例外を送出せずNone",
              _a1_exc2 is None and _a1_r2 is None)
finally:
    f._http_post_json = _orig_a1_post_json
    f.IMAGE_OUT_DIR = _orig_a1_out_dir
    if _a1_td2 is not None:
        _a1_td2.cleanup()

# --- (3) 先頭エントリが壊れている(非str/空/デコード不能)が2件目が有効なbase64 ->
#     2件目を回収してPathを返し、書き込まれたバイト列も一致する ---
_a1_valid_b64 = _a1_base64.b64encode(b"PNGDATA-a1111-valid").decode("ascii")
_a1_bad_firsts = [None, 123, {"x": 1}, [], "", "abc"]
_a1_td3 = None
try:
    _a1_td3 = _a1_tempfile.TemporaryDirectory()
    f.IMAGE_OUT_DIR = _a1_Path(_a1_td3.name)
    for _bad_first in _a1_bad_firsts:
        f._http_post_json = lambda url, payload, timeout, _v=_bad_first: {
            "images": [_v, _a1_valid_b64]
        }
        _a1_r3 = f.generate_image_a1111("prompt", "")
        check(f"generate_image_a1111: 壊れた先頭エントリ({_bad_first!r})を飛ばし2件目を回収(Pathを返す)",
              _a1_r3 is not None and isinstance(_a1_r3, _a1_Path))
        check(f"generate_image_a1111: 回収したPathのバイト列が2件目のデコード結果と一致({_bad_first!r})",
              _a1_r3 is not None and _a1_r3.read_bytes() == b"PNGDATA-a1111-valid")
    # 先頭が壊れていて、2件目がdata URIプレフィックス付きでも回収できる
    f._http_post_json = lambda url, payload, timeout: {
        "images": [None, f"data:image/png;base64,{_a1_valid_b64}"]
    }
    _a1_r3b = f.generate_image_a1111("prompt", "")
    check("generate_image_a1111: 壊れた先頭エントリを飛ばしdata URIプレフィックス付き2件目も回収できる",
          _a1_r3b is not None and _a1_r3b.read_bytes() == b"PNGDATA-a1111-valid")
finally:
    f._http_post_json = _orig_a1_post_json
    f.IMAGE_OUT_DIR = _orig_a1_out_dir
    if _a1_td3 is not None:
        _a1_td3.cleanup()

# --- (4) 回帰: 全エントリが壊れている(非str/空/デコード不能) -> 例外を送出せずNoneを
#     返し、IMAGE_OUT_DIR配下に空の書き込みも発生しない ---
_a1_td4 = None
try:
    _a1_td4 = _a1_tempfile.TemporaryDirectory()
    f.IMAGE_OUT_DIR = _a1_Path(_a1_td4.name)
    f._http_post_json = lambda url, payload, timeout: {
        "images": [None, 123, {}, [], "", "abc"]
    }
    _a1_exc4 = None
    try:
        _a1_r4 = f.generate_image_a1111("prompt", "")
    except Exception as _e:
        _a1_exc4 = _e
        _a1_r4 = "__RAISED__"
    check("generate_image_a1111: 全エントリ壊れていても例外を送出しない(回帰)",
          _a1_exc4 is None)
    check("generate_image_a1111: 全エントリ壊れていればNoneを返す(回帰)",
          _a1_r4 is None)
    check("generate_image_a1111: 全エントリ壊れていてもIMAGE_OUT_DIR配下にファイルが書き込まれない",
          list(_a1_Path(_a1_td4.name).iterdir()) == [])
finally:
    f._http_post_json = _orig_a1_post_json
    f.IMAGE_OUT_DIR = _orig_a1_out_dir
    if _a1_td4 is not None:
        _a1_td4.cleanup()

# --- (5) 回帰: 単一の正常なbase64エントリ(data URIプレフィックスあり/なし)では従来通り
#     Pathを返し、書き込まれたバイト列もデコード結果とバイト単位で一致する ---
_a1_td5 = None
try:
    _a1_td5 = _a1_tempfile.TemporaryDirectory()
    f.IMAGE_OUT_DIR = _a1_Path(_a1_td5.name)
    _a1_solo_b64 = _a1_base64.b64encode(b"PNGDATA-a1111-solo").decode("ascii")

    f._http_post_json = lambda url, payload, timeout: {"images": [_a1_solo_b64]}
    _a1_r5a = f.generate_image_a1111("solo prompt", "neg")
    check("generate_image_a1111: 単一正常エントリ(プレフィックスなし)の回帰確認(Pathを返す)",
          _a1_r5a is not None and isinstance(_a1_r5a, _a1_Path))
    check("generate_image_a1111: 単一正常エントリ(プレフィックスなし)のバイト列一致(回帰)",
          _a1_r5a is not None and _a1_r5a.read_bytes() == b"PNGDATA-a1111-solo")
    check("generate_image_a1111: 単一正常エントリの保存先ディレクトリがIMAGE_OUT_DIR配下(回帰)",
          _a1_r5a is not None and _a1_r5a.parent == f.IMAGE_OUT_DIR)

    f._http_post_json = lambda url, payload, timeout: {
        "images": [f"data:image/png;base64,{_a1_solo_b64}"]
    }
    _a1_r5b = f.generate_image_a1111("solo prompt", "neg")
    check("generate_image_a1111: 単一正常エントリ(data URIプレフィックスあり)の回帰確認(Pathを返す)",
          _a1_r5b is not None and isinstance(_a1_r5b, _a1_Path))
    check("generate_image_a1111: 単一正常エントリ(data URIプレフィックスあり)のバイト列一致(回帰)",
          _a1_r5b is not None and _a1_r5b.read_bytes() == b"PNGDATA-a1111-solo")
finally:
    f._http_post_json = _orig_a1_post_json
    f.IMAGE_OUT_DIR = _orig_a1_out_dir
    if _a1_td5 is not None:
        _a1_td5.cleanup()

# --- (6) 回帰: imagesが空リスト/キー自体が欠落 -> Noneを返す ---
_a1_td6 = None
try:
    _a1_td6 = _a1_tempfile.TemporaryDirectory()
    f.IMAGE_OUT_DIR = _a1_Path(_a1_td6.name)
    f._http_post_json = lambda url, payload, timeout: {"images": []}
    check("generate_image_a1111: imagesが空リストならNoneを返す(回帰)",
          f.generate_image_a1111("prompt", "") is None)
    f._http_post_json = lambda url, payload, timeout: {}
    check("generate_image_a1111: imagesキー自体が欠落していればNoneを返す(回帰)",
          f.generate_image_a1111("prompt", "") is None)
finally:
    f._http_post_json = _orig_a1_post_json
    f.IMAGE_OUT_DIR = _orig_a1_out_dir
    if _a1_td6 is not None:
        _a1_td6.cleanup()

check("generate_image_a1111: テスト後にf._http_post_jsonが元の状態へ復元されている",
      f._http_post_json == _orig_a1_post_json)
check("generate_image_a1111: テスト後にf.IMAGE_OUT_DIRが元の状態へ復元されている",
      f.IMAGE_OUT_DIR == _orig_a1_out_dir)

# ---------- generate_image_a1111 / generate_image_comfyui: 同一秒ファイル名衝突による
# 上書き・アーティファクト消失の防止 (2026-07-25) ----------
# 両保存関数とも保存パスをtime.strftime('%Y%m%d_%H%M%S')という秒精度のタイムスタンプ
# のみから組み立てていた。generate_image_a1111 (fugu_local.py ~L1971) はfugu_{ts}.png
# に他の一意化要素が無く、generate_image_comfyui (~L2060) もfugu_{ts}_{filename}で
# ComfyUI側filenameが同一秒内に重複しうる(同一seed/prefixの再実行やサーバ再起動後の
# 内部カウンタリセットなど)。build_pptxはPPTX_MAX_IMAGES=4枚まで generate_image() を
# 連続呼び出しし、返ったPathをimgs[idx]へ格納してスライドへ埋め込む。同一秒内に2回
# 生成が完了すると、旧実装では2回目のout.write_bytes()が1回目のファイルを無言で
# 上書きし、既に生成できていた画像を消したままスライドに同じ画像を重複表示していた。
# これはiter77(良い方を回収する)・iter139/144(壊れたエントリでも後続の有効な
# アーティファクトを取りこぼさない)と同じ「既に生成できたアーティファクトを無言で
# 失わない」系列の欠陥である。現行の8GB GPU機ではSDXL生成が数秒かかるため同一秒内
# 衝突は稀だが、apply_high_vram_profileが想定する96GB高VRAM環境の高速SDXLでは現実的
# になる。修正はuuid4由来の8桁hex(数字/英小文字のみ、パス区切りや".."を含まない
# ファイルシステム安全な文字列)をtime.strftime由来の文字列とは独立に採番し、両保存
# パスへ挟み込むことで秒精度に依存しない一意性を保証する。
# 以下のテストはf.time.strftimeを定数へ固定(=秒精度タイムスタンプ源を凍結)しつつ
# 一意化に使うuuidは実物のまま動かし、同一秒内で2回連続生成しても返り値のPathが
# 別々になり、どちらのファイルも上書きされず両方ディスクに残ることを検証する。
# f._http_post_json / urllib.request.urlopen のみをモックし、実A1111/ComfyUI/GPU/
# ネットワーク呼び出しは一切発生させない。f.IMAGE_OUT_DIRはtempfile.TemporaryDirectory
# に差し替え、テスト後は全てのグローバルをtry/finallyで元に戻す。エントリ走査ロジック
# (iter139/144のskip-and-recoverループ)・isinstance(data, dict)/isinstance(images, list)
# ガード・base64デコード・IMAGE_OUT_DIR.mkdir(parents=True, exist_ok=True)の位置関係・
# 失敗時にNoneを返す契約には一切触れない。

# --- a1111: 同一秒(time.strftime固定)で2回連続生成しても返り値Pathが別々になり、
#     1回目のファイルが2回目の書き込みで上書きされない ---
_a1_td7 = None
try:
    _a1_td7 = _a1_tempfile.TemporaryDirectory()
    f.IMAGE_OUT_DIR = _a1_Path(_a1_td7.name)
    _orig_a1_strftime = f.time.strftime
    f.time.strftime = lambda *a, **k: "20260725_120000"
    try:
        _a1_img1 = _a1_base64.b64encode(b"PNGDATA-a1111-collision-1").decode("ascii")
        _a1_img2 = _a1_base64.b64encode(b"PNGDATA-a1111-collision-2").decode("ascii")

        f._http_post_json = lambda url, payload, timeout: {"images": [_a1_img1]}
        _a1_rc1 = f.generate_image_a1111("prompt A", "")
        f._http_post_json = lambda url, payload, timeout: {"images": [_a1_img2]}
        _a1_rc2 = f.generate_image_a1111("prompt B", "")

        check("generate_image_a1111: 同一秒(time.strftime固定)でも2回の生成が別Pathを返す(衝突回避)",
              _a1_rc1 is not None and _a1_rc2 is not None and _a1_rc1 != _a1_rc2)
        check("generate_image_a1111: 衝突回避後も両方のファイルがディスク上に存在する",
              _a1_rc1 is not None and _a1_rc2 is not None
              and _a1_rc1.exists() and _a1_rc2.exists())
        check("generate_image_a1111: 1件目のファイルが2件目の書き込みで上書きされていない",
              _a1_rc1 is not None and _a1_rc1.read_bytes() == b"PNGDATA-a1111-collision-1")
        check("generate_image_a1111: 2件目のファイルのバイト列も正しい(1件目のデータが漏れていない)",
              _a1_rc2 is not None and _a1_rc2.read_bytes() == b"PNGDATA-a1111-collision-2")
        check("generate_image_a1111: 返り値Pathのファイル名がfugu_で始まり.pngで終わる(両方)",
              _a1_rc1 is not None and _a1_rc1.name.startswith("fugu_")
              and _a1_rc1.name.endswith(".png")
              and _a1_rc2 is not None and _a1_rc2.name.startswith("fugu_")
              and _a1_rc2.name.endswith(".png"))
        check("generate_image_a1111: 返り値PathはIMAGE_OUT_DIR直下(両方)",
              _a1_rc1 is not None and _a1_rc1.parent == f.IMAGE_OUT_DIR
              and _a1_rc2 is not None and _a1_rc2.parent == f.IMAGE_OUT_DIR)
    finally:
        f.time.strftime = _orig_a1_strftime
finally:
    f._http_post_json = _orig_a1_post_json
    f.IMAGE_OUT_DIR = _orig_a1_out_dir
    if _a1_td7 is not None:
        _a1_td7.cleanup()

# --- comfyui: 同一秒(time.strftime固定)かつComfyUI側filenameも同一(same_name.png)の
#     まま2回連続生成しても返り値Pathが別々になり、1回目のファイルが上書きされない ---
_cj_td5 = None
try:
    _cj_td5 = _cj_tempfile.TemporaryDirectory()
    f.IMAGE_OUT_DIR = _cj_Path(_cj_td5.name)
    f.COMFYUI_CKPT = "dummy-ckpt-preset"
    _orig_cj_strftime = f.time.strftime
    f.time.strftime = lambda *a, **k: "20260725_120000"
    try:
        _cj_view_calls5 = []
        f._http_post_json = lambda url, payload, timeout: {"prompt_id": "pid-5a"}
        urllib.request.urlopen = _make_comfy_urlopen(
            {"pid-5a": {"outputs": {
                "9": {"images": [{"filename": "same_name.png", "subfolder": "", "type": "output"}]}
            }}},
            {"same_name.png": b"PNGDATA-comfy-collision-1"},
            _cj_view_calls5,
        )
        _cj_rc1 = f.generate_image_comfyui("prompt A", "")

        f._http_post_json = lambda url, payload, timeout: {"prompt_id": "pid-5b"}
        urllib.request.urlopen = _make_comfy_urlopen(
            {"pid-5b": {"outputs": {
                "9": {"images": [{"filename": "same_name.png", "subfolder": "", "type": "output"}]}
            }}},
            {"same_name.png": b"PNGDATA-comfy-collision-2"},
            _cj_view_calls5,
        )
        _cj_rc2 = f.generate_image_comfyui("prompt B", "")

        check("generate_image_comfyui: 同一秒+同一ComfyUI filenameでも2回の生成が別Pathを返す(衝突回避)",
              _cj_rc1 is not None and _cj_rc2 is not None and _cj_rc1 != _cj_rc2)
        check("generate_image_comfyui: 衝突回避後も両方のファイルがディスク上に存在する",
              _cj_rc1 is not None and _cj_rc2 is not None
              and _cj_rc1.exists() and _cj_rc2.exists())
        check("generate_image_comfyui: 1件目のファイルが2件目の書き込みで上書きされていない",
              _cj_rc1 is not None and _cj_rc1.read_bytes() == b"PNGDATA-comfy-collision-1")
        check("generate_image_comfyui: 2件目のファイルのバイト列も正しい(1件目のデータが漏れていない)",
              _cj_rc2 is not None and _cj_rc2.read_bytes() == b"PNGDATA-comfy-collision-2")
        check("generate_image_comfyui: 返り値PathのファイルにComfyUI側filenameが引き続き埋め込まれる(両方)",
              _cj_rc1 is not None and _cj_rc1.name.endswith("_same_name.png")
              and _cj_rc2 is not None and _cj_rc2.name.endswith("_same_name.png"))
        check("generate_image_comfyui: 返り値PathはIMAGE_OUT_DIR直下(両方)",
              _cj_rc1 is not None and _cj_rc1.parent == f.IMAGE_OUT_DIR
              and _cj_rc2 is not None and _cj_rc2.parent == f.IMAGE_OUT_DIR)
    finally:
        f.time.strftime = _orig_cj_strftime
finally:
    urllib.request.urlopen = _orig_cj_urlopen
    f._http_post_json = _orig_cj_post_json
    f.COMFYUI_CKPT = _orig_cj_ckpt
    f.IMAGE_OUT_DIR = _orig_cj_out_dir
    if _cj_td5 is not None:
        _cj_td5.cleanup()

check("generate_image_a1111/comfyui衝突テスト: テスト後にf.time.strftimeが元の状態へ復元されている",
      f.time.strftime == _orig_a1_strftime == _orig_cj_strftime)
check("generate_image_a1111/comfyui衝突テスト: テスト後にurllib.request.urlopenが元の状態へ復元されている",
      urllib.request.urlopen == _orig_cj_urlopen)
check("generate_image_a1111/comfyui衝突テスト: テスト後にf._http_post_jsonが元の状態へ復元されている",
      f._http_post_json == _orig_a1_post_json)
check("generate_image_a1111/comfyui衝突テスト: テスト後にf.IMAGE_OUT_DIRが元の状態へ復元されている",
      f.IMAGE_OUT_DIR == _orig_a1_out_dir)

# ---------- extract_json dict-or-None契約の4呼び出し元への波及回帰
# (_critic_judge/second_opinion/_sd_prompt_from_request/plan_pptx_images, 2026-07-25) ----------
# 背景: 上の「extract_json: dict-or-None契約の強制」ブロックは extract_json 自体の
# 契約を直接ロックしたが、そもそもこの修正が必要だったのは `extract_json(raw) or {}`
# の直後に無条件で `.get(...)` する4箇所の呼び出し元 ―
# _critic_judge(L2662)/second_opinion(L2710、共に精度優先の critique/verify_single
# ゲート)/_sd_prompt_from_request(L1944)/plan_pptx_images(L4494) ― が、モデルの
# raw 出力が json.loads に成功する非object値(`[1,2,3]`/`42`/`true`/`"x"`)のとき
# 旧実装では truthy な非dict値をそのまま `.get` に渡して AttributeError で
# ターン全体を落としていたため。extract_json 側を直したことで、この4箇所は
# 何も変更していないのに自動的に「extract_json(raw) が None になり `or {}` で
# 空dictへ縮退 → 各関数のドキュメント通りの既定値」という安全な経路に落ちる
# ようになったはずである。ここではその波及を直接確認する(= extract_json だけを
# 直せば4箇所全部が閉じる、というタスクの前提そのものの検証)。
_orig_ej4_ask = f.ask
_orig_ej4_proposers = f.PROPOSERS
_orig_ej4_second_opinion_model = f.SECOND_OPINION_MODEL
_orig_ej4_disabled_flag = f._SECOND_OPINION_DISABLED
_orig_ej4_translate = f.IMAGE_TRANSLATE_PROMPT
try:
    # (A) _critic_judge / second_opinion: 精度優先のcritique/verify_singleゲート。
    #     second_opinion を「有効」経路に固定する(PROPOSERSにSECOND_OPINION_MODELを含める)。
    f.SECOND_OPINION_MODEL = "phi4-mini"
    f.PROPOSERS = ["phi4-mini", "qwen3:4b"]
    f._SECOND_OPINION_DISABLED = False

    for _ej4_raw, _ej4_label in (
        ('[1,2,3]', "リスト"), ('42', "整数"), ('true', "真偽値"), ('"x"', "文字列"),
    ):
        f.ask = lambda *a, _r=_ej4_raw, **k: _r
        _ej4_exc = None
        try:
            _cj_ok, _cj_issue = f._critic_judge("q", "a", think=False)
        except Exception as _exc:
            _ej4_exc = _exc
        check(f"_critic_judge: トップレベル非object JSON({_ej4_label})はAttributeErrorを送出しない",
              _ej4_exc is None)
        check(f"_critic_judge: トップレベル非object JSON({_ej4_label})は既定(True,'')に縮退する",
              _ej4_exc is None and _cj_ok is True and _cj_issue == "")

        f.ask = lambda *a, _r=_ej4_raw, **k: _r
        _ej4_exc = None
        try:
            _so_ok, _so_issue = f.second_opinion("q", "a")
        except Exception as _exc:
            _ej4_exc = _exc
        check(f"second_opinion: トップレベル非object JSON({_ej4_label})はAttributeErrorを送出しない",
              _ej4_exc is None)
        check(f"second_opinion: トップレベル非object JSON({_ej4_label})は既定(True,'')に縮退する",
              _ej4_exc is None and _so_ok is True and _so_issue == "")

    # (B) _sd_prompt_from_request: 同じ非dict JSONは (user_request, '') へ縮退する。
    f.IMAGE_TRANSLATE_PROMPT = True
    for _ej4_raw, _ej4_label in (
        ('[1,2,3]', "リスト"), ('42', "整数"), ('true', "真偽値"), ('"x"', "文字列"),
    ):
        f.ask = lambda *a, _r=_ej4_raw, **k: _r
        _ej4_exc = None
        try:
            _sd_r = f._sd_prompt_from_request("猫の絵")
        except Exception as _exc:
            _ej4_exc = _exc
        check(f"_sd_prompt_from_request: トップレベル非object JSON({_ej4_label})は"
              "AttributeErrorを送出しない",
              _ej4_exc is None)
        check(f"_sd_prompt_from_request: トップレベル非object JSON({_ej4_label})は"
              "(user_request,'')へ縮退する",
              _ej4_exc is None and _sd_r == ("猫の絵", ""))

    # (C) plan_pptx_images: 同じ非dict JSONは {} へ縮退する。
    _ej4_slides = [{"title": "Slide 1", "bullets": ["b1"]}]
    for _ej4_raw, _ej4_label in (
        ('[1,2,3]', "リスト"), ('42', "整数"), ('true', "真偽値"), ('"x"', "文字列"),
    ):
        f.ask = lambda *a, _r=_ej4_raw, **k: _r
        _ej4_exc = None
        try:
            _pi_r = f.plan_pptx_images("Title", _ej4_slides)
        except Exception as _exc:
            _ej4_exc = _exc
        check(f"plan_pptx_images: トップレベル非object JSON({_ej4_label})は"
              "AttributeErrorを送出しない",
              _ej4_exc is None)
        check(f"plan_pptx_images: トップレベル非object JSON({_ej4_label})は{{}}へ縮退する",
              _ej4_exc is None and _pi_r == {})
finally:
    f.ask = _orig_ej4_ask
    f.PROPOSERS = _orig_ej4_proposers
    f.SECOND_OPINION_MODEL = _orig_ej4_second_opinion_model
    f._SECOND_OPINION_DISABLED = _orig_ej4_disabled_flag
    f.IMAGE_TRANSLATE_PROMPT = _orig_ej4_translate

check("extract_json非object回帰: テスト後にaskが元へ復元されている", f.ask == _orig_ej4_ask)
check("extract_json非object回帰: テスト後にPROPOSERSが元へ復元されている",
      f.PROPOSERS == _orig_ej4_proposers)
check("extract_json非object回帰: テスト後にSECOND_OPINION_MODELが元へ復元されている",
      f.SECOND_OPINION_MODEL == _orig_ej4_second_opinion_model)
check("extract_json非object回帰: テスト後にIMAGE_TRANSLATE_PROMPTが元へ復元されている",
      f.IMAGE_TRANSLATE_PROMPT == _orig_ej4_translate)

# ---------- normalize_answer: \left(/\right) 等の体裁マクロ除去 (2026-07-25, iter140の残課題) ----------
# iteration 140（\( \) \[ \] 剥がし）は "\left(/\right)（実括弧を包む別マクロ）...は
# 対象外のまま維持する" と明記して \left/\right をスコープ外に据え置いていた。しかし
# grep でも \left/\right はどこでも処理されておらずテストも皆無で、\left(3, 4\right) の
# ような順序対・区間・集合表記は素の (3,4) 等とは別の投票クラスに分裂したままになる
# (na.lower() 不一致・Fraction 例外 -> gotcha #6 の math_verify フォールバックに依存、
# gotcha #7 の自己整合性投票が最も嫌う票割れ)。iteration 13/22/24/30/78/122/134/136/
# 140/148/160 と同系統の「ファストパスで拾えない票を拾う」姉妹修正としてここで閉じる。
check("normalize_answer: \\left(3, 4\\right) -> (3,4) (順序対、区切り剥がし+iter160カンマ空白除去)",
      f.normalize_answer(r"\left(3, 4\right)") == "(3,4)")
check("normalize_answer: \\left[2, 5\\right) -> [2,5) (半開区間)",
      f.normalize_answer(r"\left[2, 5\right)") == "[2,5)")
check("normalize_answer: \\left\\{1, 2\\right\\} -> {1,2} (中括弧デリミタは保持、"
      "エスケープ綴りはiter166でbareへ正規化)",
      f.normalize_answer(r"\left\{1, 2\right\}") == "{1,2}")
check("normalize_answer: \\left\\{1, 2\\right\\} と素の \\{1,2\\} が同一クラスに合流する",
      f.normalize_answer(r"\left\{1, 2\right\}") == f.normalize_answer(r"\{1,2\}"))

# \b 境界の回帰: \leftarrow/\rightarrow/\leftrightarrow/\Leftrightarrow のような矢印マクロは
# "left"/"right" の直後が英字で単語境界が生じないため（大文字マクロはそもそも小文字専用の
# 正規表現に一致しない）、バイト単位で不変のまま保たれなければならない。
check("normalize_answer: \\leftarrow は不変(\\bガードで巻き込まれない)",
      f.normalize_answer(r"\leftarrow") == r"\leftarrow")
check("normalize_answer: \\rightarrow は不変(\\bガードで巻き込まれない)",
      f.normalize_answer(r"\rightarrow") == r"\rightarrow")
check("normalize_answer: \\leftrightarrow は不変(\\bガードで巻き込まれない)",
      f.normalize_answer(r"\leftrightarrow") == r"\leftrightarrow")
check("normalize_answer: \\Leftrightarrow は不変(大文字マクロは対象外)",
      f.normalize_answer(r"\Leftrightarrow") == r"\Leftrightarrow")
check("normalize_answer: 'x \\to y' 中の裸の 'to' も不変(無関係の混入確認)",
      f.normalize_answer(r"x \to y") == r"x \to y")

# 既存回帰: 素の括弧・角括弧・frac・単純な数値表記はバイト単位で不変のまま
check("normalize_answer: 素の(3,4)は不変(over-stripなし)",
      f.normalize_answer("(3,4)") == "(3,4)")
check("normalize_answer: 素の[2,5)は不変(over-stripなし)",
      f.normalize_answer("[2,5)") == "[2,5)")
check("normalize_answer: 3.14は不変(over-stripなし)",
      f.normalize_answer("3.14") == "3.14")
check("normalize_answer: -5は不変(over-stripなし)",
      f.normalize_answer("-5") == "-5")
check("normalize_answer: 1/2は不変(over-stripなし)",
      f.normalize_answer("1/2") == "1/2")
check("normalize_answer: 1,234 -> 1234(桁区切り回帰、over-stripなし)",
      f.normalize_answer("1,234") == "1234")
check("normalize_answer: 50%は不変(over-stripなし)",
      f.normalize_answer("50%") == "50%")
check("normalize_answer: \\frac{1}{2} -> 1/2 (iter122回帰、over-stripなし)",
      f.normalize_answer(r"\frac{1}{2}") == "1/2")
check("normalize_answer: \\frac{1}{-2} -> -1/2 (iter148回帰、over-stripなし)",
      f.normalize_answer(r"\frac{1}{-2}") == "-1/2")

check("extract_final_answer: \\boxed{\\left(3, 4\\right)} -> (3,4) (boxed math分岐がnormalize_answer経由)",
      f.extract_final_answer(r"\boxed{\left(3, 4\right)}", "math") == "(3,4)")


def _t_left_right_fastpath_equiv(calls):
    result = f.answers_equivalent(r"\left(3,4\right)", "(3,4)")
    check("answers_equivalent: \\left(3,4\\right) と (3,4) が正規化高速パスで一致(math_verify不使用)",
          result is True)
    check("answers_equivalent: \\left/\\right 併合はmath_verify.parseを一切呼ばない(高速パス経由)",
          len(calls["parse_args"]) == 0)


_run_with_math_verify_stub(None, "parse", _t_left_right_fastpath_equiv)

_left_right_votes = [r"\left(3, 4\right)", "(3,4)", "(3,4)"]
_lr_top, _lr_count, _lr_classes = f.vote_answers(_left_right_votes)
check("vote_answers: \\left(3, 4\\right)/(3,4)/(3,4) の3票が単一クラスに集約される(票割れしない)",
      _lr_count == 3 and len(_lr_classes) == 1)

# ---------- normalize_answer: エスケープ済み集合中括弧 \{ \} の綴り正規化 (2026-07-25, iter166) ----------
# iteration 140/164 は \{/\} を「集合の区切り文字として保持する、剥がさない」と明示的に
# スコープ外にしていた（区切り文字を残す判断自体は正しい）。しかしそのエスケープの綴り違い
# （\boxed{\{1,2\}} の '\{1,2\}' と \boxed{{1,2}} の '{1,2}'）は残ったままで、同じ集合値が
# 別々の投票クラスに分裂していた（自己整合性投票 gotcha #7）。ここでは中括弧そのものは
# 残したまま、バックスラッシュのエスケープだけを剥がして素の { } に統一する。
check("normalize_answer: \\{1,2\\} -> {1,2} (中括弧は保持、エスケープ綴りをbareに正規化)",
      f.normalize_answer(r"\{1,2\}") == "{1,2}")
check("normalize_answer: \\{1, 2, 3\\} -> {1,2,3} (iter160の数字間カンマ空白除去も後段で継続動作)",
      f.normalize_answer(r"\{1, 2, 3\}") == "{1,2,3}")
check("normalize_answer: \\{1,2\\} と素の {1,2} がbyte-for-byteで同一クラスに合流する",
      f.normalize_answer(r"\{1,2\}") == f.normalize_answer("{1,2}"))
check("normalize_answer: \\{1, 2, 3\\} と素の {1,2,3} がbyte-for-byteで同一クラスに合流する",
      f.normalize_answer(r"\{1, 2, 3\}") == f.normalize_answer("{1,2,3}"))

# regression: 削除ではなく置換であることの確認 -- 中括弧自体は残り、素のタプル/スカラーには
# 誤って一致しない(制約: REPLACE, do not DELETE)。
check("normalize_answer: \\{1,2\\} の中括弧は削除されず残る(素の1,2への誤併合なし)",
      f.normalize_answer(r"\{1,2\}") == "{1,2}" and f.normalize_answer(r"\{1,2\}") != "1,2")


def _t_escaped_brace_fastpath_equiv(calls):
    result = f.answers_equivalent(r"\{1,2\}", "{1,2}")
    check("answers_equivalent: \\{1,2\\} と {1,2} が正規化高速パスで一致(math_verify不使用)",
          result is True)
    check("answers_equivalent: エスケープ済み中括弧の併合はmath_verify.parseを一切呼ばない(高速パス経由)",
          len(calls["parse_args"]) == 0)


_run_with_math_verify_stub(None, "parse", _t_escaped_brace_fastpath_equiv)

_escaped_brace_votes = [r"\{1,2\}", "{1,2}", "{1,2}"]
_eb_top, _eb_count, _eb_classes = f.vote_answers(_escaped_brace_votes)
check("vote_answers: \\{1,2\\}/{1,2}/{1,2} の3票が単一クラスに集約される(票割れしない)",
      _eb_count == 3 and len(_eb_classes) == 1)

check("extract_final_answer: \\boxed{\\{1,2\\}} -> {1,2} (boxed math分岐、エスケープ綴り経由)",
      f.extract_final_answer(r"\boxed{\{1,2\}}", "math") == "{1,2}")
check("extract_final_answer: \\boxed{{1,2}} -> {1,2} (boxed math分岐、bare綴り経由)",
      f.extract_final_answer(r"\boxed{{1,2}}", "math") == "{1,2}")
check("extract_final_answer: エスケープ/bare両ルートが同一の投票トークンに収束する",
      f.extract_final_answer(r"\boxed{\{1,2\}}", "math") ==
      f.extract_final_answer(r"\boxed{{1,2}}", "math"))

# regression(バイト単位で不変): 値を持たない区切り文字剥がし(iter140の\(\)/\[\])・
# \frac 変換(iter122)・\textbf 外殻剥がし(iter23)・\left/\right剥がし(iter164)・
# 素の値は今回の変更で一切影響を受けない。
check("normalize_answer: \\(5\\) -> 5 (iter140回帰、over-stripなし)",
      f.normalize_answer(r"\(5\)") == "5")
check("normalize_answer: $5$ -> 5 (回帰、over-stripなし)",
      f.normalize_answer("$5$") == "5")
check("normalize_answer: \\frac{1}{2} -> 1/2 (iter122回帰、over-stripなし)",
      f.normalize_answer(r"\frac{1}{2}") == "1/2")
check("normalize_answer: \\textbf{B} -> B (iter23回帰、over-stripなし)",
      f.normalize_answer(r"\textbf{B}") == "B")
check("normalize_answer: \\left(3, 4\\right) -> (3,4) (iter164回帰、over-stripなし)",
      f.normalize_answer(r"\left(3, 4\right)") == "(3,4)")
check("normalize_answer: 素のA/5/1/2/50%/-5は不変(over-stripなし)",
      f.normalize_answer("A") == "A"
      and f.normalize_answer("5") == "5"
      and f.normalize_answer("1/2") == "1/2"
      and f.normalize_answer("50%") == "50%"
      and f.normalize_answer("-5") == "-5")

# ---------- main(): CLI引数解析 / 質問ソース解決 / office_attached ルーティング ----------
# 2026-07-26: main() (fugu_local.py ~L5093) は本スクリプト唯一のエントリポイントだが
# これまで直接のオフラインテストが無かった。--file/--rag/--search によるコンテキスト
# 注入の有無、office_attached (Office文書→Proposer C 主軸ルーティング) の設定、
# --session/--no-history による履歴の読み込み・永続化制御など、精度に直結する分岐を
# 多数抱える「配線」部分であり、抽出/投票/reader/saver/plan/history/search の各経路は
# 既に厳密にテスト済みなのに main() だけが手つかずだった、という純粋なカバレッジの
# 空白を埋める。sys.argv を実際に差し替えて argparse の解析経路ごと駆動する
# (parse_args 自体は本物を使い、再実装はしない)。
#
# モック対象: f.setup (実サーバー疎通を避け True 固定)、f.ask_fugu / f.repl
# (実パイプライン・モデル呼び出しを避ける)、f.load_history_file (実
# ~/.fugu_history.json / HISTORY_FILE には一切触れずフェイク履歴を返す)、
# --file 系のみ f.read_file_text (python-docx/openpyxl/pdfplumber 等の実ライブラリに
# 依存しないよう固定文字列を返す)。SESSION_SAVE / _HISTORY / RAG_DIRS / sys.argv /
# sys.stdin はブロック終了後にまとめて元の状態へ復元し、復元自体もアサートする。
#
# 副次的に判明した特性(表面化のみ・修正はしない。iters 66/71/110のsurface-don't-fix
# 方針に従う):
#  (a) [2026-07-26 iteration 186で対応済み] --out はインタラクティブ分岐(質問なし+
#      isatty、L5171-5172のrepl()呼び出し)には一切転送されない。main()がrepl()へ渡す
#      引数はuse_search/rag_dirs/history_fileの3つだけで、repl()自身のシグネチャにも
#      out_file引数が無いため、`--out result.md` を質問なしで指定してもrepl()自体には
#      依然として転送されない(この設計・repl()の引数は変更しない)。iteration 185時点
#      ではこれがエラーにも警告にもならず完全に黙って無視されていたが、iteration 186で
#      この分岐にのみ「--outは対話モードでは無視される」旨とsave <path>コマンドへの
#      誘導を表示するcp932安全な警告printを追加し、可視化した
#      (surface-don't-swallow方針、gotcha #8の精神)。下のテスト(10)/(10b)を参照。
#  (b) パイプ入力(stdin)分岐(L5177)のask_fugu呼び出しはoffice_attachedを明示的に
#      渡さない(ask_fugu側の既定値Falseに暗黙依存)。--fileを経ない経路なので実害は
#      無いが、--file経路(常にoffice_attachedを明示するL5170)とは非対称なキーワード
#      渡し方になっている。
#  (c) --fileのテキストはサイズ上限が一切無く、抽出結果がどれだけ長くてもnum_ctxを
#      意識したトランケートをせずそのままask_fugu(question=...)へ渡す。
# (b)(c)はmain()自体には手を入れない。(a)はiteration 186で対話分岐に警告printを
# 1文追加したのみで、質問経路・パイプ入力経路・repl()呼び出しの引数は一切変更していない。
import tempfile as _cli_tempfile

_orig_cli_setup = f.setup
_orig_cli_ask_fugu = f.ask_fugu
_orig_cli_repl = f.repl
_orig_cli_load_history = f.load_history_file
_orig_cli_read_file_text = f.read_file_text
_orig_cli_session_save = f.SESSION_SAVE
_orig_cli_history = f._HISTORY
_orig_cli_rag_dirs = f.RAG_DIRS
_orig_cli_argv = sys.argv
_orig_cli_stdin = sys.stdin

# fugu_local.py L5148 の _OFFICE_SUFFIXES はmain()内ローカル変数でモジュールからは
# 参照できないため、同じリテラルをテスト側でも直接保持する(ソース確認済み)。
_CLI_OFFICE_SUFFIXES = (".docx", ".doc", ".xlsx", ".xls", ".pdf", ".pptx", ".ppt")


class _CliFakeStdin:
    """sys.stdin を置き換える最小フェイク。isatty()/read()のみ提供する。"""
    def __init__(self, isatty, text=""):
        self._isatty = isatty
        self._text = text

    def isatty(self):
        return self._isatty

    def read(self):
        return self._text


class _CliRunResult:
    """_cli_run()の戻り値。(ask_calls, repl_calls, load_hist_calls)の3要素タプルとして
    アンパック可能(__iter__)にしつつ、2026-07-26 iteration 186で追加した.stdout属性で
    捕捉した標準出力全文にもアクセスできるようにする(--out無視警告の文言アサート用)。
    既存(iteration 185)の `a, b, c = _cli_run(...)` という呼び出し側コードは変更不要。"""
    def __init__(self, ask_calls, repl_calls, load_hist_calls, stdout_text):
        self.ask_calls = ask_calls
        self.repl_calls = repl_calls
        self.load_hist_calls = load_hist_calls
        self.stdout = stdout_text

    def __iter__(self):
        return iter((self.ask_calls, self.repl_calls, self.load_hist_calls))


def _cli_run(argv_tail, stdin_isatty=True, stdin_text="", history_return=()):
    """main()を1回駆動し _CliRunResult (ask_fugu呼び出し記録, repl呼び出し記録,
    load_history_file呼び出し記録, 捕捉した標準出力全文) を返す。sys.argv/sys.stdinは
    呼び出し内で復元する(main()内の例外有無に関わらず)。標準出力はcontextlib.
    redirect_stdoutで捕捉し、戻り値の.stdout属性から参照できる
    (2026-07-26 iteration 186: --out無視警告の文言をアサートするために追加。
    従来のask_calls/repl_calls/load_hist_callsの3要素タプルとしてのアンパックは
    _CliRunResult.__iter__により従来通り動作する)。"""
    ask_calls = []
    repl_calls = []
    load_hist_calls = []

    def _fake_ask_fugu(question, **kwargs):
        ask_calls.append({"question": question, "kwargs": kwargs})

    def _fake_repl(**kwargs):
        repl_calls.append(kwargs)

    def _fake_load_history(path):
        load_hist_calls.append(path)
        return list(history_return)

    f.setup = lambda: True
    f.ask_fugu = _fake_ask_fugu
    f.repl = _fake_repl
    f.load_history_file = _fake_load_history
    sys.argv = ["fugu_local.py"] + list(argv_tail)
    sys.stdin = _CliFakeStdin(stdin_isatty, stdin_text)
    _stdout_buf = io.StringIO()
    try:
        with contextlib.redirect_stdout(_stdout_buf):
            f.main()
    finally:
        sys.argv = _orig_cli_argv
        sys.stdin = _orig_cli_stdin
    return _CliRunResult(ask_calls, repl_calls, load_hist_calls, _stdout_buf.getvalue())


try:
    f.RAG_DIRS = []

    # --- (1) 位置引数のみの質問: 素の質問文がそのままask_fuguへ渡り、replは呼ばれない ---
    f.SESSION_SAVE = True
    f._HISTORY = []
    _ask1, _repl1, _lh1 = _cli_run(["これはテスト質問です"])
    check("main: 位置引数の質問はask_fuguに1回だけ渡る", len(_ask1) == 1)
    check("main: 位置引数の質問文がそのまま渡る(strip等の加工なし)",
          bool(_ask1) and _ask1[0]["question"] == "これはテスト質問です")
    check("main: use_searchの既定はFalse",
          bool(_ask1) and _ask1[0]["kwargs"].get("use_search") is False)
    check("main: --rag無し/RAG_DIRS空ならrag_dirsはNone",
          bool(_ask1) and _ask1[0]["kwargs"].get("rag_dirs") is None)
    check("main: --out無しならout_fileはNone",
          bool(_ask1) and _ask1[0]["kwargs"].get("out_file") is None)
    check("main: --session無しならhistory_fileはHISTORY_FILE",
          bool(_ask1) and _ask1[0]["kwargs"].get("history_file") == f.HISTORY_FILE)
    check("main: --file無しの経路ではoffice_attachedはFalse",
          bool(_ask1) and _ask1[0]["kwargs"].get("office_attached") is False)
    check("main: 質問がある場合replは呼ばれない", len(_repl1) == 0)
    check("main: --no-history無しならload_history_fileが1回呼ばれる", len(_lh1) == 1)
    check("main: load_history_fileにHISTORY_FILEが渡る",
          bool(_lh1) and _lh1[0] == f.HISTORY_FILE)

    # --- (2) --file <.txt>: read_file_textの戻り値がstripされてquestionになり、
    #         .txtはOffice拡張子ではないのでoffice_attachedはFalse ---
    f.SESSION_SAVE = True
    f._HISTORY = []
    with _cli_tempfile.TemporaryDirectory() as _cli_dir2:
        _cli_txt = f.Path(_cli_dir2) / "task.txt"
        _cli_txt.write_text("placeholder", encoding="utf-8")
        f.read_file_text = lambda path: "  body text  \n"
        _ask2, _repl2, _lh2 = _cli_run(["--file", str(_cli_txt)])
    check("main: --file(.txt)の質問はread_file_textの戻り値をstripしたもの",
          bool(_ask2) and _ask2[0]["question"] == "body text")
    check("main: --file(.txt)はOffice拡張子でないのでoffice_attachedはFalse",
          bool(_ask2) and _ask2[0]["kwargs"].get("office_attached") is False)
    check("main: --fileの質問経路でもreplは呼ばれない", len(_repl2) == 0)

    # --- (3) --file <Office拡張子>: office_attachedがTrueになる(全_CLI_OFFICE_SUFFIXES) ---
    f.SESSION_SAVE = True
    f._HISTORY = []
    _office_results = {}
    for _suf in _CLI_OFFICE_SUFFIXES:
        with _cli_tempfile.TemporaryDirectory() as _cli_dir3:
            _cli_office_fp = f.Path(_cli_dir3) / ("doc" + _suf)
            _cli_office_fp.write_bytes(b"placeholder")
            f.read_file_text = lambda path: "Office文書の本文です"
            _ask3, _repl3, _lh3 = _cli_run(["--file", str(_cli_office_fp)])
        _office_results[_suf] = (_ask3, _repl3)
    check("main: 全Office拡張子でask_fuguが1回呼ばれる",
          all(len(a) == 1 for a, r in _office_results.values()))
    check("main: 全Office拡張子でoffice_attachedがTrueになる",
          all(a[0]["kwargs"].get("office_attached") is True
              for a, r in _office_results.values()))
    check("main: 全Office拡張子でreplは呼ばれない",
          all(len(r) == 0 for a, r in _office_results.values()))

    # --- (4) --file <存在しないパス>: 早期returnし、ask_fugu/replとも呼ばれない ---
    f.SESSION_SAVE = True
    f._HISTORY = []
    with _cli_tempfile.TemporaryDirectory() as _cli_dir4:
        _cli_missing = f.Path(_cli_dir4) / "does_not_exist.txt"
        f.read_file_text = lambda path: "呼ばれないはず"
        _ask4, _repl4, _lh4 = _cli_run(["--file", str(_cli_missing)])
    check("main: --fileが存在しないパスならask_fuguは呼ばれない(早期return)", len(_ask4) == 0)
    check("main: --fileが存在しないパスならreplも呼ばれない", len(_repl4) == 0)

    # --- (5) --file <存在するが抽出結果が空/空白のみ>: 早期returnし、ask_fuguは呼ばれない ---
    f.SESSION_SAVE = True
    f._HISTORY = []
    with _cli_tempfile.TemporaryDirectory() as _cli_dir5:
        _cli_empty = f.Path(_cli_dir5) / "empty.txt"
        _cli_empty.write_text("placeholder", encoding="utf-8")
        f.read_file_text = lambda path: "   \n  "
        _ask5, _repl5, _lh5 = _cli_run(["--file", str(_cli_empty)])
    check("main: --fileの抽出結果が空白のみならask_fuguは呼ばれない(抽出失敗の早期return)",
          len(_ask5) == 0)
    check("main: --fileの抽出結果が空白のみでもreplは呼ばれない", len(_repl5) == 0)

    # --- (c) 副次的に判明した特性: --fileのテキストにサイズ上限が無いことの特性化 ---
    # 2026-07-26 iteration 192で対応: トランケート自体はしない(surface, don't fix。
    # gotcha #7と同じ精度優先の思想)。「何文字あっても丸ごと通す」契約は不変のまま
    # 固定するが、抽出テキストがMODEL_NUM_CTXを超える場合に限り、main()の--file分岐
    # (ask_fugu呼び出しの直前)でnum_ctxオーバーフローの恐れを知らせる警告printを
    # 追加した。以下でその警告の有無と、質問本文・各種kwargsが一切変化しない
    # ことの双方を検証する(詳細な閾値境界・小ファイル・他分岐への無影響は
    # 直後の(iter192-a/b/c)ブロックで追加検証する)。
    f.SESSION_SAVE = True
    f._HISTORY = []
    with _cli_tempfile.TemporaryDirectory() as _cli_dir_big:
        _cli_big_fp = f.Path(_cli_dir_big) / "big.txt"
        _cli_big_fp.write_text("placeholder", encoding="utf-8")
        _cli_big_text = "A" * 50000
        f.read_file_text = lambda path: _cli_big_text
        _res_big = _cli_run(["--file", str(_cli_big_fp)])
    _ask_big, _repl_big, _lh_big = _res_big
    check("main(特性化・c): --fileの抽出テキストはサイズ上限なしで丸ごとquestionになる",
          bool(_ask_big) and len(_ask_big[0]["question"]) == 50000
          and _ask_big[0]["question"] == _cli_big_text)
    check("main(iter192): 50000文字(>MODEL_NUM_CTX)の--fileはoverflow警告を出す",
          "警告" in _res_big.stdout and "big.txt" in _res_big.stdout)
    check("main(iter192): overflow警告はnum_ctxの値に言及する",
          str(f.MODEL_NUM_CTX) in _res_big.stdout)
    check("main(iter192): overflow警告があってもoffice_attachedは変化しない(.txtなのでFalse)",
          bool(_ask_big) and _ask_big[0]["kwargs"].get("office_attached") is False)

    # --- (iter192-a) 閾値の境界確認: MODEL_NUM_CTX+1文字なら警告、ちょうど
    #     MODEL_NUM_CTX文字(境界そのもの)なら警告なし(len(text) > MODEL_NUM_CTXという
    #     厳密な不等号の実装を直接検証する) ---
    f.SESSION_SAVE = True
    f._HISTORY = []
    with _cli_tempfile.TemporaryDirectory() as _cli_dir_over:
        _cli_over_fp = f.Path(_cli_dir_over) / "over.txt"
        _cli_over_fp.write_text("placeholder", encoding="utf-8")
        _cli_over_text = "B" * (f.MODEL_NUM_CTX + 1)
        f.read_file_text = lambda path: _cli_over_text
        _res_over = _cli_run(["--file", str(_cli_over_fp)])
    check("main(iter192): MODEL_NUM_CTX+1文字ならoverflow警告が出る",
          "警告" in _res_over.stdout and "over.txt" in _res_over.stdout)
    check("main(iter192): 境界を超えても抽出テキストは全文そのままquestionに渡る",
          bool(_res_over.ask_calls)
          and _res_over.ask_calls[0]["question"] == _cli_over_text)

    f.SESSION_SAVE = True
    f._HISTORY = []
    with _cli_tempfile.TemporaryDirectory() as _cli_dir_eq:
        _cli_eq_fp = f.Path(_cli_dir_eq) / "eq.txt"
        _cli_eq_fp.write_text("placeholder", encoding="utf-8")
        _cli_eq_text = "C" * f.MODEL_NUM_CTX
        f.read_file_text = lambda path: _cli_eq_text
        _res_eq = _cli_run(["--file", str(_cli_eq_fp)])
    check("main(iter192): ちょうどMODEL_NUM_CTX文字(境界=超えていない)なら警告なし",
          "警告" not in _res_eq.stdout)
    check("main(iter192): 境界ちょうどでも質問本文は従来通り全文渡る",
          bool(_res_eq.ask_calls) and _res_eq.ask_calls[0]["question"] == _cli_eq_text)

    # --- (iter192-b) 小さいファイル: 警告が出ず、挙動はbyte-for-byteで従来通り ---
    f.SESSION_SAVE = True
    f._HISTORY = []
    with _cli_tempfile.TemporaryDirectory() as _cli_dir_small:
        _cli_small_fp = f.Path(_cli_dir_small) / "small.txt"
        _cli_small_fp.write_text("placeholder", encoding="utf-8")
        _cli_small_text = "  小さいファイルの本文です  \n"
        f.read_file_text = lambda path: _cli_small_text
        _res_small = _cli_run(["--file", str(_cli_small_fp)])
    _ask_small, _repl_small, _lh_small = _res_small
    check("main(iter192): 小さい--fileはoverflow警告が出ない", "警告" not in _res_small.stdout)
    check("main(iter192): 小さい--fileの質問は従来通りstripされた本文(挙動不変)",
          bool(_ask_small) and _ask_small[0]["question"] == "小さいファイルの本文です")
    check("main(iter192): 小さい--fileでもreplは呼ばれない", len(_repl_small) == 0)

    # --- (iter192-c) 警告は"print"のみであり、out_file / history_file / office_attached
    #     の転送やルーティング・repl呼び出し有無には一切影響しない(control flow不変)
    #     ことの確認。Office拡張子 + --out + --session を同時指定した大容量ファイルで検証 ---
    f.SESSION_SAVE = True
    f._HISTORY = []
    with _cli_tempfile.TemporaryDirectory() as _cli_dir_office_big:
        _cli_office_big_fp = f.Path(_cli_dir_office_big) / "bigdoc.pdf"
        _cli_office_big_fp.write_bytes(b"placeholder")
        _cli_office_big_text = "D" * 50000
        f.read_file_text = lambda path: _cli_office_big_text
        _res_office_big = _cli_run(["--file", str(_cli_office_big_fp),
                                     "--out", "out192.md", "--session", "sess192.json"])
    _ask_ob, _repl_ob, _lh_ob = _res_office_big
    check("main(iter192-c): overflow警告があってもask_fuguは1回だけ呼ばれる(control flow不変)",
          len(_ask_ob) == 1)
    check("main(iter192-c): overflow警告があってもoffice_attachedはTrueのまま(.pdf)",
          bool(_ask_ob) and _ask_ob[0]["kwargs"].get("office_attached") is True)
    check("main(iter192-c): overflow警告があってもout_fileはそのままask_fuguへ転送される",
          bool(_ask_ob) and _ask_ob[0]["kwargs"].get("out_file") == "out192.md")
    check("main(iter192-c): overflow警告があってもhistory_fileはPath化されて転送される",
          bool(_ask_ob) and _ask_ob[0]["kwargs"].get("history_file") == f.Path("sess192.json"))
    check("main(iter192-c): overflow警告があってもreplは呼ばれない", len(_repl_ob) == 0)
    check("main(iter192-c): overflow警告付きでも抽出テキスト全文がそのままquestionになる",
          bool(_ask_ob) and _ask_ob[0]["question"] == _cli_office_big_text)

    def _iter192_cp932_ok(s):
        """文字列がcp932(Windowsコンソールの既知の落とし穴#4)へ例外なくエンコード
        できるかを確認するだけのヘルパー。encode自体が失敗したら偽を返す(送出しない)。"""
        try:
            s.encode("cp932")
            return True
        except UnicodeEncodeError:
            return False

    check("main(iter192): overflow警告文を含む標準出力全体がcp932でエンコード可能(gotcha #4)",
          _iter192_cp932_ok(_res_office_big.stdout))

    # --- (6) --out <path> + 位置引数の質問: out_fileがそのままask_fuguへ転送される ---
    f.SESSION_SAVE = True
    f._HISTORY = []
    _res6 = _cli_run(["質問6", "--out", "result.md"])
    _ask6, _repl6, _lh6 = _res6
    check("main: --outの値がask_fuguのout_fileへ転送される",
          bool(_ask6) and _ask6[0]["kwargs"].get("out_file") == "result.md")
    # 2026-07-26 iteration 186 回帰確認: --out+質問ありの経路は今回の変更対象外
    # (対話分岐のみに警告を追加した)ため、replは呼ばれず、対話分岐限定の
    # --out無視警告も出ないはず。
    check("main(iter186回帰): --out+質問ありではreplは呼ばれない", len(_repl6) == 0)
    check("main(iter186回帰): --out+質問ありでは対話分岐の--out無視警告は出ない",
          "[警告]" not in _res6.stdout)

    # --- (7) --search + --rag a b: use_search=True, rag_dirs=['a','b'] ---
    f.SESSION_SAVE = True
    f._HISTORY = []
    _ask7, _repl7, _lh7 = _cli_run(["質問7", "--search", "--rag", "a", "b"])
    check("main: --searchでuse_search=Trueがask_fuguへ渡る",
          bool(_ask7) and _ask7[0]["kwargs"].get("use_search") is True)
    check("main: --rag a b でrag_dirs=['a','b']がask_fuguへ渡る",
          bool(_ask7) and _ask7[0]["kwargs"].get("rag_dirs") == ["a", "b"])

    # --- (8a) --no-history: SESSION_SAVE=Falseになり、load_history_fileは呼ばれず、
    #          _HISTORYも書き換わらない(elseブランチ自体がスキップされるため) ---
    f.SESSION_SAVE = True
    f._HISTORY = ["sentinel_no_history_untouched"]
    _ask8a, _repl8a, _lh8a = _cli_run(["--no-history", "質問8a"])
    check("main: --no-historyでSESSION_SAVEがFalseになる", f.SESSION_SAVE is False)
    check("main: --no-historyならload_history_fileは呼ばれない", len(_lh8a) == 0)
    check("main: --no-historyなら_HISTORYは書き換わらない(elseブランチ丸ごとスキップ)",
          f._HISTORY == ["sentinel_no_history_untouched"])
    check("main: --no-historyでもask_fuguは通常通り呼ばれる", len(_ask8a) == 1)

    # --- (8b) --no-history無し: load_history_fileが呼ばれ、戻り値がそのまま_HISTORYになる ---
    f.SESSION_SAVE = True
    f._HISTORY = []
    _fake_hist_8b = [{"role": "user", "content": "以前の質問"},
                     {"role": "assistant", "content": "以前の回答"}]
    _ask8b, _repl8b, _lh8b = _cli_run(["質問8b"], history_return=_fake_hist_8b)
    check("main: --no-history無しならload_history_fileが呼ばれる", len(_lh8b) == 1)
    check("main: load_history_fileの戻り値がそのまま_HISTORYになる",
          f._HISTORY == _fake_hist_8b)

    # --- (9) --session <path>: hfile==Path(path)がload_history_file/ask_fugu双方に
    #         forwardされる(対話分岐ではreplにも渡る) ---
    f.SESSION_SAVE = True
    f._HISTORY = []
    _ask9, _repl9, _lh9 = _cli_run(["質問9", "--session", "custom_session.json"])
    check("main: --sessionの値がPath化されてload_history_fileへ渡る",
          bool(_lh9) and _lh9[0] == f.Path("custom_session.json"))
    check("main: --sessionの値がPath化されてask_fuguのhistory_fileへ渡る",
          bool(_ask9) and _ask9[0]["kwargs"].get("history_file") == f.Path("custom_session.json"))

    f.SESSION_SAVE = True
    f._HISTORY = []
    _ask9i, _repl9i, _lh9i = _cli_run(["--session", "custom_session2.json"],
                                       stdin_isatty=True)
    check("main: 質問なし+isatty()でのreplにも--sessionのPathがhistory_fileとして渡る",
          bool(_repl9i) and _repl9i[0].get("history_file") == f.Path("custom_session2.json"))

    # --- (10) 質問なし + stdin.isatty()=True: replが(use_search, rag_dirs, history_file)
    #          の3引数のみで1回呼ばれ、ask_fuguは呼ばれない。
    #          特性(a)の確認: --outを付けてもrepl()呼び出し自体には一切現れない
    #          (repl()のシグネチャ・呼び出し引数は変更しない、という設計判断は不変)。
    #          2026-07-26 iteration 186: この場合のみ「--outは対話モードでは無視される」
    #          旨とsaveコマンドへの誘導を表示する警告が標準出力に出るようになった
    #          (repl()呼び出し自体・ask_fuguが呼ばれないことは従来通り) ---
    f.SESSION_SAVE = True
    f._HISTORY = []
    _res10 = _cli_run(["--out", "ignored_output.md"], stdin_isatty=True)
    _ask10, _repl10, _lh10 = _res10
    check("main: 質問なし+isatty()=Trueならreplが1回だけ呼ばれる", len(_repl10) == 1)
    check("main: 質問なし+isatty()=Trueならask_fuguは呼ばれない", len(_ask10) == 0)
    check("main: repl呼び出しの引数はuse_search/rag_dirs/history_fileの3つだけ",
          bool(_repl10) and set(_repl10[0].keys()) == {"use_search", "rag_dirs", "history_file"})
    check("main(特性・a): --outを指定してもrepl()呼び出しにout_fileは現れない"
          "(repl()自体は変更しない設計)",
          bool(_repl10) and "out_file" not in _repl10[0])
    check("main(iter186): --out指定+対話モードでは--out無視の警告がstdoutに出る",
          "--out" in _res10.stdout and "ignored_output.md" in _res10.stdout)
    check("main(iter186): 警告はsaveコマンドへの誘導を含む", "save" in _res10.stdout)

    # --- (10b) 質問なし + stdin.isatty()=True + --out無し: (10)と同じくreplは3引数
    #           のみで1回呼ばれるが、--outを指定していないので警告は出ない(回帰確認、
    #           iteration 186で追加した警告がargs.outの真偽値で正しくガードされている
    #           ことの確認) ---
    f.SESSION_SAVE = True
    f._HISTORY = []
    _res10b = _cli_run([], stdin_isatty=True)
    _ask10b, _repl10b, _lh10b = _res10b
    check("main(iter186回帰): --out無し+質問なし+isatty()=Trueでもreplが1回だけ呼ばれる",
          len(_repl10b) == 1)
    check("main(iter186回帰): --out無し+質問なし+isatty()=Trueならreplの引数は3つだけ",
          bool(_repl10b)
          and set(_repl10b[0].keys()) == {"use_search", "rag_dirs", "history_file"})
    check("main(iter186回帰): --out無しの対話分岐では--out無視警告は出ない",
          "[警告]" not in _res10b.stdout)

    # --- (11) 質問なし + stdin.isatty()=False + read()が非空文字列を返す:
    #          stripされた文字列がask_fuguのquestionになる。
    #          特性(b)の確認: この経路のask_fugu呼び出しにはoffice_attachedキーワード
    #          自体が渡されない(ask_fugu側の既定Falseに暗黙依存) ---
    f.SESSION_SAVE = True
    f._HISTORY = []
    _ask11, _repl11, _lh11 = _cli_run([], stdin_isatty=False,
                                       stdin_text="  パイプ入力の質問です  \n")
    check("main: パイプ入力はstripされてask_fuguのquestionになる",
          bool(_ask11) and _ask11[0]["question"] == "パイプ入力の質問です")
    check("main: パイプ入力経路ではreplは呼ばれない", len(_repl11) == 0)
    check("main(特性化・b): パイプ入力経路のask_fugu呼び出しはoffice_attachedを明示しない",
          bool(_ask11) and "office_attached" not in _ask11[0]["kwargs"])

    # --- (11b) --out指定 + パイプ入力(isatty()=False)で非空: 質問経路(6)と同様に
    #           out_fileがそのままask_fuguへ転送される。2026-07-26 iteration 186の
    #           対話分岐限定の警告追加が、他の分岐(パイプ入力)には一切影響しない
    #           ことの回帰確認 ---
    f.SESSION_SAVE = True
    f._HISTORY = []
    _res11b = _cli_run(["--out", "x.md"], stdin_isatty=False,
                        stdin_text="パイプ入力+out質問\n")
    _ask11b, _repl11b, _lh11b = _res11b
    check("main(iter186回帰): --out+パイプ入力ではask_fuguが1回呼ばれる",
          len(_ask11b) == 1)
    check("main(iter186回帰): --out+パイプ入力ではout_fileがそのまま転送される",
          bool(_ask11b) and _ask11b[0]["kwargs"].get("out_file") == "x.md")
    check("main(iter186回帰): --out+パイプ入力ではreplは呼ばれない", len(_repl11b) == 0)
    check("main(iter186回帰): --out+パイプ入力では対話分岐限定の警告は出ない",
          "[警告]" not in _res11b.stdout)

    # --- (12) 質問なし + stdin.isatty()=False + read()が空/空白のみ:
    #          「質問が入力されませんでした」経路を通り、ask_fugu/replとも呼ばれない ---
    f.SESSION_SAVE = True
    f._HISTORY = []
    _ask12, _repl12, _lh12 = _cli_run([], stdin_isatty=False, stdin_text="   \n  ")
    check("main: 空/空白のみのパイプ入力ならask_fuguは呼ばれない", len(_ask12) == 0)
    check("main: 空/空白のみのパイプ入力ならreplも呼ばれない", len(_repl12) == 0)
finally:
    f.setup = _orig_cli_setup
    f.ask_fugu = _orig_cli_ask_fugu
    f.repl = _orig_cli_repl
    f.load_history_file = _orig_cli_load_history
    f.read_file_text = _orig_cli_read_file_text
    f.SESSION_SAVE = _orig_cli_session_save
    f._HISTORY = _orig_cli_history
    f.RAG_DIRS = _orig_cli_rag_dirs
    sys.argv = _orig_cli_argv
    sys.stdin = _orig_cli_stdin

check("main: テスト後にf.setupが復元されている", f.setup is _orig_cli_setup)
check("main: テスト後にf.ask_fuguが復元されている", f.ask_fugu is _orig_cli_ask_fugu)
check("main: テスト後にf.replが復元されている", f.repl is _orig_cli_repl)
check("main: テスト後にf.load_history_fileが復元されている",
      f.load_history_file is _orig_cli_load_history)
check("main: テスト後にf.read_file_textが復元されている",
      f.read_file_text is _orig_cli_read_file_text)
check("main: テスト後にSESSION_SAVEが復元されている", f.SESSION_SAVE == _orig_cli_session_save)
check("main: テスト後に_HISTORYが復元されている", f._HISTORY == _orig_cli_history)
check("main: テスト後にRAG_DIRSが復元されている", f.RAG_DIRS == _orig_cli_rag_dirs)
check("main: テスト後にsys.argvが復元されている", sys.argv == _orig_cli_argv)
check("main: テスト後にsys.stdinが復元されている", sys.stdin is _orig_cli_stdin)

# ---------- repl(): 対話コマンドループ (exit/quit/reset/search on|off/save/空行/EOF) ----------
# repl() (L5382-5431) は対話モードの唯一のエントリポイントだが、save_history_file単体
# (会話履歴永続化セクション)とmain()からのrepl()ディスパッチ有無(iteration 185/186、
# 直上のmain()セクション)以外、ループ本体のコマンド解析と毎ターンの
# ask_fugu(q, use_search=..., rag_dirs=..., history_file=hfile)呼び出しには一度も
# 直接のテストが無かった。resetの永続化・search on/offのトグル伝播・history_fileの
# 転送はいずれも複数ターン会話の文脈維持に直結する精度クリティカルな経路であり、
# ここが壊れても他のテストでは検知できない。ここではrepl()を直接駆動し、
# コマンド分岐とask_fuguへの引数伝播を固定する。
#
# builtins.input を有限のスクリプト応答フェイクに差し替え、尽きたら明示的に
# EOFErrorを送出させることでループの終了を必ず保証する(ハング防止)。
# f.ask_fugu/f.save_history_fileは呼び出し記録フェイクに差し替え、本物の
# Ollama/conduct/ネットワークには一切触れない。urllib.request.urlopenと
# f.subprocess.runは「呼ばれたら即AssertionError」の番人に差し替える
# (gotcha #8 / iteration 38・76・104のtripwire流儀を踏襲)。全呼び出しでrepl()には
# 明示的な一時history_file=を渡し、本物の ~/.fugu_history.json (f.HISTORY_FILE) が
# テスト前後で一切作成/変更されないことも確認する。
# 触れないもの: 実際のOllama通信(/api/chat・num_ctx固定・thinkリトライ)、
# cp932 reconfigure、OLLAMA_MAX_LOADED_MODELS、math_verifyタイムアウト引数、
# solve_verifiableのSC投票内部ロジックはいずれも本セクションでは呼び出さない
# (ask_fugu自体を丸ごとフェイクに差し替えているため、その内部は素通りする)。
#
# 2026-07-27追記(iter205でiter204の特性化(a)/(b)を修正、iter157/172の
# 「stale characterization pinを反転させる」流儀を踏襲): 'save <path>'分岐には
# 元々2つの特性化(未修正のまま固定)があった。
#   (a) 末尾スペース無しの素の'save'(4文字)は low.startswith("save ")(5文字、
#       末尾スペース必須)に一致せず、コマンドとして扱われないままループ末尾の
#       ask_fugu(q, ...)へ落ちていた。つまり'save'という文字列そのものを質問として
#       通常のMoAパイプラインに渡してしまい、その応答が_HISTORYを汚染していた。
#   (b) さらに、'save '(末尾スペースのみ、パス無しのつもり)を入力しても、
#       repl()の先頭で q = input(...).strip() が入力全体を丸ごとstripする
#       ため、末尾の空白だけの"save "は判定前に"save"(4文字)へ潰れてしまい、
#       結局(a)と同じ経路に収束していた。「保存先パスを指定してください」
#       ガイダンス分岐(q[5:].strip()が空文字になる経路)はqが既にstrip済みで
#       ある(=末尾に空白を残せない)という不変条件と数学的に両立せず、実際の
#       input()経由では到達不能だった。
# iter205でlow == "save"の専用分岐を追加し、(a)/(b)とも同じガイダンス表示+
# continueに修正した(ask_fuguへのディスパッチなし、_HISTORY変更なし)。
# 'save <path>'(パス付き)自体の契約(iteration 182/186: パス大文字小文字保持・
# force=True・成功/失敗メッセージの分岐)には一切触れていない。
# 以下(7a)/(7b)は旧・特性化(未修正)テストをiter205で修正後の挙動へ反転したもの。
import builtins as _repl_builtins
import tempfile as _repl_tempfile

_orig_repl_input = _repl_builtins.input
_orig_repl_ask_fugu = f.ask_fugu
_orig_repl_save_history_file = f.save_history_file
_orig_repl_history = f._HISTORY
_orig_repl_urlopen = urllib.request.urlopen
_orig_repl_subprocess_run = f.subprocess.run

_repl_real_hist_path = f.HISTORY_FILE
_repl_real_hist_existed_before = _repl_real_hist_path.exists()
_repl_real_hist_bytes_before = (
    _repl_real_hist_path.read_bytes() if _repl_real_hist_existed_before else None)


def _repl_no_network_urlopen(*a, **k):
    raise AssertionError("repl: モック漏れで実urlopen(ネットワーク)が呼ばれた")


def _repl_no_subprocess_run(*a, **k):
    raise AssertionError("repl: モック漏れで実subprocess.runが呼ばれた")


def _repl_make_scripted_input(script):
    """有限のスクリプト応答を順に返すinput()フェイク。尽きたら明示的にEOFErrorを
    送出しループの終了を保証する('finite scripted iterator whose exhaustion
    raises EOFError'の要件)。script中の要素がBaseExceptionのクラスまたは
    インスタンスならそれをそのまま送出する(KeyboardInterrupt注入用)。"""
    _it = iter(list(script))

    def _fake_input(prompt=""):
        try:
            item = next(_it)
        except StopIteration:
            raise EOFError("repl test: scripted input exhausted")
        if isinstance(item, BaseException):
            raise item
        if isinstance(item, type) and issubclass(item, BaseException):
            raise item()
        return item

    return _fake_input


_repl_ask_calls = []
_repl_save_calls = []


def _repl_fake_ask_fugu(question, **kwargs):
    _repl_ask_calls.append({"question": question, "kwargs": kwargs})
    return "FAKE_ANSWER"


def _repl_fake_save_history_file(*args, **kwargs):
    _repl_save_calls.append({"args": args, "kwargs": kwargs})
    return True


def _repl_run(script, **repl_kwargs):
    """input()を有限スクリプトのフェイクに差し替えてf.repl(**repl_kwargs)を1回
    駆動する。呼び出し前に_repl_ask_calls/_repl_save_callsをクリアしてから実行し、
    (ask_calls, save_calls, stdout_text, raised_exc)を返す(raised_excは
    repl()内で捕捉されず外に漏れた例外。Noneなら伝播なし)。f.ask_fugu/
    f.save_history_file自体の差し替えは、この関数を呼ぶtry節側で一度だけ行う。"""
    _repl_ask_calls.clear()
    _repl_save_calls.clear()
    _repl_builtins.input = _repl_make_scripted_input(script)
    _stdout_buf = io.StringIO()
    raised_exc = None
    try:
        with contextlib.redirect_stdout(_stdout_buf):
            f.repl(**repl_kwargs)
    except BaseException as _e:
        raised_exc = _e
    return (list(_repl_ask_calls), list(_repl_save_calls),
            _stdout_buf.getvalue(), raised_exc)


try:
    f.ask_fugu = _repl_fake_ask_fugu
    f.save_history_file = _repl_fake_save_history_file
    urllib.request.urlopen = _repl_no_network_urlopen
    f.subprocess.run = _repl_no_subprocess_run

    with _repl_tempfile.TemporaryDirectory() as _repl_dir:
        _repl_hfile = f.Path(_repl_dir) / "repl_history.json"

        # --- (1) exit/quit(小文字)は即終了し、ask_fuguは呼ばれない ---
        for _cmd in ("exit", "quit"):
            f._HISTORY = []
            _ask, _save, _out, _exc = _repl_run([_cmd], history_file=_repl_hfile)
            check(f"repl: '{_cmd}'は即座にループを終了しask_fuguを呼ばない", len(_ask) == 0)
            check(f"repl: '{_cmd}'はrepl()内で例外を送出しない(外に伝播しない)", _exc is None)

        # --- (1b) 大文字小文字混在でも同様(low = q.lower()での判定のため) ---
        for _cmd in ("EXIT", "Quit", "QUIT", "eXiT"):
            f._HISTORY = []
            _ask, _save, _out, _exc = _repl_run([_cmd], history_file=_repl_hfile)
            check(f"repl: 大文字小文字混在'{_cmd}'もask_fuguを呼ばずに終了する",
                  len(_ask) == 0 and _exc is None)

        # --- (2) 空/空白のみの入力はask_fuguを呼ばずスキップし、次の入力へ進む ---
        f._HISTORY = []
        _ask, _save, _out, _exc = _repl_run(["   ", "\t", "", "exit"],
                                             history_file=_repl_hfile)
        check("repl: 空/空白のみの入力はすべてスキップされask_fuguは呼ばれない",
              len(_ask) == 0)
        check("repl: 空白スキップ後も例外なく'exit'まで到達しループが終了する",
              _exc is None)

        # --- (3) reset: _HISTORYを同一オブジェクトのまま空にし、save_history_fileを
        #         path=hfileでちょうど1回、force=Trueを渡さずに呼ぶ ---
        f._HISTORY = [{"role": "user", "content": "以前の質問"},
                      {"role": "assistant", "content": "以前の回答"}]
        _hist_ref_before = f._HISTORY
        _hist_id_before = id(f._HISTORY)
        _ask, _save, _out, _exc = _repl_run(["reset", "exit"], history_file=_repl_hfile)
        check("repl: resetは_HISTORYを同一オブジェクトのまま(id不変)空にする",
              id(f._HISTORY) == _hist_id_before)
        check("repl: reset後、_HISTORY(同一オブジェクト参照)の中身が空リストになる",
              f._HISTORY == [] and _hist_ref_before == [])
        check("repl: resetはask_fuguを呼ばない", len(_ask) == 0)
        check("repl: resetはsave_history_fileをちょうど1回呼ぶ", len(_save) == 1)
        check("repl: resetのsave_history_file呼び出しは_HISTORYそのもの"
              "(同一オブジェクト)を第1引数として渡す",
              bool(_save) and _save[0]["args"]
              and _save[0]["args"][0] is f._HISTORY)
        check("repl: resetのsave_history_file呼び出しはpath=hfileを渡す",
              bool(_save) and _save[0]["kwargs"].get("path") == _repl_hfile)
        check("repl: resetのsave_history_file呼び出しはforce=Trueを渡さない"
              "(既存のforce=False既定でのno-op可能な仕様を変えない、既知仕様の固定。"
              "force=Trueへ強制するのはこのイテレーションの対象外)",
              bool(_save) and "force" not in _save[0]["kwargs"])
        check("repl: reset実行後にクリア完了メッセージが出力される", "クリア" in _out)

        # --- (4) search on/off: 複数ターンにまたがってトグル状態が保持される ---
        f._HISTORY = []
        _ask, _save, _out, _exc = _repl_run(
            ["search on", "質問A", "search off", "質問B", "search on", "質問C", "exit"],
            use_search=False, rag_dirs=None, history_file=_repl_hfile)
        check("repl: search on/off/onを挟んだ3つの質問でask_fuguが3回呼ばれる",
              len(_ask) == 3)
        check("repl: 'search on'直後の質問はuse_search=Trueでディスパッチされる",
              len(_ask) == 3 and _ask[0]["kwargs"].get("use_search") is True)
        check("repl: 'search off'直後の質問はuse_search=Falseに戻る",
              len(_ask) == 3 and _ask[1]["kwargs"].get("use_search") is False)
        check("repl: 再度'search on'した質問はuse_search=Trueに戻る"
              "(トグル状態はセッション内の複数ターン間で保持される)",
              len(_ask) == 3 and _ask[2]["kwargs"].get("use_search") is True)
        check("repl: search on/offのON/OFF表示がそれぞれ出力される",
              "ON" in _out and "OFF" in _out)

        # --- (5) 通常の質問: ask_fuguをちょうど1回、question+
        #         use_search/rag_dirs/history_fileのみのkwargsで呼ぶ
        #         (out_file/office_attachedは渡さない、iteration 186設計の固定) ---
        f._HISTORY = []
        _rag_dirs_e = ["dirA", "dirB"]
        _ask, _save, _out, _exc = _repl_run(
            ["これは通常の質問です", "exit"],
            use_search=True, rag_dirs=_rag_dirs_e, history_file=_repl_hfile)
        check("repl: 通常の質問はask_fuguをちょうど1回呼ぶ", len(_ask) == 1)
        check("repl: 質問文がそのまま(加工なしで)ask_fuguへ渡る",
              bool(_ask) and _ask[0]["question"] == "これは通常の質問です")
        check("repl: kwargsはuse_search/rag_dirs/history_fileの3つだけ",
              bool(_ask)
              and set(_ask[0]["kwargs"].keys()) == {"use_search", "rag_dirs", "history_file"})
        check("repl: use_searchはrepl()呼び出し時点の現在のトグル値と一致する",
              bool(_ask) and _ask[0]["kwargs"].get("use_search") is True)
        check("repl: rag_dirsはrepl()に渡したrag_dirsがそのまま渡る(同一オブジェクト)",
              bool(_ask) and _ask[0]["kwargs"].get("rag_dirs") is _rag_dirs_e)
        check("repl: history_fileは明示的に渡した一時ファイルパスがそのまま渡る",
              bool(_ask) and _ask[0]["kwargs"].get("history_file") == _repl_hfile)

        # --- (5b) 'save'で始まるが'save'単体でも'save '(スペース)プレフィックスでも
        #          ない通常の質問(例: "saving money tips")は、iter205で追加した
        #          low == "save"分岐にも既存のlow.startswith("save ")分岐にも
        #          一致せず、従来通りask_fuguへ通常ディスパッチされる
        #          (誤って新分岐に飲み込まれていないことの確認) ---
        f._HISTORY = []
        _ask, _save, _out, _exc = _repl_run(["saving money tips", "exit"],
                                             history_file=_repl_hfile)
        check("repl(iter205回帰防止): 'save'で始まるだけの通常の質問"
              "('saving money tips')はsaveコマンドと誤認識されずask_fuguへ渡る",
              len(_ask) == 1 and _ask[0]["question"] == "saving money tips")
        check("repl(iter205回帰防止): 'saving money tips'ではsave_history_fileは"
              "呼ばれない", len(_save) == 0)

        # --- (6) save <path>: 大文字小文字混在パスがそのまま(q[5:]、lowerされたlowでは
        #         ない)保持されforce=Trueで保存される ---
        f._HISTORY = [{"role": "user", "content": "保存対象"}]
        _ask, _save, _out, _exc = _repl_run(["save MixedCase/Path.json", "exit"],
                                             history_file=_repl_hfile)
        check("repl: 'save <path>'はask_fuguを呼ばない", len(_ask) == 0)
        check("repl: 'save <path>'はsave_history_fileをちょうど1回呼ぶ", len(_save) == 1)
        check("repl: saveコマンドのpathは大文字小文字をそのまま保持する"
              "(元のqから切り出す。lower()されたlowからではない。f.Path()での比較は"
              "Windows上でのセパレータ正規化(/ -> \\)を許容しつつ大文字小文字の"
              "保持だけを見る)",
              bool(_save)
              and _save[0]["kwargs"].get("path") == f.Path("MixedCase/Path.json")
              and "MixedCase" in str(_save[0]["kwargs"].get("path"))
              and "mixedcase" not in str(_save[0]["kwargs"].get("path")))
        check("repl: saveコマンドはforce=Trueで呼ぶ",
              bool(_save) and _save[0]["kwargs"].get("force") is True)
        check("repl: save成功メッセージにパスの大文字小文字を保持したまま出力される",
              "MixedCase/Path.json" in _out)

        # --- (7a) iter205修正確認(旧・特性化(a)の反転): 末尾スペース無しの素の
        #          'save'はコマンドとして認識され、ask_fuguへは一切ディスパッチ
        #          されずガイダンスだけ表示してcontinueする ---
        f._HISTORY = []
        _ask, _save, _out, _exc = _repl_run(["save", "exit"], history_file=_repl_hfile)
        check("repl(iter205修正、旧・特性化(a)を反転): 末尾スペース無しの素の'save'は"
              "ask_fuguへディスパッチされない(_HISTORYを汚染する質問実行を防止)",
              len(_ask) == 0)
        check("repl(iter205修正): 素の'save'はsave_history_fileも呼ばない"
              "(forced saveは発火しない)", len(_save) == 0)
        check("repl(iter205修正): 素の'save'は'save <path>'の使い方ガイダンスを表示する",
              "save <path>" in _out and "保存先パスを指定してください" in _out)
        check("repl(iter205修正): 素の'save'入力後もループが継続し次のexitで正常終了する"
              "(例外が外へ伝播しない)", _exc is None)

        # --- (7b) iter205修正確認(旧・特性化(b)の反転): 'save '(末尾スペースのみ、
        #          パス無しのつもり)も、repl()冒頭のq = input(...).strip()で
        #          'save'(4文字)へ潰れてから判定されるため、(7a)と全く同じ
        #          ガイダンス表示+continueの経路になる(ask_fuguディスパッチなし) ---
        f._HISTORY = []
        _ask, _save, _out, _exc = _repl_run(["save ", "exit"], history_file=_repl_hfile)
        check("repl(iter205修正、旧・特性化(b)を反転): 'save '(末尾スペースのみ)も"
              "外側のstrip()で'save'に潰れ、(7a)と同じくask_fuguへディスパッチされない",
              len(_ask) == 0)
        check("repl(iter205修正): 'save '(末尾スペースのみ)でもforced save"
              "(save_history_file)は呼ばれない", len(_save) == 0)
        check("repl(iter205修正): 'save '(末尾スペースのみ)でも使い方ガイダンスが表示される",
              "save <path>" in _out and "保存先パスを指定してください" in _out)

        # --- (7c) 'SAVE'(大文字)も既存のlow(小文字化済み)判定により(7a)と
        #          同一に扱われる('search on/off'等、他コマンドとの大文字小文字
        #          非依存の一貫性を確認) ---
        f._HISTORY = []
        _ask, _save, _out, _exc = _repl_run(["SAVE", "exit"], history_file=_repl_hfile)
        check("repl(iter205修正): 大文字小文字混在'SAVE'も素の'save'と同様に"
              "ask_fuguへディスパッチされない",
              len(_ask) == 0 and len(_save) == 0)
        check("repl(iter205修正): 'SAVE'でも使い方ガイダンスが表示される",
              "save <path>" in _out)

        # --- (8) EOFError/KeyboardInterrupt はrepl()の外へ伝播せず、クリーンに終了する ---
        f._HISTORY = []
        _ask, _save, _out, _exc = _repl_run([EOFError], history_file=_repl_hfile)
        check("repl: input()からのEOFErrorはrepl()の外へ伝播しない", _exc is None)
        check("repl: EOFError発生時はask_fuguを呼ばない", len(_ask) == 0)
        check("repl: EOFError発生時に終了メッセージが出力される", "終了します" in _out)

        f._HISTORY = []
        _ask, _save, _out, _exc = _repl_run([KeyboardInterrupt], history_file=_repl_hfile)
        check("repl: input()からのKeyboardInterruptはrepl()の外へ伝播しない", _exc is None)
        check("repl: KeyboardInterrupt発生時はask_fuguを呼ばない", len(_ask) == 0)
        check("repl: KeyboardInterrupt発生時にも終了メッセージが出力される",
              "終了します" in _out)

        # --- (8b) 'exit'を打たずスクリプトが尽きた場合も、自動送出されるEOFErrorで
        #          クリーンに終了する(全シナリオ共通のハング防止機構そのものの確認) ---
        f._HISTORY = []
        _ask, _save, _out, _exc = _repl_run(["質問のみ、exitなし"], history_file=_repl_hfile)
        check("repl: スクリプト尽き(自動EOFError)でもrepl()は例外を送出しない",
              _exc is None)
        check("repl: スクリプトが尽きる前の質問は通常通りask_fuguへ渡る",
              len(_ask) == 1 and _ask[0]["question"] == "質問のみ、exitなし")

finally:
    _repl_builtins.input = _orig_repl_input
    f.ask_fugu = _orig_repl_ask_fugu
    f.save_history_file = _orig_repl_save_history_file
    f._HISTORY = _orig_repl_history
    urllib.request.urlopen = _orig_repl_urlopen
    f.subprocess.run = _orig_repl_subprocess_run

check("repl: テスト後にbuiltins.inputが復元されている", _repl_builtins.input is _orig_repl_input)
check("repl: テスト後にf.ask_fuguが復元されている", f.ask_fugu is _orig_repl_ask_fugu)
check("repl: テスト後にf.save_history_fileが復元されている",
      f.save_history_file is _orig_repl_save_history_file)
check("repl: テスト後にf._HISTORYが復元されている", f._HISTORY == _orig_repl_history)
check("repl: テスト後にurllib.request.urlopenが復元されている",
      urllib.request.urlopen == _orig_repl_urlopen)
check("repl: テスト後にf.subprocess.runが復元されている",
      f.subprocess.run == _orig_repl_subprocess_run)
check("repl: 実履歴ファイル(~/.fugu_history.json)はテスト前後で存在状態が変化しない"
      "(全呼び出しに明示的な一時history_file=を渡し、save_history_file自体もフェイクの"
      "ため一切触れていないはず)",
      _repl_real_hist_path.exists() == _repl_real_hist_existed_before)
if _repl_real_hist_existed_before:
    check("repl: 実履歴ファイルの内容もテスト前後でバイト単位不変",
          _repl_real_hist_path.read_bytes() == _repl_real_hist_bytes_before)

print()
if _FAILS:
    print(f"FAILED: {len(_FAILS)} 件 -> {_FAILS}")
    raise SystemExit(1)
print("ALL PASSED")
