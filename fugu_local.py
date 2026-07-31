"""
Local Fugu-style Orchestrator (RTX 4060 Laptop 8GB VRAM / RAM 48GB / i7-13700H 向け)

本家 Sakana Fugu との違いを埋めるための版。
古典的な静的 MoA（＝Sakana でいう "Fusion"）は「固定のプロポーザー全員が毎回走り、
固定のアグリゲーターが統合する」だけだが、Fugu の肝は指揮者(Conductor)LLM 自身が
「単体で十分か / 誰を何体使うか / 何回反復するか / 追加ラウンドが要るか」を
"動的に" 決めるところにある。

そこでこの版では:
  1) Conductor が質問を見て実行プランを JSON で出す（単体 or 合議、使うモデル、ラウンド数）
  2) 簡単な質問は 1 モデルで即答（MoA のオーバーヘッドを払わない）
  3) 単体回答が弱ければ Critic が検知して合議へ "エスカレーション"
  4) 合議後も不十分なら上限付きで "再帰的に" 追加ラウンド
を行う。全てローカルモデルだけで完結。
"""

import os
import re
import sys
import json
import time
import shutil
import tempfile
import argparse
import subprocess
import urllib.request
import urllib.parse
import concurrent.futures  # 並列処理用（8GB では既定で逐次）
from pathlib import Path

# Windows の cp932 コンソール/パイプでは ⚠ ✓ ❌ ⤴ ↻ 等の表示記号が encode できず、
# print 自体が UnicodeEncodeError で落ちる（実測 2026-07-04: 保険1の通知 print がクラッシュし、
# 空返答の救済パスそのものが死んだ）。記号が化けるのは許容し、encode 不能文字は置換して続行する。
for _stream in (sys.stdout, sys.stderr):
    if _stream is not None and hasattr(_stream, "reconfigure"):
        _stream.reconfigure(errors="replace")

# 推論は Ollama native /api/chat を urllib で直接叩く（依存ゼロ）。
# 旧版は openai クライアント + /v1 互換エンドポイントを使っていたが、/v1 は num_ctx を
# 無視してモデル最大 context(例: qwen3 は 262144)を確保しようとし、8GB VRAM では
# KV キャッシュ確保に失敗して llama-server がクラッシュする（500）。native /api/chat なら
# options.num_ctx がリクエスト単位で効くため、これで context を安全域に固定する。

# ==================================================
# 設定
# ==================================================

OLLAMA_URL = os.environ.get("OLLAMA_URL", "http://localhost:11434")

# --- モデルの役割 ---
# Conductor + Critic: qwen3:4b（軽量・高速・JSON安定。ルーティング専用、VRAM常駐最小化）
# Proposers: qwen3-coder:30b(コード特化MoE) / phi4(数学・PINN・物理) / gpt-oss:20b(汎用推論MoE)
#   8GB VRAM 環境では大型モデルは RAM オフロード実行（Ollama 自動制御、48GB RAM で吸収）。
#   RTX 4090 24GB 移行後はすべて VRAM 内で動作。
# Aggregator: qwen3-coder:30b（コード統合に最適）
# JP Aggregator: qwen3:4b（日本語は実績ある qwen3 で確実に処理）
# ドキュメント "3大AIオールスター" 構成: A=GPT / B=Claude / C=Gemini + D=理数専門家。
# 順序はペルソナ A,B,C,D に対応させる（gemma-4-26b-a4b の実タグは gemma4:26b）。
DESIRED_PROPOSERS = ["gpt-oss:20b", "qwen3-coder:30b", "gemma4:26b", "qwen3.6:35b"]
DESIRED_AGGREGATOR = "qwen3-coder:30b"
DESIRED_CONDUCTOR = "qwen3:4b"
FALLBACK_MODEL = "qwen3:4b"

# --- ペルソナ層（3大AIオールスター）---
# Conductor は selected_proposers に「Proposer A」〜「Proposer D」のペルソナ名を出す。
# ここで実モデルへ解決する（_resolve_proposer 参照）。モデル未導入なら解決結果から除外される。
PERSONA_MODELS = {
    "Proposer A": "gpt-oss:20b",      # ChatGPT(GPT)の存在
    "Proposer B": "qwen3-coder:30b",  # Claudeの存在
    "Proposer C": "gemma4:26b",       # Geminiの存在
    "Proposer D": "qwen3.6:35b",      # 理数・物理・PINN 専門家（2026-07-11: phi4 → qwen3.6:35b
                                      # へ更新。思考型 MoE(A3B) で数学・推論が大幅に強い）
}
# 各 proposer に注入する人格プロンプト（PROPOSER_SYS の前に前置し「個性」を再現する）
PERSONA_IDENTITY = {
    "gpt-oss:20b":     "あなたは『ChatGPT(GPT)の存在』。バランス感覚に優れ、一般的な対話と文章の骨組み作りを担当する。",
    "qwen3-coder:30b": "あなたは『Claudeの存在』。高度なプログラミング、厳密な論理チェック、コードの自己修復を担当する。",
    "gemma4:26b":      "あなたは『Geminiの存在』。RAG(Office文書)のコンテキスト分析、大量ドキュメントとWeb検索結果の集約を担当する。",
    "qwen3.6:35b":     "あなたは理数・物理・PINN(物理情報ニューラルネット)・偏微分方程式の専門家。厳密に段階を追って考える。",
}
MODEL_TO_PERSONA = {v: k for k, v in PERSONA_MODELS.items()}

# 日本語質問では qwen3-coder の日本語品質が未検証のため、実績ある qwen3:4b に切替える。
JP_AGGREGATOR = "qwen3:4b"

# 統合役の強化（2026-07-11）: 統合は「最終の誤り検出」なので思考モデルにやらせるのが筋。
# - コードを含む統合 → qwen3-coder:30b（従来どおり。実行検証タグの扱いに実績）
# - 日本語の記述式 → qwen3.6:35b（思考ON・日本語堅牢。JPサニティ不合格なら "" にして qwen3:4b 維持）
# - 英語の記述式 → gpt-oss:20b（think:"high" が MODEL_CONFIG で自動適用される）
JP_AGGREGATOR_STRONG = "qwen3.6:35b"
AGGREGATOR_REASONING = "gpt-oss:20b"

# 自己評価バイアス対策: Conductor/Critic が qwen3:4b のため、別系統モデルで独立チェックする。
# gpt-oss:20b は OpenAI 系で qwen3 と出自が異なるため second opinion に適切。
SECOND_OPINION_MODEL = "gpt-oss:20b"

# Conductor がモデル選抜に使う「各プロポーザーの得意分野」ヒント（ペルソナ付き）
PROPOSER_PROFILES = {
    "gpt-oss:20b":     "ChatGPT(GPT)の存在。バランス・一般的な対話・文章の骨組み (OpenAI OSS MoE・3.6B active・思考high対応)",
    "qwen3-coder:30b": "Claudeの存在。高度なプログラミング・厳密な論理チェック・自己修復 (SWE-bench最強クラスMoE)",
    "gemma4:26b":      "Geminiの存在。RAG(Office文書)分析・大量ドキュメント・Web検索結果の集約 (26B)",
    "qwen3.6:35b":     "理数・物理・PINN・偏微分方程式・アルゴリズム証明に強い思考型 (35B MoE A3B, 2026-02世代)",
}

# --- 生成パラメータ ---
# 大型モデルの RAM オフロード時に1コール最大2時間を確保（8192tok × 1tok/s の保険）。
REQUEST_TIMEOUT = 7200
PROPOSER_TEMP = 0.7       # 多様性のため高め
AGGREGATOR_TEMP = 0.5
CONDUCTOR_TEMP = 0.1      # 判定はブレを抑える

# proposer の思考(thinking)制御。None=モデル既定(qwen3/gemma は思考ON=高品質だが遅い)、
# False=思考を無効化して高速化(思考トークンは最終回答に使わず破棄されるため、8GB逐次では
# 大きな時間短縮になる)。品質とのトレードオフなので既定は None。非思考モデル(phi4-mini)に
# False を渡しても無害(実測でエラーなし)。
PROPOSER_THINK = None       # None=モデル既定。think:true を送ると qwen3-coder/phi4 等が
                            # 400 "does not support thinking" で即失敗するため送らない。
                            # gpt-oss:20b は think:true 対応済み（実測 2026-07-05）だが、
                            # 個別制御より既定委任の方が安全。

# --- モデル別推論設定（精度最優先の中核。2026-07-11 追加）---
# think:  None=モデル既定 / False=無効 / True=有効 / "low"|"medium"|"high"=gpt-oss 系の
#         思考量段階指定（Ollama 0.31.2 で "high" が gpt-oss:20b に効くことを実測済み。
#         high は数学・推論の精度が跳ね上がる最大のレバー）。
#         PROPOSER_THINK が None 以外ならそちらが最優先（eval の一括 OFF 等を維持）。
# num_ctx: モデル別コンテキスト長。従来は全モデル一律 MODEL_NUM_CTX=8192 だったが、
#         それは「VRAM 常駐する小型モデルの KV 上限」の話。大型モデルは重みも KV も
#         RAM オフロードで動くため 32768 まで拡大できる（思考が 8k を超えて打ち切られる
#         事故＝done_reason=length・本文空、を根本から防ぐ）。qwen3:4b 等の VRAM 常駐組は
#         8192 を維持して高速なまま使う。
# num_predict: 生成上限（思考トークン込み）。num_ctx 拡大に合わせて引き上げる。
# 【Phase 0 実測 2026-07-12, RTX 4060 8GB / RAM48GB / Ollama 0.31.2】
# - gpt-oss:20b think:"high" は動作OK。ただし num_ctx で速度が激変:
#   8192→14.9 tok/s / 32768→3.4 tok/s（KV が RAM オフロードされ帯域律速）。
# - qwen3.6:35b think:true は num_predict=8192 だと長い思考で使い切り本文ゼロ(done=length)に
#   なる実測。num_predict は厚めに、num_ctx は 16384 で throughput を確保するのが妥当。
# - VibeThinker-3B は 11.2 tok/s・VRAM 常駐、<think> タグは content 内(strip_think が除去)。
# 方針: 思考モデルは num_ctx=16384（8192 だと AIME の思考が入り切らず、32768 は遅すぎる中間点）、
# num_predict は打ち切り回避のため厚め。非思考の大型は 16384 で十分。
MODEL_CONFIG = {
    "gpt-oss:20b":     {"think": "high", "num_ctx": 16384, "num_predict": 14336},
    "gpt-oss:120b":    {"think": "high", "num_ctx": 12288, "num_predict": 8192},
    "qwen3.6:35b":     {"think": True,  "num_ctx": 16384, "num_predict": 14336},
    "NitrAI/VibeThinker-3B": {"num_ctx": 16384, "num_predict": 14336},
    "qwen3-coder:30b": {"num_ctx": 16384, "num_predict": 12288},
    "gemma4:26b":      {"num_ctx": 16384, "num_predict": 12288},
}


def model_cfg(model, key, default=None):
    """MODEL_CONFIG からモデル別設定を引く（無ければ default）。"""
    return MODEL_CONFIG.get(model, {}).get(key, default)


# ==================================================
# 大 VRAM プロファイル（将来の 96GB 等の環境向け・一発切り替え）
# ==================================================
# 8GB ラップトップは「大型モデルを RAM/NVMe にオフロードして逐次で回す」制約下にあるが、
# nk108 は本手法が有効なら VRAM 96GB 級の環境で実験予定。そこでは全モデルが VRAM 常駐でき、
# 制約が一変する（並列プロポーザー可・context 大幅拡大・SC のサンプル数を大量に増やせる・
# 120b arbiter も高速）。環境変数 FUGU_HIGH_VRAM=1 で下記を一括適用する（コード改変不要）。
#   PowerShell: $env:FUGU_HIGH_VRAM=1 ; python fugu_local.py ...
# 値は 96GB を想定した保守的な既定。より大きい環境ならさらに引き上げてよい。
def apply_high_vram_profile():
    """VRAM 潤沢環境向けに設定を一括で引き上げる。setup() 冒頭で env 判定して呼ぶ。"""
    global MODEL_CONFIG, PARALLEL_PROPOSERS, MODEL_NUM_CTX
    global SC_INITIAL, SC_STEP, SC_MAX, SC_CHEAP_VOTES, MODEL_KEEP_ALIVE, ARBITER_MODEL
    print("[setup] FUGU_HIGH_VRAM=1 → 大VRAMプロファイルを適用します")
    # 全モデル常駐前提: context を広げ生成上限も引き上げる（KV が VRAM に載るため安全）
    for m, cfg in MODEL_CONFIG.items():
        cfg["num_ctx"] = 65536
        cfg["num_predict"] = 32768
    MODEL_NUM_CTX = 32768
    # 96GB なら複数モデルを同時常駐でき、プロポーザー並列が効く（8GB では逆効果だった）
    PARALLEL_PROPOSERS = True
    MODEL_KEEP_ALIVE = "30m"          # 常駐維持でロード/アンロードの往復を消す
    # サンプルを大量に回せる＝自己一貫性の精度が上がる主レバー
    SC_INITIAL, SC_STEP, SC_MAX = 12, 8, 48
    SC_CHEAP_VOTES = 16               # VibeThinker を大量票に（多様性の底上げ）
    # 96GB では 65GB の 120b が VRAM 常駐できるため、拮抗時の裁定をローカル最上位知能に任せる
    # （8GB 既定は NVMe ページング回避で qwen3.6:35b）
    ARBITER_MODEL = "gpt-oss:120b"
    print(f"[setup] high-vram: num_ctx=65536 parallel=ON SC(init={SC_INITIAL},max={SC_MAX}) "
          f"cheap_votes={SC_CHEAP_VOTES} arbiter={ARBITER_MODEL}")


def proposer_think_for(model):
    """proposer の think 解決: グローバル PROPOSER_THINK(≠None) > MODEL_CONFIG > モデル既定。"""
    if PROPOSER_THINK is not None:
        return PROPOSER_THINK
    return model_cfg(model, "think")


def proposer_predict_for(model):
    """proposer/aggregator の生成上限: MODEL_CONFIG > 役割既定。"""
    return model_cfg(model, "num_predict", NUM_PREDICT_PROPOSER)

# --- 生成長の上限（暴走保険）---
# 未指定だと思考モデルの生成が無制限で、実測では deepseek-r1 が統合 1 回に ~4100 トークン
# (411秒) 生成した例がある。精度優先のため「打ち切り」が起きない余裕を持たせた上限とし、
# タイトな時間予算にはしない。目的は無限の暴走を有界にすることだけ。
# 【実測の教訓 2026-07-04】上限は思考トークンも消費する。長いコード/証明の統合では
# 思考だけで 5120 を食い尽くし「done_reason=length・本文空」で終わる事象が 3 回発生
# （保険2が救済）。上限は num_ctx=8192 から入力(~2k)を引いた範囲で最大限に取り、
# ask() 側で「思考中に打ち切られて本文空」を __ERROR__ として可視化する。
NUM_PREDICT_PROPOSER = 8192      # 時間無制限・打ち切りゼロのため上限最大化
NUM_PREDICT_AGGREGATOR = 8192    # 統合も上限最大化
NUM_PREDICT_JUDGE = 768          # Conductor/Critic の高速JSON(think=False)
NUM_PREDICT_JUDGE_THINK = 6144   # Critic 再検算（思考トークンに十分な余裕）

# --- Fugu 風オーケストレーションの挙動 ---
MAX_ROUNDS = 4            # 時間無制限・精度優先のため反復を増やす
ADAPTIVE_ESCALATION = True  # 単体回答が弱いと合議へ格上げ
ALLOW_RECURSION = True      # 合議後、批評 → 必要なら追加ラウンド

# --- コード実行検証（主用途: コード生成の自律修正ループ）---
# 回答中の ```python ブロックを実際に subprocess で実行し、失敗したら traceback を
# 次ラウンドの修正ヒントとして渡す。LLM の自己審査と違い実行結果は決定的なので、
# コードに関しては最強の Critic になる。エラーが残る限り MAX_ROUNDS_CODE まで
# 修正ラウンドを繰り返す（nk108 の方針: 時間をかけてでも精度優先）。
# 注意: 生成コードをこのマシンで直接実行する。信頼できる自分の質問にだけ使うこと。
CODE_EXECUTION = True
CODE_EXEC_TIMEOUT = 15      # 秒。input() 待ちや無限ループはタイムアウトで失敗扱い
MAX_ROUNDS_CODE = 8         # コード修正は時間をかけて完全に直す

# --- 表示 ---
SHOW_PLAN = True          # Conductor の判断を表示（Fugu 風の動作を可視化）
SHOW_PROPOSALS = True     # 各提案を表示（think は除去して表示）
SHOW_BASELINE = False     # 比較用の単体直答。学習用に True でも可

# --- 計測（フェーズ3用）---
# True にすると ask() が各呼び出しの (label, model, 秒) を _TIMINGS に記録する。
# 段階別（conductor/proposer/aggregator/critic）の所要時間を可視化するための軽量フック。
SHOW_TIMING = False
_TIMINGS = []

# --- 実行時フラグ ---
_SECOND_OPINION_DISABLED = False  # second opinion が未インストール時に True にセット

# --- セッション永続化 ---
# 会話履歴を JSON ファイルに保存し、次回起動時に復元する。
HISTORY_FILE: Path = Path.home() / ".fugu_history.json"
SESSION_SAVE = True          # False にすると永続化を無効化（--no-history フラグでも制御）
MAX_HISTORY_TURNS_SAVED = 50 # ファイルに保存する最大往復数（古い順に削除）

# --- Web 検索 ---
# duckduckgo_search パッケージ（pip install duckduckgo_search）が入っていれば
# フル検索結果を取得。未インストール時は DuckDuckGo Instant Answer API（urllib 内蔵）
# にフォールバック（インスタント回答のみ・件数少ない）。
WEB_SEARCH_MAX_RESULTS = 5       # 1 クエリあたりの取得件数
WEB_SEARCH_SNIPPET_CHARS = 400   # 各スニペットの文字数上限
WEB_SEARCH_TIMEOUT = 15          # 秒

# --- 反復リサーチ ---
# 1 回の検索では具体的事実（型番・アーキテクチャ名等）が欠けたまま、モデルが古い学習知識で
# 穴埋めする事故が起きる（実測 2026-07-06: RTX 5090 のアーキテクチャを Hopper と誤答）。
# そこで Conductor(qwen3:4b) に「十分な事実が集まったか」を判定させ、不足なら不足点を狙った
# 追加クエリを生成して検索を繰り返す。
SEARCH_MAX_ROUNDS = 3            # リサーチ反復の上限（十分と判定されたら早期終了）
SEARCH_CONTEXT_CHARS = 4000      # 質問に注入する検索コンテキスト上限（num_ctx=8192 の安全域）

# --- RAG（ローカル文書検索）---
# RAG_DIRS に 1 つ以上のディレクトリを指定すると、質問と関連するチャンクを
# 自動抽出してプロポーザーへのコンテキストとして注入する。
# CLI: --rag /path/to/docs  または  --rag dir1 dir2 ...
RAG_DIRS: list = []              # 空 = 無効
RAG_CHUNK_CHARS = 600            # チャンクサイズ（文字数）
RAG_CHUNK_OVERLAP = 100          # チャンク間のオーバーラップ文字数
RAG_TOP_K = 3                    # 注入する上位チャンク数
RAG_EXTENSIONS = {
    # テキスト・コード
    ".txt", ".md", ".rst", ".csv", ".tsv", ".json", ".jsonl", ".yaml", ".yml",
    ".toml", ".ini", ".cfg", ".env",
    ".py", ".ipynb", ".r", ".m", ".jl",
    ".js", ".ts", ".jsx", ".tsx", ".mjs",
    ".go", ".rs", ".c", ".cpp", ".h", ".hpp", ".cs", ".java", ".kt", ".swift",
    ".rb", ".php", ".sh", ".bat", ".ps1", ".sql", ".graphql",
    ".html", ".htm", ".xml", ".svg", ".css", ".scss",
    ".tex", ".bib",
    # ドキュメント（要ライブラリ）
    ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt",
    ".odt", ".ods", ".odp",
}

# --- 画像生成（ローカル Stable Diffusion / ComfyUI 連携）---
# Ollama はテキスト専用のため、画像生成は別バックエンドへ委譲する。Conductor が
# use_image_generation=true を出すと MoA を経由せず画像エージェントへバイパスする
# （ドキュメントの特殊ルーティング #1）。AUTOMATIC1111 stable-diffusion-webui
# (/sdapi/v1/txt2img・簡潔) を主とし、ComfyUI (/prompt+/history+/view) をフォールバック。
IMAGE_BACKEND = "auto"                       # "auto" | "a1111" | "comfyui" | "off"
A1111_URL = "http://127.0.0.1:7860"          # stable-diffusion-webui 既定ポート
COMFYUI_URL = "http://127.0.0.1:8188"        # ComfyUI 既定ポート
IMAGE_OUT_DIR = Path.home() / "fugu_images"  # 生成画像の保存先
IMAGE_STEPS = 30
IMAGE_WIDTH = 1024   # SDXL ネイティブ解像度（RTX 4060 8GB で実測OK。SD1.5 利用時は 512〜768 推奨）
IMAGE_HEIGHT = 1024
IMAGE_TIMEOUT = 600                          # 秒（ローカルGPUでの生成待ち上限）
IMAGE_TRANSLATE_PROMPT = True                # 日本語要求を英語SDプロンプトへ変換（qwen3:4b・MoA無効時のフォールバック）
COMFYUI_CKPT = ""                            # ComfyUI 用チェックポイント名（空ならサーバ既定を自動取得）
IMAGE_PROMPT_MOA = True                       # True: 画像プロンプトを LLM 群(proposers)で起草→統合。False: qwen3単独翻訳
IMAGE_PROMPT_PANEL = 2                        # プロンプト起草に使う proposer 数の上限（精度と速度の折衷）

# --- PowerPoint 生成（画像入りスライド）---
# make_pptx または --out X.pptx で、MoA が作った本文をスライド化し、内容連動で画像を埋め込む。
PPTX_OUT_DIR = Path.home() / "fugu_pptx"     # make_pptx 時の既定保存先
PPTX_MAX_SLIDES = 12                          # 生成する最大スライド数（タイトル除く）
PPTX_MAX_IMAGES = 4                           # 生成する最大画像枚数（タイトル画像含む・内容連動）
PPTX_MAX_BULLETS = 7                          # 1 スライドの最大箇条書き数

# --- 会話履歴 ---
# (user, assistant) ペアを保持し、古い交換を削除して num_ctx に収める。
# 14B モデル + num_ctx=8192 の場合、入力の余裕は ~2000 トークン（~6000文字）程度。
# 保守的に 4000 文字を上限とし、超えた分を先頭ペアから削除する。
MAX_HISTORY_CHARS = 4000
_HISTORY: list = []   # グローバル会話履歴

# --- VRAM 対策 ---
# 逐次実行でも Ollama は keep_alive でモデルを常駐させ続けるため、複数モデルだと
# 呼び出しごとにロード/アンロード（数GBのディスク読み込み）が多発して遅くなる。
# 最も効くのは「同時ロードは 1 体」を強制する環境変数（サーバ起動前に設定）:
#   Windows(PowerShell): $env:OLLAMA_MAX_LOADED_MODELS=1 ; ollama serve
#   Linux / mac:         OLLAMA_MAX_LOADED_MODELS=1 ollama serve
# 下の keep_alive をコード側から渡したい場合のみ文字列を設定（例 "0"=即アンロード, "5m"）。
# 既定 None は「渡さない」＝互換性リスクなし。
MODEL_KEEP_ALIVE = None

# 【重要】コンテキスト長。未指定だと Ollama がモデル最大(qwen3=262144 等)を確保しようとし、
# 8GB VRAM では KV キャッシュが破綻して runner がクラッシュする。実測で 8192 なら
# deepseek-r1:7b / gemma4:e2b-it-qat / qwen3:4b / phi4-mini いずれも VRAM に収まり安定。
# MoA のアグリゲータは「質問＋複数提案＋推論」で入力が伸びるので、これ以上は下げない方がよい。
MODEL_NUM_CTX = 8192

# 2026-07-02 実測: e2b-it-qat 置換後は 3 プロポーザー同時常駐が可能（qwen3 3.9 + phi4 3.7 +
# gemma-qat 1.9 = 9.5GB 表示でもクラッシュなし）だが、並列はむしろ遅い
# （warm・think=False で逐次 177.8s vs 並列 208.2s = x0.85）。GPU 演算が 1 基で奪い合いに
# なるため。よって False を維持する（安全性の問題ではなく速度メリットが無い）。
PARALLEL_PROPOSERS = False

PROPOSERS = []
AGGREGATOR = None
CONDUCTOR = None

# ==================================================
# Ollama ブートストラップ
# ==================================================


# ==================================================
# セッション永続化
# ==================================================

def load_history_file(path: Path = None) -> list:
    """JSON ファイルから会話履歴を読み込む。ファイルが無い/壊れている場合は空リストを返す。"""
    path = path or HISTORY_FILE
    if not path.exists():
        return []
    try:
        # 2026-07-23: このマシンのコンソールが cp932 である既知の落とし穴 #4 と
        # 同種の環境要因で、セッションファイルが Shift-JIS エディタでの開き直し
        # や部分的な破損により非UTF-8バイト列を含むことがある。その場合
        # encoding="utf-8" のみの read_text は UnicodeDecodeError を送出し、
        # 直下の except Exception: pass に捕まって [] を返す＝JSON構造自体は
        # 健全でも会話履歴全体を無条件に失う。iteration 47 で
        # _save_as_markdown/_save_as_text/_save_as_html の読み戻しに適用した
        # errors="replace" パターンをここにも適用し、読めないバイトだけを
        # 置換文字 (U+FFFD) に落として JSON をパース可能にし、読み取れる
        # エントリを保持する（精度優先: 履歴を丸ごと捨てない）。書き込み側の
        # save_history_file の encoding="utf-8" は変更しない。
        data = json.loads(path.read_text(encoding="utf-8", errors="replace"))
        if isinstance(data, list):
            # 最新 MAX_HISTORY_TURNS_SAVED 往復分のみ保持
            # 2026-07-23: iteration 66 で指摘され未修正だったバグを修正。
            # キーの有無だけでなく値の型も str であることを検証する。
            # {"role": "user", "content": null}（あるいは content が数値/list/dict）
            # のような壊れたセッションファイルのエントリはキーの存在チェックだけでは
            # 通過してしまい、_HISTORY に混入した後、次ターンで _trim_history()
            # (L949 付近) が sum(len(m["content"]) for m in history) を実行した際に
            # len(None) / len(123) で TypeError を送出しターンごとクラッシュする。
            # ここで文字列以外の値を持つエントリを弾き、「壊れたファイルは空/縮退
            # リストへ degrade する」という本関数の既存契約を値レベルまで拡張する。
            msgs = [m for m in data
                    if isinstance(m, dict)
                    and isinstance(m.get("role"), str)
                    and isinstance(m.get("content"), str)]
            return msgs[-(MAX_HISTORY_TURNS_SAVED * 2):]
    except Exception:
        pass
    return []


def save_history_file(history: list, path: Path = None, force: bool = False):
    """会話履歴を JSON ファイルに保存する。

    SESSION_SAVE=False (--no-history) 時は既定では何もしない（自動/受動的な
    永続化はユーザーのオプトアウトを尊重する）。ただし force=True を渡すと、
    'save <path>' のようなユーザー明示的なエクスポート操作として
    SESSION_SAVE の値に関わらず書き込みを行う。

    戻り値: 実際にファイルへ書き込めた場合のみ True。SESSION_SAVE=False かつ
    force=False で書き込みをスキップした場合、または書き込み中に例外が
    発生した場合は False（例外は送出しない）。
    """
    if not force and not SESSION_SAVE:
        return False
    path = path or HISTORY_FILE
    try:
        path.write_text(
            json.dumps(history[-(MAX_HISTORY_TURNS_SAVED * 2):],
                       ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        return True
    except Exception as e:
        print(f"   [履歴保存エラー: {e}]")
        return False


# ==================================================
# Slack 完了通知
# ==================================================

# Slack Incoming Webhook URL。環境変数 FUGU_SLACK_WEBHOOK に設定すると
# ask_fugu() 完了時（成功・失敗とも）に通知を送る。未設定なら何もしない。
# 1問数分〜十数分かかるため、離席していても完了が分かるようにする。
SLACK_WEBHOOK_URL = os.environ.get("FUGU_SLACK_WEBHOOK", "")
SLACK_NOTIFY_TIMEOUT = 10    # 秒。通知は本処理を止めない
SLACK_Q_PREVIEW = 200        # 通知に載せる質問の文字数上限
SLACK_A_PREVIEW = 500        # 通知に載せる回答の文字数上限


def _slack_truncate(text: str, limit: int) -> str:
    text = (text or "").strip()
    return text if len(text) <= limit else text[:limit] + "…"


def notify_slack(question: str, answer: str, elapsed: float):
    """完了通知を Slack Incoming Webhook へ送る。失敗しても本処理には影響させない。"""
    if not SLACK_WEBHOOK_URL:
        return
    ok = not (answer or "").startswith("__ERROR__")
    icon = ":white_check_mark:" if ok else ":x:"
    head = f"Fugu {'完了' if ok else '失敗'} ({elapsed} 秒)"
    text = (
        f"{icon} *{head}*\n"
        f"*Q:* {_slack_truncate(question, SLACK_Q_PREVIEW)}\n"
        f"*A:* {_slack_truncate(answer, SLACK_A_PREVIEW)}"
    )
    req = urllib.request.Request(
        SLACK_WEBHOOK_URL,
        data=json.dumps({"text": text}, ensure_ascii=False).encode("utf-8"),
        headers={"Content-Type": "application/json"},
        method="POST",
    )
    try:
        urllib.request.urlopen(req, timeout=SLACK_NOTIFY_TIMEOUT)
        print("   [Slack 通知を送信しました]")
    except Exception as e:
        print(f"   [Slack 通知エラー: {e}]")


# ==================================================
# Web 検索
# ==================================================

def _resolve_ddgs_class():
    """DDGS クラスを解決するだけ（クエリは実行しない＝import 段階のみ）。
    ddgs（後継、pip install ddgs）→ duckduckgo_search（旧名）の順に import を試み、
    両方とも未インストールの場合のみ ImportError を送出する。"""
    try:
        from ddgs import DDGS  # 後継パッケージ（pip install ddgs）
    except ImportError:
        from duckduckgo_search import DDGS  # 旧名（非推奨）
    return DDGS


def _ddg_full(query: str, max_results: int) -> list:
    """ddgs パッケージ（旧 duckduckgo_search）を使ってフル検索結果を返す。"""
    DDGS = _resolve_ddgs_class()
    results = []
    with DDGS() as ddgs:
        for r in ddgs.text(query, max_results=max_results):
            # 2026-07-25: _ddg_full は ddgs/duckduckgo_search がインストール済みの
            # 通常運用時に必ず通るプライマリ検索経路であるにもかかわらず、フォール
            # バック側の _ddg_instant（iter103/111/112/113/138/139で「壊れた外部
            # ペイロードの1件が全体を握り潰さない」よう isinstance ガードで段階的に
            # 固められてきた）とは非対称に、r.get("body")/r.get("title")/
            # r.get("href") を型チェックなしで直接呼んでいた。ddgs.text() が yield
            # する行が dict でない場合、r.get(...) は即座に AttributeError を送出し、
            # _search_raw の `except Exception: return []` がそれをクエリ全体の空
            # リストへ丸ごと潰すため、同じクエリ内の他の正常な行まで全て失われる
            # （1件の壊れた行が全件を道連れにする、iter113/iter138と同じ問題）。
            # ここでも同じ作法に倣い、行(r)が dict でなければ例外を出さず continue
            # して後続の有効な行を救済する。また title/body/href が str でない場合
            # （list/dict/int 等）は、そのrepr文字列をモデルの検索コンテキストへ
            # ノイズとして混入させないため str() 変換はせず空文字列として扱う
            # （_ddg_instant の末端型ガード、2026-07-25追加分と同じ考え方）。
            if not isinstance(r, dict):
                continue
            title = r.get("title", "")
            if not isinstance(title, str):
                title = ""
            body = r.get("body")
            if not isinstance(body, str):
                body = ""
            href = r.get("href", "")
            if not isinstance(href, str):
                href = ""
            snippet = (body or "")[:WEB_SEARCH_SNIPPET_CHARS]
            results.append(f"[{title}]\n{snippet}\nSource: {href}")
    return results


def _ddg_instant(query: str, max_results: int) -> list:
    """DuckDuckGo Instant Answer API (urllib のみ、フォールバック)。"""
    url = ("https://api.duckduckgo.com/?" +
           urllib.parse.urlencode({"q": query, "format": "json",
                                   "no_redirect": "1", "no_html": "1"}))
    req = urllib.request.Request(url, headers={"User-Agent": "fugu-local/1.0"})
    try:
        with urllib.request.urlopen(req, timeout=WEB_SEARCH_TIMEOUT) as r:
            data = json.loads(r.read().decode("utf-8"))
    except Exception:
        return []
    # 2026-07-24: DuckDuckGo Instant Answer API はクエリ次第で JSON のトップレベルが
    # dict ではない（配列・文字列・数値など、いずれも json.loads は例外にしない）ことが
    # あり、また "RelatedTopics" キーが存在していても値が null のことがある。
    # dict.get(key, default) の default は「キーが存在しない」場合にのみ使われる仕様の
    # ため、{"RelatedTopics": null} では data.get("RelatedTopics", []) は [] ではなく
    # None を返し、続く `for t in None` が TypeError になる。この関数は
    # _search_raw() の L544 相当 `except ImportError: return _ddg_instant(...)` の
    # 内側から try に包まれず直接呼ばれているため、その TypeError は _search_raw の
    # 「失敗時は空リスト（呼び出し側を止めない）」という契約（イテレーション75で確立）
    # をすり抜けて呼び出し元（web_search/research_search/build_context）まで伝播し、
    # ターン全体を落としてしまう。ここで dict 以外のトップレベルは即座に [] を返し、
    # RelatedTopics も非 list（null・dict・str 等）なら [] に丸めることで、
    # イテレーション85/89で整備した整形済みペイロードの挙動は一切変えずに、
    # 壊れたペイロードでも例外を出さないようにする。
    if not isinstance(data, dict):
        return []
    results = []
    # 2026-07-24: _ddg_full は (r.get("body") or "")[:WEB_SEARCH_SNIPPET_CHARS] で
    # 各スニペットを必ず切り詰めているが、この Instant Answer フォールバックは
    # Abstract / RelatedTopics の Text を無制限に追加していた。DuckDuckGo の
    # Abstract は数KBに及ぶことがあり、research_search の body 組み立てループは
    # 先頭アイテムが SEARCH_CONTEXT_CHARS を超えると
    # `if not body: body = item[:SEARCH_CONTEXT_CHARS]` で break するため、巨大な
    # Abstract 1件が他の収集済み事実を全て握り潰してしまう。_ddg_full と同じ上限を
    # 適用して対称にする（Source 行は切り詰めない）。
    # 2026-07-25: イテレーション103はコンテナ形状（トップレベル data / RelatedTopics）
    # のみを isinstance で丸め、末端（Abstract 文字列本体・直接トピックの Text・
    # ネストされた Topics[].Text）はここに来るまで一度も型を見ておらず、
    # `if data.get("Abstract"):` のような真偽値チェックの直後でいきなり
    # `[:WEB_SEARCH_SNIPPET_CHARS]` によるスライスを行っていた。DDG は壊れた
    # ペイロードで Abstract/Text に int・float・bool・list・dict を返すことがあり、
    # int/float/bool は「添字操作できません」の TypeError、dict はスライス構文上
    # 例外にはならないが該当キーの意味を持たない値になり、list は例外を出さず
    # スライスできてしまうが results にリストオブジェクトが混入し、後段の
    # research_search の `re.search(r'Source: ...')` や `'\n\n'.join(...)` を
    # 壊す。コンテナ側の非list/非dict丸め（イテレーション103・111・112・113の
    # 「isinstance で判定し例外を出さず読み飛ばす」系列と同じ考え方）と同様に、
    # ここでも末端が str でなければスライスせずそのリーフだけを読み飛ばす
    # （str への str() 変換はしない。list/dict の repr をそのままノイズとして
    # モデルの検索コンテキストに混入させないため）。_search_raw() は
    # `except ImportError:` の内側から try に包まずこの関数を直接呼んでいる
    # （イテレーション103のコメント参照）ため、ここで拾わない例外は
    # 「失敗時は空リスト」という never-raise 契約をすり抜けてターン全体を落とす。
    if data.get("Abstract") and isinstance(data["Abstract"], str):
        abstract = data["Abstract"][:WEB_SEARCH_SNIPPET_CHARS]
        results.append(f"[{data.get('AbstractTitle', '')}]\n{abstract}\n"
                       f"Source: {data.get('AbstractURL', '')}")
    rel = data.get("RelatedTopics")
    if not isinstance(rel, list):
        rel = []
    for t in rel:
        if isinstance(t, dict) and t.get("Text"):
            # トップレベルに Text がある「直接トピック」はこちらを優先し、
            # 仮に Topics も併存していても展開しない（下の分岐と二重追加しない
            # ための優先順位。直接トピックの既存挙動を保つ）。
            # Text が truthy な非文字列（例: int/float/bool/list/dict）の場合は
            # スライスせずこのリーフだけを読み飛ばす（上のコメント参照）。
            # ブランチの選択自体（direct-Text 優先）は変えない。
            if isinstance(t["Text"], str):
                results.append(t["Text"][:WEB_SEARCH_SNIPPET_CHARS])
        elif isinstance(t, dict) and isinstance(t.get("Topics"), list):
            # 2026-07-24: RelatedTopics には、トップレベル Text を持つ直接トピックと
            # {"Name": "カテゴリ名", "Topics": [...]} 形式の「グループ化トピック」が
            # 混在する。従来は isinstance(t, dict) and t.get("Text") しか見ておらず、
            # グループ化エントリ配下にネストされた事実（Topics 配列の各要素）が無条件
            # で握り潰されていた（イテレーション85がこの関数を整備した際に残った
            # 既知の欠落）。DDG のネストは1階層のみなので、ここでも1階層だけを
            # フラット化する（再帰はしない。壊れた形状は例外を出さず読み飛ばす）。
            for nested in t["Topics"]:
                # 2026-07-25: nested["Text"] も同様に truthy 非文字列を読み飛ばす
                # （上の Abstract/direct-Text と同じ理由）。
                if (isinstance(nested, dict) and nested.get("Text")
                        and isinstance(nested["Text"], str)):
                    results.append(nested["Text"][:WEB_SEARCH_SNIPPET_CHARS])
                if len(results) >= max_results:
                    break
        if len(results) >= max_results:
            break
    return results[:max_results]


def _search_raw(query: str, max_results: int = None) -> list:
    """1 クエリ分の検索結果をリストで返す。失敗時は空リスト（呼び出し側を止めない）。"""
    max_results = max_results or WEB_SEARCH_MAX_RESULTS
    # 2026-07-23: 「ライブラリ解決(import)」と「クエリ実行(DDGS().text())」を別の try
    # に分離する。旧実装は _ddg_full 内部の import 段階とクエリ実行段階を1つの try に
    # まとめて except ImportError で受けていたため、ddgs/duckduckgo_search が正しく
    # インストール済みでも、それが内部で遅延 import する primp/lxml 等のバックエンドが
    # 実行時に ImportError を送出するケースまで「ライブラリ未インストール」と誤認して
    # いた。その結果、誤った pip install 警告を出した上に Instant Answer フォールバック
    # （下の警告コメントの通り、事実系クエリで古い知識を返す事故の温床）に倒れてしまう。
    # ここでは _resolve_ddgs_class() の import 失敗（＝本当に未インストール）だけを
    # Instant Answer フォールバックの対象とし、解決後の実行時エラーは ImportError で
    # あっても他の実行時例外と同じ扱い（[] を返すのみ・Instant Answer には倒さない）にする。
    try:
        _resolve_ddgs_class()
    except ImportError:
        # Instant Answer API は事実系クエリでほぼ空を返す。無警告だと「検索したのに
        # 古い知識で回答」する事故になる（実測 2026-07-06: 最新GPUで1世代前を回答）。
        print("   [警告: ddgs 未インストールのため Instant Answer フォールバック中。"
              "pip install ddgs でフル検索が有効になります]")
        return _ddg_instant(query, max_results)
    except Exception as e:
        # 2026-07-26: _resolve_ddgs_class() は `from ddgs import DDGS` /
        # `from duckduckgo_search import DDGS` を実行するため、そのパッケージの
        # トップレベル __init__ 由来で ImportError 以外（壊れた/非推奨パッケージの
        # RuntimeError・OSError、Python バージョンガード、壊れたネイティブ依存等）
        # も送出しうる。ここを except ImportError だけで受けていると、この
        # 非ImportError が _search_raw の外（research_search → build_context →
        # ask_fugu）まで無捕捉で伝播し、本関数のドキュメント契約「失敗時は空リスト
        # （呼び出し側を止めない）」（上記docstring）と、このコードベースが
        # イテレーション75/103/111/138で繰り返し守ってきた「呼び出し側を絶対に
        # 落とさない」不変条件に違反してターン全体を失う。イテレーション83の判断
        # （解決/実行時の失敗を、事実系クエリで古い知識を返しがちな Instant Answer
        # フォールバックに倒すと本当の失敗が古い回答でマスクされる）に従い、ここも
        # Instant Answer へは倒さず [] を返すだけにする。except Exception なので
        # KeyboardInterrupt/SystemExit は従来通り素通りする。
        print(f"   [Web検索エラー(解決段階): {e}]")
        return []

    try:
        return _ddg_full(query, max_results)
    except Exception as e:
        print(f"   [Web検索エラー: {e}]")
        return []


def web_search(query: str, max_results: int = None) -> str:
    """Web 検索 1 回分をフォーマット済み文字列で返す（後方互換用の単発検索）。"""
    results = _search_raw(query, max_results)
    if not results:
        return ""
    return "## Web Search Results (DuckDuckGo)\n" + "\n\n".join(results)


# 十分性判定（Conductor と同じ think=False + スキーマ拘束パターン）
RESEARCH_SYS = (
    "You are a research assistant judging web search results. "
    "Given a user question and accumulated search results, decide whether the results "
    "contain enough SPECIFIC and up-to-date facts (exact product names, architecture "
    "names, versions, dates, numbers) to answer the question accurately without "
    "guessing. Snippets that merely mention the topic are NOT sufficient. "
    "If not sufficient, state what is missing and give up to 3 NEW search queries "
    "targeting the missing facts. Use different keywords than previous queries; "
    "include English queries for technical topics. Return ONLY JSON."
)

RESEARCH_SCHEMA = {
    "type": "object",
    "properties": {
        "sufficient": {"type": "boolean"},
        "missing": {"type": "string"},
        "queries": {"type": "array", "items": {"type": "string"}},
    },
    "required": ["sufficient", "queries"],
}


def research_search(question: str) -> str:
    """十分な事実が集まるまで検索を反復するリサーチループ。
    件数固定の単発検索と違い、Conductor が不足を判定して追加クエリで掘り下げる。"""
    results = []      # フォーマット済み結果（重複排除済み）
    seen = set()      # Source URL による重複排除
    tried = set()     # 実行済みクエリ（同一クエリの再実行を防ぐ）
    queries = [question]

    for rnd in range(1, SEARCH_MAX_ROUNDS + 1):
        for q in queries:
            q = str(q).strip()
            if not q or q.lower() in tried:
                continue
            tried.add(q.lower())
            print(f"   [Web検索 R{rnd}: {q[:60]}]")
            for item in _search_raw(q):
                m = re.search(r"Source: (\S+)", item)
                key = m.group(1) if m else item[:80]
                if key in seen:
                    continue
                seen.add(key)
                results.append(item)

        if rnd == SEARCH_MAX_ROUNDS:
            break

        # 十分性判定（qwen3:4b、~15s）。判定不能時は安全側＝そこで打ち切り（結果は使う）
        joined = "\n\n".join(results)[:SEARCH_CONTEXT_CHARS]
        raw = ask(
            CONDUCTOR,
            [{"role": "system", "content": RESEARCH_SYS},
             {"role": "user", "content": (
                 f"Question:\n{question}\n\n"
                 f"Previous queries: {sorted(tried)}\n\n"
                 f"Search results so far:\n{joined or '(no results)'}\n\n"
                 "Return ONLY the JSON judgement.")}],
            CONDUCTOR_TEMP,
            think=False, fmt=RESEARCH_SCHEMA,
            num_predict=NUM_PREDICT_JUDGE, label="research",
        )
        j = extract_json(raw)
        if not isinstance(j, dict) or j.get("sufficient"):
            break
        missing = str(j.get("missing", ""))[:120]
        # 2026-07-24: 従来の `j.get("queries") or []` は "queries" が falsy
        # (None/[]/""/0 等)の場合のみ [] に丸めるトリックで、int/float/bool のような
        # 真値だが非反復可能(non-iterable)な値はそのまま通してしまっていた。その場合は
        # 直後のリスト内包表記 `for x in ...` 自体が TypeError を送出する。さらに
        # truthy な str/dict はそのまま通ると「1文字ずつ」「キーごと」に分解された
        # 意味のない検索クエリを発行してしまう。research_search は build_context
        # (L1285付近) から ask_fugu (L3519付近) まで無防備(try/exceptなし)に
        # 呼ばれているため、この TypeError は Conductor の計画や既に完了した検索
        # ラウンドの結果ごとターン全体を落とす。イテレーション103の _ddg_instant
        # (非list RelatedTopics)・イテレーション111の plan_pptx_images (非list
        # images) と同じ isinstance 判定で丸める方式に倣い、"queries" が非list
        # (文字列/None/int/float/bool/dict 等)であれば真偽・型に関わらず必ず []
        # に倒す(`or []` の truthiness トリックには戻さない)。
        raw_queries = j.get("queries")
        if not isinstance(raw_queries, list):
            raw_queries = []
        queries = [str(x) for x in raw_queries if str(x).strip()][:3]
        if not queries:
            break
        print(f"   [リサーチ継続 R{rnd + 1}: 不足={missing or '(詳細なし)'}]")

    if not results:
        return ""
    # 注入上限で切る（結果の区切り単位）。
    # 2026-07-25: 従来はここで最初に上限超過に当たった結果で break していたため、
    # それ以降の結果（後続ラウンドで Conductor が指定した不足事実を埋める、小さく
    # 有益な結果であっても）がまとめて捨てられていた。先頭1件だけで上限超過して
    # body が空になる問題はイテレーション38の特性テスト(H)で発見・固定され、
    # イテレーション39で「空のまま返さず先頭結果を切り詰める」フォールバックとして
    # 対処済みだが、「途中の1件が大きくて超過→それ以降が全部捨てられる」という
    # 根本原因は残っていた。ここを break から continue（スキップして走査継続）に
    # 変え、順序を保ったまま残り予算に収まる結果だけを拾う greedy first-fit
    # パッキングにする。精度優先（精度優先・時間は気にしない）のため、同一の
    # SEARCH_CONTEXT_CHARS 予算内で実際に注入できる検索事実を最大化する。
    # 予算チェック式 `len(body) + len(item) > SEARCH_CONTEXT_CHARS` は変更前と同じ
    # （区切り"\n\n"の2文字はチェックに含めない）。この2文字は次の候補のチェック時に
    # body の一部として自然に繰り込まれ、最後に採用された結果の末尾の"\n\n"は最終的な
    # rstrip() で必ず除去されるため、最終的な body の長さが SEARCH_CONTEXT_CHARS を
    # 超えることはない（off-by-two回避のため、このチェック式自体は変更しないこと）。
    body = ""
    for item in results:
        if len(body) + len(item) > SEARCH_CONTEXT_CHARS:
            continue
        body += item + "\n\n"
    if not body:
        # 全結果が単独でも上限を超える場合の最終フォールバック（イテレーション38/39と
        # 同じ契約）：空のまま返さず、先頭結果を上限まで切り詰めてでも必ず注入する。
        body = results[0][:SEARCH_CONTEXT_CHARS] + "\n\n"
    header = (
        f"## Web Search Results (取得日: {time.strftime('%Y-%m-%d')})\n"
        "重要: 以下はあなたの学習データより新しい一次情報である。学習知識と矛盾する場合は"
        "必ず検索結果を優先すること。検索結果に書かれていない具体的事実"
        "（型番・アーキテクチャ名・日付・数値など）は推測で断定しないこと。\n\n"
    )
    return header + body.rstrip()


# ==================================================
# RAG（ローカル文書検索）
# ==================================================

# ==================================================
# ユニバーサルファイル読み込み
# ==================================================

def _read_pdf(path: Path) -> str:
    """PDF からテキストを抽出。pdfplumber → pypdf → PyPDF2 の順で試行。
    2026-07-23: 各ブロックは従来 except ImportError のみで、下位ライブラリへの
    フォールスルーは「上位ライブラリが未インストール」の場合にしか起きなかった。
    pdfplumber 等が実際にはインストール済みでも、暗号化/破損/パーサ固有のエッジケースで
    実行時に例外を送出するPDFに対しては例外がそのまま _read_pdf の外へ伝播し、
    read_file_text（iter53）の呼び出し側ガードがそれを握りつぶして ""
    を返してしまい、pypdf/PyPDF2 なら救えたはずのテキストがPDF丸ごとRAG/--file
    コンテキストから失われていた（精度優先の方針に反する）。全リーダー関数の
    書き換えを試みた iter51 は行き詰まったスタック案件だが、これは「呼び出し側で
    クラッシュさせない」話ではなく「下位ライブラリへフォールスルーしてテキストを
    救済する」話であり別角度の問題。iter41-44 の graceful degradation の方針に合わせ、
    ImportError に加えて except Exception（bare except にはしない。
    KeyboardInterrupt/SystemExit は握りつぶさず伝播させる）でも次候補へ
    フォールスルーさせ、ライブラリが import 自体には成功したのに失敗した場合は
    cp932セーフな警告（ファイル名+例外型のみ、絵文字等は使わない）で可視化する。

    2026-07-25: 上記iter83の修正は自身のコメントが明示する通り「実行時例外」の
    場合にしか適用されず、pdfplumberが例外を出さずに空文字列/空白のみを返す
    ケース（フォントエンコーディング等に起因する既知の仕様上の癖で、実際には
    読める有効なPDFでも起こりうる）は未対応のまま残っていた。iter41-44の
    graceful degradation方針とiter83の例外フォールスルーの系譜を継ぎ、ここでは
    「そのライブラリでの抽出処理自体は完走したが結果が空だった」場合も次候補を
    試すようにする（＝iter83が閉じ損ねた「例外ではなく結果が空」という穴を塞ぐ）。
    どの層も import すらできなかった場合（従来のpip installが必要なケース）と、
    1層以上がimportに成功し完走したが全て空だった場合（スキャンPDF等、本当に
    中身が無いケース）とを区別し、前者のみ従来のpip install通知を返す。"""
    # 少なくとも1層がimportに成功し抽出処理自体は完走したか（例外で終わった層は
    # 数えない）。全層がこれを満たさなければ「未インストール」、満たすが全層空
    # だった場合は「読める中身が無いPDF」として扱いを区別する(2026-07-25)。
    _pdf_any_tier_completed = False
    # pdfplumber (最高品質)
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            pages = [p.extract_text() or "" for p in pdf.pages]
        text = "\n\n".join(pages)
        _pdf_any_tier_completed = True
        if text.strip():
            return text
    except ImportError:
        pass
    except Exception as exc:
        print(f"[_read_pdf] pdfplumber抽出失敗のため次候補へフォールバック: {path.name} ({type(exc).__name__})")
    # pypdf (軽量・新しい)
    try:
        import pypdf
        with open(path, "rb") as f:
            r = pypdf.PdfReader(f)
            text = "\n\n".join(p.extract_text() or "" for p in r.pages)
        _pdf_any_tier_completed = True
        if text.strip():
            return text
    except ImportError:
        pass
    except Exception as exc:
        print(f"[_read_pdf] pypdf抽出失敗のため次候補へフォールバック: {path.name} ({type(exc).__name__})")
    # PyPDF2 (旧名称)
    try:
        import PyPDF2
        with open(path, "rb") as f:
            r = PyPDF2.PdfReader(f)
            text = "\n\n".join(p.extract_text() or "" for p in r.pages)
        _pdf_any_tier_completed = True
        if text.strip():
            return text
    except ImportError:
        pass
    except Exception as exc:
        print(f"[_read_pdf] PyPDF2抽出失敗のため次候補へフォールバック: {path.name} ({type(exc).__name__})")
    # 2026-07-25: 1層以上が完走していれば(全て空でも)「未インストール」ではないので
    # pip install通知は返さず、空文字列(スキャンPDF等、本当に中身が無いケース)を返す。
    if _pdf_any_tier_completed:
        return ""
    return f"[PDF: {path.name} — テキスト抽出には pdfplumber or pypdf が必要: pip install pdfplumber]"


def _read_docx(path: Path) -> str:
    """Word (.docx) からテキストを抽出。"""
    try:
        import docx

        # 2026-07-25: ネスト表の再帰の深さに上限を設け、病的/破損した文書
        # (自己参照的・異常に深いネスト等)で再帰が終わらなくなることを防ぐ。
        _DOCX_NESTED_TABLE_MAX_DEPTH = 6

        def _table_rows_text(tbl, _depth=0):
            # 2026-07-24 (iter91): python-docxのrow.cellsは表のグリッド座標の数だけ
            # _Cellを返すが、水平方向(gridSpan)にマージされたセルは全ての被マージ
            # 座標で同一の<w:tc>要素を共有する。そのためc.textはマージされた
            # テキストを座標の数だけ重複して返し、例えば2列にまたがるヘッダーセルは
            # 'Header\tHeader'のように重複抽出されてしまう(単純なタブ結合のみだった
            # 従来コードにはこの重複排除が無かった)。ここでは行内で既出の<w:tc>を
            # id(getattr(c, '_tc', None))で追跡し、同一セルの2回目以降の出現を
            # スキップすることで、マージされたセルのテキストを行につき1回だけ
            # 出力する。_tcが取得できない場合(将来のpython-docx実装変化等)は
            # 安全側に倒し重複排除を行わず従来通りの挙動にフォールバックする。
            # なお本修正は行内(水平/gridSpan)の重複排除のみに限定しており、
            # 垂直マージ(行をまたぐvMerge)の重複排除は意図的にスコープ外としている
            # (将来のフォローアップ課題)。
            # 同種の「python-docx/python-pptxの属性欠落・仕様により文字列を
            # 取りこぼす/重複する」系統の修正としてiter87(_read_pptxの表/グループ
            # 抽出)の直系であり、iter82では非マージの単純な表しかテストされて
            # いなかった穴を埋めるもの。iter93(下記)でも本ロジックはbyte-for-byte
            # 温存しており、順序修正後の本文経路/従来フォールバック経路の両方から
            # 同一関数として呼ばれる。
            #
            # 2026-07-25: python-docxの_Cell.textは、そのセル直下の<w:p>段落だけを
            # 連結して返し、セル内にネストされた<w:tbl>(表の中の表)は一切含めない。
            # そのためセルの中にさらに表が入っている文書(よくある「表内表」構成)は、
            # そのネスト表の内容が丸ごと_read_docxの出力から――ひいてはRAG検索
            # (_load_rag_chunks)や--fileで各proposerに渡る全文コンテキストからも――
            # 無言で欠落していた。iter91・iter93はどちらもこれを認識しつつ
            # 「ネストした表の再帰抽出は意図的にスコープ外(将来のフォローアップ)」と
            # 明記して先送りしていたもので、本修正がそのフォローアップを完了させる。
            # python-docxはBlockItemContainerとして_Cell.tablesでネスト表の一覧を
            # 公開しているため、セル自身のテキストを出力した直後に、そのセルが
            # 実際に出力対象だった場合(=水平マージの初出セルである場合)に限り
            # ネスト表へ追加的に再帰する。水平マージでスキップされた被マージセルの
            # 座標では再帰しない(ネスト表の二重処理を避けるため)。ネスト表の行は
            # 「そのネスト表を含む行」自身のタブ結合済みテキストの直後にまとめて
            # 追加し、行の途中に混在させない(フラットな逐次追記という意図的な配置)。
            # getattr(cell, "tables", [])で取得するため、将来python-docxの属性が
            # 欠落/改名されても例外化せず現状の(ネスト表を無視する)挙動に安全側で
            # 縮退する。vMerge(垂直マージ)の重複排除とpptx(セルが表をネストできない)
            # は引き続き明示的にスコープ外。
            rows_text = []
            for row in tbl.rows:
                seen_tc_ids = set()
                cells_text = []
                nested_rows = []
                for c in row.cells:
                    tc = getattr(c, "_tc", None)
                    if tc is not None:
                        tc_id = id(tc)
                        if tc_id in seen_tc_ids:
                            continue
                        seen_tc_ids.add(tc_id)
                    cells_text.append(c.text)
                    if _depth < _DOCX_NESTED_TABLE_MAX_DEPTH:
                        for nested_tbl in getattr(c, "tables", []):
                            nested_rows.extend(
                                _table_rows_text(nested_tbl, _depth + 1)
                            )
                rows_text.append("\t".join(cells_text))
                rows_text.extend(nested_rows)
            return rows_text

        # 2026-07-24 (iter123): docx.Document() 自体は import docx の成功/失敗とは
        # 別に、破損したzip(BadZipFile)やレガシーバイナリ.docを.docx拡張子のまま
        # 開いた場合(PackageNotFoundError)等、実行時例外を送出しうる。従来はここが
        # except ImportError のみで守られていたため、そうした例外は_read_docxの
        # 外へそのまま伝播し、read_file_text（iter53）の外側ガードに握りつぶされて
        # "" になるだけで、なぜ失敗したかの情報が失われていた。iter83(_read_pdf)/
        # iter84(_read_excel)がpdf/excelで行った「実行時例外もopen/parse範囲だけで
        # 狭く捕捉し、可視化されたnotice文字列に変換する」修正をここにも適用する。
        # bare exceptにはせずException限定でKeyboardInterrupt/SystemExitは伝播させる。
        # 直下の except ImportError: pass は既存のバイト単位のまま温存し、
        # 新しいnoticeは既存の "python-docx が必要" notice とは別の文言にして、
        # _is_lib_missing_notice の「1行・[...]・pip install を含む」という構造的
        # 判定には意図的にヒットしない（"pip install" を含まない）ようにする。
        # ライブラリはインストール済みで読めているのに"未インストール"と誤判定され
        # RAG側で扱いを誤らないようにするための意図的な区別。
        try:
            doc = docx.Document(str(path))
        except Exception as exc:
            print(f"[_read_docx] Document()での読み込みに失敗: {path.name} ({type(exc).__name__})")
            return f"[DOCX: {path.name} — 読み込みエラー ({type(exc).__name__}): 破損しているか非対応形式の可能性があります]"
        parts = []

        # 2026-07-24 (iter93): 従来はdoc.paragraphsを全件処理してからdoc.tablesを
        # 全件処理する2パス構成だった。本文中の任意の位置(例:導入文 -> データ表 ->
        # 結論文)に表が挟まる文書では、これが「導入文・結論文」に続けて「表」が
        # 文末にまとめて出力される形になり、本文の論理的な読み順を破壊していた。
        # 表とそれを説明する地の文が同じチャンクに絶対に同居できなくなるため、
        # RAG検索(_load_rag_chunks)にも--file全文コンテキストにも直接の精度劣化
        # (精度優先・時間は気にしない、の方針に反する)として効いてくる。
        # python-docxのdoc.paragraphs/doc.tablesは本文中の出現順序を保持しないため
        # (どちらも本文内の該当要素だけを別々に集めたビュー)、ここでは
        # doc.element.body（本文の直下の子要素、<w:p>と<w:tbl>が出現順に並ぶ）を
        # 直接歩き、python-docx公式に知られる CT_P/CT_Tbl の isinstance判定で
        # 種別を振り分けて、出現順のまま段落・表を交互に出力する。これは
        # iter87(_read_pptxの表/グループ抽出)・iter91(本関数の水平マージ重複排除)と
        # 同系統の「ライブラリの属性/仕様の落とし穴で本文が欠落・変形する」問題の
        # 修正であり、iter82が固定した「非マージの単純な表」テストや上記iter91の
        # マージ表テストは、どちらも表を段落の後ろに追加するフィクスチャのため
        # 本文順序が変わらず影響を受けない。doc.paragraphs/doc.tables同様、本関数は
        # 表セル内にネストした段落・表までは再帰しない(トップレベルのみ、iter91が
        # 明示的にスコープ外とした挙動を継続)。
        try:
            from docx.oxml.table import CT_Tbl
            from docx.oxml.text.paragraph import CT_P
            from docx.table import Table
            from docx.text.paragraph import Paragraph

            for child in doc.element.body.iterchildren():
                if isinstance(child, CT_P):
                    para = Paragraph(child, doc)
                    if para.text.strip():
                        parts.append(para.text)
                elif isinstance(child, CT_Tbl):
                    tbl = Table(child, doc)
                    parts.extend(_table_rows_text(tbl))
        except Exception as exc:
            # python-docxの内部実装(CT_P/CT_Tbl/doc.element.body等)が将来の版で
            # 変わり本文順の走査自体が失敗しても、read_file_text/_load_rag_chunksが
            # 期待するgraceful degradation契約(iter42/53)を守るため、クラッシュ
            # させず段落->表の従来の2パス挙動へ安全側にフォールバックする。
            print(f"[_read_docx] 本文順序抽出に失敗したため段落->表の従来順にフォールバック: "
                  f"{path.name} ({type(exc).__name__})")
            parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text)
            for tbl in doc.tables:
                parts.extend(_table_rows_text(tbl))

        return "\n".join(parts)
    except ImportError:
        pass
    return f"[DOCX: {path.name} — python-docx が必要: pip install python-docx]"


def _read_excel(path: Path) -> str:
    """Excel (.xlsx/.xls) を CSV ライクなテキストに変換。openpyxl -> pandas/xlrd の順で試行。
    2026-07-24 (iter84): read_file_text は .xlsx と .xls の両方をこの関数へディスパッチ
    するが、従来は各ブロックが except ImportError のみで、下位ライブラリへの
    フォールスルーは「上位ライブラリが未インストール」の場合にしか起きなかった。
    openpyxl は legacy な .xls（バイナリ形式）を一切読めず、
    openpyxl.load_workbook() は openpyxl.utils.exceptions.InvalidFileException という
    実行時例外（ImportError ではない）を送出する。そのため .xls は
    read_file_text で公式にディスパッチされる入力形式であるにもかかわらず、
    pandas+xlrd がインストール済みで読めるはずでも、フォールスルーが起きず例外が
    そのまま外へ伝播し、read_file_text（iter53）/ _load_rag_chunks（iter42）の
    呼び出し側ガードがそれを握りつぶしてスプレッドシート丸ごとがRAG/--file
    コンテキストから静かに失われていた（精度優先の方針に反する）。同じ隙間は
    破損/openpyxl未対応の .xlsx が pandas 側で再試行されないことにも当てはまる。
    これは iter83 が _read_pdf に対して行った修正（下位ライブラリへフォールスルーして
    テキストを救済する）の直系の姉妹修正であり、iter41-44 の graceful degradation
    方針の延長でもある。ImportError に加えて except Exception（bare except には
    しない。KeyboardInterrupt/SystemExit は握りつぶさず伝播させる）でも次候補へ
    フォールスルーさせ、ライブラリが import 自体には成功したのに失敗した場合は
    cp932セーフな警告（ファイル名+例外型のみ、絵文字等は使わない）で可視化する。"""
    # openpyxl (高速・.xlsx専用。.xlsは読めずInvalidFileExceptionを送出する)
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        # 2026-07-24: read_only=True のワークブックは内部で循環参照を持ち、単純な
        # 参照カウントだけでは解放されない（iter82のテストがこれを踏み、後始末で
        # gc.collect() を挟まないと WinError 32 でtempディレクトリの削除に失敗した
        # のが直接の証拠）。close() を呼ばないまま関数を抜けるとzipのファイル
        # ハンドルがGCまで開いたままになり、Windowsでは同じファイルの後続の
        # 読み書き/移動をブロックし、_load_rag_chunksが多数の.xlsxを読む場面では
        # ハンドルリークにもなる。抽出（parts構築）が終わった後にfinallyでclose()
        # し、close()自体が失敗しても既に取れているテキストは握りつぶさず返す。
        # 本体の except ImportError / except Exception によるiter84の
        # フォールスルー動作（.xls等でopenpyxlが読めない場合にpandasへ進む挙動）は
        # 変更しない。
        try:
            parts = []
            for ws in wb.worksheets:
                parts.append(f"[Sheet: {ws.title}]")
                for row in ws.iter_rows(values_only=True):
                    row_str = "\t".join("" if v is None else str(v) for v in row)
                    if row_str.strip():
                        parts.append(row_str)
            result = "\n".join(parts)
        finally:
            try:
                wb.close()
            except Exception:
                # close失敗で成功した抽出結果をフォールバック扱いにしない
                pass
        return result
    except ImportError:
        pass
    except Exception as exc:
        print(f"[_read_excel] openpyxl抽出失敗のため次候補へフォールバック: {path.name} ({type(exc).__name__})")
    # pandas/xlrd (.xls を含む幅広い形式を読める)
    try:
        import pandas as pd
        xl = pd.ExcelFile(str(path))
        # 2026-07-24: pd.ExcelFile も内部でzipハンドルを開いたままにするため、
        # 上のopenpyxl分岐と同じ理由（iter82のgc.collect()回避策/iter84の
        # フォールスルー方針）でclose()を明示する。
        try:
            parts = []
            for sheet in xl.sheet_names:
                df = xl.parse(sheet)
                parts.append(f"[Sheet: {sheet}]\n{df.to_csv(index=False)}")
            result = "\n\n".join(parts)
        finally:
            try:
                xl.close()
            except Exception:
                pass
        return result
    except ImportError:
        pass
    except Exception as exc:
        print(f"[_read_excel] pandas抽出失敗のため次候補へフォールバック: {path.name} ({type(exc).__name__})")
    return f"[Excel: {path.name} — openpyxl or pandas が必要: pip install openpyxl]"


# 2026-07-26: グループシェイプ(GroupShape)のネスト再帰に上限を設ける。iter157で
# 姉妹関数_read_docxのネスト表再帰(L847の_DOCX_NESTED_TABLE_MAX_DEPTH)に加えた
# 深さ上限と同じ理由で、こちらは`elif hasattr(sh, "shapes"): for sub in sh.shapes:
# out.extend(_pptx_shape_texts(sub))`が無制限に自身を呼び直す構造になっており、
# 病的/破損した<p:grpSp>の異常に深いネスト連鎖(あるいは自己参照的な構造)を
# 与えられるとPythonのsys.recursionlimitを超えてRecursionErrorを送出しうる。
# _read_pptx側のtry節はimport pptx失敗によるImportErrorしか捕捉していないため
# (L1101-1117)、このRecursionErrorはそのまま_read_pptxの外へ伝播し、
# read_file_text(iter53)・_load_rag_chunks(iter42)の広いexcept Exceptionガードに
# 握りつぶされて、PowerPointファイル全体がRAG/--fileコンテキストから無言で
# 丸ごと脱落する。これはiter87が救済した「表/グループ内容が静かに読み飛ばされる」
# 精度事故と同じ害のクラスであり、精度優先・時間は気にしないの方針に反する。
# 上限はiter157の6(表ネストは実務上ほぼ無い)より大幅に緩い50を採用する:
# 実在するPPTXの図解グループは数階層程度しかネストしないため、この上限が現実の
# デッキを切り詰めることは実質無い一方、sys.recursionlimit(既定1000)には
# 十分な余裕を残す。上限に達した時点でそれ以上は再帰せず静かに打ち切る
# (iter157と同じ方針。bare exceptは使わずKeyboardInterrupt/SystemExitも
# そもそも送出していないので握りつぶす対象にならない)。
_PPTX_GROUP_MAX_DEPTH = 50


def _pptx_shape_texts(sh, _depth=0) -> list:
    """PPTXの1シェイプから抽出できるテキスト断片を出現順のリストで返す。
    2026-07-24 (iter87): python-pptxの表シェイプ(GraphicFrame)とグループシェイプ
    (GroupShape)には `.text` 属性そのものが存在しない（hasattr(sh, "text") が False
    になる）。そのため_read_pptxの従来コード
    `[sh.text for sh in slide.shapes if hasattr(sh, "text") and sh.text.strip()]`
    は表・グループ内の文字列を一切拾えず静かに読み飛ばしていた。実データのPPTXでは
    表が情報の大半を占めることが多く、その内容がRAG/--fileコンテキストへ一切届かず
    モデルが古い学習知識で答えてしまう精度事故（精度優先・時間は気にしないの方針に
    反する）。iter82は_read_pptxの成功時抽出にテストカバレッジを整えたが対象は平文
    テキストボックスのみで、表/グループ内容は当時未検証・未対応のまま残っていた。
    iter83(_read_pdf)/iter84(_read_excel)の「呼び出し元が静かに失っていたコンテンツを
    救済する」系統の修正の直系。表は他のシェイプ種別には無い属性のため
    getattr(sh, "has_table", False)で存在確認してからsh.tableを読み、行は
    _read_docxの表取り扱い(L708-710)と同じ規約でセルをタブ結合・行を改行結合する。
    グループはGroupShapeが`.text`を持たず`.shapes`だけを持つことをhasattr()で
    ダックタイピング判定し、ネストしたグループにも再帰する。他シェイプ種別に無い
    属性を無条件で呼ぶと例外になるため、必ずgetattr/hasattrで存在確認してから読む。
    2026-07-26: 上の_PPTX_GROUP_MAX_DEPTHの説明どおり、_depthが上限に達したら
    それ以上グループの中へは再帰しない(現実的な深さのデッキは全く影響を受けず、
    病的に深いデッキだけがRecursionErrorでの全体脱落から部分抽出に縮退する)。"""
    out = []
    if hasattr(sh, "text") and sh.text.strip():
        out.append(sh.text)
    elif getattr(sh, "has_table", False):
        for row in sh.table.rows:
            row_text = "\t".join(cell.text for cell in row.cells)
            if row_text.strip():
                out.append(row_text)
    elif hasattr(sh, "shapes"):
        if _depth < _PPTX_GROUP_MAX_DEPTH:
            for sub in sh.shapes:
                out.extend(_pptx_shape_texts(sub, _depth + 1))
    return out


def _read_pptx(path: Path) -> str:
    """PowerPoint (.pptx) からテキストを抽出。"""
    try:
        from pptx import Presentation
        # 2026-07-24 (iter123): _read_docx の Document() と同じ隙間が Presentation()
        # にもある。破損したzip(BadZipFile)やレガシーバイナリ.pptを.pptx拡張子の
        # まま開いた場合(PackageNotFoundError)等の実行時例外が、従来は
        # except ImportError のみでは捕捉されず外へ伝播していた。iter83/84/直上の
        # _read_docx修正と同じ方針で、open/parseの1行だけを狭くException限定で
        # 捕捉し、可視化されたnotice文字列に変換する（bare exceptにはしない。
        # KeyboardInterrupt/SystemExitは伝播させる）。直下の except ImportError: pass
        # は既存のバイト単位のまま温存。新noticeは"pip install"を含まないため
        # _is_lib_missing_notice の構造判定には意図的にヒットしない(未インストール
        # と誤判定させないための意図的な区別)。
        try:
            prs = Presentation(str(path))
        except Exception as exc:
            print(f"[_read_pptx] Presentation()での読み込みに失敗: {path.name} ({type(exc).__name__})")
            return f"[PPTX: {path.name} — 読み込みエラー ({type(exc).__name__}): 破損しているか非対応形式の可能性があります]"
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for sh in slide.shapes:
                # 2026-07-26: 1シェイプ単位で抽出を単離する。_pptx_shape_texts内部の
                # hasattr(sh, "text")/getattr(sh, "has_table", False)によるダック
                # タイピング分岐はAttributeErrorしか吸収しない（Python 3の仕様で
                # hasattrはAttributeError以外の例外をそのまま伝播させる）ため、
                # 破損/非対応の<p:sp>や<p:graphicFrame>(表)のプロパティアクセスが
                # 別の例外（壊れたXML由来のlxml例外等）を送出すると、この
                # forループ、ひいては_read_pptx自体の外まで例外がそのまま伝播していた。
                # _read_pptxには_read_pdf/_read_excelと違って代替ライブラリへの
                # フォールスルーが無いため、1シェイプの異常だけでデッキ全体
                # （他の全スライド・同スライドの他の全シェイプ・下のtry/exceptで
                # 個別に保護済みのスピーカーノートまで）がread_file_text(iter53)の
                # 広いexcept Exceptionに握りつぶされて""へ丸ごと脱落していた。
                # 既にこの関数内で個別保護されているスピーカーノート抽出(iter98)や
                # _read_docxの本文走査(iter93)と同じskip-bad-part-keep-the-restの
                # 方針で、1シェイプだけを読み飛ばして残り全部を救済する
                # (bare exceptではなくException限定。KeyboardInterrupt/SystemExitは
                # ここで握りつぶさず伝播させる)。_pptx_shape_texts自体のhasattr/
                # getattr分岐や_PPTX_GROUP_MAX_DEPTH再帰上限(iter177)はここでは
                # 一切変更しない。
                try:
                    texts.extend(_pptx_shape_texts(sh))
                except Exception as exc:
                    print(f"[_read_pptx] シェイプ抽出に失敗したためスキップ: {path.name} ({type(exc).__name__})")
            if texts:
                parts.append(f"[Slide {i}]\n" + "\n".join(texts))
            # 2026-07-24 (iter98): スピーカーノートはslide.shapesの走査対象外なので、
            # 上のループでは一切拾えない。ノートには箇条書き本文が要約している詳細な
            # 説明が書かれていることが多く、これがRAG/--fileコンテキストへ届かないと
            # プロポーザーが古い学習知識で答えてしまう精度事故になる（iter87で救済した
            # 表/グループ抽出と同系統の「静かに落としているコンテンツを拾う」修正で、
            # 精度優先・時間は気にしないの方針に基づく）。ただしslide.notes_slideは
            # 触れた時点で（存在しなければ）python-pptxが空のノートスライドをその場で
            # 生成してしまう副作用がある（NotesSlide.notes_slideのdocstring通り）ため、
            # 必ずslide.has_notes_slideで存在確認してからのみslide.notes_slideへ触れる。
            # 壊れた/読めないノート1件のせいでそのスライドの本文抽出まで失わないよう
            # （iter72のskip-bad-part-keep-the-rest方針、iter42/53の段階的劣化と同様）、
            # ノート読み取りだけを個別のtry/exceptで保護する。bareではなくException限定
            # なのでKeyboardInterrupt/SystemExitはここで握りつぶさず伝播する。
            try:
                if slide.has_notes_slide:
                    notes_tf = slide.notes_slide.notes_text_frame
                    if notes_tf is not None:
                        notes_text = notes_tf.text.strip()
                        if notes_text:
                            parts.append(f"[Slide {i} Notes]\n{notes_text}")
            except Exception:
                pass
        return "\n\n".join(parts)
    except ImportError:
        pass
    return f"[PPTX: {path.name} — python-pptx が必要: pip install python-pptx]"


# 2026-07-22: ライブラリ未インストール通知（_read_pdf/_read_docx/_read_excel/_read_pptx が
# 抽出失敗時に返す1行メッセージ）だけを検出する専用パターン。以前は「'[' で始まる」だけで
# 判定しており、_read_excel の成功時出力 "[Sheet: ...]" や _read_pptx の成功時出力
# "[Slide N]" まで RAG から除外してしまっていた（= 抽出できた Excel/PPTX が丸ごと欠落する
# 精度事故）。通知文字列は「1行・'[' で始まり ']' で終わり・'pip install' を含む」という
# 安定した構造を持つため、それだけで狭く判定する。
_LIB_MISSING_NOTICE_RE = re.compile(r"^\[[^\n]*pip install[^\n]*\]$")


def _is_lib_missing_notice(text: str) -> bool:
    """read_file_text の抽出結果が「ライブラリ未インストール」通知そのものかを判定。
    通知は _read_pdf/_read_docx/_read_excel/_read_pptx がフォールバック時に返す
    1行文字列のみで、他の正常な抽出結果（Excel/PPTX の先頭が '[' の行等）とは
    「pip install を含む単一行」という点で区別できる。"""
    stripped = text.strip()
    if "\n" in stripped:
        return False
    return bool(_LIB_MISSING_NOTICE_RE.match(stripped))


def _decode_text_bytes(data: bytes) -> str:
    """入力ファイルのバイト列を str にデコードする（utf-8 -> cp932 -> replace の順）。

    2026-07-24: 既知の落とし穴 #4（このマシンのコンソールが cp932/Shift-JIS で
    あること）と同根の環境要因で、ローカルに保存された非UTF-8（Shift-JIS/cp932）の
    .txt/.csv/.html 等がこのマシンには普通に存在する。従来 read_file_text の
    汎用テキスト分岐と _read_html は encoding="utf-8", errors="replace" を
    無条件に使っていたため、cp932 ファイルを読むと日本語部分が「全て」
    U+FFFD（置換文字）に化け、その文字化けがそのまま精度critical な
    RAG/--file コンテキストへ注入されていた（精度優先・時間は気にしないの
    方針に反する重大な劣化）。iteration 47（_save_as_markdown/_save_as_text/
    _save_as_html の読み戻し）・iteration 70（会話履歴JSON読み込み）で導入した
    errors="replace" は「fugu自身がUTF-8で書いたファイルの読み戻し」用の保険で
    あり、他所由来ファイルの元エンコーディングを救済するものではない。
    ここでは stdlib のみで完結する最小限のデコードラダーを用意し、
    (1) utf-8 厳密デコードを試す（クリーンなUTF-8ファイルは従来と完全に
    バイト同一の結果になる。utf-8-sig ではなく素の utf-8 を使うため BOM の
    扱いも従来の read_text(encoding="utf-8") と同一）、(2) 失敗したら cp932
    厳密デコードを試す（成功すれば日本語を含む cp932 ファイルを正しく復元
    できる）、(3) それも失敗したら最終手段として utf-8 + errors="replace"
    （＝これまでの挙動そのもの、退行なし）にフォールバックする。

    呼び出し元は path.read_text(encoding=...) ではなく path.read_bytes() で
    取得した生バイト列を渡す設計のため、read_text() が標準で行う普遍改行
    変換（newline=None: ファイル上の "\r\n" や単独の "\r" を "\n" に変換）を
    ここで明示的に再現する。これを省くと、Windows上で write_text() が
    書き出す "\r\n" 改行がデコード後もそのまま残ってしまい、read_text()を
    使っていた従来の出力とバイト単位で一致しなくなる（回帰）。
    この関数は例外を送出しない。"""
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError:
        try:
            text = data.decode("cp932")
        except UnicodeDecodeError:
            text = data.decode("utf-8", errors="replace")
    return text.replace("\r\n", "\n").replace("\r", "\n")


def _read_html(path: Path) -> str:
    """HTML からタグを除去してテキストを返す（stdlib html.parser 使用）。"""
    from html.parser import HTMLParser

    # 2026-07-26: 旧実装は handle_data 1回ごとに data.strip() して self.parts
    # に積み、最後に "\n".join(parts) していた。これは「テキストノードの区切り
    # =意味的な区切り」という誤った前提に基づいており、<b>/<strong>/<a>/<sub>/
    # <sup>/<span>/<code> のようなインライン要素は子孫テキストを別ノードとして
    # 分割するだけで文/フレーズの継続を意味しない。結果、'<p>The <b>quick</b>
    # brown fox</p>' は 'The\nquick\nbrown fox' に、日本語の
    # '機械<strong>学習</strong>技術' は '機械\n学習\n技術' に断片化し、
    # ノード間の実際の空白（'The ' の末尾スペース等）もstrip()で失われていた。
    # これは iter71 が直接カバレッジを付けた _read_html の「タグ除去」自体は
    # 正しいが、iter94 の decode ラダー（_decode_text_bytes、下記で不変のまま
    # 使用）で正しく日本語復元できても、その後のインライン分割で本文の連続性が
    # 壊れるという別問題。特に致命的なのは、この断片化が iter179 で導入した
    # 日本語（CJK）バイグラムトークナイザ（_tokenize: 非ASCII連続列を隣接2文字
    # ずつのバイグラムに分解し rag_search の再現率を確保する方式）を無力化する
    # 点: '機械\n学習' のように改行を挟んでしまうと '機械'→'学習' の境界を跨ぐ
    # バイグラム '械学' が二度と生成されず、RAG検索でヒットしなくなる
    # （本プロジェクトの主要言語である日本語で --file/RAGコンテキストの再現率が
    # 下がる）。'10<sup>3</sup>' が '10\n3' になる数値表記の破壊も同根。
    # 対処: ブロックレベルタグ（p/div/li/tr/td/th/table/ul/ol/h1-6/br/hr/
    # section/article/header/footer/blockquote/pre 等）の開始・終了時にのみ
    # 区切り記号を挿入し、インライン/未知タグでは何も挿入しない。テキスト
    # ノード自体は data.strip() せず生のまま連結する（ノード内・ノード間の
    # 本来の空白を保持するため）。最後にまとめて行単位へ分割し、各行を
    # strip() して空行を除去する（1行内の単一スペースは行の内部にあるので
    # 保持され、既存の完全一致回帰 'Hello UTF-8 world. こんにちは。' 等は
    # 崩れない）。td/th/tr/li/br は引き続き区切りを出すため、
    # '<td>A</td><td>B</td>' や 'A<br>B' が1行に潰れる回帰は起きない。
    # なお <sup>/<sub> はインライン扱いのため 'The best' の '10^3' が視覚的に
    # 上付き/下付きだった情報自体はプレーンテキスト抽出では原理的に復元
    # できない（'103' という連結された数字列になる）——これは仕様上の限界で
    # あり、本修正が解決するのは「分断されない」ことまで。
    _BLOCK_TAGS = frozenset({
        "address", "article", "aside", "blockquote", "br", "caption",
        "dd", "details", "dl", "dt", "div", "figcaption", "figure",
        "footer", "form", "h1", "h2", "h3", "h4", "h5", "h6", "header",
        "hr", "li", "main", "nav", "ol", "p", "pre", "section", "summary",
        "table", "tbody", "td", "tfoot", "th", "thead", "tr", "ul",
    })

    class _Strip(HTMLParser):
        def __init__(self):
            super().__init__()
            self.parts = []
            self._skip = False
        def handle_starttag(self, tag, attrs):
            if tag in ("script", "style"):
                self._skip = True
            if tag in _BLOCK_TAGS:
                self.parts.append("\n")
        def handle_endtag(self, tag):
            if tag in ("script", "style"):
                self._skip = False
            if tag in _BLOCK_TAGS:
                self.parts.append("\n")
        def handle_data(self, data):
            # 生のまま(strip しない)積む: インライン要素をまたぐノード間の
            # 空白（例: 'The ' の末尾スペース）を保持するため。
            if not self._skip and data:
                self.parts.append(data)

    # 2026-07-24: 落とし穴 #4 (cp932コンソール) と同根の環境要因で、ローカルの
    # .html/.htm がShift-JIS(cp932)保存されていることがある。旧来の
    # encoding="utf-8", errors="replace" 一本槍だと日本語本文が全滅してU+FFFD
    # 化けがそのままRAG/--fileコンテキストに載っていた。_decode_text_bytes()
    # のutf-8→cp932→replaceラダーで復元する（iter47/iter70のerrors="replace"
    # 適用箇所とは目的が異なる: あちらはfugu自身がUTF-8で書いたファイルの
    # 読み戻し用の保険で、こちらは他所由来ファイルの元エンコーディングの救済）。
    raw = _decode_text_bytes(path.read_bytes())
    p = _Strip()
    p.feed(raw)
    # ブロック区切り("\n")と生テキストを結合してから行単位に正規化する:
    # 各行をstrip()して空行を落とす。インライン要素の内側/またぎでは区切りが
    # 一切挿入されていないため、同じ行内でテキストがそのまま連結される。
    joined = "".join(p.parts)
    lines = (line.strip() for line in joined.split("\n"))
    return "\n".join(line for line in lines if line)


def _read_ipynb(path: Path) -> str:
    """Jupyter Notebook からコードセルとマークダウンセルを抽出。"""
    # 2026-07-25: 落とし穴 #4 (cp932コンソール) と同根の環境要因で、ローカルに
    # 保存された .ipynb がShift-JIS(cp932)保存されていることがある。notebookの
    # JSON構造（波括弧・引用符・キー名）はすべてASCIIのため、cp932保存の
    # notebookでもjson.loads自体は成功し、iteration 72/113が導入したセル単位/
    # トップレベルの構造ガードも問題なく通過してしまう。だが従来の
    # path.read_text(encoding="utf-8", errors="replace")（本関数の解析用と、
    # 下のexceptフォールバック用の2箇所で使用）はデコード時点で各セルの
    # 'source'に含まれる日本語を全てU+FFFD（置換文字）へ潰してから渡して
    # いたため、json.loadsが成功して構造ガードも通ってもセル内容そのものは
    # 既に文字化けした状態のまま、精度criticalなRAG/--fileコンテキストへ
    # 注入されていた（他の多くのリーダーと異なり.ipynbはJSONなので
    # 「パース自体は成功するのに中身だけ化ける」という気づきにくい劣化）。
    # iteration 94が_read_html/read_file_text汎用テキスト分岐に導入した
    # utf-8→cp932→replaceの_decode_text_bytes()ラダーをここにも適用する
    # （iteration 70のerrors="replace"はfugu自身がUTF-8で書いたファイルの
    # 読み戻し用の保険であり、他所由来ファイルの元エンコーディングまでは
    # 救済しないという整理は変わらない）。生バイト列はpath.read_bytes()で
    # 一度だけ読み、_decode_text_bytes()で一度だけデコードし、そのデコード
    # 結果をjson.loadsにも、非JSON時のexceptフォールバック（生テキスト
    # 返却）にも共用する（旧実装のようにexcept節でpath.read_text()を
    # 再度読み直すことはしない）。iteration 72（セル単位: source文字列/
    # list[str]/その他の正規化、空白sourceのスキップ、非strリスト要素の
    # 除外）とiteration 113（トップレベル: 非dictなnbは''へ、非listな
    # cellsは[]へ強制変換）の構造ガードは、このデコード経路変更とは独立に
    # 一切変更せずそのまま維持する。
    raw = _decode_text_bytes(path.read_bytes())
    try:
        nb = json.loads(raw)
        # 2026-07-24: json.loads自体は成功しても、トップレベルがdictでない
        # ケース（例: [{"cell_type": "code"}] のような配列、あるいは裸の数値・
        # 文字列）や、nb["cells"]がtruthyな非list値（例: {"cells": 42} や
        # {"cells": {"0": {...}}} のようなdict）が存在しうる。旧実装は前者で
        # nb.get(...)がAttributeErrorを、後者でfor cell in <non-list>が
        # TypeErrorを送出し、いずれも本関数の外側except Exceptionに落ちて
        # notebook全体の生JSON（metadata・execution_count・base64画像出力を
        # 含む）がそのままRAG/--fileコンテキストへ丸ごと注入されていた。
        # これはiteration 71がtest-onlyで発見してフラグし、iteration 72が
        # セル単位の壊れたsourceについてのみ対処した「1件の破損で全体を
        # 道連れにする」劣化のトップレベル構造版であり、対処が漏れていた
        # ものである。iteration 103/111/112と同じ非list強制変換（truthy
        # チェックのトリックに頼らず必ず既定値へ倒す）の作法に倣い、構造が
        # 不正な場合は例外を送出させず空文字列/空contextへ落とす（精度優先：
        # 生JSON混入より整形済みの空の方がまし）。json.loadsが失敗した場合
        # （真に不正なJSON）は従来通り外側exceptで生テキストへフォールバック
        # する経路を変えない。
        if not isinstance(nb, dict):
            return ""
        cells = nb.get("cells", [])
        if not isinstance(cells, list):
            cells = []
        parts = []
        # 2026-07-26: コードセルのstream出力(stdout/stderr)1件あたりの取り込み文字数
        # 上限。RAG_CHUNK_CHARS/SEARCH_CONTEXT_CHARS等と同様、暴走したprintループ
        # (例: 無限に近いループ内print)がnum_ctxを圧迫しないための保護であり、下の
        # コードセル処理でのみ参照するため関数ローカルに留める（本iterationの変更を
        # _read_ipynb関数の外へ一切出さないため）。
        _IPYNB_STREAM_OUTPUT_CAP = 4000
        for cell in cells:
            # 2026-07-23 (iter72): nbformat仕様上 'source' はstrまたはlist[str]だが、
            # 壊れた/書きかけのnotebookでは source: null や、list内に非str要素
            # （例: 42）が混入することがある。旧実装は"".join(cell.get("source", []))を
            # 無条件に呼んでおり、これがTypeError/AttributeErrorを送出すると本関数の
            # 外側try/exceptで捕捉され、「1セルの壊れたsourceでnotebook全体の生JSON
            # 全文がそのままRAG/--fileコンテキストへ丸ごと注入される」という劣化を
            # 招いていた（iteration 71がtest-onlyで発見し、特性検証テストとして
            # ピン留めした上で将来のiterationでの修正候補と明示的にフラグしていたもの）。
            # ここではセル単位で壊れたsourceだけを安全にスキップし、良好な他セルの
            # 構造化抽出（精度優先：整形済みcontext > 生JSON）を守る。cell自体が
            # dictでない場合も同様にスキップする。これはiteration 42/53と同じ
            # 「1件の破損で全体を道連れにしない」graceful degradationパターン。
            if not isinstance(cell, dict):
                continue
            ct = cell.get("cell_type", "")
            raw_src = cell.get("source", [])
            if isinstance(raw_src, str):
                src = raw_src
            elif isinstance(raw_src, list):
                src = "".join(s for s in raw_src if isinstance(s, str))
            else:
                src = ""
            if not src.strip():
                continue
            if ct == "code":
                parts.append(f"```python\n{src}\n```")
                # 2026-07-26: 従来はコードセルの'source'(入力コード)のみを抽出し、
                # 'outputs'配列(実行結果)は完全に無視していた。データ分析notebookでは
                # print(...)が出力する実際の数値・結果こそが質問への回答に直結する
                # 事実であることが多く、それが精度criticalなRAG/--fileコンテキストから
                # 黙って欠落していた（精度優先・時間は気にしない、の方針に反する）。
                # iteration 188はoutput_type: stream/execute_result/display_data
                # (text/plain)の全種類を一度に扱おうとして3回試みたが行き詰まり
                # 断念した。iteration 194はあえてスコープを"stream"(stdout/stderr)
                # 出力1種類のみへ絞り込み、最も単純な単一形状・最も一般的・最も
                # 価値が高いケースに限定することで、iter188の轍を踏まずに小さく
                # 検証可能な変更に収めた。iteration 195(本iteration)はiter194の
                # 縮小方針を踏襲しつつ、次に価値が高い単一形状としてexecute_result/
                # display_dataの'data'辞書のうち'text/plain'キーのみを追加で対象と
                # する（下のstreamループの直後を参照）。'data'辞書内のtext/htmlや
                # image/png等の非'text/plain'MIME、およびerrorのtracebackは、
                # オーバーサイトではなく意図的に対象外のまま据え置く（base64画像等の
                # 非テキストMIMEを扱うには別途の慎重な設計が必要なため、iter188の
                # 「一度に全部」の轍を踏まないよう最小スコープを維持する）。'source'の
                # 正規化(iteration 72: str->そのまま、list->str要素のみjoin、
                # その他->空、空白のみはskip)と全く同じパターンを各streamの'text'
                # にも適用し、cellがdictでない場合のskip(iteration 72)・cells/nbの
                # トップレベル構造ガード(iteration 113: 非dictなnbは''へ、非listな
                # cellsは[]へ)はこの変更でも一切変更しない。'outputs'自体がtruthyな
                # 非list(例: 42)や、outputsの各要素が非dictの場合も、iteration 113の
                # cells非list強制変換・iteration 72の非dictセルskipと同じ作法
                # (`or []`のtruthinessトリックには頼らずisinstanceチェックで明示的に
                # 空/skipへ倒す)を踏襲する。暴走したprintループが巨大な出力を吐いて
                # num_ctxを圧迫しないよう、1コードセルあたりの出力文字数に上限を
                # 設け、超過分は切り詰めマーカーを付けて明示する。
                outputs = cell.get("outputs", [])
                if not isinstance(outputs, list):
                    outputs = []
                stream_chunks = []
                for out in outputs:
                    if not isinstance(out, dict):
                        continue
                    if out.get("output_type") != "stream":
                        continue
                    raw_out_text = out.get("text", [])
                    if isinstance(raw_out_text, str):
                        out_text = raw_out_text
                    elif isinstance(raw_out_text, list):
                        out_text = "".join(t for t in raw_out_text if isinstance(t, str))
                    else:
                        out_text = ""
                    if not out_text.strip():
                        continue
                    stream_chunks.append(out_text)
                if stream_chunks:
                    combined_out = "".join(stream_chunks)
                    if len(combined_out) > _IPYNB_STREAM_OUTPUT_CAP:
                        combined_out = (
                            combined_out[:_IPYNB_STREAM_OUTPUT_CAP]
                            + "\n...(出力が長いため切り詰め)"
                        )
                    parts.append(f"[Notebook stdout/stderr output]\n{combined_out}")
                # 2026-07-26 (iteration 195): stream(標準出力/標準エラー)に続き、
                # execute_result/display_dataの'data'辞書のうち'text/plain'キーのみを
                # 追加抽出する。iter188はstream/execute_result/display_data(text/plain)
                # の全種類を一度に扱おうとして行き詰まり断念し、iter194はあえて
                # streamのみへスコープを縮小した。本iterationはiter194の縮小方針を
                # そのまま踏襲し、次に価値が高い単一形状(戻り値やDataFrameのrepr等、
                # セルの"計算結果"を表すtext/plain)のみを追加する。data分析notebookでは
                # print(...)によるstream出力だけでなく、セル末尾の式の評価結果や
                # 変数のreprこそが質問への回答に直結することが多く、これも精度
                # criticalなRAG/--fileコンテキストから黙って欠落していた。'data'辞書は
                # image/png等のbase64エンコード画像やtext/html等、text/plain以外の
                # MIMEキーを同時に含みうるが、それらは意図的に対象外のまま据え置く
                # （オーバーサイトではない。base64画像を扱うには別途の慎重な設計が
                # 必要なため、iter188の「一度に全部」の轍を踏まないよう最小スコープを
                # 維持する）。'error'のtracebackも同様に対象外のまま据え置く。
                # 正規化パターン(str->そのまま、list->str要素のみjoin、その他->空、
                # 空白のみはskip)は'source'(iteration 72)・stream 'text'
                # (iteration 194)と全く同一のものを'data'内'text/plain'にも適用し、
                # `or []`等のtruthinessトリックには頼らずisinstanceチェックのみで
                # 判定する。'outputs'非list・outputs内非dict要素・cell非dict等の
                # 構造ガード(iteration 72/113)は上のstreamループと同じ変数
                # (outputs)を再利用するのみで一切変更しない。暴走した出力が
                # num_ctxを圧迫しないよう、上のstream出力と同じ
                # _IPYNB_STREAM_OUTPUT_CAPを流用して上限を設ける。
                result_chunks = []
                for out in outputs:
                    if not isinstance(out, dict):
                        continue
                    if out.get("output_type") not in ("execute_result", "display_data"):
                        continue
                    data = out.get("data", {})
                    if not isinstance(data, dict):
                        continue
                    raw_result_text = data.get("text/plain", "")
                    if isinstance(raw_result_text, str):
                        result_text = raw_result_text
                    elif isinstance(raw_result_text, list):
                        result_text = "".join(
                            t for t in raw_result_text if isinstance(t, str)
                        )
                    else:
                        result_text = ""
                    if not result_text.strip():
                        continue
                    result_chunks.append(result_text)
                if result_chunks:
                    combined_result = "".join(result_chunks)
                    if len(combined_result) > _IPYNB_STREAM_OUTPUT_CAP:
                        combined_result = (
                            combined_result[:_IPYNB_STREAM_OUTPUT_CAP]
                            + "\n...(出力が長いため切り詰め)"
                        )
                    parts.append(f"[Notebook result output]\n{combined_result}")
                # 2026-07-26 (iter196): stream(iter194)・execute_result/display_dataの
                # 'text/plain'(iter195)に続き、output_type=='error'の'ename'/'evalue'の
                # みを追加抽出する。従来はコードセルの保存済み出力がerrorであっても
                # ```pythonフェンスのみが出力され、セルが失敗した事実が一切示されない
                # まま提案者に渡っていた。これは提案者がコードが成功したものと誤って
                # 推論しうる精度上の欠陥であり、iter188がstream/execute_result/
                # display_data/errorの全種類を一度に扱おうとして行き詰まり断念して以来、
                # iter194/195が意図的に対象外のまま据え置いてきた最後の1種
                # （iter195のコメント:「'error'のtracebackも同様に対象外のまま据え置く」）
                # を、同じ縮小スコープの作法で補う。'traceback'(list of str、ANSIカラー
                # エスケープシーケンスを含む)は本iterationでも意図的に対象外のまま
                # 据え置く(iter195からのdeferralをそのまま継承する)。ANSIエスケープの
                # 除去・正規化や、'data'辞書内のtext/html・image/png等の非'text/plain'
                # MIME(iter195で既に対象外と決めたもの)の扱いには別途の慎重な設計が
                # 必要であり、iter188の「一度に全部」の轍を踏まないよう'ename'/'evalue'
                # という単一の最小形状のみへスコープを保つ。正規化パターン
                # (str->そのまま、list->str要素のみjoin、その他->空、strip後に両方
                # 空白ならそのoutputはskip)はsource(iter72)・stream 'text'(iter194)・
                # 'text/plain'(iter195)と全く同一のものをename/evalueにも適用し、非str値
                # (int/dict/None等)をstr()で強制変換して混入させることは絶対にしない。
                # 'outputs'非list・outputs内非dict要素・cell非dict等の構造ガード
                # (iter72/113)は上の2ループと同じ変数(outputs)を再利用するのみで一切
                # 変更しない。暴走した出力(極端に長いevalue)がnum_ctxを圧迫しないよう、
                # 上と同じ_IPYNB_STREAM_OUTPUT_CAPを流用して上限を設ける。
                error_chunks = []
                for out in outputs:
                    if not isinstance(out, dict):
                        continue
                    if out.get("output_type") != "error":
                        continue
                    raw_ename = out.get("ename", "")
                    if isinstance(raw_ename, str):
                        ename = raw_ename
                    elif isinstance(raw_ename, list):
                        ename = "".join(e for e in raw_ename if isinstance(e, str))
                    else:
                        ename = ""
                    raw_evalue = out.get("evalue", "")
                    if isinstance(raw_evalue, str):
                        evalue = raw_evalue
                    elif isinstance(raw_evalue, list):
                        evalue = "".join(e for e in raw_evalue if isinstance(e, str))
                    else:
                        evalue = ""
                    if not ename.strip() and not evalue.strip():
                        continue
                    if ename.strip() and evalue.strip():
                        error_chunks.append(f"{ename}: {evalue}")
                    elif ename.strip():
                        error_chunks.append(ename)
                    else:
                        error_chunks.append(evalue)
                if error_chunks:
                    combined_error = "\n".join(error_chunks)
                    if len(combined_error) > _IPYNB_STREAM_OUTPUT_CAP:
                        combined_error = (
                            combined_error[:_IPYNB_STREAM_OUTPUT_CAP]
                            + "\n...(出力が長いため切り詰め)"
                        )
                    parts.append(f"[Notebook error]\n{combined_error}")
            elif ct == "markdown":
                parts.append(src)
        return "\n\n".join(parts)
    except Exception:
        return raw


# テキストとして直接読めない拡張子
_BINARY_SKIP = {
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp",
    ".mp3", ".mp4", ".wav", ".avi", ".mov",
    ".zip", ".tar", ".gz", ".7z", ".rar",
    ".exe", ".dll", ".so", ".dylib",
    ".bin", ".dat", ".pkl", ".pt", ".pth", ".onnx",
}


def read_file_text(path: Path) -> str:
    """あらゆるファイルからテキストを抽出する。
    対応形式: テキスト・コード類 / PDF / Word / Excel / PowerPoint / HTML / Notebook。
    バイナリ（画像・動画・アーカイブ等）はスキップして空文字を返す。"""
    suffix = path.suffix.lower()
    if suffix in _BINARY_SKIP:
        return ""
    # 2026-07-23 (iter53): _read_pdf/_read_docx/_read_excel/_read_pptx/_read_html は
    # それぞれ import 文と実際のパース処理を同じ try/except ImportError で包んでいるため、
    # ライブラリ自体は入っているがファイルが壊れている/パスワード付き/旧形式バイナリ等の
    # 場合（zipfile.BadZipFile, PDFSyntaxError 等の ImportError 以外の例外）はここまで
    # 素通しで伝播していた。iter42 は _load_rag_chunks 側（RAG経由の呼び出し）だけを
    # 個別に保護したが、main() の --file 呼び出し（`read_file_text(fp).strip()`）は
    # 無防備なままで、壊れたOffice/PDFファイルを渡すとCLI全体がクラッシュしていた
    # （本関数のdocstringが約束する「空文字を返す」というgraceful degradation契約に
    # 反する）。全リーダー関数の書き換えを試みたiter51は行き詰まったため、ここでは
    # ディスパッチ関数側だけを薄く保護し、失敗時は空文字にフォールバックしつつ
    # 警告を可視化する（黙って握りつぶさない）。
    try:
        if suffix == ".pdf":
            return _read_pdf(path)
        if suffix in {".docx", ".doc"}:
            return _read_docx(path)
        if suffix in {".xlsx", ".xls"}:
            return _read_excel(path)
        if suffix in {".pptx", ".ppt"}:
            return _read_pptx(path)
        if suffix in {".html", ".htm"}:
            return _read_html(path)
        if suffix == ".ipynb":
            return _read_ipynb(path)
    except Exception as exc:
        print(f"[read_file_text] 読み込み失敗のためスキップ: {path} ({type(exc).__name__})")
        return ""
    # その他: テキストとして読む（コード・設定ファイル・Markdown など）
    try:
        # 2026-07-24: 落とし穴 #4 (cp932コンソール) と同根の環境要因で、この
        # 汎用テキスト分岐が読む .txt/.md/.csv/.py/.json 等がローカルで
        # Shift-JIS(cp932)保存されていることがある。旧来の
        # encoding="utf-8", errors="replace" 一本槍では日本語が全滅し、
        # 精度critical なRAG/--fileコンテキストに文字化けがそのまま注入
        # されていた。_decode_text_bytes() のutf-8→cp932→replaceラダーで
        # 復元する（iter47/iter70のerrors="replace"適用箇所とは目的が異なる:
        # あちらはfugu自身がUTF-8で書いたファイルの読み戻し用の保険で、
        # こちらは他所由来ファイルの元エンコーディングの救済）。
        raw_bytes = path.read_bytes()
        # 2026-07-25: read_file_textのdocstringは「バイナリはスキップして
        # 空文字を返す」契約（iter53/125のgraceful-degradation方針と同根）を
        # 約束しているが、実際にこれを支えているのは_BINARY_SKIP（約30拡張子
        # のみの小さなdenylist、上記）だけである。.npy/.h5/.parquet/
        # .safetensors/.gguf/.sqlite/.db/.woff/.ttf/.class/.wasm/.pyc や
        # 拡張子なしバイナリなど_BINARY_SKIP未収載のバイナリはこの汎用分岐まで
        # 落ちてきて、_decode_text_bytes()（iter94のutf-8→cp932→replace
        # ラダー、下記関数）は例外を送出しない設計のため、文字化けした
        # 「ゴミテキスト」がそのまま返っていた。このゴミは--file経路では
        # main()で質問全文そのものになり（下流フィルタなし）、RAG経路では
        # _load_rag_chunks（iter42のファイル単位隔離、下記）を通じて精度
        # criticalなコンテキストへチャンク注入される。RAG_EXTENSIONS（本ファイル
        # L271）は本来アローリストとして使う想定に見えるが実際は未使用の
        # dead codeであり、ここをアローリスト化する対処も検討したが、
        # .log/.conf/Dockerfile/READMEや拡張子なしテキストなど正当なNUL非
        # 含有テキストまで拡張子未収載を理由に巻き添えで弾いてしまい、
        # 精度優先・時間は気にしないの方針に反するrecall低下となるため
        # 採用しない。代わりにgit等が使う標準的なバイナリ判定手法である
        # 「生バイト列中のNUL(0x00)の有無」で判定する: NULはcp932/Shift-JIS
        # を含むテキストエンコーディングでは正当な文字の一部として現れない
        # ため、iter94のcp932救済ラダー（NUL非含有入力）を一切変更せず、
        # 真のバイナリ（NUL含有）だけを追加で弾ける。
        if b"\x00" in raw_bytes:
            print(f"[read_file_text] バイナリ検出(NULバイト)のためスキップ: {path}")
            return ""
        return _decode_text_bytes(raw_bytes)
    except Exception as exc:
        # 2026-07-24 (iter125): 上のsuffix-dispatch分岐(iter53)はここに来る前に
        # 例外を捕捉して警告を出すが、この汎用テキストフォールバック分岐は
        # 従来 except Exception: return "" のみで、権限エラー(PermissionError)や
        # Path.read_bytes自体の失敗、_decode_text_bytes内のデコード失敗が
        # 無警告のまま握りつぶされていた（_load_rag_chunksの個別ファイルガードは
        # 常に警告を出すのに対し、このパスだけ非対称だった）。read_file_textの
        # 「例外→空文字」という契約自体は変えず、可視化のための警告printのみ追加する。
        print(f"[read_file_text] 読み込み失敗のためスキップ: {path} ({type(exc).__name__})")
        return ""


def _load_rag_chunks(dirs: list) -> list:
    """指定ディレクトリ群からファイルを読み込み、
    (filepath, chunk_text) のリストを返す。"""
    chunks = []
    for d in dirs:
        p = Path(d)
        if not p.is_dir():
            print(f"   [RAG] ディレクトリが見つかりません: {d}")
            continue
        for fp in sorted(p.rglob("*")):
            if not fp.is_file():
                continue
            if fp.suffix.lower() in _BINARY_SKIP:
                continue
            # 2026-07-22 (iter42): read_file_text() 呼び出しを裸のまま置くと、破損/未対応の
            # 単一ファイル（壊れた.xlsx、パスワード付き.xlsx、python-docxに渡した旧形式.doc、
            # 壊れた.pptx/.pdf、不正な.htmlなど）が ImportError 以外の例外を送出した瞬間に
            # _load_rag_chunks -> _get_rag_chunks -> rag_search -> build_context まで伝播し、
            # 質問のたびにRAGコンテキストが丸ごと失われていた（1ファイルの破損がRAG全体を
            # 道連れにする精度事故）。精度優先・時間は気にしない方針のもと、iter41の
            # graceful-degradation方針（ImportError限定の握りつぶしを超えて、失敗は
            # スキップして続行する）を踏襲し、ここで1ファイル単位に隔離する。
            try:
                text = read_file_text(fp)
            except Exception as _rag_read_exc:
                print(f"   [RAG] 読み込み失敗のためスキップ: {fp} ({type(_rag_read_exc).__name__})")
                continue
            # 2026-07-22: 以前は text.startswith("[") で判定しており、成功時に "[Sheet: ...]"
            # (_read_excel) や "[Slide 1]" (_read_pptx) で始まる正常な抽出結果まで
            # 誤ってスキップしていた（Excel/PPTX が RAG から常時欠落する精度事故。
            # 精度優先・時間は気にしない方針に反する）。ライブラリ未インストール通知
            # （1行・pip install を含む）だけを狭く検出してスキップする。
            if not text or _is_lib_missing_notice(text):
                continue
            # チャンク分割（オーバーラップ付き）
            # 2026-07-25: RAG_CHUNK_CHARS/RAG_CHUNK_OVERLAP は上（L268-269）で
            # 「チャンクサイズ」「オーバーラップ文字数」と明記された調整可能な
            # モジュール定数。従来は step を
            # `start += RAG_CHUNK_CHARS - RAG_CHUNK_OVERLAP` と直接計算しており、
            # 将来だれかがオーバーラップを厚くしようとして
            # RAG_CHUNK_OVERLAP >= RAG_CHUNK_CHARS に設定すると step が 0 以下になり、
            # `while start < len(text)` の start が二度と前進しないまま chunks へ
            # 追記し続ける無限ループ（ハング）になる。これは
            # _load_rag_chunks -> _get_rag_chunks -> rag_search -> build_context に
            # 直結しており、以後すべての質問で応答が一切返らなくなる、"遅い"より
            # 悪い完全なフリーズになる。iter 111 (plan_pptx_images の非list ガード)・
            # iter 132 (_resolve_proposer の非ハッシュ可能値ガード) と同じ
            # 「発生確率は低いが起きれば致命的」クラスの防御として、step を
            # 最低1文字分の前進に固定する。既定値(100 < 600)では
            # max(1, 500) == 500 となり、従来のチャンク境界・オーバーラップ・
            # 個数・順序は一切変わらない（strict no-op）。
            start = 0
            step = max(1, RAG_CHUNK_CHARS - RAG_CHUNK_OVERLAP)
            while start < len(text):
                end = start + RAG_CHUNK_CHARS
                chunks.append((str(fp), text[start:end]))
                start += step
    return chunks


_RAG_CHUNKS: list = []   # キャッシュ（初回のみ読み込み）
_RAG_DIRS_LOADED: list = []


def _get_rag_chunks(dirs: list) -> list:
    global _RAG_CHUNKS, _RAG_DIRS_LOADED
    if dirs != _RAG_DIRS_LOADED:
        _RAG_CHUNKS = _load_rag_chunks(dirs)
        _RAG_DIRS_LOADED = list(dirs)
        print(f"   [RAG] {len(_RAG_CHUNKS)} チャンク読み込み完了（{len(dirs)} ディレクトリ）")
    return _RAG_CHUNKS


def _tokenize(text: str) -> set:
    """英字・数字・日本語を混在テキストから別々に抽出してトークンセットを返す。
    例: 'PINNについて' → {'pinn', 'につ', 'つい', 'いて'}
    (単純 \\w+ では 'pinnについて' 1トークンになる)。

    2026-07-26: 非ASCII連続列は以前、丸ごと1トークンにしていた
    （iter38で _tokenize('PINNについて') == {'pinn','について'} として
    固定化されていた挙動）。しかしこれだと、クエリの部分フレーズが
    チャンク側のより長い非ASCII連続列と「完全一致」しない限り
    _score_chunk の集合オーバーラップが常に0になり、'機械学習の手法に
    ついて' のような現実的な日本語クエリで rag_search が実質常に ''
    を返す＝日本語RAG再現率がほぼゼロという精度上の問題があった
    （本プロジェクトの主要言語である日本語で --file / RAG 経由の文書が
    proposer に一切届かない）。形態素解析器（MeCab等）を追加せずに
    対処する標準的な手法として、非ASCII連続列を隣接2文字の文字
    バイグラムに分解する（Lucene/Elasticsearch の CJKBigramFilter と
    同様の方式）。1文字だけの連続列はそのまま1文字トークンとする。
    ASCII側の抽出ロジックは一切変更していないため、英語/ASCIIの
    トークン化・rag_searchの挙動は完全に不変（CJKバイグラムはASCII
    クエリトークンと一致し得ない）。"""
    lower = text.lower()
    # ASCII: 英字・数字・アンダースコア（変更なし・byte-for-byte同一）
    tokens = set(re.findall(r'[a-z0-9_]+', lower))
    # 非ASCII連続列（日本語・CJK など）→ 隣接2文字ずつのバイグラムに分解
    for run in re.findall(r'[^\x00-\x7f\s]+', text):
        if len(run) == 1:
            tokens.add(run)
        else:
            tokens.update(run[i:i + 2] for i in range(len(run) - 1))
    return tokens - {''}


def _score_chunk(query_tokens: set, chunk: str) -> float:
    """クエリトークンとチャンクのキーワード重複スコア（TF-IDF 簡易版）。"""
    chunk_tokens = _tokenize(chunk)
    if not chunk_tokens:
        return 0.0
    overlap = len(query_tokens & chunk_tokens)
    return overlap / (len(chunk_tokens) ** 0.5 + 1) * 100


def rag_search(question: str, dirs: list = None, top_k: int = None) -> str:
    """ローカル文書をキーワード検索して上位チャンクをフォーマット済み文字列で返す。
    dirs が空（RAG_DIRS も空）なら空文字を返す。"""
    dirs = dirs or RAG_DIRS
    if not dirs:
        return ""
    top_k = top_k or RAG_TOP_K
    chunks = _get_rag_chunks(dirs)
    if not chunks:
        return ""
    query_tokens = _tokenize(question)
    scored = [(path, chunk, _score_chunk(query_tokens, chunk))
              for path, chunk in chunks]
    scored.sort(key=lambda x: x[2], reverse=True)
    top = scored[:top_k]
    if not top or top[0][2] == 0:
        return ""
    # 2026-07-22: 上位 top_k 件の中にスコア0（クエリトークンと無重複＝キーワード的に
    # 無関係）なチャンクが混ざっていても、以前は best さえ 0 でなければそのまま
    # 全件をプロンプトに注入していた。関係ないチャンクは各 proposer の回答を誤誘導
    # しうるノイズなので、精度優先（時間は気にしない）の方針に基づき score>0 の
    # チャンクのみを残す。降順順序と top_k 上限はそのまま、集合を縮小するだけ。
    top = [t for t in top if t[2] > 0]
    parts = []
    for path, chunk, score in top:
        parts.append(f"[Source: {Path(path).name}]\n{chunk.strip()}")
    return "## Relevant Document Context (RAG)\n\n" + "\n\n---\n\n".join(parts)


def build_context(question: str, use_search: bool = False,
                  rag_dirs: list = None) -> str:
    """Web検索 + RAG の結果を組み合わせてコンテキスト文字列を返す。
    空の場合は空文字を返す（質問がそのまま使われる）。"""
    parts = []
    if use_search:
        print("   [Web検索中...]")
        s = research_search(question)
        if s:
            n = s.count("Source:")
            print(f"   [Web検索: 計 {n} 件収集 ({len(s)} 文字)]")
            parts.append(s)
        else:
            print("   [警告: Web検索の結果が 0 件でした。回答はモデルの学習知識のみに"
                  "基づきます（最新情報は反映されません）]")
    rag_result = rag_search(question, dirs=rag_dirs or RAG_DIRS)
    if rag_result:
        parts.append(rag_result)
    return "\n\n".join(parts)


def _with_context(question: str, context: str) -> str:
    """コンテキストがあれば質問に前置する。"""
    if not context:
        return question
    return f"{context}\n\n---\n\n{question}"


def _trim_history(history):
    """_HISTORY が MAX_HISTORY_CHARS を超えたら古い (user, assistant) ペアを先頭から削除する。"""
    # 2026-07-23: ガードを `>= 2` から `> 2` に修正。
    # ask_fugu は直近の [user, assistant] ペアを追記した「直後」に本関数を呼ぶため、
    # 旧ガード `len(history) >= 2` だと履歴がその最新ペア1組だけ（長さ2）に
    # なった状態でも、文字数がまだ MAX_HISTORY_CHARS（4000。詳細なMoA回答は
    # コード/証明/節見出し込みですぐ超える）を超えていればループに入り、
    # pop(0) を2回叩いて最新ペアごと消してしまっていた。結果、直前に生成した
    # ばかりの回答が跡形もなく消え、次の行で「[会話履歴: 0 往復保持中]」と出て
    # 以降のフォローアップ質問が文脈ゼロで飛ぶ。docstring 通り「古いペアを
    # 先頭から削除する」のが目的であり、最新ペアまで削るのは仕様外。
    # `> 2` にすることで、削除対象が残り最新ペアのみになった時点でループを
    # 止め、多ターン対話の文脈（精度に直結）を必ず1ペアは残すようにする。
    while (sum(len(m["content"]) for m in history) > MAX_HISTORY_CHARS
           and len(history) > 2):
        history.pop(0)
        if history and history[0]["role"] == "assistant":
            history.pop(0)


def server_up():
    try:
        with urllib.request.urlopen(f"{OLLAMA_URL}/api/tags", timeout=5) as r:
            return r.status == 200
    except Exception:
        return False


def ensure_server():
    if server_up():
        return True
    if shutil.which("ollama") is None:
        print("⚠ ollama が見つかりません。https://ollama.com からインストールしてください。")
        return False
    print("[setup] Ollama サーバーが見つからないので起動を試みます…")
    try:
        subprocess.Popen(
            ["ollama", "serve"],
            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
        )
    except Exception as e:
        print("[setup] 自動起動に失敗:", e)
    for _ in range(30):
        if server_up():
            print("[setup] 起動を確認しました。")
            return True
        time.sleep(0.5)
    print("⚠ Ollama に接続できません。別ターミナルで `ollama serve` を起動してください。")
    return False


def installed_models():
    # 2026-07-25 (iter152): 素のリスト内包表記 [m["name"] for m in ...] は、
    # /api/tags の応答中の1件でも壊れている（要素がdictでない、"name"キーが
    # 無い、"name"が非文字列、"name"が空文字列）と内包表記の途中でその場で
    # 例外を出し、外側の except Exception: return [] が正常な要素も含めて
    # 導入済みモデル一覧を丸ごと空にしてしまっていた。これはiter103/111/
    # 112/113/139と同じ「1件の壊れた要素が全件を道連れにする」失敗形。
    # installed_models()はpull()で自己修復しない2箇所の判定に直結しており、
    # 空リストへの縮退は静かに精度を落とす: _arbitrateのis_installed(
    # ARBITER_MODEL, installed_models())は最上位知性モデルgpt-oss:120bを
    # 裁定チェーンへ加えるか否かを決めており、空リストだと裁定が弱い
    # フォールバックモデルへ静かに格下げされる（gotcha #7: SC投票の
    # tie-break劣化）。solve_verifiableのis_installed(SC_CHEAP_MODEL,
    # installed_models())も同様に安価な追加投票の有無を決めており、空
    # リストだと投票パネルが静かに薄くなる。ここでは壊れた要素だけを
    # 読み飛ばし、正常な要素は元の順序で回収する。
    try:
        with urllib.request.urlopen(f"{OLLAMA_URL}/api/tags", timeout=5) as r:
            data = json.loads(r.read().decode())
        if not isinstance(data, dict):
            return []
        models = data.get("models", [])
        if not isinstance(models, list):
            return []
        names = []
        for m in models:
            if not isinstance(m, dict):
                continue
            name = m.get("name")
            if isinstance(name, str) and name:
                names.append(name)
        return names
    except Exception:
        return []


def is_installed(model, inst):
    # 呼び出しは厳密タグで行うので、原則は厳密一致で判定（旧 startswith の誤検知を回避）。
    # タグ無し指定のときだけ :latest を許容する。
    cands = {model}
    if ":" not in model:
        cands.add(model + ":latest")
    return any(n in cands for n in inst)


def pull(model):
    print(f"\n[setup] 取得開始: {model}")
    print("※初回のみダウンロードが発生します。")
    try:
        rc = subprocess.run(["ollama", "pull", model]).returncode
    except FileNotFoundError:
        print("❌ ollama コマンドが見つかりません。")
        return False
    if rc == 0:
        print("[setup] 完了:", model)
        return True
    print(f"❌ 取得失敗: {model} (code {rc}) — スキップします。")
    return False


def resolve_models():
    inst = installed_models()
    pool = []
    for m in DESIRED_PROPOSERS:
        if is_installed(m, inst):
            print("[setup] OK (proposer)", m)
            pool.append(m)
        elif pull(m):
            pool.append(m)

    # アグリゲーター
    agg = DESIRED_AGGREGATOR
    if not (is_installed(agg, installed_models()) or agg in pool):
        if not pull(agg):
            print(f"⚠ アグリゲーター {agg} の取得に失敗 → プール先頭で代用します。")
            agg = pool[0] if pool else None

    # コンダクター（プロポーザー兼務なら追加ロード不要）
    cond = DESIRED_CONDUCTOR
    if not (is_installed(cond, installed_models()) or cond in pool or cond == agg):
        if not pull(cond):
            print(f"⚠ コンダクター {cond} の取得に失敗 → プール先頭で代用します。")
            cond = pool[0] if pool else agg
    if cond is None:
        cond = pool[0] if pool else agg

    # 全滅時の保険
    if not pool:
        print("[setup] 利用可能なモデルが無いため保険を取得します。")
        if pull(FALLBACK_MODEL):
            pool = [FALLBACK_MODEL]
            agg = agg or FALLBACK_MODEL
            cond = cond or FALLBACK_MODEL

    return pool, agg, cond

# ==================================================
# 推論ヘルパ
# ==================================================

# 2026-07-27: ask() の一過性失敗(HTTP 500等)リトライ予算を、固定2回・sleep(2)固定から
# 有界な指数バックオフに拡張する。
# 背景: このマシンは OLLAMA_MAX_LOADED_MODELS=1 (gotcha#5, GPUが8GB1枚のため)であり、
# SCバッチ境界・MoAパネル切替・arbiter呼び出しのたびに13~23GBのモデルをNVMeから
# 再ロードする。ask()自身のコメントが述べる「ロード直後の一過性500」が起きる
# windowは、この再ロード時間そのものであり2秒よりずっと広い。
# 旧予算(for attempt in (1, 2) と sleep(2)固定)は、iteration 9
# (145f285: think-strip再送がforループ境界で握り潰される別バグを修正した際も
# 「2回」の予算自体は意図的に変更しなかった)と iteration 35 (gotcha#1/#2の
# /api/chat・num_ctxピン留めのテスト整備と合わせて「ちょうど2回」をテストで
# 明示的に固定した)の2回にわたって、意図してそのまま据え置かれてきた値である。
# しかし、1回の再試行で救えなかった失敗は __ERROR__ として上位へ返り、
# _sc_sample は answer=None として扱うため、main_cot_count() の SC_MAX 投票予算
# (gotcha#7: 自己一貫性投票は数学/選択式問題の精度に直結する経路)が1票分
# 永続的に目減りする。MoA提案者の脱落・arbiterのフォールバックも同根。
# 精度優先・時間は気にしない方針の下、再ロードwindowを覆うだけ有界に広げるのは
# 純粋に精度側にプラスなので、ここを拡張する。
ASK_RETRY_ATTEMPTS = 4        # 旧: 固定2回 (iteration 9 / 35 が意図的にテストで固定していた値)
ASK_RETRY_BACKOFF = (2, 5, 10)  # 各要素はその回の失敗直後に待つ秒数。len == ASK_RETRY_ATTEMPTS - 1


def ask(model, messages, temperature, think=None, fmt=None, label=None, num_predict=None,
        num_ctx=None):
    """Ollama native /api/chat を叩く。num_ctx を必ず options で渡して context を安全域に固定する
    （/v1 互換エンドポイントは num_ctx を無視するため使わない）。失敗時は __ERROR__: を返す。
    一過性の失敗(HTTP 500 等)には ASK_RETRY_ATTEMPTS 回まで、ASK_RETRY_BACKOFF の間隔で
    再試行する(2026-07-27: 固定2回・sleep(2)固定から有界バックオフへ拡張。詳細はモジュール
    冒頭の ASK_RETRY_ATTEMPTS/ASK_RETRY_BACKOFF 付近のコメントを参照)。
    num_predict: 生成トークン上限（None=無制限）。思考モデルの暴走保険として役割別に渡す。
    num_ctx: コンテキスト長の明示指定。None なら MODEL_CONFIG > MODEL_NUM_CTX の順で解決。
      「必ず明示 pin する」不変条件は維持（未指定だと Ollama がモデル最大を確保して 8GB VRAM で
      クラッシュするため）。

    think: None=モデル既定(MODEL_CONFIG があればそれを適用) / False=無効 / True=有効 /
      "low"|"medium"|"high"=gpt-oss 系の思考量段階指定（0.31.2 実測で有効）。
      注意: native /api/chat の "think" パラメータは効くが、プロンプトに書く "/no_think" は
      このルートでは無視される（実測）。gemma4(e2b/e2b-it-qat) は think=true/false とも
      正常に受け付ける（2026-07-02 実測。思考は message.thinking に分離される）。
      非thinkingモデル phi4-mini も False を渡して無害（実測）。
    fmt: Ollama の "format"。"json" か JSON スキーマ(dict)を渡すと構造化出力を強制できる。
      Conductor/Critic では think=False + スキーマの併用が要点:
      think=False 単独だと qwen3 は思考を content に地の文で垂れ流して JSON が壊れるが、
      スキーマを与えると enum 値まで含めて妥当な JSON に拘束され、かつ高速（実測 ~14s）。"""
    if think is None:
        think = model_cfg(model, "think")   # 呼び出し側が未指定ならモデル別設定を適用
    payload = {
        "model": model,
        "messages": messages,
        "stream": False,
        "options": {"temperature": temperature,
                    "num_ctx": num_ctx or model_cfg(model, "num_ctx", MODEL_NUM_CTX)},
    }
    if num_predict is not None:
        payload["options"]["num_predict"] = num_predict
    if think is not None:
        payload["think"] = think
    if fmt is not None:
        payload["format"] = fmt
    if MODEL_KEEP_ALIVE is not None:
        payload["keep_alive"] = MODEL_KEEP_ALIVE  # 既定 None では渡さない
    req = urllib.request.Request(
        f"{OLLAMA_URL}/api/chat",
        data=json.dumps(payload).encode("utf-8"),
        headers={"Content-Type": "application/json"},
        method="POST",
    )
    t0 = time.time()
    out = "__ERROR__: unreachable"

    def _do_call(request):
        with urllib.request.urlopen(request, timeout=REQUEST_TIMEOUT) as r:
            body = json.loads(r.read().decode("utf-8"))
        msg = body.get("message") or {}
        # think 分離型モデルは thinking が別フィールドに来る場合があるので content のみ採用。
        result = (msg.get("content") or "").strip()
        # 思考が num_predict を食い尽くして本文ゼロで打ち切られた場合(実測: 空返答の
        # 根本原因)は、沈黙の空文字ではなく明示的なエラーにして上位のフォールバックを
        # 確実に発動させる。本文が一部でも出ていればそのまま使う。
        if not result and body.get("done_reason") == "length":
            result = ("__ERROR__: truncated by num_predict during thinking "
                      "(no content was generated)")
        return result

    for attempt in range(1, ASK_RETRY_ATTEMPTS + 1):
        try:
            out = _do_call(req)
            break
        except urllib.error.HTTPError as e:
            err_body = e.read().decode("utf-8", errors="replace")
            # think:true を送ったが非対応モデル（例: qwen3-coder, phi4）→ think なしで即リトライ。
            # これは設定ミスの安全網。通常は PROPOSER_THINK=None で発生しないはず。
            if e.code == 400 and "does not support thinking" in err_body and "think" in payload:
                payload.pop("think")
                req = urllib.request.Request(
                    f"{OLLAMA_URL}/api/chat",
                    data=json.dumps(payload).encode("utf-8"),
                    headers={"Content-Type": "application/json"},
                    method="POST",
                )
                # 重要(バグ修正): このブランチはかつて continue で外側の for に戻り、
                # 「次のイテレーション」が実際に送信されることに依存していた。しかし
                # attempt==1 が一過性の 500（ロード直後）で、attempt==2 で初めて
                # 「thinking非対応」400 が出た場合、continue しても for (1, 2) は
                # 既に尽きており、組み直したリクエストは一度も送信されずに
                # ループを抜けて __ERROR__: think_stripped_retry がそのまま最終値に
                # なってしまっていた（SC投票/提案が黙って1票失われる）。
                # 修正: think を pop 済みなのでこの分岐は高々1回しか到達し得ない
                # （無限ループの可能性なし）。よってここでその場で確定的に1回だけ
                # 追加送信し、一過性リトライの残り予算（sleep(2)して attempt を
                # 消費する経路）は一切消費しない。
                try:
                    out = _do_call(req)
                except urllib.error.HTTPError as e2:
                    err_body2 = e2.read().decode("utf-8", errors="replace")
                    out = f"__ERROR__: {e2} {err_body2}"
                except Exception as e2:
                    out = f"__ERROR__: {e2}"
                break  # think は pop済みでこの分岐は再発し得ないため、ここで確定終了
            out = f"__ERROR__: {e} {err_body}"
            if attempt < ASK_RETRY_ATTEMPTS:
                time.sleep(ASK_RETRY_BACKOFF[attempt - 1])
        except Exception as e:
            out = f"__ERROR__: {e}"
            if attempt < ASK_RETRY_ATTEMPTS:
                # 一過性の失敗(ロード直後の500等)向け。ASK_RETRY_ATTEMPTS 回失敗したら諦める
                time.sleep(ASK_RETRY_BACKOFF[attempt - 1])
    if SHOW_TIMING:
        _TIMINGS.append((label or "?", model, round(time.time() - t0, 1)))
    return out


# deepseek-r1 / qwen3 などの思考ログを除去。
# 注: Gemma 4 の思考は 2026-07-02 実測で /api/chat が message.thinking に分離して返す
#     （content には混入しない）ため、Gemma 用の除去パターン追加は不要。
#     gemma4:e2b / e2b-it-qat とも think=true/false パラメータが正常に効くことも確認済み。
_THINK_PATTERNS = [
    re.compile(r"<think>.*?</think>", re.DOTALL | re.IGNORECASE),
    re.compile(r"<thinking>.*?</thinking>", re.DOTALL | re.IGNORECASE),
]
# 2026-07-22: num_predict 打ち切り（gotcha #2）で思考モデルの出力が
# '<think>途中の推論...' のまま閉じタグが来ずに切れることがある。上の
# _THINK_PATTERNS は非貪欲マッチのため閉じタグが無いと一切マッチせず、
# 破棄されるべき生の思考過程がそのまま「回答」として extract_final_answer
# まで漏れ、抽出失敗 → nums[-1] フォールバックで思考中の中間値が SC 投票の
# 1票として数えられてしまう(extract_boxed の #2214 と同根の失敗モード)。
# 対策として、バランス除去の後に「閉じタグの無い開始タグ」を検出したら、
# その開始タグ以降を末尾まで丸ごと切り捨てる（開始タグより前のテキストは
# 保持）。無投票の方が誤投票より安全という方針（精度優先）に従う。
_UNTERMINATED_THINK_OPEN = re.compile(r"<think(?:ing)?>", re.IGNORECASE)


def strip_think(text):
    if not text:
        return text
    for pat in _THINK_PATTERNS:
        text = pat.sub("", text)
    m = _UNTERMINATED_THINK_OPEN.search(text)
    if m:
        text = text[:m.start()]
    return text.strip()


def extract_json(text):
    """モデル出力から最初の JSON オブジェクトを頑健に抽出（失敗時 None）。"""
    if not text:
        return None
    text = strip_think(text)
    # 1) そのまま
    # 2026-07-25: 旧実装はここで json.loads(text) の結果を型を問わず素通しで
    # return していた。docstring は「最初の JSON オブジェクト」（dict）を約束し、
    # 2)/3) は実際に dict-or-None 契約を守っているのに、1) だけはトップレベルが
    # 妥当な JSON でありさえすれば list/int/float/bool/str/None も丸ごと返して
    # しまい、契約が破られていた。呼び出し側の一部（_critic_judge L2662,
    # second_opinion L2710, _sd_prompt_from_request L1944, plan_pptx_images
    # L4494）は `extract_json(raw) or {}` の後に無条件で `.get(...)` するため、
    # モデルが `[{"ok": true}]` や `true`/`42`/`"text"` のようなトップレベル非
    # object な（しかし妥当な）JSON を出力すると、truthy な非dict値がそのまま
    # 通り抜けて `.get` で AttributeError を送出し、精度が最も重要な
    # critique/verify ゲートを含むターン全体をクラッシュさせていた。
    # これは iteration 103/111/112/113/138/139 で修正してきた「妥当だが型が
    # 想定と違うモデル出力」への防御的対処と同種のクラス。
    # 修正: 1) はパース結果が isinstance(dict) の場合のみ return し、それ以外
    # （list/int/float/bool/str/None）は return も raise もせず 2)/3) にフォール
    # スルーする。これにより例えば `[{"ok": true}]` は 3) の波括弧スキャナが
    # ネストした {"ok": true} を回収できるようになる（副次的な改善）。
    try:
        result = json.loads(text)
        if isinstance(result, dict):
            return result
    except Exception:
        pass
    # 2) ```json ... ``` フェンス内
    m = re.search(r"```(?:json)?\s*(\{.*?\})\s*```", text, re.DOTALL)
    if m:
        try:
            return json.loads(m.group(1))
        except Exception:
            pass
    # 3) 最初のバランスの取れた { ... } を波括弧の深さで走査して探す
    # 2026-07-22: 旧実装は re.search(r"\{.*\}", ..., re.DOTALL) という貪欲マッチで、
    # 最初の '{' から最後の '}' までを一括で span にしていた。本文中に地の文の
    # 集合記法 "{1,2,3}" や末尾の "{x}"、あるいは2つ目のJSONオブジェクトなど
    # 「余分な波括弧」が存在すると、その span 全体は JSON として不正な形になり
    # json.loads が例外を投げて None を返していた（本来 docstring が約束する
    # 「最初の JSON オブジェクト」を回収可能なのに握りつぶす）。
    # これは conduct() のルーティングプランと research_search() の
    # RESEARCH_SCHEMA 充足判定の両方を壊し、None が来た側で default_plan() への
    # 劣化や、リサーチの誤った早期終了を引き起こしていた。
    # 修正: 文字列中の状態（ダブルクォート内かどうか・直前のバックスラッシュに
    # よるエスケープ）を追跡しながら、'{' を見つけるたびに深さカウントで対応する
    # '}' を探し、最初に json.loads が成功した部分文字列を返す。
    n = len(text)
    i = 0
    while i < n:
        if text[i] == "{":
            depth = 0
            in_string = False
            escape = False
            j = i
            while j < n:
                ch = text[j]
                if in_string:
                    if escape:
                        escape = False
                    elif ch == "\\":
                        escape = True
                    elif ch == '"':
                        in_string = False
                else:
                    if ch == '"':
                        in_string = True
                    elif ch == "{":
                        depth += 1
                    elif ch == "}":
                        depth -= 1
                        if depth == 0:
                            candidate = text[i:j + 1]
                            try:
                                return json.loads(candidate)
                            except Exception:
                                break  # この候補は失敗、次の '{' から再探索
                j += 1
            # depth が閉じ切らなかった（切り詰められた/不正）場合もここに来る
        i += 1
    return None

# ==================================================
# コード実行検証（決定的 Critic）
# ==================================================


def extract_code(text):
    """回答から最初の ```python コードブロックを抽出（無ければ None）。
    言語タグ無しフェンスも Python とみなして拾う（proposer には python タグを指示済み）。
    2026-07-22: 旧実装は re.search(r"```(?:python|py)?[ \t]*\n(.*?)```") を使っており、
    先行する非python フェンス（```json/```bash/```text/```output 等）の開始タグに
    マッチできず、re.search が前方走査してそのブロックの「閉じフェンス」を開始フェンス
    と誤認し、2つのブロックの間にあるプロース（地の文）や本文を「コード」として誤抽出
    していた。これにより code_check が非コードを実行して見せかけの実行失敗を報告し、
    無駄な修復ラウンド(MAX_ROUNDS_CODE)を消費したり、PoT の投票が壊れたりしていた。
    修正: 全てのフェンスブロックを走査し、言語タグが python/py/python3 または
    タグ無し(bare)の「最初のブロック」の本文をそのまま(strip無し)返す。それ以外の
    タグ(json/bash/sh/text/output/js 等)のブロックは読み飛ばす。該当ブロックが
    無ければ None。"""
    if not text:
        return None
    # 2026-07-22: CommonMark の info string 仕様では、フェンス開始行の残り全体では
    # なく「最初の空白区切りトークン」だけが言語タグであり、それ以降は任意の
    # メタデータ（例: ```python {.line-numbers} や ```python title="sol.py"）。
    # 旧実装は strip した info string 全体を ("", "python", "py", "python3") と
    # 比較していたため、```python {.line-numbers} は lang=='python {.line-numbers}'
    # となって不一致になり、正当な python ブロックが読み飛ばされ、extract_code が
    # 後続の無関係なブロックを拾うか None を返していた（本関数のブロック選択自体は
    # iteration 7 で修正済みで、ここではその値の比較方法のみを直す）。
    # 修正: info string の最初のトークンのみを言語タグとして比較する（空/空白のみの
    # info string は従来通りタグ無し("")として bare フェンス扱い）。受理集合は
    # ("", "python", "py", "python3") のまま変更しない。json/bash 等、最初のトークン
    # が非python なフェンスは装飾の有無に関わらず引き続き読み飛ばされる。
    # なお _extract_code_for_output（iteration 18 の多段優先ロジック）はこのイテレー
    # ションではあえて触れていない — 同種の info-string 対応は別途の追随課題とする。
    for m in re.finditer(r"```([^\n`]*)\n(.*?)```", text, re.DOTALL):
        info = m.group(1).strip().lower()
        lang = info.split(None, 1)[0] if info else ""
        if lang in ("", "python", "py", "python3"):
            return m.group(2)
    return None


def run_python(code, timeout=None, stdout_only=False):
    """コードを一時ファイル経由で subprocess 実行する。(ok: bool, output: str) を返す。
    ok は exit code 0。既定(stdout_only=False)では stdout+stderr 結合の末尾を返す
    （traceback を修正ヒントに使うため）。stdout_only=True かつ成功時(returncode==0)
    は stdout のみを返す（sympy/numpy の DeprecationWarning 等が stderr に出て
    末尾行が汚染され、PoT の投票が壊れるのを防ぐため）。ただし失敗時(returncode!=0)
    は stdout_only の値によらず常に stdout+stderr の結合を返す — code-repair loop が
    traceback を見えるようにするため。"""
    timeout = timeout or CODE_EXEC_TIMEOUT
    if os.environ.get("FUGU_SANDBOX") == "1":
        try:
            import fugu_sandbox  # opt-in: 既定経路はこの分岐に一切入らない
        except ImportError:
            fugu_sandbox = None
        if fugu_sandbox is not None:
            res = fugu_sandbox.SubprocessSandbox(timeout=timeout).run(code)
            out = res.stdout if (stdout_only and res.ok) else res.output
            return res.ok, out.strip()[-2000:]
    fd, path = tempfile.mkstemp(suffix=".py")
    try:
        with os.fdopen(fd, "w", encoding="utf-8") as fh:
            fh.write(code)
        # 2026-07-22: stdin=DEVNULL を明示。指定しないと子プロセスは親の stdin を
        # そのまま継承し、LLM 生成コードが input() を呼ぶケース（PoT の _sc_sample や
        # code_check の検証でよくある）でハングしてしまう — repl() では対話中の親の
        # stdin を子に奪われる事故にもなり、挙動は TTY/pipe/closed のいずれかで非決定的。
        # DEVNULL にしておけば input() は即座に EOFError を送出して fail-fast する。
        # input() を要求するコードはどのみちオフラインでは正しい答えを返せないので、
        # 正解率への影響はない（無効票 / NG のままで変わらない）— 純粋にハング・
        # stdin 汚染・非決定性の除去。
        r = subprocess.run(
            [sys.executable, "-X", "utf8", path],
            capture_output=True, text=True, encoding="utf-8", errors="replace",
            timeout=timeout, stdin=subprocess.DEVNULL,
        )
        if stdout_only and r.returncode == 0:
            out = (r.stdout or "").strip()
        else:
            out = ((r.stdout or "") + (r.stderr or "")).strip()
        return r.returncode == 0, out[-2000:]
    except subprocess.TimeoutExpired:
        return False, f"TIMEOUT: code did not finish within {timeout}s (infinite loop, not input() — stdin is DEVNULL so input() now fails fast with EOFError instead)"
    except Exception as e:
        return False, f"runner error: {e}"
    finally:
        try:
            os.unlink(path)
        except OSError:
            pass


def code_check(answer):
    """回答中の Python コードを実行して検証する。問題があればエラー要約(str)、
    問題なし・コードなし・機能OFF なら None を返す。"""
    if not CODE_EXECUTION:
        return None
    code = extract_code(answer)
    if not code:
        return None
    ok, out = run_python(code)
    if ok:
        return None
    return f"code execution FAILED:\n{out[-800:]}"

# ==================================================
# 画像生成（ローカル SD / ComfyUI へバイパス）
# ==================================================


def _sd_prompt_from_request(user_request):
    """ユーザーの画像要求を SD 用プロンプトへ変換する。IMAGE_TRANSLATE_PROMPT なら
    qwen3:4b で英語の高品質プロンプト（+ネガティブ）へ翻訳する（SD系は英語で精度が出る）。
    戻り値: (prompt, negative)。"""
    if not IMAGE_TRANSLATE_PROMPT:
        return user_request, ""
    sys = (
        "You convert a user's image request into a Stable Diffusion prompt. "
        'Output ONLY JSON: {"prompt": "...", "negative": "..."}. '
        "prompt: a concise comma-separated ENGLISH prompt with quality tags "
        "(e.g. 'masterpiece, best quality, highly detailed'). "
        "negative: common negatives (e.g. 'lowres, bad anatomy, blurry, watermark'). "
        "No prose, no thinking."
    )
    raw = ask(
        CONDUCTOR,
        [{"role": "system", "content": sys},
         {"role": "user", "content": user_request}],
        CONDUCTOR_TEMP, think=False,
        fmt={"type": "object",
             "properties": {"prompt": {"type": "string"},
                            "negative": {"type": "string"}},
             "required": ["prompt"]},
        num_predict=512, label="img-prompt",
    )
    j = extract_json(raw) or {}
    prompt = str(j.get("prompt") or user_request).strip()
    negative = str(j.get("negative") or "").strip()
    return prompt, negative


def _http_post_json(url, payload, timeout):
    req = urllib.request.Request(
        url, data=json.dumps(payload).encode("utf-8"),
        headers={"Content-Type": "application/json"}, method="POST",
    )
    with urllib.request.urlopen(req, timeout=timeout) as r:
        return json.loads(r.read().decode("utf-8"))


def _backend_up(url):
    """GET で 200 が返れば True（バックエンドの疎通確認）。"""
    try:
        with urllib.request.urlopen(url, timeout=3) as r:
            return r.status == 200
    except Exception:
        return False


def generate_image_a1111(prompt, negative=""):
    """AUTOMATIC1111 stable-diffusion-webui の txt2img API で生成し保存パスを返す。"""
    import base64
    import uuid
    payload = {"prompt": prompt, "negative_prompt": negative,
               "steps": IMAGE_STEPS, "width": IMAGE_WIDTH, "height": IMAGE_HEIGHT}
    data = _http_post_json(f"{A1111_URL}/sdapi/v1/txt2img", payload, IMAGE_TIMEOUT)
    # 2026-07-25 (iter144): iter139でComfyUI側(generate_image_comfyui)のoutputs/image
    # エントリ走査に施したskip-and-recoverと対称の穴がこちら(A1111側)にも残っていた。
    # data.get("images")はdata自体がdictである前提、images[0]はimagesがnon-emptyな
    # listである前提、images[0].split(",", 1)は先頭要素がstrである前提で、いずれも
    # 無保証。非dictなdata・truthyだが非listなimages・非str/空/デコード不能な先頭
    # エントリはAttributeError/TypeError/binascii.Errorを送出し、呼び出し元
    # generate_imageの外側except Exceptionまで伝播、後続に有効な画像が残っていても
    # 生成結果ごと握り潰してNoneになっていた。iter103/111/112/113/138（非list/非dict
    # の強制truthy変換に頼らない既定値フォールバックとentry単位skip）と同じ作法に
    # 倣い、dataが非dictならNone、imagesが非listなら[]へ矯正し、先頭決め打ちではなく
    # リストを走査して非str/空/デコードまたは書き込み失敗エントリを1件ずつskipして
    # 後続の有効な画像を回収する。
    if not isinstance(data, dict):
        return None
    images = data.get("images")
    if not isinstance(images, list):
        images = []
    for img in images:
        if not isinstance(img, str) or not img:
            continue
        try:
            # data URI プレフィックスが付く場合があるのでカンマ以降を取る
            blob = base64.b64decode(img.split(",", 1)[-1])
            IMAGE_OUT_DIR.mkdir(parents=True, exist_ok=True)
            # 2026-07-25: time.strftime('%Y%m%d_%H%M%S')は秒単位までしか刻まない。
            # build_pptxはPPTX_MAX_IMAGES=4枚まで本関数を連続呼び出しし、SDXL生成が
            # 同一秒内に完了すると旧実装のfugu_{ts}.pngは2枚目以降が1枚目と同じ
            # パスに解決し、後発のout.write_bytes()が先発の画像を無言で上書きして
            # 消していた。8GB GPUの現行機では数秒かかるため滅多に起きないが、
            # apply_high_vram_profileが想定する96GB高VRAM環境の高速SDXLでは
            # 同一秒内の衝突が現実的になる。iter77/139/144と同系統の
            # 「既に生成できたアーティファクトを無言で失わない」原則に倣い、
            # uuid4由来の8桁hexを付与して秒精度に依存しない一意性を保証する。
            uniq = uuid.uuid4().hex[:8]
            out = IMAGE_OUT_DIR / f"fugu_{time.strftime('%Y%m%d_%H%M%S')}_{uniq}.png"
            out.write_bytes(blob)
            return out
        except Exception:
            continue
    return None


def generate_image_comfyui(prompt, negative=""):
    """ComfyUI に最小 txt2img ワークフローを投げて生成し保存パスを返す。
    COMFYUI_CKPT が空ならサーバの利用可能チェックポイント先頭を自動採用する。"""
    import uuid
    ckpt = COMFYUI_CKPT
    if not ckpt:
        try:
            with urllib.request.urlopen(
                    f"{COMFYUI_URL}/object_info/CheckpointLoaderSimple", timeout=5) as r:
                info = json.loads(r.read().decode("utf-8"))
            ckpt = info["CheckpointLoaderSimple"]["input"]["required"]["ckpt_name"][0][0]
        except Exception as e:
            print(f"   [ComfyUI: チェックポイント取得に失敗: {e}]")
            return None
    client_id = uuid.uuid4().hex
    wf = {
        "3": {"class_type": "KSampler", "inputs": {
            "seed": int(time.time()) % (2 ** 32), "steps": IMAGE_STEPS, "cfg": 7,
            "sampler_name": "euler", "scheduler": "normal", "denoise": 1,
            "model": ["4", 0], "positive": ["6", 0], "negative": ["7", 0],
            "latent_image": ["5", 0]}},
        "4": {"class_type": "CheckpointLoaderSimple", "inputs": {"ckpt_name": ckpt}},
        "5": {"class_type": "EmptyLatentImage", "inputs": {
            "width": IMAGE_WIDTH, "height": IMAGE_HEIGHT, "batch_size": 1}},
        "6": {"class_type": "CLIPTextEncode", "inputs": {"text": prompt, "clip": ["4", 1]}},
        "7": {"class_type": "CLIPTextEncode", "inputs": {"text": negative, "clip": ["4", 1]}},
        "8": {"class_type": "VAEDecode", "inputs": {"samples": ["3", 0], "vae": ["4", 2]}},
        "9": {"class_type": "SaveImage", "inputs": {"filename_prefix": "fugu", "images": ["8", 0]}},
    }
    try:
        resp = _http_post_json(f"{COMFYUI_URL}/prompt",
                               {"prompt": wf, "client_id": client_id}, 30)
        pid = resp["prompt_id"]
    except Exception as e:
        print(f"   [ComfyUI: プロンプト投入に失敗: {e}]")
        return None
    # 完了までポーリング
    deadline = time.time() + IMAGE_TIMEOUT
    hist = None
    while time.time() < deadline:
        try:
            with urllib.request.urlopen(f"{COMFYUI_URL}/history/{pid}", timeout=10) as r:
                h = json.loads(r.read().decode("utf-8"))
            if pid in h:
                hist = h[pid]
                break
        except Exception:
            pass
        time.sleep(2)
    if not hist:
        print("   [ComfyUI: 生成がタイムアウトしました]")
        return None
    for node in hist.get("outputs", {}).values():
        for img in node.get("images", []):
            # 2026-07-25 (iter139): 隣接するsubfolder/typeは.get()で守られている
            # 一方、filenameだけはimg["filename"]の直接dictアクセスが2箇所に
            # 残っていた。ComfyUI /history の壊れたエントリ（filenameキー欠落・
            # null・またはimages配列内の非dict要素、例えば文字列/数値）に対して
            # 旧実装はKeyError/TypeErrorを送出し、それが本関数を丸ごと巻き込んで
            # 呼び出し元generate_imageの外側except Exceptionまで伝播、以降の
            # ノード・エントリに有効な画像があっても生成結果ごと握り潰して
            # Noneを返していた。iter77（良い方を回収する）・iter103/111/112
            # （非list/非dictの強制truthy変換に頼らない既定値フォールバック）・
            # iter113/iter138（外部由来ペイロードの1件の破損で全体を道連れに
            # しないためのentry単位skip）と同じ作法に倣い、ここでもfilenameを
            # img.get()で一度だけ読み、非dictエントリ・欠落/None/空/非str
            # filenameは例外を出さずcontinueして後続の有効なエントリを救済する。
            if not isinstance(img, dict):
                continue
            filename = img.get("filename")
            if not isinstance(filename, str) or not filename:
                continue
            q = urllib.parse.urlencode({"filename": filename,
                                        "subfolder": img.get("subfolder", ""),
                                        "type": img.get("type", "output")})
            try:
                with urllib.request.urlopen(f"{COMFYUI_URL}/view?{q}", timeout=30) as r:
                    blob = r.read()
            except Exception:
                continue
            IMAGE_OUT_DIR.mkdir(parents=True, exist_ok=True)
            # 2026-07-25: time.strftime('%Y%m%d_%H%M%S')は秒単位までしか刻まない。
            # ComfyUI側のfilenameは同一seed/同一prefixのSaveImageノードが同一秒内に
            # 複数回走ると衝突しうる（サーバ再起動でComfyUI内部のカウンタがリセット
            # される場合を含む）ため、fugu_{ts}_{filename}だけでは一意性が保証
            # できない。A1111側(generate_image_a1111、同日付コメント参照)と対称の
            # 同一秒上書き問題であり、iter77/139/144と同じ「既に生成できた
            # アーティファクトを無言で失わない」原則に倣い、uuid4由来の8桁hexを
            # ComfyUI側filenameの前に挟んで秒精度に依存しない一意性を保証する。
            uniq = uuid.uuid4().hex[:8]
            out = IMAGE_OUT_DIR / f"fugu_{time.strftime('%Y%m%d_%H%M%S')}_{uniq}_{filename}"
            out.write_bytes(blob)
            return out
    return None


# --- LLM 群による画像プロンプト起草 ---
IMAGE_PROMPT_SYS = (
    "あなたは Stable Diffusion(SDXL) のプロンプトエンジニアです。"
    "与えられた要望(または回答内容)を、高品質な画像を生成するための英語プロンプトへ変換します。"
    'JSON のみ出力: {"prompt": "...", "negative": "..."}. '
    "prompt: カンマ区切りの英語。主題・画風・構図・光・品質タグ"
    "(masterpiece, best quality, highly detailed 等)を含める。"
    "negative: 典型的なネガティブ(lowres, bad anatomy, blurry, watermark, text 等)。"
    "散文や思考は出力しない。"
)
_IMG_PROMPT_SCHEMA = {
    "type": "object",
    "properties": {"prompt": {"type": "string"}, "negative": {"type": "string"}},
    "required": ["prompt"],
}


def moa_image_prompt(base_text, panel=None):
    """LLM 群(proposers)が SDXL プロンプト候補を起草し、qwen3:4b が最良の1つへ統合する。
    候補が得られなければ None（呼び出し側が単独翻訳へフォールバック）。"""
    models = [m for m in (panel or []) if m in PROPOSERS] or PROPOSERS[:IMAGE_PROMPT_PANEL]
    models = models[:IMAGE_PROMPT_PANEL]
    cands = []
    for m in models:
        raw = ask(m, [{"role": "system", "content": IMAGE_PROMPT_SYS},
                      {"role": "user", "content": base_text}],
                  PROPOSER_TEMP, think=False, fmt=_IMG_PROMPT_SCHEMA,
                  num_predict=512, label="img-moa")
        j = extract_json(raw)
        if isinstance(j, dict) and j.get("prompt"):
            cands.append((m, j))
    if not cands:
        return None
    if len(cands) == 1:
        j = cands[0][1]
        return str(j.get("prompt")).strip(), str(j.get("negative") or "").strip()
    listing = "\n".join(
        f"[{MODEL_TO_PERSONA.get(m, m)}] prompt={j.get('prompt')} | negative={j.get('negative', '')}"
        for m, j in cands)
    merge_sys = (IMAGE_PROMPT_SYS + " 以下は複数の専門家が起草した候補です。"
                 "最も的確で高品質な1つの SDXL プロンプトへ統合しなさい。")
    raw = ask(CONDUCTOR, [{"role": "system", "content": merge_sys},
                          {"role": "user", "content": base_text + "\n\n候補:\n" + listing}],
              CONDUCTOR_TEMP, think=False, fmt=_IMG_PROMPT_SCHEMA,
              num_predict=512, label="img-merge")
    j = extract_json(raw)
    if isinstance(j, dict) and j.get("prompt"):
        return str(j.get("prompt")).strip(), str(j.get("negative") or "").strip()
    j = cands[0][1]
    return str(j.get("prompt")).strip(), str(j.get("negative") or "").strip()


def author_image_prompt(base_text, panel=None):
    """画像プロンプトを決める。IMAGE_PROMPT_MOA なら LLM 群で起草→統合、
    失敗時や無効時は qwen3:4b の単独翻訳へフォールバック。戻り値 (prompt, negative)。"""
    if IMAGE_PROMPT_MOA:
        p = moa_image_prompt(base_text, panel)
        if p:
            return p
    return _sd_prompt_from_request(base_text)


# --- バックエンド検出と生成コア ---
def _detect_backend():
    """使用するバックエンドを返す（"a1111" / "comfyui" / None）。"""
    if IMAGE_BACKEND in ("a1111", "comfyui"):
        return IMAGE_BACKEND
    if IMAGE_BACKEND == "auto":
        if _backend_up(f"{A1111_URL}/sdapi/v1/sd-models"):
            return "a1111"
        if _backend_up(f"{COMFYUI_URL}/system_stats"):
            return "comfyui"
    return None


def generate_image(prompt, negative=""):
    """プロンプトから実際に画像を生成し保存 Path を返す（失敗は None）。
    PowerPoint ビルダーからも直接使う低レベル API。"""
    backend = _detect_backend()
    if backend is None:
        return None
    try:
        if backend == "a1111":
            return generate_image_a1111(prompt, negative)
        return generate_image_comfyui(prompt, negative)
    except Exception as e:
        print(f"   [画像生成エラー ({backend}): {e}]")
        return None


def handle_image_generation(user_request, *, panel=None, prompt=None, negative=None):
    """画像を生成し、人間可読な結果テキスト（保存パス）または __ERROR__ を返す。
    prompt/negative 未指定なら LLM 群で起草する。panel はプロンプト起草に使う proposer 群。"""
    if IMAGE_BACKEND == "off":
        return "__ERROR__: 画像生成は無効化されています (IMAGE_BACKEND='off')。"
    if _detect_backend() is None:
        return ("__ERROR__: 画像生成バックエンドが見つかりません。\n"
                f"  AUTOMATIC1111 を {A1111_URL} で、または ComfyUI を {COMFYUI_URL} で"
                "起動してください（IMAGE_BACKEND で明示指定も可）。")
    if prompt is None:
        prompt, negative = author_image_prompt(user_request, panel=panel)
    print(f"   [画像生成プロンプト] {prompt[:120]}")
    out = generate_image(prompt, negative or "")
    if not out:
        return "__ERROR__: 画像生成に失敗しました（バックエンドが画像を返しませんでした）。"
    msg = (f"画像を生成しました。\n"
           f"- 保存先: {out}\n"
           f"- prompt: {prompt}\n")
    if negative:
        msg += f"- negative: {negative}\n"
    return msg

# ==================================================
# プロンプト
# ==================================================

# ChatGPT / Gemini のような読みやすい提示スタイル（最終回答に付与する）
PRESENTATION_STYLE = (
    "\n\n【回答の体裁（ChatGPT / Gemini 風）】\n"
    "- 最初に結論・要点を1〜2文で述べ、その後に詳細を続ける。\n"
    "- Markdown を適切に使う: 見出し(##)、箇条書き(-)、番号付き手順、重要語の**強調**、"
    "コードは``` で囲む、比較や一覧は表を使う。\n"
    "- 丁寧で分かりやすい対話的なトーン。冗長な前置きや『AIとして』等の自己言及は避ける。\n"
    "- 長い回答は見出しで構造化し、必要なら最後に一言まとめを付ける。質問と同じ言語で書く。"
)

PROPOSER_SYS = (
    "You are one expert in a panel. Answer the user's question as accurately and "
    "concretely as you can. Think step-by-step internally to avoid math or logical errors. "
    "Briefly show key reasoning if it helps. Respond in the same language as the question.\n"
    "If the question asks for code: put ONE complete, runnable program in a single "
    "```python code block, and end that same block with a few assert-based self-tests "
    "that verify it (the code will be executed to check it works). No input() calls."
)

AGGREGATOR_SYS = (
    "You are the aggregator of a Mixture-of-Agents system. Several independent "
    "models each answered the SAME question; their answers are given below, "
    "anonymized as Answer A, B, C, ...\n\n"
    "Do NOT merely summarize or average them. Produce ONE answer that is more "
    "accurate and complete than any single one, by reasoning critically:\n"
    "1. Where they AGREE: a signal of likely correctness — but verify it is not a shared mistake.\n"
    "2. Where they DISAGREE or CONTRADICT: treat it as a red flag. Decide which side is actually correct using your own reasoning, not majority vote.\n"
    "3. Detect and DISCARD errors, hallucinations, unsupported claims, even if several answers share them.\n"
    "4. PRESERVE the strongest reasoning and the most useful specifics; combine them.\n"
    "5. If all answers are weak, override them with a better one.\n"
    "6. If an answer carries an [Execution check: ...] tag, that is GROUND TRUTH from "
    "actually running its code — prefer code that PASSED; never base the final answer on "
    "code that FAILED without fixing the reported error.\n\n"
    "Output only the final polished answer for the user. Do not mention the other answers or this process.\n"
    "CRITICAL LANGUAGE RULE: Write the ENTIRE final answer in exactly the same language as the "
    "question. Do NOT mix in words or characters from any other language. For a Japanese question, "
    "use natural Japanese only — never Chinese-only words/characters (e.g. 个, 至少, 确, 说) and no "
    "stray English words; use the Japanese equivalents (個, 少なくとも, など)."
    + PRESENTATION_STYLE
)

CONDUCTOR_SYS = (
    "あなたは Fugu オーケストレーションシステムの最高司令塔(Conductor)です。"
    "ユーザーの入力(および RAG で読み込まれた Office ファイルや Web 検索結果)を分析し、"
    "最適な実行プランを厳密な JSON で出力します。あなた自身は質問に答えません。\n\n"
    "【専門家チームの布陣】\n"
    "- Proposer A (ChatGPT/GPTの存在): バランス・一般的な対話・文章の骨組み担当。\n"
    "- Proposer B (Claudeの存在): 高度なプログラミング・厳密な論理チェック・コード自己修復担当。\n"
    "- Proposer C (Geminiの存在): RAG(Officeファイル)のコンテキスト分析・大量ドキュメント・"
    "Web 検索結果の集約担当。\n"
    "- Proposer D (理数専門家): 数学・物理・PINN・偏微分方程式・アルゴリズム証明担当。\n\n"
    "【特殊ルーティング指示】\n"
    "1. 画像生成の要素がある要求では use_image_generation=true にすること。さらに:\n"
    "   - 『絵を描いて』『イラストを作って』のように画像だけが目的なら image_only=true"
    "(テキスト回答は不要)。\n"
    "   - 『〜を説明して図も作って』のようにテキスト回答も画像も必要なら image_only=false とし、"
    "mode と selected_proposers は通常どおり選ぶ(本文は MoA が作り、その内容から画像を生成する)。\n"
    "   画像が不要な通常の質問は use_image_generation=false, image_only=false。\n"
    "2. 『パワポ』『スライド』『プレゼン』『PowerPoint』『資料を作って』等のスライド作成要求では "
    "make_pptx=true とし、mode='moa' で本文(見出し・箇条書き構成)を作ること。画像はスライド内容に応じて"
    "自動生成されるので use_image_generation は false で良い。\n"
    "3. Office ファイル(.docx/.xlsx/.pdf 等)が添付・指定され、その解説・分析を求めている場合は、"
    "必ず mode='moa' とし、selected_proposers に必ず 'Proposer C' を含めて主軸に据えること。\n"
    "4. 質問が『最新』『今』『現在』『2025年』『2026年』『価格』『相場』『型番』『バージョン』など、"
    "時間とともに変わる事実・時事情報を含む場合は search_required=true にすること。"
    "自分の知識だけで答えられると過信しないこと(学習データは古くなる)。"
    "純粋に不変な知識(数学・定義・確立した歴史的事実)のみ false。\n"
    "5. task_type を分類すること: 最終答が数値・式で一意に定まる問題は 'math'、"
    "選択肢(A/B/C/D等)から選ぶ問題は 'mcq'、コード作成・デバッグは 'code'、"
    "事実知識・解説は 'knowledge'、文章作成は 'writing'、雑談・その他は 'chat'。"
    "証明や『なぜ』の説明は答えが一意の数値でないため 'math' ではなく 'knowledge'。\n\n"
    "【精度優先の原則(最重要・厳守)】\n"
    "- mode='single' を許可するのは、挨拶・単純な事実確認・短い定義・簡単な計算のみ。\n"
    "- コード生成/実装/デバッグ・数学/証明・論理パズル・多段推論・比較/要約/翻訳/説明は、"
    "たとえ簡単そうに見えても必ず mode='moa' とすること。単一モデルが見落とす誤りを"
    "複数 Proposer の相互チェックで潰すのが目的。特にコードと証明は例外なく moa とし 3〜4 体選ぶ。\n"
    "- 迷ったら moa を選ぶ(精度優先。計算コストは二の次)。\n"
    f"- rounds は通常 1。明確に洗練が必要な難問だけ 2 以上(最大 {MAX_ROUNDS})。\n\n"
    "selected_proposers には上記のペルソナ名(\"Proposer A\"〜\"Proposer D\")だけを使うこと。"
    "質問の内容と各 Proposer の得意分野を照合して選ぶこと。\n"
    "思考や散文を一切出さず、以下の JSON オブジェクトのみを出力すること。"
)

CRITIC_SYS = (
    "You are a strict reviewer. Given a question and a candidate answer, judge whether "
    "the answer is correct, complete and well-reasoned. "
    'Output ONLY JSON: {"ok": true or false, "issue": "short reason if not ok"}. '
    "Be conservative: only set ok=false when there is a real, identifiable problem. "
    "No prose, no thinking."
)

# Ollama の format に渡す JSON スキーマ。enum で値域まで拘束し、think=False と併用して
# Conductor/Critic を高速・確実にする（freeform "json" だと mode に "mathematical_proof" 等の
# enum 外の値を捏造するため、スキーマで縛るのが要点）。
CONDUCTOR_SCHEMA = {
    "type": "object",
    "properties": {
        "mode": {"type": "string", "enum": ["single", "moa"]},
        "task_type": {"type": "string",
                      "enum": ["math", "code", "mcq", "knowledge", "writing", "chat"]},
        "selected_proposers": {"type": "array", "items": {"type": "string"}},
        "rounds": {"type": "integer", "minimum": 1, "maximum": MAX_ROUNDS},
        "use_image_generation": {"type": "boolean"},
        "image_only": {"type": "boolean"},
        "make_pptx": {"type": "boolean"},
        "search_required": {"type": "boolean"},
        "reason": {"type": "string"},
    },
    "required": ["mode", "task_type", "selected_proposers", "rounds",
                 "use_image_generation", "search_required", "reason"],
}

CRITIC_SCHEMA = {
    "type": "object",
    "properties": {
        "ok": {"type": "boolean"},
        "issue": {"type": "string"},
    },
    "required": ["ok"],
}

# ==================================================
# Conductor（動的プランニング）
# ==================================================


def build_proposer_desc():
    """ペルソナ名付きで、導入済みモデルのみ列挙する。"""
    lines = []
    for label, model in PERSONA_MODELS.items():
        if model in PROPOSERS:
            prof = PROPOSER_PROFILES.get(model, "汎用")
            lines.append(f"- {label} ({model}): {prof}")
    return "\n".join(lines)


def _resolve_proposer(name):
    """ペルソナ名 or モデル名を実モデル名へ解決する（未導入・未知なら None）。
    'Proposer A' / 'proposer a' / 'A' / 実モデル名 の緩い表記を許容。"""
    # 2026-07-25: CONDUCTOR_SCHEMA は selected_proposers の要素を
    # items:{type:string} で拘束しているが、本ファイル各所のコメントが
    # 明記する通りこの強制は完全ではない（"スキーマ強制でも稀に JSON が
    # 崩れる"）。list/dict のような非文字列・非ハッシュ可能な要素が
    # 紛れ込むと、直後の `name in PERSONA_MODELS`（dict メンバーシップ
    # テスト）が isinstance/str() 変換を経る前に TypeError: unhashable
    # type を送出し、これが唯一の呼び出し元 validate_plan を経て
    # conduct -> ask_fugu/fugu_answer まで無捕捉で伝播し、ターン全体を
    # クラッシュさせて計算済みの回答を失っていた。iter 103
    # (_ddg_instant の非list RelatedTopics)、iter 111
    # (plan_pptx_images の非list images)、iter 112
    # (research_search の非list queries)、iter 113 (_read_ipynb の
    # 非dict/非list cells) と同じ「壊れたスキーマ制約付きプランは例外を
    # 出さず既定値へフォールバックさせる」作法に倣い、非文字列要素は
    # ここで即座に None へ倒す（精度優先: フォールバックの方がターン
    # 喪失よりまし）。
    if not isinstance(name, str):
        return None
    if name in PERSONA_MODELS:
        m = PERSONA_MODELS[name]
        return m if m in PROPOSERS else None
    if name in PROPOSERS:
        return name
    key = str(name).strip().lower()
    for label, model in PERSONA_MODELS.items():
        if key in (label.lower(), label.lower().replace("proposer ", "")):
            return model if model in PROPOSERS else None
    return None


def _persona_str(model):
    """表示用に 'Proposer X (model)' 形式へ整形する。"""
    p = MODEL_TO_PERSONA.get(model)
    return f"{p} ({model})" if p else str(model)


# コード/証明/論理は単一モデルが誤りやすいので、Conductor が single を選んでも moa へ
# 格上げする決定的ガードレール。小型 Conductor(qwen3:4b) が精度優先ルールを取りこぼす件の
# 安全網（旧スキーマの complexity フィールドが担っていた難易度分類の代替）。
# 精度優先・計算コスト非優先の方針に沿う。画像生成バイパスは対象外。
_HARD_SIGNALS = re.compile(
    r"実装|コード|プログラム|関数|クラス|アルゴリズム|デバッグ|バグ|ソート|計算量|"
    r"証明|背理法|論理的|パズル|嘘つき|"
    r"\bcode\b|\bimplement|\bfunction\b|\balgorithm|\bdebug|\bprove\b|\bproof\b|```",
    re.IGNORECASE)


def _apply_accuracy_guardrails(question, plan):
    """精度優先: コード/証明/論理の質問で single が選ばれたら moa へ格上げする。"""
    if plan.get("image_only"):
        return plan  # 画像のみはテキスト側の格上げ不要
    if plan["mode"] == "single" and _HARD_SIGNALS.search(question or ""):
        plan["mode"] = "moa"
        if len(plan.get("selected_proposers", [])) < 2:
            plan["selected_proposers"] = PROPOSERS[:3]
        plan["reason"] = "[guardrail: code/proof→moa] " + plan.get("reason", "")
    return plan


# 出力形態（画像/PowerPoint）の決定的ガードレール。小型 Conductor(qwen3:4b) は
# image_only / make_pptx の取りこぼしが多いため、明確なキーワードで補正する。
_PPTX_SIGNALS = re.compile(
    r"パワポ|パワーポイント|スライド|プレゼン(テーション)?|power\s*point|\bpptx\b|スライド資料|プレゼン資料",
    re.IGNORECASE)
_IMAGE_SIGNALS = re.compile(
    r"絵を描|イラスト|描いて|画像を?(生成|作|つく)|図[をもの]?(作|生成|描|示)|"
    r"イメージ図|図解|概念図|挿絵|ダイアグラム|"
    r"\bdraw\b|illustration|\bpicture\b|diagram|generate.{0,10}image",
    re.IGNORECASE)
# 画像に加えテキスト解説も要る＝イラスト付き回答（image_only=False）の手掛かり
_TEXT_TASK_SIGNALS = re.compile(
    r"説明|解説|教えて|まとめ|要約|について|比較|分析|解析|理由|方法|とは|手順|"
    r"explain|describe|summar|analy|\bwhy\b|\bhow\b",
    re.IGNORECASE)


def _apply_routing_guardrails(question, plan):
    """出力形態を補正: PowerPoint / 画像のみ / イラスト付き回答 を明確なキーワードで確定する。"""
    q = question or ""
    if _PPTX_SIGNALS.search(q):
        plan["make_pptx"] = True
        plan["image_only"] = False
        plan["mode"] = "moa"
        if len(plan.get("selected_proposers", [])) < 2:
            plan["selected_proposers"] = PROPOSERS[:3]
        return plan
    if _IMAGE_SIGNALS.search(q):
        plan["use_image_generation"] = True
        # 2026-07-24: make_pptx=True ⇒ image_only=False は validate_plan
        # (iteration 40, 2227-2228/2240-2241) が確立した不変条件。ここは
        # validate_plan の「後」に conduct() から呼ばれる（2320行目）ため、
        # plan["make_pptx"] を見ずに image_only を立てるとその不変条件を
        # 再び壊してしまう（コンダクタが '発表資料'/'デッキ' 等 _PPTX_SIGNALS
        # 非一致のPPTX同義語で make_pptx=True にした上で、質問が画像シグナルに
        # マッチし解説シグナルには非マッチな場合、make_pptx=True かつ
        # image_only=True という矛盾が復活する）。ask_fugu はルート1(画像)を
        # ルート2(PowerPoint)より先に判定するため、その矛盾は要求された
        # PowerPoint を黙って握りつぶし画像のみを返す形で表面化する。
        # テキスト解説も求めている＝イラスト付き / 画像だけ＝image_only
        # （ただし make_pptx が既に True なら image_only は常に False）
        plan["image_only"] = (not plan.get("make_pptx")) and not bool(_TEXT_TASK_SIGNALS.search(q))
    return plan


# task_type の決定的ガードレール（小型 Conductor の分類ミス対策。2026-07-11 追加）。
# math/mcq は solve_verifiable(自己一貫性投票) へルートされるため、誤って自由記述問題を
# math にすると答えの抽出ができない。確実なキーワードでのみ確定させ、証明・説明系は外す。
_MCQ_SIGNALS = re.compile(
    r"which of the following|次のうち|選びなさい|選べ|正しいものを|適切なものを|"
    r"(?:^|\n)\s*\(?A[).:：]\s*\S.*\n\s*\(?B[).:：]\s*\S",
    re.IGNORECASE)
_MATH_TASK_SIGNALS = re.compile(
    r"\\boxed|求めよ|求めなさい|を計算し|を求め|何通り|余りを|剰余|確率を求?|"
    r"\bhow many\b|\bcompute\b|\bcalculate\b|"
    r"\bfind the (?:number|value|sum|remainder|area|probability|smallest|largest)\b|"
    r"answer with (?:the )?number",
    re.IGNORECASE)
_CODE_TASK_SIGNALS = re.compile(
    r"実装|コード|プログラム|関数を書|クラスを書|デバッグ|"
    r"\bimplement|\bwrite (?:a |the )?(?:function|program|code|class|script)\b|```",
    re.IGNORECASE)
# 2026-07-26: _MATH_TASK_SIGNALS には英語の数学トリガー（\bhow many\b, \bcompute\b,
# \bcalculate\b, \bfind the ...\b 等）が揃っているのに、この降格用セーフティバルブは
# 日本語の説明系動詞（証明/説明して/解説して/なぜ）しか見ておらず、英語の explain/
# describe/why が抜けていた（非対称）。結果、"How many SOLID principles are there?
# Explain each." のように \bhow many\b で math 判定されつつ説明を求める問いが降格されず、
# solve_verifiable の自己一貫性投票 + extract_final_answer に流れて説明が黙って
# 落ちる回帰があった。ここでは explain/describe/why の3語だけを追加する（bare \bhow\b は
# \bhow many\b と衝突し正当な計数問題まで誤って降格するため絶対に追加しないこと。
# gotcha #7: 自己一貫性投票は精度優先の要）。
_FREEFORM_SIGNALS = re.compile(
    r"証明|\bprove\b|\bproof\b|説明して|解説して|なぜ|"
    r"\bexplain\b|\bdescribe\b|\bwhy\b",
    re.IGNORECASE)


def _apply_tasktype_guardrails(question, plan):
    """task_type を決定的キーワードで補正する。確実なシグナルが無ければ Conductor の分類を尊重。"""
    q = question or ""
    t = str(plan.get("task_type") or "").lower()
    if t not in ("math", "code", "mcq", "knowledge", "writing", "chat"):
        t = ""
    if _MCQ_SIGNALS.search(q):
        t = "mcq"
    elif _CODE_TASK_SIGNALS.search(q):
        t = "code"
    elif _MATH_TASK_SIGNALS.search(q):
        t = "math"
    # 証明・説明は最終答が一意の数値/文字にならないため投票では解けない → MoA 経路へ
    if t == "math" and _FREEFORM_SIGNALS.search(q):
        t = "knowledge"
    plan["task_type"] = t or "chat"
    return plan


# 2026-07-26: Office 添付ファイル(.docx/.xlsx/.pdf 等)の決定的ガードレール。
# CONDUCTOR_SYS の【特殊ルーティング指示】3(本ファイル 2547-2548 行目)は
# 「Office ファイルが添付されその解説・分析を求めている場合は、必ず mode='moa' とし
# selected_proposers に必ず 'Proposer C' を含めて主軸に据えること」を明記し、
# conduct() はさらに自然文のヒント（'[注記] ... 特殊ルーティング指示 #2 ...'、
# conduct() 内 hint 変数）まで添えている。だが本ファイル各所のガードレール群
# ―― PPTX/画像出力形態は _apply_routing_guardrails（iter 36）、コード/証明の
# single→moa 格上げは _apply_accuracy_guardrails（iter 97）、task_type 誤分類補正は
# 直上の _apply_tasktype_guardrails（iter 64、二重JSON失敗フォールバックにも適用）
# ―― が繰り返し記録している通り、小型 Conductor(qwen3:4b) はプロンプト中の
# ルーティング指示文を取りこぼすことが実測されている。office_attached は
# これまで conduct() のヒント文言生成以外に一切参照されておらず（grep 済み）、
# Conductor がヒントを無視すると .docx/.xlsx/.pdf 添付の解説・分析タスクが
# mode='single' のまま単一の非専門モデルへ回されたり、RAG/Office 文書専門家
# gemma4:26b (Proposer C) が selected_proposers から抜け落ちたまま実行される
# ――このフラグが守るはずの文書解析経路そのものがサイレントに劣化する。
# 他のガードレール群と同じ「プロンプトのヒントとは別に、フラグ/キーワードで
# 決定的に確定させる」作法（belt-and-suspenders）に倣い、office_attached を
# 直接の決定的シグナルとして mode='moa' を強制し、導入済みであれば
# Proposer C を主軸(先頭)に確実に加える。プロンプトのヒント自体は削らない
# （二重の安全網として両方を維持する）。
def _apply_office_guardrail(plan, office_attached):
    """office_attached=True を決定的ガードレールとして扱う。
    mode='moa' を強制し、導入済みなら Proposer C (PERSONA_MODELS 由来の実モデル名)を
    selected_proposers の先頭に確実に加える(重複排除・4件上限は validate_plan の
    models[:4] と揃える)。office_attached=False は完全な no-op（plan を一切変更しない）。
    image_only=True の画像専用プランは _apply_accuracy_guardrails と同じ理由で
    テキスト側 MoA 強制の対象外（画像だけの回答にテキストパネルを割り当てても無意味）。
    べき等: 2回適用しても mode/selected_proposers/reason は変化しない。"""
    if not office_attached:
        return plan
    if plan.get("image_only"):
        return plan
    plan["mode"] = "moa"
    tag = "[guardrail: office→moa+ProposerC] "
    reason = str(plan.get("reason", ""))
    if not reason.startswith(tag):
        plan["reason"] = tag + reason
    c_model = PERSONA_MODELS.get("Proposer C")
    if c_model and c_model in PROPOSERS:
        models = list(plan.get("selected_proposers") or [])
        if c_model in models:
            models.remove(c_model)
        models.insert(0, c_model)
        deduped = []
        for m in models:
            if m not in deduped:
                deduped.append(m)
        plan["selected_proposers"] = deduped[:4]
    return plan


def default_plan():
    return {
        "mode": "moa",
        "task_type": "",
        "selected_proposers": PROPOSERS[:3],
        "rounds": 1,
        "use_image_generation": False,
        "image_only": False,
        "make_pptx": False,
        "search_required": False,
        "reason": "fallback (planner unavailable)",
        "_fallback": True,
    }


def validate_plan(p):
    base = default_plan()
    if not isinstance(p, dict):
        return base
    out = dict(base)
    out["_fallback"] = False

    out["use_image_generation"] = bool(p.get("use_image_generation", False))
    out["image_only"] = bool(p.get("image_only", False))
    out["make_pptx"] = bool(p.get("make_pptx", False))
    out["search_required"] = bool(p.get("search_required", False))

    # make_pptx / illustrated(text+image) は本文が要るので image_only を無効化
    if out["make_pptx"]:
        out["image_only"] = False

    # 2026-07-22: スキーマ制約付き Conductor が image_only=True かつ
    # use_image_generation=False という矛盾した plan を出すことがある
    # （_apply_routing_guardrails は image_only を use_image_generation=True と
    # セットでしか立てないため、この矛盾を修復するガードレールが他に存在しない）。
    # 放置すると (1) fugu_answer の SC投票ゲート（math/mcq で精度に直結、
    # gotcha#7）が image_only により無効化され、(2) MoA提案パネルが
    # IMAGE_PROMPT_PANEL=2 まで縮小され、(3) それでいて use_image_generation=False
    # なので画像自体は一枚も生成されない、という三重の劣化が起きる。
    # image_only は「画像を生成してテキストを省く」の意味なので、画像を
    # 生成しないなら image_only も強制的に False にする。
    if not out["use_image_generation"]:
        out["image_only"] = False

    mode = str(p.get("mode", "moa")).lower()
    out["mode"] = "single" if mode == "single" else "moa"

    t = str(p.get("task_type", "")).lower()
    out["task_type"] = t if t in ("math", "code", "mcq",
                                  "knowledge", "writing", "chat") else ""

    # selected_proposers（ペルソナ名 or モデル名）を実モデル名へ解決・重複排除
    models = []
    raw_props = p.get("selected_proposers")
    if isinstance(raw_props, list):
        for name in raw_props:
            m = _resolve_proposer(name)
            if m and m not in models:
                models.append(m)

    # use_image_generation は非排他フラグ。テキスト側の mode/proposers は通常どおり解決する。
    if out["image_only"]:
        # 画像のみ: テキスト提案は不要（プロンプト起草用に models は保持してよい）
        out["selected_proposers"] = models[:IMAGE_PROMPT_PANEL]
    elif out["mode"] == "single":
        out["selected_proposers"] = (models[:1] or PROPOSERS[:1])
    else:
        out["selected_proposers"] = (models[:4] or PROPOSERS[:3])

    try:
        r = int(p.get("rounds", 1))
    except Exception:
        r = 1
    out["rounds"] = max(1, min(MAX_ROUNDS, r))

    out["reason"] = str(p.get("reason", ""))[:200]
    return out


def conduct(question, history=None, office_attached=False):
    desc = build_proposer_desc()
    hist_note = ""
    if history:
        recent = history[-4:]  # 直近 2 往復分をテキストとして埋め込む
        lines = [
            ("[User]" if m["role"] == "user" else "[Assistant]")
            + ": " + m["content"][:200] + ("..." if len(m["content"]) > 200 else "")
            for m in recent
        ]
        hist_note = "\n\n直近の会話(参考):\n" + "\n".join(lines) + "\n"
    hint = ""
    if office_attached:
        hint = ("\n[注記] Office 文書(.docx/.xlsx/.pdf 等)が添付されています。"
                "特殊ルーティング指示 #2 を適用し、mode='moa' かつ selected_proposers に "
                "'Proposer C' を含めてください。\n")
    user = (
        f"利用可能な Proposer とその強み:\n{desc}\n"
        f"{hist_note}"
        f"{hint}"
        f"\nユーザーの質問:\n{question}\n\n"
        "JSON プランのみを返すこと。"
    )
    msgs = [{"role": "system", "content": CONDUCTOR_SYS},
            {"role": "user", "content": user}]
    raw = ask(
        CONDUCTOR, msgs, CONDUCTOR_TEMP,
        think=False,   # 思考は無効化（+スキーマ拘束）でプランJSONを高速・確実に得る
        fmt=CONDUCTOR_SCHEMA,
        num_predict=NUM_PREDICT_JUDGE,
        label="conductor",
    )
    plan = extract_json(raw)
    if plan is None:
        # スキーマ強制でも稀に JSON が崩れる(実測 ~1/10 程度の一過性)。固定フォールバック
        # プランより正しいプランの方が良いので、1 回だけ引き直す。
        raw = ask(
            CONDUCTOR, msgs, CONDUCTOR_TEMP,
            think=False, fmt=CONDUCTOR_SCHEMA,
            num_predict=NUM_PREDICT_JUDGE, label="conductor",
        )
        plan = extract_json(raw)
    plan = _apply_routing_guardrails(question, validate_plan(plan))
    plan = _apply_accuracy_guardrails(question, plan)
    plan = _apply_office_guardrail(plan, office_attached)
    return _apply_tasktype_guardrails(question, plan), raw


def _critic_judge(question, answer, think):
    """Critic 1 回分の呼び出し。(ok, issue)。think=True は再検算（思考込みで遅いが正確）。"""
    raw = ask(
        CONDUCTOR,
        [{"role": "system", "content": CRITIC_SYS},
         {"role": "user", "content": (
             f"Question:\n{question}\n\nCandidate answer:\n{answer}\n\n"
             "Return ONLY JSON."
         )}],
        CONDUCTOR_TEMP,
        think=think,
        fmt=CRITIC_SCHEMA,
        num_predict=(NUM_PREDICT_JUDGE_THINK if think else NUM_PREDICT_JUDGE),
        label="critic",
    )
    # 2026-07-22: __ERROR__ センチネル（ask() の通信/モデル失敗）と、空文字や
    # パース不能だが正常な出力とを区別する。後者は gpt-oss:20b の思考予算切れで
    # 本文が空になる既知ケースのため ok=True 既定を維持するが、前者は critic 呼び出し
    # そのものが失敗しているだけで「回答に問題なし」を意味しない。ここで黙って
    # ok=True にすると verify_single() の最終審判（think=True critic）が事実上
    # 無審査で通ってしまい、精度優先の方針に反するため ok=False にしてエスカレーション
    # させる（呼び出し元は MoA パネルへフォールバックできる）。
    if strip_think(raw).startswith("__ERROR__"):
        return False, f"critic call failed: {strip_think(raw)}"[:200]
    p = extract_json(raw) or {}
    return bool(p.get("ok", True)), str(p.get("issue", ""))[:200]


def critique(question, answer):
    """回答の十分性を 2 段階で判定。(ok: bool, issue: str) を返す。
    1段目: think=False + スキーマで高速判定。ok ならそこで確定（高速パス維持）。
    2段目: 1段目が NG のときだけ think=True で再検算して最終判定。
    think=False の Critic は頭の中で再計算ができず、正答 '700' を誤って NG にする
    偽エスカレーション(310秒浪費)が 2026-07-03 のフル評価で実測されたための対策。"""
    ok, _issue = _critic_judge(question, answer, think=False)
    if ok:
        return True, ""
    return _critic_judge(question, answer, think=True)


def second_opinion(question, answer):
    """自己評価バイアス対策の独立チェック。Conductor(qwen3) とは別系統の
    SECOND_OPINION_MODEL に同じ審査をさせる。(ok, issue) を返す。
    phi4-mini は非thinkingモデルなので think パラメータは送らない。"""
    global _SECOND_OPINION_DISABLED
    if SECOND_OPINION_MODEL not in PROPOSERS:
        if not _SECOND_OPINION_DISABLED:
            print(f"   ? second_opinion モデル {SECOND_OPINION_MODEL} が見つかりません "
                  f"→ 自己評価バイアス対策が無効化されます。verify_single() は思考ON再検算を必須にします。")
            _SECOND_OPINION_DISABLED = True
        return True, ""
    raw = ask(
        SECOND_OPINION_MODEL,
        [{"role": "system", "content": CRITIC_SYS},
         {"role": "user", "content": (
             f"Question:\n{question}\n\nCandidate answer:\n{answer}\n\n"
             "Return ONLY JSON."
         )}],
        CONDUCTOR_TEMP,
        fmt=CRITIC_SCHEMA,
        # gpt-oss:20b は MODEL_CONFIG で think:"high" が自動適用されるため、思考が予算を
        # 食い尽くして本文空(=ok既定)になるのを防ぐべく思考込みの上限を使う。
        num_predict=NUM_PREDICT_JUDGE_THINK,
        label="critic2",
    )
    # 2026-07-22: _critic_judge と同様、__ERROR__ センチネル（second opinion モデルの
    # 通信/モデル失敗）は「空文字/非JSONだが正常出力」（gpt-oss:20b の思考予算切れ既定）
    # とは別物として扱う。second_opinion は自己評価バイアス対策の独立チェックであり、
    # ここが黙って ok=True になると verify_single() でバイアス対策が機能しないまま
    # 高速パスが通ってしまう。エラー時は ok=False にして think=True 再検算に回す。
    if strip_think(raw).startswith("__ERROR__"):
        return False, f"critic call failed: {strip_think(raw)}"[:200]
    p = extract_json(raw) or {}
    return bool(p.get("ok", True)), str(p.get("issue", ""))[:200]


def verify_single(question, answer):
    """単体回答の採用可否。(ok, issue) を返す。
    高速チェック 2 系統（qwen3 think=False と、別系統 phi4-mini による独立チェック＝
    自己評価バイアス対策）を先に走らせ、どちらかが疑義を出したときだけ
    qwen3 think=True の再検算を最終審判にする。

    second_opinion が無効化されている場合は、高速チェック 1 系統だけになり、疑義が
    ある場合は即座に think=True 再検算を通す（手厚い保護）。

    実測(2026-07-03): 高速チェックはどちらも正答 '700' を誤って NG にすることがあるが、
    思考ONの再検算は正しく ok と判定した。逆に誤答は再検算が明確な理由付きで NG にする。
    コード回答はまず実行(決定的)で検証し、失敗なら LLM 審査を待たず即 NG。"""
    code_issue = code_check(answer)
    if code_issue:
        return False, code_issue
    ok1, issue1 = _critic_judge(question, answer, think=False)
    ok2, issue2 = second_opinion(question, answer)

    if _SECOND_OPINION_DISABLED:
        doubt = issue1
        if ok1:
            return True, ""
    else:
        if ok1 and ok2:
            return True, ""
        doubt = issue1 if not ok1 else f"second opinion ({SECOND_OPINION_MODEL}): {issue2}"

    ok3, issue3 = _critic_judge(question, answer, think=True)
    if ok3:
        return True, ""
    return False, (issue3 or doubt)

# ==================================================
# 提案・統合
# ==================================================


def proposer_sys_for(model):
    """モデルのペルソナ人格を PROPOSER_SYS の前に前置したシステムプロンプトを返す。"""
    identity = PERSONA_IDENTITY.get(model)
    return f"{identity}\n{PROPOSER_SYS}" if identity else PROPOSER_SYS


def get_single_proposal(model, question, reference, issue=None, history=None):
    """issue: Critic の指摘。history: 過去の会話履歴（コンテキスト継続用）。"""
    history = history or []
    sys_prompt = proposer_sys_for(model)
    if reference is None:
        content = question
        if issue:
            content = (f"{question}\n\n(Note: a previous attempt was flagged by a reviewer: "
                       f"{issue} — avoid that pitfall.)")
        msgs = (
            [{"role": "system", "content": sys_prompt}]
            + history
            + [{"role": "user", "content": content}]
        )
    else:
        note = (f"A reviewer flagged this issue with the draft: {issue}\n\n"
                if issue else "")
        msgs = (
            [{"role": "system", "content": sys_prompt}]
            + history
            + [{"role": "user", "content": (
                f"Question:\n{question}\n\n"
                f"A draft answer from the panel:\n{reference}\n\n"
                f"{note}"
                "Improve it: fix errors, add missing points, make it more "
                "accurate. Output your improved answer only."
            )}]
        )
    # 2026-07-23: ここは長らく think=PROPOSER_THINK（生のグローバル、既定 None）を直接渡して
    # おり、num_predict だけ proposer_predict_for(model) で MODEL_CONFIG 対応させたのに think
    # が取り残されていた欠落サイト。これだと gpt-oss:20b/qwen3.6:35b が MoA 提案でも
    # think:"high"/True を一度も受け取れず、_sc_sample（SC 経路、proposer_think_for 使用済み）
    # とだけ非対称だった。proposer_think_for に揃えて解決（PROPOSER_THINK override 優先は温存）。
    return model, ask(model, msgs, PROPOSER_TEMP, think=proposer_think_for(model),
                      num_predict=proposer_predict_for(model), label="proposer")


def get_proposals(models, question, reference=None, issue=None, history=None):
    """Conductor が選んだ models のみで提案を集める。history: 会話コンテキスト。
    多様性維持: reference がある回でも先頭 1 体は元の質問から新規に回答する。"""
    jobs = [(m, (None if (reference is not None and i == 0) else reference))
            for i, m in enumerate(models)]
    if PARALLEL_PROPOSERS:
        # 2026-07-23: as_completed() は完了順で future を返すため、返り値の
        # (model, answer) 順序が実行毎に非決定的になり、aggregator の
        # "Answer A/B/C" ラベル割当がブレて MoA 統合の再現性が失われる。加えて
        # 多様性契約（jobs[0] は reference=None の新規回答を先頭に置く）が
        # aggregator に見える位置として保証されなくなる。このパスは 8GB GPU
        # の既定では休眠中（gotcha #5: OLLAMA_MAX_LOADED_MODELS=1 により逐次
        # 実行が正しい既定）だが、apply_high_vram_profile が有効化する 96GB
        # 実験用の高VRAM構成ではまさにここが使われるため、投入順
        # (futs リストの順序 = jobs の順序) で結果を集めて決定的にする。
        # 全ジョブを先に submit してから結果収集するので並列度・wall-clock は
        # 変わらない（as_completed をやめても同時実行性は失われない）。
        with concurrent.futures.ThreadPoolExecutor(max_workers=max(1, len(jobs))) as ex:
            futs = [ex.submit(get_single_proposal, m, question, ref, issue, history)
                    for m, ref in jobs]
            out = [fut.result() for fut in futs]
        return out
    return [get_single_proposal(m, question, ref, issue, history) for m, ref in jobs]


def use_jp_aggregator(text):
    """qwen3 を統合役にすべき質問かどうか。ひらがな/カタカナに加え、漢字のみの短文
    （「東京都の人口は?」等）も対象にする。漢字だけでは日本語と中国語を確実に区別
    できないが、qwen3 は中国語も堅牢なので qwen3 側に倒すのが安全
    （deepseek-r1 の言語混入の既知問題を踏まない）。"""
    t = text or ""
    return bool(re.search(r"[぀-ヿ]", t)) or bool(re.search(r"[㐀-鿿]", t))


def pick_aggregator(question, has_code=False):
    """統合役の選定。コード付き→AGGREGATOR(qwen3-coder)、日本語→強い思考型JPモデル
    （未導入なら従来の qwen3:4b）、それ以外→思考型 AGGREGATOR_REASONING（未導入なら既定）。"""
    if has_code:
        return AGGREGATOR
    if use_jp_aggregator(question):
        if JP_AGGREGATOR_STRONG and JP_AGGREGATOR_STRONG in PROPOSERS:
            return JP_AGGREGATOR_STRONG
        if JP_AGGREGATOR in PROPOSERS or JP_AGGREGATOR == CONDUCTOR:
            return JP_AGGREGATOR
        return AGGREGATOR
    if AGGREGATOR_REASONING in PROPOSERS:
        return AGGREGATOR_REASONING
    return AGGREGATOR


def aggregate(question, proposals):
    # 【修正】アグリゲーターに渡す前に各提案の think を除去（文脈圧迫と審判の混乱を防ぐ）
    good = [(m, strip_think(a)) for (m, a) in proposals if not a.startswith("__ERROR__")]
    if not good:
        return "__ERROR__: 全プロポーザーが失敗しました（モデル/Ollama/VRAM を確認）。"
    # コード付き提案には実行結果の証拠を添える。どの案が実際に動くかは決定的情報なので、
    # アグリゲータの取捨選択の判断材料として最強（AGGREGATOR_SYS のルール6が対応）。
    # 【2026-07-22 修正】以前はここで `good` 自体をタグ付き版で上書きしていたため、
    # 保険2（統合失敗時に good から直接返す）が [Execution check: ...] タグや生の
    # トレースバックをユーザー向け回答に漏らしていた。`good` はクリーンな
    # (model, strip_think(answer)) のまま保持し、タグ付きビューは別変数
    # （annotated）に持たせて、アグリゲータへの block 文字列構築にのみ使う。
    annotated = good
    # 【2026-07-24】保険2（下記）が「コードが実際に FAILED したと判明している提案」を
    # 避けられるよう、このアノテーションループで既に呼んでいる code_check() の結果を
    # good のインデックスに紐づけて保持する。code_check/run_python を再実行はしない。
    failed_idxs = set()
    if CODE_EXECUTION:
        annotated = []
        for i, (m, ans) in enumerate(good):
            if extract_code(ans):
                issue = code_check(ans)
                if issue is not None:
                    failed_idxs.add(i)
                tag = ("[Execution check: PASSED]" if issue is None
                       else f"[Execution check: FAILED]\n{issue}")
                ans = f"{ans}\n\n{tag}"
            annotated.append((m, ans))

    labels = [chr(ord("A") + i) for i in range(len(annotated))]
    block = "\n\n".join(
        f"Answer {lab}:\n{ans}" for lab, (_m, ans) in zip(labels, annotated)
    )
    user = f"Question:\n{question}\n\n{block}"

    def _run(model, think=None):
        return ask(
            model,
            [{"role": "system", "content": AGGREGATOR_SYS},
             {"role": "user", "content": user}],
            AGGREGATOR_TEMP,
            think=think,
            num_predict=model_cfg(model, "num_predict", NUM_PREDICT_AGGREGATOR),
            label="aggregator",
        )

    def _bad(out):
        return out.startswith("__ERROR__") or not strip_think(out).strip()

    primary = pick_aggregator(question,
                              has_code=any(extract_code(a) for _m, a in good))
    out = _run(primary)

    # 【保険1】空/エラーなら qwen3 の think=False で再統合。空返答の根本原因は
    # 「思考が num_predict を食い尽くして本文ゼロ」(2026-07-04 実測)なので、
    # 思考を切って全予算を本文に回すのが確実。primary が qwen3 だった場合にも有効。
    if _bad(out):
        print(f"   ⚠ アグリゲータ({primary})が空返答/失敗 → {JP_AGGREGATOR}(think=False) で再統合")
        out = _run(JP_AGGREGATOR, think=False)

    # 【保険2】それでもダメなら、Critic が ok と判定した提案をそのまま返す。
    # 3体分の正しい提案が手元にあるのに空回答で失点するのが最悪ケースなので、それを塞ぐ。
    if _bad(out):
        print("   ⚠ 統合に失敗 → 提案から直接選択します")
        # 【2026-07-24 修正】critique() は LLM によるレビューのみでコードを実行しない。
        # そのため以前はここで good を単純に先頭から走査しており、CODE_EXECUTION の
        # アノテーションループで「実行して FAILED と判明済み」の提案でも critique() が
        # ok と判定すればそのまま最終回答として返してしまい得た
        # （AGGREGATOR_SYS ルール6「FAILED のコードを最終回答の根拠にしない」に反する）。
        # iteration 9 で good はクリーンな (model, strip_think(answer)) のまま保持する
        # ようにしたが、それだけでは保険2の選択ロジック自体は直せていなかった。
        # ここでは failed_idxs（上の実行済み証拠、再実行はしない）を使って、まず
        # 「コードが FAILED していない」提案（コード無し提案も含む）に候補を絞り、
        # 全提案が FAILED だった場合のみ元の全候補（good, 元の順序）にフォールバックする。
        candidates = [(m, a) for i, (m, a) in enumerate(good) if i not in failed_idxs]
        if not candidates:
            candidates = good
        for _m, a in candidates:
            ok, _issue = critique(question, a)
            if ok:
                return a
        return max(candidates, key=lambda x: len(x[1]))[1]

    return out

# ==================================================
# 自己一貫性投票（Self-Consistency + PoT）
# ==================================================
# 答えが機械照合できるタスク(math / mcq)では、1 回の MoA 統合より
# 「k 回独立に解かせて最終答を抽出し多数決」の方が確実に強い（Self-Consistency）。
# さらに math では「Python を書かせて実行し、その出力を 1 票にする」PoT 票を混ぜ、
# 計算ミス系の誤答を機械的に排除する。時間無制限・精度最優先の方針の中核機能（2026-07-11）。

# nk108 方針: 時間は無制限・精度最優先。SC は「サンプルを増やすほど当たる」ので上限を高めに取る。
SC_ENABLED = True
SC_INITIAL = 6          # 第1バッチの CoT サンプル数（精度優先で厚め）
SC_STEP = 4             # 過半数が取れないときの追加サンプル数
SC_MAX = 20             # 主力 CoT サンプルの上限（PoT・安価票は別枠）。時間無制限方針で高め
# 全会一致判定 (cnt == n) は抽出成功サンプルのみで n を数えるため、thinking の
# num_predict 打ち切りで __ERROR__ になる／PoT コードが実行失敗する／\boxed{} が
# 出ない、といった抽出失敗が第1バッチで多発すると n=1 (残り全滅) でも「全会一致」
# 扱いになり、事実上 k=1 で確定してしまう（精度優先方針に反する縮退）。過半数側は
# 既に n>=4 の下限があるため、全会一致側にも同じ考え方の下限を設ける（2026-07-21）。
SC_MIN_VOTES = 3        # 全会一致で確定してよい最小サンプル数（これ未満は追加サンプリングへ）
# 2026-07-22: _arbitrate に同時提示する同数タイ候補の上限。3-way以上の拮抗も正しく
# 全候補を裁定役に見せるための変更（下記参照）だが、病的に多い同数タイで num_ctx
# (gotcha #2: 8192/16384 に固定)を溢れさせないよう上限で保護する。超過分は
# 黙って捨てず _arbitrate 内でログに出す。
ARBITRATE_MAX_CANDIDATES = 4
SC_TEMP = 0.7           # 多様性確保（投票の独立性）
SC_POT = True           # math で PoT(Python 実行)票を混ぜる
SC_POT_TIMEOUT = 90     # PoT コードの実行タイムアウト秒（総当たり解法に余裕を持たせる）
REASONING_MODELS = ["gpt-oss:20b", "qwen3.6:35b"]  # SC の主力（導入済みのものだけ使われる）
SC_CHEAP_MODEL = "NitrAI/VibeThinker-3B"  # VRAM 常駐の量産サンプラー（3B・高速）
SC_CHEAP_VOTES = 0      # 安価票の数。VibeThinker の AIME ミニ実測で合格したら 6〜12 へ引き上げる
# 票が拮抗したときの最終審判。8GB 環境では gpt-oss:120b(65GB) が RAM48+VRAM8=56GB を超え
# NVMe ページングで1裁定に数十分〜数時間かかり得るため、qwen3.6:35b（思考型・理数最強格）で裁く
# （2026-07-12 nk108 決定）。120b は FUGU_HIGH_VRAM=1（96GB 環境）で解禁される。
# 失敗/空/未導入なら _arbitrate が REASONING_MODELS へ自動フォールバックする。
ARBITER_MODEL = "qwen3.6:35b"

SC_PROMPT_MATH = (
    "Solve the problem step by step, rigorously. Verify your result before answering. "
    "At the very end, put ONLY the final answer in \\boxed{}."
)
SC_PROMPT_MCQ = (
    "Solve the problem step by step. Compare all choices before deciding. "
    "At the very end, output ONLY the letter of the correct choice in \\boxed{} "
    "(for example \\boxed{B})."
)
SC_PROMPT_POT = (
    "Solve the problem by writing ONE complete Python program. "
    "Prefer exact arithmetic (integers, fractions, sympy) over floats. Brute force is fine. "
    "The program must print ONLY the final answer on its last line. "
    "Wrap it in a single ```python block. No input() calls."
)

# 2026-07-22: 全角数字/A-E に加えて、CJK 寄りのプロポーザ (qwen/gemma 系) が
# しばしば出力する Unicode 記号の等価表記もここで正規化する:
#   U+2212 (MINUS SIGN) / U+FF0D (fullwidth hyphen-minus) -> '-'
#   U+FF0E (fullwidth full stop)                          -> '.'
#   U+FF0F (fullwidth solidus)                            -> '/'
#   U+FF0C (fullwidth comma)                               -> ','
# これらを潰さないと "−5"(U+2212) と "-5" は na.lower()==nb.lower() でも
# Fraction() でも一致せず（Fraction は U+2212 を拒否する）、answers_equivalent
# は math_verify 頼みになる。math_verify が失敗すると本来同じ答えが
# vote_answers で票が2系統に割れ、誤答がプルラリティを取ったり無駄なサンプル
# 消費・仲裁が発生する（精度優先の自己整合性投票が崩れる）。曖昧な en-dash
# (U+2013) / em-dash (U+2014) は区間表記等と衝突しうるため、意図的にここでは
# マッピングしない。
# 2026-07-25: 全角パーセント（U+FF05 ％）も同じ理由で '%' へ正規化する。CJK 寄りの
# プロポーザ（qwen/gemma 系）が百分率の答えを全角数字ごと「５０％」のように書くと、
# 全角数字自体は本テーブルで半角化されても ％ は素通りして "50％" のまま残り、
# 半角の "50%" と別の投票クラスに分裂する（下の normalize_answer 側の '\%' -> '%'
# 正規化＝iteration 13/22/78/122 と同系統の姉妹修正、詳細は normalize_answer 内の
# コメント参照）。str.maketrans は変換元・変換先の文字数を一致させる必要があるため、
# 変換元の末尾に '％'、変換先の末尾に対応する '%' を追加する。
_FW_TRANS = str.maketrans(
    "０１２３４５６７８９ＡＢＣＤＥａｂｃｄｅ−－．／，％",
    "0123456789ABCDEabcde--./,%",
)


def extract_boxed(text):
    """最後の \\boxed{...} の中身を波括弧の対応を数えて取り出す（無ければ None）。"""
    if not text:
        return None
    # \boxed{ の出現位置を前から全部集めておく（rfind 一発ではなく）。
    positions = []
    start = 0
    while True:
        idx = text.find("\\boxed{", start)
        if idx == -1:
            break
        positions.append(idx)
        start = idx + 1
    if not positions:
        return None

    def _scan(idx):
        """idx にある \\boxed{ を depth 走査。(閉じたか, 中身) を返す。"""
        i = idx + len("\\boxed{")
        depth = 1
        out = []
        while i < len(text) and depth > 0:
            c = text[i]
            if c == "{":
                depth += 1
            elif c == "}":
                depth -= 1
                if depth == 0:
                    break
            out.append(c)
            i += 1
        return depth == 0, "".join(out)

    # 2026-07-22 (iteration 12, iteration 11 の続き): 末尾に一番近い \boxed{ から
    # 順に手前へ遡り、波括弧が閉じている（＝打ち切られていない）最初の候補を採用する。
    # thinking モデルが「先に完結した \boxed{回答} を出してから、2 個目の
    # \boxed{...} を書き始めた直後に num_predict/num_ctx で打ち切られる」ケースが
    # 実際に観測されており（gotcha #2）、これまでは text.rfind() で最後の
    # \boxed{ だけを見て depth>0 なら即 None を返していたため、直前に確定していた
    # 正しい答えまで一緒に捨てていた。solve_verifiable の多数決は 1 票でも多い方が
    # 精度に効く（gotcha #7 / 精度優先・時間は気にしない）ため、閉じている
    # \boxed{} が手前に見つかるならそれを採用して票を救出する。
    # 一方、iteration 11 で修正した「未確定の残骸を答えとして返さない」動作は
    # そのまま維持する: 見つかった \boxed{ が一つも閉じていない場合は、
    # 従来どおり None を返す（無投票の方が誤答票より安全）。
    # 2026-07-22 (iteration 25, iteration 11/23 の続き): 閉じてはいるが中身が
    # 空/空白のみの \boxed{}（\boxed{} や \boxed{ }）を「無投票」と同一視して
    # 手前へ遡り続ける。以前は closed かどうかしか見ておらず、末尾の \boxed{}
    # が空でも即座に return content.strip() or None で None を返していたため、
    # その手前にあった \boxed{5} のような確定済みの正しい票まで一緒に捨てて
    # いた。extract_final_answer の math 分岐はここで None を受け取ると
    # 文中の数値を拾う最終フォールバックへ落ちるため、単なる無投票のはずが
    # 誤投票（gotcha #7 の自己整合性投票で最も避けたいケース）に変わって
    # しまう。空 \boxed{} は「無投票」として読み飛ばし、閉じていて中身が
    # 空でない最初の（末尾に近い）候補が見つかった時点でそれを採用する
    # last-wins 方式は維持する。iteration 11 の「一つも閉じていなければ
    # None」という安全側の挙動、および num_ctx 打ち切り絡みの gotcha #2 の
    # 前提もそのまま変えていない。
    for idx in reversed(positions):
        closed, content = _scan(idx)
        if not closed:
            continue
        ans = content.strip()
        if ans:
            return ans
        # 閉じているが空/空白のみ → 無投票として扱い、手前を探索し続ける。
    return None


def normalize_answer(ans):
    """投票用の答え正規化: 全角→半角、$ と桁区切り除去、宣言前置きの除去、外殻 \\text{} 剥がし。"""
    if ans is None:
        return ""
    s = str(ans).strip().translate(_FW_TRANS)
    # 2026-07-25: LaTeX でバレの '%' はコメント開始文字として解釈されてしまうため、
    # 学習済みの行儀の良いモデルは百分率を \boxed{50\%} のようにエスケープして書く。
    # 一方、素の散文回答は \boxed{50%} のようにエスケープなしで書き、_FW_TRANS
    # （直上）が全角数字は正規化しても全角パーセント ％ をここまで手つかずのままにしていた
    # 旧版では qwen/gemma 系 CJK プロポーザが「５０％」を出すこともあり、同じ50%という
    # 値が \%（LaTeX）/ %（素の散文）/ ％（全角）の最大3系統の投票クラスに分裂して
    # vote_answers の集計を薄める（自己整合性投票 gotcha #7 が票割れに最も弱い箇所）。
    # \boxed{} 経路（extract_final_answer の math 分岐、L3033 付近）はここで
    # normalize_answer をそのまま素通しするため数値コア正規表現の対象外であり、この
    # スペルの違いを吸収できる唯一の場所がここになる。値は変えず綴りだけを '%' に
    # 揃える点で iteration 13（全角記号）/22（末尾カンマ）/78/122（\frac 数値正規化）と
    # 同系統の姉妹修正。パーセント記号自体を落とすと 50% が 50 という別の値に化けて
    # しまう（精度劣化）ため、ここでは絶対に除去しない。
    # 2026-07-25: \(...\)（inline 数式モード）/ \[...\]（display 数式モード）はエスケープ
    # 済みの LaTeX 数式区切り文字であり、直前で剥がしている '$'（同じく数式モードの区切り）・
    # '\!'・'\,'（幅なしスペース系マクロ）と同様、値を一切持たない体裁トークンに過ぎない。
    # まだこの4つだけ剥がしていなかったため、プロポーザーが \boxed{\(5\)} や
    # \[x+1\]、「answer is \(42\)」のように答えをこれらで包んで書くと、素の "5"/"x+1"/"42"
    # 票とは別の "\(5\)" 等という投票クラスに分裂し、na.lower()/Fraction の高速パスにも
    # 乗らない（自己整合性投票 gotcha #7 が最も嫌う票割れ）。値を変えない純粋な区切り
    # 記号なので、$ と対称的にここで剥がす。\left(/\right)（実括弧を包む別マクロ）・
    # エスケープ済み中括弧 \{/\}（集合記法）・素の ( )/[ ]・\frac 等の値を持つマクロは
    # 対象外のまま維持する。iteration 13/22/78/122/134/136 と同系統の姉妹修正。
    s = s.replace("\\%", "%").replace("$", "").replace("\\!", "").replace("\\,", "") \
        .replace("\\(", "").replace("\\)", "").replace("\\[", "").replace("\\]", "").strip()
    # 2026-07-25: エスケープ済み中括弧 \{ / \} は集合記法の区切り文字であり、直上のブロック
    # （iteration 140、\( \) \[ \] 剥がし）も直下のブロック（iteration 164、\left/\right 剥がし）
    # も、コメントで明示的に「\{/\} は集合の区切り文字として保持する、剥がさない」とスコープ外に
    # していた。区切り文字そのものを残す判断は正しい（剥がすと \{1,2\} が "1,2" になり、素の
    # タプル/スカラー票と値を変えて誤って一致してしまう）が、そのエスケープの綴り違いは残った
    # ままだった。多様性重視の SC パネルでモデル A が \boxed{{1,2}}（extract_boxed は '{1,2}'
    # を返す）、モデル B が \boxed{\{1,2\}}（extract_boxed は '\{1,2\}' を返す）のように同じ
    # 集合を異なるエスケープ綴りで書くと、na.lower() も Fraction() もこの2つを一致させられず、
    # 同一の値なのに別々の投票クラスに分裂してしまう（自己整合性投票 gotcha #7 が最も嫌う
    # 票割れ）。ここでは中括弧そのものは絶対に削除せず（削除は禁止 — \{1,2\} を "1,2" に
    # してしまうと集合を素のタプル/スカラーと誤って同一視する）、バックスラッシュのエスケープ
    # だけを外して素の { } に統一する。iteration 13/22/78/122/134/136/148/160 と同系統の
    # 「表記の綴り違いを吸収して票をまとめる」姉妹修正で、iteration 140/164 が意図的に
    # 先送りにしていた \{\}（集合デリミタのエスケープ綴り）のギャップを埋める。
    s = s.replace("\\{", "{").replace("\\}", "}")
    # 2026-07-25: \left(/\right) 等は実括弧・角括弧・波括弧の「自動サイズ調整」だけを行う
    # 体裁マクロで、値は一切持たない。iteration 140（直前のブロック、\( \) \[ \] 剥がし）
    # はコメントで明示的にこれをスコープ外としたまま据え置いており（"\left(/\right)
    # （実括弧を包む別マクロ）...は対象外のまま維持する"）、grep でも \left/\right は
    # どこでも処理されておらずテストも皆無だった。プロポーザーは MATH-500 形式の
    # 順序対・区間・集合を \left(3, 4\right) / \left[2, 5\right) / \left\{1, 2\right\}
    # のように書くことが多く、素の (3,4) 等と綴りだけが違う別の投票クラスに分裂する。
    # na.lower() は不一致、Fraction("\left(3,4\right)") は例外で math_verify に落ちるが
    # gotcha #6 の通り Windows 上では信頼できず、自己整合性投票（gotcha #7）が最も
    # 嫌う票割れをそのまま放置してしまう。\left/\right の文字列だけを剥がし、内側の
    # 実区切り文字 ( ) [ ] \{ \} はそのまま残す（\{/\} を含む集合記法は剥がさない）。
    # \b は必須: \leftarrow/\rightarrow/\leftrightarrow のような矢印マクロは "left"/
    # "right" の直後が英字（word文字）で続き単語境界が生じないため、\b が誤って
    # 巻き込むのを防ぐ（正規表現は \left/\right の小文字専用で \Leftrightarrow 等の
    # 大文字マクロにもそもそも一致しない）。ここで剥がした後の空白・カンマは後段の
    # \s+ 畳み込み（下）と iteration 160 の数字間カンマ空白除去にそのまま乗るため、
    # \left(3, 4\right) -> (3, 4) -> (3,4) と、素の (3,4) と完全一致するようになる。
    # \frac/\sqrt/\operatorname 等の値を持つマクロ、素の ( )/[ ]、エスケープ済み中括弧
    # \{/\} は対象外のまま維持する。iteration 13/22/24/30/78/122/134/136/140/148/160 と
    # 同系統の「ファストパスで拾えない票を拾う」姉妹修正で、iteration 140 が明示的に
    # 先送りにしたギャップを埋める。
    s = re.sub(r"\\(?:left|right)\b", "", s)
    s = re.sub(r"^(?:the\s+)?(?:final\s+)?(?:answer|答え|正解)\s*(?:is|[:：は])?\s*",
               "", s, flags=re.IGNORECASE)
    # 桁区切りの除去。"11,\! 111,\! 100" のように区切り後に空白が入る表記（MATH-500 の
    # 正解表記で実在、2026-07-12 実測）も潰すため \s* を許容する
    s = re.sub(r"(?<=\d),\s*(?=\d{3}\b)", "", s)   # 12,345 / 12, 345 → 12345
    # 2026-07-22: 末尾の「,」を句読点と同様に除去する。extract_final_answer の数値抽出正規表現
    # -?\d[\d,]*(?:\.\d+)?(?:\s*/\s*\d+)? は [\d,]* が桁区切りでない末尾のカンマ（例:
    # "the answer is 42," の "42,"）も貪欲に飲み込んでしまい、上の桁区切り除去（3桁が後続する
    # カンマのみ対象）では取り切れない。末尾カンマ付きの文字列のまま投票させると、
    # answers_equivalent の正規化一致に失敗して票を落とす（Fraction("42,") は例外）か、
    # 万一 math_verify のフォールバックでも拾えず誤った文字列がそのまま最終回答として出力される。
    # 「無投票／誤投票より綺麗な1票」の方針（精度優先）に従い、末尾の「,」は句読点と同様に
    # ここで落とす。内部（末尾以外）のカンマは rstrip の性質上そのまま保持される。
    s = s.rstrip("。．.,").strip()
    s = re.sub(r"\s+", " ", s)
    # 2026-07-25: 数字に挟まれたカンマの直後の空白を除去する。(3, 4) と (3,4) は同じ
    # 順序対/区間/座標の値だが、上の桁区切り正規表現 (?<=\d),\s*(?=\d{3}\b) はカンマの
    # 直後にちょうど3桁が続く場合しか吸収しない（"4)" は \d{3} に一致しない）ため、
    # このままでは "(3, 4)" と "(3,4)" が別々の投票クラスに分裂する。MATH-500 では
    # 順序対・区間・座標形式の正解が珍しくなく、いずれも na.lower()/Fraction の高速
    # パスに乗らないため math_verify のフォールバックに落ちるが、gotcha #6 の通り
    # math_verify は Windows 上で parsing_timeout/timeout_seconds のマルチプロセス
    # タイムアウト実装に問題があり信頼できない。自己整合性投票（gotcha #7）は票割れに
    # 最も弱いため、ここで吸収して1票にまとめる。カンマの両側を厳密に数字だけに限定
    # することで、"yes, it is" のような散文カンマや (x, y) のような記号的タプルの
    # カンマは対象外のままとし、値を変えずに表記だけを揃える（精度優先）。上の桁区切り
    # 正規表現・空白畳み込みより後段に置くことで "12, 345" -> "12345" の桁区切り処理を
    # 妨げない。iteration 13/22/24/30/78/122/134/136/140/148 と同系統の姉妹修正。
    s = re.sub(r"(?<=\d),\s+(?=\d)", ",", s)
    # 2026-07-23: \text/\mathrm 以外にも \textbf/\mathbf/\boldsymbol などの
    # 「見た目だけ」を変える体裁マクロで最終回答を装飾するプロポーザーがいる。
    # mcq では \boxed{\textbf{B}} を extract_final_answer の先頭選択肢文字
    # 正規表現が読めず None（無投票）になり、math では \boxed{\mathbf{42}} が
    # 文字列 "\mathbf{42}" のまま投票されて answers_equivalent の
    # na.lower()/Fraction 系ファストパスに乗らず、素の "42" 票と別の投票クラスに
    # 分裂してしまう（自己整合性投票 gotcha #7 での票落ち・票割れ）。これらは
    # 値を変えない純粋な体裁指定なので剥がしても精度は落とさず、むしろ本来1つ
    # であるべき票をまとめられる（精度優先・時間は気にしない）。\frac/\sqrt/
    # \operatorname 等の値を変えうるマクロは対象外のまま維持する。\mathbf{\text{D}}
    # のような入れ子に対応するため、固定回数の上限（暴走防止）付きループで剥がす。
    _wrap_re = re.compile(
        r"\\(?:text|mathrm|textbf|textit|texttt|textsf|textnormal|"
        r"mathbf|mathit|mathsf|mathtt|boldsymbol|bm)\{(.*)\}"
    )
    for _ in range(4):
        m = _wrap_re.fullmatch(s)
        if not m:
            break
        s = m.group(1).strip()
    # 2026-07-24: \frac{a}{b}/\dfrac{a}{b}/\tfrac{a}{b}（純粋に数値のみ、非入れ子）を
    # "a/b" 形に正規化する。iteration 108 で同じ修正を試みて3回スタックしたまま未着手
    # だったギャップ（gotcha #7 の自己整合性投票がここに直結）で、\boxed{\frac{1}{2}} と
    # \boxed{0.5}/\boxed{1/2} が本来同じ値なのに別の投票クラスに分裂し、票の合算を
    # 過小評価する（math_verify に頼れば拾えるが遅く、_FW_TRANS/Fraction 比較等の
    # 既存ファストパス優先方針に反する）。前回の停滞を踏まえてスコープを最小化し、
    # ネストした中括弧（\frac{\frac{1}{2}}{3} 等）や変数を含む場合（\frac{x}{2} 等）は
    # 正規表現で捕捉せずそのまま素通りさせ、既存の math_verify フォールバックに委ねる
    # （クラッシュさせない・誤った値を作らない）。iteration 13/22/78 の姉妹的な
    # normalize_answer 修正群と同じ「ファストパスで拾えない票を拾う」方針に沿う。
    _frac_re = re.compile(
        r"^(-?)\\[dt]?frac\{\s*(-?\d+(?:\.\d+)?)\s*\}\{\s*(-?\d+(?:\.\d+)?)\s*\}$"
    )
    m = _frac_re.match(s)
    if m:
        sign, num, den = m.group(1), m.group(2), m.group(3)
        # 2026-07-25: iteration 122 は分子・先頭の符号だけを neg へ XOR しており、
        # 分母の符号は無視して den をそのまま埋め込んでいた（同コメント末尾
        # 「numerator/leading sign」限定のスコープ通り）。しかし _frac_re の分母グループ
        # は (-?\d+...) で先頭マイナスを許容するため、\frac{1}{-2}（=-1/2）は
        # "1/-2" に、\frac{-1}{-2}（=+1/2）は "-1/-2" になってしまう。Fraction() は
        # 分母に符号が付いた文字列を受け付けないため、これらは下の Fraction 高速パスに
        # 乗れずすり抜け、-1/2・\frac{-1}{2}・-0.5 という「素直な」表記の票とは別の
        # 投票クラスに分裂して自己整合性投票（gotcha #7）の票を薄める。iteration
        # 13/22/24/30/78/122/134/136/140 と同系統の「ファストパスで拾えない票を拾う」
        # 修正として、分母の符号も neg へ XOR してから den から取り除く。正の分母
        # （den.startswith("-") が False）の場合は lstrip が no-op のため、既存の
        # iteration 122/134/136/140 の正規化結果は一切変わらない。
        neg = (sign == "-") ^ num.startswith("-") ^ den.startswith("-")
        num = num.lstrip("-")
        den = den.lstrip("-")
        s = f"{'-' if neg else ''}{num}/{den}"
    # 2026-07-27 (iteration 214、iteration 210 が3回スタックしたまま未着手だったギャップの
    # 縮小スコープでのリトライ): PoT サンプルは stdout に sympy 綴り（sqrt(2)、pi）を出力する
    # 一方、CoT サンプルは \boxed{} に LaTeX 綴り（\sqrt{2}、\pi）を書く。na.lower() は
    # 不一致、Fraction() は例外で、同じ無理数の答えが math_verify フォールバックに落ちる
    # （gotcha #6: Windows では parsing_timeout/timeout_seconds のマルチプロセス実装に
    # 問題があり信頼できないため、正規化側の高速パスで吸収できる形は吸収しておきたい）。
    # 素の値が変わらず同じ2系統に分裂するのは自己整合性投票（gotcha #7）が最も嫌う
    # 票割れであり、iteration 122 の _frac_re と全く同じ「アンカー付き全文一致のみ」の
    # 形で最小スコープに絞ってリトライする。answers_equivalent（L4062-4064 付近）は
    # ここで正規化した文字列をそのまま $na$/$nb$ として math_verify に渡すため、もし
    # \pi をグローバルな部分文字列置換で扱うと 2\pi や \frac{\pi}{4} のような複合式を
    # 2pi や \frac{pi}{4} に書き換えてしまい、LaTeX パーサはこれを 2*p*i と読んで今日
    # 通っている等価判定を壊しかねない。そのため両ルールとも fullmatch 相当の
    # ^...$ アンカーで「値全体がちょうどこの形」の場合だけに限定し、2\sqrt{3} や
    # \sqrt{\sqrt{2}}、\sqrt{x}、\sqrt[3]{8}、\sqrt2（波括弧なし）、\sqrt{2}/2、2\pi、
    # \pi/2、\frac{\pi}{4} のような複合/入れ子/係数付き/変数入りの形は一切触れず
    # バイト単位で素通りのまま維持する（この位置は _wrap_re の展開・_frac_re の後段の
    # ため、\textbf{\sqrt{2}} のような体裁マクロ入れ子も既に剥がされた後に評価される）。
    _sqrt_re = re.compile(r"^(-?)\\sqrt\{\s*(\d+(?:\.\d+)?)\s*\}$")
    m = _sqrt_re.match(s)
    if m:
        sign, radicand = m.group(1), m.group(2)
        s = f"{sign}sqrt({radicand})"
    _pi_re = re.compile(r"^(-?)\\pi$")
    m = _pi_re.match(s)
    if m:
        s = f"{m.group(1)}pi"
    return s


def extract_final_answer(text, task_type="math"):
    """回答テキストから最終答を抽出する（見つからなければ None）。
    優先順: \\boxed{} > 「答え/Answer」宣言 > 最後の数値。mcq は選択肢文字 A-E を返す。"""
    if not text:
        return None
    text = strip_think(text)
    boxed = extract_boxed(text)
    if task_type == "mcq":
        if boxed:
            # 2026-07-21: 「A-E のどれかを本文中どこでも最後の1文字」ではなく、boxed 内容の
            # 先頭にある選択肢文字だけを拾う。答えの文字は慣例的に boxed の先頭に来るため、
            # \boxed{C, because it is the largest} のような散文混じりでも "Because" の B を
            # 誤って拾わない。先頭にマッチしなければ（\boxed{None of the above} 等）誤った
            # 文字を返さず、下の宣言パターン探索 → 最終的に None（無投票、誤投票より安全）に
            # フォールスルーさせる。
            # 2026-07-24: \(? は ASCII 括弧しか許容しておらず、CJK寄りのプロポーザー
            # （qwen/gemma 系）が好む全角括弧 \boxed{（A）} を通すと、normalize_answer/
            # _FW_TRANS（iter 13）は全角数字・A-E・マイナス・小数点・スラッシュ・カンマは
            # 正規化するが全角括弧（U+FF08/FF09）は意図的に対象外のため （ がそのまま残り、
            # \(? は幅ゼロでマッチ済み、続く ([A-E]) が （ に一致できず extract_final_answer
            # が None を返して自己整合性投票（gotcha #7）から正当な1票が無投票のまま
            # 静かに失われていた。[(（]? に広げて全角括弧も許容し、この票落ちを回復する
            # （無投票より正しい1票、精度優先・時間は気にしない）。iter 3 の \b 境界ガードと
            # iter 26 の複数文字競合時の棄権ロジックはそのまま維持する。
            # 2026-07-25: SC_PROMPT_MCQ は選択肢文字を \boxed{} に入れるよう指示しているが、
            # 上のコメント群（iter 3/26/102/109）が繰り返し記録している通り CJK 寄りの
            # プロポーザー（qwen/gemma 系）は無視して散文で答えを書く。その散文寄りの答えを
            # 「answer is **B**」「単独行の *C*／__A__」のように Markdown の強調記号で
            # 装飾するのは LLM の出力で非常によくある癖で、strip_think も normalize_answer も
            # '*'/'_' を除去しないため、装飾された選択肢文字はこの3か所のどれにもマッチせず
            # None（無投票）になり、自己整合性投票（gotcha #7）から正当な1票が静かに
            # 失われていた。文字に直接隣接する強調記号だけを [*_]{0,2} で束縛して許容し
            # （'*'/'_' をテキスト全体から一律に剥がすと math 分岐の下付き添字等を壊すため
            # 厳禁）、この票落ちを回復する。iter 3 の \b／(?![A-Za-z]) 境界ガードは強調記号が
            # 非単語文字であるため自然に維持され（'**Bee**' の 'B' の直後は依然 'e' で
            # 境界不成立のまま）、iter 26 の複数文字競合時の棄権ロジックもそのまま変えない。
            m = re.match(r"[*_]{0,2}[(（]?\s*([A-E])\b", normalize_answer(boxed).upper())
            if m:
                return m.group(1)
        # 2026-07-22: 連結詞（is/：/は）を省略可にしていたせいで、「answer A」のような
        # 単なる言及（例:「Note that answer A was a common distractor.」）まで宣言と
        # 誤認していた。本物の宣言は「answer is B」「答え：B」のように連結詞を伴うのが
        # 通例なので、下の math 宣言ブランチ（L2334 付近）と同じく連結詞を必須にする。
        # 2026-07-24: 上の boxed 分岐と同じ理由（全角括弧はカッコ内はA-E数字と異なり
        # _FW_TRANS 未対応、iter 13 参照）で、宣言パターン「答えは（B）です」・単独行
        # パターン「（C）」のどちらも \(?/\)? のままだと全角括弧を素通りできず None を
        # 返し無投票になる。[(（]?/[)）]? に広げて ASCII と同様に許容する。iter 3 の
        # \b・先頭文字ガードと iter 26 の複数文字競合時 None 化（下の len(letters) > 1）は
        # 変更しない。
        # 2026-07-24: iter 102 は上記の全角括弧対応時に宣言パターン（1行下）の文字クラスは
        # [A-EＡ-Ｅ] に広げたが、こちらの単独行パターンは [A-E] のまま ASCII 限定に
        # 取り残されていた。そのため CJK 寄りのプロポーザー（qwen/gemma 系、iter 102 が
        # 想定した対象そのもの）が \boxed{} を無視して「Ｃ」「（Ｃ）」のように全角文字だけを
        # 単独行で答えると、boxed 分岐（該当なし）・宣言パターン（answer/答え/正解 の
        # 連結詞なし）・この単独行パターン（[A-E] は U+FF23 に不一致）のいずれにも拾われず
        # None となり、自己整合性投票（gotcha #7）の正当な1票が静かに失われていた。
        # 宣言パターンと同じ [A-EＡ-Ｅ] に揃えて票落ちを解消する。.translate(_FW_TRANS)
        # （iter 13）が全角→ASCII正規化を担うため誤投票のリスクはなく、iter 26 の
        # 複数文字競合時の棄権ロジックもそのまま維持する。
        # 2026-07-25: boxed 分岐と同じ理由（上のコメント参照、iter 3/26/102/109 系列）で、
        # 宣言パターン「answer is **B**」・単独行パターン「*C*」「__A__」も Markdown の
        # 強調記号を素通りできず None（票落ち）になっていた。文字に直接隣接する
        # [*_]{0,2} のみを括弧の外側（既存の [(（]?/[)）]? の外）に許容し、iter 3 の
        # \b／(?![A-Za-z]) 境界ガード（強調記号は非単語文字なので '**Bee**' の誤爆は
        # 引き続き起きない）・iter 26 の複数文字競合時 None 化・単独行パターンの
        # ^...$ MULTILINE アンカー（'**A. Introduction**' のような見出し行を拾わない）は
        # いずれもそのまま維持する。
        for pat in (r"(?:answer|答え|正解)\s*(?:is|[:：は])\s*[*_]{0,2}[(（]?([A-EＡ-Ｅ])[)）]?[*_]{0,2}(?![A-Za-z])",
                    r"^\s*[*_]{0,2}[(（]?([A-EＡ-Ｅ])[)）]?[*_]{0,2}\s*(?:が正解|です)?\s*$"):
            ms = re.findall(pat, text, re.IGNORECASE | re.MULTILINE)
            if ms:
                letters = {m.translate(_FW_TRANS).upper() for m in ms}
                # 2026-07-22: 連結詞を必須にしても、「the answer is B; oh wait, the
                # answer: A」のような言い直し・訂正では複数の宣言が異なる文字を指した
                # まま両方マッチしうる。ここで機械的に ms[-1]（最後のマッチ）を採用すると、
                # 訂正前/ディストラクタ側の文字を確信ありの1票として投票してしまい、
                # 自己整合性投票（gotcha #7）を汚染する。「無投票 > 誤投票」の方針
                # （精度優先・時間は気にしない）に従い、競合する文字が混在する場合は
                # ここで確定させず None を返して棄権する。全て同一文字で一致する場合のみ
                # その文字を返す。
                if len(letters) > 1:
                    return None
                return letters.pop()
        return None
    if boxed:
        return normalize_answer(boxed) or None
    ms = re.findall(r"(?:final answer|answer|答え|正解)\s*(?:is|[:：は])\s*([^\n]{1,60})",
                    text, re.IGNORECASE)
    if ms:
        # 2026-07-22: iteration 26 で MCQ 宣言ブランチ（上の L2326-2345 付近）に適用した
        # 「複数宣言が競合する場合は無投票（None）とする」修正の math 版（gotcha #7）。
        # 以前はここで ms[-1] だけを見ていたため、「答えは24 …実は答えは12」のような
        # 言い直し・訂正や、ディストラクタに触れる文でも最後の宣言だけを機械的に採用し、
        # 訂正前/誤った値を確信ありの1票として self-consistency 投票（solve_verifiable）に
        # 混入させていた。ここでは各宣言から数値部を抽出して候補リストを作り
        # （空/抽出不能な宣言は候補に数えず単にスキップする）、answers_equivalent で
        # 最後の候補と相互比較する。全候補が同値（例: 1/2 と 0.5、1,000 と 1000）なら
        # 単一の票として確定させ、一つでも非同値な候補が混在すれば None を返して棄権する
        # （無投票 > 誤投票、精度優先・時間は気にしない）。
        cands = []
        for raw in ms:
            cand = normalize_answer(raw)
            if not cand:
                continue
            # 2026-07-27: mcq 宣言分岐（iter 173、この関数の上のほう L3795 以降のコメント
            # 参照）に施した Markdown 強調記号の束縛許容がこちらの math 宣言分岐には
            # 適用されておらず非対称なままだった。「The answer is **42**.」のように
            # normalize_answer 後も cand が "**42**" のまま Markdown の強調記号を
            # 保持していると、直後の数値コア re.match（下）は '*' の位置で先頭マッチに
            # 失敗し、m が None のまま cand（"**42**" という装飾込みの生文字列）自体が
            # そのまま候補に append されていた。これは二重に有害で、(1) Fraction("**42**")
            # は例外になり、math_verify も Windows では信頼できない（gotcha #6）ため
            # "42" という素の票と永久に合流しない別クラスの誤投票として
            # self-consistency 投票（solve_verifiable、gotcha #7）を汚染し、(2) さらに
            # 悪いことに「answer is **12**」と「answer is 12」のように同じ値を指す
            # 2つの宣言が cands=['**12**','12'] という非同値な候補ペアに分裂してしまい、
            # 全会一致のはずの正当な1票が iteration 30 の複数候補非同値棄権ロジック
            # （すぐ下の any(not answers_equivalent(...))）に誤って引っかかって
            # None（無投票）に化けてしまう。iter 173 と同じ設計方針（強調記号は
            # 値に隣接する境界だけを束縛して剥がし、テキスト全体から一律に '*'/'_' を
            # 剥がすことは a_1 のような下付き添字や a*b のような乗算表記を壊すため
            # 厳禁のまま、というのがそこのコメントの要旨）をここでも踏襲し、cand の
            # 先頭・末尾に限定して [*_]{1,3} を剥がす。normalize_answer 自体（iteration
            # 13/22/24/78/122/134/136/140/148/160/164 の蓄積）や iteration 30 の
            # 非同値棄権ロジック、iteration 136 の '%' 保持はいずれも変更しない。
            cand = re.sub(r"^[*_]{1,3}\s*|\s*[*_]{1,3}$", "", cand)
            if not cand:
                continue
            # 「700 円です」のような後置き単位・助詞を落とす: 数値で始まるなら数値部のみ
            # 2026-07-22: 整数部の文字クラスを \d[\d,]* から「桁区切りとして妥当な
            # \d{1,3}(?:,\d{3})+ か、カンマなしの \d+」の二択へ厳格化する（iteration 13/22/24 の
            # 続き）。旧パターンはカンマを桁区切りかどうか問わず貪欲に飲み込むため、宣言文が
            # 「答えは 1,2,3 です」のようにカンマ区切りの数値列で終わると数値部が "1,2,3" と
            # いう1トークンに誤って結合される。normalize_answer の桁区切り除去（L2306、後続
            # ちょうど3桁のみ対象）はこれを吸収できず、"1,2,3" のまま answers_equivalent へ渡ると
            # Fraction("1,2,3") が例外になり、自己整合性投票（solve_verifiable、gotcha #7）へ
            # 単独の誤投票クラスとして混入する。「無投票 > 誤投票」の方針（精度優先・時間は
            # 気にしない）に従い、正当な桁区切り（1,234 / 1,234,567 等）は1トークンのまま保つ
            # 一方、桁区切りとして不正なカンマ列は数値ごとに分離させる。
            # 2026-07-25 (iteration 136): この数値コア正規表現は '%' を含めていなかった
            # ため、"The final answer is 50%" のような非boxed宣言では cand が
            # normalize_answer 経由で "50%"（iteration 134 が '%' を保持するよう修正済み、
            # 50% を 50 という別の値に化けさせない = 精度劣化を防ぐ）になっていても、
            # ここで数値部だけを切り出す際に '%' を落として "50" にしてしまっていた。
            # 一方 \boxed{50\%} 側は normalize_answer(boxed) をそのまま返す boxed 分岐
            # （この関数の上のほう）を通るため '%' が残る。同じ50%という値が
            # boxed 経路の '50%' 票と宣言経路の '50' 票という別クラスに分裂し、
            # 自己整合性投票（solve_verifiable、gotcha #7）が本来1票にまとまるはずの
            # 票を薄めてしまう（無投票ならまだしも、ここでは意味の異なる整数値へ
            # 誤って投票する点がより悪い）。数値本体に「直後」（間に \s* を挟まない）で
            # 続く '%' だけを任意で拾う %? を末尾に足し、"50 %"のような離れた%は
            # 従来通り拾わない（新たな別クラスを作らない）。iteration 13/22/24 の
            # 桁区切り厳格化・符号クラス拡張と同系統の姉妹修正であり、iteration 30 の
            # 複数候補の非同値棄権ロジック（すぐ上）はそのまま維持する。
            # 「無投票／誤投票より綺麗な1票」の方針（精度優先・時間は気にしない）。
            m = re.match(r"-?(?:\d{1,3}(?:,\d{3})+|\d+)(?:\.\d+)?(?:\s*/\s*\d+)?%?", cand)
            cands.append(m.group(0).replace(" ", "") if m else cand)
        if cands:
            if any(not answers_equivalent(c, cands[-1]) for c in cands[:-1]):
                return None
            return cands[-1]
    # 2026-07-22: 最後の数値フォールバックは RAW text に対して ASCII の "-?" だけを符号として
    # 拾っていたため、直前の符号が Unicode マイナス U+2212（−）や全角ハイフンマイナス U+FF0D
    # （－）だと拾えず、負の答えが正の値として投票されてしまう（例:「結果は −5。」→ "5"）。
    # 宣言ブランチ（上の boxed / answer-is 分岐）は normalize_answer で全角→半角変換
    # （iteration 13 の _FW_TRANS）を先に済ませてから数値部を切り出すため符号が保持されるのに対し、
    # ここは正規化前の生テキストに正規表現をかけるため非対称になっていた。符号クラスに
    # U+2212/U+FF0D を追加して拾えるようにし、変換自体は既存の normalize_answer(nums[-1]) に
    # 任せる（二重正規化はしない）。
    # 2026-07-22: ここも上の宣言ブランチ（L2383 付近）と同じ理由で整数部を
    # \d{1,3}(?:,\d{3})+ | \d+ に厳格化する（iteration 13/22/24 と同じ抽出経路の姉妹修正）。
    # \boxed{} も宣言文もない本文が「the roots are 1,2,3」や「the point is (1,2)」のように
    # カンマ区切りの数値列で終わる場合、旧パターンの [\d,]* は "1,2,3" / "1,2" ごと1トークンとして
    # 貪欲に飲み込んでしまい、本来拾うべき最後の数値（3 / 2）ではなく Fraction() が例外を
    # 投げる壊れたトークンを最終フォールバック票として返していた。この誤投票は
    # solve_verifiable の自己整合性投票（gotcha #7）を汚染するため、「無投票 > 誤投票」
    # （精度優先・時間は気にしない）の方針に従い、桁区切りとして妥当なカンマ（1,234 や
    # 1,234,567）のみ1トークンにまとめ、そうでないカンマ区切りは個々の数値へ分離する。
    # 2026-07-25 (iteration 136): 上の宣言ブランチと同じ理由で、\boxed{} も「答え/正解/
    # answer」宣言も無く「so the probability is 50%.」のように percent で終わる本文だと、
    # この最後の数値フォールバックが '%' を含めていなかったため "50" を返し、
    # \boxed{50\%}（iteration 134 が '%' 保持へ修正済み）が返す "50%" 票と別クラスに
    # 分裂していた。数値本体に「直後」（\s* を挟まない）で続く '%' だけを任意で拾う
    # %? を末尾に足して merge する。findall は文中の各数値トークンを個別に拾うため、
    # 「increased by 50% to reach 75」のように途中の percent は自分の直後にしか
    # 付与されず、末尾の別の数値（75）にまで '%' が誤って伝播することはない。
    # iteration 13/22/24 の桁区切り厳格化・符号クラス拡張、iteration 30 の宣言側の
    # 棄権ロジックと同系統の姉妹修正（無投票／誤投票より綺麗な1票、精度優先・
    # 時間は気にしない）。
    nums = re.findall(r"[-−－]?(?:\d{1,3}(?:,\d{3})+|\d+)(?:\.\d+)?(?:\s*/\s*\d+)?%?", text)
    return normalize_answer(nums[-1]) if nums else None


# 2026-07-27: MCQ の自己一貫性投票（gotcha #7）で観測されている票落ちパターンへの
# 対処、第一段。プロポーザーが選択肢の文字（A-E）ではなく選択肢の「値」そのものを
# \boxed{} に入れてしまうことがあり（例:「4^2 は？ A) 12 B) 16 C) 20 D) 24」に対し
# \boxed{16} と回答）、この場合 extract_final_answer(text, 'mcq') は boxed 先頭の
# [A-EＡ-Ｅ] しか読まないため None を返し、正当な1票が自己整合性投票から静かに
# 失われる。この「値を箱に入れてしまう」化けそのものは stuck した iteration 207/209
# が「パーサーの実装」と「_sc_sample への配線」を同じイテレーションで両方やろうと
# して行き詰まった経緯があり、直近では代表文（サンプル群の勝者テキスト）側で
# デッキレベルの \boxed 再整合を行う対策（上の fugu_answer、SC 投票結果を
# \boxed{} で包み直して再抽出とのずれを解消する修正）が入ったが、これは
# _sc_sample が個々のサンプルから投票用の値を抽出する時点（vote_answers に渡す
# 前）の票落ちそのものは救えていない。
# 以下はその救済の下ごしらえとして、question_text 中の選択肢列（_MCQ_SIGNALS,
# 上のL2952-2955、が要求するのと同じ最小限の行頭マーカー形状）を解析し、value が
# その選択肢本文のどれと一意に一致するかを answers_equivalent 経由で判定する、
# 副作用のない純粋関数。_sc_sample/extract_final_answer への配線（実際にこの
# 関数を投票経路で呼ぶこと）は本イテレーションのスコープ外であり、次イテレーション
# へ明示的に先送りする（本関数は現時点でどこからも呼ばれていない）。
def map_value_to_choice(value, question_text):
    """value を question_text 中の選択肢文字(A-E)へ一意にマップできればその文字を返す。
    マップできなければ None（無投票 > 誤投票、精度優先・時間は気にしない）。

    選択肢列の認識形状は _MCQ_SIGNALS と同じ行頭アンカー限定: 行頭（または直前が
    改行）、任意の空白、任意の '(' / 全角 '（'、選択肢文字(A-EまたはＡ-Ｅ)、続けて
    ')' '.' ':' '：' '）' '．' のいずれか一つ。選択肢本文は同じ行の残り部分のみを
    対象とし、複数行にまたがる選択肢本文は本イテレーションのスコープ外とする
    （_MCQ_SIGNALS 自体も単一行前提であり、それに合わせている）。小文字マーカー
    （a) b) 等）は「a) の場合は」のような地の文の箇条書き様の記述を選択肢と
    誤認するリスクを避けるため今回は非対応とし、大文字ASCII/全角のみを認識する
    （プロポーザーが小文字マーカーの選択肢を書くことは実運用上ほぼ無い）。

    副作用なし（ask()/ネットワーク/subprocess/ファイルI/Oなし、モジュールグローバル
    変更なし）で、任意のテキスト入力に対して例外を投げない（bare except は使わず、
    純粋な re.finditer とループだけで構成する）。
    """
    if not value:
        return None
    # 値そのものが既に単独の選択肢文字（大小文字・全角含む）なら救済不要。
    # ここで選択肢本文と比較してしまうと誤爆のリスクがあるため触れずに None を返す
    # （iteration 129 の answers_equivalent 内 A-E ガードと同じ判断: 精度優先）。
    nv = normalize_answer(value)
    if not nv:
        return None
    if re.fullmatch(r"[A-Ea-e]", nv):
        return None
    text = question_text or ""
    # マーカー: 行頭 or 直前が改行、空白、任意の開き括弧、選択肢文字、閉じマーカー。
    # 本文は同じ行の残り（\S から行末まで、'.' はデフォルトで改行に一致しないため
    # 複数行へは絶対に広がらない）のみを取る。
    option_re = re.compile(r"(?:^|\n)[ \t]*[(（]?([A-EＡ-Ｅ])[)\.:：）．][ \t]*(\S.*)")
    seen = {}
    for m in option_re.finditer(text):
        letter = m.group(1).translate(_FW_TRANS).upper()
        if letter in seen:
            continue  # 選択肢列がプロンプト中で再掲されている場合は先勝ち
        seen[letter] = m.group(2)
    # _MCQ_SIGNALS と同じ最小限の証拠（A・Bが両方そろって初めてMCQの選択肢列とみなす）
    if len(seen) < 2 or "A" not in seen or "B" not in seen:
        return None
    matches = [letter for letter, body in seen.items() if answers_equivalent(value, body)]
    if len(matches) != 1:
        return None
    return matches[0]


def answers_equivalent(a, b):
    """2つの答えが数学的に同値か。正規化一致 → 分数/小数の数値一致 → math_verify の順で判定。"""
    na, nb = normalize_answer(a), normalize_answer(b)
    if not na or not nb:
        return False
    if na.lower() == nb.lower():
        return True
    # 2026-07-25 (iteration 129): iteration 107 が指摘したが当時は未修正のまま残した懸念
    # への対応（gotcha #6/#7参照）。MCQ の SC サンプルは単一の選択肢文字 A-E であり、ここまで
    # 到達した時点で na.lower()==nb.lower() は不一致（=文字として異なる）ことが確定している。
    # このまま下へ抜けると Fraction('A') は例外で次段へフォールスルーし、最終的に
    # math_verify フォールバックまで到達しうる。だが sympy は 'E' を自然対数の底
    # (Euler's number)、'I' を虚数単位として数式解釈するため、選択肢文字を数式として解析
    # するのは意味論的に不健全であり、万一「同値」と誤判定されれば solve_verifiable の
    # 自己整合性投票（gotcha #7）で本来別クラスの票が併合され多数決が汚染されかねない。
    # この開発環境には math_verify が未インストールで実挙動を検証できない（import失敗の
    # 例外はexceptでFalseに握り潰されるだけ）ため、決定的なガードとして「両方とも単一の
    # 選択肢文字 A-E」の形をしている場合は Fraction/math_verify に一切渡さず、ここで確定的に
    # 判定して返す（math_verify呼び出し自体・gotcha #6のparsing_timeout=None/
    # timeout_seconds=Noneは変更しない）。na.lower()==nb.lower()が既に不一致なのでこの分岐の
    # 戻り値は常にFalseになるが、意図を明示するため case-insensitive 比較として書く。
    if re.fullmatch(r"[A-Ea-e]", na) and re.fullmatch(r"[A-Ea-e]", nb):
        return na.upper() == nb.upper()
    try:
        from fractions import Fraction
        if Fraction(na.replace(" ", "")) == Fraction(nb.replace(" ", "")):
            return True
    except Exception:
        pass
    try:
        # タイムアウトは無効化して呼ぶ: math_verify の既定タイムアウトは multiprocessing を
        # 使い、Windows では __main__ 再import でハンドルエラーを撒き散らす（実測 2026-07-11）。
        # 答えは短い文字列（<=80字）なので sympy が固まるリスクは実用上無視できる。
        import logging as _logging
        _logging.getLogger("math_verify").setLevel(_logging.ERROR)
        from math_verify import parse as _mv_parse, verify as _mv_verify
        return bool(_mv_verify(_mv_parse(f"${na}$", parsing_timeout=None),
                               _mv_parse(f"${nb}$", parsing_timeout=None),
                               timeout_seconds=None))
    except Exception:
        return False


def vote_answers(answers):
    """答えリストを同値クラスへ集約し (最多答, その票数, クラス一覧) を返す。
    クラス一覧は [[代表答え, 票数], ...] 票数降順。答えが無ければ (None, 0, [])。"""
    classes = []
    for a in answers:
        if not a:
            continue
        for c in classes:
            if answers_equivalent(a, c[0]):
                c[1] += 1
                break
        else:
            classes.append([a, 1])
    if not classes:
        return None, 0, []
    classes.sort(key=lambda c: -c[1])
    return classes[0][0], classes[0][1], classes


def _sc_sample(model, question, task_type, pot=False, history=None):
    """SC の 1 サンプル。(answer, full_text) を返す（抽出/実行失敗は answer=None）。"""
    if pot:
        sysp = SC_PROMPT_POT
    else:
        sysp = SC_PROMPT_MCQ if task_type == "mcq" else SC_PROMPT_MATH
    raw = ask(model,
              [{"role": "system", "content": sysp}] + list(history or [])
              + [{"role": "user", "content": question}],
              SC_TEMP, think=proposer_think_for(model),
              num_predict=proposer_predict_for(model), label="sc")
    if raw.startswith("__ERROR__"):
        return None, raw
    text = strip_think(raw)
    if pot:
        code = extract_code(text)
        if not code:
            return None, text
        ok, out = run_python(code, timeout=SC_POT_TIMEOUT, stdout_only=True)
        out = (out or "").strip()
        if not ok or not out:
            return None, text
        ans = out.splitlines()[-1].strip()
        if not ans:
            return None, text
        # 2026-07-25: PoT分岐(ここ)とCoT分岐(直下の extract_final_answer)で答えの抽出方法が
        # 非対称だった。CoT側はextract_final_answerがtextから\boxed{}を剥がしてから
        # normalize_answerへ渡すが、PoT側はstdout最終行(iteration 4)をnormalize_answerへ
        # 素通しするだけで\boxed{}を剥がしていなかった。数学寄りにチューニングされた
        # モデルはコード内でprint(f'\\boxed{{{ans}}}')のように反射的に\boxed{}で答えを
        # 包むことがあり、その場合stdout最終行が丸ごと"\boxed{42}"になる。normalize_answer
        # は$・\%・\(\)\[\]・\text/\mathbf系ラッパは剥がすが\boxed{}は対象外なので、
        # 素の"42"というCoT側の投票クラスとは別クラスに割れ、gotcha #7の自己整合性投票
        # （CoTの計算ミスをPoTで裏取り/上書きするための仕組み）でPoT票が丸ごと死票化して
        # いた。ここで最終行に\boxed{が含まれる場合のみextract_boxedで中身を取り出し、
        # 取れた場合だけそれをnormalize_answerへ渡す。取れなかった場合（末尾が
        # \boxed{や\boxed{}のように閉じていない/空、iteration 11/23/25と同じ安全側判定）は
        # 従来どおり生の行をnormalize_answerへ渡す。値を捏造しない・既存の票を壊さない
        # 「票を拾えるところだけ拾う」方針は iteration 11/23/25/78/122/134/136/140/148の
        # 票救出系修正列と同じ（精度優先・時間は気にしない）。長さガード(len>80)は
        # 従来どおり最終的に投票される値に対して適用する（生の行ではなく剥がした後の値）。
        vote_src = ans
        if "\\boxed{" in ans:
            unwrapped = extract_boxed(ans)
            if unwrapped is not None:
                vote_src = unwrapped
        if not vote_src or len(vote_src) > 80:
            return None, text
        return (normalize_answer(vote_src) or None), text + f"\n\n[PoT execution output]\n{out[-500:]}"
    return extract_final_answer(text, task_type), text


def _representative_text(samples, answer):
    """勝ったクラスの代表解答テキスト（CoT 優先、無ければ PoT、最後は最長）を返す。"""
    fallback = None
    for s in samples:
        if s["answer"] and answers_equivalent(s["answer"], answer):
            if not s["pot"]:
                return s["text"]
            fallback = fallback or s["text"]
    if fallback:
        return fallback
    if samples:
        return max(samples, key=lambda s: len(s.get("text") or "")).get("text") or ""
    return ""


def _arbitrate(question, task_type, samples, classes):
    """票が拮抗した上位 2 クラスの代表解答を突き合わせて裁定する。
    戻り値: (裁定answer, 裁定役自身の解答テキスト) のタプル。全裁定役が失敗/空なら None。
    裁定役の優先順: ARBITER_MODEL（導入済みなら。既定 gpt-oss:120b の最上位知能）→
    それが失敗/空/未導入なら REASONING_MODELS の先頭へ堅牢にフォールバック。
    120b は 65GB で RAM(48GB)+VRAM(8GB) を超え NVMe ページングで非常に遅い可能性があるが、
    精度優先方針で「拮抗時だけ最上位が裁く」価値を取る。ダメでも degrade して止まらない。"""
    inst = installed_models()
    chain = []
    if ARBITER_MODEL and is_installed(ARBITER_MODEL, inst):
        chain.append(ARBITER_MODEL)
    for m in REASONING_MODELS:                 # 120b が失敗しても軽い思考モデルで必ず裁く
        if m in PROPOSERS and m not in chain:
            chain.append(m)
    if not chain:
        return None
    # 2026-07-22: 呼び出し側の拮抗判定(classes[0][1]==classes[1][1])は「トップ2クラスが
    # 同数」であることしか見ておらず、3クラス以上が同数タイになる N 択拮抗（例: 票数
    # [2,2,2]）もここに到達しうる。従来は classes[:2] で常に先頭2クラスしか裁定役に
    # 見せておらず、3番目以降の同数クラス（それが正解かもしれない）が黙って握りつぶされ、
    # プロンプトも "Two candidate solutions disagree" と決め打ちだった。ここではトップ
    # 票数と同数のクラスを全て（ただし num_ctx 保護のため上限 ARBITRATE_MAX_CANDIDATES
    # 件までに制限し、超過分は省略せずログに出す）候補として提示する。
    # 2件のみの通常拮抗では従来と完全に同じ挙動(先頭2件)になる。
    top_count = classes[0][1]
    tied = [c for c in classes if c[1] == top_count]
    if len(tied) > ARBITRATE_MAX_CANDIDATES:
        omitted = tied[ARBITRATE_MAX_CANDIDATES:]
        tied = tied[:ARBITRATE_MAX_CANDIDATES]
        omitted_desc = ", ".join(str(c[0]) for c in omitted)
        print(f"   [SC] {len(omitted)}件の同数タイ候補は上限のため裁定役に提示されません: {omitted_desc}")
    # 2026-07-25: 従来はタイ候補ごとに samples を先頭から走査し、canon と
    # answers_equivalent な最初のサンプルをそのまま裁定役へ見せていた(下のfor内break)。
    # add_batch(~L3731-3738)は各バッチの末尾でそのバッチのPoTサンプルを1件だけ追加する
    # ため、samples内では「あるバッチのPoTサンプル」が「後続バッチのCoTサンプル」より
    # 先に並ぶ。その結果、同じ答えにCoT一致サンプルが後から存在していても、先に並んだ
    # PoTサンプルの方(コード＋'[PoT execution output]'という実行結果ブロック)がそのまま
    # 「代表解答」として裁定役に渡り、「各候補の誤りを指摘して裁定せよ」という指示に対し
    # 自然言語の思考過程ではなくコードしか見せられない弱い入力になっていた。
    # _representative_text(iteration 2/55, L3594)はユーザー向け代表解答の選出で全く同じ
    # 状況をCoT優先(pot=Falseを優先、無ければ最初に一致したPoT、リスト順序に非依存)で
    # 既に解決しており、L3471-3472のコメントが指摘する通り _arbitrate 側のrep選出は
    # それとは無関係な別contractのまま放置されていた。ここでも同じCoT>PoT優先を適用する。
    # _representative_text の「一致なし→全サンプル中最長」フォールバックはここでは
    # 発火しえない: tied の各 canon は vote_answers/classes 経由でsamples中の実在の
    # answerそのものから作られたクラス代表であり、必ず>=1件のanswers_equivalentな
    # サンプルが存在するため(不一致候補やその場しのぎの最長テキストが紛れ込むことはない)。
    # strip_thinkと[:3000]切り詰め(num_ctx保護、gotcha #2)は選ばれたテキストへ従来通り適用する。
    reps = []
    for canon, _cnt in tied:
        rep_text = _representative_text(samples, canon)
        reps.append((canon, strip_think(rep_text or "")[:3000]))
    listing = "\n\n".join(
        f"### Candidate {chr(ord('A') + i)} (final answer: {c})\n{t}"
        for i, (c, t) in enumerate(reps))
    # 2026-07-22: iteration 16 でヘッダー行（"{len(reps)} candidate solutions
    # disagree:"）は候補数に依存しない表現に直したが、本文の指示文は "Carefully
    # check both, find the flaw in the wrong one" のまま2択決め打ちで残っていた。
    # 3/4択タイでは候補が3-4件（reps）提示されるのに "both" / "the wrong one"（単数）
    # と言われ、3件目以降への精査が手薄になる恐れがあった。「each candidate」
    # 「the incorrect one(s)」という候補数非依存の表現に統一する。2択の場合も
    # 意味は従来の "check both / find the flaw in the wrong one" と同等。
    # 2026-07-24: _arbitrate は math/mcq 両方の拮抗解消で共用される（呼び出し元は
    # solve_verifiable）が、末尾の出力形式指示はここまで math 前提の "put ONLY the
    # correct final answer in \boxed{}" 一本しか無かった。extract_final_answer(text,
    # 'mcq')（iter 3 で確立・iter 26/102 で誤爆修正済み）は \boxed{} の中身が選択肢
    # 文字 A-E で始まる場合しか採用せず、計算値や選択肢本文（\boxed{7} や
    # \boxed{Paris} 等）は None（無投票）になる。math 寄りの文言に引きずられた裁定役が
    # それを箱に入れると、下のループが全裁定役で ans=None のまま尽き、_arbitrate 全体が
    # None を返して mcq の拮抗が黙って MoA フォールバックへ劣化していた（iter 16/45 は
    # math 側の文言調整のみで、この mcq 側の出力形式齟齬は未対応だった）。
    # task_type=='mcq' のときだけ「\boxed{} には選択肢の文字(A-Eのいずれか1文字)だけを
    # 入れよ」と明示する。推論自体（"solve the problem yourself if needed"）は禁止せず、
    # 制約するのは最後の boxed トークンのみ。math/その他の文言は従来と完全に同一
    # （バイト単位で不変。回帰テストで固定）。
    if task_type == "mcq":
        final_instruction = (
            "Carefully check each candidate, find the flaw(s) in the incorrect one(s), "
            "and solve the problem yourself if needed. At the very end, put ONLY the "
            "single choice letter (A, B, C, D, or E) in \\boxed{} -- not the option's "
            "wording or a computed value, just that one letter.")
    else:
        final_instruction = (
            "Carefully check each candidate, find the flaw(s) in the incorrect one(s), "
            "and solve the problem yourself if needed. At the very end, put ONLY the "
            "correct final answer in \\boxed{}.")
    prompt = (f"Problem:\n{question}\n\n"
              f"{len(reps)} candidate solutions disagree:\n\n{listing}\n\n"
              + final_instruction)
    for arb in chain:
        print(f"   [SC] 票が拮抗 → {arb} が裁定します")
        raw = ask(arb, [{"role": "user", "content": prompt}], 0.1,
                  num_predict=model_cfg(arb, "num_predict", 8192), label="arbiter")
        # 2026-07-22: ask() は失敗時 raw を '__ERROR__: HTTP Error 500 ...' のような
        # 文字列で返す（line ~1079）。ここでチェックせずに extract_final_answer へ渡すと、
        # math タスクの最終数値フォールバック（line ~2299）がエラーメッセージ中の
        # '500'/'429'/'400' を「裁定役の最終解答」として誤採用してしまい、拮抗投票が
        # 誤った自信満々の数値に化ける。_sc_sample（iter 4）・ask() 自身（iter 9）・
        # _critic_judge/second_opinion（iter 15）で修正済みの同種バグがここでは未対応
        # だったので、raw.startswith('__ERROR__') と空/空白応答を同様に弾き、次の
        # 裁定役へフォールバックする。
        if not raw or not raw.strip() or raw.startswith("__ERROR__"):
            print(f"   [SC] {arb} の裁定が空/抽出不能 → 次の裁定役へ")
            continue
        text = strip_think(raw)
        ans = extract_final_answer(text, task_type)
        if ans:
            return ans, text
        print(f"   [SC] {arb} の裁定が空/抽出不能 → 次の裁定役へ")
    return None


def solve_verifiable(question, task_type="math", history=None):
    """Self-Consistency + PoT で math/mcq を解く。
    戻り値: {"answer", "text", "votes", "n_samples"}。票が全く得られなければ None
    （呼び出し側が通常の MoA へフォールバックする）。"""
    models = [m for m in REASONING_MODELS if m in PROPOSERS]
    if not models:
        models = list(PROPOSERS[:2])
    if not models:
        return None
    cheap_ok = (SC_CHEAP_VOTES > 0 and SC_CHEAP_MODEL
                and is_installed(SC_CHEAP_MODEL, installed_models()))
    samples = []

    def add(model, pot=False):
        ans, text = _sc_sample(model, question, task_type, pot=pot, history=history)
        samples.append({"answer": ans, "text": text, "model": model, "pot": pot})
        kind = "PoT" if pot else "CoT"
        print(f"   [SC {len(samples)}] {model} ({kind}) -> {ans if ans else '(抽出失敗)'}")

    def main_cot_count():
        return sum(1 for s in samples if not s["pot"] and s["model"] != SC_CHEAP_MODEL)

    # 【重要】OLLAMA_MAX_LOADED_MODELS=1 では毎サンプルでモデルを切り替えると 13〜23GB の
    # 再ロードが多発して致命的に遅い。そこで「モデルごとにまとめて」サンプリングし再ロードを
    # 最小化する（多様性は temp=0.7 の複数サンプルで確保）。各モデルから同数ずつ引く。
    def add_batch(n):
        per = max(1, n // len(models))
        for m in models:
            for _ in range(per):
                add(m)
        # PoT は先頭モデルがロード済みのうちに末尾で実行（追加ロードなし）
        if SC_POT and task_type == "math":
            add(models[0], pot=True)

    add_batch(SC_INITIAL)
    if cheap_ok:                       # 安価票は最後にまとめて（VibeThinker を1回ロード）
        for _ in range(SC_CHEAP_VOTES):
            add(SC_CHEAP_MODEL)

    while True:
        answers = [s["answer"] for s in samples if s["answer"]]
        top, cnt, classes = vote_answers(answers)
        n = len(answers)
        # 確定条件: 全会一致（ただし n < SC_MIN_VOTES の疑似全会一致は不可）、
        # または 4 票以上で過半数
        if top is not None and n > 0 and (
            (cnt == n and n >= SC_MIN_VOTES) or (n >= 4 and cnt * 2 > n)
        ):
            break
        if main_cot_count() >= SC_MAX:
            break
        head = [(c[0], c[1]) for c in classes[:3]]
        print(f"   [SC] 票が割れています {head} → 追加サンプリング")
        add_batch(SC_STEP)

    answers = [s["answer"] for s in samples if s["answer"]]
    top, cnt, classes = vote_answers(answers)
    if top is None:
        return None
    rep = None
    if len(classes) >= 2 and classes[0][1] == classes[1][1]:
        arb_result = _arbitrate(question, task_type, samples, classes)
        if arb_result:
            top, rep = arb_result
            # 2026-07-22: 裁定役（_arbitrate）は既存の票クラスと無関係な「第三の答え」を
            # 返すことがある（例: 拮抗した {'1','2'} に対し裁定役が '3' を新規提示）。
            # 従来はここで cnt/classes を再計算せず、拮抗していた旧トップの票数
            # （classes[0][1]）をそのまま流用していたため、
            #   - res['votes'] に裁定結果の答えが載らず、実際は0票の新答えなのに
            #     まるで敗者候補が「無投票」であるかのような矛盾した内訳になる
            #   - 直後の「[SC] 確定: {top} (票 {cnt}/...)」ログが、敗者候補（旧トップ）の
            #     票数を裁定結果の答えの票数であるかのように誤表示する
            # という報告面のバグがあった。ここでは裁定後の top に対応する真の票数を
            # classes から同値判定で引き直し（一致するクラスが無ければ 0 票）、
            # votes 辞書にも裁定結果の答えを必ずキーとして載せる。
            # なお SC_MIN_VOTES の床判定（下）は rep is not None の間は素通りする既存の
            # 挙動のままで変更していない。
            match = next((c for c in classes if answers_equivalent(top, c[0])), None)
            cnt = match[1] if match else 0
            if match is None:
                classes = classes + [[top, cnt]]
            elif match[0] != top:
                # 2026-07-22 (iteration 10 の続き): 上の一致判定は answers_equivalent な
                # クラスを見つけられるが、そのクラスの代表表記（match[0]）が裁定役の
                # 返した文字列（top）と食い違うことがある（例: 拮抗クラスの代表が '1/2'
                # で裁定役は '0.5' と書く、'1000' vs '1,000'、'012' vs '12' など、
                # 分数⇄小数や桁区切りの書き直しは裁定役がよくやる）。
                # 旧コードは「match is None or match[0] != top」を一括りにして
                # 新規クラス [top, cnt] を無条件追加していたため、同値のはずの票が
                # 旧代表表記のキーと裁定後表記のキーの二つに分裂して計上され
                # （例: res['votes'] == {'1/2': 3, '0.5': 3}）、合計票数が実際の
                # 有効票数の2倍になる「truthful でない votes」を生んでいた。これは
                # iteration 10 が退治したはずの矛盾内訳バグの兄弟ケースにあたる。
                # ここでは新規クラスを追加せず、一致したクラス自身のキーを裁定後の
                # 表記に書き換えるだけにして、同じ票が二重に数えられないようにする。
                classes = [([top, cnt] if c is match else c) for c in classes]
    # 2026-07-21: ループ内の早期確定条件（cnt==n and n>=SC_MIN_VOTES / n>=4 and cnt*2>n）は
    # SC_MIN_VOTES 未満の疑似全会一致を弾くが、それは while ループの break 条件だけの話。
    # SC_MAX 消化で抜けた場合（多くのサンプルが __ERROR__/抽出失敗/\boxed{}欠落 等で無効票になった
    # ケース）はここを素通りしてしまい、1〜2票しか残っていない「勝者」をそのまま確定扱いで返して
    # いた。理由は違えど中身は同じ疑似全会一致問題なので、最終returnにも同じ床（floor）をかける。
    # ただし裁定（_arbitrate）が成功して rep が既に埋まっている場合は、裁定役が新たに出した
    # answer/text をそのまま尊重し、票数に関わらずここでは弾かない。
    if rep is None and cnt < SC_MIN_VOTES:
        print(f"   [SC] 確定票が {cnt} 票のみ (< SC_MIN_VOTES={SC_MIN_VOTES}) → MoA フォールバックへ")
        return None
    if rep is None:
        rep = _representative_text(samples, top)
    print(f"   [SC] 確定: {top}  (票 {cnt}/{len(answers)}, サンプル計 {len(samples)})")
    return {"answer": top, "text": rep,
            "votes": {c[0]: c[1] for c in classes}, "n_samples": len(samples)}


# ==================================================
# Fugu 風オーケストレーション本体
# ==================================================


def _print_plan(plan):
    tag = " (フォールバック)" if plan.get("_fallback") else ""
    print(f"\n🎼 Conductor の判断{tag}:")
    if plan.get("make_pptx"):
        print("   output      = PowerPoint (画像は内容連動で自動生成)")
    if plan.get("use_image_generation"):
        kind = "画像のみ" if plan.get("image_only") else "テキスト+イラスト"
        print(f"   image_gen   = True ({kind})")
    if plan.get("image_only"):
        # テキスト提案は走らないので mode 表示は省略
        pass
    elif plan["mode"] == "single":
        sel = plan.get("selected_proposers") or PROPOSERS[:1]
        model = sel[0] if sel else AGGREGATOR
        print(f"   mode        = single ({_persona_str(model)})")
    else:
        labels = [_persona_str(m) for m in plan.get("selected_proposers", [])]
        print(f"   mode        = moa {labels}  rounds={plan['rounds']}")
    print(f"   search_req  = {plan.get('search_required', False)}")
    if plan.get("task_type"):
        print(f"   task_type   = {plan['task_type']}")
    if plan.get("reason"):
        print(f"   reason      = {plan['reason']}")


def fugu_answer(question, plan=None, history=None):
    """事前に conduct() で得た plan に従って回答を生成する。
    plan は validate_plan 済み（selected_proposers は実モデル名で解決済み）。
    plan=None のときは内部で conduct() を実行する（eval など単体呼び出し向けの後方互換）。"""
    history = history or []
    if plan is None:
        plan, _raw = conduct(question, history=history)
        if SHOW_PLAN:
            _print_plan(plan)

    # ---------- 検証可能タスク（math/mcq）: 自己一貫性投票で解く ----------
    if (SC_ENABLED and plan.get("task_type") in ("math", "mcq")
            and not plan.get("image_only") and not plan.get("make_pptx")
            and not plan.get("use_image_generation")):
        print(f"   [SC] 検証可能タスク({plan['task_type']}) → 自己一貫性投票で解く")
        res = solve_verifiable(question, plan["task_type"], history=history)
        if res and res.get("answer"):
            txt = res.get("text") or ""
            # 裁定で答えが差し替わった場合など、本文の結論と投票結果がずれたら明示する
            body_ans = extract_final_answer(txt, plan["task_type"])
            if not (body_ans and answers_equivalent(body_ans, res["answer"])):
                # 2026-07-27: bench_fugu の主要config 'fugu'（run_fugu）は grade_item に
                # answer_value=None を渡し、fugu_answer の戻り値「テキスト」から
                # extract_final_answer で答えを再抽出して採点する（res["answer"] は
                # 直接見ない）。これまでの素のプレーンテキスト注記
                # 「(自己一貫性投票による最終解答: X)」は「答え/正解/answer」宣言の
                # 連結詞パターンにも mcq の単独行パターンにも一致せず無視される一方、
                # extract_boxed は「末尾に一番近い、閉じている \boxed{}」を採用する
                # last-balanced-box-wins（本関数の上、iteration 12/25 のコメント参照 —
                # この修正はその順序に依存するカップリングである）。そのため本文側に
                # 残った古い \boxed{16} 等の方が拾われ続け、PoT専業の代表文・
                # 空/抽出不能な代表文・\boxed{16} vs 勝者 'B' のような value-boxed MCQ
                # 代表文で再抽出結果が投票結果とずれていた（stuck iteration 207 が
                # per-sample で個別に救おうとしていた MCQ value-box 化けの、デッキ
                # レベルでの根治でもある）。注記自体を \boxed{} で包み、last-wins の
                # 並びを利用して注記側を「最後の、閉じた box」にすることで、再抽出が
                # 必ず投票結果と一致するようにする（iteration 19 がこの分岐のテスト
                # カバレッジ、gotcha #7: 自己一貫性投票は精度優先・時間は気にしない の
                # 中核パスであり、ここでの再抽出ずれはその投票結果を静かに握りつぶす）。
                # ただし答えの値自体が波括弧の対応が崩れた病的な文字列（例 "}{"）だと
                # \boxed{} 化した注記そのものが壊れうるため、クラッシュさせず・注記を
                # 消してもしまわないよう、往復確認（extract_final_answer で読み戻して
                # answers_equivalent か）に失敗した場合だけ従来の素のプレーンテキスト
                # 注記へフォールバックする。
                boxed_note = f"\n\n(自己一貫性投票による最終解答: \\boxed{{{res['answer']}}})"
                plain_note = f"\n\n(自己一貫性投票による最終解答: {res['answer']})"
                try:
                    roundtrip = extract_final_answer(txt + boxed_note, plan["task_type"])
                    ok = bool(roundtrip) and answers_equivalent(roundtrip, res["answer"])
                except Exception:
                    ok = False
                txt += boxed_note if ok else plain_note
            return txt
        print("   [SC] 投票不成立 → 通常の合議へフォールバック")

    seed_answer = None  # エスカレーション時、単体回答を捨てずに合議の初期ドラフトにする
    seed_issue = None   # Critic の指摘も合議側へ伝える

    # ---------- 単体モード ----------
    if plan["mode"] == "single":
        sel = plan.get("selected_proposers") or PROPOSERS[:1]
        model = sel[0] if sel else (PROPOSERS[0] if PROPOSERS else AGGREGATOR)
        ans = strip_think(ask(
            model,
            ([{"role": "system", "content": proposer_sys_for(model) + PRESENTATION_STYLE}]
             + history
             + [{"role": "user", "content": question}]),
            PROPOSER_TEMP,
            # 2026-07-23: get_single_proposal と同じ欠落パターン（gotcha 該当箇所として明記）。
            # 単体モードの ask も think=PROPOSER_THINK の生グローバルを直渡ししており、隣の
            # num_predict=proposer_predict_for(model) だけがモデル別設定に対応していた非対称を
            # 解消。proposer_think_for(model) で _sc_sample と同じ解決順（PROPOSER_THINK override
            # > MODEL_CONFIG > モデル既定）に統一する。
            think=proposer_think_for(model),
            num_predict=proposer_predict_for(model),
            label="single",
        ))
        if ans.startswith("__ERROR__"):
            print("   (単体モデル失敗 → 合議へ切替)")
            plan["mode"] = "moa"
            plan["selected_proposers"] = PROPOSERS[:3]
        elif ADAPTIVE_ESCALATION:
            # 高速チェック2系統 + 疑義があれば思考ON再検算（verify_single 参照）
            ok, issue = verify_single(question, ans)
            if ok:
                return ans
            print(f"   ⤴ 単体回答に難あり（{issue}）→ 合議へエスカレーション")
            seed_answer, seed_issue = ans, issue
            plan["mode"] = "moa"
            plan["selected_proposers"] = PROPOSERS[:3]
            plan["rounds"] = max(1, plan["rounds"])
        else:
            return ans

    # ---------- 合議(MoA)モード：選抜した分だけ、必要なら再帰的に反復 ----------
    models = plan["selected_proposers"] or PROPOSERS[:3]
    planned = min(MAX_ROUNDS, max(1, plan["rounds"]))
    reference = seed_answer  # エスカレーションなら単体回答を初期ドラフトとして再利用
    issue_hint = seed_issue
    final = None
    r = 0
    while True:
        proposals = get_proposals(models, question, reference, issue_hint, history=history)
        if SHOW_PROPOSALS:
            mode = "並列" if PARALLEL_PROPOSERS else "逐次"
            print(f"\n--- ラウンド {r + 1}: 各提案（{mode}・{len(models)}体） ---")
            for m, a in proposals:
                print(f"[{m}]\n{strip_think(a)}\n")
        final = aggregate(question, proposals)
        # 2026-07-24: reference には think を持ち越さない。aggregate の生出力を
        # そのまま次ラウンドの reference にすると、get_single_proposal が
        # 'A draft answer from the panel:\n{reference}' として <think>...</think>
        # の内部思考をプロポーザーへ「ドラフト回答」として提示してしまい改善を誤誘導
        # する上、8192/16384 に固定した num_ctx（gotcha #2）を think ブロックが圧迫し
        # 本来の質問文/ドラフトが切り詰められかねない（精度優先＝gotcha #7 に反する）。
        # aggregate の結果は L2321、ask_fugu の最終出力は L3196 で同様に strip_think
        # 済みにしている慣例に合わせ、ここでも strip_think 済みの fin を reference に
        # 使う。返り値の final 自体は生のまま返す（ask_fugu 側で最終的に strip_think）。
        fin = strip_think(final)
        reference = fin  # 次ラウンドは今回の統合結果(think除去済み)を土台に改善
        issue_hint = None  # 指摘は消費済み。以降のチェックが新しい指摘を設定する
        r += 1

        # コード回答は実行検証で誤りが機械的に見つかるため、修正ラウンドの上限を広げる
        limit = (MAX_ROUNDS_CODE if (CODE_EXECUTION and extract_code(fin))
                 else MAX_ROUNDS)
        if r >= limit:
            break

        # 全プロポーザーが失敗した場合は即打ち切り。
        # これ以上ラウンドを重ねても同じ失敗が繰り返されるだけで、
        # Critic が「エラーメッセージは不十分」と正しく判定するため
        # MAX_ROUNDS 分だけ無駄なループが発生する（実測: x8 proposer + x12 critic 呼び出し）。
        if fin.startswith("__ERROR__"):
            break

        # 続行判断1（決定的・最優先）: コードを実行し、失敗なら traceback を修正ヒントに
        # して追加ラウンド。これが「自律的にコードを直し続ける」ループの本体。
        code_issue = code_check(fin)
        if code_issue:
            issue_hint = code_issue
            tail = code_issue.strip().splitlines()[-1][:80]
            print(f"   ↻ コード実行に失敗 → 修正ラウンド {r + 1}（{tail}）")
            continue

        # 続行判断2: 計画分がまだ残っていれば続行。消化済みなら Critic に委ねる（再帰）。
        if r < planned:
            need_more = True
        elif ALLOW_RECURSION:
            ok, issue = critique(question, fin)
            need_more = not ok
            if need_more:
                issue_hint = issue  # 何が不十分かを次ラウンドの提案へ伝える
                print(f"   ↻ 品質不足のため追加ラウンド（{issue}）")
            else:
                print("   ✓ 十分な品質と判断 → 反復を打ち切り")
        else:
            need_more = False
        if not need_more:
            break

    return final

# ==================================================
# 実行制御
# ==================================================

_READY = False


def setup():
    global PROPOSERS, AGGREGATOR, CONDUCTOR, _READY
    if _READY:
        return True
    if not ensure_server():
        return False
    if os.environ.get("FUGU_HIGH_VRAM") in ("1", "true", "True"):
        apply_high_vram_profile()
    print("[setup] ローカルモデル構成を確認します…")
    PROPOSERS, AGGREGATOR, CONDUCTOR = resolve_models()
    if not PROPOSERS or AGGREGATOR is None or CONDUCTOR is None:
        print("利用可能なモデルを用意できませんでした。")
        return False
    persona_lines = "\n".join(
        f"    {label} = {model}"
        + ("" if model in PROPOSERS else "  [未導入]")
        for label, model in PERSONA_MODELS.items()
    )
    print(f"""
===================================================
 🐡 Local Fugu-style MoA Orchestrator (3大AIオールスター)
  conductor : {CONDUCTOR}   (動的に委譲を決定)
  proposers :
{persona_lines}
  aggregator: {AGGREGATOR}
  image_gen : backend={IMAGE_BACKEND}  (a1111={A1111_URL} / comfyui={COMFYUI_URL})
  max_rounds: {MAX_ROUNDS}  escalation: {ADAPTIVE_ESCALATION}  recursion: {ALLOW_RECURSION}
  mode      : {"並列" if PARALLEL_PROPOSERS else "逐次"}
===================================================
 OLLAMA_MAX_LOADED_MODELS=1 は恒久設定済み（ユーザー環境変数）。
""")
    _READY = True
    return True


def ask_fugu(question, baseline=SHOW_BASELINE, *,
             use_search=False, rag_dirs=None, out_file=None,
             history_file=None, office_attached=False):
    """質問を Fugu パイプラインで処理する。
    use_search: True なら Web 検索を行いコンテキストに注入する（Conductor が
      search_required=true を出した場合も自動で有効化される）。
    rag_dirs: ローカル文書ディレクトリ（省略時は RAG_DIRS グローバル設定を使用）。
    out_file: 回答を保存するファイルパス（.md 推奨）。
    history_file: 永続化に使う JSON ファイルパス（省略時は HISTORY_FILE）。
    office_attached: Office 文書が添付されている旨を Conductor へ伝えるヒント。
    """
    global _HISTORY
    if not setup():
        return None
    t0 = time.time()

    # --- Conductor プランを先に取得（検索要否・画像生成・Office ルーティングを決める）---
    print("\n[Fugu] Conductor がオーケストレーションを開始します...")
    plan, _raw = conduct(question, history=list(_HISTORY),
                         office_attached=office_attached)
    if SHOW_PLAN:
        _print_plan(plan)

    panel = plan.get("selected_proposers") or PROPOSERS[:IMAGE_PROMPT_PANEL]

    # --- 経路1: 画像のみ（テキスト回答不要・LLM群がプロンプト起草）---
    if plan.get("use_image_generation") and plan.get("image_only"):
        print("\n[Fugu] 画像のみ生成（LLM群がプロンプトを起草）...")
        result = handle_image_generation(question, panel=panel)
        elapsed = round(time.time() - t0, 1)
        print("\n===== 画像生成結果 =====")
        # 2026-07-24: handle_image_generation が失敗時に返す内部センチネル
        # '__ERROR__: ...' をそのまま print(result) していたため、コンソールに
        # 機械向けマーカーが生のまま漏出していた。aggregate()（iteration 9）、
        # _critic_judge/second_opinion（iteration 15）、_arbitrate（iteration 20）、
        # および経路3のイラスト付き回答（iteration 99）で対処した「内部センチネル
        # をユーザ向け出力に漏らさない」バグと同種。ここではコンソール表示だけを
        # 人間可読な文言に置き換える。notify_slack への通知（iteration 119 の
        # 失敗アイコン判定用）・_save_answer_to_file のゲート（iteration 80 の
        # エラー時未保存）・関数の戻り値は、従来どおり生の result
        # （'__ERROR__' 始まりで失敗を示す）のまま変更しない。
        if result.startswith("__ERROR__"):
            note = result[len("__ERROR__"):].lstrip(":").strip()
            print(f"画像生成に失敗しました: {note}" if note else "画像生成に失敗しました")
        else:
            print(result)
        print(f"\n(所要 {elapsed} 秒)")
        notify_slack(question, result, elapsed)
        if out_file and not result.startswith("__ERROR__"):
            _save_answer_to_file(question, result, elapsed, out_file, context="")
        return result

    # --- コンテキスト構築（Web検索 + RAG）。検索は CLI フラグ or Conductor 判断で有効化 ---
    do_search = use_search or plan.get("search_required", False)
    context = build_context(question, use_search=do_search,
                            rag_dirs=rag_dirs or RAG_DIRS)
    question_with_ctx = _with_context(question, context)

    if baseline:
        print("\n===== 単体ベースライン（aggregator モデル直答） =====")
        base = ask(
            AGGREGATOR,
            [{"role": "system", "content": PROPOSER_SYS},
             {"role": "user", "content": question_with_ctx}],
            AGGREGATOR_TEMP,
            label="baseline",
        )
        print(strip_think(base))

    # --- 本文を MoA で生成 ---
    final = strip_think(
        fugu_answer(question_with_ctx, plan, history=list(_HISTORY)) or ""
    )
    text_answer = final  # 履歴にはテキスト本文のみ保存する

    # --- 経路2: PowerPoint（本文をスライド化し内容連動で画像を埋め込む）---
    if plan.get("make_pptx") and not final.startswith("__ERROR__"):
        print("\n[Fugu] 本文を PowerPoint 化します（画像は内容連動で自動生成）...")
        pptx_out = out_file if (out_file and str(out_file).lower().endswith(
            (".pptx", ".ppt"))) else None
        deck = build_pptx(question, final, pptx_out)
        final = text_answer + f"\n\n---\n## 生成した PowerPoint\n- 保存先: {deck}"
        # 2026-07-26: out_file を無条件に None化すると、--out が .pptx/.ppt 以外
        # （例: notes.md）の場合に pptx_out は None のまま渡され build_pptx は
        # 既定の PPTX_OUT_DIR に保存する一方、ユーザーが明示指定した --out 先への
        # 保存は下の汎用 _save_answer_to_file がスキップされて silently 消えて
        # いた（iteration 186 で表面化した repl 経路の --out 取りこぼしと同種）。
        # ここは pptx_out が実際に消費された（＝out_file が .pptx/.ppt だった）
        # 場合のみ None化し、それ以外は out_file を残して下の汎用保存に委ね、
        # 高コストな MoA 回答を保存ステップで失わない（iteration 41-47/80と同種の原則）。
        if pptx_out is not None:
            out_file = None  # 保存はここで完結（下の汎用保存は行わない）

    # --- 経路3: イラスト付き回答（本文＋回答内容から画像生成）---
    elif (plan.get("use_image_generation") and not plan.get("image_only")
          and not final.startswith("__ERROR__")):
        print("\n[Fugu] 回答内容からイラストを生成します...")
        base = f"{question}\n\n[回答の要点]\n{text_answer[:800]}"
        img = handle_image_generation(base, panel=panel)
        # 2026-07-24: handle_image_generation が失敗時に返す内部センチネル
        # '__ERROR__: ...' をそのまま final に連結すると、final は
        # text_answer から始まるため直後の全ての `final.startswith("__ERROR__")`
        # 判定（コンソール表示・notify_slack・履歴保存・_save_answer_to_file）を
        # すり抜けてしまい、内部向けマーカーがユーザ向け回答・Slack通知・保存
        # ファイルにそのまま漏出していた。aggregate()（iteration 9）、
        # _critic_judge/second_opinion（iteration 15）、_arbitrate（iteration 20）
        # で対処した「内部センチネル/タグをユーザ向け出力に漏らさない」バグと
        # 同種。ここでは img がセンチネルなら人間可読な失敗ノートに置き換える
        # （プレフィックスを剥がすだけで、テキスト本文自体は失敗として扱わない）。
        # iteration 73 で確立した「_HISTORY にはクリーンな text_answer のみを
        # 積む」分離は変更しない（下の履歴追記は従来通り text_answer のまま）。
        if img.startswith("__ERROR__"):
            note = img[len("__ERROR__"):].lstrip(":").strip()
            img = f"(画像生成に失敗しました: {note})" if note else "(画像生成に失敗しました)"
        final = text_answer + "\n\n---\n## 生成画像\n" + img

    elapsed = round(time.time() - t0, 1)
    print("\n===== 最終回答 =====")
    if final.startswith("__ERROR__"):
        print("生成に失敗しました:", final)
    else:
        print(final)
    print(f"\n(所要 {elapsed} 秒)")
    notify_slack(question, final, elapsed)

    # --- 会話履歴を更新（エラーでなければ記録・永続化）---
    if not final.startswith("__ERROR__"):
        _HISTORY.append({"role": "user", "content": question})   # 元の質問を保存
        # 2026-07-23: L3098 のコメント「履歴にはテキスト本文のみ保存する」の
        # 意図と実装が乖離していたのを修正。旧実装は final（PPTX 経路では
        # '## 生成した PowerPoint / 保存先: <deck>'、イラスト経路では
        # '## 生成画像 / <img markup-or-status>' が本文に追記された成果物付き
        # 文字列）をそのまま履歴に積んでおり、次ターンの Conductor/proposers が
        # ファイルパスや画像生成ステータスを「前回回答の実質的内容」として誤読
        # し得る状態だった（iteration 59/67/70 で対処した複数ターン間の忠実性
        # 劣化と同種の問題）。加えて MAX_HISTORY_CHARS の予算も無駄に消費する。
        # 修正: 履歴には text_answer（成果物注記より前のクリーンな本文）を積む。
        # 画像・PPTX 経路を使わない通常パスでは text_answer == final のため
        # 挙動は変わらない。戻り値・コンソール出力・notify_slack・
        # _save_answer_to_file は従来通り final（成果物付き）を使い続ける。
        _HISTORY.append({"role": "assistant", "content": text_answer})
        _trim_history(_HISTORY)
        save_history_file(_HISTORY, path=history_file)
        print(f"   [会話履歴: {len(_HISTORY) // 2} 往復保持中]")

    # --- ファイル出力 ---
    if out_file and not final.startswith("__ERROR__"):
        _save_answer_to_file(question, final, elapsed, out_file,
                             context=context)

    return final


# コード系拡張子（コードブロックを抽出してそのまま書き出す）
_CODE_EXTENSIONS = {
    ".py", ".js", ".ts", ".jsx", ".tsx", ".mjs", ".go", ".rs",
    ".c", ".cpp", ".h", ".hpp", ".cs", ".java", ".kt", ".swift",
    ".rb", ".php", ".sh", ".bat", ".ps1", ".sql", ".r", ".m", ".jl",
}


def _extract_code_for_output(answer: str, suffix: str) -> str:
    """回答から対象言語のコードブロックを抽出して返す。
    見つからない場合は回答全体からマークダウン装飾を除いたテキストを返す。

    2026-07-22: iteration-7 の extract_code と同じ誤抽出クラスの修正。
    旧実装は言語指定ありフェンスを re.search(rf"```{lang}[ \t]*\n(.*?)```") で
    検索し（```python3 の '3' が [ \t]*\n にマッチせず python3 ブロックを取り
    こぼす）、言語指定なしフォールバックも re.search(r"```(?:\w+)?[ \t]*\n(.*?)```")
    による前方走査だったため、```json/```text/```output 等の非コードブロックが
    先行すると、その閉じフェンスを誤って開始フェンスとみなし、2ブロック間の
    プロースや先行ブロックの中身をコードとして誤抽出していた。
    修正: extract_code と同一のフェンス正規表現 re.finditer(r"```([^\n`]*)\n(.*?)```")
    で全ブロックを一度に収集し、(1) suffix の対象言語タグ一致 → (2) タグ無し
    (bare) → (3) 既知の非コードタグ以外 の優先順で最初に見つかったブロックの
    本文を返す。該当ブロックが無ければ従来通りフェンス無しフォールバックを返す。

    2026-07-22: iteration 28 で extract_code (L1192, iteration 18 由来の関数) に
    施した info-string 修正をここにも追随。CommonMark の info string 仕様では
    フェンス開始行の残り全体ではなく「最初の空白区切りトークン」だけが言語タグで、
    それ以降は任意のメタデータ（例: ```python title="sol.py" や
    ```python {.line-numbers}）。旧実装は strip した info string 全体を lang として
    比較していたため、(1) 装飾付き python フェンスが lang=='python title="sol.py"'
    などとなって langs との一致 (1) に失敗し、tier-1 の優先扱いを受けられず tier-3
    （非コードタグ以外なら何でも）に落ちてしまい、(2) 装飾付き非コードフェンス
    （例: ```json {.line-numbers}）が lang=='json {.line-numbers}' となって
    _NON_CODE_TAGS の "json" と不一致になり、tier-3 でそのまま採用され、後続の
    本物の python ブロックより先に JSON/プロースが _save_as_code で .py/.js に
    書き出されてしまっていた。修正: info string の最初のトークンのみを言語タグ
    として全3段の判定に使う（extract_code と同じ考え方）。受理集合・優先順位・
    _NON_CODE_TAGS 自体は変更しない。"""
    lang_map = {
        ".py": ["python", "py", "python3"],
        ".js": ["javascript", "js"],
        ".ts": ["typescript", "ts"],
        ".go": ["go"],
        ".rs": ["rust"],
        ".c": ["c"],
        ".cpp": ["cpp", "c++"],
        ".cs": ["csharp", "cs"],
        ".java": ["java"],
        ".rb": ["ruby", "rb"],
        ".sh": ["bash", "sh", "shell"],
        ".sql": ["sql"],
        ".r": ["r"],
        # 2026-07-23: _CODE_EXTENSIONS (L3100, --out がコードとして書き出す25拡張子)
        # と lang_map の同期漏れ。上の13エントリしか無く、残り12拡張子
        # (.jsx .tsx .mjs .h .hpp .kt .swift .php .bat .ps1 .m .jl) は
        # langs が空集合になるため tier-1 (対象言語タグ一致) が絶対に発火せず、
        # 複数フェンスの回答で非対象言語のブロック(例: 使い方説明の```bash)が
        # 先行すると tier-3 (非コードタグ以外の最初のブロック) でそれが誤って
        # 採用され、意図した言語のブロックより先に書き出されていた
        # (iteration 7/18/28/29/56 が繰り返し修正してきたのと同じ
        # ブロック誤選択バグクラスだが、このタグ集合の穴自体は今回まで未対応で、
        # 既存の code_out: テストは全て .py/.c という対応済み拡張子だけを使って
        # いたため検出されずにいた)。tier-1 は「優先されるブロックを選ぶだけ」
        # で受理判定を誤って拡げても非対象ブロックを誤採用する副作用は無いため、
        # .h (C/C++/Objective-C 共用ヘッダ) と .m (MATLAB/Objective-C/Octave) の
        # ように言語が一意に決まらない拡張子は、単一言語を憶測するのではなく
        # 妥当なタグを広めに全部含める。
        ".jsx": ["jsx", "javascript", "js"],
        ".tsx": ["tsx", "typescript", "ts"],
        ".mjs": ["javascript", "js", "mjs"],
        ".h": ["c", "cpp", "c++", "objc", "objective-c", "objectivec"],
        ".hpp": ["cpp", "c++", "hpp"],
        ".kt": ["kotlin", "kt"],
        ".swift": ["swift"],
        ".php": ["php"],
        ".bat": ["bat", "batch", "cmd"],
        ".ps1": ["powershell", "ps1", "pwsh"],
        ".m": ["matlab", "objc", "objective-c", "objectivec", "octave"],
        ".jl": ["julia"],
    }
    # 非コードとみなす既知のドキュメント系タグ（保守的なスキップリスト）
    _NON_CODE_TAGS = {
        "json", "text", "txt", "output", "console", "log",
        "yaml", "yml", "xml", "csv", "markdown", "md", "diff", "ini", "toml",
    }
    langs = {l.lower() for l in lang_map.get(suffix, [])}

    def _fence_lang(info: str) -> str:
        info = info.strip().lower()
        return info.split(None, 1)[0] if info else ""

    blocks = [(_fence_lang(m.group(1)), m.group(2))
              for m in re.finditer(r"```([^\n`]*)\n(.*?)```", answer, re.DOTALL)]

    # (1) 対象言語タグと一致する最初のブロック
    for lang, body in blocks:
        if lang in langs:
            return body
    # (2) タグ無し(bare)の最初のブロック
    for lang, body in blocks:
        if lang == "":
            return body
    # (3) 既知の非コードタグ以外の最初のブロック
    for lang, body in blocks:
        if lang not in _NON_CODE_TAGS:
            return body
    # フェンスなし: マークダウン見出し行を除いた本文を返す
    # 2026-07-23: 旧実装は l.startswith("#") で '#' から始まる行を無条件に
    # 全削除しており、docstring が意図する「マークダウン見出し(ATX heading)の
    # 除去」の範囲を大きく超えていた。この関数はコード拡張子ファイルへの
    # 保存(--out file.<ext>)にも使われるため、フェンス無しの生コードに含まれる
    # C の #include/#define/#pragma、シェバン行 #!/usr/bin/env python、
    # Rust の属性 #[derive(Debug)] まで「見出し」として問答無用に削除され、
    # コンパイル不能な壊れたファイルが書き出されていた（iteration 7/18/28/29 が
    # 繰り返し修正してきたのはこの上のフェンス選択側で、この素通しの兄弟分岐は
    # 今回まで手つかずだった）。CommonMark の ATX heading 仕様では、見出しは
    # 行頭の '#' が1〜6個続いた直後が空白または行末である行に限られる
    # （'#include' のように '#' の直後に空白なしで文字が続く行は見出しではない）。
    # そこで削除対象を正規表現 ^#{1,6}(?:\s|$) に限定し、真の見出し行
    # (# Title, ## Section 等)だけを除去して、上記のような '#' 始まりの
    # 非見出しコード行は保持する。列0アンカー(インデント行を見出し扱いしない)は
    # 従来通り維持する。
    _ATX_HEADING_RE = re.compile(r"^#{1,6}(?:\s|$)")
    lines = [l for l in answer.splitlines() if not _ATX_HEADING_RE.match(l)]
    return "\n".join(lines).strip()


def _save_as_markdown(out: Path, question: str, answer: str,
                      elapsed: float, context: str):
    """Markdown 形式で追記保存。"""
    from datetime import datetime
    ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    block = f"## Q ({ts})\n\n{question}\n\n"
    if context:
        block += (f"<details><summary>Context (search/RAG)</summary>\n\n"
                  f"{context}\n\n</details>\n\n")
    block += f"## A\n\n{answer}\n\n*所要: {elapsed}s*\n\n---\n\n"
    # 2026-07-23: --out が既存ファイルを指す追記保存で、その既存ファイルが
    # cp932/Shift_JIS 等の非UTF-8バイト列を含む場合（このマシンのコンソールが
    # cp932 であることに起因する既知の落とし穴 #4 と同種の環境要因）、
    # encoding="utf-8" のみの read_text は UnicodeDecodeError を送出し、
    # _save_answer_to_file 経由で保存ステップ全体が異常終了して、せっかく
    # 計算し終えた回答が失われてしまう。これは iteration 41-44 で
    # _save_as_excel/_docx/_pdf に対して行った「保存を絶対にクラッシュさせず
    # 回答を失わない」という degrade-gracefully 修正と同じバグクラスであり、
    # errors="replace" を付けることで読めない既存バイトだけを置換文字に
    # 落とし、正常な（UTF-8な）既存内容とこれから書く新規回答は従来通り
    # そのまま保持する。書き込み側の encoding="utf-8" はそのまま維持。
    existing = out.read_text(encoding="utf-8", errors="replace") if out.exists() else ""
    out.write_text(existing + block, encoding="utf-8")


def _save_as_text(out: Path, question: str, answer: str, elapsed: float):
    """プレーンテキスト形式で追記保存。"""
    from datetime import datetime
    ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    block = f"[{ts}]\nQ: {question}\n\nA:\n{answer}\n\n(所要 {elapsed}s)\n{'='*60}\n\n"
    # 2026-07-23: _save_as_markdown と同じ理由（落とし穴 #4 のcp932環境で
    # 既存 --out ファイルが非UTF-8バイトを含むと read_text がクラッシュし、
    # iteration 41-44 の office savers 同様に回答保存が丸ごと失われる）で
    # errors="replace" を追加。書き込みは encoding="utf-8" のまま不変。
    existing = out.read_text(encoding="utf-8", errors="replace") if out.exists() else ""
    out.write_text(existing + block, encoding="utf-8")


def _save_as_code(out: Path, answer: str):
    """コード拡張子のファイルとして保存。コードブロックを抽出して書き出す。"""
    code = _extract_code_for_output(answer, out.suffix.lower())
    out.write_text(code + "\n", encoding="utf-8")


def _save_as_html(out: Path, question: str, answer: str, elapsed: float):
    """HTML 形式で保存。"""
    from datetime import datetime
    import html as _html
    ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    q_esc = _html.escape(question)
    a_lines = []
    # 2026-07-22: 旧実装は開始・終了の ``` フェンスを両方とも "<pre><code>" に
    # マップしており "</code></pre>" を一度も出力しないため、コードブロックを
    # 含む回答が "<pre><code>...<pre><code>" という入れ子・未クローズの不整合
    # HTML になっていた。さらにコード本文の行は else 節に落ちて末尾に "<br>"
    # が付与され、整形済みのはずのコードが崩れていた。開始/終了フェンスを
    # in_code フラグで区別し、コード本文は <br> を付けず改行 "\n" を保持した
    # まま escape して <pre><code>...</code></pre> の中に入れる。
    in_code = False
    for line in answer.splitlines():
        # 2026-07-23: フェンス判定が line.startswith("```") のまま列0固定
        # だったため、LLM が番号付き/箇条書きリストの中に```pythonブロックを
        # 2〜4スペースでインデントして出力する（よくあるケース）と、行頭に
        # 空白があるだけでフェンスとして認識されず in_code に入れなかった。
        # 結果、```python/``` の行自体がプレーンテキストとして escape され
        # そのまま文字列として表示され、コード本文の各行も else 節に落ちて
        # 余計な <br> が付与された状態で崩れて出力されていた。これは
        # extract_boxed（iteration 11）・strip_think（iteration 16）・
        # _save_as_html のタグ整合性そのもの（iteration 37）・_parse_slides
        # （iteration 50）で既に対処済みのフェンス未検出/不整合と同種の
        # バグクラスで、本関数だけインデント方向が未対応だった。他の全ての
        # フェンス処理（extract_code, _extract_code_for_output, _parse_slides）
        # と同じ line.strip().startswith("```") に揃えることで、インデント
        # されたフェンスも列0のフェンスと同様に単一の <pre><code>…
        # </code></pre> へバランスよく変換されるようにする。
        if line.strip().startswith("```"):
            if in_code:
                a_lines.append("</code></pre>")
                in_code = False
            else:
                a_lines.append("<pre><code>")
                in_code = True
        elif in_code:
            a_lines.append(_html.escape(line) + "\n")
        else:
            a_lines.append(_html.escape(line) + "<br>")
    if in_code:
        # フェンスが奇数個（閉じ忘れ）の場合でも <pre>/<code> を確実に閉じる
        a_lines.append("</code></pre>")
    body = (f"<h2>Q <small>({ts})</small></h2>\n<p>{q_esc}</p>\n"
            f"<h2>A</h2>\n<p>{''.join(a_lines)}</p>\n"
            f"<hr><p><small>所要: {elapsed}s</small></p>\n")
    existing_body = ""
    if out.exists():
        # 2026-07-23: _save_as_markdown/_save_as_text と同じ理由（落とし穴 #4 の
        # cp932環境で既存 --out ファイルが非UTF-8バイトを含むケース）で
        # errors="replace" を追加。既存<body>のマージ読み戻しがクラッシュして
        # iteration 41-44 の office savers と同じバグクラスで回答保存が失われる
        # のを防ぐ。書き込みは encoding="utf-8" のまま不変。
        content = out.read_text(encoding="utf-8", errors="replace")
        m = re.search(r"(<body>)(.*?)(</body>)", content, re.DOTALL)
        if m:
            existing_body = m.group(2)
    html_content = (f"<!DOCTYPE html>\n<html lang='ja'><head>"
                    f"<meta charset='UTF-8'><title>Fugu Output</title></head>\n"
                    f"<body>\n{existing_body}{body}</body></html>")
    out.write_text(html_content, encoding="utf-8")


def _save_as_pdf(out: Path, question: str, answer: str, elapsed: float):
    """PDF 形式で保存。fpdf2 が必要。未インストール時は .md にフォールバック。"""
    try:
        from fpdf import FPDF
    except ImportError:
        md_path = out.with_suffix(".md")
        _save_as_markdown(md_path, question, answer, elapsed, "")
        print(f"   [PDF生成には fpdf2 が必要 (pip install fpdf2)。代わりに保存: {md_path}]")
        return md_path

    # 2026-07-22: fpdf2 には 'DejaVu' という名前で事前登録された組み込み
    # Unicode フォントは存在しない（add_font での明示登録が必要）。
    # set_font("DejaVu") は FPDFException を送出し、直下の except Exception
    # で Helvetica にフォールバックするが、Helvetica はコア latin-1 フォント
    # のため、既定言語である日本語などの非ASCII文字を multi_cell/cell に渡す
    # と FPDFUnicodeEncodingException（ImportError ではない ただの Exception）
    # が送出される。従来は except ImportError だけを捕捉していたため例外が
    # そのまま伝播し、_save_as_pdf の呼び出し元 _save_answer_to_file 自体に
    # ガードがなく、回答保存ステップ全体が異常終了して回答を失っていた
    # （iteration 41 の _save_as_excel の IllegalCharacterError 修正、
    # iteration 43 の _save_as_docx の制御文字 ValueError 修正と同じ
    # バグクラス）。ここでは Unicode フォントの登録・同梱は行わず（環境依存
    # のため別対応とする）、PDF 構築・出力段階の失敗を捕捉して既存の .md
    # フォールバックへ安全に降格させるに留める。
    try:
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        # Unicode フォント: fpdf2 は内蔵 DejaVu を使用
        try:
            pdf.set_font("DejaVu", size=12)
        except Exception:
            pdf.set_font("Helvetica", size=12)
        from datetime import datetime
        ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        for heading, text in [("Q", question), ("A", answer)]:
            pdf.set_font_size(14)
            pdf.cell(0, 10, f"{heading} ({ts})" if heading == "Q" else heading, ln=True)
            pdf.set_font_size(11)
            for line in text.splitlines():
                pdf.multi_cell(0, 7, line or " ")
            pdf.ln(5)
        pdf.output(str(out))
        return
    except Exception as e:
        md_path = out.with_suffix(".md")
        _save_as_markdown(md_path, question, answer, elapsed, "")
        print(f"   [PDF保存に失敗しました ({e!r})。代わりに保存: {md_path}]")
        return md_path


def _save_as_docx(out: Path, question: str, answer: str, elapsed: float):
    """Word (.docx) 形式で保存。python-docx が必要。未インストール時は .md にフォールバック。"""
    try:
        import docx as _docx
    except ImportError:
        md_path = out.with_suffix(".md")
        _save_as_markdown(md_path, question, answer, elapsed, "")
        print(f"   [DOCX保存には python-docx が必要 (pip install python-docx)。代わりに保存: {md_path}]")
        return md_path

    # 2026-07-22: LLM の回答（および question）には稀にフォームフィード(\x0c)、
    # ANSIエスケープ(\x1b)、NUL 等の制御文字が混入する。これらは XML 1.0
    # 仕様上不正な文字であり、python-docx (lxml) の add_paragraph/add_heading
    # に渡すと ValueError が送出される。従来は except ImportError だけを
    # 捕捉していたため例外がそのまま伝播し、_save_answer_to_file 全体が
    # 異常終了して回答保存に失敗していた（iteration 41 の _save_as_excel の
    # IllegalCharacterError 修正と同じバグクラス）。python-docx に渡す前に
    # 各文字列から XML 不正制御文字 (0x00-0x08, 0x0B, 0x0C, 0x0E-0x1F) だけを
    # 除去することで実データ（本文）は保ったまま正常な .docx を書き出す。
    # 万一それでも保存に失敗した場合は、既存の .md フォールバックへ安全に
    # 降格させ、決してここで異常終了させない。
    try:
        _illegal_xml_re = re.compile(r"[\x00-\x08\x0B\x0C\x0E-\x1F]")
        from datetime import datetime
        doc = _docx.Document()
        ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        doc.add_heading(f"Q ({ts})", level=1)
        doc.add_paragraph(_illegal_xml_re.sub("", question))
        doc.add_heading("A", level=1)
        in_code = False
        code_lines = []
        for line in answer.splitlines():
            line = _illegal_xml_re.sub("", line)
            if line.startswith("```"):
                if in_code:
                    doc.add_paragraph("\n".join(code_lines), style="No Spacing")
                    code_lines = []
                    in_code = False
                else:
                    in_code = True
            elif in_code:
                code_lines.append(line)
            else:
                doc.add_paragraph(line) if line.strip() else doc.add_paragraph("")
        doc.add_paragraph(f"所要: {elapsed}s")
        doc.save(str(out))
        return
    except Exception as e:
        md_path = out.with_suffix(".md")
        _save_as_markdown(md_path, question, answer, elapsed, "")
        print(f"   [DOCX保存に失敗しました ({e!r})。代わりに保存: {md_path}]")
        return md_path


def _save_as_excel(out: Path, answer: str):
    """回答中のCSVライクな表を Excel (.xlsx) として保存。openpyxl が必要。"""
    # 2026-07-22: LLM の回答には稀にフォームフィード(\x0c)、ANSIエスケープ
    # (\x1b)、NUL 等の制御文字が混入する。これらは XML 1.0 の仕様上
    # 不正な文字であり、ws.append() の時点で
    # openpyxl.utils.exceptions.IllegalCharacterError（ImportError ではない
    # ただの Exception）が送出され、従来は except ImportError だけを
    # 捕捉していたため例外がそのまま伝播し、_save_answer_to_file 全体が
    # 異常終了して回答保存に失敗していた。ws.append() する前に各セル文字列
    # から XML 不正制御文字 (0x00-0x08, 0x0B, 0x0C, 0x0E-0x1F) だけを除去
    # することで実データ（表の中身）は保ったまま正常な .xlsx を書き出す。
    # 万一それでも保存に失敗した場合は、既存の .csv フォールバックへ
    # 安全に降格させ、決してここで異常終了させない。
    try:
        import openpyxl
    except ImportError:
        txt_path = out.with_suffix(".csv")
        txt_path.write_text(answer, encoding="utf-8")
        print(f"   [Excel保存には openpyxl が必要 (pip install openpyxl)。代わりに保存: {txt_path}]")
        return txt_path

    try:
        _illegal_xml_re = re.compile(r"[\x00-\x08\x0B\x0C\x0E-\x1F]")
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Fugu Output"
        for line in answer.splitlines():
            if line.strip():
                _stripped_line = line.strip()
                # 2026-07-26: LLM の回答が Markdown 表（例: '| Name | Age |'）の場合、
                # 従来の re.split(r"[,\t|]", line) では行頭/行末の '|' がそれぞれ
                # 空文字列の列を生み（例: ['', 'Name', 'Age', '']）、区切り線の行
                # '| --- | --- |' がそのままゴミの1データ行として書き込まれ、さらに
                # セル内の桁区切りカンマ（例: '| 1,234 | total |'）まで誤って列区切り
                # とみなされ数値データが '1' と '234' に分断されてしまっていた。これは
                # iteration 41-47・68 の保存系ハードニング（IllegalCharacterError対策
                # 等）と同種の、表形式パース漏れによる保存内容の破損バグ。行の前後を
                # '|' で囲まれた Markdown 表の行だけは、外側の '|' を1つずつ剥がした
                # うえで中身を '|' のみで分割し（カンマ/タブでは分割しない）、
                # '---'/':---'/'---:'/':---:' のような区切り線の行は列として追加しない。
                # それ以外の行（Markdown 表でない行）は従来通り
                # re.split(r"[,\t|]", line) のままで挙動を一切変えない。
                if _stripped_line.startswith("|") and _stripped_line.endswith("|"):
                    inner = _stripped_line[1:-1]
                    cols = [_illegal_xml_re.sub("", c.strip()) for c in inner.split("|")]
                    if cols and all(re.match(r"^:?-+:?$", c) for c in cols):
                        continue
                    ws.append(cols)
                else:
                    cols = [_illegal_xml_re.sub("", c.strip()) for c in re.split(r"[,\t|]", line)]
                    ws.append(cols)
        wb.save(str(out))
        return
    except Exception as e:
        txt_path = out.with_suffix(".csv")
        txt_path.write_text(answer, encoding="utf-8")
        print(f"   [Excel保存に失敗しました ({e!r})。代わりに保存: {txt_path}]")
        return txt_path


# ==================================================
# PowerPoint 生成（画像入りスライド）
# ==================================================

def _parse_slides(answer):
    """Markdown 回答をスライド構造 [{'title':..,'bullets':[..]}] へ分解する。
    見出し(#〜####)でスライドを区切り、箇条書き/段落を bullets にする。"""
    slides = []
    cur = None
    in_code = False

    # 2026-07-23: ```フェンスの数が奇数（＝閉じ忘れ/開き忘れで打ち切られた
    # proposer 出力）だと、旧実装の単純なトグルでは in_code が文書の残り
    # 全体で True に固定されたままになる。すると以降の全ての '#'/'##' 見出し
    # が新規スライドを作れず（"if m and not in_code" のガードに阻まれ）、かつ
    # 強調記号除去（[*_`#]+ の除去）もスキップされるため、見出しが
    # "## 節タイトル" というリテラルな文字列のまま1枚のスライドの箇条書きに
    # 押し込まれ、複数節あるはずのデッキが1枚に崩壊してしまう。これは
    # フェンスの閉じ忘れ/不対応によって以降の処理状態が汚染され続ける同種の
    # バグで、extract_boxed（iteration 11）・strip_think（iteration 16）・
    # _save_as_html（iteration 37）で既に対処済みのバグクラスだが、
    # _parse_slides だけは未対応だった。ここでは事前に本文中の```で始まる
    # 行を数え、その総数が奇数＝最後の1個が対になっていない場合、その最後の
    # 1個だけをトグル対象から除外する（フェンスとして扱わずコード開始/終了
    # 処理をしない）。それより前の、正しく対になったコードブロックは従来
    # 通りコードモードを維持するので、そのブロック内の '# コメント' が見出し
    # に昇格することはない。フェンスが偶数個（＝全て正しく対応）の入力では
    # unpaired_idx が -1 のままとなり、挙動は一切変わらない。
    lines = answer.splitlines()
    fence_idxs = [i for i, ln in enumerate(lines) if ln.rstrip().strip().startswith("```")]
    unpaired_idx = fence_idxs[-1] if len(fence_idxs) % 2 == 1 else -1

    for i, ln in enumerate(lines):
        s = ln.rstrip()
        if s.strip().startswith("```"):
            if i != unpaired_idx:
                in_code = not in_code
            continue
        m = re.match(r"^\s*#{1,4}\s+(.*)$", s)
        if m and not in_code:
            if cur is not None:
                slides.append(cur)
            # 2026-07-23: 旧実装 re.sub(r"[*_`#]+", "", t) は行中の *, _, `, # を
            # 位置・対応関係を一切見ず無差別に全削除しており、Markdown装飾
            # （**太字**、`インラインコード`）と、数式の演算子/指数（m*c_0**2 の
            # * や **）・識別子中のアンダースコア（snake_case の a_1、do_thing()）・
            # 地の文のハッシュ（C#, #123 等）を区別できなかった。検証例:
            # "E = m*c_0**2" が "E = mc02" に、"`do_thing()`" が "dothing()" に
            # 化けていた。スライド本文はそのまま .pptx へ書き出される最終成果物
            # であり、装飾の見た目より内容の忠実性を優先すべき（精度優先・
            # 時間は気にしない）。ここでは対になった `...` インラインコード区間
            # だけをバッククォート込みで除去して中身（アンダースコア含む）は
            # そのまま残す精密な正規表現に限定し、対になっていない単独の
            # `/*/_/# はいずれも曖昧なので削除せず残す（安全な方向＝落とすより
            # 残す）。過度に広い正規表現を用途を満たす最小限のパターンへ精密化
            # する同種の修正は iteration 26（MCQ宣言の連結子）・28（数式宣言の
            # 桁区切り）・30（3桁区切りカンマ）・56（_extract_code_for_output の
            # ATX見出し限定除去）で既出のバグクラス。本関数自体は iteration 50
            # （奇数フェンスの未対応対策）で一度手を入れているが、この強調記号
            # 除去の過剰マッチは今回まで未対応だった。次の本文側の同一置換
            # （t = re.sub(...)）にも同じ理由で同じ精密化を適用する。
            cur = {"title": re.sub(r"`([^`\n]+)`", r"\1", m.group(1)).strip()[:80], "bullets": []}
            continue
        t = s.strip()
        if not t:
            continue
        if not in_code:
            t = re.sub(r"^\s*(?:[-*+]|\d+[.)])\s+", "", t)  # 箇条書き記号除去
            # 2026-07-23: タイトル側と同じ精密化（上のコメント参照）。対になった
            # `...` のみ除去し、*/_/# や対になっていないバッククォートは
            # 演算子/識別子/地の文の可能性があるため残す（フィデリティ優先）。
            t = re.sub(r"`([^`\n]+)`", r"\1", t).strip()    # 強調記号除去（対のバッククォートのみ）
        if not t:
            continue
        if cur is None:
            cur = {"title": "概要", "bullets": []}
        cur["bullets"].append(t[:200])
    if cur is not None:
        slides.append(cur)
    return slides


def _deck_title(question, slides):
    # 2026-07-22: `if question` は生文字列を真偽判定するため、空白のみの質問
    # （例 "   " / "\n\n"）は truthy のまま素通りし、strip()後は空文字列に
    # なって splitlines() が [] を返し、[0] で IndexError が発生していた。
    # 一度リストにしてからガード付きで先頭要素を取り出す。
    _q_lines = (question or "").strip().splitlines()
    q = _q_lines[0] if _q_lines else ""
    if 0 < len(q) <= 40:
        return q
    if slides and slides[0].get("title"):
        return slides[0]["title"]
    return "プレゼンテーション"


_PPTX_IMG_SCHEMA = {
    "type": "object",
    "properties": {"images": {"type": "array", "items": {
        "type": "object",
        "properties": {"index": {"type": "integer"}, "prompt": {"type": "string"}},
        "required": ["index", "prompt"]}}},
    "required": ["images"],
}


def plan_pptx_images(title, slides):
    """タイトル+各スライドの見出しから、画像が効果的なスライドを選び英語SDプロンプトを割り当てる。
    戻り値 {index: prompt}（index 0=タイトルのヒーロー画像 / 1..=各スライド）。最大 PPTX_MAX_IMAGES。"""
    outline = f"Title (index 0): {title}\n" + "\n".join(
        f"Slide {i + 1}: {s['title']} — {'; '.join(s['bullets'][:3])}"
        for i, s in enumerate(slides))
    sys = (
        "You plan illustrative images for a slide deck. Include index 0 (a title hero image) AND "
        f"as many conceptual slides as add value, aiming for {PPTX_MAX_IMAGES} images total when the "
        "content allows. Only SKIP a slide if it is purely a numeric table, code, or a bare list of "
        "figures. For each chosen slide give a vivid English Stable Diffusion prompt with quality tags "
        "that visually represents that slide's topic. "
        'Output ONLY JSON: {"images":[{"index":int,"prompt":str}]}. No prose, no thinking.'
    )
    raw = ask(CONDUCTOR, [{"role": "system", "content": sys},
                          {"role": "user", "content": outline}],
              CONDUCTOR_TEMP, think=False, fmt=_PPTX_IMG_SCHEMA,
              num_predict=768, label="pptx-img-plan")
    j = extract_json(raw) or {}
    out = {}
    # 2026-07-24: 従来の `j.get("images") or []` は "images" が falsy(None/{}/[]/""/0
    # 等)の場合のみ [] に丸めるトリックで、int/float/bool のような真値だが
    # 非反復可能(non-iterable)な値はそのまま通してしまっていた。その場合は
    # 直後の `for it in ...` 自体が TypeError を送出し、内側の try/except
    # (int(it.get("index")) 用)より手前=for文そのもので落ちるため捕捉されない。
    # plan_pptx_images は build_pptx (L4162付近) から XML安全化 try/except
    # (iteration68) の外で呼ばれ、build_pptx 自体も ask_fugu (L3545付近) から
    # 無防備に呼ばれているため、この TypeError は計算済み(数学/MCQでは
    # solve_verifiable の自己整合性投票済み)の MoA 回答ごとターン全体を落とす
    # (iteration41-47/68/80と同種の「高コスト回答喪失」障害)。イテレーション110は
    # このケースをテストコメント(旧case 8)で発見済みだったが、当時はテストのみの
    # 変更に限定されており修正は見送られていた。イテレーション103の _ddg_instant
    # における非list RelatedTopics 補正(isinstance判定で丸める方式)と同じやり方
    # で、"images" が非list(文字列/None/int/float/bool 等)であれば真偽・型に
    # 関わらず必ず [] に倒す(`or []` の truthiness トリックには戻さない)。
    imgs = j.get("images")
    if not isinstance(imgs, list):
        imgs = []
    for it in imgs:
        if not isinstance(it, dict):
            continue
        try:
            idx = int(it.get("index"))
        except Exception:
            continue
        p = str(it.get("prompt") or "").strip()
        if p and idx not in out and 0 <= idx <= len(slides):
            out[idx] = p
        if len(out) >= PPTX_MAX_IMAGES:
            break
    return out


def build_pptx(question, answer, out_path=None):
    """MoA 回答をスライド化し、内容連動で画像を埋め込んだ .pptx を生成して Path を返す。
    python-pptx 不在時は .md にフォールバック。画像バックエンド不在時はテキストのみで生成。"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        md = (Path(out_path).with_suffix(".md") if out_path
              else PPTX_OUT_DIR / f"fugu_{time.strftime('%Y%m%d_%H%M%S')}.md")
        md.parent.mkdir(parents=True, exist_ok=True)
        _save_as_markdown(md, question, answer, 0.0, "")
        print(f"   [PowerPoint には python-pptx が必要。代わりに保存: {md}]")
        return md

    if out_path is None:
        PPTX_OUT_DIR.mkdir(parents=True, exist_ok=True)
        out_path = PPTX_OUT_DIR / f"fugu_{time.strftime('%Y%m%d_%H%M%S')}.pptx"
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    raw_slides = _parse_slides(answer) or [{"title": "概要", "bullets": [answer[:400]]}]
    # 箇条書きを PPTX_MAX_BULLETS 単位に分割
    slides = []
    for s in raw_slides:
        bullets = s["bullets"] or [""]
        for i in range(0, len(bullets), PPTX_MAX_BULLETS):
            slides.append({"title": s["title"] + ("（続き）" if i else ""),
                           "bullets": bullets[i:i + PPTX_MAX_BULLETS]})
    slides = slides[:PPTX_MAX_SLIDES]
    title = _deck_title(question, raw_slides)

    # 画像計画 → 生成
    imgs = {}
    if _detect_backend() is not None and IMAGE_BACKEND != "off":
        plan = plan_pptx_images(title, slides)
        had_zero = 0 in plan
        plan.setdefault(0, None)  # タイトルには必ずヒーロー画像
        items = list(plan.items())
        if not had_zero and len(items) > PPTX_MAX_IMAGES:
            # 2026-07-23: dict は挿入順を保持するため、plan_pptx_images が
            # index 0 を含めずに（LLM が上の「タイトルには必ずヒーロー画像」
            # 指示を無視して）ちょうど PPTX_MAX_IMAGES 件返した場合、直上の
            # setdefault(0, None) は 0 を末尾に追加するだけになる。従来は
            # 直後の [:PPTX_MAX_IMAGES] スライスがその末尾の 0 を切り捨てて
            # おり、タイトルスライドにヒーロー画像が入らないまま不変条件が
            # 静かに破られていた。0 を先頭に固定してから残りを詰め直すことで
            # 画像総数は PPTX_MAX_IMAGES のまま（枠を1つタイトルへ再割当）0
            # の生成を保証する。0 が既に予算内にある場合／plan が定員未満の
            # 場合はこの分岐に入らず、従来と完全に同じ順序のまま処理する。
            items = [(0, plan[0])] + [kv for kv in items if kv[0] != 0]
        print(f"   [PPTX画像: {min(len(plan), PPTX_MAX_IMAGES)} 枚を生成します...]")
        for idx, pr in items[:PPTX_MAX_IMAGES]:
            if pr:
                path = generate_image(pr, "")
            else:
                base = title if idx == 0 else slides[idx - 1]["title"]
                p2, n2 = author_image_prompt(base)
                path = generate_image(p2, n2)
            if path:
                imgs[idx] = str(path)

    # 2026-07-23: デッキ構築〜保存の間、LLM 回答由来の制御文字
    # (NUL 0x00, ESC 0x1B 等) が add_textbox 経由で run.text に渡ると
    # python-pptx/lxml が ValueError ('All strings must be XML compatible:
    # no NULL bytes or control characters') を送出する。従来は
    # except ImportError しか捕捉しておらず、この ValueError がそのまま
    # _save_answer_to_file、さらには ask_fugu の make_pptx 経路まで伝播し、
    # 計算済み（math/mcq では SC投票済みの）回答ごと失われていた。
    # iteration 41 (_save_as_excel の IllegalCharacterError)・43
    # (_save_as_docx の ValueError)・44 (_save_as_pdf の
    # FPDFUnicodeEncodingException) と同じバグクラスの修正として、
    # add_textbox の choke point（タイトルスライド・コンテンツスライドの
    # 見出し・箇条書き全てがここを通る）で XML 不正制御文字のみを除去し
    # （実データは維持）、かつデッキ構築〜保存全体を broad except で囲んで
    # 失敗時は既存の .md フォールバックへ安全に降格させる。画像生成
    # (plan_pptx_images/generate_image/add_image) は既に自前で失敗を
    # 許容しており（backend 不在時は None を返す／embed エラーは
    # add_image が catch する）、このガードはデッキ構築・保存のみを対象と
    # し、画像生成の挙動（backend 不在時も画像無しでデッキが構築される）
    # は変更しない。
    _illegal_xml_re = re.compile(r"[\x00-\x08\x0B\x0C\x0E-\x1F]")
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)   # 16:9
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        def add_textbox(slide, text, left, top, width, height, size, bold=False):
            tb = slide.shapes.add_textbox(left, top, width, height)
            tf = tb.text_frame
            tf.word_wrap = True
            first = True
            for line in (text if isinstance(text, list) else [text]):
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                run = p.add_run()
                raw = ("• " + line) if isinstance(text, list) else line
                run.text = _illegal_xml_re.sub("", raw)
                run.font.size = Pt(size)
                run.font.bold = bold
            return tb

        def add_image(slide, path, left, top, width):
            try:
                slide.shapes.add_picture(path, left, top, width=width)
            except Exception as e:
                print(f"   [PPTX画像埋込エラー: {e}]")

        # タイトルスライド
        s0 = prs.slides.add_slide(blank)
        if 0 in imgs:
            add_image(s0, imgs[0], Inches(4.17), Inches(2.5), Inches(5.0))
            add_textbox(s0, title, Inches(0.7), Inches(0.6), Inches(12.0), Inches(1.3), 40, True)
            add_textbox(s0, "Fugu MoA 生成", Inches(0.7), Inches(1.9), Inches(12.0), Inches(0.6), 18)
        else:
            add_textbox(s0, title, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.6), 44, True)
            add_textbox(s0, "Fugu MoA 生成", Inches(0.9), Inches(4.2), Inches(11.5), Inches(0.7), 20)

        # コンテンツスライド
        for i, s in enumerate(slides, start=1):
            sl = prs.slides.add_slide(blank)
            add_textbox(sl, s["title"], Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.0), 30, True)
            has_img = i in imgs
            body_w = Inches(7.0) if has_img else Inches(12.1)
            add_textbox(sl, s["bullets"], Inches(0.6), Inches(1.6), body_w, Inches(5.4), 18)
            if has_img:
                add_image(sl, imgs[i], Inches(7.9), Inches(1.7), Inches(4.9))

        prs.save(str(out_path))
        return out_path
    except Exception as e:
        md = out_path.with_suffix(".md")
        _save_as_markdown(md, question, answer, 0.0, "")
        print(f"   [PPTX保存に失敗しました ({e!r})。代わりに保存: {md}]")
        return md


def _save_answer_to_file(question: str, answer: str, elapsed: float,
                         path: str, context: str = ""):
    """回答を --out で指定した拡張子に合わせた形式で保存する。
    .py/.js 等のコード拡張子 → コード抽出して書き出し
    .md/.txt → テキスト形式追記
    .pdf     → fpdf2 で生成（未インストール時 .md にフォールバック）
    .docx    → python-docx で生成（未インストール時 .md にフォールバック）
    .xlsx    → openpyxl で生成（未インストール時 .csv にフォールバック）
    .pptx    → python-pptx で画像入りスライド生成（未インストール時 .md にフォールバック）
    .html    → HTML で生成
    その他   → Markdown で保存
    """
    out = Path(path)
    # 2026-07-23: --out に存在しないサブディレクトリ（例: reports/answer.md）
    # を指定すると、.md/.txt/.py 等は out.write_text が、.docx/.xlsx 等は
    # 各 lib.save が FileNotFoundError を送出し、ask_fugu 側は無捕捉のため
    # MoA/SC投票まで完了した高コストな計算済み回答がトレースバックと共に
    # 丸ごと失われていた。さらに office 系フォールバック（.docx/.xlsx 失敗時
    # の .md/.csv 代替書き込み）も同じ存在しないディレクトリに書こうとして
    # 二重に失敗し、出力が完全に消える。build_pptx は既に自前で
    # out_path.parent.mkdir を呼んでいる（L3729、フォールバック分岐にも
    # 別途 mkdir あり）が、他の拡張子には無かった。iteration 41-47・68
    # （_save_as_excel の IllegalCharacterError、_save_as_docx の
    # ValueError、_save_as_pdf の FPDFUnicodeEncodingException 等、保存段で
    # 計算済み回答を失わないための一連の修正）と同じバグクラスであるため、
    # 全 suffix 分岐・その場フォールバック双方をこの一箇所で保護する。
    out.parent.mkdir(parents=True, exist_ok=True)
    suffix = out.suffix.lower()
    actual = out  # 実際に書かれたファイル（フォールバック時に変わる可能性あり）

    if suffix in _CODE_EXTENSIONS:
        _save_as_code(out, answer)
    elif suffix == ".txt":
        _save_as_text(out, question, answer, elapsed)
    elif suffix in {".pdf"}:
        result = _save_as_pdf(out, question, answer, elapsed)
        if result:
            actual = result
    elif suffix in {".docx", ".doc"}:
        result = _save_as_docx(out, question, answer, elapsed)
        if result:
            actual = result
    elif suffix in {".xlsx", ".xls"}:
        result = _save_as_excel(out, answer)
        if result:
            actual = result
    elif suffix in {".pptx", ".ppt"}:
        result = build_pptx(question, answer, out)
        if result:
            actual = result
    elif suffix in {".html", ".htm"}:
        _save_as_html(out, question, answer, elapsed)
    else:
        # .md またはその他 → Markdown
        _save_as_markdown(out, question, answer, elapsed, context)

    print(f"   [回答を保存しました: {actual}]")


def repl(use_search=False, rag_dirs=None, history_file=None):
    global _HISTORY
    hfile = history_file or HISTORY_FILE
    flags = []
    if use_search:
        flags.append("Web検索ON")
    if rag_dirs or RAG_DIRS:
        dirs = rag_dirs or RAG_DIRS
        flags.append(f"RAG:{','.join(str(d) for d in dirs)}")
    if flags:
        print(f"   [{', '.join(flags)}]")
    print("コマンド: 'exit'/'quit' で終了  'reset' で会話履歴クリア  "
          "'search on/off' で Web検索切替  'save <path>' で履歴エクスポート")
    while True:
        try:
            q = input("\nUser> ").strip()
        except (EOFError, KeyboardInterrupt):
            print("\n終了します。")
            break
        if not q:
            continue
        low = q.lower()
        if low in ("exit", "quit"):
            break
        if low == "reset":
            _HISTORY.clear()
            save_history_file(_HISTORY, path=hfile)
            print("   [会話履歴をクリアしました]")
            continue
        if low == "search on":
            use_search = True
            print("   [Web検索: ON]")
            continue
        if low == "search off":
            use_search = False
            print("   [Web検索: OFF]")
            continue
        # 2026-07-27: iteration 204で確認・特性化された(a)/(b)の修正
        # (iteration 182/186の'save <path>'契約には一切触れない)。
        # (a) 末尾スペース無しの素の'save'(4文字)はlow.startswith("save ")
        #     (5文字、末尾スペース必須)に一致せず、従来はどのコマンドにも
        #     一致しないままループ末尾のask_fugu(q, ...)へ落ちていた。つまり
        #     'save'という文字列そのものがフルのConductor+MoA/SCパイプラインへの
        #     質問として実行され、その(意味のない)応答が(user:'save',
        #     assistant:<応答>)として_HISTORYに追加され、以降の全ターンの
        #     Conductor/proposerコンテキストを汚染してしまう
        #     (iteration 59/67/70/73と同じ複数ターン文脈汚染バグ種)。
        # (b) 'save '(末尾スペースのみ、パス指定なしのつもり)も、repl()冒頭の
        #     q = input(...).strip()で入力全体が丸ごとstripされるため、判定前に
        #     'save'(4文字)へ潰れて結局(a)と同じ経路に収束していた。この結果、
        #     直後のL5421-5423相当「保存先パスを指定してください」ガイダンスは
        #     qが既にstrip済み(=末尾に空白を残せない)という不変条件と両立せず、
        #     実際のinput()経由では到達不能なデッドコードだった。
        # 修正: low(既存の小文字化済み変数、'SAVE'/'Save'等も統一的に扱う)が
        # 'save'単体に完全一致する場合を明示コマンド分岐として捕捉し、使用方法の
        # 案内だけ出してcontinueする。ask_fuguへは絶対にディスパッチせず、
        # _HISTORYも一切変更しない。これにより(a)/(b)双方が同じ分岐で解消される。
        if low == "save":
            print("   [保存先パスを指定してください: save <path>]")
            continue
        if low.startswith("save "):
            save_path = q[5:].strip()
            if not save_path:
                # 上のlow == "save"分岐により、qは既にstrip済みという不変条件から
                # ここへは実際のinput()経由では到達し得ない(iteration 204で確認済み
                # の到達不能デッドコード)。安全側の防御としてそのまま残す。
                print("   [保存先パスを指定してください: save <path>]")
                continue
            ok = save_history_file(_HISTORY, path=Path(save_path), force=True)
            if ok:
                print(f"   [履歴を保存しました: {save_path}]")
            else:
                print(f"   [履歴の保存に失敗しました: {save_path}]")
            continue
        ask_fugu(q, use_search=use_search, rag_dirs=rag_dirs,
                 history_file=hfile)


def main():
    parser = argparse.ArgumentParser(
        description="Local Fugu-style MoA オーケストレーター",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""使用例:
  python fugu_local.py                               # 対話モード
  python fugu_local.py "91は素数ですか？"             # 一問一答
  python fugu_local.py --file task.txt               # .txt から質問を読む
  python fugu_local.py --file spec.pdf               # PDF から読む（pdfplumber 要）
  python fugu_local.py --file task.py   --out fix.py  # Python コードを修正して .py で保存
  python fugu_local.py --file report.md --out out.md  # Markdown → Markdown
  python fugu_local.py --file spec.docx --out result.pdf  # Word → PDF（fpdf2 要）
  python fugu_local.py --search "最新のS&P500は？"    # Web検索あり
  python fugu_local.py --rag ./docs "PINNを実装して"  # RAGあり
  python fugu_local.py --no-history "一時的な質問"    # 履歴を使わない
  python fugu_local.py --session ./project.json       # プロジェクト専用履歴

対応ライブラリ（pip install で追加）:
  PDF読込: pdfplumber  PDF書出: fpdf2
  Word:    python-docx  Excel: openpyxl  PowerPoint: python-pptx""",
    )
    parser.add_argument("question", nargs="?",
                        help="質問文（省略時は対話モード）")
    parser.add_argument("--file", "-f", metavar="PATH",
                        help="質問をテキストファイルから読む")
    parser.add_argument("--out", "-o", metavar="PATH",
                        help="回答を保存するファイル（拡張子で形式自動選択: "
                             ".md/.txt/.py/.js/.pdf/.docx/.xlsx/.html 等）")
    parser.add_argument("--search", "-s", action="store_true",
                        help="Web 検索を有効化してコンテキストに注入する")
    parser.add_argument("--rag", "-r", nargs="+", metavar="DIR",
                        help="RAG 用ドキュメントディレクトリ（複数指定可）")
    parser.add_argument("--no-history", action="store_true",
                        help="セッション永続化を無効化（履歴を読まず保存もしない）")
    parser.add_argument("--session", metavar="PATH",
                        help=f"会話履歴ファイルパス（既定: {HISTORY_FILE}）")
    args = parser.parse_args()

    # --- セッション設定 ---
    global SESSION_SAVE, _HISTORY
    hfile = Path(args.session) if args.session else HISTORY_FILE
    if args.no_history:
        SESSION_SAVE = False
    else:
        _HISTORY = load_history_file(hfile)
        if _HISTORY:
            print(f"[session] 会話履歴を読み込みました: {len(_HISTORY) // 2} 往復 ({hfile})")

    # --- RAG ディレクトリ設定 ---
    rag_dirs = args.rag or (RAG_DIRS if RAG_DIRS else None)

    if not setup():
        return

    # --- 質問の取得 ---
    _OFFICE_SUFFIXES = {".docx", ".doc", ".xlsx", ".xls", ".pdf", ".pptx", ".ppt"}
    question = None
    office_attached = False
    if args.file:
        fp = Path(args.file)
        if not fp.exists():
            print(f"ファイルが見つかりません: {args.file}")
            return
        question = read_file_text(fp).strip()
        if not question:
            print(f"ファイルからテキストを抽出できませんでした: {args.file}")
            return
        office_attached = fp.suffix.lower() in _OFFICE_SUFFIXES
        print(f"[file] {fp.name} ({fp.suffix}) から {len(question)} 文字を読み込みました"
              + ("  [Office→Proposer C 主軸]" if office_attached else ""))
        # 2026-07-26 (iteration 192): gotcha #2 ―― num_ctx は全モデルで明示pinが必須
        # (非thinkingモデルはMODEL_NUM_CTX=8192、thinkingモデルでもMODEL_CONFIG側の
        # 上限は16384。32768はこの8GB GPUでは遅すぎるため常用しない)。iteration 185で
        # 表面化した特性(c)(「--fileのテキストはサイズ上限が一切無く、抽出結果が
        # どれだけ長くてもnum_ctxを意識したトランケートをせずそのままask_fugu()へ
        # 渡す」)はそのまま残っている。抽出テキストがnum_ctxを超えると、Ollama側で
        # 入力が黙って切り詰められ、本文末尾に置かれがちな指示・質問そのものが
        # 失われて精度が劣化しうるのに、従来はユーザーへの合図が一切無かった。
        # ここではトランケート・チャンク分割・サンプリングは一切行わない
        # (精度優先・時間は気にしない。gotcha #7と同じ思想で、正確さを文脈窓に
        # 収める都合で犠牲にしない)。抽出テキスト全文は従来通りask_fugu()へ渡した
        # うえで、オーバーフローがほぼ確実なときにのみ可視の警告を追加する。
        # 閾値はモジュールグローバルMODEL_NUM_CTX(既定8192。FUGU_HIGH_VRAM=1なら
        # 32768)をそのまま「文字数」の目安として使う保守的な見積もりで、
        # 1文字≒1トークンという最も密な想定（日本語の全角文字は実際にはこれより
        # トークン効率が悪いことが多く、本当はもっと早く溢れる）かつシステム
        # プロンプト・生成分の余地をゼロ扱いする、つまり「これを超えたら絶対に
        # 収まらない」と言い切れるラインだけを基準にする。そのため中規模ファイル
        # では誤報が出ず、本当に危険なケースだけを確実に警告できる。
        if len(question) > MODEL_NUM_CTX:
            print(f"[警告] {fp.name}: 抽出テキストが{len(question)}文字あり、"
                  f"num_ctx({MODEL_NUM_CTX}トークン)を超える可能性が高いため、"
                  f"モデルへの入力の一部(末尾の質問・指示を含む可能性があります)が"
                  f"Ollama側で黙って切り詰められる恐れがあります。"
                  f"本文はトランケートせずそのまま渡します。")
    elif args.question:
        question = args.question

    # --- 実行 ---
    if question:
        ask_fugu(question, use_search=args.search,
                 rag_dirs=rag_dirs, out_file=args.out, history_file=hfile,
                 office_attached=office_attached)
    elif sys.stdin.isatty():
        # 2026-07-26: iteration 185 で表面化した特性(a) ―― --out はこの対話分岐(repl())
        # には一切転送されず、repl()自身にもout_fileパラメータが無い設計のため、
        # `--out result.md` を質問なしで指定すると黙って無視されていた(エラー・警告なし)。
        # repl()にout_fileを追加してターン毎に自動保存する案は「どのターンを保存するか・
        # 上書きか追記か」が曖昧なため見送り(対話中の `save <path>` コマンドで既に手動
        # エクスポート可能)。surface-don't-swallow方針(gotcha #8, iters 66/71/110)に
        # 従い、ここでは無視される旨を可視化する警告のみ追加する
        # (repl()呼び出し自体・引数は不変のまま)。
        if args.out:
            print(f"[警告] --out {args.out} は対話モードでは無視されます。"
                  f"回答を保存するには対話中に `save {args.out}` のように"
                  f"saveコマンドを使ってください。")
        repl(use_search=args.search, rag_dirs=rag_dirs, history_file=hfile)
    else:
        # パイプ入力: stdin を質問として読む
        q = sys.stdin.read().strip()
        if q:
            ask_fugu(q, use_search=args.search,
                     rag_dirs=rag_dirs, out_file=args.out, history_file=hfile)
        else:
            print("質問が入力されませんでした。")
            parser.print_help()


if __name__ == "__main__":
    main()
